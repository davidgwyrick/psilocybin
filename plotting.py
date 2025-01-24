#Base
import json, os, time, sys
from os.path import join
import numpy as np
import pandas as pd
import io

#Plot
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
from ssm.plots import gradient_cmap
from matplotlib.gridspec import GridSpec
from PIL import Image

#Datashader
import datashader as ds
from datashader.mpl_ext import dsshow

#Colors
color_names=['amber','red','windows blue','faded green',
             'dusty purple',
             'orange',
             'steel blue',
             'pink',
             'greyish',
             'mint',
             'clay',
             'light cyan',
             'forest green',
             'pastel purple',
             'salmon',
             'dark brown',
             'lavender',
             'pale green',
             'dark red',
             'gold',
             'dark teal',
             'rust',
             'fuchsia',
             'pale orange',
             'cobalt blue',
             'mahogany',
             'cloudy blue',
             'dark pastel green',
             'dust',
             'electric lime',
             'fresh green',
             'light eggplant',
             'nasty green']
color_palette = sns.xkcd_palette(color_names)
cc = sns.xkcd_palette(color_names)
cmap = gradient_cmap(color_palette)

#Directories
PlotDir = '/home/david.wyrick/projects/zap-n-zip/plots'

#Define behavioral states
behavior_ranges = {0: [0,1], 1: [1,500]}#, 3:[30,500]}
behavior_dict = { 0: 'rest (<1cm/s)', 1: 'active (>1cm/s)'}
behavior_strs2 = ['rest','active']
behavior_strs = list(behavior_dict.values())
nBehaviors = len(behavior_strs)

#Define windows to calculate firing rate
stim_window_dict = {'spontaneous': [], 'evoked': [],'pre-rebound': [], 'rebound': [], 'post-rebound': [], 'visual': []}
stim_strs = ['spontaneous','evoked','pre-rebound','rebound','post-rebound','visual']
nStimTypes = len(stim_strs)
evoked_windows = [[.002,.025],[0.025,0.075],[.075,.3],[.3,1]]
evoked_strings = ['evoked','pre-rebound','rebound','post-rebound']

##------------------------------------------
def using_datashader(ax, x, y,vmax=2000):

    df = pd.DataFrame(dict(x=x, y=y))
    dsartist = dsshow(
        df,
        ds.Point("x", "y"),
        ds.count(),
        vmin=1,
        vmax=vmax,
        norm="log",
        aspect="auto",
        ax=ax,
        cmap='inferno'
    )

    plt.colorbar(dsartist)

##------------------------------------------
def adjust_spines(ax, spines=['left','bottom'],outward=5):
    for loc, spine in ax.spines.items():
        if loc in spines:
            spine.set_position(('outward', outward))  # outward by 10 points
            # spine.set_smart_bounds(True)
        else:
            spine.set_color('none')  # don't draw spine
        # if loc == 'left':
        #     spine.set_bounds(*ax.get_ylim())
        # if loc == 'bottom':
        #     spine.set_bounds(*ax.get_xlim())

    # turn off ticks where there is no spine
    if 'left' in spines:
        ax.yaxis.set_ticks_position('left')
    else:
        # no yaxis ticks
        ax.yaxis.set_ticks([])

    if 'bottom' in spines:
        ax.xaxis.set_ticks_position('bottom')
    else:
        # no xaxis ticks
        ax.xaxis.set_ticks([])

##------------------------------------------
def save_fig_to_pptx(fig,prs,slide=None,position=None):

    slide_W = 10; slide_H = 7.5

    # Save the Matplotlib figure to an in-memory buffer with tight bounding box
    image_buffer = io.BytesIO()
    plt.savefig(image_buffer, format='png',facecolor='white',dpi=300,bbox_inches='tight')#,pad_inches=0.1)
    image_buffer.seek(0)
    
    #Get figure size
    image = Image.open(image_buffer)
    width, height = image.size
    fig_W_inches = width / 300
    fig_H_inches = height / 300
    
    margin_pad = 0.1
    #Creat new slide if none given
    if slide is None:
        slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Resize the figure if it's larger than the slide dimensions
    if fig_W_inches > slide_W or fig_H_inches > slide_H:
        # Calculate the scaling factor to fit within the slide
        scale_x = slide_W / fig_W_inches
        scale_y = slide_H / fig_H_inches

        # Use the smaller scaling factor to maintain the aspect ratio
        scale = min(scale_x, scale_y)

        # Set the new figure size
        fig.set_size_inches(fig_W_inches * scale, fig_H_inches * scale)
        fig_W_inches = fig.get_figwidth()
        fig_H_inches = fig.get_figheight()

    # Calculate the position and size to center the figure on the slide
    if position is None:
        left = (slide_W - fig_W_inches) / 2  # Center horizontally
        top = (slide_H - fig_H_inches) / 2  # Center vertically
    else:
        left = position[0]; top = position[1]

    #Convert to float point length (?)
    left = left * Inches(1) 
    top = top * Inches(1) 
    width = fig_W_inches * Inches(1)
    height = fig_H_inches * Inches(1)

    # Add the figure from the in-memory buffer to the slide
    pic = slide.shapes.add_picture(image_buffer, left, top, width, height)
    
    # Close the in-memory buffer
    image_buffer.close()
    image.close()

##------------------------------------------
def plot_CCM_results(corr_mat,fcf,flow,takens,vmax_fcf=None,dir_clims=None,vmax_takens=None,title=None,fname_suffix=None,save=False,ticks=None,labels=None,boundaries=None,scatter_compare=False):

    indy = np.where(np.isnan(corr_mat)); corr_mat[indy] = 0
    indy = np.where(np.isnan(fcf)); fcf[indy] = 0
    indy = np.where(np.isnan(flow)); flow[indy] = 0
    indy = np.where(np.isnan(takens)); takens[indy] = 0

    fig, axes = plt.subplots(1,5,figsize=(20,4),gridspec_kw={'wspace': 0.4,'hspace':0.4})
    plt.suptitle(title,fontsize=16,fontweight='bold',y=0.95)

    if vmax_fcf is None:
        vmax1 = np.round(np.nanpercentile(np.abs(corr_mat.ravel()),99),2)
        vmax2 = np.round(np.nanpercentile(np.abs(fcf.ravel()),99),2)
        vmax_fcf = np.max([vmax1,vmax2])
        
    ax = axes[0]
    ax.set_title('|Correlation|',fontsize=12)#,fontweight='bold')
    sns.heatmap(np.abs(corr_mat),square=True,annot=False,cmap='viridis',vmin=0, vmax=vmax_fcf,ax=ax,cbar_kws={'shrink':0.5,'ticks':[0,vmax_fcf]},rasterized=True) #,'label': '|Correlation| 'cmap=sns.color_palette("vlag", as_cmap=True),center=0,ax=ax,cbar_kws={'shrink':0.5,'label': 'Correlation'})

    ax = axes[1]
    ax.set_title('Functional causal flow',fontsize=12)#,fontweight='bold')
    sns.heatmap(fcf,square=True,annot=False,cmap='viridis',vmin=0, vmax=vmax_fcf,ax=ax,cbar_kws={'shrink':0.5,'ticks':[0,vmax_fcf]},rasterized=True) #,'label': 'FCF'

    ax = axes[2]
    if scatter_compare:
        x = corr_mat.ravel()
        y = fcf.ravel()
        mask = (~np.isnan(x)) & (~np.isnan(y)) &(~np.isinf(x)) & (~np.isinf(y))
        using_datashader(ax, x[mask], y[mask])
        ax.plot([-0.5,1],[-0.5,1],'-k')
        ax.set_xlim([-0.5,1])
        ax.set_ylim([-0.5,1])
        ax.set_xticks(np.arange(-0.5,1.1,0.25)); ax.set_xticklabels(np.arange(-0.5,1.1,0.25))
        ax.set_yticks(np.arange(-0.5,1.1,0.25)); ax.set_yticklabels(np.arange(-0.5,1.1,0.25))
        ax.vlines(0,-0.5,1,color='r'); ax.hlines(0,-0.5,1,color='r')
        ax.set_xlabel('Correlation')
        ax.set_ylabel('Functional causal flow')
        
    else:
        vmax = np.round(np.nanpercentile(fcf - np.abs(corr_mat),98),2)
        vmin = np.round(np.nanpercentile(fcf - np.abs(corr_mat),2),2)
        ax.set_title('FCF - |Correlation|',fontsize=12)#,fontweight='bold')
        sns.heatmap(fcf - np.abs(corr_mat),square=True,annot=False,cmap='bwr',center=0, vmin=vmin,vmax=vmax,ax=ax,cbar_kws={'shrink':0.5,'ticks':[vmin,0,vmax]},rasterized=True) #,'label': 'FCF - |Corr|'

    ax = axes[3]
    ax.set_title('Directionality',fontsize=12)
    if dir_clims is None:
        vmax = np.round(np.nanpercentile(flow,97.5),2); vmin = np.round(np.nanpercentile(flow,2.5),2) #-1*vmax
    else:
        vmax = dir_clims[1]; vmin = dir_clims[0]
    # pdb.set_trace()
    sns.heatmap(flow,square=True,annot=False,cmap='RdBu_r',center=0,vmin=vmin,vmax=vmax,ax=ax,cbar_kws={'shrink':0.5,'ticks':[vmin,vmax]},rasterized=True) #,'label': 'Directionality'
    
    ax = axes[4]
    ax.set_title('Embedding\nDimensionality',fontsize=12)#,fontweight='bold')
    if vmax_takens is None:
        vmax_takens = np.round(np.nanpercentile(np.abs(takens.ravel()),98))
    sns.heatmap(takens,square=True,annot=False,cmap='rocket',ax=ax,vmin=0,vmax=vmax_takens,cbar_kws={'shrink':0.5,'ticks':[0,vmax_takens]},rasterized=True) #,'label': 'Takens'

    if ticks is None:
        for ax in axes:
            ax.set_xticks([]);ax.set_yticks([])
    else:
        for ii, ax in enumerate(axes):
            ax.set_xticks(np.array(ticks)+0.5);ax.set_yticks(np.array(ticks)+0.5)
            ax.set_xticklabels(labels,rotation=90,fontsize=8)
            ax.set_yticklabels(labels,fontsize=8,rotation=0)
            # if ii == 0:
            #     ax.set_yticklabels(labels,fontsize=8,rotation=0)
            # else:
            #     ax.set_yticklabels([])
        
    if boundaries is not None:
        for ii, ax in enumerate(axes):
            if (ii == 2) | (ii == 3):
                c = 'k'
            else:
                c = 'w'
            ax.vlines(boundaries,*ax.get_ylim(),color=c,lw=1,alpha=0.5)
            ax.hlines(boundaries,*ax.get_xlim(),color=c,lw=1,alpha=0.5)
    ax = axes[4]
    ax.autoscale(tight=True)
        
    if save:
        if fname_suffix is None:
            fname_suffix = 'mystery-plot'
        plt.savefig(os.path.join(PlotDir,f'FCF_groups_{fname_suffix}.pdf'))
    else:
        return fig

##------------------------------------------
def visualize_matrix(mat,ax=None,cmap='viridis',clims=None,center=None,cbar=True,cbar_ax=None,plot_ylabel=True,title=None,ticks=None,labels=None,boundaries=None,cbar_label=None):
    N = mat.shape[0]

    if ax is None:
        fig, ax = plt.subplots(figsize=(5,5))

    if clims is None:
        vmax = np.round(np.nanpercentile(mat,97.5),2); vmin = np.round(np.nanpercentile(mat,2.5),2) #-1*vmax
    else:
        vmin = clims[0]
        vmax = clims[1]
    if title is not None:
        ax.set_title(title,fontsize=12,fontweight='bold')
    sns.heatmap(mat,square=True,cmap=cmap,center=center,vmin=vmin,vmax=vmax,ax=ax,cbar=cbar,cbar_ax=cbar_ax,cbar_kws={'shrink':0.5,'label': cbar_label,'ticks':[vmin,0,vmax]},rasterized=True)
    
    #Plot labels
    if ticks is None:
        ax.set_xticks([]);ax.set_yticks([])
    else:
        ax.set_xticks(np.array(ticks)+0.5);ax.set_yticks(np.array(ticks)+0.5)
        ax.set_xticklabels(labels,rotation=90,fontsize=8)
        if plot_ylabel:
            ax.set_yticklabels(labels,rotation=0,fontsize=8)
        else:
            ax.set_yticklabels([])

    if boundaries is not None:
        if (cmap == 'bwr') | (cmap == 'RdBu_r')| (cmap == 'vlag') | (cmap == 'coolwarm'):
            c = 'k'
        else:
            c = 'w'
        # c = 'k'
        ax.vlines(boundaries,0,mat.shape[0],color=c,lw=1,ls='-',alpha=1)
        ax.hlines(boundaries,0,mat.shape[1],color=c,lw=1,ls='-',alpha=1)
    if ax is None:
        return fig


def plot_boxplot_comparison_spont(dFR_spont_df, i, j, spont_filename_list, areas_sub, groups_sub, g = None, pval_thresh = 0.05):

    epoch_i = spont_filename_list[i]
    epoch_j = spont_filename_list[j]
    if g is None:
        x = 'group'
        sub_df = dFR_spont_df.loc[(dFR_spont_df.epoch_i == epoch_i) & (dFR_spont_df.epoch_j == epoch_j)]
    else:
        x = 'area'
        sub_df = dFR_spont_df.loc[(dFR_spont_df.epoch_i == epoch_i) & (dFR_spont_df.epoch_j == epoch_j) & (dFR_spont_df.group == g)]

    #Determine unique boxes
    uniq_boxes = np.unique(sub_df[x])
    nBoxes = len(uniq_boxes)

    #Plot modulation index
    fig, axes = plt.subplots(2,1,figsize=(10,8))
    ax = axes[0]
    ax.set_title(f'MI: "{epoch_i}" vs "{epoch_j}"')
    gs = sns.boxplot(x=x,y='MI',hue='behavior',palette=sns.color_palette('Reds',5),order=uniq_boxes,hue_order=behavior_strs,data=sub_df,legend=False,ax=ax)
    ax.set_xlabel(x)
    xlim = ax.get_xlim(); ylim = ax.get_ylim()
    ax.hlines(0,-0.5,nBoxes-0.5,color='k',zorder=0)
    ax.set_xlim(xlim)
    ax.set_xlabel('')
    uniq_behaviors = np.unique(sub_df['behavior'])

    for ii, b in enumerate(uniq_boxes):
        for jj, w in enumerate(uniq_behaviors):
            if x == 'area':
                sub_sub_df = sub_df.loc[(sub_df.area == b) & (sub_df.behavior == w)]
            else:
                sub_sub_df = sub_df.loc[(sub_df.group == b) & (sub_df.behavior == w)]
            y = sub_sub_df['MI'].values
            if len(y) < 2:
                continue
            res = pg.wilcoxon(y)
            pval = res['p-val'][0]
            if pval < pval_thresh/len(uniq_behaviors):
                ax.text(ii-0.335+jj*0.4,0.925,'*',fontsize=20,fontweight='bold',color='k')

    #Plot delta firing rate
    # fig, ax = plt.subplots(figsize=(10,4))
    ax = axes[1]
    ax.set_title(f'\u0394-FR: "{epoch_i}" vs "{epoch_j}"')
    gs = sns.boxplot(x=x,y='dFR',hue='behavior',palette=sns.color_palette('Reds',5),order=uniq_boxes,hue_order=behavior_strs,data=sub_df,legend=False,ax=ax)
    ax.set_xlabel(x)

    tmp = sub_df['dFR'].values
    mask = (~np.isnan(tmp)) & (~np.isinf(tmp))
    ymin = np.round(np.nanpercentile(tmp[mask],2.5))
    ymax = np.round(np.nanpercentile(tmp[mask],97.5))
    if np.isnan(ymin) | np.isnan(ymax) | np.isinf(ymin) | np.isinf(ymax):
        print('inf or nan values in data')
    else:
        ax.set_ylim([ymin,ymax])
    xlim = ax.get_xlim(); ylim = ax.get_ylim()
    ax.hlines(0,-0.5,nBoxes-0.5,color='k',zorder=0)
    ax.set_xlim(xlim)
    ax.set_ylabel(f'\u0394-FR (Hz)')
    uniq_behaviors = np.unique(sub_df['behavior'])
    for ii, b in enumerate(uniq_boxes):
        if x == 'area':
            n = np.sum(areas_sub == b)
        else:
            n = np.sum(groups_sub == b)
        ax.text(ii-.15,ylim[0],n)
        for jj, w in enumerate(uniq_behaviors):
            if x == 'area':
                sub_sub_df = sub_df.loc[(sub_df.area == b) & (sub_df.behavior == w)]
            else:
                sub_sub_df = sub_df.loc[(sub_df.group == b) & (sub_df.behavior == w)]
            y = sub_sub_df['dFR'].values
            if len(y) < 2:
                continue
            res = pg.wilcoxon(y)
            pval = res['p-val'][0]
            if pval < pval_thresh/len(uniq_behaviors):
                ax.text(ii-0.335+jj*0.4,ylim[1]-6,'*',fontsize=20,fontweight='bold',color='k')

    return fig

##------------------------------------------
def plot_raster(spikes, ts, plot_tuple, time_bin = 0.1,cmap='gray_r',rereference_t=False,clims=None, title=None, run_tuple = None, pupil_tuple = None, eeg_tuple = None, stim_log = None, time_to_plot = 30, tStart = 0):
    
    if len(plot_tuple) == 11:
        #Ordered by area
        boundaries, ticks, labels, celltypes, durations, layers, areas, groups, mesogroup, supergroups, order_by_group = plot_tuple
        sort_by_area = True
    else:
        #Ordered by probe
        boundaries, ticks, labels = plot_tuple
        sort_by_area = False

    time_bin_ms = time_bin*1000
    T, N = spikes.shape
    if sort_by_area:
        spikes = spikes[:,order_by_group]

    if clims is None:
        vmax = np.nanpercentile(spikes,99)
        vmin = 0
    else:
        vmin = clims[0]
        vmax = clims[1]
    # Display a spike raster of the image
    tmp = [150, 800, 150]
    fig, axes =  plt.subplots(len(tmp),1,figsize=(12,12),gridspec_kw={'height_ratios':np.array(tmp)/np.sum(tmp)})
    plt.suptitle(title,y=0.925)

    #Define time window
    tslice = slice(int(tStart/time_bin),int((tStart+time_to_plot)/time_bin))
    ts_sub = ts[tslice]

    #Params
    xticks = np.arange(0,int((time_to_plot+1)/time_bin),int(15/time_bin))
    tlabel = xticks*time_bin
    xticks = np.arange(ts_sub[0],ts_sub[-1]+1,15)

    ax = axes[1]
    T, N = spikes[tslice].shape
    spk_plot = ax.imshow(spikes[tslice].T, aspect='auto',vmax=vmax,vmin=vmin, cmap=cmap)

    xticks = np.arange(0,T+1,int(T/3)); xticks[-1] -= 1
    ax.set_xticks(xticks)
    if rereference_t:
        t0 = ts_sub[0]
    else:
        t0 = 0
    ax.set_xticklabels(np.round(ts_sub[xticks]-t0))
    
    # ax.set_xticks([])
    ax.set_yticks(ticks,minor=True)
    ax.set_yticks(boundaries,minor=False)
    ax.set_yticklabels(labels,minor=True)
    ax.set_yticklabels([],minor=False)

    if sort_by_area:
        max_HZ = vmax/time_bin
        ax.set_title(f'Raster plot sorted by group')
        # ax.set_title(f'Raster plot sorted by group, time_bin = {time_bin_ms}ms, scale: [{vmin}, {vmax}] spikes')
        ax.set_ylabel("Group")
    else:
        ax.set_title(f'Raster plot sorted by probe, time_bin = {time_bin_ms}ms, scale: [{vmin}, {vmax}] spikes')
        ax.set_ylabel("Probe")

    #Plot pupil and running
    if run_tuple is not None:
        run_ts, run_signal, time_bin_r = run_tuple
        indy = np.where((run_ts >= ts_sub[0]) & (run_ts <= ts_sub[-1]))[0]
        ax = axes[-1]; ax.autoscale(tight=True)
        run_plot, = ax.plot(run_ts[indy],run_signal[indy],lw=0.6,color='k')
        ylim = ax.get_ylim()
        # xticks = np.arange(0,len(indy)+1,int(len(indy)/3)); xticks[-1] -= 1
        # ax.set_xticks(xticks)
        # ax.set_xticklabels(np.round(run_ts[indy][xticks]/60,2))

        tmp =np.arange(0,len(indy),int(len(indy)/3)); xticks = run_ts[indy[tmp]]
        ax.set_xticks(xticks)
        ax.set_xticklabels(np.round(xticks-t0,1))
        ax.autoscale(tight=True)
    else:
        run_plot = None
        ylim = [0,1]
    ax = axes[-1]
    ax.set_ylabel("Speed (cm/s)"); ax.set_xlabel("Time (seconds)")
    ax.set_title('Behavior')

    if pupil_tuple is not None:
        pupil_ts, pupil_radius, time_bin_p = pupil_tuple
        indy = np.where((pupil_ts >= ts_sub[0]) & (pupil_ts <= ts_sub[-1]))[0]
        ax2 = ax.twinx()
        pup_plot, = ax2.plot(pupil_ts[indy],pupil_radius[indy],lw=0.6,color=cc[1])
        ax2.set_ylabel("Pupil radius (pix)",color=cc[1])
        # xticks = np.arange(pupil_ts[indy[0]],pupil_ts[indy[-1]],int(len(indy)/3)); xticks[-1] -= 1
        # xticks = [pupil_ts[indy[0]],pupil_ts[indy[0]],pupil_ts[indy[-1]]
        tmp =np.arange(0,len(indy),int(len(indy)/3)); xticks = pupil_ts[indy[tmp]]
        ax2.set_xticks(xticks)
        ax2.set_xticklabels(np.round(xticks-t0,1))
        ax2.autoscale(tight=True)
        # pup_plot = None
    else:
        pup_plot = None

    #Plot EEG 
    if eeg_tuple is not None:
        eeg_ts, eeg_data, time_bin_e = eeg_tuple

        indy = np.where((eeg_ts >= ts_sub[0]) & (eeg_ts <= ts_sub[-1]))[0]
        ax = axes[0]; ax.autoscale(tight=True)
        # if plot_eeg_heatmap:
        eeg_plot = ax.imshow(eeg_data[indy].T,aspect='auto', vmin=-150,vmax=150,cmap='RdBu_r')
        ax.set_title('Re-referenced EEG, scale: [-150, 150] \u03BCV')

        xticks = np.arange(0,len(indy)+1,int(len(indy)/3)); xticks[-1] -= 1
        ax.set_xticks(xticks)
        ax.set_xticklabels(np.round(eeg_ts[indy][xticks]-t0,1))
    else:
        eeg_plot = None
        # eeg_plot = ax.plot(eeg_ts[indy],eeg_data[indy,10],lw=0.6,color='k',label='Channel 10')
        # ax.set_title('Re-referenced EEG, channel 10')
        # ax.set_ylabel("\u03BC V")
        # adjust_spines(ax)
        # ax.plot(eeg_ts[indy],eeg_data[indy,5],lw=0.6,color=cc[1],label='Channel 25')
        # ax.plot(eeg_ts[indy],eeg_data[indy,12],lw=0.6,color=cc[2],label='Channel 12')
        # ax.legend()
        # vmin = np.round(np.nanpercentile(eeg_data[indy,:],2))
        # vmax = np.round(np.nanpercentile(eeg_data[indy,:],98))
        # sns.heatmap(eeg_data[indy].T,ax=ax,cmap='RdBu_r',center=0,vmin=vmin,vmax=vmax,cbar=False)

    if stim_log is not None:
        stim_evoked_times = stim_log['onset'].values
        stim_types = stim_log['stim_type'].values
        indy = np.where((stim_evoked_times >= ts_sub[0]) & (stim_evoked_times <= ts_sub[-1]))[0]
        if len(indy) > 0:
            axes[-1].vlines(stim_evoked_times[indy],ylim[0],ylim[1],color=cc[1],lw=2,alpha=1)    
    
    return fig, (axes, eeg_plot,spk_plot, run_plot, pup_plot)

##------------------------------------------
## fig = usrplt.plot_behavior((run_ts,run_signal),(pupil_ts,pupil_radius),f'{mID} {rec_name}',(evoked_time_window_list,evoked_type_list),(injection_times,injection_types,injection_colors))
def plot_behavior(run_tuple,pupil_tuple=None,title_str=None,evoked_tuple=None,inj_tuple=None,iso_tuple=None,behavior_map=None,epoch_dict=None):

    time_bin_r = 0.01
    time_bin_p = 1/30
    #Expand tuples
    run_ts,run_signal = run_tuple
    if pupil_tuple is not None:
        pupil_ts,pupil_radius = pupil_tuple
        plot_pupil = True
    else:
        plot_pupil = False

    if inj_tuple is not None:
        injection_times, injection_types, injection_colors = inj_tuple
    else:
        injection_times = None

    if iso_tuple is not None:
        iso_induction_times, induction_colors = iso_tuple
    else:
        iso_induction_times = None

    if evoked_tuple is not None:
        evoked_time_window_list, evoked_type_list = evoked_tuple
    else:
        evoked_time_window_list = None

    if behavior_map is None:
        fig, ax = plt.subplots(figsize=(10,4))

        #Plot running speed
        ax.plot(run_ts/60,run_signal,'-k',lw=1)
        ax.set_xlabel('Time (min)')
        ax.set_ylabel('Running speed (cm/s)')

        if injection_times is not None:
            for ii, t in enumerate(injection_times):
                ax.axvline(t/60,color=injection_colors[ii],ls='--',lw=3,label=f'{injection_types[ii]} injection',zorder=4)
            ax.legend()

        if iso_induction_times is not None:
            for ii, t in enumerate(iso_induction_times):
                ax.axvline(t/60,color=induction_colors[ii],lw=3,ls='--',zorder=4)
            
        if plot_pupil:
            ax2 = ax.twinx()
            ax2.plot(pupil_ts/60,pupil_radius,'-b',lw=1)
            ax2.set_ylabel('Pupil radius (pixels)',color='b')
            plt.suptitle('Behavioral measures')
        else:
            plt.suptitle('Running speed')

        ylim = ax.get_ylim()
        if evoked_time_window_list is not None:
            for tW, stim_type in zip(evoked_time_window_list,evoked_type_list):
                if stim_type == 'biphasic':
                    c = cc[8]
                else:
                    c = cc[0]
                ax.fill_between(np.array(tW)/60,[ylim[0],ylim[0]],[ylim[1],ylim[1]],color=c,alpha=0.25,label=stim_type)
        if epoch_dict is not None:
            for i, (tW, e) in epoch_dict.items():
                ax.vlines(np.array(tW)/60,*ylim,color=cc[8],lw=2,ls='-',zorder=4)

    else:

        T = len(run_ts)
        tMin = T*time_bin_r/60
        fig, axes = plt.subplots(2,1,figsize=(12,5),gridspec_kw={'hspace':0.5})
        
        ax = axes[0]; ax.set_title('Running signal')
        ax2 = ax.twinx()
        sns.heatmap(behavior_map.reshape(1,-1),vmax=nBehaviors,vmin=0,ax=ax,cbar=False,cmap=sns.color_palette('Reds',nBehaviors),cbar_kws={'ticks':np.arange(nBehaviors)})

        ax2.plot(run_signal,color='k')
        ax2.set_ylabel('Speed (cm/s)')
        ax.set_yticks([])
        ax.set_xticks([0,T/2,T]); ax.set_xticklabels([0,np.round(tMin/2,2),np.round(tMin,2)],rotation=0)#;ax.set_xlabel('Time (mins)')

        ylim = ax2.get_ylim()
        if injection_times is not None:
            for i in range(2):
                ax2.vlines(injection_times[i]/time_bin_r,*ylim,color=injection_colors[i],lw=3,ls='--',zorder=4)
        if iso_induction_times is not None:
            for i in range(2):
                ax.vlines(iso_induction_times[i],*ylim,color=induction_colors[i],lw=3,ls='--',zorder=4)

        if plot_pupil:
            ax = axes[1]; ax.set_title('Pupil radius')
            ax.set_ylabel('Normalized pupil radius')    
            ax.plot(pupil_ts,pupil_radius,'-b')
            ax.autoscale(tight=True)
            T = pupil_ts[-1]
            ax.set_xticks([0,T/2,T]); ax.set_xticklabels([0,np.round(T/60/2,2),np.round(T/60,2)],rotation=0);ax.set_xlabel('Time (mins)')

            ylim = ax.get_ylim()
            if injection_times is not None:
                for i in range(2):
                    ax.vlines(injection_times[i],*ylim,color=injection_colors[i],lw=3,ls='--',zorder=4)
            if iso_induction_times is not None:
                for i in range(2):
                    ax.vlines(iso_induction_times[i],*ylim,color=induction_colors[i],lw=3,ls='--',zorder=4)

        ax = axes[1]

        ylim = ax.get_ylim()
        if evoked_time_window_list is not None:
            for tW, stim_type in zip(evoked_time_window_list,evoked_type_list):
                if stim_type == 'biphasic':
                    c = cc[8]
                else:
                    c = cc[0]
                ax.fill_between(np.array(tW),[ylim[0],ylim[0]],[ylim[1],ylim[1]],color=c,alpha=0.25,label=stim_type)
        if epoch_dict is not None:
            for i, (tW, e) in epoch_dict.items():
                ax.vlines(np.array(tW),*ylim,color=cc[8],lw=2,ls='-',zorder=4)


    if title_str is not None:
        plt.suptitle(title_str)
    return fig


    



# def plot_boxplot_comparison_spont2(dFR_df, epoch_i,epoch_j, g = None, pval_thresh = 0.05):
#     bar_width=0.8
#     if g is None:
#         x = 'group'
#         sub_df = dFR_df.loc[(dFR_df.epoch_i == epoch_i) & (dFR_df.epoch_j == epoch_j)]
#     else:
#         x = 'area'
#         sub_df = dFR_df.loc[(dFR_df.epoch_i == epoch_i) & (dFR_df.epoch_j == epoch_j) & (dFR_df.group == g)]

#     #Determine unique boxes
#     uniq_boxes = np.unique(sub_df[x])
#     nBoxes = len(uniq_boxes)

#     #Plot modulation index
#     fig, axes = plt.subplots(2,1,figsize=(10,8))

#     if x == 'area':
#         plt.suptitle(f'Group {g}: "{epoch_i}" vs "{epoch_j}"')
#     else:
#         plt.suptitle(f'All groups "{epoch_i}" vs "{epoch_j}"')

#     axes[0].set_title(f'Modulation index')
#     axes[1].set_title(f'\u0394-Firing rate:')
#     for ii, var in enumerate(['MI','dFR']):
#         ax = axes[ii]
    
#         gs = sns.boxplot(x=x,y=var,hue='behavior',palette=sns.color_palette('Reds',nBehaviors),order=uniq_boxes,hue_order=behavior_strs,data=sub_df,width=bar_width,legend=False,ax=ax)
#         ax.set_xlabel(x)
#         xlim = ax.get_xlim(); ylim = ax.get_ylim()
#         ax.hlines(0,-0.5,nBoxes-0.5,color='k',zorder=0)
#         ax.set_xlim(xlim)
#         ax.set_xlabel('')
#         usrplt.adjust_spines(ax)
#         for jj, b in enumerate(uniq_boxes):
#             n = len(sub_df.loc[(sub_df[x] == b)])
#             if ii == 0:
#                 ax.text(jj-.15,ylim[0],n)
#             for kk, w in enumerate(behavior_strs):
#                 sub_sub_df = sub_df.loc[(sub_df[x] == b) & (sub_df.behavior == w)]
#                 y = sub_sub_df[var].values
#                 if (len(y) < 2) | (np.all(y == 0)):
#                     continue
                
#                 if np.all(y == 0):
#                     continue
#                 res = pg.wilcoxon(y)
#                 pval = res['p-val'][0]
#                 x_pos = jj-0.45+kk*bar_width/len(behavior_strs)
#                 if pval < pval_thresh/len(behavior_strs):
#                     ax.text(x_pos,0.9*ylim[1],'*',fontsize=20,fontweight='bold',color='k')

#     if x == 'area':
#         axes[1].set_xlabel('Area')
#     else:
#         axes[1].set_xlabel('Group')

#     return fig

# def plot_boxplot_comparison_spont(dFR_df, epoch_i,epoch_j, g = None, behav_str = None, pval_thresh = 0.05):
#     bar_width=0.8
#     if g is None:
#         x = 'group'
#         sub_df = dFR_df.loc[(dFR_df.epoch_i == epoch_i) & (dFR_df.epoch_j == epoch_j) & (dFR_df.mwu_significant == 1)]
#     else:
#         x = 'area'
#         sub_df = dFR_df.loc[(dFR_df.epoch_i == epoch_i) & (dFR_df.epoch_j == epoch_j) & (dFR_df.group == g) & (dFR_df.mwu_significant == 1)]

#     #Determine unique boxes
#     uniq_boxes = np.unique(sub_df[x])
#     nBoxes = len(uniq_boxes)

#     #Plot modulation index
#     fig, axes = plt.subplots(2,1,figsize=(10,8))

    
#     if x == 'area':
#         plt.suptitle(f'Behavior: {behav_str}, Group {g}: "{epoch_i}" vs "{epoch_j}"')
#     else:
#         plt.suptitle(f'Behavior: {behav_str}, All groups "{epoch_i}" vs "{epoch_j}"')

#     axes[0].set_title(f'Modulation index')
#     axes[1].set_title(f'\u0394-Firing rate:')
#     for ii, var in enumerate(['abs_MI','dFR']):
#         ax = axes[ii]
    
#         gs = sns.barplot(x=x,y=var,hue='mod_group',palette=sns.xkcd_palette(['silver','green']),errorbar='se',order=uniq_boxes,hue_order=['neg_mod','pos_mod'],data=sub_df,width=bar_width,legend=False,ax=ax)
#         ax.set_xlabel(x)
#         xlim = ax.get_xlim(); ylim = ax.get_ylim()
#         ax.hlines(0,-0.5,nBoxes-0.5,color='k',zorder=0)
#         ax.set_xlim(xlim)
#         ax.set_xlabel('')
#         usrplt.adjust_spines(ax)
#         for jj, b in enumerate(uniq_boxes):
#             n1 = len(sub_df.loc[(sub_df[x] == b) & (sub_df.mod_group == 'neg_mod')])
#             n2 = len(sub_df.loc[(sub_df[x] == b) & (sub_df.mod_group == 'pos_mod')])
#             if ii == 0:
#                 ax.text(jj-.15,ylim[0],f'{n1}/{n2}',fontsize=10)

#             for kk, g in enumerate(['neg_mod','pos_mod']):
#                 sub_sub_df = sub_df.loc[(sub_df[x] == b) & (sub_df.mod_group == g)]
#                 y = sub_sub_df[var].values
#                 if (len(y) < 2) | (np.all(y == 0)):
#                     continue
                
#                 if np.all(y == 0):
#                     continue
#                 res = pg.wilcoxon(y)
#                 pval = res['p-val'][0]
#                 x_pos = jj-0.3+kk*bar_width/2
#                 if pval < pval_thresh/len(behavior_strs):
#                     ax.text(x_pos,0.9*ylim[1],'*',fontsize=12,fontweight='bold',color='k')

#     if x == 'area':
#         axes[1].set_xlabel('Area')
#     else:
#         axes[1].set_xlabel('Group')

#     return fig

# def plot_boxplot_comparison_evoked(dFR_df, epoch_i,epoch_j, b_str, g = None, pval_thresh = 0.05):
#     bar_width=0.8
#     if g is None:
#         x = 'group'
#         sub_df = dFR_df.loc[(dFR_df.epoch_i == epoch_i) & (dFR_df.epoch_j == epoch_j) & (dFR_df.behavior == b_str)]
#     else:
#         x = 'area'
#         sub_df = dFR_df.loc[(dFR_df.epoch_i == epoch_i) & (dFR_df.epoch_j == epoch_j) & (dFR_df.behavior == b_str) & (dFR_df.group == g)]

#     #Determine unique boxes
#     uniq_boxes = np.unique(sub_df[x])
#     nBoxes = len(uniq_boxes)

#     #Plot modulation index
#     fig, axes = plt.subplots(2,1,figsize=(10,8))
#     # plt.suptitle(f'Behavior: {b_str}\n"{epoch_i}" vs "{epoch_j}"')

#     axes[0].set_title(f'Modulation index')
#     axes[1].set_title(f'\u0394-Firing rate:')
#     if x == 'area':
#         plt.suptitle(f'Group: {g}, Behavior: {b_str}\n"{epoch_i}" vs "{epoch_j}"')
#     else:
#         plt.suptitle(f'All groups Behavior: {b_str}\n"{epoch_i}" vs "{epoch_j}"')
    
#     for ii, var in enumerate(['MI','dFR']):
#         ax = axes[ii]
    
#         gs = sns.boxplot(x=x,y=var,hue='window',palette=sns.color_palette('Purples',4),order=uniq_boxes,hue_order=window_strs,data=sub_df,width=bar_width,legend=False,ax=ax)

#         if var == 'dFR':
#             tmp = sub_df['dFR'].values
#             mask = (~np.isnan(tmp)) & (~np.isinf(tmp))
#             ymin = np.round(np.nanpercentile(tmp[mask],2.5))
#             ymax = np.round(np.nanpercentile(tmp[mask],97.5))
#             if np.isnan(ymax) | np.isinf(ymax):
#                 print('inf or nan values in data')
#             else:
#                 ax.set_ylim([ymin,ymax])
#         ax.set_xlabel(x)
#         xlim = ax.get_xlim(); ylim = ax.get_ylim()
#         ax.hlines(0,-0.5,nBoxes-0.5,color='k',zorder=0)
#         ax.set_xlim(xlim)
#         ax.set_xlabel('')
#         usrplt.adjust_spines(ax)

#         counter = 0
#         for jj, b in enumerate(uniq_boxes):
#             n = len(sub_df.loc[(sub_df[x] == b) & (sub_df.window == window_strs[0])])
#             if ii == 1:
#                 ax.text(jj-.15,ylim[0],n)
#             for kk, w in enumerate(window_strs):
#                 x_pos = jj-0.45+kk*bar_width/len(window_strs)
#                 # ax.text(x_pos,0.9*ylim[1],counter,fontsize=10,color='k')
#                 # counter += 1
#                 sub_sub_df = sub_df.loc[(sub_df[x] == b) & (sub_df.window == w)]
#                 y = sub_sub_df[var].values

#                 if (len(y) < 2) | (np.all(y == 0)):
#                     continue
#                 if np.all(y == 0):
#                     continue
#                 res = pg.wilcoxon(y)
#                 pval = res['p-val'][0]

#                 if pval < pval_thresh/len(window_strs):
#                     ax.text(x_pos,0.9*ylim[1],'*',fontsize=20,fontweight='bold',color='k')

#     if x == 'area':
#         axes[1].set_xlabel('Area')
#     else:
#         axes[1].set_xlabel('Group')
    
#     return fig
