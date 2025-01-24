base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
base_dir_local = '/data/tiny-blue-dot/zap-n-zip/EEG_exp'; base_dir = base_dir_server
ProjDir = '/home/david.wyrick/projects/zap-n-zip/'
ServDir = '/data/projects/zap-n-zip/'
BehavDir = '/home/david.wyrick/projects/zap-n-zip/plots/behavior'

#Basea
import gspread
from os.path import join
from glob import glob
import json, os, time, sys, argparse
import numpy as np
import pandas as pd
import scipy.stats as st
import scipy.signal as sig
import pingouin as pg
from scipy.ndimage import gaussian_filter
from scipy.interpolate import interp1d
from itertools import combinations
from delay_embedding import evaluation as E
from tqdm import tqdm

#Plot
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
from matplotlib.gridspec import GridSpec
from matplotlib.lines import Line2D   

#Project
from tbd_eeg.data_analysis.eegutils import EEGexp
import tbd_eeg.data_analysis.Utilities.utilities as tbd_util
from tbd_eeg.data_analysis.Utilities.behavior_movies import Movie

#Allen
from allensdk.core.mouse_connectivity_cache import MouseConnectivityCache
mcc = MouseConnectivityCache(resolution=10)
str_tree = mcc.get_structure_tree()

#Read in allen CCF
ccfsum = pd.read_csv('/home/david.wyrick/projects/zap-n-zip/data/mouse_connectivity/ontology_v2.csv')

#User
sys.path.append(ProjDir)
import util
import plotting as usrplt

#Single cell stuff
import elephant.statistics as est
from quantities import ms, s, Hz
from neo.core import SpikeTrain

#Define behavioral states
behavior_ranges = {0: [0,1], 1: [1,15], 2: [15,500]}#, 3:[30,500]}
behavior_dict = {0: 'rest (0-1cm/s)', 1: 'walk (1-15cm/s)', 2: 'run (>15cm/s)'}
behavior_strs2 = ['rest','walk','run']

behavior_strs = list(behavior_dict.values())
nBehaviors = len(behavior_strs)

#Define windows to calculate firing rate
windows_of_interest = [[.002,.025],[.075,.3],[.3,1],[1,4.5]]
window_strs = ['evoked (2-25ms)','rebound (75-300ms)','post-rebound (0.3-1s)','ISI (1-4.5s)']
nWindows = len(window_strs)

gc = gspread.service_account() # need a key file to access the account (step 2) 
sh = gc.open('Templeton-log_exp') # open the spreadsheet 

df = pd.DataFrame(sh.sheet1.get()) 
exp_table = df.T.set_index(0).T # put it in a nicely formatted dataframeexp_table.iloc[10:]
exp_table = exp_table.set_index('mouse_name')

##===== ============================ =====##
##===== Parse Command Line Arguments =====##
parser = argparse.ArgumentParser(description='Synchrony STTC')

##===== Data Options =====##
parser.add_argument('--mID',type=str, default='mouse669117',
                    help='mouse to perform analysis on')

parser.add_argument('--rec_name',type=str, default='pilot_aw_psi_2023-03-30_11-37-07',
                    help='experiment to perform analysis on')

parser.add_argument('--synchrony_window_ms',type=int, default=50,
                    help='synchrony_window_ms')

parser.add_argument('--time_bin_ms',type=int, default=100,
                    help='time_bin_ms')

parser.add_argument('--fr_thresh',type=float, default=1,
                    help='Firing rate threshold for neurons to include in analysis')

parser.add_argument('--window_t_min',type=int, default=10,
                    help='Window width (mins) to segment data into')

def calculate_STTC_DW(spk_train_i,spk_train_j,dt,T):
    
    def run_T(spk_train,dt,T):

        #Add extra boundary spikes
        spk_train = np.concatenate(([np.min(spk_train)-3*dt],spk_train))
        spk_train_buff = np.vstack([spk_train-dt,spk_train+dt]).T
        l_buff = spk_train_buff[1:,0]; r_buff = spk_train_buff[:-1,1]

        #Determine where spike buffers overlap
        mask = l_buff < r_buff
        indy_l = np.where(~mask)[0]+1
        indy_r = np.concatenate((indy_l[1:]-1,[len(mask)]))

        #Change first boundary spike back
        t_rolling = []
        for l, r in zip(indy_l,indy_r):
            if r - l == 1:
                t_rolling.append(spk_train_buff[l])
                t_rolling.append(spk_train_buff[r])
            else:
                t_rolling.append([spk_train_buff[l,0],spk_train_buff[r,1]])
        t_rolling = np.array(t_rolling)

        #return the proportion of total recording time which lies within +/- delta-t of any spike in A
        t_prop = np.sum(np.diff(t_rolling,axis=1))/T
        return t_prop, t_rolling
 
    #Calculate the fraction of the total recording time that is delta-t away from a spike in each spike train
    T_i, t_rolling_i = run_T(spk_train_i,dt,T)
    T_j, t_rolling_j = run_T(spk_train_j,dt,T)

    #Get the proportion of spikes from spk_train_i that lie within +/- delta-t of any spike in spk_train_j
    P_i = np.sum([np.any((s > t_rolling_j[:,0]) & (s < t_rolling_j[:,1])) for s in spk_train_i])/len(spk_train_i)

    #Get the proportion of spikes from spk_train_j that lie within +/- delta-t of any spike in spk_train_i
    P_j = np.sum([np.any((s > t_rolling_i[:,0]) & (s < t_rolling_i[:,1])) for s in spk_train_j])/len(spk_train_j)


    STTC = 0.5*((P_i-T_j)/(1-P_i*T_j) + (P_j-T_i)/(1-P_j*T_i))
    return STTC

#Yasmeen's code
def calculate_STTC_YN(spike_times_1, spike_times_2, dtv, Time):
 
    N1v = len(spike_times_1)
    N2v = len(spike_times_2)

    dt = dtv
    N1 = N1v
    N2 = N2v
 
    def run_P(N1, N2, dt, spike_times_1, spike_times_2):
 
        Nab = 0
        j = 0
        for i in range(0,N1):
            while j <= N2 - 1:
                if abs(spike_times_1[i] - spike_times_2[j]) <= dt:
                    Nab += 1
                    break
               
                # stops if the spike time is larger than that the initial one.
                # a way to compare all spikes in proximity of the spike 1.
                elif spike_times_2[j] > spike_times_1[i]:
                    break
                else:
                    j += 1
 
        return Nab
 
 
    def run_T(N1v, dtv, startv, endv, spike_times_1):
        dt = dtv
        start = startv
        endvv = endv
        tempN = N1v
 
        # print(tempN)
        # print(len(spike_times_2))
 
        time_A = 2 * tempN * dt
 
        # If there is just one spike, want to avid any miscalculation
        # that stem from over estimation of time covered in T.
        if tempN == 1:
 
            if spike_times_1[0] - start < dt:
                time_A = time_A - start + spike_times_1[0] - dt
            elif spike_times_1[0] + dt > endvv:
                time_A = time_A - spike_times_1[0] - dt + endvv

        else:
            i = 0                        
            while i < tempN - 1:
                diff = spike_times_1[i+1] - spike_times_1[i]
                if diff < 2 * dt:
                    time_A = time_A - 2 * dt + diff
                i += 1 
            if spike_times_1[0] - start < dt:
                time_A = time_A - start + spike_times_1[0] - dt

            if endvv - spike_times_1[tempN - 1] < dt:
                time_A = time_A - spike_times_1[tempN - 1] - dt + endvv
 
        return time_A
 
 
    if N1 == 0 or N2 == 0:
        index = None
 
    else:
        T = Time[1] - Time[0]
        TA = run_T(N1, dt, Time[0], Time[1], spike_times_1)
        TA = TA / T
        TB = run_T(N2, dt, Time[0], Time[1], spike_times_2)
        TB = TB / T
        PA = run_P(N1, N2, dt, spike_times_1, spike_times_2)
        PA = PA / N1
        PB = run_P(N2, N1, dt, spike_times_2, spike_times_1)
        PB = PB / N2
        index = (0.5 * ((PA - TB) / (1 - TB * PA))) + (0.5 * ((PB - TA) / (1 - TA * PB)))
   
    tileCoef = index
 
    return tileCoef

if __name__ == '__main__':

    # Parse the arguments
    args = parser.parse_args()
    mID = args.mID
    rec_name = args.rec_name
    fr_thresh = args.fr_thresh
    synchrony_window_ms = args.synchrony_window_ms
    synchrony_window = synchrony_window_ms/1000

    time_bin_ms = args.time_bin_ms
    time_bin = time_bin_ms/1000

    #Segement data into 15 minute windows
    window_t_min = args.window_t_min
    window_t = window_t_min*60
    
    #Extract row from templeton dataframe
    exp_df = exp_table.loc[exp_table.exp_name == rec_name]
    if type(exp_df) == pd.core.series.Series:
        exp_df = exp_df.to_frame().T

    drug_type = exp_df['drug'].values[0]
    stim_type = exp_df['stimulation'].values[0]
    print(f'Experiment type: {stim_type}, {drug_type}')

    #Define directories
    SaveDir = join(ServDir,'results','synchrony',mID,rec_name)
    if not os.path.exists(SaveDir):
        os.makedirs(SaveDir)

    PlotDir = join(SaveDir,'plots')
    if not os.path.exists(PlotDir):
        os.makedirs(PlotDir)

    #%% Upload the whole experiment and generate the global clock
    file_name = os.path.join(base_dir_server,mID,rec_name,'experiment1','recording1')
    exp = EEGexp(file_name, preprocess=False, make_stim_csv=False)

    #Load probe data
    probe_unit_data, probe_info, total_units = tbd_util.get_neuropixel_data(exp)
    probe_list = [x.replace('_sorted', '') for x in exp.experiment_data if 'probe' in x]

    #Get recording start time
    probei = probe_list[0]
    open_ephys_start = np.round(np.min(probe_unit_data[probei]['spike_times']))
    open_ephys_end = np.round(np.max(probe_unit_data[probei]['spike_times']))

    #For saline & psilocybin experiments, get injection times and types of injection
    if drug_type in ['saline', 'psilocybin']:
        injection_times = [float(exp_df['First injection time (s)'].values[0]),
                        float(exp_df['Second injection time (s)'].values[0])]
        
        #Determine injection type
        if 'psi' in rec_name:
            injection_types = ['sal1','psi']
        else:
            injection_types = ['sal1', 'sal2']
    else:
        injection_times = None

    #For isoflurane experiments, get iso level
    if drug_type == 'isoflurane':
        iso_level, iso_times = exp.load_analog_iso()
        iso_induction_times, iso_maintenance_times = exp.load_iso_times()

    # extract the timestamps of the selected stimuli
    try:
        stim_log = pd.read_csv(exp.stimulus_log_file)
        stim_exists = True
    except:
        stim_exists = False
        
    # Load running speed
    run_file = os.path.join(base_dir,mID,rec_name,'experiment1','recording1','running_signal.npy')
    if os.path.exists(run_file):
        # print('File exists')
        run_signal = np.load(run_file)

        ts_file = os.path.join(base_dir,mID,rec_name,'experiment1','recording1','running_timestamps.npy')
        if os.path.exists(ts_file):
            run_ts = np.load(ts_file)
        else:
            ts_file = os.path.join(base_dir,mID,rec_name,'experiment1','recording1','running_timestamps_master_clock.npy')
            run_ts = np.load(ts_file)

    else:
        print('\tRunning file does not exist')
        run_signal, run_ts = exp.load_running()
        np.save(os.path.join(base_dir,mID,rec_name,'experiment1','recording1','running_signal.npy'),run_signal)
        np.save(os.path.join(base_dir,mID,rec_name,'experiment1','recording1','running_timestamps.npy'),run_ts)

    #Apply savitzky-golay filter to running trace to make it more differentiable
    fw = 500; run_time_bin = 0.01
    run_signal = sig.savgol_filter(run_signal,int(fw/run_time_bin/1000),3)
    acceleration = np.gradient(run_signal)
    f_run = interp1d(run_ts,run_signal)
    f_acc = interp1d(run_ts,acceleration)
    
    #Create gaussian smoothed running signal to condition on
    fw = 750
    run_signal_s = gaussian_filter(run_signal,int(fw/run_time_bin/1000))
    f_run_s = interp1d(run_ts,run_signal_s)

    #Load pupil data
    pupil_csv = os.path.join(base_dir,mID,rec_name,'experiment1','recording1',f'Pupil_{rec_name}.csv')

    try:
        table = pd.read_csv(pupil_csv)
        pupil_radius = table['Largest_Radius'].values

        #Pupil master clock
        pupil_ts = Movie(filepath=exp.pupilmovie_file,
                        sync_filepath=exp.sync_file,
                        sync_channel='eyetracking'
                        ).sync_timestamps
        plot_pupil = True

        #Ensure timeseries are same length
        t = np.min([len(pupil_ts),len(pupil_radius)])
        pupil_ts = pupil_ts[:t]
        pupil_radius = pupil_radius[:t]

        pupil_radius_z = st.zscore(pupil_radius,nan_policy='omit')

        #Interpolate to equal time bins & remove outliers
        indy = np.where(~np.isnan(pupil_radius) & ~np.isinf(pupil_radius) & (np.abs(pupil_radius_z) < 3))[0]
        f_pupil = interp1d(pupil_ts[indy],pupil_radius[indy])
        pupil_time_bin = 1/30
        pupil_ts = np.arange(np.nanmin(pupil_ts),np.nanmax(pupil_ts),pupil_time_bin)
        pupil_radius = f_pupil(pupil_ts)

        #Apply savitzky-golay filter to pupil trace to make it more differentiable
        fw = 500
        pupil_radius_s = sig.savgol_filter(pupil_radius,int(fw/pupil_time_bin/1000),3)
        
        #Interpolate running signal to pupil_ts
        run_signal_p = f_run(pupil_ts)
        run_signal_p_s = f_run_s(pupil_ts)
        acceleration_p = f_acc(pupil_ts)

    except:
        print('\t No Pupil ?!')
        pupil_ts = np.array([np.nan])
        pupil_radius =np.array([np.nan])
        plot_pupil = False

    # Create a new PowerPoint presentation to save figures to
    prs = Presentation()

    #Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f'mID = {mID}; rec_name = {rec_name}'
    slide.placeholders[1].text = f'Synchrony measurements spontaneous blocks'

    if drug_type in ['saline', 'psilocybin']:
        tb = 180
        if injection_times[0] < 200:
            tb = 60
        macro_windows = [[open_ephys_start,injection_times[0]-tb]]; macro_name_list = ['Pre-injection']
        macro_windows.append([injection_times[0]+60,injection_times[1]-180]); macro_name_list.append(f'post_{injection_types[0]}_inj')
        macro_windows.append([injection_times[1]+60,open_ephys_end]); macro_name_list.append(f'post_{injection_types[1]}_inj')
        if drug_type == 'saline':
            cmap_macro = sns.color_palette('Greens',3)
        else:
            cmap_macro = np.concatenate((sns.color_palette('Blues',2),sns.color_palette('Reds',1)))
    elif drug_type == 'isoflurane':
        macro_windows = [[open_ephys_start,iso_induction_times[0]-120]]; macro_name_list = ['Pre-iso']
        macro_windows.append([iso_induction_times[0],iso_induction_times[1]]); macro_name_list.append(f'iso-ind')
        macro_windows.append([iso_induction_times[1]+60*5,open_ephys_end]); macro_name_list.append(f'post-iso')
        cmap_macro = sns.color_palette('Purples',3)

    elif drug_type == 'urethane':
        macro_windows = [[open_ephys_start,open_ephys_end/3]]; macro_name_list = ['urethane_1/3']
        macro_windows.append([open_ephys_end/3,open_ephys_end*2/3]); macro_name_list.append('urethane_2/3')
        macro_windows.append([open_ephys_end*2/3,open_ephys_end]); macro_name_list.append('urethane_3/3')
        cmap_macro = sns.color_palette('Oranges',3)

    if plot_pupil:
        if drug_type == 'urethane':
            fig, ax = plt.subplots(figsize=(6,6))
            for jj, tW in enumerate(macro_windows):
                indy = np.where((pupil_ts >= tW[0]) & (pupil_ts < tW[1]) & (~np.isnan(pupil_radius) ))[0]
                sns.histplot(pupil_radius[indy],element='step',stat='density',fill=False,binwidth=2,ax=ax,label=f'{macro_name_list[jj]}',color=cmap_macro[jj],lw=2,ls='-',common_norm=False)
            ax.legend()
            usrplt.adjust_spines(ax)
            ax.set_xlabel('Pupil size (pixels)',fontsize=12)  

        else:
            fig, axes = plt.subplots(1,2,figsize=(10,5))

            for jj, tW in enumerate(macro_windows):
                indy = np.where((pupil_ts >= tW[0]) & (pupil_ts < tW[1]) & (~np.isnan(pupil_radius) & (run_signal_p_s < 1)))[0]
                sns.histplot(pupil_radius[indy],element='step',stat='density',fill=False,binwidth=2,ax=axes[1],label=f'{macro_name_list[jj]}_rest',color=cmap_macro[jj],lw=2,ls='--',common_norm=False)
                
                indy = np.where((pupil_ts >= tW[0]) & (pupil_ts < tW[1]) & (~np.isnan(pupil_radius) & (run_signal_p_s > 1)))[0]
                sns.histplot(pupil_radius[indy],element='step',stat='density',fill=False,binwidth=2,ax=axes[1],label=f'{macro_name_list[jj]}_active',color=cmap_macro[jj],lw=3,common_norm=False)
                
                indy = np.where((run_ts >= tW[0]) & (run_ts < tW[1]) & (~np.isnan(run_signal_s)))[0]
                sns.histplot(run_signal_s[indy],element='step',stat='density',fill=False,ax=axes[0],label=f'{macro_name_list[jj]}',color=cmap_macro[jj],lw=3,common_norm=False)
        
            legend_elements = [Line2D([0], [0], color='k', lw=2, label='rest',ls='--'),Line2D([0], [0], color='k', lw=2, label='active',ls='-')]
            axes[1].legend(handles=legend_elements)
            axes[0].legend()     
            for ax in axes:
                usrplt.adjust_spines(ax)
            axes[0].set_xlabel('Running speed (cm/s)',fontsize=12)
            axes[1].set_xlabel('Pupil size (pixels)',fontsize=12)        

        plt.suptitle(f'{mID}, {rec_name}',fontsize=16,fontweight='bold')
        usrplt.save_fig_to_pptx(fig, prs)
        plt.savefig(os.path.join(PlotDir,f'pupil_distributions_{mID}_{rec_name}.pdf'),facecolor='white',dpi=300,bbox_inches='tight')
        plt.savefig(os.path.join(BehavDir,f'pupil_distributions_{mID}_{rec_name}.pdf'),facecolor='white',dpi=300,bbox_inches='tight')

    ##----- Determine windows to calculate synchrony measures -----##
    fig, axes = plt.subplots(2,1,figsize=(10,6),gridspec_kw={'height_ratios':[4,2]})
    ax = axes[0]
    plt.suptitle(f'{mID}, {rec_name}')

    ax.plot(run_ts/60,run_signal,'-k')
    ax.set_ylabel('Running speed')

    if drug_type in ['saline', 'psilocybin']:
        ax.vlines(np.array(injection_times)/60,*ax.get_ylim(),color=usrplt.cc[0],lw=2,zorder=4)

    if plot_pupil:
        ax2 = ax.twinx()
        ax2.plot(pupil_ts/60,pupil_radius,'-r')
        ax2.set_ylabel('Pupil size',color='r')
        
    #Create landmark list to create windows relative to
    landmarks = [open_ephys_start,open_ephys_end]

    #Add periods of evoked activity if it exists
    evoked_time_window_list = []
    evoked_type_list = []
    if stim_exists:
        for s in np.unique(stim_log['sweep']):
            for t in np.unique(stim_log.loc[stim_log.sweep == s]['stim_type']):
                sub_df = stim_log.loc[(stim_log.sweep == s) & (stim_log.stim_type == t)]
                tS = sub_df['onset'].min(); tE = sub_df['offset'].max()
                
                evoked_time_window_list.append([tS,tE])
                landmarks.extend([tS,tE])
                evoked_type_list.append(t)
        evoked_time_window_array = np.array(evoked_time_window_list)
                
    #Add injection times to create windows relative to those
    if drug_type in ['saline', 'psilocybin']:
        landmarks.extend(list(injection_times))

    if drug_type == 'isoflurane':
        #Add induction times plus buffer
        landmarks.extend(iso_induction_times)

        t5 = iso_induction_times[1]+5*60
        landmarks.append(t5)
        if stim_exists:
            ##Select first evoked block after iso induction
            ss = np.where(evoked_time_window_array[:,0] > iso_induction_times[-1])[0][0]
            t15 = evoked_time_window_list[ss][0]- window_t
            if t15 > t5:
                landmarks.append(t15) 
    landmarks = sorted(landmarks)

    ## Finish plotting
    ax = axes[1]
    if drug_type == 'isoflurane':
        if plot_pupil:
            t_mm = np.min([np.max(iso_times),np.max(pupil_ts)])
        else:
            t_mm = np.max(iso_times)
        indy = np.where(iso_times < t_mm)[0]
        ax.plot(iso_times[indy]/60,iso_level[indy],'-k')
        ax.set_ylabel('Iso level')
    ax.vlines(np.array(landmarks)/60,0,6,color='g',ls='-')

    for ax in axes:
        ax.autoscale(tight=True)
        usrplt.adjust_spines(ax)

    ##----- Create windows relative to landmarks -----##
    time_window_list = []
    for ii, l1 in enumerate(landmarks[:-1]):
        #Get next landmark and define window with 2 points
        l2 = landmarks[ii+1]

        if injection_times is not None:
            #Add buffer to injection times
            if l1 in injection_times:
                l1 += 60
            if l2 in injection_times:
                if (l2 - 60*3) < 60:
                    l2 -= 60*2
                else:
                    l2 -= 60*3
        tW = [l1,l2]

        #If window is already defined as stimulas evoked window
        #or if the window is less than the desired window length, just add to list as is
        if (tW in evoked_time_window_list) | (np.diff(tW)[0] <  window_t):
            time_window_list.append(tW)

        #Else create more windows between these 2 landmarks at the desired window width
        else:
            spacing = np.arange(l1,l2,window_t)
            for jj, t1 in enumerate(spacing[:-1]):
                t2 = spacing[jj+1]
                time_window_list.append([t1,t2])
            if spacing[-1] < l2:
                time_window_list.append([spacing[-1],l2])

    #Elimate windows that are not long enough
    keep = []
    for ll, tW in enumerate(time_window_list):
        if np.diff(tW)[0] > 60*2:
            keep.append(ll)
    time_window_array = np.array(time_window_list)
    time_window_list = list(time_window_array[keep])

    #Define periods of spontaneous activity
    filename_list = []; window_type_list = []
    for s, tW in enumerate(time_window_list):
        
        if list(tW) in evoked_time_window_list:
            indy = np.where(evoked_time_window_list == tW)[0][0]
            window_type = evoked_type_list[indy]
        else:
            window_type = 'spont'
        window_type_list.append(window_type)
        if drug_type in ['saline', 'psilocybin']:
            if tW[0] < injection_times[0]:
                filename_list.append(f'{window_type}-{s:02d}_pre-inj')
            elif (tW[0] >= injection_times[0]) & (tW[0] < injection_times[1]):
                filename_list.append(f'{window_type}-{s:02d}_post-{injection_types[0]}-inj')
            else:
                filename_list.append(f'{window_type}-{s:02d}_post-{injection_types[1]}-inj')
        elif drug_type == 'isoflurane':
            if tW[0] < iso_induction_times[0]:
                filename_list.append(f'{window_type}-{s:02d}_pre-iso')
            elif (tW[0] >= iso_induction_times[0]) & (tW[0] < iso_induction_times[1]):
                filename_list.append(f'{window_type}-{s:02d}_iso-ind')
            else:
                filename_list.append(f'{window_type}-{s:02d}_post-iso')
        else:
            t1 = int(tW[0]/60)
            t2 = int(tW[1]/60)
            filename_list.append(f'{window_type}-{s:02d}_{t1}-{t2}')

    spont_indy = []
    evoked_indy = []
    for ii, (epoch, tW) in enumerate(zip(filename_list,time_window_list)):
        print(f'\t{epoch:25s} = [{tW[0]:6.1f}, {tW[1]:6.1f}] -> {np.diff(tW)[0]/60:.1f} minutes')
        if 'biphasic' in epoch:
            c = usrplt.cc[1]
            evoked_indy.append(ii)
        elif 'circle' in epoch:
            c = usrplt.cc[2]
            evoked_indy.append(ii)
        else:
            c = usrplt.cc[ii]
            spont_indy.append(ii)
        ax.fill_between(np.array(tW)/60,0,6,color=c,lw=1.5,zorder=0,alpha=0.5)
        ax.text(tW[0]/60,ii*0.5,epoch,fontsize=8)
    usrplt.save_fig_to_pptx(fig, prs)
    plt.savefig(os.path.join(PlotDir,f'behavioral_traces_{mID}_{rec_name}.pdf'),facecolor='white',dpi=300,bbox_inches='tight')
    plt.savefig(os.path.join(BehavDir,f'behavioral_traces_{mID}_{rec_name}.pdf'),facecolor='white',dpi=300,bbox_inches='tight')

    spont_time_window_list = [time_window_list[ii] for ii in spont_indy]
    spont_filename_list = [filename_list[ii] for ii in spont_indy]
    evoked_time_window_list = [time_window_list[ii] for ii in evoked_indy]
    evoked_filename_list = [filename_list[ii] for ii in evoked_indy]


    # nPsi = np.sum([1 for f in spont_filename_list if 'psi' in f])
    # nSal = len(spont_filename_list) - nPsi
    # cmap_spont = np.concatenate((sns.color_palette('Blues',nSal),sns.color_palette('Reds',nPsi)))

    # nPsi = np.sum([1 for f in evoked_filename_list if 'psi' in f])
    # nSal = len(evoked_filename_list) - nPsi
    # cmap_evoked = np.concatenate((sns.color_palette('Blues',nSal),sns.color_palette('Reds',nPsi)))

    #Read in neuropixel data for spontaneous periods
    data_list, ts_list, neuron_indices, plot_tuple, FR_perblock = util.bin_spiking_data(probe_unit_data, time_window_list, time_bin=time_bin,fr_thresh=fr_thresh)
    
    if len(plot_tuple) == 9:
        boundaries_group, ticks_group, labels_group, celltypes, durations, areas, groups, supergroups, order_by_group = plot_tuple
    else:
        print('Experiment has no area information. Check to see if histology is completed. Exiting')
        exit()

    areas_sub = areas[neuron_indices]
    groups_sub = groups[neuron_indices]
    N = len(neuron_indices)

    #Plot rasters for each epoch
    neurons_per_probe = [len(np.unique(probe_unit_data[probei]['units'])) for probei in probe_list]
    neurons_per_probe_cum = np.cumsum(neurons_per_probe)
    for jj, tWindow in enumerate(time_window_list):
        ts = ts_list[jj]
        spikes = data_list[jj]
        T, N = spikes.shape
        
        # Display a spike raster of the image
        tmp = np.concatenate((neurons_per_probe,[150]))
        fig, axes =  plt.subplots(len(tmp),1,figsize=(12,12),gridspec_kw={'height_ratios':np.array(tmp)/np.sum(tmp)})
        plt.suptitle(filename_list[jj],y=0.925)

        #Define time window
        if T*time_bin < 60:
            tStart = 0
        else:
            tStart = 30# 5*60
        time_to_plot = 30 #sec
        tslice = slice(int(tStart/time_bin),int((tStart+time_to_plot)/time_bin))
        ts_sub = ts[tslice]

        #Params
        xticks = np.arange(0,int((time_to_plot+1)/time_bin),int(15/time_bin))
        tlabel = xticks*time_bin
        xticks = np.arange(ts_sub[0],ts_sub[-1]+1,15)

        for ii, probei in enumerate(probe_list):
            ax = axes[ii]; ax.set_title(probei)
            if ii == 0:
                nslice = slice(0,neurons_per_probe[0])
            else:
                nslice = slice(neurons_per_probe_cum[ii-1],neurons_per_probe_cum[ii])

            ax.imshow(spikes[tslice,nslice].T, aspect='auto',vmax=4,vmin=0, cmap='gray_r')
            ax.set_xticks([])

            uniq_groups, uniq_indices, num_groups = np.unique(groups[nslice],return_index=True, return_counts=True)
            group_labels = uniq_groups[np.argsort(uniq_indices)]
            nNeurons_area = num_groups[np.argsort(uniq_indices)]

            boundaries = np.concatenate(([0],np.cumsum(nNeurons_area)))
            yticks = boundaries[:-1] + np.diff(boundaries)/2

            ax.set_yticks(yticks,minor=True)
            ax.set_yticks(boundaries,minor=False)
            ax.set_yticklabels(group_labels,minor=True)
            ax.set_yticklabels([],minor=False)
            ax.set_ylabel("Unit")

        #Plot pupil and running
        indy = np.where((run_ts >= ts_sub[0]) & (run_ts <= ts_sub[-1]))[0]
        ax = axes[-1]; ax.autoscale(tight=True)
        ax.plot(run_ts[indy],run_signal[indy],lw=0.6,color='k')
        ax.set_ylabel("Speed (cm/s)"); ax.set_xlabel("Time (s)")
        # ax.set_xticks([]); 
        ax.set_title('Behavior')

        if plot_pupil:
            indy = np.where((pupil_ts >= ts_sub[0]) & (pupil_ts <= ts_sub[-1]))[0]
            ax = ax.twinx()
            ax.plot(pupil_ts[indy],pupil_radius[indy],lw=0.6,color=usrplt.cc[1])
            ax.set_ylabel("Pupil radius (pix)") 

        usrplt.save_fig_to_pptx(fig, prs)
        plt.savefig(os.path.join(PlotDir,f'raster_{filename_list[jj]}.pdf'),facecolor='white',dpi=300,bbox_inches='tight')

    #Define neurons we want to look at
    nCond = len(time_window_list)
    groups_of_interest = ['MO','SSp','SM-TH']
    FR_overall = np.concatenate([probe_unit_data[probei]['firing_rate'] for probei in probe_list])
    FR_mask1 = FR_overall > fr_thresh

     # nslice = np.concatenate([np.where((groups == g))[0] for g in groups_of_interest])
    nslice = np.concatenate([np.where((groups == g) & (FR_overall > 1))[0] for g in groups_of_interest])
    N = len(nslice)
    print(f'{N} neurons in groups of interest')

    #Get labels
    uniq_groups, uniq_indices, num_groups = np.unique(groups[nslice],return_index=True, return_counts=True)
    labels_MO = uniq_groups[np.argsort(uniq_indices)]
    nNeurons_area = num_groups[np.argsort(uniq_indices)]
    boundaries_MO = np.concatenate(([0],np.cumsum(nNeurons_area)))
    ticks_MO = boundaries_MO[:-1] + np.diff(boundaries_MO)/2

    np.savez(join(SaveDir,f'plotting_{mID}_{rec_name}.npz'),labels_MO=labels_MO,boundaries_MO=boundaries_MO,ticks_MO=ticks_MO,filename_list=filename_list,time_window_list=time_window_list)
    # exit()
    #Preallocate
    synchrony_mat_DW = np.zeros((nCond,N,N))
    synchrony_mat_YN = np.zeros((nCond,N,N))
    synchrony_mat_CC = np.zeros((nCond,N,N))
    correlation_mat = np.zeros((nCond,N,N))

    tmp_list_DW = []
    tmp_list_YN = []
    tmp_list_CC = []

    import ctypes
    so_file = "/home/david.wyrick/projects/zap-n-zip/src/spike_time_tiling_coefficient.so"
    lib = ctypes.CDLL(so_file)
    
    # Define argument and return types for the functions
    lib.run_P.argtypes = [ctypes.c_int, ctypes.c_int, ctypes.c_double, ctypes.POINTER(ctypes.c_double), ctypes.POINTER(ctypes.c_double)]
    lib.run_P.restype = ctypes.c_double

    lib.run_T.argtypes = [ctypes.c_int, ctypes.c_double, ctypes.c_double, ctypes.c_double, ctypes.POINTER(ctypes.c_double)]
    lib.run_T.restype = ctypes.c_double

    lib.run_sttc.argtypes = [ctypes.POINTER(ctypes.c_int), ctypes.POINTER(ctypes.c_int), ctypes.POINTER(ctypes.c_double), ctypes.POINTER(ctypes.c_double), ctypes.POINTER(ctypes.c_double), ctypes.POINTER(ctypes.c_double)]
    lib.run_sttc.restype = ctypes.c_float


    for iS, tWindow in enumerate(time_window_list):
        if iS in spont_indy:
            exp_type = 'spont'
        else:
            exp_type = 'evoked'

        print('\n',exp_type, filename_list[iS])

        spktime_dict = {}
        index = 0
        for probei in probe_list:
            tmp_dict = util.get_spike_time_dict(probe_unit_data[probei]['spike_times'], probe_unit_data[probei]['spike_clusters'], probe_unit_data[probei]['units'], tWindow)
            for k in tmp_dict.keys():
                spktime_dict[index+k] = tmp_dict[k]
            index += len(probe_unit_data[probei]['units'])

        #Calculate correlation matrix for neurons in 'MO', 'SSp' and 'SM-TH'
        data = data_list[iS][:,nslice]
        correlation_mat[iS] = E.correlation_FC(data)

        fig = usrplt.visualize_matrix(np.abs(correlation_mat[iS]),title=filename_list[iS],ticks=ticks_MO,labels=labels_MO,boundaries=boundaries_MO,cbar_label='| Correlation |')
        usrplt.save_fig_to_pptx(fig, prs)

        ijs = list(combinations(np.arange(N),2))
        pair_indices = list(combinations(nslice,2))
        print(f'\t {len(ijs)} pairs to calculate synchrony')

        for ij, pair in tqdm(zip(ijs,pair_indices)):
            i,j = ij; iN, jN = pair

            t_start=tWindow[0];t_stop=tWindow[1]; T = t_stop-t_start
            indy = np.where((spktime_dict[iN] > t_start) & (spktime_dict[iN] < t_stop))[0]
            spk_train_i = spktime_dict[iN][indy]
            indy = np.where((spktime_dict[jN] > t_start) & (spktime_dict[jN] < t_stop))[0]
            spk_train_j = spktime_dict[jN][indy]

            # #Calculate using my code
            # sttc = calculate_STTC_DW(spk_train_i,spk_train_j, synchrony_window, T)
            # synchrony_mat_DW[iS,i,j] = sttc
            # synchrony_mat_DW[iS,j,i] = sttc

            # combo = f'{groups[iN]}-{groups[jN]}'
            # tmp_list_DW.append((i,j,areas[iN],groups[iN],areas[jN],groups[jN],combo,exp_type,filename_list[iS],sttc))

            #Calculate using yasmins code
            sttc = calculate_STTC_YN(spk_train_i,spk_train_j, synchrony_window, np.array(tWindow))
            synchrony_mat_YN[iS,i,j] = sttc
            synchrony_mat_YN[iS,j,i] = sttc

            combo = f'{groups[iN]}-{groups[jN]}'
            tmp_list_YN.append((i,j,areas[iN],groups[iN],areas[jN],groups[jN],combo,exp_type,filename_list[iS],sttc))

            N1v = int(len(spk_train_i)) #Total number of spikes in spk_train_i
            N2v = int(len(spk_train_j)) #Total number of spikes in spk_train_j
            Time = np.array([t_start,t_stop], dtype=np.float64) #Window spikes are within
            dt = np.float64(synchrony_window)
            sttc = lib.run_sttc(
                ctypes.pointer(ctypes.c_int(N1v)),
                ctypes.pointer(ctypes.c_int(N2v)),
                ctypes.pointer(ctypes.c_double(dt)),
                Time.ctypes.data_as(ctypes.POINTER(ctypes.c_double)),
                spk_train_i.ctypes.data_as(ctypes.POINTER(ctypes.c_double)),
                spk_train_j.ctypes.data_as(ctypes.POINTER(ctypes.c_double))
            )
            synchrony_mat_CC[iS,i,j] = sttc
            synchrony_mat_CC[iS,j,i] = sttc
            
            combo = f'{groups[iN]}-{groups[jN]}'
            tmp_list_CC.append((i,j,areas[iN],groups[iN],areas[jN],groups[jN],combo,exp_type,filename_list[iS],sttc))

        
        fig, axes = plt.subplots(1,2,figsize=(10,5))
        plt.suptitle(filename_list[iS])
        usrplt.visualize_matrix(synchrony_mat_YN[iS],ax=axes[0],center=0,cmap='RdBu_r',title='Yasmins code',ticks=ticks_MO,labels=labels_MO,boundaries=boundaries_MO,cbar_label='Synchrony')
        usrplt.visualize_matrix(synchrony_mat_CC[iS],ax=axes[1],center=0,cmap='RdBu_r',title='Code from paper',ticks=ticks_MO,labels=labels_MO,boundaries=boundaries_MO,cbar_label='Synchrony')
        usrplt.save_fig_to_pptx(fig, prs)
        plt.savefig(os.path.join(PlotDir,f'synchrony_{filename_list[iS]}.pdf'),facecolor='white',dpi=300,bbox_inches='tight')

    np.savez(join(SaveDir,f'synchrony_mats_{mID}_{rec_name}.npz'),correlation_mat=correlation_mat,synchrony_mat_DW=synchrony_mat_DW,synchrony_mat_YN = synchrony_mat_YN,synchrony_mat_CC=synchrony_mat_CC)
    
    synchrony_df = pd.DataFrame(np.stack(tmp_list_CC),columns=['i','j','area_i','group_i','area_j','group_j','combo','exp_type','epoch','val'])
    synchrony_df = synchrony_df.astype({'i':int,'j':int,'area_i':str,'group_i':str,'area_j':str,'group_j':str,'combo':str,'exp_type':str,'epoch':str,'val':float})
    synchrony_df.to_hdf(join(SaveDir,f'synchrony_df_CC_{mID}_{rec_name}.h5'),'df')


    fig, ax = plt.subplots(figsize=(10,4))
    plt.suptitle('Spontaneous epochs')
    sns.boxplot(x='combo',y='val',hue='epoch',ax=ax,hue_order=spont_filename_list,data=synchrony_df.loc[synchrony_df.exp_type == 'spont'])
    usrplt.save_fig_to_pptx(fig, prs)

    fig, ax = plt.subplots(figsize=(10,4))
    plt.suptitle('Evoked epochs')
    sns.boxplot(x='combo',y='val',hue='epoch',ax=ax,hue_order=evoked_filename_list,data=synchrony_df.loc[synchrony_df.exp_type == 'evoked'])
    usrplt.save_fig_to_pptx(fig, prs)



    synchrony_df = pd.DataFrame(np.stack(tmp_list_YN),columns=['i','j','area_i','group_i','area_j','group_j','combo','exp_type','epoch','val'])
    synchrony_df = synchrony_df.astype({'i':int,'j':int,'area_i':str,'group_i':str,'area_j':str,'group_j':str,'combo':str,'exp_type':str,'epoch':str,'val':float})
    synchrony_df.to_hdf(join(SaveDir,f'synchrony_df_YN_{mID}_{rec_name}.h5'),'df')


    fig, ax = plt.subplots(figsize=(10,4))
    plt.suptitle('Spontaneous epochs')
    sns.boxplot(x='combo',y='val',hue='epoch',ax=ax,hue_order=spont_filename_list,data=synchrony_df.loc[synchrony_df.exp_type == 'spont'])
    usrplt.save_fig_to_pptx(fig, prs)

    fig, ax = plt.subplots(figsize=(10,4))
    plt.suptitle('Evoked epochs')
    sns.boxplot(x='combo',y='val',hue='epoch',ax=ax,hue_order=evoked_filename_list,data=synchrony_df.loc[synchrony_df.exp_type == 'evoked'])
    usrplt.save_fig_to_pptx(fig, prs)

    prs.save(join(SaveDir,f'synchrony_figs_{mID}_{rec_name}.pptx'))
    print('DONE!!!')
