base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
base_dir_local = '/data/tiny-blue-dot/zap-n-zip/EEG_exp'; base_dir = base_dir_server
ProjDir = '/home/david.wyrick/projects/zap-n-zip/'
ServDir = '/data/projects/zap-n-zip/'
BehavDir = '/home/david.wyrick/projects/zap-n-zip/plots/behavior'

#Basea
import gspread
from os.path import join
from glob import glob
import json, os, time, sys, argparse
import numpy as np
import pandas as pd
import scipy.stats as st
import scipy.signal as sig
import pingouin as pg
from scipy.ndimage import gaussian_filter
from scipy.interpolate import interp1d
from itertools import combinations

#Plot
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
from matplotlib.gridspec import GridSpec
from matplotlib.lines import Line2D   

#Project
from tbd_eeg.data_analysis.eegutils import EEGexp
import tbd_eeg.data_analysis.Utilities.utilities as tbd_util
from tbd_eeg.data_analysis.Utilities.behavior_movies import Movie

#Allen
from allensdk.core.mouse_connectivity_cache import MouseConnectivityCache
mcc = MouseConnectivityCache(resolution=10)
str_tree = mcc.get_structure_tree()

#Read in allen CCF
ccfsum = pd.read_csv('/home/david.wyrick/projects/zap-n-zip/data/mouse_connectivity/ontology_v2.csv')

#User
sys.path.append(ProjDir)
import util
import plotting as usrplt

#Single cell stuff
import elephant.statistics as est
from quantities import ms, s, Hz
from neo.core import SpikeTrain

#Define behavioral states
# behavior_ranges = {0: [0,1], 1: [1,10], 2: [10,30], 3:[30,500]}
# behavior_dict = {0: 'rest (0-1cm/s)', 1: 'walk (1-10cm/s)', 2: 'shuffle (10-30cm/s)', 3: 'run (>30cm/s)'}
# behavior_strs2 = ['rest','walk','shuffle','run']

behavior_ranges = {0: [0,1], 1: [1,15], 2: [15,500]}#, 3:[30,500]}
behavior_dict = {0: 'rest (0-1cm/s)', 1: 'walk (1-15cm/s)', 2: 'run (>15cm/s)'}
behavior_strs2 = ['rest','walk','run']

# behavior_ranges = {0: [0,1], 1: [1,500]}#, 3:[30,500]}
# behavior_dict = {0: 'rest (0-1cm/s)', 1: 'active (>1cm/s)'}
# behavior_strs2 = ['rest','active']

behavior_strs = list(behavior_dict.values())
nBehaviors = len(behavior_strs)

#Define windows to calculate firing rate
windows_of_interest = [[.002,.025],[.075,.3],[.3,1],[1,4.5]]
window_strs = ['evoked (2-25ms)','rebound (75-300ms)','post-rebound (0.3-1s)','ISI (1-4.5s)']
nWindows = len(window_strs)

gc = gspread.service_account() # need a key file to access the account (step 2) 
sh = gc.open('Templeton-log_exp') # open the spreadsheet 

df = pd.DataFrame(sh.sheet1.get()) 
exp_table = df.T.set_index(0).T # put it in a nicely formatted dataframeexp_table.iloc[10:]
exp_table = exp_table.set_index('mouse_name')

##===== ============================ =====##
##===== Parse Command Line Arguments =====##
parser = argparse.ArgumentParser(description='single-cell-metrics')

##===== Data Options =====##
parser.add_argument('--mID',type=str, default='mouse666193',
                    help='mouse to perform analysis on')

parser.add_argument('--rec_name',type=str, default='pilot_aw_psi_2023-02-16_10-55-48',
                    help='experiment to perform analysis on')

parser.add_argument('--time_bin_ms',type=int, default=50,
                    help='time_bin_ms')

parser.add_argument('--fr_thresh',type=float, default=1,
                    help='Firing rate threshold for neurons to include in analysis')

parser.add_argument('--calculate_evoked',type=int, default=1,
                    help='blah')

parser.add_argument('--window_t_min',type=int, default=15,
                    help='Window width (mins) to segment data into')


def plot_boxplot_comparison_spont(dFR_df, epoch_i,epoch_j, g = None, pval_thresh = 0.05):
    bar_width=0.8
    if g is None:
        x = 'group'
        sub_df = dFR_df.loc[(dFR_df.epoch_i == epoch_i) & (dFR_df.epoch_j == epoch_j)]
    else:
        x = 'area'
        sub_df = dFR_df.loc[(dFR_df.epoch_i == epoch_i) & (dFR_df.epoch_j == epoch_j) & (dFR_df.group == g)]

    #Determine unique boxes
    uniq_boxes = np.unique(sub_df[x])
    nBoxes = len(uniq_boxes)

    #Plot modulation index
    fig, axes = plt.subplots(2,1,figsize=(10,8))

    if x == 'area':
        plt.suptitle(f'Group {g}: "{epoch_i}" vs "{epoch_j}"')
    else:
        plt.suptitle(f'All groups "{epoch_i}" vs "{epoch_j}"')

    axes[0].set_title(f'Modulation index')
    axes[1].set_title(f'\u0394-Firing rate:')
    for ii, var in enumerate(['MI','dFR']):
        ax = axes[ii]
    
        gs = sns.boxplot(x=x,y=var,hue='behavior',palette=sns.color_palette('Reds',nBehaviors),order=uniq_boxes,hue_order=behavior_strs,data=sub_df,width=bar_width,legend=False,ax=ax)
        ax.set_xlabel(x)
        xlim = ax.get_xlim(); ylim = ax.get_ylim()
        ax.hlines(0,-0.5,nBoxes-0.5,color='k',zorder=0)
        ax.set_xlim(xlim)
        ax.set_xlabel('')
        usrplt.adjust_spines(ax)
        for jj, b in enumerate(uniq_boxes):
            n = len(sub_df.loc[(sub_df[x] == b)])
            if ii == 0:
                ax.text(jj-.15,ylim[0],n)
            for kk, w in enumerate(behavior_strs):
                sub_sub_df = sub_df.loc[(sub_df[x] == b) & (sub_df.behavior == w)]
                y = sub_sub_df[var].values
                if (len(y) < 2) | (np.all(y == 0)):
                    continue
                
                if np.all(y == 0):
                    continue
                res = pg.wilcoxon(y)
                pval = res['p-val'][0]
                x_pos = jj-0.45+kk*bar_width/len(behavior_strs)
                if pval < pval_thresh/len(behavior_strs):
                    ax.text(x_pos,0.9*ylim[1],'*',fontsize=20,fontweight='bold',color='k')

    if x == 'area':
        axes[1].set_xlabel('Area')
    else:
        axes[1].set_xlabel('Group')

    return fig

def plot_boxplot_comparison_evoked(dFR_df, epoch_i,epoch_j, b_str, g = None, pval_thresh = 0.05):
    bar_width=0.8
    if g is None:
        x = 'group'
        sub_df = dFR_df.loc[(dFR_df.epoch_i == epoch_i) & (dFR_df.epoch_j == epoch_j) & (dFR_df.behavior == b_str)]
    else:
        x = 'area'
        sub_df = dFR_df.loc[(dFR_df.epoch_i == epoch_i) & (dFR_df.epoch_j == epoch_j) & (dFR_df.behavior == b_str) & (dFR_df.group == g)]

    #Determine unique boxes
    uniq_boxes = np.unique(sub_df[x])
    nBoxes = len(uniq_boxes)

    #Plot modulation index
    fig, axes = plt.subplots(2,1,figsize=(10,8))
    # plt.suptitle(f'Behavior: {b_str}\n"{epoch_i}" vs "{epoch_j}"')

    axes[0].set_title(f'Modulation index')
    axes[1].set_title(f'\u0394-Firing rate:')
    if x == 'area':
        plt.suptitle(f'Group: {g}, Behavior: {b_str}\n"{epoch_i}" vs "{epoch_j}"')
    else:
        plt.suptitle(f'All groups Behavior: {b_str}\n"{epoch_i}" vs "{epoch_j}"')
    
    for ii, var in enumerate(['MI','dFR']):
        ax = axes[ii]
    
        gs = sns.boxplot(x=x,y=var,hue='window',palette=sns.color_palette('Purples',4),order=uniq_boxes,hue_order=window_strs,data=sub_df,width=bar_width,legend=False,ax=ax)

        if var == 'dFR':
            tmp = sub_df['dFR'].values
            mask = (~np.isnan(tmp)) & (~np.isinf(tmp))
            ymin = np.round(np.nanpercentile(tmp[mask],2.5))
            ymax = np.round(np.nanpercentile(tmp[mask],97.5))
            if np.isnan(ymax) | np.isinf(ymax):
                print('inf or nan values in data')
            else:
                ax.set_ylim([ymin,ymax])
        ax.set_xlabel(x)
        xlim = ax.get_xlim(); ylim = ax.get_ylim()
        ax.hlines(0,-0.5,nBoxes-0.5,color='k',zorder=0)
        ax.set_xlim(xlim)
        ax.set_xlabel('')
        usrplt.adjust_spines(ax)

        counter = 0
        for jj, b in enumerate(uniq_boxes):
            n = len(sub_df.loc[(sub_df[x] == b) & (sub_df.window == window_strs[0])])
            if ii == 1:
                ax.text(jj-.15,ylim[0],n)
            for kk, w in enumerate(window_strs):
                x_pos = jj-0.45+kk*bar_width/len(window_strs)
                # ax.text(x_pos,0.9*ylim[1],counter,fontsize=10,color='k')
                # counter += 1
                sub_sub_df = sub_df.loc[(sub_df[x] == b) & (sub_df.window == w)]
                y = sub_sub_df[var].values

                if (len(y) < 2) | (np.all(y == 0)):
                    continue
                if np.all(y == 0):
                    continue
                res = pg.wilcoxon(y)
                pval = res['p-val'][0]

                if pval < pval_thresh/len(window_strs):
                    ax.text(x_pos,0.9*ylim[1],'*',fontsize=20,fontweight='bold',color='k')

    if x == 'area':
        axes[1].set_xlabel('Area')
    else:
        axes[1].set_xlabel('Group')
    
    return fig

if __name__ == '__main__':

    ##------------------------------------------
    # Parse the arguments
    args = parser.parse_args()
    mID = args.mID
    rec_name = args.rec_name
    fr_thresh = args.fr_thresh
    time_bin_ms = args.time_bin_ms
    time_bin = time_bin_ms/1000
    calculate_evoked = bool(args.calculate_evoked)

    #Segement data into X minute windows
    window_t_min = args.window_t_min
    window_t = window_t_min*60
    
    #Extract row from templeton dataframe
    exp_df = exp_table.loc[exp_table.exp_name == rec_name]
    if type(exp_df) == pd.core.series.Series:
        exp_df = exp_df.to_frame().T

    drug_type = exp_df['drug'].values[0]
    stim_type = exp_df['stimulation'].values[0]
    print(f'Experiment type: {stim_type}, {drug_type}')

    #Define directories
    SaveDir = join(ServDir,'results','FR_20240412',mID,rec_name)
    if not os.path.exists(SaveDir):
        os.makedirs(SaveDir)

    PlotDir = join(SaveDir,'plots')
    if not os.path.exists(PlotDir):
        os.makedirs(PlotDir)

    #%% Upload the whole experiment and generate the global clock
    file_name = os.path.join(base_dir_server,mID,rec_name,'experiment1','recording1')
    exp = EEGexp(file_name, preprocess=False, make_stim_csv=False)

    #Load probe data
    probe_unit_data, probe_info, total_units = tbd_util.get_neuropixel_data(exp)
    probe_list = [x.replace('_sorted', '') for x in exp.experiment_data if 'probe' in x]

    #Get recording start time
    probei = probe_list[0]
    open_ephys_start = np.round(np.min(probe_unit_data[probei]['spike_times']))
    open_ephys_end = np.round(np.max(probe_unit_data[probei]['spike_times']))

    ##------------------------------------------
    ##TODO: Once Lydia starts recording psilocybin/saline experiments where she does 2 messages per injection, this code will have to change
    #For saline & psilocybin experiments, get injection times and types of injection
    if drug_type in ['saline', 'psilocybin']:
        injection_times = [float(exp_df['First injection time (s)'].values[0]),
                        float(exp_df['Second injection time (s)'].values[0])]
        injection_time_windows = None
        #Determine injection type
        if 'psi' in rec_name:
            injection_types = ['sal1','psi']
        else:
            injection_types = ['sal1', 'sal2']

    elif drug_type == 'ketanserin+psilocybin':
        injection_time_windows = np.array([np.array(exp_df['First injection time (s)'].values[0].split(','),dtype=float),
                                np.array(exp_df['Second injection time (s)'].values[0].split(','),dtype=float)])
        
        #Take second time in each window as "injection time"
        injection_times = [injection_time_windows[0,1],injection_time_windows[1,1]]

        injection_types = ['ket','psi']
    else:
        injection_time_windows = None
        injection_times = None

    #For isoflurane experiments, get iso level
    if drug_type == 'isoflurane':
        iso_level, iso_times = exp.load_analog_iso()
        iso_induction_times, iso_maintenance_times = exp.load_iso_times()
    else:
        iso_induction_times = None

    #Extract the timestamps of the selected stimuli
    try:
        stim_log = pd.read_csv(exp.stimulus_log_file)

        if 'circle' in np.unique(stim_log['stim_type']):
            vis_stim_exists = True
        else:
            vis_stim_exists = False
        if 'biphasic' in np.unique(stim_log['stim_type']):
            stim_exists = True
        else:
            stim_exists = False
    except:
        stim_log = None
        stim_exists = False
        vis_stim_exists = False

    ##------------------------------------------
    #Load behavior
    run_ts, _, run_signal, run_signal_s, pupil_ts, pupil_radius, run_signal_p, run_signal_p_s, plot_pupil = util.get_behavioral_data(exp, mID, rec_name)
    f_run = interp1d(run_ts,run_signal)
    f_run_s = interp1d(run_ts,run_signal_s)

    #Get time windows for each epoch
    epoch_list, time_window_list = util.define_epochs_of_interest([open_ephys_start,open_ephys_end], drug_type, window_t_min=window_t_min, injection_times=injection_times,injection_time_windows=injection_time_windows, iso_induction_times=iso_induction_times, stim_log=stim_log)
    filename_list = epoch_list
    for e, tW in zip(epoch_list, time_window_list):
        print(f'{e}: {np.diff(tW)[0]/60:.1f} min')

    ##------------------------------------------
    # Create a new PowerPoint presentation to save figures to
    prs = Presentation()

    #Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f'mID = {mID}; rec_name = {rec_name}'
    slide.placeholders[1].text = f'Firing rate changes due to psilocybin in evoked & spontaneous blocks'

    if drug_type in ['saline', 'psilocybin']:
        tb = 180
        if injection_times[0] < 200:
            tb = 60
        macro_windows = [[open_ephys_start,injection_times[0]-tb]]; macro_name_list = ['pre-inj']
        macro_windows.append([injection_times[0]+60,injection_times[1]-180]); macro_name_list.append(f'post_{injection_types[0]}_inj')
        macro_windows.append([injection_times[1]+60,open_ephys_end]); macro_name_list.append(f'post_{injection_types[1]}_inj')
        if drug_type == 'saline':
            cmap_macro = sns.xkcd_palette(['silver','dark sky blue','cobalt blue'])
        else:
            cmap_macro = sns.xkcd_palette(['silver','dark sky blue','darkish red'])
    
    elif drug_type == 'ketanserin+psilocybin':
        macro_windows = [[open_ephys_start,injection_time_windows[0,0]]]; macro_name_list = ['pre-inj']
        macro_windows.append([injection_time_windows[0,1],injection_time_windows[1,0]]); macro_name_list.append(f'post_{injection_types[0]}_inj')
        macro_windows.append([injection_time_windows[1,1],open_ephys_end]); macro_name_list.append(f'post_{injection_types[1]}_inj')
        cmap_macro = sns.xkcd_palette(['silver','goldenrod','darkish red'])
    
    elif drug_type == 'isoflurane':
        macro_windows = [[open_ephys_start,iso_induction_times[0]-120]]; macro_name_list = ['pre-iso']
        macro_windows.append([iso_induction_times[0],iso_induction_times[1]]); macro_name_list.append(f'iso-ind')
        macro_windows.append([iso_induction_times[1]+60*5,open_ephys_end]); macro_name_list.append(f'post-iso')
        cmap_macro = sns.xkcd_palette(['silver','light teal','teal'])

    elif drug_type == 'urethane':
        macro_windows = [[open_ephys_start,open_ephys_end/3]]; macro_name_list = ['urethane_1/3']
        macro_windows.append([open_ephys_end/3,open_ephys_end*2/3]); macro_name_list.append('urethane_2/3')
        macro_windows.append([open_ephys_end*2/3,open_ephys_end]); macro_name_list.append('urethane_3/3')
        cmap_macro = sns.color_palette('Oranges',3)

    ##------------------------------------------
    if plot_pupil:
        if drug_type == 'urethane':
            fig, ax = plt.subplots(figsize=(6,6))
            for jj, tW in enumerate(macro_windows):
                indy = np.where((pupil_ts >= tW[0]) & (pupil_ts < tW[1]) & (~np.isnan(pupil_radius) ))[0]
                sns.histplot(pupil_radius[indy],element='step',stat='density',fill=False,binwidth=2,ax=ax,label=f'{macro_name_list[jj]}',color=cmap_macro[jj],lw=2,ls='-',common_norm=False)
            ax.legend()
            usrplt.adjust_spines(ax)
            ax.set_xlabel('Pupil size (pixels)',fontsize=12)  

        else:
            fig, axes = plt.subplots(1,3,figsize=(12,4)); plt.suptitle(f'{mID}, {rec_name}',y=1.05)
            pupil_time_bin = 1/30
            for jj, tW in enumerate(macro_windows):
                
                #Plot pupil size during rest periods
                indy = np.where((pupil_ts >= tW[0]) & (pupil_ts < tW[1]) & (~np.isnan(pupil_radius) & (run_signal_p_s < 1)))[0]
                n = len(indy); t = (n*pupil_time_bin)/60
                sns.kdeplot(pupil_radius[indy],ax=axes[0],label=f'{macro_name_list[jj]}_rest: {t:.1f}min',color=cmap_macro[jj],lw=3)

                #Plot pupil size during active periods
                indy = np.where((pupil_ts >= tW[0]) & (pupil_ts < tW[1]) & (~np.isnan(pupil_radius) & (run_signal_p_s > 1)))[0]
                n = len(indy); t = (n*pupil_time_bin)/60
                sns.kdeplot(pupil_radius[indy],ax=axes[1],label=f'{macro_name_list[jj]}_active: {t:.1f}min',color=cmap_macro[jj],lw=3)
                
                #Plot running speed 
                indy = np.where((run_ts >= tW[0]) & (run_ts < tW[1]) & (~np.isnan(run_signal)))[0]
                sns.kdeplot(run_signal[indy],ax=axes[2],label=macro_name_list[jj],color=cmap_macro[jj],lw=3)

            vmin = np.nanpercentile(pupil_radius,2); vmax = np.nanpercentile(pupil_radius,98)
            axes[0].set_title('Pupil during rest')
            axes[1].set_title('Pupil during active')
            # axes[0].legend()
            for ax in axes[:-1]:
                usrplt.adjust_spines(ax)
                ax.set_xlabel('Normalized pupil radius')
                ax.set_xlim([vmin,vmax])
                ax.legend(loc=2)

            ax = axes[-1]
            usrplt.adjust_spines(ax)
            ax.set_xlabel('Running speed (cm/s)')
            ax.set_title('Locomotion')
            # import pdb; pdb.set_trace()
        
        usrplt.save_fig_to_pptx(fig, prs)
        plt.savefig(os.path.join(PlotDir,f'pupil_distributions_{mID}_{rec_name}.pdf'),facecolor='white',dpi=300,bbox_inches='tight')
        plt.savefig(os.path.join(BehavDir,f'pupil_distributions_{mID}_{rec_name}.pdf'),facecolor='white',dpi=300,bbox_inches='tight')

    ##------------------------------------------
    fig, ax = plt.subplots(figsize=(10,4))
    plt.suptitle(f'{mID}, {rec_name}')

    ax.plot(run_ts/60,run_signal,'-k')
    ax.set_ylabel('Running speed')

    if drug_type in ['saline', 'psilocybin','ketanserin+psilocybin']:
        ax.vlines(np.array(injection_times)/60,*ax.get_ylim(),color=usrplt.cc[0],lw=2,zorder=4)

    if plot_pupil:
        ax2 = ax.twinx()
        ax2.plot(pupil_ts/60,pupil_radius,'-r')
        ax2.set_ylabel('Pupil size',color='r')  
    usrplt.save_fig_to_pptx(fig, prs)

    ##------------------------------------------
    #Subselect epochs where we have spontaneous and biphasic data
    spont_indy = []; evoked_indy = []; indy = []
    for ii, (epoch, tW) in enumerate(zip(filename_list,time_window_list)):
        if 'biphasic' in epoch:
            evoked_indy.append(ii)
            indy.append(ii)
        elif ('circle' in epoch) | ('natural' in epoch):
            pass
        else:
            spont_indy.append(ii)
            indy.append(ii)
    
    # import pdb; pdb.set_trace()
    spont_time_window_list = [time_window_list[ii] for ii in spont_indy]
    spont_filename_list = [filename_list[ii] for ii in spont_indy]
    evoked_time_window_list = [time_window_list[ii] for ii in evoked_indy]
    evoked_filename_list = [filename_list[ii] for ii in evoked_indy]

    time_window_list = [time_window_list[ii] for ii in indy]
    filename_list = [filename_list[ii] for ii in indy]

    ##------------------------------------------
    #Read in neuropixel data for spontaneous periods
    data_list, ts_list, neuron_indices, plot_tuple, FR_perblock = util.bin_spiking_data(probe_unit_data, time_window_list, time_bin=time_bin,fr_thresh=fr_thresh)
    
    if len(plot_tuple) == 10:
        boundaries_group, ticks_group, labels_group, celltypes, durations, layers, areas, groups, supergroups, order_by_group = plot_tuple
    else:
        print('Experiment has no area information. Check to see if histology is completed. Exiting')
        exit()
    
    ##------------------------------------------
    ## Determine running behaviors during spontaneous periods
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    behav_list = []; spont_mask_list = []
    for jj, tWindow in enumerate(time_window_list):
        print(f'{filename_list[jj]}')
        if jj > 7:
            if jj == 8:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
            fig_num = jj - 8
        else:
            fig_num = jj
        position = [(fig_num//4)*5, (fig_num%4)*1.875]
        
        ts = ts_list[jj]
        T = len(ts)
        tMin = T*time_bin/60
        run_spont = f_run_s(ts) #Use "smoothed" running trace to interpolate to spike bins
        behavior_map = np.zeros(run_spont.shape)

        for key, b_range in behavior_ranges.items():
            indy = np.where((run_spont >= b_range[0]) & (run_spont < b_range[1]))[0]
            behavior_map[indy] = key
            t = len(indy)*time_bin
            print(f'\t: {behavior_dict[key]:10s} -> {t:.1f}s')

        x = behavior_map
        transitions = np.where(np.diff(np.concatenate(([x[0]+1],x,[x[-1]+1]))))[0]

        #Ensure behavioral bouts are at least 5 seconds
        min_t_bout = 3
        for ii, t1 in enumerate(transitions[:-1]):
            t2 = transitions[ii+1]
             
            if (t2-t1)*time_bin < min_t_bout:
                behavior_map[t1:t2] = np.nan
        behav_list.append(behavior_map)

        fig, ax = plt.subplots(figsize=(5,1))
        ax.set_title(filename_list[jj])
        ax2 = ax.twinx()
        sns.heatmap(behavior_map.reshape(1,-1),vmax=nBehaviors-1,vmin=0,ax=ax,cbar=False,cmap=sns.color_palette('Reds',nBehaviors),cbar_kws={'ticks':np.arange(nBehaviors)})
        
        ax2.plot(f_run(ts),color='k')
        ax2.plot(f_run_s(ts),color='b')
        ax2.set_ylabel('Speed (cm/s)')
        ax.set_yticks([])
        ax.set_xticks([0,T/2,T]); ax.set_xticklabels([0,np.round(tMin/2,2),np.round(tMin,2)],rotation=0);ax.set_xlabel('Time (mins)')
        usrplt.save_fig_to_pptx(fig, prs,slide=slide,position=position)

        if 'spont' in filename_list[jj]:
            spont_periods = np.ones((T))
            spont_mask_list.append(spont_periods)
        elif 'biphasic' in filename_list[jj]:
            block = filename_list[jj]
            stim, epoch = block.split('_')
            stim_type, _, sweep = stim.split('-')
            sweep = int(sweep)
            stim_log_b = stim_log.loc[(stim_log.stim_type == 'biphasic') & (stim_log.sweep == sweep)]

            #Mask ISI periods
            spont_periods = np.zeros((T))
            buffer = 0.05
            onsets = stim_log_b['onset'].values
            for ii, on in enumerate(onsets[:-1]):
                t1 = on + 1
                t2 = onsets[ii+1] - buffer
                indy = np.where((ts >= t1) & (ts < t2))[0]
                spont_periods[indy] = 1
            spont_mask_list.append(spont_periods)

    # if stim_exists:
    #     buffer = 0.05
    #     stim_log_b = stim_log.loc[stim_log.stim_type == 'biphasic']
    #     isi_periods_per_sweep = []
    #     for s in np.unique(stim_log_b['sweep']):
    #         isi_list = []; behav_isi_list = []
    #         sub_df = stim_log_b.loc[(stim_log_b.sweep == s)]

    #         onsets = sub_df['onset'].values
    #         for ii, on in enumerate(onsets[:-1]):
    #             t1 = on + 1
    #             t2 = onsets[ii+1] - buffer
    #             isi_list.append([t1,t2])
    #             mr = np.mean(run_signal[(run_ts >= t1) & (run_ts < t2)])
    #             behav = -1
    #             for key, b_range in behavior_ranges.items():
    #                 if (mr >= b_range[0]) & (mr < b_range[1]):
    #                     behav = key
    #                     break
    #             behav_isi_list.append(key)
    #         behav_per_sweep.append(behav_isi_list)
    #         isi_periods_per_sweep.append(isi_list)
    # # import pdb; pdb.set_trace()

    ##------------------------------------------
    ## Calculate firing rates
    T, nNeurons = data_list[0].shape
    nCond = len(time_window_list)
    FR_spont = np.full((nCond,nNeurons,nBehaviors),np.nan)
    nBins_per_behav_spont = np.zeros((nCond,nBehaviors))

    for jj, tWindow in enumerate(time_window_list):
        ts = ts_list[jj]
        data = data_list[jj]
        behavior_map = behav_list[jj]
        spont_periods = spont_mask_list[jj]
  
        for key, behav_str in behavior_dict.items():
            indy = np.where((behavior_map == key) & (spont_periods == 1))[0]
            if len(indy) == 0:
                continue
            FR_spont[jj,:,key] = np.mean(data[indy],axis=0)/time_bin
            nBins_per_behav_spont[jj,key] = len(indy)


    #Save time windows
    np.savez(join(SaveDir,f'time_windows.npz'),time_window_list=time_window_list,filename_list=filename_list,nBins_per_behav_spont=nBins_per_behav_spont,areas=areas)

    #Create combinations that we are interested in comparing
    cbs = combinations(filename_list,2)
    ijs = combinations(np.arange(len(filename_list)),2)
    combos_spont = []
    for cb,ij in zip(cbs,ijs):
        epoch_i = cb[0];i = ij[0]
        epoch_j = cb[1];j = ij[1]
        if ('pre' in epoch_i) | ('post-sal1' in epoch_i) | ('post-ket' in epoch_i):
            if ('post' in epoch_j):
                print(ij,cb)
                combos_spont.append(ij)

    #Convert firing rate array into dataframe for easy plotting & saving
    for indy, fsuffix in zip([neuron_indices,np.arange(nNeurons)],[f'fr-thresh-{fr_thresh}','ALL']):
        #Compare FR changes due to saline or psilocybin injections vs no-injection
        tmp_list = []
        for j, tWindow in enumerate(time_window_list):
            epoch_j = filename_list[j]
            tW_j = tWindow[0] + np.diff(tWindow)[0]

            #Determine the time elapsed since 2nd injection or isoflurane induction
            if drug_type in ['saline', 'psilocybin']:
                tW_afterinj = (tW_j - injection_times[1])/60
            elif drug_type == 'isoflurane':
                tW_afterinj = (tW_j - iso_induction_times[1])/60
            else:
                tW_afterinj = tW_j

            blocklabel = epoch_j
            if ('post' in epoch_j):
                block_windows = np.arange(0,181,15)
                for l,r in zip(block_windows[:-1],block_windows[1:]):
                    if (tW_afterinj >= l) & (tW_afterinj < r):
                        blocklabel = f'{l:d}-{r:d}'
                        break

            print(f'\nEpoch: {filename_list[j]}')
            for key, behav_str in behavior_dict.items():
                #Select neurons FR for the epoch & behavior 
                FR_j = FR_spont[j,indy,key]
                N = FR_j.shape[0]

                t_behav = nBins_per_behav_spont[j,key]*time_bin
                #Check to make sure each epoch has 10 seconds of data
                bin_thresh = 10/time_bin
                if (nBins_per_behav_spont[j,key] < bin_thresh):
                    print(f'\t{behav_str:20s}: Not enough data to calculate FR')
                    continue
                else:
                    print(f'\t{behav_str:20s}: {t_behav:.1f} seconds of data to calculate FR')

                tmp_list.append(np.stack((np.repeat(mID,N),np.repeat(rec_name,N),indy,celltypes[indy], durations[indy], areas[indy],groups[indy],supergroups[indy],np.repeat(behav_str,N),np.repeat(epoch_j,N),np.repeat(tW_afterinj,N),np.repeat(blocklabel,N),FR_j)))
        FR_spont_df = pd.DataFrame(np.hstack(tmp_list).T,columns=['mID','rec_name','neuron_index','celltype','duration','area','group','supergroup','behavior','epoch','tW_after','blocklabel','FR'])
        FR_spont_df = FR_spont_df.astype({'mID':str,'rec_name':str,'neuron_index': int,'celltype':str,'duration':float,'area':str,'group':str,'supergroup':str,'behavior':str,'epoch':str,'tW_after':float,'blocklabel':str,'FR':float})
        FR_spont_df.to_hdf(join(SaveDir,f'spont_FR_{fsuffix}_{rec_name}.h5'),'df')

        if len(combos_spont) == 0:
            continue
        #Compare FR changes due to saline or psilocybin injections vs no-injection
        #Or Pre vs Post ISO induction
        tmp_list = []
        for cb in combos_spont:
            i = cb[0]
            j = cb[1]

            tW_i = time_window_list[i][0] + np.diff(time_window_list[i])[0]
            tW_j = time_window_list[j][0] + np.diff(time_window_list[j])[0]
            tW_diff = (tW_j - tW_i)/60 #Calculate how far away center of windows are from each other
            
            #Determine the time elapsed since 2nd injection or isoflurane induction
            if drug_type in ['saline', 'psilocybin']:
                tW_afterinj = (tW_j - injection_times[1])/60
            elif drug_type == 'isoflurane':
                tW_afterinj = (tW_j - iso_induction_times[1])/60
            else:
                tW_afterinj = tW_j

            block_windows = np.arange(0,181,15)
            
            if 'post-sal1' in filename_list[j]:
                blocklabel = 'post-sal1'
            else:
                blocklabel = 'error'
                for l,r in zip(block_windows[:-1],block_windows[1:]):
                    if (tW_afterinj >= l) & (tW_afterinj < r):
                        blocklabel = f'{l:d}-{r:d}'
                        break
            
            for key, behav_str in behavior_dict.items():
                #Select neurons FR for the epoch & behavior 
                FR_i = FR_spont[i,indy,key]; epoch_i = filename_list[i]
                FR_j = FR_spont[j,indy,key]; epoch_j = filename_list[j]

                #Check to make sure each epoch has 5 seconds of data
                bin_thresh = 10/time_bin
                if (nBins_per_behav_spont[i,key] < bin_thresh) | (nBins_per_behav_spont[j,key] < bin_thresh):
                    continue
                
                dFR = FR_j - FR_i
                MI = (FR_j - FR_i)/(FR_j + FR_i)
                N = MI.shape[0]

                tmp_list.append(np.stack((np.repeat(mID,N),np.repeat(rec_name,N),indy,celltypes[indy],durations[indy],areas[indy],groups[indy],supergroups[indy],np.repeat(behav_str,N),np.repeat(epoch_i,N),np.repeat(epoch_j,N),np.repeat(tW_diff,N),np.repeat(tW_afterinj,N),np.repeat(blocklabel,N),FR_i,FR_j,dFR,MI)))
            
        dFR_spont_df = pd.DataFrame(np.hstack(tmp_list).T,columns=['mID','rec_name','neuron_index','celltype','duration','area','group','supergroup','behavior','epoch_i','epoch_j','tW_diff','tW_after','blocklabel','FR_i','FR_j','dFR','MI'])
        dFR_spont_df = dFR_spont_df.astype({'mID':str,'rec_name':str,'neuron_index':int,'celltype':str,'duration':float,'area':str,'group':str,'supergroup':str,'behavior':str,'epoch_i':str, 'epoch_j':str,'tW_diff':float,'tW_after':float,'blocklabel':str,'FR_i':float,'FR_j':float,'dFR':float,'MI':float})
        dFR_spont_df.to_hdf(join(SaveDir,f'delta-spont_FR_{fsuffix}_{rec_name}.h5'),'df')

    #Create colormap for epochs
    nCond = len(filename_list)
    if drug_type == 'saline':
        cmap = sns.color_palette('Greens',nCond)
    elif drug_type == 'psilocybin':
        nPsi = np.sum(['psi' in fname for fname in filename_list])
        nSal = nCond - nPsi
        if nSal == 0:
            cmap = sns.color_palette('Reds',nCond)
        elif nPsi == 0:
            cmap = sns.color_palette('Blues',nCond)
        else:
            cmap = np.concatenate((sns.color_palette('Blues',nSal),sns.color_palette('Reds',nPsi)))
    elif drug_type == 'isoflurane':
        nPre = np.sum(['pre' in fname for fname in filename_list])
        nPost = nCond - nPre
        if nPre == 0:
            cmap = sns.color_palette('Oranges',nCond)
        elif nPost == 0:
            cmap = sns.color_palette('Blues',nCond)
        else:
            cmap = np.concatenate((sns.color_palette('Blues',nPre),sns.color_palette('Oranges',nPost)))

    elif drug_type == 'ketanserin+psilocybin':
        nPre = np.sum(['pre' in fname for fname in filename_list])
        nKet = np.sum(['ket' in fname for fname in filename_list])
        nPsi = np.sum(['psi' in fname for fname in filename_list])

        cmap = np.concatenate((sns.color_palette('Blues',nPre),sns.color_palette('Oranges',1),sns.color_palette('Reds',nPsi)))
    else:
        cmap = sns.color_palette('Oranges',nCond)

    #Plot FR per group
    uniq_groups = np.unique(FR_spont_df['group'])
    for ii, g in enumerate(uniq_groups):
        sub_df = FR_spont_df.loc[(FR_spont_df.group == g)]
        uniq_boxes = np.unique(sub_df['behavior'])

        fig, ax = plt.subplots(figsize=(10,4))
        ax.set_title(f'Group: {g}, n = {np.sum(groups == g)}')
        sns.boxplot(x='behavior',y='FR',hue='epoch',data=sub_df,order=uniq_boxes,hue_order=filename_list,palette=cmap,legend=False)
        ymax = np.nanpercentile(sub_df['FR'],97)
        tmp = sub_df['FR'].values
        mask = (~np.isnan(tmp)) & (~np.isinf(tmp))
        ymax = np.round(np.nanpercentile(tmp[mask],97.5))
        if np.isnan(ymax) | np.isinf(ymax):
            print('inf or nan values in data')
        else:
            ax.set_ylim([-2,ymax])

        usrplt.save_fig_to_pptx(fig, prs)
        plt.savefig(join(PlotDir,f'spont_firing_rate_across_epochs_group-{g}.pdf'),facecolor='white',dpi=300,bbox_inches='tight')
        plt.close(fig)

        uniq_areas = np.unique(sub_df['area'])
        uniq_boxes = np.unique(sub_df['behavior'])
        for a in uniq_areas:
            sub_df = FR_spont_df.loc[(FR_spont_df.group == g) & (FR_spont_df.area == a)]

            fig, ax = plt.subplots(figsize=(10,4))
            ax.set_title(f'Group: {g}, area {a}, n = {np.sum(areas == a)}')
            sns.boxplot(x='behavior',y='FR',hue='epoch',data=sub_df,order=uniq_boxes,hue_order=filename_list,palette=list(cmap),legend=False)
            ymax = np.nanpercentile(sub_df['FR'],97)
            tmp = sub_df['FR'].values
            mask = (~np.isnan(tmp)) & (~np.isinf(tmp))
            ymax = np.round(np.nanpercentile(tmp[mask],97.5))
            if np.isnan(ymax) | np.isinf(ymax):
                print('inf or nan values in data')
            else:
                ax.set_ylim([-2,ymax])
            usrplt.save_fig_to_pptx(fig, prs)

    #Plot change in firing rate for each possible combination that we are interested in
    dFR_spont_df_rest = dFR_spont_df.loc[dFR_spont_df.behavior == 'rest (0-1cm/s)']
    pval_thresh = 0.05
    for cb in combos_spont:
        i = cb[0]; epoch_i = filename_list[i]
        j = cb[1]; epoch_j = filename_list[j]
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f'"{epoch_i}" vs "{epoch_j}"'

        fig = plot_boxplot_comparison_spont(dFR_spont_df_rest, epoch_i, epoch_j, pval_thresh = 0.05)
        usrplt.save_fig_to_pptx(fig, prs,slide=slide)
        plt.savefig(join(PlotDir,f'delta-{epoch_i}_vs_{epoch_j}.pdf'),facecolor='white',dpi=300,bbox_inches='tight')
        plt.close(fig)

        for g in uniq_groups:
            sub_df = dFR_spont_df_rest.loc[(dFR_spont_df_rest.epoch_i == epoch_i) & (dFR_spont_df_rest.epoch_j == epoch_j) & (dFR_spont_df_rest.group == g)]
            if len(sub_df) < 5:
                continue
            fig =  plot_boxplot_comparison_spont(dFR_spont_df_rest, epoch_i,epoch_j, g = g, pval_thresh = 0.05)
            usrplt.save_fig_to_pptx(fig, prs)
            plt.close(fig)

    # Save the PowerPoint presentation
    prs.save(join(SaveDir,f'firing_rate_spont_figs_{rec_name}.pptx'))
    print('DONE!!!')
    print(calculate_evoked)
    if (not calculate_evoked) | (not stim_exists):
        exit()

    print('Calculate evoked firing rate changes')
    # Create a new PowerPoint presentation to save figures to
    prs2 = Presentation()
    slide = prs2.slides.add_slide(prs2.slide_layouts[0])
    slide.shapes.title.text = f'mID = {mID}; rec_name = {rec_name}'
    slide.placeholders[1].text = f'Firing rate changes due to psilocybin in evokedblocks'

    ##------------------ CALCULATE EVOKED FIRING RATES ------------------##
    def get_evoked_FRs(sweep, sub_df):
        
        #Get minimum inter-stimulus-interval
        ISI = sub_df['onset'].values[1:] - sub_df['offset'].values[:-1]
        min_ISI = np.round(np.min(ISI),1)

        #Get stimulation times
        event_times = sub_df['onset'].values
        stim_amp = sub_df['parameter'].values

        #Get evoked spike counts centered around stimulation times
        plot_before = 0;plot_after = min_ISI
        spk_list = []; time_bin_1ms = 1/1000
        for probei in probe_list:
            evoked_spks, bins = util.get_evoked_spike_counts(probe_unit_data[probei]['spike_times'], probe_unit_data[probei]['spike_clusters'], probe_unit_data[probei]['units'], event_times, plot_before, plot_after, time_bin_1ms)
            spk_list.append(evoked_spks)
            spk_ts_trial = bins[:-1] + time_bin_1ms/2
            
        evoked_ts_list = []
        for e in event_times:
            evoked_ts_list.append(e+spk_ts_trial)
        evoked_ts = np.array(evoked_ts_list)

        #Reshape
        spk_counts_evoked = np.concatenate(spk_list)
        spk_counts_evoked = np.transpose(spk_counts_evoked,[1,2,0])

        #Get running for evoked periods
        run_signal_evo, run_ts_evo, run_ts_trial = util.get_running_during_evoked_period(run_signal, run_ts, event_times, plot_before, plot_after)
        nTrials, _ = run_signal_evo.shape

        behavior_map = np.zeros((nTrials,len(windows_of_interest)))
        #Classify behavior for each window/trial
        for ii, tW in enumerate(windows_of_interest):
            if ii == 3:
                tW[1] = min_ISI
            indy = np.where((run_ts_trial >= tW[0]) & (run_ts_trial <= tW[1]))[0]
            for iTrial in range(nTrials):
                run_sub = run_signal_evo[iTrial,indy]
                mean_speed = np.mean(run_sub)
                for key, b_range in behavior_ranges.items():
                    if (mean_speed >= b_range[0]) & (mean_speed < b_range[1]):
                        behavior_map[iTrial,ii] = key

        fig, axes = plt.subplots(1,2,figsize=(6,5),gridspec_kw={'width_ratios':[15,1]})
        # plt.suptitle(evoked_filename_list[sweep])
        sns.heatmap(behavior_map,vmax=nBehaviors,vmin=0,ax=axes[0],cbar_ax=axes[1],cmap=sns.color_palette('Reds',nBehaviors),cbar_kws={'ticks':np.arange(nBehaviors)+0.5})

        ax = axes[0]
        ax.set_xlabel('Time window')
        ax.set_ylabel('Trial #')
        ax.set_xticklabels(window_strs,rotation=30)

        ax = axes[1]
        ax.set_yticklabels(behavior_strs)

        nTrials, nBins, nNeurons = spk_counts_evoked.shape
        #Sum spikes in each window for all neurons
        spk_counts_per_trial = np.zeros((nTrials,len(windows_of_interest),nNeurons))
        for ii, tW in enumerate(windows_of_interest):
            indy = np.where((spk_ts_trial >= tW[0]) & (spk_ts_trial <= tW[1]))[0]
            spk_counts_per_trial[:,ii] = np.sum(spk_counts_evoked[:,indy],axis=1)

        #Calculate mean firing rate for each window and possible behavior
        FR = np.full((nNeurons,nBehaviors+1,nWindows),np.nan)
        nTrials_per_behavior = np.zeros((nBehaviors,nWindows))
        for key, behav_str in behavior_dict.items():
            for ii, tW in enumerate(windows_of_interest):
                if ii == 3:
                    tW[1] = min_ISI
                trial_mask = behavior_map[:,ii] == key
                nTrials_per_behavior[key,ii] = np.sum(trial_mask)
                if np.sum(trial_mask) < 1:
                    continue
                tElapsed = np.diff(tW)
                FR[:,key,ii] = np.mean(spk_counts_per_trial[trial_mask,ii],axis=0)/tElapsed

        #Calculate firing rate irrespective of behavior
        for ii, tW in enumerate(windows_of_interest):
            if ii == 3:
                tW[1] = min_ISI
            tElapsed = np.diff(tW)
            FR[:,-1,ii] = np.mean(spk_counts_per_trial[:,ii],axis=0)/tElapsed
        return FR, nTrials_per_behavior, behavior_map, min_ISI, fig
    
    stim_log_biphasic = stim_log.loc[stim_log.stim_type == 'biphasic']
    # biphasic_indy = [ii for ii, e in enumerate(evoked_filename_list) if 'biphasic' in e]
    # evoked_filename_list = [evoked_filename_list[ii] for ii in biphasic_indy]
    # evoked_time_window_list = [evoked_time_window_list[ii] for ii in biphasic_indy]


    FR_df_list = []; FR_list = []
    nTrials_per_behavior_list = []
    behavior_list = []; min_ISI_list = []
    for sweep  in np.unique(stim_log_biphasic['sweep']):
        sub_df = stim_log.loc[(stim_log.stim_type == 'biphasic') & (stim_log['sweep'] == sweep)]
        FR, nTrials_per_behavior, behavior_map, min_ISI, fig = get_evoked_FRs(sweep,sub_df)
        FR_list.append(FR); min_ISI_list.append(min_ISI)
        behavior_list.append(behavior_map)
        nTrials_per_behavior_list.append(nTrials_per_behavior)
        usrplt.save_fig_to_pptx(fig, prs2) 
        plt.close(fig)

    behavior_strs = np.concatenate((behavior_strs,['ALL']))
    behavior_strs2 = np.concatenate((behavior_strs2,['All-trials']))
    
    print('Unique sweeps: ', np.unique(stim_log_biphasic['sweep']))
    print('evoked_filename_list: ', evoked_filename_list)
    #Get combinations we interested in comparing
    cbs = combinations(evoked_filename_list,2)
    ijs = combinations(np.arange(len(evoked_filename_list)),2)
    combos_evoked = []
    for cb,ij in zip(cbs,ijs):
        epoch_i = cb[0];i = ij[0]
        epoch_j = cb[1];j = ij[1]
        if ('pre' in epoch_i) | ('post-sal1' in epoch_i):
            if ('post' in epoch_j):
                print(ij,cb)
                combos_evoked.append(ij)
    if len(combos_evoked) == 0:
        skip_evoked_comparison = True
    else:
        skip_evoked_comparison = False

    #Convert firing rate array into dataframe for easy plotting & saving
    for indy, fsuffix in zip([neuron_indices,np.arange(nNeurons)],[f'fr-thresh-{fr_thresh}','ALL']):

        min_trials_behavior = 2
        tmp_list = []
        for i, FR in enumerate(FR_list):
            nTrials_per_behavior = nTrials_per_behavior_list[i]
            for j, w_str in enumerate(window_strs):
                for k, behav_str in behavior_dict.items():
                    
                    #Check to make sure each window has 2 trials 
                    if (nTrials_per_behavior[k,j] < min_trials_behavior):
                        continue

                    #Select neurons FR for the particular window + behavior
                    FR_sub = FR[indy,k,j]
                    N = FR_sub.shape[0]
                    tmp_list.append(np.stack((np.repeat(mID,N),np.repeat(rec_name,N),indy,celltypes[indy],durations[indy],areas[indy],groups[indy],supergroups[indy],np.repeat(behav_str,N),np.repeat(w_str,N),np.repeat(evoked_filename_list[i],N),FR_sub)))
                
                #Select neurons FR for the particular window + behavior
                FR_sub = FR[indy,-1,j]
                N = FR_sub.shape[0]
                tmp_list.append(np.stack((np.repeat(mID,N),np.repeat(rec_name,N),indy,celltypes[indy],durations[indy],areas[indy],groups[indy],supergroups[indy],np.repeat('ALL',N),np.repeat(w_str,N),np.repeat(evoked_filename_list[i],N),FR_sub)))

        FR_evoked_df = pd.DataFrame(np.hstack(tmp_list).T,columns=['mID','rec_name','neuron_index','celltype','duration','area','group','supergroup','behavior','window','epoch','FR'])
        FR_evoked_df = FR_evoked_df.astype({'mID':str,'rec_name':str,'neuron_index':int,'celltype':str,'duration':float,'area':str,'group':str,'supergroup':str,'behavior':str,'window':str,'epoch':str,'FR':float})
        FR_evoked_df.to_hdf(join(SaveDir,f'evoked_FR_{fsuffix}_{rec_name}.h5'),'df')
        
        if skip_evoked_comparison:
            continue
        #Loop over other epochs
        tmp_list = []
        for cb in combos_evoked:
            i = cb[0]
            j = cb[1]
            
            FR_i = FR_list[i]
            epoch_i = evoked_filename_list[i]
            nTrials_per_behavior_i = nTrials_per_behavior_list[i]

            FR_j = FR_list[j]
            epoch_j = evoked_filename_list[j]
            nTrials_per_behavior_j = nTrials_per_behavior_list[j]

            tW_i = evoked_time_window_list[i][0] + np.diff(evoked_time_window_list[i])[0]
            tW_j = evoked_time_window_list[j][0] + np.diff(evoked_time_window_list[j])[0]
            tW_diff = tW_j - tW_i #Calculate how far away center of windows are from each other

            #Determine the time elapsed since 2nd injection or isoflurane induction
            if drug_type in ['saline', 'psilocybin']:
                tW_afterinj = (tW_j - injection_times[1])/60
            elif drug_type == 'isoflurane':
                tW_afterinj = (tW_j - iso_induction_times[1])/60
            else:
                tW_afterinj = tW_j

            blocklabel = epoch_j
            if ('post' in epoch_j):
                block_windows = np.arange(0,181,15)
                for l,r in zip(block_windows[:-1],block_windows[1:]):
                    if (tW_afterinj >= l) & (tW_afterinj < r):
                        blocklabel = f'{l:d}-{r:d}'
                        break
  
            #Loop over behavior and window
            for k, behav_str in behavior_dict.items():
                for l, w_str in enumerate(window_strs):
                    if (nTrials_per_behavior_i[k,l] < min_trials_behavior) | (nTrials_per_behavior_j[k,l] < min_trials_behavior):
                        continue
                
                    #Select neurons FR for the window & behavior 
                    FR_i_sub = FR_i[indy,k,l] 
                    FR_j_sub = FR_j[indy,k,l] 

                    #Calculate change in firing rate
                    dFR = FR_j_sub - FR_i_sub
                    MI = (FR_j_sub - FR_i_sub)/(FR_j_sub + FR_i_sub)
                    N = MI.shape[0]

                    tmp_list.append(np.stack((np.repeat(mID,N),np.repeat(rec_name,N),indy,celltypes[indy],durations[indy],areas[indy],groups[indy],supergroups[indy],np.repeat(behav_str,N),np.repeat(w_str,N),np.repeat(epoch_i,N),np.repeat(epoch_j,N),np.repeat(tW_diff,N),np.repeat(tW_afterinj,N),np.repeat(blocklabel,N),FR_i_sub,FR_j_sub,dFR,MI)))
                    #Loop over window
            for l, w_str in enumerate(window_strs):
                #Select neurons FR for the window & behavior 
                FR_i_sub = FR_i[indy,-1,l] 
                FR_j_sub = FR_j[indy,-1,l] 

                #Calculate change in firing rate
                dFR = FR_j_sub - FR_i_sub
                MI = (FR_j_sub - FR_i_sub)/(FR_j_sub + FR_i_sub)
                N = MI.shape[0]

                tmp_list.append(np.stack((np.repeat(mID,N),np.repeat(rec_name,N),indy,celltypes[indy],durations[indy],areas[indy],groups[indy],supergroups[indy],np.repeat('ALL',N),np.repeat(w_str,N),np.repeat(epoch_i,N),np.repeat(epoch_j,N),np.repeat(tW_diff,N),np.repeat(tW_afterinj,N),np.repeat(blocklabel,N),FR_i_sub,FR_j_sub,dFR,MI)))

        dFR_evoked_df = pd.DataFrame(np.hstack(tmp_list).T,columns=['mID','rec_name','neuron_index','celltype','duration','area','group','supergroup','behavior','window','epoch_i','epoch_j','tW_diff','tW_afterinj','blocklabel','FR_i','FR_j','dFR','MI'])
        dFR_evoked_df = dFR_evoked_df.astype({'mID':str,'rec_name':str,'neuron_index':int,'celltype':str,'duration':float,'area':str,'group':str,'supergroup':str,'behavior':str,'window':str,'epoch_i':str, 'epoch_j':str,'tW_diff':float,'tW_afterinj':float,'blocklabel':str,'FR_i':float,'FR_j':float,'dFR':float,'MI':float})
        dFR_evoked_df.to_hdf(join(SaveDir,f'evoked_delta-FR_{fsuffix}_{rec_name}.h5'),'df')

    behavior_dict[4] = 'ALL'
    #Create colormap for epochs
    nCond = len(evoked_filename_list)
    if drug_type == 'saline':
        cmap = sns.color_palette('Greens',nCond)
    elif drug_type == 'psilocybin':
        nPsi = np.sum(['psi' in fname for fname in evoked_filename_list])
        nSal = nCond - nPsi
        if nSal == 0:
            cmap = sns.color_palette('Reds',nPsi)
        else:
            cmap = np.concatenate((sns.color_palette('Blues',nSal),sns.color_palette('Reds',nPsi)))
    elif drug_type == 'isoflurane':
        nPre = np.sum(['pre' in fname for fname in evoked_filename_list])
        nPost = nCond - nPre
        if nPre == 0:
            cmap = sns.color_palette('Oranges',nPost)
        else:
            cmap = np.concatenate((sns.color_palette('Blues',nPre),sns.color_palette('Oranges',nPost)))
    else:
        cmap = sns.color_palette('Oranges',nCond)

    uniq_groups = np.unique(FR_evoked_df['group'])
    for g in uniq_groups:
        for jj, b in enumerate(behavior_strs):
            b2 = behavior_strs2[jj]
            sub_df = FR_evoked_df.loc[(FR_evoked_df['group'] == g) & (FR_evoked_df['behavior'] == b)]

            fig, ax = plt.subplots(figsize = (10,4))
            ax.set_title(f'Group: {g}, n = {np.sum(groups == g)}, behavior = {b}')
            sns.boxplot(x='window',y='FR',hue='epoch',palette=cmap,data=sub_df,legend=False)

            ymax = np.nanpercentile(sub_df['FR'],97)
            tmp = sub_df['FR'].values
            mask = (~np.isnan(tmp)) & (~np.isinf(tmp))
            ymin = np.round(np.nanpercentile(tmp[mask],2.5))
            ymax = np.round(np.nanpercentile(tmp[mask],97.5))
            if np.isnan(ymax) | np.isinf(ymax):
                print('inf or nan values in data')
            else:
                ax.set_ylim([-2,ymax])
  
            usrplt.save_fig_to_pptx(fig, prs2)
            plt.savefig(join(PlotDir,f'evoked_firing_rate_across_epochs_group-{g}_{b2}.pdf'),facecolor='white',dpi=300,bbox_inches='tight')
            plt.close(fig)

    if skip_evoked_comparison:
        print('No evoked firing rate comparisons to be had')
        prs2.save(join(SaveDir,f'firing_rate_evoked_figs_{rec_name}.pptx'))
        print('DONE!!!')
        exit()

    for cb in combos_evoked:
        i = cb[0]
        j = cb[1]
            
        FR_i = FR_list[i]
        epoch_i = evoked_filename_list[i]
        nTrials_per_behavior_i = nTrials_per_behavior_list[i]

        FR_j = FR_list[j]
        epoch_j = evoked_filename_list[j]
        nTrials_per_behavior_j = nTrials_per_behavior_list[j]
        for k, b_str in behavior_dict.items():
            sub_df = dFR_evoked_df.loc[(dFR_evoked_df.epoch_i == epoch_i) & (dFR_evoked_df.epoch_j == epoch_j) & (dFR_evoked_df.behavior == b_str)]
            if len(sub_df) < 5:
                continue

            fig =  plot_boxplot_comparison_evoked(dFR_evoked_df, epoch_i,epoch_j, b_str, g = None, pval_thresh = 0.05)
            usrplt.save_fig_to_pptx(fig, prs2)
            plt.close(fig)

        # for g in uniq_groups:
        #     for k, b_str in behavior_dict.items():


        #         sub_df = dFR_evoked_df.loc[(dFR_evoked_df.epoch_i == epoch_i) & (dFR_evoked_df.epoch_j == epoch_j) & (dFR_evoked_df.group == g) & (dFR_evoked_df.behavior == b_str)]
        #         if len(sub_df) < 5:
        #             continue

        #         fig =  plot_boxplot_comparison_evoked(dFR_evoked_df, epoch_i,epoch_j, b_str, g = g, pval_thresh = 0.05)
        #         usrplt.save_fig_to_pptx(fig, prs2)
        #         plt.close(fig)

    prs.save(join(SaveDir,f'firing_rate_evoked_figs_{rec_name}.pptx'))
    print('DONE!!!')

