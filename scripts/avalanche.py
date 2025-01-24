#Basea
import gspread
from os.path import join
from glob import glob
import json, os, time, sys, argparse
import numpy as np
import xarray as xr
import pandas as pd
import scipy.stats as st
from tqdm.notebook import trange, tqdm
import ray
import pingouin as pg
from scipy.ndimage import gaussian_filter
from scipy.signal import decimate
import io

#Plot
PlotDir = '/home/david.wyrick/projects/zap-n-zip/plots'
import matplotlib as mpl
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
from ssm.plots import gradient_cmap, white_to_color_cmap
from matplotlib.gridspec import GridSpec

#Project
from tbd_eeg.data_analysis.eegutils import EEGexp
import tbd_eeg.data_analysis.Utilities.utilities as util
from tbd_eeg.data_analysis.Utilities.behavior_movies import Movie
from tbd_eeg.data_analysis.Utilities.utilities import find_nearest_ind

#SSM
import ssm
from sklearn.model_selection import StratifiedKFold

#Allen
from allensdk.brain_observatory.ecephys.lfp_subsampling.subsampling import remove_lfp_offset
from allensdk.core.mouse_connectivity_cache import MouseConnectivityCache

mcc = MouseConnectivityCache(resolution=10)
str_tree = mcc.get_structure_tree()

#Read in allen CCF
ccfsum = pd.read_csv('/home/david.wyrick/projects/zap-n-zip/data/mouse_connectivity/ontology_v2.csv')

#User
from util import *
from plotting import *

#Single cell stuff
import elephant.statistics as est
from quantities import ms, s, Hz
from neo.core import SpikeTrain

base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
base_dir_local = '/data/tiny-blue-dot/zap-n-zip/EEG_exp' #local
base_dir = base_dir_server
ProjDir = '/home/david.wyrick/projects/zap-n-zip/'
ServDir = '/data/projects/zap-n-zip/'

#Define behavioral states
behavior_dict = {0: 'rest (0-1cm/s)', 1: 'walk (1-5cm/s)',2: 'slow run (5-20cm/s)',3:'run (20-50cm/s)', 4: 'fast run (>50cm/s)'}
behavior_strs = list(behavior_dict.values())
nBehaviors = len(behavior_strs)

#Define windows to calculate firing rate
windows_of_interest = [[.002,.025],[.075,.3],[.3,1],[1,4.5]]
window_strs = ['evoked (2-25ms)','rebound (75-300ms)','post-rebound (0.3-1s)','ISI (1-4.5s)']
nWindows = len(window_strs)

gc = gspread.service_account() # need a key file to access the account (step 2) 
sh = gc.open('Templeton-log_exp') # open the spreadsheet 

df = pd.DataFrame(sh.sheet1.get()) 
exp_table = df.T.set_index(0).T # put it in a nicely formatted dataframeexp_table.iloc[10:]
exp_table = exp_table.iloc[9:].set_index('mouse_name')

##===== ============================ =====##
##===== Parse Command Line Arguments =====##
parser = argparse.ArgumentParser(description='single-cell-metrics')

##===== Data Options =====##
parser.add_argument('--mouseID',type=str, default='mouse669117',
                    help='mouse to perform analysis on')

parser.add_argument('--rec_name',type=str, default='pilot_aw_psi_2023-03-30_11-37-07',
                    help='experiment to perform analysis on')

parser.add_argument('--time_bin_ms',type=int, default=100,
                    help='time_bin_ms')

parser.add_argument('--fr_thresh',type=float, default=2,
                    help='Firing rate threshold for neurons to include in analysis')

parser.add_argument('--calculate_evoked',type=int, default=1,
                    help='Firing rate threshold for neurons to include in analysis')

if __name__ == '__main__':

    # Parse the arguments
    args = parser.parse_args()
    mouse_name = args.mouseID 
    mID = args.mouseID
    rec_name = args.rec_name
    fr_thresh = args.fr_thresh
    time_bin_ms = args.time_bin_ms
    time_bin = time_bin_ms/1000
    calculate_evoked = bool(args.calculate_evoked)
    
    # mID = 'mouse669117'; rec_name = 'pilot_aw_2023-03-29_11-09-15'
    # mID = 'mouse669117'; rec_name = 'pilot_aw_psi_2023-03-30_11-37-07'
    # mID = 'mouse673449'; rec_name = 'aw_psi_2023-04-19_11-23-26'
    # mID = 'mouse673449'; rec_name = 'aw_psi_d2_2023-04-20_10-05-31'

    #Define directories
    SaveDir = join(ProjDir,'results','FR',mID,rec_name)
    if not os.path.exists(SaveDir):
        os.makedirs(SaveDir)

    PlotDir = join(SaveDir,'plots')
    if not os.path.exists(PlotDir):
        os.makedirs(PlotDir)

    #Get injection times and types of injection
    injection_times = [float(exp_table.loc[exp_table.exp_name == rec_name]['First injection time'].values[0]),
                    float(exp_table.loc[exp_table.exp_name == rec_name]['Second injection time'].values[0])]

    if 'psi' in rec_name:
        injection_types = ['sal','psi']
        cmap = np.concatenate((sns.color_palette('Blues',2),sns.color_palette('Reds',1)))
    else:
        injection_types = ['sal1', 'sal2']
        cmap = sns.color_palette('Blues',3)

    #%% Upload the whole experiment and generate the global clock
    file_name = os.path.join(base_dir_server,mID,rec_name,'experiment1','recording1')
    exp = EEGexp(file_name, preprocess=False, make_stim_csv=False)

    #Load probe data
    probe_unit_data, probe_info, total_units = util.get_neuropixel_data(exp)
    probe_list = [x.replace('_sorted', '') for x in exp.experiment_data if 'probe' in x]

    #Get recording start time
    probei = probe_list[0]
    open_ephys_start = np.round(np.min(probe_unit_data[probei]['spike_times']))
    open_ephys_end = np.round(np.max(probe_unit_data[probei]['spike_times']))

    # extract the timestamps of the selected stimuli
    stim_log = pd.read_csv(exp.stimulus_log_file)
    
    #Determine time periods of evoked activity
    time_window_list = []
    evoked_filename_list = []
    nSweeps = len(np.unique(stim_log['sweep']))
    for s in np.unique(stim_log['sweep']):
        sub_df = stim_log.loc[stim_log.sweep == s]
        tS = sub_df['onset'].min(); tE = sub_df['offset'].max()
        
        time_window_list.append([tS,tE])
        if tS < injection_times[0]:
            evoked_filename_list.append(f'sweep-{s:02d}_pre-inj')
        elif (tS > injection_times[0]) & (tS < injection_times[1]):
            evoked_filename_list.append(f'sweep-{s:02d}_post-{injection_types[0]}-inj')
        else:
            evoked_filename_list.append(f'sweep-{s:02d}_post-{injection_types[1]}-inj')

    #Get sweep times
    sweep_time_windows = np.array(time_window_list)

    #Get inter-sweep windows 
    inter_sweep_windows = np.stack((sweep_time_windows[:-1,1],sweep_time_windows[1:,0])).T

    #Add period before sweep 0
    inter_sweep_windows = np.vstack(([open_ephys_start,sweep_time_windows[0,0]],inter_sweep_windows))

    #Change second spont epoch to just after 1st injection to before second injection
    inter_sweep_windows[1,0] = injection_times[0] + 30
    inter_sweep_windows[1,1] = injection_times[1] - 30

    #Define periods of spontaneous activity
    spont_filename_list = []
    for s, tW in enumerate(inter_sweep_windows):

        if tW[0] < injection_times[0]:
            spont_filename_list.append(f'spont-{s:02d}_pre-inj')
        elif (tW[0] >= injection_times[0]) & (tW[0] < injection_times[1]):
            spont_filename_list.append(f'spont-{s:02d}_post-{injection_types[0]}-inj')
        else:
            spont_filename_list.append(f'spont-{s:02d}_post-{injection_types[1]}-inj')
    
    # Load running speed
    run_file = os.path.join(base_dir,mID,rec_name,'experiment1','recording1','running_signal.npy')
    if os.path.exists(run_file):
        # print('File exists')
        run_signal = np.load(run_file)

        ts_file = os.path.join(base_dir,mID,rec_name,'experiment1','recording1','running_timestamps.npy')
        if os.path.exists(ts_file):
            run_ts = np.load(ts_file)
        else:
            ts_file = os.path.join(base_dir,mID,rec_name,'experiment1','recording1','running_timestamps_master_clock.npy')
            run_ts = np.load(ts_file)
    else:
        print('\tRunning file does not exist')
        run_signal, run_ts = exp.load_running()
        np.save(os.path.join(base_dir,mID,rec_name,'experiment1','recording1','running_signal.npy'),run_signal)
        np.save(os.path.join(base_dir,mID,rec_name,'experiment1','recording1','running_timestamps.npy'),run_ts)

    #Load pupil data
    pupil_csv = os.path.join(base_dir,mID,rec_name,'experiment1','recording1',f'Pupil_{rec_name}.csv')

    try:
        table = pd.read_csv(pupil_csv)
        pupil_radius = table['Largest_Radius'].values

        #Pupil master clock
        pupil_ts = Movie(filepath=exp.pupilmovie_file,
                        sync_filepath=exp.sync_file,
                        sync_channel='eyetracking'
                        ).sync_timestamps
        plot_pupil = True
    except:
        pupil_ts = np.array([np.nan])
        pupil_radius =np.array([np.nan])
        plot_pupil = False

    #Smooth running signal with gaussian 
    run_signal_smoothed = gaussian_filter(run_signal,10)
    acceleration_smoothed = np.gradient(run_signal_smoothed)

    if 'psi' in rec_name:
        cmap = np.concatenate((sns.color_palette('Blues',2),sns.color_palette('Reds',1)))
    else:
        cmap = sns.color_palette('Blues',3)

    # Create a new PowerPoint presentation to save figures to
    prs = Presentation()

    #Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f'mID = {mID}; rec_name = {rec_name}'
    slide.placeholders[1].text = f'Firing rate changes due to psilocybin in evoked & spontaneous blocks'

    #Plot behavior
    fig = plt.figure(figsize=(9,6))
    plt.suptitle(f'{mID}, {rec_name}')
    gs = GridSpec(2, 3, figure=fig,height_ratios=[2,1.5])

    ax = fig.add_subplot(gs[1, :])
    ax.plot(run_ts/60,run_signal_smoothed,'-k',lw=1)
    ax.set_xlabel('Time (mins)')
    ax.set_ylabel('Speed (cm/s)')
    ax.autoscale(tight=True)

    # ax = fig.add_subplot(gs[2, :])
    ax.vlines(injection_times[0]/60,*ax.get_ylim(),color=cmap[1],lw=3,zorder=4)
    ax.vlines(injection_times[1]/60,*ax.get_ylim(),color=cmap[2],lw=3,zorder=4)

    for ii, tW in enumerate(sweep_time_windows):
        ax.vlines(tW/60,*ax.get_ylim(),color=cc[0],lw=1.5,zorder=4)

    axes = []
    for ii, (ts, sig) in enumerate(zip([run_ts,run_ts,pupil_ts],[run_signal,acceleration_smoothed,pupil_radius])):
        ax = fig.add_subplot(gs[0, ii]); axes.append(ax)
        ax.set_title(f'{mID}, {rec_name}')
        indy = np.where((ts <= injection_times[0]))[0]
        sns.kdeplot(sig[indy],ax=ax,label='No-injection period',color=cmap[0],lw=1.5,common_norm=False)

        indy = np.where((ts >= injection_times[0]) & (ts <= injection_times[1]))[0]
        sns.kdeplot(sig[indy],ax=ax,label=f'After {injection_types[0]} injection',color=cmap[1],lw=1.5,common_norm=False)

        indy = np.where((ts >= injection_times[1]))[0]
        sns.kdeplot(sig[indy],ax=ax,label=f'After {injection_types[1]} injection',color=cmap[2],lw=1.5,common_norm=False)
  
    axes[0].set_title('Running speed')
    axes[1].set_title('Acceleration');axes[1].set_xlim([-0.5,0.5]); axes[1].legend()
    axes[2].set_title('Pupil radius'); axes[2].set_xlim([0,140])

    save_fig_to_pptx(fig, prs)
    plt.savefig(join(PlotDir,f'behavior_{rec_name}.pdf'))

    #Read in neuropixel data
    spont_data_list, spont_ts_list, neuron_indices, plot_tuple = get_neuropixel_data(probe_unit_data,inter_sweep_windows, time_bin=0.1,fr_thresh=fr_thresh)
    boundaries, yticks, labels, areas, groups, order_by_group = plot_tuple

    areas_sub = areas[neuron_indices]
    groups_sub = groups[neuron_indices]
    N = len(neuron_indices)

    #Create new slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Lets plot some raster plots for fun'

    tf = body_shape.text_frame
    tf.text = '30seconds is being shown'

    p = tf.add_paragraph()
    p.text = f'The grayscale has a max of 50Hz'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Rasters are organized by probe'
    p.level = 1

    #Plot rasters for each epoch
    jj = 0
    neurons_per_probe = [len(np.unique(probe_unit_data[probei]['units'])) for probei in probe_list]
    neurons_per_probe_cum = np.cumsum(neurons_per_probe)
    for jj, tWindow in enumerate(inter_sweep_windows):
        ts = spont_ts_list[jj]
        spikes = spont_data_list[jj]
        T, N = spikes.shape
        
        # Display a spike raster of the image
        tmp = np.concatenate((neurons_per_probe,[150]))
        fig, axes =  plt.subplots(len(tmp),1,figsize=(12,12),gridspec_kw={'height_ratios':np.array(tmp)/np.sum(tmp)})
        plt.suptitle(spont_filename_list[jj],y=0.925)

        #Define time window
        if T*time_bin < 60:
            tStart = 0
        else:
            tStart = 30# 5*60
        time_to_plot = 30 #sec
        tslice = slice(int(tStart/time_bin),int((tStart+time_to_plot)/time_bin))
        ts_sub = ts[tslice]

        #Params
        xticks = np.arange(0,int((time_to_plot+1)/time_bin),int(15/time_bin))
        tlabel = xticks*time_bin
        xticks = np.arange(ts_sub[0],ts_sub[-1]+1,15)

        for ii, probei in enumerate(probe_list):
            ax = axes[ii]; ax.set_title(probei)
            if ii == 0:
                nslice = slice(0,neurons_per_probe[0])
            else:
                nslice = slice(neurons_per_probe_cum[ii-1],neurons_per_probe_cum[ii])
            print(np.percentile(spikes[tslice,nslice].T,97))
            ax.imshow(spikes[tslice,nslice].T, aspect='auto',vmax=5,vmin=0, cmap='gray_r')
            ax.set_xticks([])

            uniq_groups, uniq_indices, num_groups = np.unique(groups[nslice],return_index=True, return_counts=True)
            group_labels = uniq_groups[np.argsort(uniq_indices)]
            nNeurons_area = num_groups[np.argsort(uniq_indices)]

            boundaries = np.concatenate(([0],np.cumsum(nNeurons_area)))
            yticks = boundaries[:-1] + np.diff(boundaries)/2

            ax.set_yticks(yticks,minor=True)
            ax.set_yticks(boundaries,minor=False)
            ax.set_yticklabels(group_labels,minor=True)
            ax.set_yticklabels([],minor=False)
            ax.set_ylabel("Unit")

        #Plot pupil and running
        indy = np.where((run_ts >= ts_sub[0]) & (run_ts <= ts_sub[-1]))[0]
        ax = axes[-1]; ax.autoscale(tight=True)
        ax.plot(run_ts[indy],run_signal[indy],lw=0.6,color='k')
        ax.set_ylabel("Speed (cm/s)"); ax.set_xlabel("Time (s)")
        # ax.set_xticks([]); 
        ax.set_title('Behavior')

        if plot_pupil:
            indy = np.where((pupil_ts >= ts_sub[0]) & (pupil_ts <= ts_sub[-1]))[0]
            ax = ax.twinx()
            ax.plot(pupil_ts[indy],pupil_radius[indy],lw=0.6,color=cc[1])
            ax.set_ylabel("Pupil radius (pix)") 

        save_fig_to_pptx(fig, prs)
        plt.savefig(os.path.join(PlotDir,f'raster_{spont_filename_list[jj]}.pdf'),facecolor='white',dpi=300,bbox_inches='tight')

    #Look at spontaneous periods first
    #Smooth running signal with gaussian 
    run_signal_smoothed = gaussian_filter(run_signal,10)
    f_run = interp1d(run_ts,run_signal_smoothed)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    num_slides_per_slide = 8
    #Determine behaviors during spontaneous periods
    spont_behav_list = []
    for jj, tWindow in enumerate(inter_sweep_windows):
        if jj > 7:
            if jj == num_slides_per_slide:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
            fig_num = jj - 8
        else:
            fig_num = jj
        position = [(fig_num//4)*5, (fig_num%4)*1.875]
        
        ts = spont_ts_list[jj]
        T = len(ts)
        tMin = T*time_bin/60
        run_spont = f_run(ts)
        behavior_map = np.zeros(run_spont.shape)
        
        indy = np.where(run_spont < 1)[0]
        behavior_map[indy] = 0

        indy = np.where((run_spont >= 1) & (run_spont < 5))[0]
        behavior_map[indy] = 1

        indy = np.where((run_spont >= 5) & (run_spont < 20))[0]
        behavior_map[indy] = 2

        indy = np.where((run_spont >= 20) & (run_spont < 50))[0]
        behavior_map[indy] = 3

        indy = np.where((run_spont >= 50))[0]
        behavior_map[indy] = 4

        print(np.unique(behavior_map,return_counts=True))
        spont_behav_list.append(behavior_map)

        fig, ax = plt.subplots(figsize=(5,1))
        ax.set_title(spont_filename_list[jj])
        ax2 = ax.twinx()
        sns.heatmap(behavior_map.reshape(1,-1),vmax=4,vmin=0,ax=ax,cbar=False,cmap=sns.color_palette('Reds',5),cbar_kws={'ticks':np.arange(5)})
        
        ax2.plot(run_spont,color='k')
        ax2.set_ylabel('Speed (cm/s)')
        ax.set_yticks([])
        ax.set_xticks([0,T/2,T]); ax.set_xticklabels([0,np.round(tMin/2,2),np.round(tMin,2)],rotation=0);ax.set_xlabel('Time (mins)')

        save_fig_to_pptx(fig, prs,slide=slide,position=position)
        # plt.savefig(join(PlotDir,f'behavior_{spont_filename_list[jj]}.pdf'),facecolor='white',dpi=300,bbox_inches='tight')

    #Create slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Lets look avalanches per group'

    tf = body_shape.text_frame
    tf.text = 'Now the x-axis will be "group"'

    p = tf.add_paragraph()
    p.text = f'This is being compared to either "{spont_filename_list[0]}" or "{spont_filename_list[1]}"'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Note that multiple areas may be contained within a group'
    p.level = 1

    #Read in neuropixel data
    spont_data_list, spont_ts_list, neuron_indices, plot_tuple = get_neuropixel_data(probe_unit_data,inter_sweep_windows, time_bin=0.005,fr_thresh=fr_thresh)
    boundaries, yticks, labels, areas, groups, order_by_group = plot_tuple

    # Save the PowerPoint presentation
    prs.save(join(PlotDir,f'firing_rate_spont_figs_{rec_name}.pptx'))
    print('DONE!!!')


