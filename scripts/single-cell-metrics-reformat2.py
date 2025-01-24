base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
base_dir_local = '/data/tiny-blue-dot/zap-n-zip/EEG_exp'; base_dir = base_dir_server
ProjDir = '/home/david.wyrick/projects/zap-n-zip/'
ServDir = '/data/projects/zap-n-zip/'
BehavDir = '/home/david.wyrick/projects/zap-n-zip/plots/behavior'

#Basea
import gspread
from os.path import join
from glob import glob
import json, os, time, sys, argparse
import numpy as np
import pandas as pd
import scipy.stats as st
import scipy.signal as sig
import pingouin as pg
from scipy.ndimage import gaussian_filter
from scipy.interpolate import interp1d
from itertools import combinations
from tqdm import tqdm
from scipy.optimize import curve_fit

#Plot
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
from matplotlib.gridspec import GridSpec
from matplotlib.lines import Line2D   

#Project
from tbd_eeg.data_analysis.eegutils import EEGexp
import tbd_eeg.data_analysis.Utilities.utilities as tbd_util
from tbd_eeg.data_analysis.Utilities.behavior_movies import Movie

#Allen
from allensdk.core.mouse_connectivity_cache import MouseConnectivityCache
mcc = MouseConnectivityCache(resolution=10)
str_tree = mcc.get_structure_tree()

#Read in allen CCF
ccfsum = pd.read_csv('/home/david.wyrick/projects/zap-n-zip/data/mouse_connectivity/ontology_v2.csv')

#User
sys.path.append(ProjDir)
import util
import plotting as usrplt

#Single cell stuff
import elephant.statistics as est
from quantities import ms, s, Hz
from neo.core import SpikeTrain

#Define behavioral states
# behavior_ranges = {0: [0,1], 1: [1,10], 2: [10,30], 3:[30,500]}
# behavior_dict = {0: 'rest (0-1cm/s)', 1: 'walk (1-10cm/s)', 2: 'shuffle (10-30cm/s)', 3: 'run (>30cm/s)'}
# behavior_strs2 = ['rest','walk','shuffle','run']

behavior_ranges = {0: [0,1], 1: [1,500]}#, 3:[30,500]}
behavior_dict = {0: 'rest (0-1cm/s)', 1: 'active (>1cm/s)'}
behavior_strs2 = ['rest','active']
behavior_dict2 = { 0: 'rest (<1cm/s)', 1: 'active (>1cm/s)', 2: 'all (>0cm/s)'}
behavior_strs = list(behavior_dict.values())
nBehaviors = len(behavior_strs)

behavior_strs = list(behavior_dict.values())
nBehaviors = len(behavior_strs)

#Define windows to calculate firing rate
windows_of_interest = [[.002,.025],[.075,.3],[.3,1],[1,4.5]]
window_strs = ['evoked (2-25ms)','rebound (75-300ms)','post-rebound (0.3-1s)','ISI (1-4.5s)']
nWindows = len(window_strs)

gc = gspread.service_account() # need a key file to access the account (step 2) 
sh = gc.open('Templeton-log_exp') # open the spreadsheet 

df = pd.DataFrame(sh.sheet1.get()) 
exp_table = df.T.set_index(0).T # put it in a nicely formatted dataframeexp_table.iloc[10:]
exp_table = exp_table.set_index('mouse_name')

##===== ============================ =====##
##===== Parse Command Line Arguments =====##
parser = argparse.ArgumentParser(description='single-cell-metrics')

##===== Data Options =====##
parser.add_argument('--mID',type=str, default='mouse728449',
                    help='mouse to perform analysis on')

parser.add_argument('--rec_name',type=str, default='aw_sal_2024-04-10_10-04-50',
                    help='experiment to perform analysis on')

parser.add_argument('--time_bin_ms',type=int, default=50,
                    help='time_bin_ms')

parser.add_argument('--fr_thresh',type=float, default=0,
                    help='Firing rate threshold for neurons to include in analysis')

parser.add_argument('--calculate_evoked',type=int, default=0,
                    help='blah')

parser.add_argument('--window_t_min',type=int, default=1,
                    help='Window width (mins) to segment data into')

#Define windows to calculate firing rate
stim_window_dict = {'spontaneous': [], 'evoked': [],'pre-rebound': [], 'rebound': [], 'post-rebound': [], 'isi': [],'visual': []}
stim_strs = ['spontaneous','evoked','pre-rebound','rebound','post-rebound','isi_biphasic','visual','isi_visual']

visual_stim = ['circle','natural_scene']
spont_types = ['spontaneous','isi_biphasic','isi_visual']
evoked_windows = [[.002,.025],[0.025,0.075],[.075,.3],[.3,1]]
evoked_strings = ['evoked','pre-rebound','rebound','post-rebound']

import warnings
warnings.filterwarnings("ignore")
def calculate_gain(fr, plot=True, titlestr=None):

    nPoints = fr.shape
    mFR = np.max(fr)

    #Build quantile-quantile plot between distribution of firing rates and assumed normal distribution of inputs
    #draw gaussian distribution of input currents
    input_currents = np.array(sorted(np.random.normal(size=nPoints)))
    fr_sorted = np.array(sorted(fr))

    if all(fr_sorted == 0):
        return [], [], 1, None
    
    #Calculate initial estimate of inflection point
    tmp = np.where(fr_sorted > 0)[0][0]
    pos_x = tmp + int((np.argmax(fr_sorted) - tmp)/2)
    h0 = input_currents[pos_x]

    #Fit sigmoid
    params_sigmoid, _ = curve_fit(util.sigmoid,input_currents,fr_sorted,p0=[2,h0,np.max(fr_sorted)],method='lm',maxfev=1*10**9)

    #Evaluate Fit
    fr_sigmoid = util.sigmoid(input_currents,*params_sigmoid)
    ss_res_sigmoid = np.sum((fr_sorted - fr_sigmoid)**2)
    ss_tot = np.sum((fr_sorted - np.mean(fr_sorted))**2)
    r2_sigmoid = 1 - ss_res_sigmoid/ss_tot
    x = np.arange(-4,5,0.01); h0 = params_sigmoid[1]; m = params_sigmoid[0]*params_sigmoid[2]/4
    
    if plot:

        fig, axes = plt.subplots(1,2,figsize=(8,4));plt.autoscale(enable=True,tight=True)
        if titlestr is not None:
            plt.suptitle(titlestr,y=0.995)

        #Plot histogram of firing rates
        ax = axes[0]
        ax.set_title('Firing rate distribution')
        ax.hist(fr,histtype='step',color=usrplt.cc[8],LineWidth=3,density=True)
        sns.kdeplot(util.sigmoid(input_currents,*params_sigmoid),ax=ax,color=usrplt.cc[2],lw=4,zorder=3,alpha=0.75)
        ax.vlines(np.mean(fr),*ax.get_ylim(),lw=3,ls='--',colors=usrplt.cc[8])
        ax.set_ylabel('Prob. Density');ax.set_xlim([0,mFR])
        ax.set_xlabel(f'Firing rate')

        #plot relationship between distribution of firing rates and assumed normal distribution of inputs
        ax = axes[1]; ax.set_title('Transfer curve')
        ax.plot(input_currents,fr_sorted,'o',color=usrplt.cc[8],fillstyle='none',ms=7)
        ax.set_ylabel('Firing rate (Hz)')
        ax.set_xlabel('Input Current (au)')

        #Plot sigmoid fit
        ax.plot(x,util.sigmoid(x,*params_sigmoid),color=usrplt.cc[2],lw=4,zorder=3,alpha=0.75, label='Sigmoid -> \u03B2: {:.2f} m_h\u2080: {:.2f} r2: {:.2f}'.format(params_sigmoid[0],m,r2_sigmoid))
        ax.plot([h0,h0],[0,util.sigmoid(h0,*params_sigmoid)],color=usrplt.cc[3],ls='--',lw=3,zorder=2)
        ax.plot(x,m*(x-h0)+util.sigmoid(h0,*params_sigmoid),color=usrplt.cc[5],ls='--',lw=3,zorder=2)
        ax.set_ylim([0,10*np.ceil(mFR/10)])
    else:
        fig = None

    ##===== Fit RELU =====##
    pos_x = np.where(fr_sorted > 0)[0][0]
    p0 = [4,input_currents[pos_x]] 
    params_relu, _ = curve_fit(util.relu,input_currents,fr_sorted,p0=p0,method='lm',maxfev=1*10**9)

    #Evaluate Fit
    fr_relu = util.relu(input_currents,*params_relu)
    ss_res_relu = np.sum((fr_sorted - fr_relu)**2)
    r2_relu = 1 - ss_res_relu/ss_tot

    if params_sigmoid[0] > 100:
        params_sigmoid[0] = np.nan
        m = np.nan
    
    if plot:
        #RELU fit
        ax = axes[1]
        ax.plot(x,util.relu(x,*params_relu),'k',linestyle='-',label='Relu -> m: {:.2f}, x0: {:.2f} -> r2: {:.2f}'.format(*params_relu,r2_relu))
        ax.legend(loc=2)
        
    ##===== Save Parameters for summary statistics =====##
    relu_params = [r2_relu,*params_relu]
    sigm_params = [r2_sigmoid,*params_sigmoid,m]
    
    return relu_params, sigm_params, 0, fig

if __name__ == '__main__':

    ##------------------------------------------
    # Parse the arguments
    args = parser.parse_args()
    mID = args.mID
    rec_name = args.rec_name
    fr_thresh = args.fr_thresh
    time_bin_ms = args.time_bin_ms
    time_bin = time_bin_ms/1000
    calculate_evoked = bool(args.calculate_evoked)

    #Segement data into X minute windows
    window_t_min = args.window_t_min
    window_t = window_t_min*60
    
    #Extract row from templeton dataframe
    exp_df = exp_table.loc[exp_table.exp_name == rec_name]
    if type(exp_df) == pd.core.series.Series:
        exp_df = exp_df.to_frame().T

    drug_type = exp_df['drug'].values[0]
    stim_type = exp_df['stimulation'].values[0]
    print(f'Experiment type: {stim_type}, {drug_type}')

    #Define directories4
    SaveDir = join(ServDir,'results','FR_reformat4',mID,rec_name)
    if not os.path.exists(SaveDir):
        os.makedirs(SaveDir)

    PlotDir = join(SaveDir,'plots')
    if not os.path.exists(PlotDir):
        os.makedirs(PlotDir)

    #%% Upload the whole experiment and generate the global clock
    file_name = os.path.join(base_dir_server,mID,rec_name,'experiment1','recording1')
    exp = EEGexp(file_name, preprocess=False, make_stim_csv=False)

    #Load probe data
    probe_unit_data, probe_info, total_units, metric_list = tbd_util.get_neuropixel_data(exp)
    probe_list = [x.replace('_sorted', '') for x in exp.experiment_data if 'probe' in x]

    #Get recording start time
    probei = probe_list[0]
    open_ephys_start = np.round(np.min(probe_unit_data[probei]['spike_times']))
    open_ephys_end = np.round(np.max(probe_unit_data[probei]['spike_times']))
    recording_length = open_ephys_end - open_ephys_start
    mm, ss = divmod(recording_length,60)
    hh, mm = divmod(mm, 60)
    print(f'{hh} hrs, {mm} minutes, {ss} seconds')
    
    ##------------------------------------------
    #Get injection times and types of injection
    if drug_type in ['saline', 'psilocybin','ketanserin+psilocybin']:
        injection_time_windows = np.array([np.array(exp_df['First injection window'].values[0].split(','),dtype=float),
                                np.array(exp_df['Second injection window'].values[0].split(','),dtype=float)])
        injection_times = np.array([exp_df['First injection time (s)'].values[0],exp_df['Second injection time (s)'].values[0]],dtype=float)

        if drug_type == 'psilocybin':
            injection_types = ['sal1','psi']
            injection_colors = sns.xkcd_palette(['dark sky blue','darkish red'])
        elif drug_type == 'saline':
            injection_types = ['sal1', 'sal2']
            injection_colors = sns.xkcd_palette(['dark sky blue','cobalt blue'])
        elif drug_type == 'ketanserin+psilocybin':
            injection_types = ['ket','psi']
            injection_colors = sns.xkcd_palette(['magenta','goldenrod'])
        inj_tuple = (injection_times,injection_types,injection_colors)
    else:
        injection_times = None
        injection_time_windows= None
        inj_tuple = None
    
    #For isoflurane experiments, get iso level
    if drug_type == 'isoflurane':
        iso_level, iso_times = exp.load_analog_iso()
        iso_induction_times, iso_maintenance_times = exp.load_iso_times()
        induction_colors = sns.xkcd_palette(['light teal','teal'])
        iso_tuple = (iso_induction_times, induction_colors)
    else:
        iso_induction_times = None
        iso_tuple = None

    ##------------------------------------------
    window_size_min = 1
    window_size = window_size_min*60
    minumum_window_size_min = 0.75
    minumum_window_size = minumum_window_size_min*60

    #Define time windows for each epoch  
    if drug_type in ['saline', 'psilocybin','ketanserin+psilocybin']:
        # epoch_windows = [[open_ephys_start,injection_time_windows[0,0]]]; epoch_names = ['pre_inj']
        # block_labels = epoch_names.copy()
        # epoch_windows.append([injection_time_windows[0,1],injection_time_windows[1,0]]); epoch_names.append(f'post_{injection_types[0]}_inj'); block_labels.append(f'post_{injection_types[0]}_inj')
        # macro_windows = epoch_windows.copy(); macro_names = epoch_names.copy()
        # macro_windows.append([injection_time_windows[1,1],open_ephys_end]); macro_names.append(f'post_{injection_types[1]}_inj')


        if (injection_time_windows[0,0] - open_ephys_start) <= window_size:
            epoch_windows = [[open_ephys_start,injection_time_windows[0,0]]]
            epoch_names = ['pre_inj']
            block_labels = ['pre_inj']

        else:
            epoch_windows = []; epoch_names = []; block_labels = []
            t0 = open_ephys_start
            iW = 0; tblock = injection_time_windows[0,0]
            while t0 < tblock:
                t1 = t0 + window_size
                if t1 > tblock:
                    t1 = tblock
                if (t1 - t0) < minumum_window_size:
                    epoch_windows[-1] = [epoch_windows[-1][0],tblock]
                    break
                epoch_windows.append([t0,t1]); epoch_names.append(f'pre_inj_{int(iW*window_size_min)}_{int((iW+1)*window_size_min)}')
                block_labels.append(f'pre_inj_{int(iW*window_size_min)}_{int((iW+1)*window_size_min)}')
                t0 = t1
                iW += 1

        epoch_windows.append([injection_time_windows[0,1],injection_time_windows[1,0]]); epoch_names.append(f'post_{injection_types[0]}_inj'); block_labels.append(f'post_{injection_types[0]}_inj')
        macro_windows = epoch_windows.copy(); macro_names = epoch_names.copy()
        macro_windows.append([injection_time_windows[1,1],open_ephys_end]); macro_names.append(f'post_{injection_types[1]}_inj')
        t_inj2 = injection_time_windows[1,1]
        t0 = t_inj2

        iW = 0
        while t0 < open_ephys_end:
            t1 = t0 + window_size
            if t1 > open_ephys_end:
                t1 = open_ephys_end
            if (t1 - t0) < minumum_window_size:
                epoch_windows[-1] = [epoch_windows[-1][0],open_ephys_end]
                break
            epoch_windows.append([t0,t1]); epoch_names.append(f'post_{injection_types[1]}_inj')
            block_labels.append(f'{int(iW*window_size_min)}_{int((iW+1)*window_size_min)}')
            t0 = t1
            iW += 1
        nCond = len(epoch_windows)
        if drug_type == 'saline':
            cmap_macro = np.concatenate((sns.xkcd_palette(['silver','dark sky blue']),sns.color_palette('Blues_r',nCond-2)))
        elif drug_type == 'psilocybin':
            cmap_macro = np.concatenate((sns.xkcd_palette(['silver','dark sky blue']),sns.color_palette('Reds_r',nCond-2)))
        elif drug_type == 'ketanserin+psilocybin':
            cmap_macro = np.concatenate((sns.xkcd_palette(['silver','magenta']),sns.color_palette('Purples_r',nCond-2)))
        
    elif drug_type == 'isoflurane':

        epoch_windows = [[open_ephys_start,iso_induction_times[0]-120]]; epoch_names = ['pre_iso']
        epoch_windows.append([iso_induction_times[0],iso_induction_times[1]]); epoch_names.append(f'iso_ind')
        macro_windows = epoch_windows.copy(); macro_names = epoch_names.copy()
        macro_windows.append([iso_induction_times[1]+60*5,open_ephys_end]); macro_names.append(f'post-iso')
        block_labels = epoch_names.copy()
        t_ind = iso_induction_times[1] + 60*5
        t0 = t_ind
        iW = 0
        while t0 < open_ephys_end:
            t1 = t0 + window_size
            if t1 > open_ephys_end:
                t1 = open_ephys_end
            if (t1 - t0) < minumum_window_size:
                epoch_windows[-1] = [epoch_windows[-1][0],open_ephys_end]
                break
            epoch_windows.append([t0,t1]); epoch_names.append(f'post_iso')
            block_labels.append(f'{int(iW*window_size_min)}_{int((iW+1)*window_size_min)}')
            t0 = t1
            iW += 1
        nCond = len(epoch_windows)
        cmap_macro = np.concatenate((sns.xkcd_palette(['silver','light teal']),sns.color_palette('Greens_r',nCond-2)))  

    elif drug_type == 'urethane':
        macro_windows = [[open_ephys_start,open_ephys_end]]; macro_names = ['urethane']
        t0 = 0
        epoch_windows = [];epoch_names = []; block_labels = []
        iW = 0
        while t0 < open_ephys_end:
            t1 = t0 + window_size
            if t1 > open_ephys_end:
                t1 = open_ephys_end
            t0_str = int(t0/60); t1_str = int(t1/60)
            if (t1 - t0) < minumum_window_size:
                epoch_windows[-1] = [epoch_windows[-1][0],open_ephys_end]
                break
            epoch_windows.append([t0,t1]) #; block_labels.append(f'{t0_str}_{t1_str}')
            epoch_names.append(f'urethane')
            block_labels.append(f'{int(iW*window_size_min)}_{int((iW+1)*window_size_min)}')
            t0 = t1
            iW += 1
        nCond = len(epoch_windows)
        cmap_macro = sns.color_palette('Oranges_r',nCond)
    nEpochs = len(epoch_windows)
    epoch_dict = {i: (tW, e) for i, (tW, e) in enumerate(zip(epoch_windows, epoch_names))}
    for ii, (tW, e) in epoch_dict.items():
        print(f'{e}, {block_labels[ii]}: {tW[0]:.1f}s - {tW[1]:.1f}s')
    pdb.set_trace()
    ##------------------------------------------
    #Load stimulus log if it exists
    try:
        stim_log = pd.read_csv(exp.stimulus_log_file)

        if ('circle' in np.unique(stim_log['stim_type'])) | ('natural_scene' in np.unique(stim_log['stim_type'])):
            vStim_exists = True
            stim_log_v = stim_log.loc[stim_log.stim_type.isin(visual_stim)]
        else:
            vStim_exists = False
        if 'biphasic' in np.unique(stim_log['stim_type']):
            eStim_exists = True
            stim_log_b = stim_log.loc[stim_log.stim_type == 'biphasic']
            uniq_currents = np.unique(stim_log_b['parameter'])
            uniq_sweeps = np.unique(stim_log_b['sweep'])
            nSweeps = len(uniq_sweeps)
            nCurrents = len(uniq_currents)
        else:
            eStim_exists = False
    except:
        stim_log = None
        eStim_exists = False
        vStim_exists = False

    ##------------------------------------------
    #Load behavior
    run_ts, run_signal, run_signal_s, acceleration, pupil_ts, pupil_radius, pupil_dxdt, plot_pupil = util.get_behavioral_data(exp, mID, rec_name,normalize=True)
    run_signal[run_signal < 0] = 0
    run_signal_s[run_signal_s < 0] = 0
    f_run = interp1d(run_ts,run_signal)
    f_run_s = interp1d(run_ts,run_signal_s)
    f_pupil = interp1d(pupil_ts,pupil_radius)
    run_signal_p = f_run(pupil_ts)

    # #Get time windows for each epoch
    # epoch_list, block_labels, time_window_list = util.define_epochs_of_interest([open_ephys_start,open_ephys_end], drug_type, window_t_min=window_t_min, injection_times=injection_times,injection_time_windows=injection_time_windows, iso_induction_times=iso_induction_times, stim_log=None)
    
    epoch_list = epoch_names
    time_window_list = epoch_windows
    filename_list = epoch_list
    # Create a new PowerPoint presentation to save figures to
    prs = Presentation()

    #Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f'mID = {mID}; rec_name = {rec_name}'
    slide.placeholders[1].text = f'Firing rate changes due to psilocybin in evoked & spontaneous blocks'

    ##------------------------------------------
    fig, axes = plt.subplots(2,1,figsize=(10,6),gridspec_kw={'height_ratios':[4,2]})
    ax = axes[0]
    plt.suptitle(f'{mID}, {rec_name}')

    ax.plot(run_ts/60,run_signal,'-k')
    ax.set_ylabel('Running speed')

    if drug_type in ['saline', 'psilocybin','ketanserin+psilocybin']:
        ax.vlines(np.array(injection_times)/60,*ax.get_ylim(),color=usrplt.cc[0],lw=2,zorder=4)

    if plot_pupil:
        ax2 = ax.twinx()
        ax2.plot(pupil_ts/60,pupil_radius,'-r')
        ax2.set_ylabel('Pupil size',color='r')

    ax = axes[1]
    if drug_type == 'isoflurane':
        if plot_pupil:
            t_mm = np.min([np.max(iso_times),np.max(pupil_ts)])
        else:
            t_mm = np.max(iso_times)
        indy = np.where(iso_times < t_mm)[0]
        ax.plot(iso_times[indy]/60,iso_level[indy],'-k')
        ax.set_ylabel('Iso level')

    for ax in axes:
        ax.autoscale(tight=True)
        usrplt.adjust_spines(ax)

    for ii, (epoch, tW) in enumerate(zip(filename_list,time_window_list)):
        print(f'\t{epoch:25s} = [{tW[0]:6.1f}, {tW[1]:6.1f}] -> {np.diff(tW)[0]/60:.1f} minutes')
        ax.fill_between(np.array(tW)/60,0,6,color=usrplt.cc[ii],lw=1.5,zorder=0,alpha=0.5)
        ax.text(tW[0]/60,ii*0.5,epoch,fontsize=8)
    usrplt.save_fig_to_pptx(fig, prs)
    plt.savefig(join(PlotDir,f'behavioral_traces_{mID}_{rec_name}.png'))
    plt.savefig(join(SaveDir,f'epochs_{mID}_{rec_name}.png'),dpi=300,bbox_inches='tight',facecolor='w')
    # prs.save(join(SaveDir,f'epochs_{rec_name}.pptx'))
    # print('DONE!!!')
    # exit()
    
    ##------------------------------------------
    #Read in neuropixel data for spontaneous periods
    data_list, ts_list, neuron_indices, plot_tuple, FR_perblock = util.bin_spiking_data(probe_unit_data, time_window_list, time_bin=time_bin,fr_thresh=fr_thresh)
    
    if len(plot_tuple) == 11:
        boundaries_group, ticks_group, labels_group, celltypes, durations, layers, areas, groups,mesogroups, supergroups, order_by_group = plot_tuple
    else:
        print('Experiment has no area information. Check to see if histology is completed. Exiting')
        exit()
    
    nNeurons = len(areas)
    ##------------------------------------------
    ## Determine running behaviors during spontaneous periods
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    behav_list = []; spont_mask_list = []
    for jj, tWindow in enumerate(time_window_list):
        e = epoch_list[jj]
        print(f'{e}: {np.diff(tWindow)[0]/60:.1f} min')
        ts = ts_list[jj]
        T = len(ts)
        tMin = T*time_bin/60
        run_spont = f_run_s(ts) #Use "smoothed" running trace to interpolate to spike bins
        behavior_map = np.zeros(run_spont.shape)

        for key, b_range in behavior_ranges.items():
            indy = np.where((run_spont >= b_range[0]) & (run_spont < b_range[1]))[0]
            behavior_map[indy] = key
            t = len(indy)*time_bin
            print(f'\t: {behavior_dict[key]:10s} -> {t:.1f}s')

        x = behavior_map
        transitions = np.where(np.diff(np.concatenate(([x[0]+1],x,[x[-1]+1]))))[0]

        #Ensure behavioral bouts are at least 5 seconds
        min_t_bout = 3
        for ii, t1 in enumerate(transitions[:-1]):
            t2 = transitions[ii+1]
             
            if (t2-t1)*time_bin < min_t_bout:
                behavior_map[t1:t2] = np.nan
        behav_list.append(behavior_map)

        spont_periods = np.ones((T))
        #Mask out periods of evoked activity
        if eStim_exists | vStim_exists:
            
            buffer = 0.05
            for ii, row in stim_log.iterrows():
                stim_type = row['stim_type']
                if stim_type == 'biphasic':
                    #Just block out 1 second after electrical stimulation
                    t1 = row['onset']  - buffer
                    t2 = row['onset'] + 1
                else:
                    #Block out whole duration of stimulus, as it is likely a visual stimulus
                    t1 = row['onset'] - buffer
                    t2 = row['offset']

                indy = np.where((ts >= t1) & (ts < t2))[0]
                if len(indy) > 0:
                    spont_periods[indy] = 0
        spont_mask_list.append(spont_periods)

    ##------------------------------------------
    ## Calculate firing rates
    ## Calculate single cell gain
    T, nNeurons = data_list[0].shape
    nCond = len(time_window_list)
    FR_spont = np.full((nCond,nNeurons,nBehaviors+1),np.nan)
    gain_relu = np.full((nCond,nNeurons,nBehaviors+1),np.nan); r2_relu = np.full((nCond,nNeurons,nBehaviors+1),np.nan)
    gain_sigm = np.full((nCond,nNeurons,nBehaviors+1),np.nan); r2_sigm = np.full((nCond,nNeurons,nBehaviors+1),np.nan)
    nBins_per_behav_spont = np.zeros((nCond,nBehaviors+1))

    idx_to_plot = np.random.choice(nNeurons,10,replace=False)
    data_distributions = []
    for jj, tWindow in enumerate(time_window_list):
        print(f'\nEpoch: { epoch_list[jj]}')
        ts = ts_list[jj]
        data = data_list[jj]
        behavior_map = behav_list[jj]
        spont_periods = spont_mask_list[jj]
  
        for key, behav_str in behavior_dict2.items():
            if key == 2:
                indy = np.where(spont_periods == 1)[0]
            else:
                indy = np.where((behavior_map == key) & (spont_periods == 1))[0]
            if len(indy) == 0:
                continue
            FR_all = data[indy]/time_bin
            FR_spont[jj,:,key] = np.mean(data[indy],axis=0)/time_bin
            nBins_per_behav_spont[jj,key] = len(indy)
            # import pdb; pdb.set_trace()

            #Calculate single cell gain
            for iN in tqdm(np.arange(nNeurons)):
                if FR_spont[jj,iN,key] < 0.5:
                    continue
                if iN in idx_to_plot:
                    plot = True
                    titlestr = f'Neuron {iN}, {areas[iN]}, {celltypes[iN]}, {behav_str}, {epoch_list[jj]}'
                    rp, sp, error, fig = calculate_gain(FR_all[:,iN],plot=plot,titlestr=titlestr)
                    usrplt.save_fig_to_pptx(fig, prs)
                else: 
                    rp, sp, error, fig = calculate_gain(FR_all[:,iN])
                
                if error == 1:
                    continue
                gain_relu[jj,iN,key] = rp[2]; r2_relu[jj,iN,key] = rp[0]
                gain_sigm[jj,iN,key] = sp[-1]; r2_sigm[jj,iN,key] = sp[0]

    #Save time windows
    np.savez(join(SaveDir,f'time_windows.npz'),time_window_list=time_window_list,filename_list=filename_list,nBins_per_behav_spont=nBins_per_behav_spont,areas=areas)

    #Create combinations that we are interested in comparing
    cbs = combinations(filename_list,2)
    ijs = combinations(np.arange(len(filename_list)),2)
    combos_spont = []
    for cb,ij in zip(cbs,ijs):
        epoch_i = cb[0];i = ij[0]
        epoch_j = cb[1];j = ij[1]
        if ('pre' in epoch_i) | ('post-sal1' in epoch_i) | ('post-ket' in epoch_i):
            if ('post' in epoch_j):
                print(ij,cb)
                combos_spont.append(ij)

    nComparisons = len(combos_spont); pval_thresh = 0.05
    uniq_uidx = np.array([f'{rec_name}_{i}' for i in range(nNeurons)])

    #Convert firing rate array into dataframe for easy plotting & saving
    indy = np.arange(nNeurons); fsuffix = 'ALL'; N = nNeurons
    tmp_list = []
    for j, tWindow in enumerate(time_window_list):
        print(f'\nEpoch: { epoch_list[j]}')
        for key, behav_str in behavior_dict2.items():
            #Check to make sure each epoch has 10 seconds of data
            bin_thresh = 10/time_bin
            if (nBins_per_behav_spont[j,key] < bin_thresh):
                print(f'\t{behav_str:20s}: Not enough data to calculate FR')
                continue
            else:
                print(f'\t{behav_str:20s}: {nBins_per_behav_spont[j,key]*time_bin:.1f} seconds of data to calculate FR')

            tmp_list.append(np.stack((np.repeat(mID,N),np.repeat(rec_name,N),np.repeat(drug_type,N),uniq_uidx,celltypes, durations, areas,layers,groups,mesogroups,supergroups,np.repeat(behav_str,N),np.repeat( epoch_list[j],N),np.repeat(block_labels[j],N),np.repeat('',N),np.repeat('spont',N),np.repeat('mean_FR',N),FR_spont[j,:,key])))
    FR_spont_df = pd.DataFrame(np.hstack(tmp_list).T,columns=['mID','rec_name','drug_type','neuron_index','celltype','duration','area','layer','group','mesogroup','supergroup','behavior','epoch','block','num_spks','stim_type','stat','val'])
    FR_spont_df = FR_spont_df.astype({'val':float})
    FR_spont_df.to_hdf(join(SaveDir,f'spont_FR_df_{rec_name}.h5'),'df')

    #Convert firing rate array into dataframe for easy plotting & saving
    indy = np.arange(nNeurons); fsuffix = 'ALL'
    tmp_list = []
    for j, tWindow in enumerate(time_window_list):
        print(f'\nEpoch: { epoch_list[j]}')
        for key, behav_str in behavior_dict2.items():
            #Check to make sure each epoch has 10 seconds of data
            bin_thresh = 10/time_bin
            if (nBins_per_behav_spont[j,key] < bin_thresh):
                print(f'\t{behav_str:20s}: Not enough data to calculate FR')
                continue
            else:
                print(f'\t{behav_str:20s}: {nBins_per_behav_spont[j,key]*time_bin:.1f} seconds of data to calculate FR')

            tmp_list.append(np.stack((np.repeat(mID,N),np.repeat(rec_name,N),np.repeat(drug_type,N),uniq_uidx,celltypes, durations, areas,layers,groups,mesogroups,supergroups,np.repeat(behav_str,N),np.repeat( epoch_list[j],N),np.repeat(block_labels[j],N),np.repeat('',N),np.repeat('spont',N),np.repeat('gain_sigma',N),r2_sigm[j,:,key],gain_sigm[j,:,key])))
            tmp_list.append(np.stack((np.repeat(mID,N),np.repeat(rec_name,N),np.repeat(drug_type,N),uniq_uidx,celltypes, durations, areas,layers,groups,mesogroups,supergroups,np.repeat(behav_str,N),np.repeat( epoch_list[j],N),np.repeat(block_labels[j],N),np.repeat('',N),np.repeat('spont',N),np.repeat('gain_relu',N),r2_relu[j,:,key],gain_relu[j,:,key])))
    gain_df = pd.DataFrame(np.hstack(tmp_list).T,columns=['mID','rec_name','drug_type','neuron_index','celltype','duration','area','layer','group','mesogroup','supergroup','behavior','epoch','block','num_spks','stim_type','stat','r2','val'])
    gain_df = gain_df.astype({'val':float})
    gain_df.to_hdf(join(SaveDir,f'gain_df_{rec_name}.h5'),'df')

    if len(combos_spont) == 0:
        print('No comparisons to make. Exiting')
        exit()

    #Compare FR changes due to saline or psilocybin injections vs no-injection
    f_list = []; g_list = []
    for cb in combos_spont:
        i = cb[0]
        j = cb[1]
        print(f'Calculating differences between {filename_list[i]} and {filename_list[j]}')

        for key, behav_str in behavior_dict2.items():

            #Select neurons FR for the epoch & behavior 
            FR_i = FR_spont[i,:,key]; epoch_i = filename_list[i]; block_i = block_labels[i]
            FR_j = FR_spont[j,:,key]; epoch_j = filename_list[j]; block_j = block_labels[j]
            gs_i = gain_sigm[i,:,key]; gs_j = gain_sigm[j,:,key]
            gr_i = gain_relu[i,:,key]; gr_j = gain_relu[j,:,key]

            #Check to make sure each epoch has 5 seconds of data
            bin_thresh = 10/time_bin
            if (nBins_per_behav_spont[i,key] < bin_thresh) | (nBins_per_behav_spont[j,key] < bin_thresh):
                continue
            
            #Calculate difference in mean firing rates
            dFR = FR_j - FR_i
            MI = (FR_j - FR_i)/(FR_j + FR_i)

            dGS = gs_j - gs_i
            dGR = gr_j - gr_i
            MI_GS = (gs_j - gs_i)/(gs_j + gs_i)
            MI_GR = (gr_j - gr_i)/(gr_j + gr_i)

            # #Get distributions for each condition
            # if key == 2:
            #     mask_i = np.where(spont_mask_list[i] == 1)[0]
            #     mask_j = np.where(spont_mask_list[j] == 1)[0]
            # else:
            #     mask_i = np.where((behav_list[i] == key) & (spont_mask_list[i] == 1))[0]
            #     mask_j = np.where((behav_list[j] == key) & (spont_mask_list[j] == 1))[0]
            # data_i = data_list[i][mask_i]
            # data_j = data_list[j][mask_j]

            # #Test for significance for each neuron
            # def statistic(x, y):
            #     return np.mean(x) - np.mean(y)
            
            for iN in tqdm(indy):

                # res = st.mannwhitneyu(data_i[:,iN],data_j[:,iN])
                # pval = res.pvalue
                # sign = int(pval < pval_thresh)
                # res = st.permutation_test([data_i[:,iN],data_j[:,iN]],statistic,vectorized=False)
 
                # pval2 = res.pvalue
                # sign2 = int(pval2 < pval_thresh)
                # sig = 1 if (sign + sign2) == 2 else 0
                # # import pdb; pdb.set_trace()
                pval = np.nan; pval2 = np.nan
                sign = 1; sign2 = 1; sig = 1; 

                f_list.append((mID,rec_name,drug_type,uniq_uidx[iN],celltypes[iN],durations[iN],areas[iN],layers[iN],groups[iN],mesogroups[iN],supergroups[iN],behav_str,epoch_i,block_i,epoch_j,block_j,'mean_FR',FR_i[iN],FR_j[iN],dFR[iN],pval,sign,pval2,sign2,sig,MI[iN]))
        
                g_list.append((mID,rec_name,drug_type,uniq_uidx[iN],celltypes[iN],durations[iN],areas[iN],layers[iN],groups[iN],mesogroups[iN],supergroups[iN],behav_str,epoch_i,block_i,epoch_j,block_j,'gain_relu',gr_i[iN],gr_j[iN],dGR[iN],r2_relu[i,iN,key],r2_relu[j,iN,key],MI_GR[iN]))
                g_list.append((mID,rec_name,drug_type,uniq_uidx[iN],celltypes[iN],durations[iN],areas[iN],layers[iN],groups[iN],mesogroups[iN],supergroups[iN],behav_str,epoch_i,block_i,epoch_j,block_j,'gain_sigma',gs_i[iN],gs_j[iN],dGS[iN],r2_sigm[i,iN,key],r2_sigm[j,iN,key],MI_GS[iN]))

    dFR_spont_df = pd.DataFrame(np.stack(f_list),columns=['mID','rec_name','drug_type','neuron_index','celltype','duration','area','layer','group','mesogroup','supergroup','behavior','epoch_i','block_i','epoch_j','block_j','stat','val_i','val_j','delta','pval1','sig1','pval2','sig2','sig','mod'])
    dFR_spont_df = dFR_spont_df.astype({'duration':float,'val_i':float,'val_j':float,'delta':float,'pval1':float,'sig1':int,'pval2':float,'sig2':int,'sig':int,'mod':float})
    dFR_spont_df.to_hdf(join(SaveDir,f'delta-spont_FR_{fsuffix}_{rec_name}.h5'),'df')

    dGain_df = pd.DataFrame(np.stack(g_list),columns=['mID','rec_name','drug_type','neuron_index','celltype','duration','area','layer','group','mesogroup','supergroup','behavior','epoch_i','block_i','epoch_j','block_j','stat','val_i','val_j','delta','r2_i','r2_j','mod'])
    dGain_df = dGain_df.astype({'duration':float,'val_i':float,'val_j':float,'delta':float,'r2_i':float,'r2_j':float,'mod':float})
    dGain_df.to_hdf(join(SaveDir,f'delta-gain_{fsuffix}_{rec_name}.h5'),'df')

    # Save the PowerPoint presentation
    prs.save(join(SaveDir,f'firing_rate_spont_figs_{rec_name}.pptx'))
    print('DONE!!!')


