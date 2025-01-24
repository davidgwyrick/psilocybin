base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
base_dir_local = '/data/tiny-blue-dot/zap-n-zip/EEG_exp'; base_dir = base_dir_server
ProjDir = '/home/david.wyrick/projects/zap-n-zip/'
ServDir = '/data/projects/zap-n-zip/'

#Base
import argparse
from glob import glob
from os.path import join
import json, os, time, sys
import ray
import gspread
import pandas as pd
import numpy as np
import itertools as it

#Powerlaw
import powerlaw as pl
from scipy.optimize import curve_fit
from sklearn.metrics import r2_score
import criticality as cr

#Scipy
import scipy.signal as sig
import scipy.stats as st
from scipy.ndimage import gaussian_filter
from scipy.interpolate import interp1d

#Plot
import matplotlib.pyplot as plt
import seaborn as sns
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation

#Project
from tbd_eeg.data_analysis.eegutils import EEGexp
import tbd_eeg.data_analysis.Utilities.utilities as tbd_util

#User
sys.path.append(ProjDir)
import util
import plotting as usrplt

#Templeton-log_exp
gc = gspread.service_account() 
sh = gc.open('Templeton-log_exp') 
exp_table = pd.DataFrame(sh.sheet1.get()).T.set_index(0).T 
exp_table = exp_table.set_index('mouse_name')
exp_table.head()

bandpass_freqs = np.array([[1, 4], [4, 8], [8,13], [13, 30], [30, 55], [65, 100], [100,200]])
bandpass_str = ['delta','theta','alpha','beta','low-gamma','high-gamma','HFO']
n_bands = len(bandpass_freqs)

## Parse Command Line Arguments ----------------------------------------
parser = argparse.ArgumentParser(description='Avalanche Analysis, REDUX ')

parser.add_argument('--mID',type=str, default='mouse735052',
                    help='mouse to perform analysis on')

parser.add_argument('--rec_name',type=str, default='aw_psi_2024-06-20_10-40-56',
                    help='experiment to perform analysis on')

parser.add_argument('--level',type=str, default='group',
                    help='Level at which to group neurons')

parser.add_argument('--time_bin_ms',type=int, default=50,
                    help='Time bin width')

parser.add_argument('--pl_distr',type=str, default='power_law',
                    help='power_law or truncated_power_law')

parser.add_argument('--perc',type=float, default=0.1,
                    help='Network threshold for determining avalanches')

parser.add_argument('--F_thresh',type=float, default=0.75,
                    help='Threshold for determining goodness of fit')

parser.add_argument('--network_threshold_type',type=str, default='constant',
                    help='Whether to use the same network threshold per epoch')

parser.add_argument('--parallel',type=int, default=1,
                    help='Parallel processing')


def get_avalanches(spks, prune_size=1,thresh=0.35,const_threshold=None):
    T, N = spks.shape

    # #Binarize spiking data
    # spks = np.array(spks > 0,dtype=int)

    #Identify threshold based on network activity
    network_activity = np.sum(spks,axis=1)
    if const_threshold is not None:
        network_thresh = const_threshold
    else: 
        network_thresh = np.nanpercentile(network_activity,thresh*100)

    #Identify transitions between avalanches
    active_time_bins = np.array(network_activity > network_thresh,dtype=bool)
    active_time_bins1 = np.array(network_activity > network_thresh,dtype=int)
    transitions = np.where(np.diff(active_time_bins))[0]
    transi_type = np.diff(active_time_bins1)[transitions]

    #identify first & last avalanche
    ii = np.where(transi_type == 1)[0][0]
    jj = np.where(transi_type == -1)[0][-1] + 1
    nAvalanches = len(transitions[ii:jj])/2
    assert len(transitions[ii:jj]) % 2 == 0

    #Get indices of avalanches
    indices = transitions[ii:jj].reshape(-1,2) + 1

    #Prune avalanches that are too short
    D_bin = indices[:,1] - indices[:,0]
    indices = indices[D_bin >= prune_size]
    
    #Get avalanche length, size, and spikes
    D = (indices[:,1] - indices[:,0])# * time_bin
    S = np.array([np.sum(network_activity[iPair[0]:iPair[1]]) for iPair in indices])

    return D, S, indices

def resample_cdf(data, cdf, num_points_per_decade=10):
    # Logarithmic spacing between x_min and x_max
    log_min, log_max = np.log10(data[0]), np.log10(data[-1])
    resampled_points = np.logspace(log_min, log_max, num=int((log_max - log_min) * num_points_per_decade))
    
    # Interpolate the CDF at the resampled points
    resampled_cdf = np.interp(resampled_points, data, cdf)
    return resampled_points, resampled_cdf

def calculate_fraction_within_bounds(resampled_cdf, lower_bound, upper_bound):
    # Check if resampled points fall within bounds
    within_bounds = (resampled_cdf >= lower_bound) & (resampled_cdf <= upper_bound)
    fraction_within_bounds = np.sum(within_bounds) / len(resampled_cdf)
    return fraction_within_bounds

# Define the power-law function
def power_law(x, exponent):
    return  x ** exponent

# Define the power-law function
def power_law_sigma(x, exponent, sigma):
    return  sigma * (x ** exponent)

def line_fit(x, slope):
    return slope*x

@ray.remote
def goodness_of_fit_p(data, xmin, xmax, distribution='power_law',pba=None):
    return goodness_of_fit(data, xmin, xmax, distribution, pba)

def goodness_of_fit(data, xmin, xmax, distribution='power_law',pba=None): 
    #Update ray progress bar
    if pba is not None:
        pba.update.remote(1)
    
    #Ensure data spans at least 1 decade
    if xmax - xmin < 10:
        return np.nan, np.nan, np.nan
    
    xmin = int(xmin); xmax = int(xmax)
    #Truncate data
    x_truncated = data[(data >= xmin) & (data <= xmax)]

    #Get powerlaw range
    pl_range = np.log10(xmax/xmin)

    #Define pdf
    pdf = np.histogram(data, bins = np.arange(1, np.max(data) + 2))[0]
    p = pdf/np.sum(pdf)

    # 0) Fit powerlaw to data to obtain tau 
    pl_data = pl.Fit(data,xmin_distribution=distribution,discrete=False,xmin=xmin,xmax=xmax,pdf_ends_at_xmax=True)

    # 1) Create cumulative distribution function (CDF) of real data
    if distribution == 'truncated_power_law':
        params = [pl_data.truncated_power_law.parameter1,pl_data.truncated_power_law.parameter2]
        # bins, actual_CDF = pl_data.truncated_power_law.cdf()
    elif distribution == 'power_law':
        params = [pl_data.power_law.parameter1]
        
    bins, actual_CDF = pl_data.cdf()
    tau = params[0]
   
    # 2) Define a theoretical CDF for a truncated power law with the same range and tau
    x = np.arange(xmin,xmax,.1)
    if distribution == 'truncated_power_law':
        pl_model = pl.Truncated_Power_Law(xmin=xmin, xmax=xmax, parameters=params)
      
    elif distribution == 'power_law':
        pl_model = pl.Power_Law(xmin=xmin, xmax=xmax, parameters=params)
        
    theoretical_CDF = pl_model.cdf(x)
    KS = pl_model.KS(x_truncated)

    # 3) Define a region delimited by upper and lower bounds defined as the theoretical CDF +0.03 and -0.03
    lower_bound = theoretical_CDF - 0.03
    upper_bound = theoretical_CDF + 0.03
  
    # 4) Resample the real CDF at 10 logarithmically spaced values per decade
    resampled_points, actual_CDF_resampled = resample_cdf(bins, actual_CDF, 1000)
    if len(resampled_points) == 0:
        return np.nan, KS, tau

    # 5) Calculate the fraction of resampled points in real CDF data that fell within +-0.03 bounds of the theoretical CDF
    eps = 1E-6
    theoretical_CDF_resampled =  pl_model.cdf(resampled_points+eps)
    lower_bound = theoretical_CDF_resampled - 0.03
    upper_bound = theoretical_CDF_resampled + 0.03
    if actual_CDF_resampled.shape[0] != theoretical_CDF_resampled.shape[0]:
        actual_CDF_resampled = actual_CDF_resampled[1:]
    F = calculate_fraction_within_bounds(actual_CDF_resampled, lower_bound, upper_bound)

    return F, KS, tau

def goodness_of_fit_plot(data, xmin, xmax, distribution='power_law', plot=True, xlabel='Size (# of spikes)'):

    xmin = int(xmin); xmax = int(xmax)
    #Truncate data
    x_truncated = data[(data >= xmin) & (data <= xmax)]

    #Get powerlaw range
    pl_range = np.log10(xmax/xmin)

    #Define pdf
    pdf = np.histogram(data, bins = np.arange(1, np.max(data) + 2))[0]
    p = pdf/np.sum(pdf)

    # 0) Fit powerlaw to data to obtain tau 
    pl_data = pl.Fit(data,xmin_distribution=distribution,discrete=False,xmin=xmin,xmax=xmax,pdf_ends_at_xmax=True)

    # 1) Create cumulative distribution function (CDF) of real data
    if distribution == 'truncated_power_law':
        params = [pl_data.truncated_power_law.parameter1,pl_data.truncated_power_law.parameter2]
        # bins, actual_CDF = pl_data.truncated_power_law.cdf()
    elif distribution == 'power_law':
        params = [pl_data.power_law.parameter1]
        
    bins, actual_CDF = pl_data.cdf()
    tau = params[0]
    
    if plot:
        fig, axes = plt.subplots(1,2,figsize=(10,5))
        #Plot PDF
        ax = axes[0]; ax.set_title('PDF')
        ax.scatter(np.arange(1,len(p)+1),p,marker='o',s=10,facecolor='none',edgecolor=usrplt.cc[2])
        ax.scatter(np.arange(1,len(p)+1)[xmin:xmax],p[xmin:xmax],marker='o',s=10,color=usrplt.cc[2])
        # ax.vlines([xmin,xmax],*ax.get_ylim(),colors=usrplt.cc[2],linestyles='--')
        ax.set_xscale('log')
        ax.set_yscale('log')
        ax.set_xlabel(xlabel)
        ax.set_ylabel('Probablity P(S)')
        usrplt.adjust_spines(ax)

        edges, hist = pl_data.pdf(original_data=True) #linear_bins
        ax.scatter(edges[:-1] + np.diff(edges)/2,hist,marker='o',s=10,color='k')

        #Plot CDF
        ax = axes[1]; ax.set_title('CDF')
        ax.plot(bins, actual_CDF, label='Empirical CDF')
        ax.set_xscale('log')
        ax.set_xlabel('Size (# of spikes)')
        ax.set_ylabel('Cumulative Probability')
        usrplt.adjust_spines(ax)

    # 2) Define a theoretical CDF for a truncated power law with the same range and tau
    x = np.arange(xmin,xmax,.1)
    if distribution == 'truncated_power_law':
        pl_model = pl.Truncated_Power_Law(xmin=xmin, xmax=xmax, parameters=params)
        if plot:
            pl_model.plot_pdf(data,ax=axes[0],color='k',ls='--',label=f'\u03C4={params[0]:.3f}, \u03BB = {params[1]:.3f}')
            axes[0].legend()

    elif distribution == 'power_law':
        pl_model = pl.Power_Law(xmin=xmin,xmax=xmax, parameters=params)
        if plot:
            pl_model.plot_pdf(data,ax=axes[0],color='k',ls='--',label=f'\u03C4={params[0]:.3f}')
            axes[0].legend()
        
    theoretical_CDF = pl_model.cdf(x)
    KS = pl_model.KS(x_truncated)

    # 3) Define a region delimited by upper and lower bounds defined as the theoretical CDF +0.03 and -0.03
    lower_bound = theoretical_CDF - 0.03
    upper_bound = theoretical_CDF + 0.03
    if plot:
        ax.plot(x, theoretical_CDF, label=f'Theoretical CDF, \u03C4 = {tau:.2f}',color='k',ls='--')
        ax.fill_between(x, lower_bound, upper_bound, color='k', alpha=0.3, label='3% Bounds')
        ax.legend()

    # 4) Resample the real CDF at 10 logarithmically spaced values per decade
    resampled_points, actual_CDF_resampled = resample_cdf(bins, actual_CDF, 1000)
    if len(resampled_points) == 0:
        return np.nan, KS, tau, None
    # if plot:
    #     ax.scatter(resampled_points, actual_CDF_resampled, label='Resampled CDF',color='r',s=10)

    # 5) Calculate the fraction of resampled points in real CDF data that fell within +-0.03 bounds of the theoretical CDF
    eps = 1E-6
    theoretical_CDF_resampled =  pl_model.cdf(resampled_points+eps)
    lower_bound = theoretical_CDF_resampled - 0.03
    upper_bound = theoretical_CDF_resampled + 0.03
    if actual_CDF_resampled.shape[0] != theoretical_CDF_resampled.shape[0]:
        actual_CDF_resampled = actual_CDF_resampled[1:]
    F = calculate_fraction_within_bounds(actual_CDF_resampled, lower_bound, upper_bound)

    if plot:
        plt.suptitle(f'bounds = [{xmin},{xmax}]; # of decades = {pl_range:.2f}; F = {F:.2f}')
        return F, KS, tau, fig
    else:
        return F, KS, tau, None
     
## Main ----------------------------------------
if __name__ == '__main__':

    ## Parse the arguments ----------------------------------------
    args = parser.parse_args()

    #Which experiment?
    mID = args.mID
    rec_name = args.rec_name
    print(f'{mID}, {rec_name}')

    #How to segment data
    time_bin_ms = args.time_bin_ms
    time_bin = time_bin_ms/1000
    level = args.level
    pl_distr = args.pl_distr
    perc = args.perc
    F_thresh = args.F_thresh
    parallel = bool(args.parallel)
    network_threshold_type = args.network_threshold_type

    ## FOLDERS ----------------------------------------
    #Create directory for saving to
    TempDir = os.path.join(ServDir,'results','avalanches_0107',mID,rec_name)
    if not os.path.exists(TempDir):
        os.makedirs(TempDir)

    tmp_list = sorted(glob(join(TempDir,f'run_*')))
    if len(tmp_list) == 0:
        curr_run = 0
    else:
        last_run = int(tmp_list[-1][-1])
        curr_run = last_run+1
    # curr_run = 7
    folder = f'run_{curr_run:02d}'
    SaveDir = os.path.join(TempDir,folder)
    if not os.path.exists(SaveDir):
        os.makedirs(SaveDir)
    
    PlotDir = join(SaveDir,'plots')
    if not os.path.exists(PlotDir):
        os.makedirs(PlotDir)

    #Save model parameters
    args_dict = args.__dict__
    args_dict['SaveDir'] = SaveDir
    with open(join(SaveDir,f"parameters_run_{curr_run}.json"), "w") as outfile:
        json.dump(args_dict, outfile)

    #Extract row from templeton dataframe
    exp_df = exp_table.loc[exp_table.exp_name == rec_name]
    if type(exp_df) == pd.core.series.Series:
        exp_df = exp_df.to_frame().T

    drug_type = exp_df['drug'].values[0]
    stim_type = exp_df['stimulation'].values[0]
    print(f'Experiment type: {stim_type}, {drug_type}')

    if 'sal2' in drug_type:
        drug_type2 = 'saline2'
    else:
        drug_type2 = drug_type
    # Powerpoint for figures
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f'{mID} {rec_name}'
    slide.placeholders[1].text = f'Avalanche analysis'

    ## Read in experiment ----------------------------------------
    #Upload the whole experiment and generate the global clock
    exp = EEGexp(os.path.join(base_dir,mID,rec_name,'experiment1','recording1'), preprocess=False, make_stim_csv=False)

    #Load probe data
    probe_unit_data, probe_info, total_units, metric_list = tbd_util.get_neuropixel_data(exp)
    probe_list = [x.replace('_sorted', '') for x in exp.experiment_data if 'probe' in x]

    #Get recording start time
    probei = probe_list[0]
    open_ephys_start = np.round(np.min(probe_unit_data[probei]['spike_times']))+10
    open_ephys_end = np.round(np.max(probe_unit_data[probei]['spike_times']))-10
    recording_length = open_ephys_end - open_ephys_start
    mm, ss = divmod(recording_length,60)
    hh, mm = divmod(mm, 60)
    print(f'{hh} hrs, {mm} minutes, {ss} seconds')

    #Bin spiking data
    tW = [open_ephys_start,open_ephys_end]
    data_list, ts_list, neuron_indices, plot_tuple, FR_perblock = util.bin_spiking_data(probe_unit_data,[tW], time_bin=time_bin,fr_thresh=0)
    spike_counts = data_list[0].T; ts = ts_list[0]
    
    if len(plot_tuple) == 11:
        boundaries_group, ticks_group, labels_group, celltypes, durations, layers, areas, groups,mesogroups, supergroups, order_by_group = plot_tuple
    else:
        print('Experiment has no area information. Check to see if histology is completed. Exiting')
        exit()

    nNeurons = len(neuron_indices)
    spike_counts_sub = spike_counts[neuron_indices]
    spike_counts_ro = spike_counts_sub[order_by_group]
    areas_sub = areas[neuron_indices]; areas_ro = areas_sub[order_by_group]
    groups_sub = groups[neuron_indices]; groups_ro = groups_sub[order_by_group]
    mesogroups_sub = mesogroups[neuron_indices]; mesogroups_ro = mesogroups_sub[order_by_group]
    supergroups_sub = supergroups[neuron_indices]; supergroups_ro = supergroups_sub[order_by_group]

    #Just sort data first
    areas = areas_ro; groups = groups_ro; mesogroups = mesogroups_ro; supergroups = supergroups_ro
    spike_counts = spike_counts_ro

    #Read in behavior
    run_ts, run_signal, run_signal_s, acceleration, pupil_ts, pupil_radius, pupil_dxdt, plot_pupil = util.get_behavioral_data(exp, mID, rec_name)
    f_run = interp1d(run_ts,run_signal); run_signal_p = f_run(pupil_ts)
    f_run_s = interp1d(run_ts,run_signal_s); run_signal_p_s = f_run_s(pupil_ts)
    f_pupil = interp1d(pupil_ts,pupil_radius)
    
    injection_times = None; injection_time_windows = None; inj_tuple = None
    iso_induction_times = None; iso_maintenance_times = None; iso_tuple = None
    #For saline & psilocybin experiments, get injection times and types of injection
    if drug_type in ['saline', 'psilocybin','ketanserin+psilocybin','urethane+psilocybin']:
        injection_time_windows = np.array([np.array(exp_df['First injection window'].values[0].split(','),dtype=float),
                                np.array(exp_df['Second injection window'].values[0].split(','),dtype=float)])
        
        #Take second time in each window as "injection time"
        injection_times = np.array([exp_df['First injection time (s)'].values[0],exp_df['Second injection time (s)'].values[0]],dtype=float)

        if drug_type in ['psilocybin','urethane+psilocybin']:
            injection_types = ['sal1','psi']
            injection_colors = sns.xkcd_palette(['dark sky blue','darkish red'])
            epoch_cmap = sns.xkcd_palette(['silver','dark sky blue','darkish red','darkish red'])

        elif drug_type == 'saline':
            injection_types = ['sal1', 'sal2']
            injection_colors = sns.xkcd_palette(['dark sky blue','cobalt blue'])
            epoch_cmap = sns.xkcd_palette(['silver','dark sky blue','cobalt blue','cobalt blue'])

        elif drug_type == 'ketanserin+psilocybin':
            injection_types = ['ket','psi']
            injection_colors = sns.xkcd_palette(['magenta','goldenrod'])
            epoch_cmap = sns.xkcd_palette(['silver','magenta','goldenrod','goldenrod'])
        inj_tuple = (injection_times, injection_types, injection_colors)

        #Define epochs
        epoch_windows = [[open_ephys_start,injection_time_windows[0,0]]]; epoch_names = ['pre-inj']; tPre = injection_time_windows[0,0]-open_ephys_start
        epoch_windows.append([injection_time_windows[0,1],injection_time_windows[1,0]]); epoch_names.append(f'post-{injection_types[0]}-inj')

        epoch_windows.append([injection_time_windows[1,1],injection_time_windows[1,1]+tPre]); epoch_names.append(f'post-{injection_types[1]}-inj_start')
        epoch_windows.append([open_ephys_end-tPre,open_ephys_end]); epoch_names.append(f'post-{injection_types[1]}-inj_end')

        #urethane+psilocybin mouse only had 1 injection
        if injection_times[0] == injection_times[1]:
            epoch_windows.pop(1); epoch_names.pop(1)


    #For isoflurane experiments, get iso level
    elif drug_type == 'isoflurane':
        iso_level, iso_times = exp.load_analog_iso()
        iso_induction_times, iso_maintenance_times = exp.load_iso_times()
        induction_colors = sns.xkcd_palette(['light teal','teal'])
        iso_tuple = (iso_induction_times, induction_colors)

        epoch_windows = [[open_ephys_start,iso_induction_times[0]-120]]; epoch_names = ['pre-iso']
        epoch_windows.append([iso_induction_times[0],iso_induction_times[1]]); epoch_names.append(f'iso-ind')
        epoch_windows.append([iso_induction_times[1]+60*5,open_ephys_end]); epoch_names.append(f'post-iso')
        epoch_cmap = sns.xkcd_palette(['silver','light teal','teal'])
    
    elif drug_type == 'urethane':
        epoch_windows = [[open_ephys_start,open_ephys_end/3]]; epoch_names = ['urethane-1']
        epoch_windows.append([open_ephys_end/3,open_ephys_end*2/3]); epoch_names.append('urethane-2')
        epoch_windows.append([open_ephys_end*2/3,open_ephys_end]); epoch_names.append('urethane-3')
        epoch_cmap = sns.color_palette('Oranges',3)

    # extract the timestamps of the selected stimuli
    try:
        stim_log = pd.read_csv(exp.stimulus_log_file)
        stim_exists = True
    except:
        stim_exists = False

    ## Determine windows ----------------------------------------
    evoked_time_window_list = []; evoked_type_list = []; evoked_tuple = None
    if stim_exists:
        for s in np.unique(stim_log['sweep']):
            for t in np.unique(stim_log.loc[stim_log.sweep == s]['stim_type']):
                sub_df = stim_log.loc[(stim_log.sweep == s) & (stim_log.stim_type == t)]
                tS = sub_df['onset'].min(); tE = sub_df['offset'].max()
                
                evoked_time_window_list.append([tS,tE])
                evoked_type_list.append(t)
        evoked_tuple = (evoked_time_window_list,evoked_type_list)
    
    #Plot behavior
    fig = usrplt.plot_behavior((run_ts,run_signal),(pupil_ts,pupil_radius),f'{mID} {rec_name}',evoked_tuple,inj_tuple,iso_tuple)
    plt.savefig(join(PlotDir,f'behavior_{rec_name}.png'),dpi=300,bbox_inches='tight',facecolor='w')
    usrplt.save_fig_to_pptx(fig, prs)
    plt.close(fig)

    ## Fit avalanche distributions ----------------------------------------
    t0_outer = time.perf_counter()
    lag_sec = 1; sigma = 1
    pl_list = []

    if level == 'supergroup':
        groups_to_iterate = np.unique(supergroups)
    elif level == 'group':
        groups_to_iterate = np.unique(groups)
    elif level == 'all':
        groups_to_iterate = ['all']

    np.savez(join(SaveDir,'time_windows.npz'),groups_to_iterate=groups_to_iterate,epoch_windows=epoch_windows,epoch_names=epoch_names,areas=areas,groups=groups,mesogroups=mesogroups,supergroups=supergroups,epoch_cmap=epoch_cmap)
    # print('poop')
    # exit()
    
    results_list = []; DCC_list = []
    #Loop over groups
    for g in groups_to_iterate:
        if level == 'supergroup':
            g_indy = np.where(supergroups == g)[0]
        elif level == 'group':
            g_indy = np.where(groups == g)[0]
        elif level == 'all':
            g_indy = np.arange(nNeurons)
        if g in ['FT']:
            continue
        nG = len(g_indy)
        print(f'Group {g} has {nG} neurons')
        if nG < 40:
            print('\tGroup too small, skipping')
            continue

        #First define threshold for network activity over the whole experiment
        data_g = spike_counts[g_indy]
        if network_threshold_type == 'constant':
            network_activity = np.sum(data_g,axis=0)
            network_threshold = np.percentile(network_activity,perc*100)
        
        #Loop over epochs
        for iW, (epoch, tW) in enumerate(zip(epoch_names,epoch_windows)):
            # if iW < 2:
            #     continue
            
            #Subselect data
            indy_t = np.where((ts >= tW[0]) & (ts < tW[1]))[0]
            data_g = spike_counts[g_indy][:,indy_t]
            if network_threshold_type == 'variable':
                network_activity = np.sum(data_g,axis=0)
                network_threshold = np.percentile(network_activity,perc*100)
            # network_activity = np.sum(data_g,axis=0)
            # network_threshold = np.percentile(network_activity,perc*100)
            nT = data_g.shape[1]
            ts_sub = ts[indy_t]
            print(f'Epoch {epoch} has {nT} time bins')

            #Get avalanches
            # D, S, avalanche_indices = get_avalanches(data_g, const_threshold=network_threshold)
            # nAvalanches = len(D)
            r = cr.get_avalanches(data_g, perc, const_threshold=network_threshold)
            S = r['S']; D = r['T']; loc = r['loc']
            # if len(D) - len(loc) == 1:
            #     loc = np.concatenate((loc,[nT-D[-1]]))
            nAvalanches = len(r['S'])
            # avalanche_indices = np.stack((loc,loc+r['T']),axis=1)
            print(f'\t{nAvalanches} avalanches detected')
            if nAvalanches < 100:
                print(f'\tNot enough avalanches detected; skipping')
                continue
            ##------------------------------------------
            fig, axes = plt.subplots(1,2,figsize=(10,5))
            plt.suptitle(f'{mID}, {rec_name}, {g}, {epoch}',y=1)
            #Loop over avalanche size and duration distributions to calculate goodness of fit
            ava_strs = ['size','duration']
            best_fit_pl = []
            for iA, A in enumerate([S, D]):
                F_thresh = args.F_thresh
                print(f'\tCalculating goodness of fit for avalanche {ava_strs[iA]} power-law distributions')
                t0 = time.perf_counter()

                #Find best fit for avalanche size/durations by searching through all possible combinations of xmin and xmax
                uniq_DS = np.unique(A)
                nDS = len(uniq_DS)
                print(f'\t\t{len(uniq_DS)} unique {ava_strs[iA]} values')
                
                #Preallocate for goodness of fit
                F =  np.full((nDS,nDS),np.nan)
                KS = np.full((nDS,nDS),np.nan)
                tau = np.full((nDS,nDS),np.nan)
                plr = np.full((nDS,nDS),np.nan)
                nAv = np.full((nDS,nDS),np.nan)

                #Loop through all combinations of unique avalanche sizes
                uniq_xmm = it.combinations(uniq_DS,2)
                uniq_ijs = it.combinations(np.arange(nDS),2)
                nCombinations = len(list(it.combinations(np.arange(nDS),2)))
                # print(f'\t\t\t{nCombinations} combinations to test')
                if parallel:
                    #Parallelize
                    ray.init(num_cpus=20,include_dashboard=True, dashboard_port=5432)
                    A_id = ray.put(A)

                    #Initiate progress bar
                    pb = util.ProgressBar(nCombinations); actor = pb.actor
                
                processes = []; process_ijs = []
                for (i,j), (xmin,xmax) in zip(uniq_ijs,uniq_xmm):
                    plr[i,j] = np.log10(xmax/xmin)
                    nAv[i,j] = len((S >= xmin) & (S <= xmax))

                    #Calculate goodness of fit
                    if parallel:
                        processes.append(goodness_of_fit_p.remote(A_id, xmin, xmax, pl_distr, actor))
                        process_ijs.append((i,j))
                    else:
                        F[i,j], KS[i,j], tau[i,j] = goodness_of_fit(A, xmin, xmax, pl_distr)

                if parallel:
                    #And then print progress bar until done
                    pb.print_until_done()

                    #Initiate parallel processing
                    results = list(ray.get(processes))

                    #Extract results
                    for (i,j), r in zip(process_ijs,results):
                        F[i,j], KS[i,j], tau[i,j] = r

                    #Shutdown parallel threads
                    ray.shutdown()

                tE = (time.perf_counter() - t0)/60
                print(f'\t\tElapsed time = {tE:.2f} mins')

                ##------------------------------------------
                ## Find best fit
                max_PLR = np.log10(np.max(A)/np.min(A))
                plr_frac = plr/max_PLR
                C = plr_frac + F

                nCandidates = np.sum((F > F_thresh) & (plr > 1))
                F_compare = F_thresh
                while nCandidates == 0:
                    F_compare -= 0.01
                    nCandidates = np.sum((F > F_compare) & (plr > 1))
                    if F_compare < 0:
                        print(f'\t\tNo good fits found; removing F threshold criterion')
                        F_compare = 0
                        skip_fit = True
                        break
                candidate_ijs = []; candidate_plr = []; candidate_KS = []
                for i, j in zip(*np.where((F > F_compare) & (plr > 1))):
                    xmin = uniq_DS[i]
                    xmax = uniq_DS[j]
                    candidate_ijs.append((i,j))
                    candidate_plr.append(plr[i,j])
                candidate_plr = np.array(candidate_plr)

                index = np.argmax(candidate_plr)
                i, j = candidate_ijs[index]
                xmin = uniq_DS[i]; xmax = uniq_DS[j]
                print(f'Best fit: bounds = [{xmin},{xmax}]; F = {F[i,j]:.2f}; # of decades = {plr[i,j]:.2f}; KS = {KS[i,j]:.2e}; tau = {tau[i,j]:.2f}')
                best_fit_pl.append(([xmin,xmax],F[i,j],KS[i,j],tau[i,j]))

                #Save to list for dataframe
                results_list.append([mID,rec_name,g,nG,epoch,ava_strs[iA],pl_distr,xmin,xmax,nAv[i,j],F[i,j],KS[i,j],plr[i,j],tau[i,j]])
                # pdb.set_trace()

                #Save goodness of fit results
                np.savez(join(SaveDir,f'goodness-of-fit_A-{ava_strs[iA]}_{g}_{epoch}.npz'),D=D,S=S,loc=loc,F=F,KS=KS,tau=tau,plr=plr,nAv=nAv,C=C,i=i,j=j,network_threshold=network_threshold)

                x = plr.ravel()
                y = F.ravel()
                c = C.ravel()

                ax = axes[iA]
                ax.set_title(f'Avalanche {ava_strs[iA]} power-law fits')
                ax.scatter(x,y,c=c,marker='.',s=5)
                ax.scatter(plr[i,j],F[i,j],marker='x',s=30,color='k')
                ax.vlines(1,0,1,colors='r',linestyles='--')
                ax.hlines(F_compare,*ax.get_xlim(),colors='k',linestyles='--',label=f'F = {F_compare:.2f}')
                ax.set_xlabel('Number of decades')
                ax.set_ylabel('Fraction of data within bounds')
                usrplt.adjust_spines(ax)
                ax.legend()
            plt.savefig(join(PlotDir,f'power-law-fitting_{ava_strs[iA]}_{g}_{epoch}.png'),dpi=300,bbox_inches='tight',facecolor='w')
            usrplt.save_fig_to_pptx(fig, prs)
            plt.close(fig)

            
            ##------------------------------------------
            #Plot power-law fits
            fig, axes = plt.subplots(1,3,figsize=(12,4))
            plt.suptitle(f'{mID}, {rec_name}, {g}, {epoch}',y=1)

            #Plot Avalanche size distribution
            xmm, F, KS, tau = best_fit_pl[0]
            xmin = int(xmm[0]); xmax = int(xmm[1])
            ax = axes[0]; ax.set_title('Avalanche size')
            pdf = np.histogram(S, bins = np.arange(1, np.max(S) + 2))[0]
            p = pdf/np.sum(pdf)

            ax.scatter(np.arange(1,len(p)+1),p,marker='o',s=10,facecolor='none',edgecolor=epoch_cmap[iW])
            ax.scatter(np.arange(1,len(p)+1)[xmin:xmax],p[xmin:xmax],marker='o',s=10,color=epoch_cmap[iW])
            ax.set_xlabel('Size (# of spikes)')
            ax.set_ylabel('Probablity P(S)')

            #Plot log spaced distribution
            pl_data = pl.Fit(S,xmin_distribution=pl_distr,discrete=False,xmin=xmin,xmax=xmax,pdf_ends_at_xmax=True)
            edges, hist = pl_data.pdf(original_data=True) #linear_bins
            ax.scatter(edges[:-1] + np.diff(edges)/2,hist,marker='o',s=10,color='k')

            #Plot power law model fit
            if pl_distr == 'truncated_power_law':
                params = [pl_data.truncated_power_law.parameter1,pl_data.truncated_power_law.parameter2]
                pl_model = pl.Truncated_Power_Law(xmin=xmin, xmax=xmax, parameters=params)
            elif pl_distr == 'power_law':
                params = [pl_data.power_law.parameter1]
                pl_model = pl.Power_Law(xmin=xmin, xmax=xmax, parameters=params)

            pl_model.plot_pdf(S,ax=axes[0],color='k',ls='--',label=f'\u03C4 = {params[0]:.3f}')
            axes[0].legend()

            #Plot Avalanche duration distribution
            xmm, F, KS, tau_t = best_fit_pl[1]
            tmin = int(xmm[0]); tmax = int(xmm[1])

            ax = axes[1]; ax.set_title('Avalanche duration')
            pdf = np.histogram(D, bins = np.arange(1, np.max(D) + 2))[0]
            p = pdf/np.sum(pdf)
            ax.scatter(np.arange(1,len(p)+1),p,marker='o',s=10,facecolor='none',edgecolor=epoch_cmap[iW])
            ax.scatter(np.arange(1,len(p)+1)[tmin:tmax],p[tmin:tmax],marker='o',s=10,color=epoch_cmap[iW])
            ax.set_xlabel('Duration (s)')
            ax.set_ylabel('Probablity P(T)')

            #Plot log spaced distribution
            pl_data = pl.Fit(D,xmin_distribution=pl_distr,discrete=False,xmin=tmin,xmax=tmax,pdf_ends_at_xmax=True)
            edges, hist = pl_data.pdf(original_data=True) #linear_bins
            ax.scatter(edges[:-1] + np.diff(edges)/2,hist,marker='o',s=10,color='k')

            #Plot power law model fit
            if pl_distr == 'truncated_power_law':
                params = [pl_data.truncated_power_law.parameter1,pl_data.truncated_power_law.parameter2]
                pl_model = pl.Truncated_Power_Law(xmin=tmin, xmax=tmax, parameters=params)
            elif pl_distr == 'power_law':
                params = [pl_data.power_law.parameter1]
                pl_model = pl.Power_Law(xmin=tmin, xmax=tmax, parameters=params)
            pl_model.plot_pdf(S,ax=axes[1],color='k',ls='--',label=f'\u03C4_t = {params[0]:.3f}')

            #Plot scaling relationship btw size and duration
            ax = axes[2]; ax.set_title('"Crackling noise" scaling law')
            ax.scatter(D,S,marker='.',s=10,color=epoch_cmap[iW])
            ax.set_xlabel('Duration (s)')
            ax.set_ylabel('Size (# of spikes)')

            #Fit empirical relationship
            indy = np.where((D >= tmin) & (D <= tmax))[0]
            try:
                params, covariance = curve_fit(power_law_sigma, D[indy], S[indy])
                beta_emp = params[0]; sigma_emp = params[1]
                r2_emp = r2_score(S[indy], power_law_sigma(D[indy], *params))

                x = np.arange(np.min(D),np.max(D),.1)
                y = sigma_emp*x**beta_emp
                ax.plot(x,y,color='y',ls='--',label=f'1/(\u03C3*\u03BD*z) = {beta_emp:.2f}; R2 = {r2_emp:.2f}')
            except:
                beta_emp = np.nan; sigma_emp = np.nan; r2_emp = np.nan

            #Calculate predicted relationship based on tau and tau_t
            beta_pred = (tau_t-1)/(tau-1)
            x = np.arange(np.min(D),np.max(D),.1)
            y = x**beta_pred

            #Fit scalar for loglog fit
            try:
                params, covariance = curve_fit(line_fit, D[indy]**beta_pred, S[indy])
                sigma_pred = params[0]
                x = np.arange(np.min(D),np.max(D),.1)
                y = sigma_pred*x**beta_pred
                r2_pred = r2_score(S[indy], sigma_pred*D[indy]**beta_pred)
                ax.plot(x,y,color='k',ls='--',label=f'(\u03C4_t - 1)/(\u03C4 - 1) = {beta_pred:.2f}; R2 = {r2_pred:.2f}')
            except:
                sigma_pred = np.nan; r2_pred = np.nan
            DCC_list.append([mID,rec_name,g,nG,epoch,beta_emp,sigma_emp,r2_emp,beta_pred,sigma_pred,r2_pred,np.abs(beta_emp-beta_pred)])
            for ax in axes:
                usrplt.adjust_spines(ax)
                ax.set_xscale('log'); ax.set_yscale('log')
                ax.legend()

            plt.savefig(join(PlotDir,f'avalanche-fit_{g}_{epoch}.png'),dpi=300,bbox_inches='tight',facecolor='w')
            usrplt.save_fig_to_pptx(fig, prs)
            plt.close(fig)

    comp_length = time.perf_counter() - t0_outer
    mm, ss = divmod(comp_length,60)
    hh, mm = divmod(mm, 60)
    print(f'Completed in {hh:.0f} hrs, {mm:.0f} minutes, {ss:.0f} seconds')

    ## Save results ----------------------------------------
    DCC_df = pd.DataFrame(DCC_list,columns=['mID','rec_name','group','N','epoch','beta_emp','sigma_emp','r2_emp','beta_pred','sigma_pred','r2_pred','diff'])
    DCC_df = DCC_df.astype({'N':int,'beta_emp':float,'sigma_emp':float,'r2_emp':float,'beta_pred':float,'sigma_pred':float,'r2_pred':float,'diff':float})
    DCC_df['drug_type'] = drug_type2
    DCC_df.to_csv(join(SaveDir,f'DCC_{rec_name}.csv'),index=False)

    pl_df = pd.DataFrame(results_list,columns=['mID','rec_name','group','N','epoch','type','distribution','xmin','xmax','n_avalanches','F','KS','plr','tau'])
    pl_df = pl_df.astype({'N':int,'xmin':int,'xmax':int,'n_avalanches':int,'F':float,'KS':float,'plr':float,'tau':float})
    pl_df['drug_type'] = drug_type2
    pl_df.to_csv(join(SaveDir,f'powerlaw_{rec_name}.csv'),index=False)

    #Create combinations that we are interested in comparing for spontaneous epochs
    cbs = it.combinations(epoch_names,2)
    ijs = it.combinations(np.arange(len(epoch_names)),2)
    combos = []
    for cb,ij in zip(cbs,ijs):
        epoch_i = cb[0];i = ij[0]
        epoch_j = cb[1];j = ij[1]
        if ('pre' in epoch_i):
            if ('post' in epoch_j):
                combos.append(ij)
                print(f'{epoch_i} vs {epoch_j}')
    nComps = len(combos)

    tmp_list = []
    for g in DCC_df['group'].unique():
        for i,j in combos:
            tau0 = DCC_df[(DCC_df['group'] == g) & (DCC_df['epoch'] == epoch_names[i])]['diff'].values[0]
            tau2 = DCC_df[(DCC_df['group'] == g) & (DCC_df['epoch'] == epoch_names[j])]['diff'].values[0]
            tmp_list.append([mID,rec_name,g,epoch_names[i],epoch_names[j],tau0,tau2,tau2-tau0])
    DCC_diff_df = pd.DataFrame(tmp_list,columns=['mID','rec_name','group','epoch_i','epoch_j','DCC_i','DCC_j','diff'])
    DCC_diff_df = DCC_diff_df.astype({'DCC_i':float,'DCC_j':float,'diff':float})

    tmp_list = []
    for g in pl_df['group'].unique():
        for a in pl_df['type'].unique():
            tmp_df = pl_df[(pl_df['group'] == g) & (pl_df['type'] == a)]
            for i, j in combos:
                tau0 = tmp_df.loc[tmp_df['epoch'] == epoch_names[i]]['tau'].values[0]
                tau2 = tmp_df.loc[tmp_df['epoch'] == epoch_names[j]]['tau'].values[0]
                tmp_list.append([mID,rec_name,g,a,epoch_names[i],epoch_names[j],tau0,tau2,tau2-tau0])

    pl_diff_df = pd.DataFrame(tmp_list,columns=['mID','rec_name','group','type','epoch_i','epoch_j','tau_i','tau_j','diff'])
    pl_diff_df = pl_diff_df.astype({'tau_i':float,'tau_j':float,'diff':float})
    DCC_diff_df['drug_type'] = drug_type2
    pl_diff_df['drug_type'] = drug_type2
    DCC_diff_df.to_csv(join(SaveDir,f'DCC_diff_{rec_name}.csv'),index=False)
    pl_diff_df.to_csv(join(SaveDir,f'powerlaw_diff_{rec_name}.csv'),index=False)

    for i, j in combos:
        c_compare = np.stack([epoch_cmap[i],epoch_cmap[j]])
        e_compare = [epoch_names[i],epoch_names[j]]

        ## Plot results ----------------------------------------
        fig, axes = plt.subplots(1,3,figsize=(12,4))
        plt.suptitle(f'{network_threshold_type} threshold for each epoch; {epoch_names[i]} vs {epoch_names[j]}',y=1)
        #Power law exponents for avalanche size
        ax = axes[0]; ax.set_title('Avalanche size')
        sns.boxplot(data=pl_df.loc[(pl_df.type == 'size') & (pl_df.epoch.isin(e_compare))],x='epoch',y='tau',ax=ax,palette=c_compare)
        # ax.boxplot(pl_df.loc[(pl_df.type == 'size') & (pl_df.epoch == epoch_names[0])]['tau'].values,positions=[0],widths=0.5,palette=epoch_cmap[0])
        tau0 = pl_diff_df.loc[(pl_diff_df.type == 'size') & (pl_diff_df.epoch_i == e_compare[0]) & (pl_diff_df.epoch_j == e_compare[1])]['tau_i']
        tau2 = pl_diff_df.loc[(pl_diff_df.type == 'size') & (pl_diff_df.epoch_i == e_compare[0]) & (pl_diff_df.epoch_j == e_compare[1])]['tau_j']
        ax.plot([0,1],[tau0,tau2],'.-',color='k',alpha=0.5,zorder=4)
        ax.set_ylabel('\u03C4')
        res = st.wilcoxon(tau0,tau2,nan_policy='omit'); pval = res.pvalue

        if pval < 0.05:
            ylim = ax.get_ylim()
            ax.text(0.5,0.95*ylim[1],f'*',fontsize=22,fontweight='bold')

        #Power law exponents for avalanche duration
        ax = axes[1]; ax.set_title('Avalanche duration')
        sns.boxplot(data=pl_df.loc[(pl_df.type == 'duration') & (pl_df.epoch.isin(e_compare))],x='epoch',y='tau',ax=ax,palette=c_compare)
        ax.set_ylabel('\u03C4_t')
        tau0 = pl_diff_df.loc[(pl_diff_df.type == 'duration') & (pl_diff_df.epoch_i == e_compare[0]) & (pl_diff_df.epoch_j == e_compare[1])]['tau_i']
        tau2 = pl_diff_df.loc[(pl_diff_df.type == 'duration') & (pl_diff_df.epoch_i == e_compare[0]) & (pl_diff_df.epoch_j == e_compare[1])]['tau_j']
        ax.plot([0,1],[tau0,tau2],'.-',color='k',alpha=0.5,zorder=4)
        res = st.wilcoxon(tau0,tau2,nan_policy='omit'); pval = res.pvalue
        if pval < 0.05:
            ylim = ax.get_ylim()
            ax.text(0.5,0.95*ylim[1],f'*',fontsize=22,fontweight='bold')

        #Deviation from criticality
        ax = axes[2]; ax.set_title('Deviation from criticality')
        sns.boxplot(data=DCC_df.loc[DCC_df.epoch.isin(e_compare)],x='epoch',y='diff',ax=ax,palette=c_compare)
        ax.set_ylabel('DCC')
        tau0 = DCC_diff_df.loc[(DCC_diff_df.epoch_i == e_compare[0]) & (DCC_diff_df.epoch_j == e_compare[1])]['DCC_i']
        tau2 = DCC_diff_df.loc[(DCC_diff_df.epoch_i == e_compare[0]) & (DCC_diff_df.epoch_j == e_compare[1])]['DCC_j']
        print(len(tau2))
        ax.plot([0,1],[tau0,tau2],'.-',color='k',alpha=0.5,zorder=4)
        res = st.wilcoxon(tau0,tau2,nan_policy='omit'); pval = res.pvalue
        if pval < 0.05:
            ylim = ax.get_ylim()
            ax.text(0.5,0.95*ylim[1],f'*',fontsize=22,fontweight='bold')
        for ax in axes:
            usrplt.adjust_spines(ax)
            # ax.set_xticklabels([epoch_names[0],epoch_names[2]])
            ax.set_xlabel('Epoch')

        plt.savefig(join(PlotDir,f'pl-diff_{network_threshold_type}-thresh_{rec_name}_{epoch_names[i]}_{epoch_names[j]}.png'),dpi=300,bbox_inches='tight',facecolor='w')    
        usrplt.save_fig_to_pptx(fig, prs)
        plt.close(fig)

    ## CLOSE files ----------------------------------------
    prs.save(join(PlotDir,f'avalanche-fits-{rec_name}.pptx'))

    




