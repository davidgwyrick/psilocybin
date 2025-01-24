base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
base_dir_local = '/data/tiny-blue-dot/zap-n-zip/EEG_exp'; base_dir = base_dir_server
ProjDir = '/home/david.wyrick/projects/zap-n-zip/'
ServDir = '/data/projects/zap-n-zip/'

#Base
import argparse
from glob import glob
from os.path import join
import json, os, time, sys
import gspread
import pandas as pd
import numpy as np
import itertools as it
from tqdm import tqdm, trange

#Scipy
import scipy.signal as sig
import scipy.stats as st
from scipy.ndimage import gaussian_filter
from scipy.ndimage import gaussian_filter1d
from scipy.interpolate import interp1d
from sklearn.linear_model import LinearRegression, PoissonRegressor

#Plot
import matplotlib.pyplot as plt
import seaborn as sns
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation

#Project
from tbd_eeg.data_analysis.eegutils import EEGexp
import tbd_eeg.data_analysis.Utilities.utilities as tbd_util

#User
sys.path.append(ProjDir)
import util
import plotting as usrplt

#Templeton-log_exp
gc = gspread.service_account() 
sh = gc.open('Templeton-log_exp') 
exp_table = pd.DataFrame(sh.sheet1.get()).T.set_index(0).T 
exp_table = exp_table.set_index('mouse_name')
exp_table.head()

bandpass_freqs = np.array([[1, 4], [4, 8], [8,13], [13, 30], [30, 55], [65, 100], [100,200]])
bandpass_str = ['delta','theta','alpha','beta','low-gamma','high-gamma','HFO']
n_bands = len(bandpass_freqs)

## Parse Command Line Arguments ----------------------------------------
parser = argparse.ArgumentParser(description='Rolling RG analysis')

parser.add_argument('--mID',type=str, default='mouse735052',
                    help='mouse to perform analysis on')

parser.add_argument('--rec_name',type=str, default='aw_psi_2024-06-20_10-40-56',
                    help='experiment to perform analysis on')

parser.add_argument('--level',type=str, default='group',
                    help='Level at which to group neurons')

parser.add_argument('--time_bin_ms',type=int, default=50,
                    help='Time bin width')

parser.add_argument('--tWindow_width_s',type=int, default=240,
                    help='Time window over which to calculate FC')

parser.add_argument('--tWindow_shift_s',type=float, default=60,
                    help='Amount of time to shift rolling window by')

parser.add_argument('--fr_thresh',type=float, default=0,
                    help='Firing rate threshold')

from scipy.optimize import curve_fit
from sklearn.metrics import r2_score
from statsmodels.tsa.stattools import acf
import powerlaw as pl

# Define the exponential function
def exponential_func(x, tau):
    return np.exp(-1 * x / tau)

# Define the power-law function
def power_law(x, exponent):
    return  x ** exponent

# Define the power-law function
def power_law_sigma(x, exponent, sigma):
    return  sigma* (x ** exponent)

def get_pairs(corr_mat):
    # Check if number of neurons is even
    n = corr_mat.shape[0]
    if n % 2 != 0:
        raise ValueError('Number of neurons must be even')

    #Sort correlation matrix
    ijs = np.unravel_index(np.argsort(corr_mat.ravel())[::-1],corr_mat.shape)
    pairs = []; cc_list = []
    paired_neurons = np.zeros(n, dtype=bool)
    for i,j in zip(*ijs):
        #Check to see if either neuron has been paired
        if paired_neurons[i] or paired_neurons[j] or i == j:
            continue
        paired_neurons[i] = True; paired_neurons[j] = True
        pairs.append((i,j)); cc_list.append(corr_mat[i,j])

    return pairs, cc_list

def calculate_acf(X, lag_sec=2, time_bin=0.05):
    n, T = X.shape
    lags = np.arange(0,lag_sec+time_bin,time_bin)

    #Calculate autocorrelation for each neuron
    autocorrs = [acf(X[i],nlags=np.ceil(lag_sec/time_bin)) for i in range(n)]
    mean_ac = np.nanmean(np.array(autocorrs),axis=0)

    xdata = lags; ydata = mean_ac
    # Fit the exponential curve
    popt, pcov = curve_fit(exponential_func, xdata, ydata)#, p0=(1, 0.1, 1))
    tau = popt[0]

    # Generate fitted data
    y_fit = exponential_func(xdata, tau)
    r2 = r2_score(ydata, y_fit)

    xdense = np.linspace(xdata[0], xdata[-1], 500)
    ydense = exponential_func(xdense, tau)

    # Plot the data and the fit
    fig, ax = plt.subplots(figsize=(6, 6))
    ax.plot(xdata, ydata, 'o', label='Data')
    ax.plot(xdense, ydense, '-', label=f'Exponential fit: \u03C4 = {tau: .2f} , R2 = {r2:.2f} ')
    ax.set_ylabel('Auto-Correlation')
    ax.set_xlabel('Lag (s)')
    ax.legend()
    usrplt.adjust_spines(ax)
    return lags, mean_ac, tau, fig #, (xdense, ydense)

# def determine_time_bin(data):

    
#     return time_bin

## Main ----------------------------------------
if __name__ == '__main__':

    ## Parse the arguments ----------------------------------------
    args = parser.parse_args()

    #Which experiment?
    mID = args.mID
    rec_name = args.rec_name
    print(f'{mID}, {rec_name}')

    #How to segment data
    time_bin_ms = args.time_bin_ms
    time_bin = time_bin_ms/1000
    tWindow_width_s = args.tWindow_width_s
    tWindow_width = tWindow_width_s #*60
    tWindow_shift = args.tWindow_shift_s
    fr_thresh = args.fr_thresh
    level = args.level

    ## FOLDERS ----------------------------------------
    #Create directory for saving to
    TempDir = os.path.join(ServDir,'results','RG1',mID,rec_name)
    if not os.path.exists(TempDir):
        os.makedirs(TempDir)

    tmp_list = sorted(glob(join(TempDir,f'run_*')))
    if len(tmp_list) == 0:
        curr_run = 0
    else:
        last_run = int(tmp_list[-1][-1])
        curr_run = last_run+1

    folder = f'run_{curr_run:02d}'
    SaveDir = os.path.join(TempDir,folder)
    if not os.path.exists(SaveDir):
        os.makedirs(SaveDir)
    
    PlotDir = join(SaveDir,'plots')
    if not os.path.exists(PlotDir):
        os.makedirs(PlotDir)

    #Save model parameters
    args_dict = args.__dict__
    args_dict['SaveDir'] = SaveDir
    with open(join(SaveDir,f"ccm_parameters_run_{curr_run}.json"), "w") as outfile:
        json.dump(args_dict, outfile)

    #Extract row from templeton dataframe
    exp_df = exp_table.loc[exp_table.exp_name == rec_name]
    if type(exp_df) == pd.core.series.Series:
        exp_df = exp_df.to_frame().T

    drug_type = exp_df['drug'].values[0]
    stim_type = exp_df['stimulation'].values[0]
    print(f'Experiment type: {stim_type}, {drug_type}')

    # Powerpoint for figures
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f'{mID} {rec_name}'
    slide.placeholders[1].text = f'Phenomenological Renormalization Group Approach'

    ## Read in experiment ----------------------------------------
    #Upload the whole experiment and generate the global clock
    exp = EEGexp(os.path.join(base_dir,mID,rec_name,'experiment1','recording1'), preprocess=False, make_stim_csv=False)

    #Load probe data
    probe_unit_data, probe_info, total_units, metric_list = tbd_util.get_neuropixel_data(exp)
    probe_list = [x.replace('_sorted', '') for x in exp.experiment_data if 'probe' in x]

    #Get recording start time
    probei = probe_list[0]
    open_ephys_start = np.round(np.min(probe_unit_data[probei]['spike_times']))+10
    open_ephys_end = np.round(np.max(probe_unit_data[probei]['spike_times']))-10
    recording_length = open_ephys_end - open_ephys_start
    mm, ss = divmod(recording_length,60)
    hh, mm = divmod(mm, 60)
    print(f'{hh} hrs, {mm} minutes, {ss} seconds')

    #Bin spiking data
    tW = [open_ephys_start,open_ephys_end]
    data_list, ts_list, neuron_indices, plot_tuple, FR_perblock = util.bin_spiking_data(probe_unit_data,[tW], time_bin=time_bin,fr_thresh=0)
    spike_counts = data_list[0].T; ts = ts_list[0]
    
    if len(plot_tuple) == 11:
        boundaries_group, ticks_group, labels_group, celltypes, durations, layers, areas, groups,mesogroups, supergroups, order_by_group = plot_tuple
    else:
        print('Experiment has no area information. Check to see if histology is completed. Exiting')
        exit()
    nNeurons = len(neuron_indices)
    spike_counts_sub = spike_counts[neuron_indices]
    spike_counts_ro = spike_counts_sub[order_by_group]
    areas_sub = areas[neuron_indices]; areas_ro = areas_sub[order_by_group]
    groups_sub = groups[neuron_indices]; groups_ro = groups_sub[order_by_group]
    mesogroups_sub = mesogroups[neuron_indices]; mesogroups_ro = mesogroups_sub[order_by_group]
    supergroups_sub = supergroups[neuron_indices]; supergroups_ro = supergroups_sub[order_by_group]

    #Just sort data first
    areas = areas_ro; groups = groups_ro; mesogroups = mesogroups_ro; supergroups = supergroups_ro
    spike_counts = spike_counts_ro


    #Read in behavior
    run_ts, run_signal, run_signal_s, acceleration, pupil_ts, pupil_radius, pupil_dxdt, plot_pupil = util.get_behavioral_data(exp, mID, rec_name)
    f_run = interp1d(run_ts,run_signal); run_signal_p = f_run(pupil_ts)
    f_run_s = interp1d(run_ts,run_signal_s); run_signal_p_s = f_run_s(pupil_ts)
    f_pupil = interp1d(pupil_ts,pupil_radius)
    
    injection_times = None; injection_time_windows = None; inj_tuple = None
    #For saline & psilocybin experiments, get injection times and types of injection
    if drug_type in ['saline', 'psilocybin','ketanserin+psilocybin','urethane+psilocybin']:
        injection_time_windows = np.array([np.array(exp_df['First injection window'].values[0].split(','),dtype=float),
                                np.array(exp_df['Second injection window'].values[0].split(','),dtype=float)])
        
        #Take second time in each window as "injection time"
        injection_times = np.array([exp_df['First injection time (s)'].values[0],exp_df['Second injection time (s)'].values[0]],dtype=float)

        if drug_type in ['psilocybin','urethane+psilocybin']:
            injection_types = ['sal1','psi']
            injection_colors = sns.xkcd_palette(['dark sky blue','darkish red'])
            macro_names = ['pre-inj','post-sal1-inj','post-psi-inj']
            cmap = sns.xkcd_palette(['silver','dark sky blue','darkish red'])
        elif drug_type == 'saline':
            injection_types = ['sal1', 'sal2']
            injection_colors = sns.xkcd_palette(['dark sky blue','cobalt blue'])
            macro_names = ['pre-inj','post-sal1-inj','post-sal2-inj']
            cmap = sns.xkcd_palette(['silver','dark sky blue','cobalt blue'])
        elif drug_type == 'ketanserin+psilocybin':
            injection_types = ['ket','psi']
            injection_colors = sns.xkcd_palette(['magenta','goldenrod'])
            macro_names = ['pre-inj','post-ket-inj','post-psi-inj']
            cmap = sns.xkcd_palette(['silver','magenta','goldenrod'])
        inj_tuple = (injection_times, injection_types, injection_colors)
    else:
        injection_times = None

    #For isoflurane experiments, get iso level
    iso_induction_times = None; iso_maintenance_times = None; iso_tuple = None
    if drug_type == 'isoflurane':
        iso_level, iso_times = exp.load_analog_iso()
        iso_induction_times, iso_maintenance_times = exp.load_iso_times()
        induction_colors = sns.xkcd_palette(['light teal','teal'])
        iso_tuple = (iso_induction_times, induction_colors)
        macro_names = ['pre-iso','iso-ind','post-iso']
        cmap = sns.xkcd_palette(['silver','light teal','teal'])

    # extract the timestamps of the selected stimuli
    try:
        stim_log = pd.read_csv(exp.stimulus_log_file)
        stim_exists = True
    except:
        stim_exists = False

    ## Determine windows ----------------------------------------
    evoked_time_window_list = []; evoked_type_list = []; evoked_tuple = None
    if stim_exists:
        for s in np.unique(stim_log['sweep']):
            for t in np.unique(stim_log.loc[stim_log.sweep == s]['stim_type']):
                sub_df = stim_log.loc[(stim_log.sweep == s) & (stim_log.stim_type == t)]
                tS = sub_df['onset'].min(); tE = sub_df['offset'].max()
                
                evoked_time_window_list.append([tS,tE])
                evoked_type_list.append(t)
        evoked_tuple = (evoked_time_window_list,evoked_type_list)
    
    #Determine rolling windows
    tW_starts = np.arange(open_ephys_start,open_ephys_end-tWindow_width,tWindow_shift)
    tW_ends = tW_starts + tWindow_width
    time_window_array = np.array((tW_starts,tW_ends)).T
    time_window_list = time_window_array.tolist()
    time_window_centers = time_window_array[:,0] + tWindow_width/2
    nWindows = len(time_window_list)
    print(f'{nWindows} windows, {tWindow_width} seconds long, separated by {tWindow_shift} sec')

    # exit()
    #Create epoch names
    epoch_list = []
    for ii, tW in enumerate(time_window_list):
        tW_center = tW[0] + (tW[1]-tW[0])
 
        #Determine whether most of the data is in evoked or spontaneous periods
        window_type = 'spont'
        for epoch_type, epoch_window in zip(evoked_type_list,evoked_time_window_list):
            if (tW_center >= epoch_window[0]) & (tW_center < epoch_window[1]):
                window_type = epoch_type
                break
        
        # epoch_list.append(f'{window_type}-{ii:03d}_no-inj')
        if drug_type in ['saline', 'psilocybin','ketanserin+psilocybin','urethane+psilocybin']: 
            if tW_center < injection_times[0]:
                epoch_list.append(f'{window_type}-{ii:03d}_pre-inj')
            elif (tW_center >= injection_times[0]) & (tW_center < injection_times[1]):
                epoch_list.append(f'{window_type}-{ii:03d}_post-{injection_types[0]}-inj')
            else:
                epoch_list.append(f'{window_type}-{ii:03d}_post-{injection_types[1]}-inj')
        elif drug_type == 'isoflurane':
            if tW_center < iso_induction_times[0]:
                epoch_list.append(f'{window_type}-{ii:03d}_pre-iso')
            elif (tW_center >= iso_induction_times[0]) & (tW_center < iso_induction_times[1]):
                epoch_list.append(f'{window_type}-{ii:03d}_iso-ind')
            else:
                epoch_list.append(f'{window_type}-{ii:03d}_post-iso')
        elif drug_type == 'urethane':
            t1 = int(tW[0]/60)
            t2 = int(tW[1]/60)
            epoch_list.append(f'urethane_{window_type}-{ii:03d}_{t1}-{t2}')
        else:
            t1 = int(tW[0]/60)
            t2 = int(tW[1]/60)
            epoch_list.append(f'{window_type}-{ii:03d}_{t1}-{t2}')
    if drug_type == 'urethane': 
        e_list = [f'urethane' for e in epoch_list]
        macro_names = ['urethane']
        cmap = sns.xkcd_palette(['orange'])
    else:
        e_list = np.array([e.split('_')[-1] for e in epoch_list])  
    #Calculate behavioral measures
    running_moments = []
    pupil_moments = []
    for tW in time_window_list:
        indy = np.where((run_ts > tW[0]) & (run_ts <= tW[1]))
        rs_m = np.nanmean(run_signal[indy])
        rs_std = st.iqr(run_signal[indy])
        running_moments.append((rs_m,rs_std))
        if plot_pupil:
            indy = np.where((pupil_ts > tW[0]) & (pupil_ts <= tW[1]))
            pd_m = np.nanmean(pupil_radius[indy])
            pd_std = st.iqr(pupil_radius[indy])
            pupil_moments.append((pd_m,pd_std))
    running_moments = np.array(running_moments)
    pupil_moments = np.array(pupil_moments)

    #Plot behavior
    fig = usrplt.plot_behavior((run_ts,run_signal),(pupil_ts,pupil_radius),f'{mID} {rec_name}',evoked_tuple,inj_tuple,iso_tuple)
    plt.savefig(join(PlotDir,f'behavior_{rec_name}.png'),dpi=300,bbox_inches='tight',facecolor='w')
    usrplt.save_fig_to_pptx(fig, prs)
    plt.close(fig)

    #Save time windows 
    np.savez(os.path.join(SaveDir,f'time_windows.npz'),time_window_list=time_window_list,epoch_list=epoch_list,running_moments=running_moments,pupil_moments=pupil_moments)
    import warnings
    warnings.filterwarnings("ignore")
    ## Calculate rolling RG ----------------------------------------
    t0_outer = time.perf_counter()
    lag_sec = 1; sigma = 1
    pl_list = []

    if level == 'supergroup':
        groups_to_iterate = np.unique(supergroups)
    elif level == 'group':
        groups_to_iterate = np.unique(groups)
    elif level == 'all':
        groups_to_iterate = ['all']

    for g in groups_to_iterate:
        if g in ['root']:
            continue
        if level == 'supergroup':
            g_indy = np.where(supergroups == g)[0]

            groups_sg = groups[g_indy]
            unique_sorted, uniq_indices = np.unique(groups_sg, return_index=True)
            ylabels = unique_sorted[np.argsort(uniq_indices)]
            yorder = []
            yticks = []; count = 0; ybounds = []
            for a in ylabels:
                indy = np.where(groups_sg == a)[0]
                yticks.append(count + len(indy)//2)
                yorder.append(indy)
                ybounds.append(count + len(indy))
                count += len(indy)
            yorder = np.concatenate(yorder)

        elif level == 'group':
            g_indy = np.where(groups == g)[0]

            areas_g = areas[g_indy]
            unique_sorted, uniq_indices = np.unique(areas_g, return_index=True)
            ylabels = unique_sorted[np.argsort(uniq_indices)]
            yorder = []
            yticks = []; count = 0; ybounds = []
            for a in ylabels:
                indy = np.where(areas_g == a)[0]
                yticks.append(count + len(indy)//2)
                yorder.append(indy)
                ybounds.append(count + len(indy))
                count += len(indy)
            yorder = np.concatenate(yorder)

        elif level == 'all':
            g_indy = np.arange(nNeurons)
            yticks = ticks_group
            ylabels = labels_group
            ybounds = boundaries_group
            yorder = np.arange(nNeurons)

        n = len(g_indy)
        print(f'Group {g} has {n} neurons')
        if n < 128:
            print('\tGroup too small, skipping')
            continue

        exponents = np.full((nWindows,4),np.nan)
        #Loop over different time blocks
        for iW, tW in enumerate(tqdm(time_window_list)):
            epoch = epoch_list[iW]
            t_indy = np.where((ts >= tW[0]) & (ts < tW[1]))[0]

            # import pdb; pdb.set_trace()
            #Get spike counts for this window
            X = spike_counts[g_indy][:,t_indy]
            n, T = X.shape
            N_g = n

            #Plot spiking activity
            vmax = np.ceil(np.nanpercentile(X.ravel(),98))
            fig, axes = plt.subplots(2,1,figsize=(8,6),gridspec_kw={'height_ratios':[6,1],'hspace':0.3})
            plt.suptitle(f'{epoch}; {n} {g} neurons',y=0.975)

            ax = axes[0]
            ax.set_title(f'\u0394T = {time_bin:.3f}s, clims = [0, {vmax}]')
            sns.heatmap(X[yorder],cmap='gray_r',vmin=0,vmax=vmax,cbar=False,ax=ax)
            ax.set_ylabel('Neural activity')
            ax.set_xlabel('Time in window (s)')

            # ax.vlines(boundaries,0,n,colors='r',linestyles='--')
            # yticks = np.arange(0,n+1,n//2)
            xticks = np.arange(0,T+1,int(60/time_bin))
            ax.set_yticks(yticks); ax.set_yticklabels(ylabels)
            if level == 'group':
                ax.hlines(ybounds,0,T,colors='r',linestyles='--',lw=1)
            ax.set_yticks(ybounds, minor=True)
            ax.set_xticks(xticks); ax.set_xticklabels(np.round(xticks*time_bin,0),rotation=30)

            ax = axes[1]
            # ax.plot(ts[t_indy],np.sum(X,axis=0),color='k')
            # ax.set_ylabel('Total activity')  
            t_indy = np.where((run_ts >= tW[0]) & (run_ts < tW[1]))[0]
            ax.plot(run_ts[t_indy]/60,run_signal[t_indy],color='k')
            ax.set_ylabel('Speed')
            ax.set_xlabel('Time in experiment (min)')
            ax.autoscale(tight=True)
            # xticks = np.arange(0,T+1,int(60/time_bin))
            if plot_pupil:
                ax2 = ax.twinx()
                t_indy = np.where((pupil_ts >= tW[0]) & (pupil_ts < tW[1]))[0]
                ax2.plot(pupil_ts[t_indy]/60,pupil_radius[t_indy],color='b')
                ax2.set_ylabel('Pupil',color='b')

            #Save figure
            plt.savefig(join(PlotDir,f'raster-plot_{g}_{epoch}.png'),dpi=300,bbox_inches='tight',facecolor='w')
            usrplt.save_fig_to_pptx(fig, prs)
            plt.close(fig)
            # exit()
            #Calculate block neuron sizes
            Ks = []; n_k = n; RG = 0
            while n > 0:
                Ks.append(2**RG)
                n_k = n // 2
                n = n_k
                RG+=1
            Ks = np.array(Ks)
            nRG = len(Ks)
            # print(f'\tRG-{RG:2d}: Block neuron sizes: {Ks}')

            #Preallocate arrays
            lags = np.arange(0,lag_sec+time_bin,time_bin)
            variance = np.full((nRG),np.nan)
            free_energy = np.full((nRG),np.nan)
            frac_silent = np.full((nRG),np.nan)
            ac_tau = np.full((nRG),np.nan)
            autocorr = np.full((nRG,len(lags)),np.nan)
                    
            #Calculate variance of actual neurons
            variance[0] = np.mean(np.var(X,axis=1))

            #Calculate free energy of actual neurons
            Sk = np.sum(X == 0, axis=1)/T
            fe = -1*np.log10(Sk)
            indy = np.where(~np.isinf(fe) & ~np.isnan(fe))[0]
            free_energy[0] = np.mean(fe[indy])
            frac_silent[0] = np.mean(Sk)

            #Calcuate autocorrelation of actual neurons
            lags, mean_ac, tau, fig = calculate_acf(X, lag_sec=lag_sec,time_bin=time_bin)
            # plt.suptitle(f'{g} neurons, RG-0',y=0.95) 
            # plt.savefig(join(PlotDir,f'autocorrelation_{g}_RG-0_{epoch}.png'),dpi=300,bbox_inches='tight',facecolor='w')
            # usrplt.save_fig_to_pptx(fig, prs)
            # plt.close(fig)

            # autocorr[index,0] = mean_ac
            ac_tau[0] = tau

            #Temporal coarse graining
            for iRG in range(1,nRG):
                n, T = X.shape
                #If odd # of neurons, remove neuron with lowest firing rate
                if n % 2 == 1:
                    FR_g = np.sum(X,axis=1)/(X.shape[1]*time_bin)
                    mask = np.ones(n,dtype=bool)
                    indy = np.argmin(FR_g)
                    mask[indy] = False
                    X = X[mask]
                n, T = X.shape

                #Calculate correlation matrix
                corr_mat = np.corrcoef(X)
                # corr_mat[np.isnan(corr_mat)] = 0
                # np.save(join('./',f'corr_mat.npy'),corr_mat)
                # fig, axes = plt.subplots(1,2,figsize=(12,6))
                # vmax = np.nanpercentile(np.abs(corr_mat),98)
                # sns.heatmap(np.abs(corr_mat),ax=axes[0],square=True,cmap='viridis',vmin=0,vmax=vmax,cbar_kws={'label':'| Correlation |','shrink':0.5})
                # axes[0].set_title('Correlation matrix')

                # sns.histplot(corr_mat[np.triu_indices(n,1)],ax=axes[1],stat='density')
                # usrplt.adjust_spines(axes[1])
                # axes[1].set_title('Correlation distribution')
                # plt.suptitle(f'{g} neurons, RG-{iRG}',y=0.95) 
                # plt.savefig(join(PlotDir,f'correlation_{g}_RG-{iRG}_{epoch}.png'),dpi=300,bbox_inches='tight',facecolor='w')
                # usrplt.save_fig_to_pptx(fig, prs)
                # plt.close(fig)

                #Perform first temporal coarse graining
                pairs = get_pairs(corr_mat.copy())

                #Combine pairs
                n_k = n//2
                X_k = np.zeros((n//2,T))
                X_norm = np.zeros((n//2,T))
                corr_rank = []
                for ii, (i,j) in enumerate(pairs):
                    corr_rank.append(corr_mat[i,j])
                    tmp = X[i] + X[j]
                    non_zero_mean = np.nanmean(tmp[tmp > 0])
                    X_k[ii] = tmp
                    X_norm[ii] = (1/non_zero_mean)*tmp

                #Calculate variance of block neurons
                variance[iRG] = np.mean(np.var(X_k,axis=1))

                #Calculate free energy of block neurons
                Sk = np.sum(X_norm == 0, axis=1)/T
                fe = -1*np.log10(Sk)
                indy = np.where(~np.isinf(fe) & ~np.isnan(fe))[0]
                free_energy[iRG] = np.mean(fe[indy])
                frac_silent[iRG] = np.mean(Sk)

                #Calcuate autocorrelation of block neurons
                lags, mean_ac, tau, fig = calculate_acf(X_norm, lag_sec=lag_sec,time_bin=time_bin)
                # plt.suptitle(f'{g} neurons, RG-{iRG}',y=0.95)
                # plt.savefig(join(PlotDir,f'autocorrelation_{g}_RG-{iRG}_{epoch}.png'),dpi=300,bbox_inches='tight',facecolor='w')
                # usrplt.save_fig_to_pptx(fig, prs)
                # plt.close(fig)

                autocorr[iRG] = mean_ac
                ac_tau[iRG] = tau

                #Reset X
                X = X_k
                n, T = X.shape

            ##------------------------------------------
            normalized_var = variance/variance[0]

            #Fit power law
            try:
                params, covariance = curve_fit(power_law, Ks, normalized_var,maxfev=10000)
                r2 = r2_score(normalized_var, power_law(Ks, *params))
                alpha = params[0]; exponents[iW,0] = alpha
                pl_list.append((iW,epoch,g,time_bin,'alpha',alpha,r2))

                fig, ax = plt.subplots(figsize=(6,6))
                plt.suptitle(f'{g}: Activity variance, {epoch}',y=0.95)
                ax.plot(Ks,normalized_var,'o')
                ax.plot(Ks,(Ks**alpha),'--k',label=f'alpha = {alpha:.3f}; R2 = {r2:.3f}')

                ax.plot(Ks,Ks**1,':k')
                ax.plot(Ks,Ks**2,':k')
                ax.legend()
                ax.set_xscale('log')
                ax.set_yscale('log')
                ax.set_ylabel('Normalized variance')
                ax.set_xlabel('K')
                usrplt.adjust_spines(ax)
                usrplt.adjust_spines(ax)

                # plt.savefig(join(PlotDir,f'variance_{g}_{epoch}.png'),dpi=300,bbox_inches='tight',facecolor='w')
                usrplt.save_fig_to_pptx(fig, prs)
                plt.close(fig)
            except:
                print(f'Could not fit power law to variance for {g}, {epoch}')

                exponents[iW,0] = np.nan
        
            ##------------------------------------------
            normalized_fe = free_energy/free_energy[0]

            #Fit power law
            indy = np.where(~np.isnan(normalized_fe) & ~np.isinf(normalized_fe))[0]
            try:
                params, covariance = curve_fit(power_law, Ks[indy], normalized_fe[indy],maxfev=10000)
                r2 = r2_score(normalized_fe[indy], power_law(Ks[indy], *params))
                alpha = params[0]; exponents[iW,1] = alpha
                pl_list.append((iW,epoch,g,time_bin,'beta',alpha,r2))

                fig, ax = plt.subplots(figsize=(6,6))
                plt.suptitle(f'{g}: Free energy, {epoch}',y=0.95)

                ax.plot(Ks,normalized_fe,'o')
                ax.plot(Ks,(Ks**alpha),'--k',label=f'alpha = {alpha:.3f}; R2 = {r2:.3f}')
                ax.plot(Ks,Ks**1,':k')
                ax.legend()

                ax.set_xscale('log')
                ax.set_yscale('log')
                ax.set_ylabel('Normalized free energy')
                ax.set_xlabel('K')
                usrplt.adjust_spines(ax)

                # plt.savefig(join(PlotDir,f'free-energy_{g}_{epoch}.png'),dpi=300,bbox_inches='tight',facecolor='w')
                usrplt.save_fig_to_pptx(fig, prs)
                plt.close(fig)
            except:
                print(f'Could not fit power law to free energy for {g}, {epoch}')
                exponents[iW,1] = np.nan

            ##------------------------------------------
            normalized_ac = ac_tau/ac_tau[1]

            #Fit power law
            try:
                params, covariance = curve_fit(power_law, Ks, normalized_ac,maxfev=10000)
                r2 = r2_score(normalized_ac, power_law(Ks, *params))
                alpha = params[0]; exponents[iW,2] = alpha
                pl_list.append((iW,epoch,g,time_bin,'zeta',alpha,r2))

                fig, ax = plt.subplots(figsize=(6,6))
                plt.suptitle(f'{g}: Autocorrelation',y=0.95)
                ax.plot(Ks,normalized_ac,'o')
                ax.plot(Ks,(Ks**alpha),'--k',label=f'alpha = {alpha:.3f}; R2 = {r2:.3f}')
                ax.legend()

                ax.set_xscale('log')
                ax.set_yscale('log')
                ax.set_ylabel('Normalized autocorrelation')
                ax.set_xlabel('K')
                usrplt.adjust_spines(ax)

                # plt.savefig(join(PlotDir,f'autocorrelation_{g}_{epoch}.png'),dpi=300,bbox_inches='tight',facecolor='w')
                usrplt.save_fig_to_pptx(fig, prs)
                plt.close(fig)
                # import pdb; pdb.set_trace()
            except:
                print(f'Could not fit power law to autocorrelation for {g}, {epoch}')
                exponents[iW,2] = np.nan

            #Save to file
            np.savez(join(SaveDir,f'RG-metrics_{g}_{epoch}.npz'),ac_tau=ac_tau,variance=variance,free_energy=free_energy,frac_silent=frac_silent,autocorr=autocorr)
        
        np.savez(join(SaveDir,f'RG-exponents_{g}.npz'),exponents=exponents,Ks=Ks,macro_names=macro_names,epoch_list=epoch_list,e_list=e_list,running_moments=running_moments,pupil_moments=pupil_moments)

        ## ----------------------------------------
        pl_strs = ['alpha','beta','zeta']
        pl_titles = ['Activity variance','Free energy','Autocorrelation']
        for iE in range(3):

            #Plot exponent timecourse
            fig, axes = plt.subplots(2,1,figsize=(11,8))#,gridspec_kw={'height_ratios':[3,1]})
            plt.suptitle(f'{N_g} {g} neurons, Ks = {Ks}',y=0.95)
            ax = axes[0]; ax.set_title(f'Power-law exponent for {pl_titles[iE]}')
            ax.plot(exponents[:,iE],label=f'{pl_strs[iE]}',lw=2)

            #Plot running speed
            ax = axes[1]; ax.set_title('Behavioral measures')
            cc = np.corrcoef(running_moments[:,0],exponents[:,iE])[0,1]
            ax.plot(run_ts/60,run_signal,'-k',lw=1,label=f'\u03C1  (speed, scaling exponent) = {cc:.2f}')
            ax.set_xlabel('Time (min)')
            ax.set_ylabel('Running speed (cm/s)')
            ax.legend()
            if injection_times is not None:
                for ii, t in enumerate(injection_times):
                    ax.axvline(t/60,color=injection_colors[ii],ls='--',lw=3)#,label=f'{injection_types[ii]} injection',zorder=4)
                ax.legend(loc=1)

            if iso_induction_times is not None:
                for ii, t in enumerate(iso_induction_times):
                    ax.axvline(t/60,color='c',ls='--',lw=3)#,label=f't-{ii} induction',zorder=4)
                ax.legend(loc=1)                 
            if plot_pupil:
                ax2 = ax.twinx()
                cc = np.corrcoef(pupil_moments[:,0],exponents[:,iE])[0,1]
                ax2.plot(pupil_ts/60,pupil_radius,'-b',lw=1,label=f'\u03C1 (pupil size, scaling exponent) = {cc:.2f}')
                ax2.set_ylabel('Pupil radius (pixels)',color='b')
                ax2.legend(loc=2)

            plt.savefig(join(PlotDir,f'RG_{pl_strs[iE]}_exponent_{g}_{rec_name}.png'),dpi=300,bbox_inches='tight',facecolor='w')
            plt.savefig(join(ProjDir,'plots','RG1',f'RG_{pl_strs[iE]}_exponent_{g}_{rec_name}.png'),dpi=300,bbox_inches='tight',facecolor='w')
            usrplt.save_fig_to_pptx(fig, prs)
            plt.close(fig)

            if drug_type != 'urethane':
                #Plot distribution of exponents per window
                fig, ax = plt.subplots(figsize=(6,6))
                # plt.suptitle(f'{g}: Power-law exponents',y=0.95)
                distr_list = [exponents[np.where(e_list == e)[0],iE] for e in macro_names]
                mean_list = [np.nanmean(d) for d in distr_list]
                sns.histplot(distr_list,ax=ax,palette=cmap,multiple='layer',stat='density',kde=True)
                for c, m, e in zip(cmap,mean_list, macro_names):
                    ax.axvline(m,color=c,ls='--',label=f'{e}: {m:.2f}')
                ax.legend()
                ax.set_xlabel('Power-law exponent')
                ax.set_title(f'Power-law exponent for {pl_titles[iE]}')
                usrplt.adjust_spines(ax)
                xlim = ax.get_xlim(); ylim = ax.get_ylim()
                # res = st.ks_2samp(distr_list[0],distr_list[2])#,nan_policy='omit')
                # pval = res.pvalue
                # if pval < 0.05:
                #     ax.text(xlim[0] + np.diff(xlim)*0.5,0.9*ylim[1],'*',color=cmap[2])

                # res = st.ks_2samp(distr_list[1],distr_list[2])#,nan_policy='omit')
                # pval = res.pvalue
                # if pval < 0.05:
                #     ax.text(xlim[0] + np.diff(xlim)*0.25,0.9*ylim[1],'*',color=cmap[2])

                plt.savefig(join(PlotDir,f'RG-distri_{pl_strs[iE]}_exponent_{g}_{rec_name}.png'),dpi=300,bbox_inches='tight',facecolor='w')
                plt.savefig(join(ProjDir,'plots','RG1',f'RG-distri_{pl_strs[iE]}_exponent_{g}_{rec_name}.png'),dpi=300,bbox_inches='tight',facecolor='w')
                usrplt.save_fig_to_pptx(fig, prs)
                plt.close(fig)
    comp_length = time.perf_counter() - t0_outer
    mm, ss = divmod(comp_length,60)
    hh, mm = divmod(mm, 60)
    print(f'Completed in {hh:.0f} hrs, {mm:.0f} minutes, {ss:.0f} seconds')
    ## CLOSE files ----------------------------------------
    prs.save(join(PlotDir,f'RG-{rec_name}.pptx'))

    




