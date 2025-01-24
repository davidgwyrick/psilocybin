base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
base_dir_local = '/data/tiny-blue-dot/zap-n-zip/EEG_exp'; base_dir = base_dir_server
ProjDir = '/home/david.wyrick/projects/zap-n-zip/'
ServDir = '/data/projects/zap-n-zip/'
BehavDir = '/home/david.wyrick/projects/zap-n-zip/plots/behavior'

#Basea
import gspread
from os.path import join
from glob import glob
import json, os, time, sys, argparse
import numpy as np
import pandas as pd
import scipy.stats as st
import scipy.signal as sig
import pingouin as pg
from scipy.ndimage import gaussian_filter
from scipy.interpolate import interp1d
from itertools import combinations

#Plot
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
from matplotlib.gridspec import GridSpec
from matplotlib.lines import Line2D   

#Project
from tbd_eeg.data_analysis.eegutils import EEGexp
import tbd_eeg.data_analysis.Utilities.utilities as tbd_util
from tbd_eeg.data_analysis.Utilities.behavior_movies import Movie

#Allen
from allensdk.core.mouse_connectivity_cache import MouseConnectivityCache
mcc = MouseConnectivityCache(resolution=10)
str_tree = mcc.get_structure_tree()

#Read in allen CCF
ccfsum = pd.read_csv('/home/david.wyrick/projects/zap-n-zip/data/mouse_connectivity/ontology_v2.csv')

#User
sys.path.append(ProjDir)
import util
import plotting as usrplt

#Single cell stuff
import elephant.statistics as est
from quantities import ms, s, Hz
from neo.core import SpikeTrain

#Define behavioral states
# behavior_ranges = {0: [0,1], 1: [1,10], 2: [10,30], 3:[30,500]}
# behavior_dict = {0: 'rest (0-1cm/s)', 1: 'walk (1-10cm/s)', 2: 'shuffle (10-30cm/s)', 3: 'run (>30cm/s)'}
# behavior_strs2 = ['rest','walk','shuffle','run']

behavior_ranges = {0: [0,1], 1: [1,15], 2: [15,500]}#, 3:[30,500]}
behavior_dict = {0: 'rest (0-1cm/s)', 1: 'walk (1-15cm/s)', 2: 'run (>15cm/s)'}
behavior_strs2 = ['rest','walk','run']

# behavior_ranges = {0: [0,1], 1: [1,500]}#, 3:[30,500]}
# behavior_dict = {0: 'rest (0-1cm/s)', 1: 'active (>1cm/s)'}
# behavior_strs2 = ['rest','active']

behavior_strs = list(behavior_dict.values())
nBehaviors = len(behavior_strs)

#Define windows to calculate firing rate
windows_of_interest = [[.002,.025],[.075,.3],[.3,1],[1,4.5]]
window_strs = ['evoked (2-25ms)','rebound (75-300ms)','post-rebound (0.3-1s)','ISI (1-4.5s)']
nWindows = len(window_strs)

gc = gspread.service_account() # need a key file to access the account (step 2) 
sh = gc.open('Templeton-log_exp') # open the spreadsheet 

df = pd.DataFrame(sh.sheet1.get()) 
exp_table = df.T.set_index(0).T # put it in a nicely formatted dataframeexp_table.iloc[10:]
exp_table = exp_table.set_index('mouse_name')

##===== ============================ =====##
##===== Parse Command Line Arguments =====##
parser = argparse.ArgumentParser(description='single-cell-metrics')

##===== Data Options =====##
parser.add_argument('--mID',type=str, default='mouse709400',
                    help='mouse to perform analysis on')

parser.add_argument('--rec_name',type=str, default='aw_ket_2024-02-01_11-12-34',
                    help='experiment to perform analysis on')

parser.add_argument('--time_bin_ms',type=int, default=100,
                    help='time_bin_ms')

parser.add_argument('--window_t_min',type=int, default=15,
                    help='Window width (mins) to segment data into')


if __name__ == '__main__':

    # Parse the arguments
    args = parser.parse_args()
    mID = args.mID
    rec_name = args.rec_name
    time_bin_ms = args.time_bin_ms
    time_bin = time_bin_ms/1000

    #Segement data into 15 minute windows
    window_t_min = args.window_t_min
    window_t = window_t_min*60
    
    #Extract row from templeton dataframe
    exp_df = exp_table.loc[exp_table.exp_name == rec_name]
    if type(exp_df) == pd.core.series.Series:
        exp_df = exp_df.to_frame().T

    drug_type = exp_df['drug'].values[0]
    stim_type = exp_df['stimulation'].values[0]
    print(f'Experiment type: {stim_type}, {drug_type}')

    #Define directories
    SaveDir = join(ServDir,'data',mID,rec_name)
    if not os.path.exists(SaveDir):
        os.makedirs(SaveDir)

    PlotDir = join(ProjDir,'pptx')
    if not os.path.exists(PlotDir):
        os.makedirs(PlotDir)

    #%% Upload the whole experiment and generate the global clock
    file_name = os.path.join(base_dir_server,mID,rec_name,'experiment1','recording1')
    exp = EEGexp(file_name, preprocess=False, make_stim_csv=False)

    #Load probe data
    probe_unit_data, probe_info, total_units = tbd_util.get_neuropixel_data(exp)
    probe_list = [x.replace('_sorted', '') for x in exp.experiment_data if 'probe' in x]

    #Get recording start time
    probei = probe_list[0]
    open_ephys_start = np.round(np.min(probe_unit_data[probei]['spike_times']))
    open_ephys_end = np.round(np.max(probe_unit_data[probei]['spike_times']))

    #For saline & psilocybin experiments, get injection times and types of injection
    if drug_type in ['saline', 'psilocybin']:
        injection_times = [float(exp_df['First injection time (s)'].values[0]),
                        float(exp_df['Second injection time (s)'].values[0])]
        
        #Determine injection type
        if 'psi' in rec_name:
            injection_types = ['sal1','psi']
        else:
            injection_types = ['sal1', 'sal2']
        injection_time_windows = None
    elif drug_type == 'ketanserin+psilocybin':
        injection_time_windows = np.array([np.array(exp_df['First injection time (s)'].values[0].split(','),dtype=float),
                                np.array(exp_df['Second injection time (s)'].values[0].split(','),dtype=float)])
        
        #Take second time in each window as "injection time"
        injection_times = [injection_time_windows[0,1],injection_time_windows[1,1]]

        injection_types = ['ket','psi']
    else:
        injection_times = None
        injection_time_windows = None

    #For isoflurane experiments, get iso level
    if drug_type == 'isoflurane':
        iso_level, iso_times = exp.load_analog_iso()
        iso_induction_times, iso_maintenance_times = exp.load_iso_times()
    else:
        iso_maintenance_times = None
        iso_induction_times = None


    # extract the timestamps of the selected stimuli
    try:
        stim_log = pd.read_csv(exp.stimulus_log_file)

        if 'circle' in np.unique(stim_log['stim_type']):
            vis_stim_exists = True
        else:
            vis_stim_exists = False
        if 'biphasic' in np.unique(stim_log['stim_type']):
            stim_exists = True
        else:
            stim_exists = False
    except:
        stim_exists = False
        vis_stim_exists = False
        stim_log = None
    #Load behavior
    run_ts, raw_run_signal, run_signal, run_signal_s, pupil_ts, pupil_radius, run_signal_p, run_signal_p_s, plot_pupil = util.get_behavioral_data(exp, mID, rec_name)
    f_run = interp1d(run_ts,run_signal)
    f_pupil = interp1d(pupil_ts,pupil_radius)

    #Get time windows for each epoch
    epoch_list, time_window_list = util.define_epochs_of_interest([open_ephys_start,open_ephys_end], drug_type, window_t_min=window_t_min, injection_times=injection_times,injection_time_windows=injection_time_windows, iso_induction_times=iso_induction_times, stim_log=stim_log)

    #Create colormap for epochs
    nCond = len(epoch_list)
    if drug_type == 'saline':
        cmap = sns.color_palette('Greens',nCond)
    elif drug_type == 'psilocybin':
        nPsi = np.sum(['psi' in fname for fname in epoch_list])
        nSal = nCond - nPsi
        if nSal == 0:
            cmap = sns.color_palette('Reds',nCond)
        elif nPsi == 0:
            cmap = sns.color_palette('Blues',nCond)
        else:
            cmap = np.concatenate((sns.color_palette('Blues',nSal),sns.color_palette('Reds',nPsi)))
    elif drug_type == 'isoflurane':
        nPre = np.sum(['pre' in fname for fname in epoch_list])
        nPost = nCond - nPre
        if nPre == 0:
            cmap = sns.color_palette('Oranges',nCond)
        elif nPost == 0:
            cmap = sns.color_palette('Blues',nCond)
        else:
            cmap = np.concatenate((sns.color_palette('Blues',nPre),sns.color_palette('Oranges',nPost)))

    elif drug_type == 'ketanserin+psilocybin':
        nPre = np.sum(['pre' in fname for fname in epoch_list])
        nKet = np.sum(['ket' in fname for fname in epoch_list])
        nPsi = np.sum(['psi' in fname for fname in epoch_list])

        cmap = np.concatenate((sns.color_palette('Blues',nPre),sns.color_palette('Oranges',1),sns.color_palette('Reds',nPsi)))
    else:
        cmap = sns.color_palette('Oranges',nCond)

    # Create a new PowerPoint presentation to save figures to
    prs = Presentation()

    #Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f'mID = {mID}; rec_name = {rec_name}'
    slide.placeholders[1].text = f'Perform Poisson GLM on spiking activity to "control" for running'
    
    #Read in neuropixel data 
    data_list, ts_list, neuron_indices, plot_tuple, FR_perblock = util.bin_spiking_data(probe_unit_data, time_window_list, time_bin=time_bin,fr_thresh=0) 
    boundaries_group, ticks_group, labels_group, celltypes, durations, layers, areas, groups, supergroups, order_by_group = plot_tuple

    #Plot running
    fig, ax = plt.subplots(figsize=(10,4))
    ax.set_title(mID)
    ax.plot(run_ts,run_signal_s,'-k',lw=0.6)
    ax2 = ax.twinx()
    ax2.plot(pupil_ts,pupil_radius,'-r',lw=0.6)
    usrplt.save_fig_to_pptx(fig,prs)

    ##-------------------------------------
    #Plot correlation with behavior
    nCond = len(data_list)
    T, N = data_list[0].shape
    nNeurons = N
    corr_behavior = np.zeros((2,nCond,N))
    running_list = []; pupil_list = []
    fig, axes = plt.subplots(1,3,figsize=(12,4))
    for ii, (ts, data) in enumerate(zip(ts_list,data_list)):
        run_cond = f_run(ts)
        pupil_cond = f_pupil(ts)
        running_list.append(run_cond); pupil_list.append(pupil_cond)

        for jj in range(N):
            corr_behavior[0,ii,jj] = np.corrcoef(run_cond,data[:,jj])[0,1]
            corr_behavior[1,ii,jj] = np.corrcoef(pupil_cond,data[:,jj])[0,1]

        sns.kdeplot(corr_behavior[0,ii],ax=axes[0],color=cmap[ii])
        sns.kdeplot(corr_behavior[1,ii],ax=axes[1],color=cmap[ii])

        axes[2].scatter(corr_behavior[0,ii],corr_behavior[1,ii],marker='.',s=10,color=cmap[ii])
    plt.suptitle('Correlation of spiking activity')
    ax = axes[0]
    axes[0].set_title('with running')
    axes[1].set_title('with pupil')

    for ax in axes:
        ax.set_xlabel('Correlation')
        ax.set_xlim([-0.75,0.75])
        usrplt.adjust_spines(ax,['bottom','left'])

    axes[2].set_xlabel('Correlation with running')
    axes[2].set_ylabel('Correlation with pupil')
    usrplt.save_fig_to_pptx(fig,prs)
 
    ##-------------------------------------
    fig, axes = plt.subplots(2,1,figsize=(10,8),sharey=True)
    ax = axes[0]
    for ii, epoch in enumerate(epoch_list):
        nN_n_corr = np.sum(corr_behavior[0,ii] <= -0.1)
        nN_p_corr = np.sum(corr_behavior[0,ii] >= +0.1)
        print(f'{epoch:20s}: {nN_p_corr:3d} positively correlated neurons & {nN_n_corr:3d} negatively correlated neurons ')
        ax.bar(ii-0.5,nN_p_corr/nNeurons,width=0.4,color=cmap[ii])
        ax.bar(ii,nN_n_corr/nNeurons,width=0.4,color=cmap[ii],hatch='//')
    ax.set_xticks(np.arange(nCond)-0.25)
    ax.set_xticklabels([],rotation=30)
    ax.set_title('Fraction of cells positively (or negatively) correlated with running')
    ax.set_ylabel('Fraction of cells')
    # ax.set_xlabel('Epoch')  

    print('Pupil')
    # fig, ax = plt.subplots(figsize=(10,4))
    ax = axes[1]
    for ii, epoch in enumerate(epoch_list):
        nN_n_corr = np.sum(corr_behavior[1,ii] <= -0.1)
        nN_p_corr = np.sum(corr_behavior[1,ii] >= +0.1)
        print(f'{epoch:20s}: {nN_p_corr:3d} positively correlated neurons & {nN_n_corr:3d} negatively correlated neurons ')
        ax.bar(ii-0.5,nN_p_corr/nNeurons,width=0.4,color=cmap[ii])
        ax.bar(ii,nN_n_corr/nNeurons,width=0.4,color=cmap[ii],hatch='//')
    ax.set_xticks(np.arange(nCond)-0.25)

    ax.set_title('Fraction of cells positively (or negatively) correlated with pupil size')
    ax.set_ylabel('Fraction of cells')
    ax.set_xlabel('Epoch')  

    for ax in axes:
        # ax.set_ylim([0,0.45])
        # ax.set_yticks([0,0.2,0.4])
        usrplt.adjust_spines(ax)
    ax.set_xticklabels(epoch_list,rotation=30)
    usrplt.save_fig_to_pptx(fig,prs)

    ##-------------------------------------
    from sklearn.linear_model import LinearRegression, PoissonRegressor
    data_list2 = []; r2_list2 = []
    corr_behavior2 = np.zeros(corr_behavior.shape)

    for ii, data in enumerate(data_list):
        spks = data.T
        N = spks.shape[0]
        run_sub = running_list[ii]
        pupil_sub = pupil_list[ii]

        tmp_list = []
        spks2 = np.zeros(spks.shape)
        X = run_sub.reshape(-1,1)
        for jj in range(N):
            y = spks[jj].copy()
            reg = PoissonRegressor().fit(X,y)
            tmp = y - reg.predict(X)
            spks2[jj] = tmp * (tmp > 0)
            tmp_list.append(reg.score(X,y))
            corr_behavior2[0,ii,jj] = np.corrcoef(run_sub,spks2[jj])[0,1]
            corr_behavior2[1,ii,jj] = np.corrcoef(pupil_sub,spks2[jj])[0,1]
        data_list2.append(spks2.T)
        r2_list2.append(tmp_list)

    ##-------------------------------------
    fig, axes = plt.subplots(1,3,figsize=(12,4))
    for ii, (ts, data) in enumerate(zip(ts_list,data_list)):

        sns.kdeplot(corr_behavior2[0,ii],ax=axes[0],color=cmap[ii])#,bins=np.arange(-1,1.01,0.05))
        sns.kdeplot(corr_behavior2[1,ii],ax=axes[1],color=cmap[ii])#,bins=np.arange(-1,1.01,0.05))

        axes[2].scatter(corr_behavior2[0,ii],corr_behavior2[1,ii],marker='.',s=10,color=cmap[ii])
    plt.suptitle('Correlation of spiking activity')
    ax = axes[0]
    axes[0].set_title('with running')
    axes[1].set_title('with pupil')

    for ax in axes:
        ax.set_xlabel('Correlation')
        ax.set_xlim([-0.75,0.75])
        usrplt.adjust_spines(ax,['bottom','left'])

    axes[2].set_xlabel('Correlation with running')
    axes[2].set_ylabel('Correlation with pupil')
    usrplt.save_fig_to_pptx(fig,prs)

    
    for ii, (ts, data) in enumerate(zip(ts_list,data_list)):
        fig, ax = plt.subplots(figsize=(8,8))
        ax.set_title(epoch_list[ii])
        ax.scatter(corr_behavior[0,ii],corr_behavior2[0,ii],marker='.',s=10,color=cmap[ii])
        ax.plot([-1,1],[-1,1],'-k')
        ax.vlines(0,*ax.get_ylim())
        ax.hlines(0,*ax.get_xlim())
        ax.set_xlim([-0.5,0.5])
        ax.set_ylim([-0.5,0.5])
        ax.set_xlabel('Correlation before')
        ax.set_ylabel('Correlation after')
        usrplt.save_fig_to_pptx(fig,prs)


    for ii, data in enumerate(data_list2):
        # print(data.shape)
        fsuffix = f'{epoch_list[ii]}_{time_bin_ms}ms-bins'
        np.savez(join(SaveDir,f'spike-counts_poissonGLM_{fsuffix}.npz'),X=data,ts=ts_list[ii])

    prs.save(join(SaveDir,f'spiking-data-preprocessing_{rec_name}.pptx'))
    print('DONE!!!')

# %%
