base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
base_dir_local = '/data/tiny-blue-dot/zap-n-zip/EEG_exp'; base_dir = base_dir_server
ProjDir = '/home/david.wyrick/projects/zap-n-zip/'
ServDir = '/data/projects/zap-n-zip/'
BehavDir = '/home/david.wyrick/projects/zap-n-zip/plots/behavior'

##------------------------------------------
##Load modules
#Base
import gspread
from os.path import join
from glob import glob
import json, os, time, sys, argparse
import numpy as np
import pandas as pd
import itertools as it
import scipy.stats as st
from scipy.interpolate import interp1d
from tqdm import tqdm

#Plot
import matplotlib.pyplot as plt
from pptx import Presentation
import seaborn as sns

#Project
from tbd_eeg.data_analysis.eegutils import EEGexp
import tbd_eeg.data_analysis.Utilities.utilities as tbd_util

#Allen
from allensdk.core.mouse_connectivity_cache import MouseConnectivityCache
mcc = MouseConnectivityCache(resolution=10)
str_tree = mcc.get_structure_tree()

#User
sys.path.append(ProjDir)
import util
import plotting as usrplt

#Define behavioral states
behavior_ranges = {0: [0,1], 1: [1,500]}#, 3:[30,500]}
behavior_dict = { 0: 'rest (<1cm/s)', 1: 'active (>1cm/s)'}
behavior_dict2 = { 0: 'rest (<1cm/s)', 1: 'active (>1cm/s)', 2: 'all (>0cm/s)'}
behavior_strs2 = ['rest','active']
behavior_strs = list(behavior_dict.values())
nBehaviors = len(behavior_strs)

#Define windows to calculate firing rate
visual_stim = ['circle','natural_scene']
spont_types = ['spontaneous','isi_biphasic','isi_visual']
evoked_windows = [[.002,.025],[0.025,0.075],[.075,.3],[.3,1]]
evoked_strings = ['evoked','pre-rebound','rebound','post-rebound']

gc = gspread.service_account() # need a key file to access the account (step 2) 
sh = gc.open('Templeton-log_exp') # open the spreadsheet 

df = pd.DataFrame(sh.sheet1.get()) 
exp_table = df.T.set_index(0).T # put it in a nicely formatted dataframeexp_table.iloc[10:]
exp_table = exp_table.set_index('mouse_name')

##------------------------------------------
##Define Command Line Arguments
parser = argparse.ArgumentParser(description='single-cell-metrics')

parser.add_argument('--mID',type=str, default='mouse730913',
                    help='mouse to perform analysis on')

parser.add_argument('--rec_name',type=str, default='aw_ket_2024-04-26_10-11-34',
                    help='experiment to perform analysis on')

parser.add_argument('--time_bin_ms',type=int, default=5,
                    help='time_bin_ms')


if __name__ == '__main__':

    ##------------------------------------------
    #Parse the arguments
    args = parser.parse_args()
    mID = args.mID
    rec_name = args.rec_name
    time_bin_ms = args.time_bin_ms
    time_bin = time_bin_ms/1000

    # #Extract row from templeton dataframe
    # exp_df = exp_table.loc[exp_table.exp_name == rec_name]
    # if type(exp_df) == pd.core.series.Series:
    #     exp_df = exp_df.to_frame().T

    # drug_type = exp_df['drug'].values[0]
    # stim_type = exp_df['stimulation'].values[0]
    # print(f'Experiment type: {stim_type}, {drug_type}')

    # #Define directories
    # TempDir = join(ServDir,'results','evoked',mID,rec_name)
    # if not os.path.exists(TempDir):
    #     os.makedirs(TempDir)
    BaseDir = join(ServDir,'results','evoked')

    # Create a new PowerPoint presentation to save figures to
    prs = Presentation()

    #Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f'All mice'
    slide.placeholders[1].text = f'Evoked PSTHs'

    for mID, row in exp_table[:96].iterrows():
        rec_name = row.exp_name
        stimulation = row.stimulation
        drug_type = row.drug

        #Extract row from templeton dataframe
        exp_df = exp_table.loc[exp_table.exp_name == rec_name]
        if type(exp_df) == pd.core.series.Series:
            exp_df = exp_df.to_frame().T

        if 'electrical' not in stimulation:
            continue
        if mID == 'mouse703062':
            continue
        if 'stim_train' in rec_name:
            continue
        SaveDir = join(BaseDir,mID,rec_name)
        if not os.path.exists(SaveDir):
            os.makedirs(SaveDir)
        print(f'Processing {mID}, {rec_name}')

        try:
            ##------------------------------------------
            #Upload the whole experiment 
            file_name = os.path.join(base_dir_server,mID,rec_name,'experiment1','recording1')
            exp = EEGexp(file_name, preprocess=False, make_stim_csv=False)

            #Load probe data
            probe_unit_data, probe_info, total_units, metric_list = tbd_util.get_neuropixel_data(exp)
            probe_list = [x.replace('_sorted', '') for x in exp.experiment_data if 'probe' in x]
            plot_tuple = util.get_group_plotting_info(probe_unit_data)
            boundaries, yticks, labels, celltypes, durations, layers, areas, groups, mesogroups, supergroups, order_by_group = plot_tuple
            nNeurons = len(areas)

            #Get recording start time
            probei = probe_list[0]
            open_ephys_start = np.round(np.min(probe_unit_data[probei]['spike_times']))
            open_ephys_end = np.round(np.max(probe_unit_data[probei]['spike_times']))

            ##------------------------------------------
            #Get injection times and types of injection
            if drug_type in ['saline', 'psilocybin','ketanserin+psilocybin']:
                injection_time_windows = np.array([np.array(exp_df['First injection window'].values[0].split(','),dtype=float),
                                        np.array(exp_df['Second injection window'].values[0].split(','),dtype=float)])
                injection_times = np.array([exp_df['First injection time (s)'].values[0],exp_df['Second injection time (s)'].values[0]],dtype=float)

                if drug_type == 'psilocybin':
                    injection_types = ['sal1','psi']
                    injection_colors = sns.xkcd_palette(['dark sky blue','darkish red'])
                elif drug_type == 'saline':
                    injection_types = ['sal1', 'sal2']
                    injection_colors = sns.xkcd_palette(['dark sky blue','cobalt blue'])
                elif drug_type == 'ketanserin+psilocybin':
                    injection_types = ['ket','psi']
                    injection_colors = sns.xkcd_palette(['magenta','goldenrod'])
                inj_tuple = (injection_times,injection_types,injection_colors)
            else:
                injection_times = None
                inj_tuple = None

            #For isoflurane experiments, get iso level
            if drug_type == 'isoflurane':
                iso_level, iso_times = exp.load_analog_iso()
                iso_induction_times, iso_maintenance_times = exp.load_iso_times()
                induction_colors = sns.xkcd_palette(['light teal','teal'])
                iso_tuple = (iso_induction_times, induction_colors)
            else:
                iso_induction_times = None
                iso_tuple = None

            ##------------------------------------------
            #Define time windows for each epoch  
            if drug_type in ['saline', 'psilocybin','ketanserin+psilocybin']:
                macro_windows = [[open_ephys_start,injection_time_windows[0,0]]]; macro_names = ['pre-inj']
                macro_windows.append([injection_time_windows[0,1],injection_time_windows[1,0]]); macro_names.append(f'post-{injection_types[0]}-inj')
                macro_windows.append([injection_time_windows[1,1],open_ephys_end]); macro_names.append(f'post-{injection_types[1]}-inj')
                if drug_type == 'saline':
                    cmap_macro = sns.xkcd_palette(['silver','dark sky blue','cobalt blue'])
                elif drug_type == 'psilocybin':
                    cmap_macro = sns.xkcd_palette(['silver','dark sky blue','darkish red'])
                elif drug_type == 'ketanserin+psilocybin':
                    cmap_macro = sns.xkcd_palette(['silver','magenta','dusty orange'])
                
            elif drug_type == 'isoflurane':

                macro_windows = [[open_ephys_start,iso_induction_times[0]-120]]; macro_names = ['pre-iso']
                macro_windows.append([iso_induction_times[0],iso_induction_times[1]]); macro_names.append(f'iso-ind')
                macro_windows.append([iso_induction_times[1]+60*5,open_ephys_end]); macro_names.append(f'post-iso')
                cmap_macro = sns.xkcd_palette(['silver','light teal','teal'])

            elif drug_type == 'urethane':
                macro_windows = [[open_ephys_start,open_ephys_end]]; macro_names = ['urethane']
                cmap_macro = sns.xkcd_palette(['darkish purple'])

            ##------------------------------------------
            #Load stimulus log if it exists
            try:
                stim_log = pd.read_csv(exp.stimulus_log_file)

                if ('circle' in np.unique(stim_log['stim_type'])) | ('natural_scene' in np.unique(stim_log['stim_type'])):
                    vStim_exists = True
                    stim_log_v = stim_log.loc[stim_log.stim_type.isin(visual_stim)]
                else:
                    vStim_exists = False

                if 'biphasic' in np.unique(stim_log['stim_type']):
                    eStim_exists = True
                    stim_log_b = stim_log.loc[stim_log.stim_type == 'biphasic']
                    uniq_currents = np.unique(stim_log_b['parameter'])
                    uniq_sweeps = np.unique(stim_log_b['sweep'])
                    nSweeps = len(uniq_sweeps)
                    nCurrents = len(uniq_currents)
                    print(f'Number of sweeps: {nSweeps}, Number of currents: {nCurrents}')
                else:
                    eStim_exists = False
            except:
                stim_log = None
                eStim_exists = False
                vStim_exists = False

            evoked_time_window_list = []
            evoked_type_list = []
            sweep_time_window_list = []
            if eStim_exists | vStim_exists:
                for s in np.unique(stim_log['sweep']):
                    for t in np.unique(stim_log.loc[stim_log.sweep == s]['stim_type']):
                        sub_df = stim_log.loc[(stim_log.sweep == s) & (stim_log.stim_type == t)]
                        tS = sub_df['onset'].min(); tE = sub_df['offset'].max()
                        
                        evoked_time_window_list.append([tS,tE])
                        evoked_type_list.append(t)

            ##------------------------------------------
            #Load behavior
            run_ts, run_signal, run_signal_s, acceleration, pupil_ts, pupil_radius, pupil_dxdt, plot_pupil = util.get_behavioral_data(exp, mID, rec_name,normalize=True)
            run_signal[run_signal < 0] = 0
            run_signal_s[run_signal_s < 0] = 0
            f_run = interp1d(run_ts,run_signal)
            f_pupil = interp1d(pupil_ts,pupil_radius)

            ##------------------------------------------
            #Plot behavior and save to pptx
            T = len(run_ts)
            time_bin_run = 0.01
            tMin = T*time_bin_run/60

            #Create behavioral map based on smoothed running speed
            behavior_map = np.zeros(run_signal_s.shape)
            for key, b_range in behavior_ranges.items():
                indy = np.where((run_signal_s >= b_range[0]) & (run_signal_s < b_range[1]))[0]
                behavior_map[indy] = key
                t = len(indy)*time_bin_run

            fig = usrplt.plot_behavior((run_ts,run_signal),(pupil_ts,pupil_radius),f'{mID} {rec_name}',(evoked_time_window_list,evoked_type_list),inj_tuple,iso_tuple,behavior_map)
            usrplt.save_fig_to_pptx(fig, prs)

            boundaries, ticks, labels, celltypes, durations, layers, areas, groups, mesogroups,supergroups, order_by_group = util.get_group_plotting_info(probe_unit_data)
            areas_ro = areas[order_by_group]
            groups_ro = groups[order_by_group]
            supergroups_ro = supergroups[order_by_group]
            celltypes_ro = celltypes[order_by_group]

            ##------------------------------------------
            #Plot evoked PSTHs
            stim_log_biphasic = stim_log.loc[stim_log.stim_type == 'biphasic']
            for sweep  in np.unique(stim_log_biphasic['sweep']):
                sub_df = stim_log.loc[(stim_log.stim_type == 'biphasic') & (stim_log['sweep'] == sweep)]

                #Get stimulation times
                event_times = sub_df['onset'].values
                stim_amp = sub_df['parameter'].values

                #Get evoked spike counts centered around stimulation times
                plot_before = 0.25; plot_after = .75; time_bin = 5/1000
                spk_list = []
                for probei in probe_list:
                    evoked_spks, bins = util.get_evoked_spike_counts(probe_unit_data[probei]['spike_times'], probe_unit_data[probei]['spike_clusters'], probe_unit_data[probei]['units'], event_times, plot_before, plot_after, time_bin)
                    spk_list.append(evoked_spks)
                    spk_ts_trial = bins[:-1] + time_bin/2
                    
                evoked_ts_list = []
                for e in event_times:
                    evoked_ts_list.append(e+spk_ts_trial)
                evoked_ts = np.array(evoked_ts_list)

                drug_state = f'{evoked_ts[0,0]/60:.2f} min into recording'
                for e, tW in zip(macro_names,macro_windows):
                    if (evoked_ts[0,0] > tW[0]) & (evoked_ts[0,0] < tW[1]):
                        drug_state = e
                        break

                #Reshape
                spk_counts_evoked = np.concatenate(spk_list)
                spk_counts_evoked = np.transpose(spk_counts_evoked,[1,2,0])
                T = spk_ts_trial.shape[0]

                #Select neurons in groups of interest
                # g_indy = np.where([g in ['MO','SM-TH','TH'] for g in groups_ro])[0]
                g_indy = []
                for ii, (a, g, ct) in enumerate(zip(areas_ro, groups_ro, celltypes_ro)):
                    if (g in ['MO','SM-TH','ANT-TH','TH','RT']):#& (ct == 'RS'):
                        g_indy.append(ii)
                g_indy = np.array(g_indy)
                # g_indy = np.where([g in ['HIP', 'HPF', 'MO', 'PALm', 'RSP', 'SM-TH', 'SS','STR', 'TH', 'VIS', 'VIS-TH'] for g in groups_ro])[0]

                #Get areas and labels
                areas_g = areas_ro[g_indy]
                uniq_a, indices = np.unique(areas_g, return_index=True)
                labels_a = uniq_a[np.argsort(indices)]
                boundaries_a = [0]; ticks_a = []
                counter = 0
                for a in labels_a:
                    n = np.sum(areas_g == a)
                    ticks_a.append(boundaries_a[-1]+n/2)
                    boundaries_a.append(boundaries_a[-1]+n)

                for iParam in np.unique(sub_df['parameter']):
                    indy = np.where(sub_df['parameter'] == iParam)[0]

                    PSTH = np.nanmean(spk_counts_evoked[indy],axis=0)
                    PSTH_sort = PSTH[:,order_by_group]
                    indy_pre = np.where(spk_ts_trial < 0)[0]
                    mean_baseline = np.nanmean(PSTH_sort[indy_pre],axis=0)
                    mean_stddev = np.nanstd(PSTH_sort[indy_pre],axis=0)
                    PSTH_z = (PSTH_sort - mean_baseline)/mean_stddev
                    T, N = PSTH_sort.shape

                    #Plot PSTH for each stimulation protocol
                    fig, ax = plt.subplots(figsize=(11,7))
                    plt.suptitle(f'{mID}, {rec_name}\nPSTH, sweep {sweep}, {iParam}uA, {drug_state}',y=0.96)

                    spk_plot = ax.imshow(PSTH_z[:,g_indy].T, aspect='auto',cmap='bwr',vmin=-5,vmax=5)
                    i0 = np.where(spk_ts_trial > 0)[0][0]-1
                    i1 = np.where(spk_ts_trial > plot_after/2)[0][0]

                    xticks = np.array([0,i0,i1,T]); xticks[-1] -= 1
                    ax.set_xticks(xticks)
                    ax.set_xticklabels(np.round(spk_ts_trial[xticks],2))
                    ax.vlines(i0,*ax.get_ylim(),color='g',lw=1)

                    ax.set_yticks(ticks_a,minor=True)
                    ax.set_yticks(boundaries_a,minor=False)
                    ax.set_yticklabels(labels_a,minor=True)
                    ax.set_yticklabels([],minor=False)
                    ax.set_ylabel('Area')
                    ax.set_xlabel('Time (s)')
                    
                    np.savez(join(SaveDir,f'PSTH_{drug_state}_sweep-{sweep}_{iParam}uA.npz'),PSTH_z=PSTH_z,PSTH_sort=PSTH_sort,PSTH=PSTH,spk_ts_trial=spk_ts_trial,areas_ro=areas_ro)
                    #Save figure to pptx
                    plt.savefig(join(SaveDir,f'PSTH_{drug_state}_{mID}_{rec_name}_sweep-{sweep}_{iParam}uA.png'),dpi=300,bbox_inches='tight',facecolor='w')
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = f'{mID}, {rec_name}'
                    usrplt.save_fig_to_pptx(fig, prs,slide)
                    plt.close(fig)
        except:
            print(f'Error in {mID}, {rec_name}')
            continue
        prs.save(join(BaseDir,f'evoked_PSTHs_motor-areas.pptx'))
        # ##------------------------------------------
        # #Save results
        # np.savez(join(SaveDir,f'single-cell-metrics_spont_{rec_name}.npz'), combos_spont=combos_spont, \
        #         moments_spont=moments_spont,pvalues_spont=pvalues_spont,mod_idx_spont=mod_idx_spont,numspks_spont=numspks_spont)
