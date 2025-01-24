base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
base_dir_local = '/data/tiny-blue-dot/zap-n-zip/EEG_exp'; base_dir = base_dir_server
ProjDir = '/home/david.wyrick/projects/zap-n-zip/'
ServDir = '/data/projects/zap-n-zip/'
BehavDir = '/home/david.wyrick/projects/zap-n-zip/plots/behavior'

#Basea
import gspread
from os.path import join
from glob import glob
import json, os, time, sys, argparse
import numpy as np
import pandas as pd
import scipy.stats as st
import pingouin as pg

#Plot
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
import seaborn as sns

#Project
from tbd_eeg.data_analysis.eegutils import EEGexp
import tbd_eeg.data_analysis.Utilities.utilities as tbd_util
from tbd_eeg.data_analysis.Utilities.behavior_movies import Movie

#Allen
from allensdk.core.mouse_connectivity_cache import MouseConnectivityCache
mcc = MouseConnectivityCache(resolution=10)
str_tree = mcc.get_structure_tree()

#Read in allen CCF
ccfsum = pd.read_csv('/home/david.wyrick/projects/zap-n-zip/data/mouse_connectivity/ontology_v2.csv')

#User
sys.path.append(ProjDir)
import util
import plotting as usrplt

#Define behavioral states
behavior_ranges = {0: [0,1], 1: [1,500]}#, 3:[30,500]}
behavior_dict = { 0: 'rest (<1cm/s)', 1: 'active (>1cm/s)'}
behavior_dict2 = { 0: 'rest (<1cm/s)', 1: 'active (>1cm/s)', 2: 'all (>0cm/s)'}
behavior_strs2 = ['rest','active']
behavior_strs = list(behavior_dict.values())
nBehaviors = len(behavior_strs)

#Define windows to calculate firing rate
stim_window_dict = {'spontaneous': [], 'evoked': [],'pre-rebound': [], 'rebound': [], 'post-rebound': [], 'isi': [],'visual': []}
stim_strs = ['spontaneous','evoked','pre-rebound','rebound','post-rebound','isi_biphasic','visual','isi_visual']

visual_stim = ['circle','natural_scene']
spont_types = ['spontaneous','isi_biphasic','isi_visual']
evoked_windows = [[.002,.025],[0.025,0.075],[.075,.3],[.3,1]]
evoked_strings = ['evoked','pre-rebound','rebound','post-rebound']

gc = gspread.service_account() 
sh = gc.open('Templeton-log_exp')

df = pd.DataFrame(sh.sheet1.get()) 
exp_table = df.T.set_index(0).T 
exp_table = exp_table.set_index('mouse_name')

th_dict = {'SM-TH': ['PO','VAL','VPL','VPM','VM'], 
'VIS-TH': ['LP','LGN','LGd','LGd-co','LGd-sh','LGd-ip'],
'ANT-TH': ['AV','AMd','AMv','AD','MD','MDm','MDc','MDl','RE','RH','CM','LD', 'CL'],
'TH': ['Eth', 'IAD', 'IGL', 'IntG', 'LGv','MGd', 'MGm', 'MGv', 'PCN', 'PF', 'PIL', 'PoT', 'SGN','SPFp', 'TH','LH'],'RT': ['RT']}

##===== ============================ =====##
##===== Parse Command Line Arguments =====##
parser = argparse.ArgumentParser(description='plot_single-cell-metrics')

##===== Data Options =====##
parser.add_argument('--run_num',type=int, default=0,
                    help='run_num')

parser.add_argument('--plot_only_sig_mod',type=int, default=0,
                    help='Plot all neurons or only significantly modulated ones')

if __name__ == '__main__':

    pval_thresh = 0.01
    # Parse the arguments
    args = parser.parse_args()
    run_num = args.run_num
    plot_only_sig_mod = bool(args.plot_only_sig_mod)
    tuples_included = []
    pval_thresh = 0.05
    FR_thresh = 1

    FR_df_list = []; dFR_df_list = []

    #Loop over experiments and read in single cell metrics
    for mID, row in exp_table.iterrows():
        rec_name = row.exp_name
        drug_type = row.drug
        stimulation = row.stimulation
        
        if 'electrical' in stimulation:
            st_type = 'electrical'
        elif 'spontaneous' in stimulation:
            st_type = 'spontaneous'

        #Read in single cell metrics
        SaveDir = join(ServDir,'results','FR_reformat4',mID,rec_name)
        fpath = join(SaveDir,f'delta-spont_FR_ALL_{rec_name}.h5')
        if os.path.exists(fpath) == False:
            continue

        if run_num == 2:
            tuples_included.append((mID,rec_name,drug_type))    
        dFR_path = join(SaveDir,f'delta-spont_FR_ALL_{rec_name}.h5')
        if os.path.exists(dFR_path):
            dFR_spont_df = pd.read_hdf(dFR_path,'df')
            dFR_spont_df = dFR_spont_df.astype({'val_i':float, 'val_j':float,'delta':float, 'pval1':float, 'sig1':float, 'pval2':float, 'sig2':float, 'sig':float, 'mod':float})

            dFR_spont_df['abs_mod'] = np.abs(dFR_spont_df['mod'])
            dFR_spont_df['abs_delta'] = np.abs(dFR_spont_df['delta'])
            dFR_spont_df['st_type'] = st_type
            
            dFR_df_list.append(dFR_spont_df)
            print(f'{mID}, {rec_name} single-cell-metrics dataframe read in')
            continue


    dFR_spont_df = pd.concat(dFR_df_list)
    dFR_spont_df = dFR_spont_df.reset_index(drop=True)

  

    ##------------------------------------------
    #Save dataframes & create save directory
    # if plot_only_sig_mod:
    #     print('Plotting only significantly modulated neurons')
    #     BaseDir = join(ServDir,'results','FR_ISI','all_sig_neurons','metrics_run_{:02d}'.format(run_num))
    # else:
    #     print('Plotting all neurons')
    #     BaseDir = join(ServDir,'results','FR_ISI','all_neurons','metrics_run_{:02d}'.format(run_num))

    #Plot only spontaneous experiments
    BaseDir = join(ServDir,'results','FR_reformat4','all_neurons','metrics_run_{:02d}'.format(run_num))
    # FR_spont_df = FR_spont_df.loc[FR_spont_df.st_type == 'spontaneous']
    # dFR_spont_df = dFR_spont_df.loc[dFR_spont_df.st_type == 'spontaneous']
    print(f'Saving dataframes to {BaseDir}')

    if not os.path.exists(BaseDir):
        os.makedirs(BaseDir)

    blocks = ['post_1st_inj','0_15','15_30','30_45','45_60','60_75','75_90']
    blocks = ['0_15','15_30','30_45','45_60','60_75','75_90']
    drug_order = ['saline','psilocybin']#,'ketanserin+psilocybin']
    cmap_macro = sns.xkcd_palette(['cobalt blue','darkish red'])#,'dusty orange'])
    stat_list = ['mean_FR']#,'CV_ISI']
    stat_title_list = ['Mean firing rate (Hz)']#,'Coefficient of variation of ISI']
    bstrs_list = ['rest','run','all']
    layer_list = ['2/3','4','5','6']
    legend = False
    mesogroup_list = ['TH_core','TH_matrix','TH_intralaminar','CTX_frontal','CTX_sensory','HIP','STR']
    CTX_mesogroup_list = np.array(['CTX_frontal','CTX_sensory'])
    group_list = ['SM-TH','VIS-TH','ANT-TH','TH','MO','ILA','PL','ACA','ORB','SSp','VIS','OLF','RSP','HIP','STR']
    CTX_group_list = ['MO','ILA','PL','ACA','ORB','SSp','VIS','OLF','RSP']
    TH_group_list = ['SM-TH','VIS-TH','ANT-TH','TH','RT','HIP','STR']
    epoch_i = 'pre_inj'; block_i = 'pre_inj'
    # epoch_i = 'post_1st_inj'; block_i = 'post_1st_inj'

    nPlots = 0

    ##------------------------------------------
    # Create a new powerPoint presentations to save figures to
    prs_list = []
    for iS, stat in enumerate(stat_list):
        prs = Presentation()

        #Add title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f'All Mice, neurons organized in "mesogroups" level'
        slide.placeholders[1].text = f'Change in {stat_title_list[iS]} due to saline or psilocybin injection'
        prs_list.append(prs)

    ##------------------------------------------
    #Plot MI over time per mesogroup, celltype, behavior
    print('Plotting mesogroup barplots over time, per celltype, behavior, and mesogroup')
    for iS, stat in enumerate(stat_list):
        for ct in ['RS','FS']:
            if ct == 'FS':
                mesogroup_list = np.array(['TH_core','TH_matrix','TH_RT','CTX_frontal','CTX_sensory','HIP','STR'])
            else:
                mesogroup_list = np.array(['TH_core','TH_matrix','TH_intralaminar','CTX_frontal','CTX_sensory','HIP','STR'])
            for iB, b_str in behavior_dict2.items():
                PlotDir = join(BaseDir,'mesogroup',bstrs_list[iB])
                if not os.path.exists(PlotDir):
                    os.makedirs(PlotDir)
                for mg in mesogroup_list:
                
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.mesogroup == mg) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                    if len(sub_df) < 5:
                        continue
                    fig, axes = plt.subplots(3,1,figsize=(10,15))
                    plt.suptitle(f'Change in {stat_title_list[iS]} due to saline or psilocybin injection\nBaseline = {block_i}, {mg}, {ct}, {b_str}',y=0.925)

                    #Fraction significantly modulated
                    ax = axes[0]
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.mesogroup == mg) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]

                    iX = 0; xticks = []
                    for i, b in enumerate(blocks):
                        for jj, d in enumerate(drug_order):
                            sub_sub_df = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == d)]

                            n = len(sub_sub_df)
                            if n == 0:
                                continue
                            nSig = np.sum(sub_sub_df['sig'].values)
                            fracSig = nSig/n
                            pos_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] > 0)])/n
                            neg_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] < 0)])/n
                            not_mod = len(sub_sub_df.loc[sub_sub_df.sig == 0])/n

                            h1 = ax.bar(iX, pos_mod, color=sns.xkcd_palette(['light red']), label='Positively Modulated')
                            h2 = ax.bar(iX, not_mod, bottom=pos_mod, color=sns.xkcd_palette(['silver']), label='Not Modulated')
                            h3 = ax.bar(iX, neg_mod, bottom=pos_mod+not_mod, color=sns.xkcd_palette(['medium blue']), label='Negatively Modulated')
                            
                            n = len(np.unique(sub_sub_df['neuron_index']))
                            ax.text(iX,0.025,f'{n}',ha='center',fontsize=8,rotation=90)
                            if jj == 1:
                                xticks.append(iX-0.5)
                            iX += 1
                        iX += 1

                    usrplt.adjust_spines(ax)
                    plt.autoscale(tight=True)
                    xticks = np.arange(1,3*len(blocks),3)-0.5
                    ax.set_xticks(xticks)
                    ax.set_xticklabels(blocks,rotation=0)
                    ax.set_ylabel('Fraction modulated')
                    ax.legend(handles=[h1,h2,h3],loc=2)

                    #Modulation index
                    ax = axes[1]
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.mesogroup == mg) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                    sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                    sns.barplot(data=sub_df,x='block_j',y='abs_mod',hue='drug_type',ax=ax,order=blocks,hue_order=drug_order,palette=cmap_macro,errorbar='se',legend=legend)
                    # ax.set_ylim([0,0.6])
                    usrplt.adjust_spines(ax)
                    ax.set_xlabel('Time after 2nd injection (min)'); ax.set_ylabel('| Modulation index |')

                    for i, b in enumerate(blocks):
                        distr_i = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[0])]['abs_mod'].values
                        distr_j = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[1])]['abs_mod'].values
                        if (len(distr_i) < 5) | (len(distr_j) < 5):
                            continue
                        res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')
                                            
                        if res.pvalue < pval_thresh:
                            ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])

                    #Modulation index (dFR_spont_df.sig == 1) &
                    ax = axes[2]
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.mesogroup == mg) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                    sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                    sns.barplot(data=sub_df,x='block_j',y='mod',hue='drug_type',ax=ax,order=blocks,hue_order=drug_order,palette=cmap_macro,errorbar='se',legend=legend)
                    # ax.set_ylim([0,0.6])
                    usrplt.adjust_spines(ax)
                    ax.set_xlabel('Time after 2nd injection (min)'); ax.set_ylabel('Modulation index')

                    for i, b in enumerate(blocks):
                        distr_i = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[0])]['mod'].values
                        distr_j = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[1])]['mod'].values
                        if (len(distr_i) < 5) | (len(distr_j) < 5):
                            continue
                        res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')             
                        if res.pvalue < pval_thresh:
                            ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])
                
                    plt.savefig(join(PlotDir,f'mesogroup_barplot_{stat}_{block_i}_{mg}_{ct}_{bstrs_list[iB]}.png'),dpi=300,bbox_inches='tight',facecolor='w')

                    prs = prs_list[iS]
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = f'{mg}, {ct}, {b_str}'
                    usrplt.save_fig_to_pptx(fig, prs,slide)
                    plt.close(fig)
                    nPlots+=1

    ##------------------------------------------
    #Plot MI over time per mesogroup, celltype, behavior, & LAYER
    print('Plotting mesogroup barplots over time, per celltype, behavior, mesogroup, and layer')
    for iS, stat in enumerate(stat_list):
        for ct in ['RS','FS']:
            for iB, b_str in behavior_dict2.items():
                PlotDir = join(BaseDir,'mesogroup',bstrs_list[iB])
                if not os.path.exists(PlotDir):
                    os.makedirs(PlotDir)
                for mg in CTX_mesogroup_list:
                    for l in layer_list:
                        if l == '2/3':
                            l_str = '2-3'
                        else:
                            l_str = l

                        sub_df = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.mesogroup == mg) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                        if len(sub_df) < 5:
                            continue
                        fig, axes = plt.subplots(3,1,figsize=(10,15))
                        plt.suptitle(f'Change in {stat_title_list[iS]} due to saline or psilocybin injection\nBaseline = {block_i}, {mg}, layer {l}, {ct}, {b_str}',y=0.925)

                        #Fraction significantly modulated
                        ax = axes[0]
                        sub_df = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.mesogroup == mg) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                        iX = 0; xticks = []
                        for i, b in enumerate(blocks):
                            for jj, d in enumerate(drug_order):
                                sub_sub_df = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == d)]
                                n = len(sub_sub_df)
                                if n == 0:
                                    continue
                                nSig = np.sum(sub_sub_df['sig'].values)
                                fracSig = nSig/n
                                pos_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] > 0)])/n
                                neg_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] < 0)])/n
                                not_mod = len(sub_sub_df.loc[sub_sub_df.sig == 0])/n

                                h1 = ax.bar(iX, pos_mod, color=sns.xkcd_palette(['light red']), label='Positively Modulated')
                                h2 = ax.bar(iX, not_mod, bottom=pos_mod, color=sns.xkcd_palette(['silver']), label='Not Modulated')
                                h3 = ax.bar(iX, neg_mod, bottom=pos_mod+not_mod, color=sns.xkcd_palette(['medium blue']), label='Negatively Modulated')
                                n = len(np.unique(sub_sub_df['neuron_index']))
                                ax.text(iX,0.025,f'{n}',ha='center',fontsize=8,rotation=90)
                                if jj == 1:
                                    xticks.append(iX-0.5)
                                iX += 1
                            iX += 1

                        usrplt.adjust_spines(ax)
                        plt.autoscale(tight=True)
                        xticks = np.arange(1,3*len(blocks),3)-0.5 #[0.5, 3.5, 6.5, 9.5,12.5,]
                        ax.set_xticks(xticks)
                        ax.set_xticklabels(blocks,rotation=0)
                        ax.set_ylabel('Fraction modulated')
                        ax.legend(handles=[h1,h2,h3],loc=2)

                        #Modulation index
                        ax = axes[1]
                        sub_df = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.mesogroup == mg) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                        sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                        sns.barplot(data=sub_df,x='block_j',y='abs_mod',hue='drug_type',ax=ax,order=blocks,hue_order=drug_order,palette=cmap_macro,errorbar='se',legend=legend)

                        usrplt.adjust_spines(ax)
                        ax.set_xlabel('Time after 2nd injection (min)'); ax.set_ylabel('| Modulation index |')

                        for i, b in enumerate(blocks):
                            distr_i = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[0])]['abs_mod'].values
                            distr_j = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[1])]['abs_mod'].values
                            if (len(distr_i) < 5) | (len(distr_j) < 5):
                                continue
                            res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')
                                                
                            if res.pvalue < pval_thresh:
                                ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])

                        #Modulation index
                        ax = axes[2]
                        sub_df = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.mesogroup == mg) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                        sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                        sns.barplot(data=sub_df,x='block_j',y='mod',hue='drug_type',ax=ax,order=blocks,hue_order=drug_order,palette=cmap_macro,errorbar='se',legend=legend)
                        usrplt.adjust_spines(ax)
                        ax.set_xlabel('Time after 2nd injection (min)'); ax.set_ylabel('Modulation index')

                        for i, b in enumerate(blocks):
                            distr_i = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[0])]['mod'].values
                            distr_j = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[1])]['mod'].values
                            if (len(distr_i) < 5) | (len(distr_j) < 5):
                                continue
                            res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')
                                                
                            if res.pvalue < pval_thresh:
                                ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])
                    
                        plt.savefig(join(PlotDir,f'mesogroup_barplot_{stat}_{block_i}_{mg}_layer-{l_str}_{ct}_{bstrs_list[iB]}.png'),dpi=300,bbox_inches='tight',facecolor='w')
                        prs = prs_list[iS]
                        slide = prs.slides.add_slide(prs.slide_layouts[1])
                        slide.shapes.title.text = f'{mg}, layer {l}, {ct}, {b_str}'
                        usrplt.save_fig_to_pptx(fig, prs,slide)
                        plt.close(fig)
                        nPlots+=1

    PlotDir = join(BaseDir,'mesogroup')
    for iS, stat in enumerate(stat_list):
        prs = prs_list[iS]
        prs.save(join(PlotDir,f'mesogroup_barplots_{stat}.pptx'))

    ##------------------------------------------
    # Create a new powerPoint presentations to save figures to
    prs_list = []
    for iS, stat in enumerate(stat_list):
        prs = Presentation()

        #Add title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f'All Mice, neurons organized in "group" level'
        slide.placeholders[1].text = f'Change in {stat_title_list[iS]} due to saline or psilocybin injection'
        prs_list.append(prs)

    ##------------------------------------------
    #Plot MI over GROUPS per time window, celltype, behavior
    print('Plotting group barplots over regions, per celltype, behavior, and time window')
    for iS, stat in enumerate(stat_list):
        for ct in ['RS','FS']:
            for iB, b_str in behavior_dict2.items():
                PlotDir = join(BaseDir,'group',bstrs_list[iB])
                if not os.path.exists(PlotDir):
                    os.makedirs(PlotDir)
                for block_j in blocks:
                    fig, axes = plt.subplots(3,1,figsize=(10,15))
                    if block_j == 'post_1st_inj':
                        plt.suptitle(f'Change in {stat_title_list[iS]} due to saline or psilocybin injection\nCelltype = {ct}, behavior = {b_str}\n {block_i} vs {block_j}',y=0.94)
                    else:
                        plt.suptitle(f'Change in {stat_title_list[iS]} due to saline or psilocybin injection\nCelltype = {ct}, behavior = {b_str}\n {block_i} vs {block_j} minute after 2nd injection',y=0.94)

                    #Fraction significantly modulated
                    ax = axes[0]
                    sub_df = dFR_spont_df.loc[(dFR_spont_df['group'].isin(group_list)) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j == block_j)]
                    iX = 0; xticks = []
                    for i, g in enumerate(group_list):
                        for jj, d in enumerate(drug_order):
                            sub_sub_df = sub_df.loc[(sub_df['group'] == g) & (sub_df.drug_type == d)]
                            n = len(sub_sub_df)
                            if n == 0:
                                continue
                            nSig = np.sum(sub_sub_df['sig'].values)
                            fracSig = nSig/n
                            pos_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] > 0)])/n
                            neg_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] < 0)])/n
                            not_mod = len(sub_sub_df.loc[sub_sub_df.sig == 0])/n

                            h1 = ax.bar(iX, pos_mod, color=sns.xkcd_palette(['light red']), label='Positively Modulated')
                            h2 = ax.bar(iX, not_mod, bottom=pos_mod, color=sns.xkcd_palette(['silver']), label='Not Modulated')
                            h3 = ax.bar(iX, neg_mod, bottom=pos_mod+not_mod, color=sns.xkcd_palette(['medium blue']), label='Negatively Modulated')
                            n = len(np.unique(sub_sub_df['neuron_index']))
                            ax.text(iX,0.025,f'{n}',ha='center',fontsize=8,rotation=90)
                            if jj == 1:
                                xticks.append(iX-0.5)
                            iX += 1
                        iX += 1

                    usrplt.adjust_spines(ax)
                    plt.autoscale(tight=True)
                    xticks = np.arange(1,len(group_list)*3,3)-0.5
                    ax.set_xticks(xticks)
                    ax.set_xticklabels(group_list,rotation=0)
                    ax.set_ylabel('Fraction modulated')
                    ax.legend(handles=[h1,h2,h3],loc=2)

                    #Modulation index
                    ax = axes[1]
                    sub_df = dFR_spont_df.loc[(dFR_spont_df['group'].isin(group_list)) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j == block_j)]
                    sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                    sns.barplot(data=sub_df,x='group',y='abs_mod',hue='drug_type',ax=ax,order=group_list,hue_order=drug_order,palette=cmap_macro,errorbar='se',legend=legend)
                    # ax.set_ylim([0,0.6])
                    usrplt.adjust_spines(ax)
                    ax.set_xlabel('Mesogroup'); ax.set_ylabel('| Modulation index |')

                    for i, mg in enumerate(group_list):
                        distr_i = sub_df.loc[(sub_df.group == mg) & (sub_df.drug_type == drug_order[0])]['abs_mod'].values
                        distr_j = sub_df.loc[(sub_df.group == mg) & (sub_df.drug_type == drug_order[1])]['abs_mod'].values
                        if (len(distr_i) < 2) | (len(distr_j) < 2):
                            continue
                        res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')                  
                        if res.pvalue < pval_thresh:
                            ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])

                    #Modulation index
                    ax = axes[2]
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.group.isin(group_list)) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j == block_j)]
                    sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                    sns.barplot(data=sub_df,x='group',y='mod',hue='drug_type',ax=ax,order=group_list,hue_order=drug_order,palette=cmap_macro,errorbar='se',legend=legend)
                    usrplt.adjust_spines(ax)
                    ax.set_xlabel('Mesogroup'); ax.set_ylabel('Modulation index')

                    for i, mg in enumerate(group_list):
                        distr_i = sub_df.loc[(sub_df.group == mg) & (sub_df.drug_type == drug_order[0])]['mod'].values
                        distr_j = sub_df.loc[(sub_df.group == mg) & (sub_df.drug_type == drug_order[1])]['mod'].values
                        if (len(distr_i) < 2) | (len(distr_j) < 2):
                            continue
                        res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')
                                            
                        if res.pvalue < pval_thresh:
                            ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])
                    plt.savefig(join(PlotDir,f'group_barplot_{stat}_{ct}_{bstrs_list[iB]}_{epoch_i}_vs_{block_j}.png'),dpi=300,bbox_inches='tight',facecolor='w')
                    prs = prs_list[iS]
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = f'{epoch_i}_vs_{block_j}, {ct}, {b_str}'
                    usrplt.save_fig_to_pptx(fig, prs,slide)
                    plt.close(fig)
                    nPlots+=1

    ##------------------------------------------
    #Plot MI over time per group, celltype, behavior
    print('Plotting group barplots over time, per celltype, behavior, and group')
    for iS, stat in enumerate(stat_list):
        for ct in ['RS','FS']:
            for iB, b_str in behavior_dict2.items():
                PlotDir = join(BaseDir,'group',bstrs_list[iB])
                if not os.path.exists(PlotDir):
                    os.makedirs(PlotDir)

                for g in group_list:
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.group == g) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                    if len(sub_df) < 5:
                        continue
                    fig, axes = plt.subplots(3,1,figsize=(10,15))
                    plt.suptitle(f'Change in {stat_title_list[iS]} due to saline or psilocybin injection\nBaseline = {block_i}, {g}, {ct}, {b_str}',y=0.925)

                    #Fraction significantly modulated
                    ax = axes[0]
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.group == g) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j.isin(blocks))]
                    iX = 0; xticks = []
                    for i, b in enumerate(blocks):
                        for jj, d in enumerate(drug_order):
                            sub_sub_df = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == d)]
                            n = len(sub_sub_df)
                            if n == 0:
                                continue
                            nSig = np.sum(sub_sub_df['sig'].values)
                            fracSig = nSig/n
                            pos_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] > 0)])/n
                            neg_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] < 0)])/n
                            not_mod = len(sub_sub_df.loc[sub_sub_df.sig == 0])/n

                            h1 = ax.bar(iX, pos_mod, color=sns.xkcd_palette(['light red']), label='Positively Modulated')
                            h2 = ax.bar(iX, not_mod, bottom=pos_mod, color=sns.xkcd_palette(['silver']), label='Not Modulated')
                            h3 = ax.bar(iX, neg_mod, bottom=pos_mod+not_mod, color=sns.xkcd_palette(['medium blue']), label='Negatively Modulated')
                            n = len(np.unique(sub_sub_df['neuron_index']))
                            ax.text(iX,0.025,f'{n}',ha='center',fontsize=8,rotation=90)
                            if jj == 1:
                                xticks.append(iX-0.5)
                            iX += 1
                        iX += 1

                    usrplt.adjust_spines(ax)
                    plt.autoscale(tight=True)
                    # ax.hlines(0.5,*ax.get_xlim(),linestyles='--',color='k',alpha=0.5)
                    xticks = np.arange(1,len(blocks)*3,3)-0.5
                    ax.set_xticks(xticks)
                    ax.set_xticklabels(blocks,rotation=0)
                    ax.set_ylabel('Fraction modulated')
                    ax.legend(handles=[h1,h2,h3],loc=2)

                    #Modulation index
                    ax = axes[1]
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.group == g) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                    sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                    sns.barplot(data=sub_df,x='block_j',y='abs_mod',hue='drug_type',ax=ax,order=blocks,hue_order=drug_order,palette=cmap_macro,errorbar='se',legend=legend)
                    usrplt.adjust_spines(ax)
                    ax.set_xlabel('Time after 2nd injection (min)'); ax.set_ylabel('| Modulation index |')

                    for i, b in enumerate(blocks):

                        distr_i = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[0])]['abs_mod'].values
                        distr_j = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[1])]['abs_mod'].values
                        if (len(distr_i) < 5) | (len(distr_j) < 5):
                            continue
                        res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')
                                            
                        if res.pvalue < pval_thresh:
                            ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])

                    #Modulation index
                    ax = axes[2]
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.group == g) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                    sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                    sns.barplot(data=sub_df,x='block_j',y='mod',hue='drug_type',ax=ax,order=blocks,hue_order=drug_order,palette=cmap_macro,errorbar='se',legend=legend)
                    usrplt.adjust_spines(ax)
                    ax.set_xlabel('Time after 2nd injection (min)'); ax.set_ylabel('Modulation index')
                    for i, b in enumerate(blocks):

                        distr_i = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[0])]['mod'].values
                        distr_j = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[1])]['mod'].values
                        if (len(distr_i) < 5) | (len(distr_j) < 5):
                            continue
                        res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')
                                            
                        if res.pvalue < pval_thresh:
                            ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])
                    
                    plt.savefig(join(PlotDir,f'group_barplot_{stat}_{block_i}_{g}_{ct}_{bstrs_list[iB]}.png'),dpi=300,bbox_inches='tight',facecolor='w')
                    prs = prs_list[iS]
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = f'{g}, {ct}, {b_str}'
                    usrplt.save_fig_to_pptx(fig, prs,slide)
                    plt.close(fig)
                    nPlots+=1

    ##------------------------------------------
    #Plot MI over time per group, celltype, behavior, & LAYER
    print('Plotting group barplots over time, per celltype, behavior, group, and layer')
    for iS, stat in enumerate(stat_list):
        for ct in ['RS','FS']:
            for iB, b_str in behavior_dict2.items():
                PlotDir = join(BaseDir,'group',bstrs_list[iB])
                if not os.path.exists(PlotDir):
                    os.makedirs(PlotDir)

                for g in CTX_group_list:
                    for l in layer_list:
                        if l == '2/3':
                            l_str = '2-3'
                        else:
                            l_str = l
                        sub_df0 = dFR_spont_df.loc[(dFR_spont_df.drug_type == drug_order[0]) & (dFR_spont_df.group == g) & (dFR_spont_df.layer == l) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                        sub_df1 = dFR_spont_df.loc[(dFR_spont_df.drug_type == drug_order[1]) & (dFR_spont_df.group == g) & (dFR_spont_df.layer == l) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                        if (len(sub_df0) < 5) | (len(sub_df1) < 5):
                            continue
                        fig, axes = plt.subplots(3,1,figsize=(10,15))
                        plt.suptitle(f'Change in {stat_title_list[iS]} due to saline or psilocybin injection\nBaseline = {block_i}, {g}, layer {l}, {ct}, {b_str}',y=0.925)

                        #Fraction significantly modulated
                        ax = axes[0]
                        sub_df = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.group == g) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                        iX = 0; xticks = []
                        for i, b in enumerate(blocks):
                            for jj, d in enumerate(drug_order):
                                sub_sub_df = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == d)]
                                n = len(sub_sub_df)
                                if n == 0:
                                    continue
                                nSig = np.sum(sub_sub_df['sig'].values)
                                fracSig = nSig/n
                                pos_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] > 0)])/n
                                neg_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] < 0)])/n
                                not_mod = len(sub_sub_df.loc[sub_sub_df.sig == 0])/n

                                h1 = ax.bar(iX, pos_mod, color=sns.xkcd_palette(['light red']), label='Positively Modulated')
                                h2 = ax.bar(iX, not_mod, bottom=pos_mod, color=sns.xkcd_palette(['silver']), label='Not Modulated')
                                h3 = ax.bar(iX, neg_mod, bottom=pos_mod+not_mod, color=sns.xkcd_palette(['medium blue']), label='Negatively Modulated')
                                n = len(np.unique(sub_sub_df['neuron_index']))
                                ax.text(iX,0.025,f'{n}',ha='center',fontsize=8,rotation=90)
                                if jj == 1:
                                    xticks.append(iX-0.5)
                                iX += 1
                            iX += 1

                        usrplt.adjust_spines(ax)
                        plt.autoscale(tight=True)
                        xticks = np.arange(1,len(blocks)*3,3)-0.5
                        ax.set_xticks(xticks)
                        ax.set_xticklabels(blocks,rotation=0)
                        ax.set_ylabel('Fraction modulated')
                        ax.legend(handles=[h1,h2,h3],loc=2)

                        #Modulation index
                        ax = axes[1]
                        sub_df = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.group == g) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                        sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                        sns.barplot(data=sub_df,x='block_j',y='abs_mod',hue='drug_type',ax=ax,order=blocks,hue_order=drug_order,palette=cmap_macro,errorbar='se',legend=legend)
                        usrplt.adjust_spines(ax)
                        ax.set_xlabel('Time after 2nd injection (min)'); ax.set_ylabel('| Modulation index |')

                        for i, b in enumerate(blocks):

                            distr_i = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[0])]['abs_mod'].values
                            distr_j = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[1])]['abs_mod'].values
                            if (len(distr_i) < 5) | (len(distr_j) < 5):
                                continue
                            res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')              
                            if res.pvalue < pval_thresh:
                                ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])

                        #Modulation index
                        ax = axes[2]
                        sub_df = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.group == g) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                        sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                        sns.barplot(data=sub_df,x='block_j',y='mod',hue='drug_type',ax=ax,order=blocks,hue_order=drug_order,palette=cmap_macro,errorbar='se',legend=legend)
                        # ax.set_ylim([0,0.6])
                        usrplt.adjust_spines(ax)
                        ax.set_xlabel('Time after 2nd injection (min)'); ax.set_ylabel('Modulation index')


                        for i, b in enumerate(blocks):

                            distr_i = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[0])]['mod'].values
                            distr_j = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[1])]['mod'].values
                            if (len(distr_i) < 5) | (len(distr_j) < 5):
                                continue
                            res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')          
                            if res.pvalue < pval_thresh:
                                ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])
                        
                        plt.savefig(join(PlotDir,f'group_barplot_{stat}_{block_i}_{g}_layer-{l_str}_{ct}_{bstrs_list[iB]}.png'),dpi=300,bbox_inches='tight',facecolor='w')
                        prs = prs_list[iS]
                        slide = prs.slides.add_slide(prs.slide_layouts[1])
                        slide.shapes.title.text = f'{g}, layer {l}, {ct}, {b_str}'
                        usrplt.save_fig_to_pptx(fig, prs,slide)
                        plt.close(fig)
                        nPlots+=1

    PlotDir = join(BaseDir,'group')
    for iS, stat in enumerate(stat_list):
        prs = prs_list[iS]
        prs.save(join(PlotDir,f'group_barplots_{stat}.pptx'))

    ##------------------------------------------
    # Create a new powerPoint presentations to save figures to
    prs_list = []
    for iS, stat in enumerate(stat_list):
        prs = Presentation()

        #Add title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f'All Mice, neurons organized in "group" level'
        slide.placeholders[1].text = f'Change in {stat_title_list[iS]} due to saline or psilocybin injection'
        prs_list.append(prs)


    ##------------------------------------------
    #Determine the shared areas that are uniq for each condition 
    uniq_areas = ['ACAd5', 'ACAd6a', 'ACAv6a', 'ACB', 'AD', 'AMd', 'AON', 'APN',
        'AV', 'CA1', 'CA2', 'CA3', 'CL', 'CP', 'DG-mo', 'DG-po', 'DG-sg',
        'DP', 'Eth', 'FRP5', 'FRP6a', 'HY', 'IAD', 'ILA1', 'ILA2/3',
        'ILA5', 'ILA6a', 'LD', 'LGd-co', 'LGd-ip', 'LGd-sh', 'LP', 'LSr',
        'LSv', 'MGd', 'MGm', 'MGv', 'MOp2/3', 'MOp5', 'MOp6a', 'MOs2/3',
        'MOs5', 'MOs6a', 'MOs6b', 'MS', 'OLF', 'ORBm1', 'ORBm2/3', 'ORBm5',
        'ORBvl6a', 'PAL', 'PCN', 'PIL', 'PL2/3', 'PL5', 'PL6a', 'PO',
        'PoT', 'ProS', 'RSPagl2/3', 'RSPagl5', 'RT', 'SGN', 'SSp-bfd1',
        'SSp-bfd2/3', 'SSp-bfd4', 'SSp-bfd5', 'SSp-bfd6a', 'SSp-bfd6b',
        'SSp-ll4', 'SSp-ll5', 'SSp-ll6a', 'SSp-ll6b', 'SSp-tr2/3',
        'SSp-tr4', 'SSp-tr5', 'SSp-tr6a', 'SSp-tr6b', 'STR', 'SUB', 'TH',
        'TTd', 'VAL', 'VISa2/3', 'VISa4', 'VISa5', 'VISp1', 'VISp2/3',
        'VISp4', 'VISp5', 'VISp6a', 'VISp6b', 'VISrl4', 'VISrl5',
        'VISrl6a', 'VISrl6b', 'VPL', 'VPM', 'ZI', 'alv', 'ar', 'ccb',
        'ccg', 'ccs', 'cing', 'em', 'fa', 'fi','fp',
        'int', 'ml', 'nan', 'or','scwm', 'st']
    
    ##------------------------------------------
    # iS = 0; stat = stat_list[iS]
    prs_list = []
    for iS, stat in enumerate(stat_list):
        #Area level barplots
        prs = Presentation()

        #Add title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f'All Mice, neurons organized in "area" level'
        slide.placeholders[1].text = f'Change in {stat_title_list[iS]} due to saline or psilocybin injection'
    prs_list.append(prs)
    

    for iS, stat in enumerate(stat_list):
        for ct in ['RS','FS']:
            for iB, b_str in behavior_dict2.items():
                PlotDir = join(BaseDir,'area',bstrs_list[iB])
                if not os.path.exists(PlotDir):
                    os.makedirs(PlotDir)
                for a in uniq_areas:
                    print(f'{ct}, {b_str}, {a}')
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.area == a) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                    if len(sub_df) < 5:
                        continue
                    fig, axes = plt.subplots(3,1,figsize=(10,15))
                    plt.suptitle(f'Change in {stat_title_list[iS]} due to saline or psilocybin injection\nBaseline = {block_i}, {a}, {ct}, {b_str}',y=0.925)

                    #Fraction significantly modulated
                    ax = axes[0]
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.area == a) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]

                    iX = 0; xticks = []
                    for i, b in enumerate(blocks):
                        for jj, d in enumerate(drug_order):
                            sub_sub_df = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == d)]

                            n = len(sub_sub_df)
                            if n == 0:
                                continue
                            nSig = np.sum(sub_sub_df['sig'].values)
                            fracSig = nSig/n
                            pos_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] > 0)])/n
                            neg_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] < 0)])/n
                            not_mod = len(sub_sub_df.loc[sub_sub_df.sig == 0])/n

                            h1 = ax.bar(iX, pos_mod, color=sns.xkcd_palette(['light red']), label='Positively Modulated')
                            h2 = ax.bar(iX, not_mod, bottom=pos_mod, color=sns.xkcd_palette(['silver']), label='Not Modulated')
                            h3 = ax.bar(iX, neg_mod, bottom=pos_mod+not_mod, color=sns.xkcd_palette(['medium blue']), label='Negatively Modulated')
                            
                            n = len(np.unique(sub_sub_df['neuron_index']))
                            ax.text(iX,0.025,f'{n}',ha='center',fontsize=8,rotation=90)
                            if jj == 1:
                                xticks.append(iX-0.5)
                            iX += 1
                        iX += 1

                    usrplt.adjust_spines(ax)
                    plt.autoscale(tight=True)
                    xticks = np.arange(1,3*len(blocks),3)-0.5
                    ax.set_xticks(xticks)
                    ax.set_xticklabels(blocks,rotation=0)
                    ax.set_ylabel('Fraction modulated')
                    ax.legend(handles=[h1,h2,h3],loc=2)

                    #Modulation index
                    ax = axes[1]
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.area == a) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                    sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                    sns.barplot(data=sub_df,x='block_j',y='abs_mod',hue='drug_type',ax=ax,order=blocks,hue_order=drug_order,palette=cmap_macro,errorbar='se',legend=legend)
                    # ax.set_ylim([0,0.6])
                    usrplt.adjust_spines(ax)
                    ax.set_xlabel('Time after 2nd injection (min)'); ax.set_ylabel('| Modulation index |')

                    for i, b in enumerate(blocks):
                        distr_i = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[0])]['abs_mod'].values
                        distr_j = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[1])]['abs_mod'].values
                        if (len(distr_i) < 5) | (len(distr_j) < 5):
                            continue
                        res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')
                                            
                        if res.pvalue < pval_thresh:
                            ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])

                    #Modulation index (dFR_spont_df.sig == 1) &
                    ax = axes[2]
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.area == a) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                    sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                    sns.barplot(data=sub_df,x='block_j',y='mod',hue='drug_type',ax=ax,order=blocks,hue_order=drug_order,palette=cmap_macro,errorbar='se',legend=legend)
                    # ax.set_ylim([0,0.6])
                    usrplt.adjust_spines(ax)
                    ax.set_xlabel('Time after 2nd injection (min)'); ax.set_ylabel('Modulation index')

                    for i, b in enumerate(blocks):
                        distr_i = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[0])]['mod'].values
                        distr_j = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[1])]['mod'].values
                        if (len(distr_i) < 5) | (len(distr_j) < 5):
                            continue
                        res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')             
                        if res.pvalue < pval_thresh:
                            ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])
                
                    a_str = a.replace('/','_')
                    plt.savefig(join(PlotDir,f'area_barplot_{stat}_{block_i}_{a_str}_{ct}_{bstrs_list[iB]}.png'),dpi=300,bbox_inches='tight',facecolor='w')

                    prs = prs_list[iS]
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = f'{a}, {ct}, {b_str}'
                    usrplt.save_fig_to_pptx(fig, prs,slide)
                    plt.close(fig)
                    nPlots+=1

    PlotDir = join(BaseDir,'area')
    for iS, stat in enumerate(stat_list):
        prs = prs_list[iS]
        prs.save(join(PlotDir,f'area_barplots_{stat}.pptx'))
    




    # ##------------------------------------------
    # #Plot scatter plot group, time window, celltype
    # PlotDir = join(BaseDir,'group','scatter')
    # if not os.path.exists(PlotDir):
    #     os.makedirs(PlotDir)
    # print('Plotting group scatter plots per region, celltype, and time window')
    # for iS, stat in enumerate(stat_list):
    #     for ct in ['RS','FS']:
    #         for g in TH_group_list:
    #             for block_j in blocks:

    #                 sub_df1 = dFR_spont_df.loc[(dFR_spont_df.group == g) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.drug_type == drug_order[0]) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j == block_j) & (dFR_spont_df.behavior == 'all (>0cm/s)')]
    #                 sub_df2 = dFR_spont_df.loc[(dFR_spont_df.group == g) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.drug_type == drug_order[1]) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j == block_j) & (dFR_spont_df.behavior == 'all (>0cm/s)')]
    #                 if (len(sub_df1) < 2) | (len(sub_df2) < 2):
    #                     continue
    #                 fig, axes = plt.subplots(3,4,figsize=(20,15),gridspec_kw={'hspace':0.35,'wspace':0.35})
    #                 if block_j == 'post_1st_inj':
    #                     block_j_str = 'Post 1st saline inj'
    #                 else:
    #                     block_j_str = f'{block_j} mins past 2nd injection'

    #                 plt.suptitle(f'Change in {stat_title_list[iS]} due to saline or psilocybin injection\nGroup: {g}, cellltype = {ct}, {block_i} vs {block_j_str}',y=0.95)

    #                 for iB, b_str in behavior_dict2.items():
    #                     sub_df1 = dFR_spont_df.loc[(dFR_spont_df.group == g) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.drug_type == drug_order[0]) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j == block_j) & (dFR_spont_df.behavior == b_str)]
    #                     sub_df2 = dFR_spont_df.loc[(dFR_spont_df.group == g) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.drug_type == drug_order[1]) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j == block_j) & (dFR_spont_df.behavior == b_str)]
    #                     sub_df1 = sub_df1.loc[sub_df1['sig'] == 1] if plot_only_sig_mod else sub_df1
    #                     sub_df2 = sub_df2.loc[sub_df2['sig'] == 1] if plot_only_sig_mod else sub_df2
    #                     if (len(sub_df1) < 2) | (len(sub_df2) < 2):
    #                         continue
                            
    #                     delta_list = []; xymin_list = [];  xymax_list = []
    #                     for ii, d in enumerate(drug_order):

    #                         sub_df = dFR_spont_df.loc[(dFR_spont_df.group == g) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.drug_type == d) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j == block_j) & (dFR_spont_df.behavior == b_str)]
    #                         sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
    #                         FR_i = sub_df['val_i'].values
    #                         FR_j = sub_df['val_j'].values
    #                         delta = sub_df['delta'].values
    #                         cts = sub_df['celltype'].values
    #                         mod = sub_df['mod'].values
    #                         delta_list.append((FR_i,FR_j,delta,mod))

    #                         nN = len(np.unique(sub_df['neuron_index'].values)); nE = len(np.unique(sub_df['rec_name'].values))
    #                         ax = axes[iB,ii]
    #                         ax.set_title(f'{stat} during {b_str}: {d} \n(n = {nN} neurons from {nE} exps)')
    #                         ax.scatter(FR_i,FR_j,s=10,marker='o',edgecolor=cmap_macro[ii],facecolor='none')

    #                         ymin = np.nanmin(FR_j); xmin = np.nanmin(FR_i)
    #                         ymax = np.nanpercentile(FR_j,98); xmax = np.nanpercentile(FR_i,98)
    #                         xymin = np.nanmin([ymin,xmin]);xymin_list.append(xymin)
    #                         xymax = np.nanmax([ymax,xmax]);xymax_list.append(xymax)

    #                     xymin = np.nanmin(xymin_list); xymax = np.nanmax(xymax_list)   
    #                     for ii, d in enumerate(drug_order):
    #                         ax = axes[iB,ii]
    #                         FR_i, FR_j, delta, mod = delta_list[ii]
                            
    #                         ax.plot([xymin,xymax],[xymin,xymax],'-k')
    #                         ax.set_xlim([xymin,xymax]); ax.set_ylim([xymin,xymax])
    #                         usrplt.adjust_spines(ax)
    #                         ax.set_xlabel(f'{block_i}')
    #                         ax.set_ylabel(f'{block_j}')
    #                         res = st.wilcoxon(FR_i,FR_j,nan_policy='omit')
    #                         if res.pvalue < pval_thresh:
    #                             ax.text(0.5*xymax,0.9*xymax,f'*',color=cmap_macro[ii],fontweight='bold',fontsize=20)

    #                     #Compare saline vs psilocybin
    #                     distr_i = delta_list[0][-1]
    #                     distr_j = delta_list[1][-1]

    #                     FR_i, FR_j, delta, mod_i = delta_list[0]
    #                     FR_i, FR_j, delta, mod_j = delta_list[1]
    #                     ax = axes[iB,2]
    #                     sns.histplot([mod_i,mod_j],palette=cmap_macro,binwidth=0.1,multiple='layer',stat='density',common_norm=False,kde=False,ax=ax)

    #                     ax.set_title('Modulation index')
    #                     ax.set_xlabel('MI')
    #                     ax.set_xlim([-1,1])
    #                     usrplt.adjust_spines(ax)
    #                     ylim = ax.get_ylim(); xlim = ax.get_xlim()
    #                     ax.vlines(np.nanmean(mod_i),*ylim,ls='-',color=cmap_macro[0],lw=2,label=f'Saline MI: {np.nanmean(mod_i):.2f}')
    #                     ax.vlines(np.nanmean(mod_j),*ylim,ls='-',color=cmap_macro[1],lw=2,label=f'Psilocybin MI: {np.nanmean(mod_j):.2f}')
    #                     ax.legend()

    #                     if (len(mod_i) > 0) & (len(mod_j) > 0):
    #                         res = st.mannwhitneyu(mod_i,mod_j,nan_policy='omit')
    #                         if res.pvalue < pval_thresh/3:
    #                             ax.text(.5,0.96*ax.get_ylim()[1],'*',fontsize=20,fontweight='bold',color='k')

    #                     ax = axes[iB,3]
    #                     sns.histplot([np.abs(mod_i),np.abs(mod_j)],palette=cmap_macro,binwidth=0.1,multiple='layer',stat='density',common_norm=False,kde=False,ax=ax)
    #                     ax.set_title('| Modulation index |')
    #                     ax.set_xlabel('|MI|')
    #                     ax.set_xlim([0,1])
    #                     usrplt.adjust_spines(ax)
    #                     ylim = ax.get_ylim(); xlim = ax.get_xlim()
    #                     ax.vlines(np.nanmean(np.abs(mod_i)),*ylim,ls='--',color=cmap_macro[0],lw=2,label=f'Saline |MI|: {np.nanmean(np.abs(mod_i)):.2f}')
    #                     ax.vlines(np.nanmean(np.abs(mod_j)),*ylim,ls='--',color=cmap_macro[1],lw=2,label=f'Psilocybin |MI|: {np.nanmean(np.abs(mod_j)):.2f}')
    #                     ax.legend()

    #                     if (len(mod_i) > 0) & (len(mod_j) > 0):
    #                         res = st.mannwhitneyu(np.abs(mod_i),np.abs(mod_j),nan_policy='omit')
    #                         if res.pvalue < pval_thresh:
    #                             ax.text(.5,0.96*ax.get_ylim()[1],'*',fontsize=20,fontweight='bold',color='k')

    #                 # pdb.set_trace()
    #                 plt.savefig(join(PlotDir,f'group_scatter_{stat}_{g}_{ct}_{block_i}_vs_{block_j}.png'),dpi=300,bbox_inches='tight',facecolor='w')
    #                 prs = prs_list[iS]
    #                 slide = prs.slides.add_slide(prs.slide_layouts[1])
    #                 slide.shapes.title.text = f'{g}, {ct}, {block_i}_vs_{block_j}'
    #                 usrplt.save_fig_to_pptx(fig, prs,slide)
    #                 plt.close(fig)
    #                 nPlots+=1

    # ##------------------------------------------
    # #Plot scatter plot group, time window, celltype, & LAYER
    # print('Plotting group scatter plots per region, celltype, layer, and time window')
    # for iS, stat in enumerate(stat_list):
    #     for ct in ['RS','FS']:
    #         for g in CTX_group_list:
    #             layers = np.unique(cells_df.loc[cells_df.group == g]['layer'].values)
    #             for l in layers:
    #                 if l == '2/3':
    #                     l_str = '2-3'
    #                 else:
    #                     l_str = l
    #                 # l = '5'; l_str = '5'
    #                 for block_j in blocks:

    #                     sub_df1 = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.group == g) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.drug_type == drug_order[0]) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j == block_j) & (dFR_spont_df.behavior == 'all (>0cm/s)')]
    #                     sub_df2 = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.group == g) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.drug_type == drug_order[1]) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j == block_j) & (dFR_spont_df.behavior == 'all (>0cm/s)')]
    #                     sub_df1 = sub_df1.loc[sub_df1['sig'] == 1] if plot_only_sig_mod else sub_df1
    #                     sub_df2 = sub_df2.loc[sub_df2['sig'] == 1] if plot_only_sig_mod else sub_df2
    #                     if (len(sub_df1) < 2) | (len(sub_df2) < 2):
    #                         continue
    #                     fig, axes = plt.subplots(3,4,figsize=(20,15),gridspec_kw={'hspace':0.35,'wspace':0.35})
    #                     if block_j == 'post_1st_inj':
    #                         block_j_str = 'Post 1st saline inj'
    #                     else:
    #                         block_j_str = f'{block_j} mins past 2nd injection'

    #                     plt.suptitle(f'Change in {stat_title_list[iS]} due to saline or psilocybin injection\nGroup: {g}, layer {l}, cellltype = {ct}, {block_i} vs {block_j_str}',y=0.95)

    #                     for iB, b_str in behavior_dict2.items():
    #                         sub_df1 = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.group == g) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.drug_type == drug_order[0]) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j == block_j) & (dFR_spont_df.behavior == b_str)]
    #                         sub_df2 = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.group == g) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.drug_type == drug_order[1]) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j == block_j) & (dFR_spont_df.behavior == b_str)]
    #                         sub_df1 = sub_df1.loc[sub_df1['sig'] == 1] if plot_only_sig_mod else sub_df1
    #                         sub_df2 = sub_df2.loc[sub_df2['sig'] == 1] if plot_only_sig_mod else sub_df2
    #                         if (len(sub_df1) < 2) | (len(sub_df2) < 2):
    #                             continue
                                
    #                         delta_list = []; xymin_list = [];  xymax_list = []
    #                         for ii, d in enumerate(drug_order):

    #                             sub_df = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.group == g) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.drug_type == d) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j == block_j) & (dFR_spont_df.behavior == b_str)]
    #                             sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
    #                             FR_i = sub_df['val_i'].values
    #                             FR_j = sub_df['val_j'].values
    #                             delta = sub_df['delta'].values
    #                             cts = sub_df['celltype'].values
    #                             mod = sub_df['mod'].values
    #                             delta_list.append((FR_i,FR_j,delta,mod))

    #                             nN = len(np.unique(sub_df['neuron_index'].values)); nE = len(np.unique(sub_df['rec_name'].values))
    #                             ax = axes[iB,ii]
    #                             ax.set_title(f'{stat} during {b_str}: {d} \n(n = {nN} neurons from {nE} exps)')
    #                             ax.scatter(FR_i,FR_j,s=10,marker='o',edgecolor=cmap_macro[ii],facecolor='none')

    #                             ymin = np.nanmin(FR_j); xmin = np.nanmin(FR_i)
    #                             ymax = np.nanpercentile(FR_j,98); xmax = np.nanpercentile(FR_i,98)
    #                             xymin = np.nanmin([ymin,xmin]);xymin_list.append(xymin)
    #                             xymax = np.nanmax([ymax,xmax]);xymax_list.append(xymax)

    #                         xymin = np.nanmin(xymin_list); xymax = np.nanmax(xymax_list)   
    #                         for ii, d in enumerate(drug_order):
    #                             ax = axes[iB,ii]
    #                             FR_i, FR_j, delta, mod = delta_list[ii]
                                
    #                             ax.plot([xymin,xymax],[xymin,xymax],'-k')
    #                             ax.set_xlim([xymin,xymax]); ax.set_ylim([xymin,xymax])
    #                             usrplt.adjust_spines(ax)
    #                             ax.set_xlabel(f'{block_i}')
    #                             ax.set_ylabel(f'{block_j}')
    #                             res = st.wilcoxon(FR_i,FR_j,nan_policy='omit')
    #                             if res.pvalue < pval_thresh:
    #                                 ax.text(0.5*xymax,0.9*xymax,f'*',color=cmap_macro[ii],fontweight='bold',fontsize=20)

    #                         #Compare saline vs psilocybin
    #                         distr_i = delta_list[0][-1]
    #                         distr_j = delta_list[1][-1]

    #                         FR_i, FR_j, delta, mod_i = delta_list[0]
    #                         FR_i, FR_j, delta, mod_j = delta_list[1]
    #                         ax = axes[iB,2]
    #                         sns.histplot([mod_i,mod_j],palette=cmap_macro,binwidth=0.1,multiple='layer',stat='density',common_norm=False,kde=False,ax=ax)

    #                         ax.set_title('Modulation index')
    #                         ax.set_xlabel('MI')
    #                         ax.set_xlim([-1,1])
    #                         usrplt.adjust_spines(ax)
    #                         ylim = ax.get_ylim(); xlim = ax.get_xlim()
    #                         ax.vlines(np.nanmean(mod_i),*ylim,ls='-',color=cmap_macro[0],lw=2,label=f'Saline MI: {np.nanmean(mod_i):.2f}')
    #                         ax.vlines(np.nanmean(mod_j),*ylim,ls='-',color=cmap_macro[1],lw=2,label=f'Psilocybin MI: {np.nanmean(mod_j):.2f}')
    #                         ax.legend()

    #                         if (len(mod_i) > 0) & (len(mod_j) > 0):
    #                             res = st.mannwhitneyu(mod_i,mod_j,nan_policy='omit')
    #                             if res.pvalue < pval_thresh/3:
    #                                 ax.text(.5,0.96*ax.get_ylim()[1],'*',fontsize=20,fontweight='bold',color='k')

    #                         ax = axes[iB,3]
    #                         sns.histplot([np.abs(mod_i),np.abs(mod_j)],palette=cmap_macro,binwidth=0.1,multiple='layer',stat='density',common_norm=False,kde=False,ax=ax)
    #                         ax.set_title('| Modulation index |')
    #                         ax.set_xlabel('|MI|')
    #                         ax.set_xlim([0,1])
    #                         usrplt.adjust_spines(ax)
    #                         ylim = ax.get_ylim(); xlim = ax.get_xlim()
    #                         ax.vlines(np.nanmean(np.abs(mod_i)),*ylim,ls='--',color=cmap_macro[0],lw=2,label=f'Saline |MI|: {np.nanmean(np.abs(mod_i)):.2f}')
    #                         ax.vlines(np.nanmean(np.abs(mod_j)),*ylim,ls='--',color=cmap_macro[1],lw=2,label=f'Psilocybin |MI|: {np.nanmean(np.abs(mod_j)):.2f}')
    #                         ax.legend()

    #                         if (len(mod_i) > 0) & (len(mod_j) > 0):
    #                             res = st.mannwhitneyu(np.abs(mod_i),np.abs(mod_j),nan_policy='omit')
    #                             if res.pvalue < pval_thresh:
    #                                 ax.text(.5,0.96*ax.get_ylim()[1],'*',fontsize=20,fontweight='bold',color='k')

    #                     plt.savefig(join(PlotDir,f'group_scatter_{stat}_{g}_layer-{l_str}_{ct}_{block_i}_vs_{block_j}.png'),dpi=300,bbox_inches='tight',facecolor='w')
    #                     prs = prs_list[iS]
    #                     slide = prs.slides.add_slide(prs.slide_layouts[1])
    #                     slide.shapes.title.text = f'{g}, layer {l}, {ct}, {block_i}_vs_{block_j}'
    #                     usrplt.save_fig_to_pptx(fig, prs,slide)
    #                     plt.close(fig)
    #                     nPlots+=1

    # PlotDir = join(BaseDir,'group')
    # for iS, stat in enumerate(stat_list):
    #     prs = prs_list[iS]
    #     prs.save(join(PlotDir,f'group_scatter_{stat}.pptx'))

    print(f'DONE: Generated {nPlots} plots, Good luck with the analysis!')

                    
                    
