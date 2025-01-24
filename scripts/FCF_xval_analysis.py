base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
base_dir_local = '/data/tiny-blue-dot/zap-n-zip/EEG_exp'; base_dir = base_dir_server
ProjDir = '/home/david.wyrick/projects/zap-n-zip/'
ServDir = '/data/projects/zap-n-zip/'

#Basea
import gspread
from os.path import join
from glob import glob
import json, os, time, sys, argparse
import numpy as np
import pandas as pd
import scipy.stats as st
import pingouin as pg
from scipy.ndimage import gaussian_filter
from scipy.interpolate import interp1d
from itertools import combinations

#Plot
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
from matplotlib.gridspec import GridSpec

#Project
from tbd_eeg.data_analysis.eegutils import EEGexp
import tbd_eeg.data_analysis.Utilities.utilities as tbd_util
from tbd_eeg.data_analysis.Utilities.behavior_movies import Movie

#Allen
from allensdk.core.mouse_connectivity_cache import MouseConnectivityCache
mcc = MouseConnectivityCache(resolution=10)
str_tree = mcc.get_structure_tree()

#Read in allen CCF
ccfsum = pd.read_csv('/home/david.wyrick/projects/zap-n-zip/data/mouse_connectivity/ontology_v2.csv')

#User
sys.path.append(ProjDir)
import util
import plotting as usrplt

gc = gspread.service_account() # need a key file to access the account (step 2) 
sh = gc.open('Templeton-log_exp') # open the spreadsheet 

df = pd.DataFrame(sh.sheet1.get()) 
exp_table = df.T.set_index(0).T # put it in a nicely formatted dataframeexp_table.iloc[10:]
exp_table = exp_table.iloc[9:-8].set_index('mouse_name')

##===== ============================ =====##
##===== Parse Command Line Arguments =====##
parser = argparse.ArgumentParser(description='FCF cross-validation analysis')

##===== Data Options =====##
parser.add_argument('--mouseID',type=str, default='mouse669117',
                    help='mouse to perform analysis on')

parser.add_argument('--rec_name',type=str, default='pilot_aw_psi_2023-03-30_11-37-07',
                    help='experiment to perform analysis on')

parser.add_argument('--last_run',type=int, default=None,
                    help='FCF run')

def calculate_average(mat,filename_list,b,cmap,axis=0,ax=None):
    tmp_list = []
    for jj, fname in enumerate(filename_list):
        tmp = np.nanmean(mat[jj],axis=axis)
        tmp_list.append(tmp)
    mat_j = np.array(tmp_list)
    mat_bluepill = np.mean(mat_j[:b],axis=0)
    mat_redpill = np.mean(mat_j[b:],axis=0)
    if ax is not None:
        sns.kdeplot(mat_bluepill,ax=ax,color=cmap[b-1],ls='-',label='Non-psi average',lw=2)
        sns.kdeplot(mat_redpill,ax=ax,color=cmap[-1],ls='-',label='Psi average',lw=2)
    return mat_bluepill, mat_redpill

if __name__ == '__main__':

    # Parse the arguments
    args = parser.parse_args()
    mouse_name = args.mouseID 
    mID = args.mouseID
    rec_name = args.rec_name
    last_run = args.last_run
    
   #Find folder with results inside
    TempDir = os.path.join(ServDir,'results','ccm_xval',mID)
    if last_run is None:
        tmp_list = sorted(glob(join(TempDir,f'{rec_name}_run_*')))
        last_run = int(tmp_list[-1][-1])
    folder = f'{rec_name}_run_{last_run:02d}'

    #Define save and plot directories
    SaveDir = os.path.join(TempDir,folder)
    PlotDir = os.path.join(SaveDir,'plots')
    if not os.path.exists(PlotDir):
        os.makedirs(PlotDir)

    #Load ccm parameters json
    fpath = glob(join(SaveDir,'ccm_parameters*.json'))[0]
    with open(fpath, 'r') as f:
        params = json.loads(f.read())

    time_bin_ms = params['time_bin_ms']
    time_bin = time_bin_ms/1000
    delay = params['delay']
    rand_proj = bool(params['rand_proj'])
    xval = bool(params['xval'])
    zscore = bool(params['zscore'])
    nKfold = params['nKfold']
    fr_thresh = params['fr_thresh'] 

    #Get time windows used for CCM
    data = np.load(os.path.join(SaveDir,f'time_windows.npz'))
    time_window_list = data['time_window_list'].tolist()
    filename_list = data['filename_list'].tolist()
    neuron_indices = data['neuron_indices']
    nNeurons = len(neuron_indices)

    #Create colormap for epochs
    nPsi = np.sum(['psi' in fname for fname in filename_list])
    nSal = len(filename_list) - nPsi
    cmap = np.concatenate((sns.color_palette('Blues',nSal),sns.color_palette('Reds',nPsi)))
    print(f'Time bin = {time_bin_ms}ms, z-score = {zscore}, x-val = {xval}, delay = {delay}, rand_proj = {rand_proj}, fr_thresh = {fr_thresh}Hz, nNeurons = {nNeurons}')
    
    #%% Upload the whole experiment and generate the global clock
    file_name = os.path.join(base_dir,mouse_name,rec_name,'experiment1','recording1')
    exp = EEGexp(file_name, preprocess=False, make_stim_csv=False) 

    #Load probe data
    probe_unit_data, probe_info, total_units = tbd_util.get_neuropixel_data(exp)
    probe_list = [x.replace('_sorted', '') for x in exp.experiment_data if 'probe' in x]

    #Bin spiking data
    data_list, ts_list, _, plot_tuple = util.bin_spiking_data(probe_unit_data, time_window_list, time_bin=time_bin,fr_thresh=fr_thresh)
    boundaries_group, ticks_group, labels_group, celltypes, durations, areas, groups, supergroups, order_by_group = plot_tuple
    boundaries_probe, ticks_probe, labels_probe = util.get_probe_plotting_info(probe_unit_data,neuron_indices)
    boundaries = boundaries_probe; ticks = ticks_probe; labels = labels_probe
    
    #Get full names
    area_fullnames = []
    for a in areas:
        if a in ['null','nan']:
            area_fullnames.append(a)
        else:
            area_fullnames.append(str_tree.get_structures_by_acronym([a])[0]['name'])
            
    #Get running signal
    run_file = join(base_dir,mouse_name,rec_name,'experiment1','recording1','running_signal.npy')
    run_signal = np.load(run_file)

    ts_file = join(base_dir,mouse_name,rec_name,'experiment1','recording1','running_timestamps.npy')
    if os.path.exists(ts_file):
        run_timestamps = np.load(ts_file)
    else:
        ts_file = join(base_dir,mouse_name,rec_name,'experiment1','recording1','running_timestamps_master_clock.npy')
        run_timestamps = np.load(ts_file)

    #Get pupil diameter
    base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
    pupil_csv = join(base_dir_server,mID,rec_name,'experiment1','recording1',f'Pupil_{rec_name}.csv')

    try:
        table = pd.read_csv(pupil_csv)
        pupil_radius = table['Largest_Radius'].values

        #Pupil master clock
        pupil_ts = Movie(filepath=exp.pupilmovie_file,
                        sync_filepath=exp.sync_file,
                        sync_channel='eyetracking'
                        ).sync_timestamps
        pupil_data = True
    except:
        pupil_ts = np.array([np.nan])
        pupil_radius =np.array([np.nan])
        pupil_data = False

    fr_list = []
    for X in data_list:
        fr_list.append(np.sum(X,axis=0)/(X.shape[0]*time_bin))
    FR_perblock = np.array(fr_list).T
    FR_perblock_sub = FR_perblock[neuron_indices]
    
    #Load FCF!!
    fsuffix = f'{rec_name}_{time_bin_ms}ms-bins_tau-{delay}'
    if zscore:
        fsuffix = f'{fsuffix}_z' 
    if xval:
        fsuffix = f'{fsuffix}_xval'
    if rand_proj:
        fsuffix = f'{fsuffix}_randproj'

    #FCF
    results = np.load(os.path.join(SaveDir,f'results_{fsuffix}_corrected.npz'))
    FCF_optimal = results['FCF'][:,:,:nNeurons,:nNeurons]
    complexity = results['complexity'][:,:,:nNeurons,:nNeurons]
    directionality = results['directionality'][:,:,:nNeurons,:nNeurons]
    correlation = results['correlation'][:,:,:nNeurons,:nNeurons]
    FCF_takens = results['FCF_takens'][:,:,:,:nNeurons,:nNeurons]
    nCond, nTakens, nKfold, N, _ = FCF_takens.shape

    #Take mean over kfolds
    FCF_optimal_mkf = np.nanmean(FCF_optimal,axis=1)
    complexity_mkf = np.nanmean(complexity,axis=1)
    directionality_mkf = np.nanmean(directionality,axis=1)
    correlation_mkf = np.nanmean(correlation,axis=1)

    #Calculate masked FCF optimal for mean kfold
    FCF_optimal_mask = FCF_optimal_mkf.copy()
    complexity_mask = complexity_mkf.copy()
    directionality_mask = directionality_mkf.copy()
    correlation_mask = correlation_mkf.copy()

    for ii, epoch in enumerate(filename_list):
        mask = FCF_optimal_mkf[ii] < 0.1
        FCF_optimal_mask[ii][mask] = np.nan
        complexity_mask[ii][mask] = np.nan
        directionality_mask[ii][mask] = np.nan
        correlation_mask[ii][mask] = np.nan

    #Calculate masked FCF optimal for each kfold
    FCF_optimal_kf_mask = FCF_optimal.copy()
    complexity_kf_mask = complexity.copy()
    directionality_kf_mask = directionality.copy()
    correlation_kf_mask = correlation.copy()

    for ii, epoch in enumerate(filename_list):
        for iK in range(nKfold):
            mask = FCF_optimal[ii,iK] < 0.1
            FCF_optimal_kf_mask[ii,iK][mask] = np.nan
            complexity_kf_mask[ii,iK][mask] = np.nan
            directionality_kf_mask[ii,iK][mask] = np.nan
            correlation_kf_mask[ii,iK][mask] = np.nan

    #Make arrays for reorder
    FCF_optimal_mkf_ro = np.zeros(FCF_optimal_mkf.shape)
    complexity_mkf_ro = np.zeros(complexity_mkf.shape)
    directionality_mkf_ro = np.zeros(directionality_mkf.shape)
    correlation_mkf_ro = np.zeros(correlation_mkf.shape)

    #Reorder by group
    for ii in range(len(filename_list)):
        FCF_optimal_mkf_ro[ii] = FCF_optimal_mkf[ii][:,order_by_group][order_by_group]
        complexity_mkf_ro[ii] = complexity_mkf[ii][:,order_by_group][order_by_group]
        directionality_mkf_ro[ii] = directionality_mkf[ii][:,order_by_group][order_by_group]
        correlation_mkf_ro[ii] = correlation_mkf[ii][:,order_by_group][order_by_group]

    # Create a new PowerPoint presentation to save figures to
    prs = Presentation()

    #Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f'mID = {mID}; rec_name = {rec_name}'
    slide.placeholders[1].text = f'FC analysis with cross-validation'

    nPairs = []
    perc_FCF = []
    for ii, epoch in enumerate(filename_list):
        tmp = FCF_optimal_mkf[ii].copy()
        mask = FCF_optimal_mkf[ii] > 0.1
        perc_FCF.append(np.nanpercentile(tmp[mask].ravel(),75))
        nPairs.append(np.sum(mask)/(N**2))

    fig, ax = plt.subplots(figsize=(8,5))
    ax.plot(nPairs,'-k')
    for jj, fname in enumerate(filename_list):
        ax.plot(jj,nPairs[jj],'o',color=cmap[jj])
    ax.set_xticks(np.arange(len(filename_list)))
    ax.set_xticklabels(filename_list,rotation=30)
    ax.set_ylabel('Fraction of pairs with FC > 0.1')
    ax.set_xlabel('Epoch')

    ax2 = ax.twinx()
    ax2.plot(perc_FCF,'-',color=usrplt.cc[8])
    for jj, fname in enumerate(filename_list):
        ax2.plot(jj,perc_FCF[jj],'o',color=cmap[jj])
    ax2.set_ylabel('75th percentile FCF of significant pairs',color=usrplt.cc[8])
    usrplt.save_fig_to_pptx(fig, prs)
    plt.savefig(join(PlotDir,f'FCF_sigpairs_{rec_name}.pdf'))

    # Calculate running per kfold
    from sklearn.model_selection import KFold

    fig, axes = plt.subplots(1,2,figsize=(10,5))
    plt.suptitle(f'{mID}: {rec_name}')
    running_moments = np.zeros((len(data_list),nKfold,3))
    pupil_moments = np.zeros((len(data_list),nKfold,3))
    for ii, data in enumerate(data_list):
        T, N = data.shape
        ts = ts_list[ii]

        k_fold = KFold(n_splits=nKfold)
        run_list = []
        pup_list = []
        for iK, (cue_indices, tar_indices) in enumerate(k_fold.split(np.arange(T))):
            #Get running during target periods
            ts_window = ts[tar_indices]
            indy = np.where((run_timestamps >= ts_window[0]) & (run_timestamps <= ts_window[-1]))[0]
            run_tmp = run_signal[indy]

            running_moments[ii,iK,0] = np.nanmean(run_tmp)
            running_moments[ii,iK,1] = np.nanstd(run_tmp)
            running_moments[ii,iK,2] = st.skew(run_tmp)

            indy = np.where((pupil_ts >= ts_window[0]) & (pupil_ts <= ts_window[-1]))[0]
            pupil_tmp = pupil_radius[indy]

            pupil_moments[ii,iK,0] = np.nanmean(pupil_tmp)
            pupil_moments[ii,iK,1] = np.nanstd(pupil_tmp)
            pupil_moments[ii,iK,2] = st.skew(pupil_tmp)
            # sns.kdeplot(run_tmp,ax=axes[0],color=cmap[ii],lw=0.75)
            # sns.kdeplot(pupil_tmp,ax=axes[1],color=cmap[ii],lw=0.75)
            run_list.append(run_tmp)
            pup_list.append(pupil_tmp)
        sns.kdeplot(np.concatenate(run_list),ax=axes[0],color=cmap[ii],lw=1.5)
        sns.kdeplot(np.concatenate(pup_list),ax=axes[1],color=cmap[ii],lw=1.5,label=filename_list[ii])
    axes[0].set_title('Locomotion')
    axes[0].set_xlabel('Running speed (cm/s)')
            
    axes[1].set_title('Pupil size')
    axes[1].set_xlabel('Pupil radius (pixels)')
    axes[1].legend()
    usrplt.save_fig_to_pptx(fig, prs)
    plt.savefig(join(PlotDir,f'behavior_distributions_{rec_name}.pdf'))

    fcf_moments = np.zeros((len(data_list),nKfold,5))
    sig_elements = np.zeros((len(data_list),nKfold))
    #Calculate mean and percentile of FCF per epoch/kfold
    for ii, data in enumerate(data_list):
        for iK in range(nKfold):
            fcf_moments[ii,iK,0] = np.nanmean(FCF_optimal[ii,iK])
            fcf_moments[ii,iK,1] = np.nanpercentile(FCF_optimal[ii,iK],75)

            tmp = np.copy(FCF_optimal[ii,iK]).ravel()
            mask = tmp > 0.1
            sig_elements[ii,iK] = np.sum(mask)/len(mask)
            fcf_moments[ii,iK,2] = np.sum(mask)/len(mask) 
            fcf_moments[ii,iK,3] = np.nanmean(tmp[mask])
            fcf_moments[ii,iK,4] = np.nanpercentile(tmp[mask],75)
    markers = ['o','v','*','s','p']

    for x_array, xlabel in zip((pupil_moments,running_moments),('pupil radius (pixels)','running speed (cm/s)')):
        for jj, ylabel in enumerate(['Mean FCF','75th percentile FCF']):
            fig, axes = plt.subplots(1,2,figsize=(10,5))
            plt.suptitle(f'{mID}, {rec_name}; Trends in FCF')
            for ii, data in enumerate(data_list):
                for iK in range(nKfold):
                    axes[0].scatter(x_array[ii,iK,0],fcf_moments[ii,iK,jj],marker='.',facecolors="None",edgecolors=cmap[ii])#,linewidths=2)#,label=fname)
                    axes[1].scatter(x_array[ii,iK,1],fcf_moments[ii,iK,jj],marker='.',facecolors="None",edgecolors=cmap[ii])#,linewidths=2)#,label=fname)

                axes[0].plot(np.nanmean(x_array[ii,:,0]),np.nanmean(fcf_moments[ii,:,jj]),'o',color=cmap[ii])#,label=fname)
                axes[1].plot(np.nanmean(x_array[ii,:,1]),np.nanmean(fcf_moments[ii,:,jj]),'o',color=cmap[ii])#,label=fname)

            x = x_array[:,:,0].ravel()
            y = fcf_moments[:,:,jj].ravel()
            r = st.linregress(x,y)

            x_plot = np.linspace(np.min(x)-1,np.max(x)+1)
            y_plot = r.slope*x_plot+r.intercept
            r2 = r.rvalue**2
            axes[0].plot(x_plot,y_plot,'-k',label=f'r2 = {r2:.3f}',zorder=0)
            
            # p, ss_res, rank,svals,rcond = np.polyfit(x, y, 2,full=True)
            # y_plot = p[0]*x_plot**2+p[1]*x_plot+p[2]
            # ss_tot = np.sum((y - np.nanmean(y))**2)
            # r2 = 1 - ss_res[0]/ss_tot
            # axes[0].plot(x_plot,y_plot,'--k',label=f'r2 = {r2:.3f}',zorder=0)
            axes[0].legend()
            x = x_array[:,:,1].ravel()
            y = fcf_moments[:,:,jj].ravel()

            r = st.linregress(x,y)
            x_plot = np.linspace(np.min(x)-1,np.max(x)+1)
            y_plot = r.slope*x_plot+r.intercept
            r2 = r.rvalue**2
            axes[1].plot(x_plot,y_plot,'-k',label=f'r2 = {r2:.3f}',zorder=0)
        
            # p, ss_res, rank,svals,rcond = np.polyfit(x, y, 2,full=True)
            # y_plot = p[0]*x_plot**2+p[1]*x_plot+p[2]
            # ss_tot = np.sum((y - np.nanmean(y))**2)
            # r2 = 1 - ss_res[0]/ss_tot
            # axes[1].plot(x_plot,y_plot,'--k',label=f'r2 = {r2:.3f}',zorder=0)
            axes[1].legend()
            ax = axes[0]
            ax.set_xlabel(f'Mean {xlabel}')
            axes[1].set_xlabel(f'Std. dev. {xlabel}')
            ax.set_ylabel(ylabel)

            usrplt.save_fig_to_pptx(fig, prs)
            plt.savefig(join(PlotDir,f'behavior-v-FCF_{jj}_{rec_name}.pdf'))

    markers = ['o','v','*','s','p']
    jj = 1
    for x_array, xlabel in zip((pupil_moments,running_moments),('pupil radius (pixels)','running speed (cm/s)')):

        fig, axes = plt.subplots(1,2,figsize=(10,5))
        plt.suptitle(f'{mID}, {rec_name}; Trends in FCF')
        for ii, data in enumerate(data_list):
            for iK in range(nKfold):
                axes[0].scatter(x_array[ii,iK,0],fcf_moments[ii,iK,2],marker='.',facecolors="None",edgecolors=cmap[ii])#,linewidths=2)#,label=fname)
                axes[1].scatter(x_array[ii,iK,0],fcf_moments[ii,iK,4],marker='.',facecolors="None",edgecolors=cmap[ii])#,linewidths=2)#,label=fname)

            axes[0].plot(np.nanmean(x_array[ii,:,0]),np.nanmean(fcf_moments[ii,:,2]),'o',color=cmap[ii])#,label=fname)
            axes[1].plot(np.nanmean(x_array[ii,:,0]),np.nanmean(fcf_moments[ii,:,4]),'o',color=cmap[ii])#,label=fname)

        x = x_array[:,:,0].ravel()
        y = fcf_moments[:,:,2].ravel()
        r = st.linregress(x,y)
        p, ss_res, rank,svals,rcond = np.polyfit(x, y, 2,full=True)

        x_plot = np.linspace(np.min(x)-1,np.max(x)+1)
        # y_plot = p[0]*x_plot**2+p[1]*x_plot+p[2]
        # ss_tot = np.sum((y - np.nanmean(y))**2)
        # r2 = 1 - ss_res[0]/ss_tot
        # axes[0].plot(x_plot,y_plot,'--k',label=f'r2 = {r2:.3f}',zorder=0)

        y_plot = r.slope*x_plot+r.intercept
        r2 = r.rvalue**2
        axes[0].plot(x_plot,y_plot,'-k',label=f'r2 = {r2:.3f}',zorder=0)
        axes[0].legend()

        x = x_array[:,:,0].ravel()
        y = fcf_moments[:,:,4].ravel()
        r = st.linregress(x,y)
        p, ss_res, rank,svals,rcond = np.polyfit(x, y, 2,full=True)

        x_plot = np.linspace(np.min(x)-1,np.max(x)+1)
        # y_plot = p[0]*x_plot**2+p[1]*x_plot+p[2]
        # ss_tot = np.sum((y - np.nanmean(y))**2)
        # r2 = 1 - ss_res[0]/ss_tot
        # axes[1].plot(x_plot,y_plot,'--k',label=f'r2 = {r2:.3f}',zorder=0)

        y_plot = r.slope*x_plot+r.intercept
        r2 = r.rvalue**2
        axes[1].plot(x_plot,y_plot,'-k',label=f'r2 = {r2:.3f}',zorder=0)
        axes[1].legend()

        ax = axes[0]
        ax.set_xlabel(f'Mean {xlabel}')
        ax.set_ylabel('Fraction of significant FC pairs (>0.1)')

        ax = axes[1]
        ax.set_xlabel(f'Mean {xlabel}')
        ax.set_ylabel('75th percentile FCF of significant pairs')

        usrplt.save_fig_to_pptx(fig, prs)
        plt.savefig(join(PlotDir,f'behavior-v-sig-FCF_{rec_name}.pdf'))
    
    np.savez(os.path.join(SaveDir,f'FCFvsBehavior_{rec_name}.npz'),pupil_moments=pupil_moments,running_moments=running_moments,fcf_moments=fcf_moments,cmap=cmap,filename_list=filename_list,)
    
    # vmax1 = np.round(np.nanpercentile(np.abs(correlation.ravel()),99),2)
    # vmax2 = np.round(np.nanpercentile(np.abs(FCF_optimal.ravel()),99),2)
    # vmax_fcf = np.round(np.max([vmax1,vmax2]),1)

    # for ii, epoch in enumerate(filename_list):
    #     fig, axes = plt.subplots(2,5,figsize=(20,8))
    #     plt.suptitle(epoch,fontsize=18,fontweight='bold')
    #     for iK in range(nKfold):
    #         usrplt.visualize_matrix(FCF_optimal[ii,iK],ax=axes[iK//5,iK%5],cmap='viridis',title=f'kfold {iK}',clims=[0,vmax_fcf],center=None,ticks=ticks,labels=labels,boundaries=boundaries)
    #     usrplt.save_fig_to_pptx(fig, prs)

    # for ii, epoch in enumerate(filename_list):
    #     fig, axes = plt.subplots(1,4,figsize=(20,5))
    #     plt.suptitle(epoch,fontsize=18,fontweight='bold')
    #     #Plot mean correlation across kfolds
    #     usrplt.visualize_matrix(np.nanmean(correlation[ii],axis=0),ax=axes[0],cmap='viridis',title='Mean Correlation',clims=[0,vmax_fcf],center=None,ticks=ticks,labels=labels,boundaries=boundaries)

    #     #Plot std. dev. correlation across kfolds
    #     usrplt.visualize_matrix(np.nanstd(correlation[ii],axis=0),ax=axes[1],cmap='rocket',title='Std. Dev Correlation',center=None,ticks=ticks,labels=labels,boundaries=boundaries)

    #     #Plot mean FCF across kfolds
    #     usrplt.visualize_matrix(np.nanmean(FCF_optimal[ii],axis=0),ax=axes[2],cmap='viridis',title='Mean FC',clims=[0,vmax_fcf],center=None,ticks=ticks,labels=labels,boundaries=boundaries)

    #     #Plot mean FCF across kfolds
    #     usrplt.visualize_matrix(np.nanstd(FCF_optimal[ii],axis=0),ax=axes[3],cmap='rocket',title='Std. Dev. FC',center=None,ticks=ticks,labels=labels,boundaries=boundaries)
        
    #     usrplt.save_fig_to_pptx(fig, prs)
        
    # fig, ax = plt.subplots(figsize=(8,8))
    # plt.suptitle(f'{mID}; all FCF pairs')
    # fcf_j_sal = []
    # fcf_j_psi = []

    # for ii, epoch in enumerate(filename_list):
    #     for iK in range(nKfold):
    #         lw = 1.5; zo = 0; ls = '-'
    #         fcf_tmp = FCF_optimal[ii,iK]#  np.arctanh(FCF_optimal[ii,iK])# - FCF_takens[jj,0]
    #         fcf_j = np.nanmean(fcf_tmp,axis=0)
    #         if ii < nSal:
    #             fcf_j_sal.append(fcf_j)
    #         else:
    #             fcf_j_psi.append(fcf_j)
    #         sns.kdeplot(fcf_j,ax=ax,color=cmap[ii],ls=ls,label=epoch,lw=lw,zorder=zo)

    # ax_in = ax.inset_axes([0.7, 0.7, 0.3, 0.3])
    # mat_bluepill = np.nanmean(fcf_j_sal,axis=0)
    # mat_redpill = np.nanmean(fcf_j_psi,axis=0)
    # sns.kdeplot(mat_bluepill,ax=ax_in,color=cmap[nSal-1],ls='-',label='Non-psi average',lw=2)
    # sns.kdeplot(mat_redpill,ax=ax_in,color=cmap[-1],ls='-',label='Psi average',lw=2)

    # ax.set_xlabel('FCF_j')
    # ax.set_title('Functional causal flow')
    # ax.set_xlim([0,0.25])
    # ax_in.set_xlim([0,0.25])
    # usrplt.save_fig_to_pptx(fig, prs)
    # plt.savefig(join(PlotDir,f'FCF_j_kfolds_{rec_name}.pdf'))

    # fig, ax = plt.subplots(figsize=(8,8))
    # plt.suptitle(f'{mID}; masked FCF')
    # fcf_j_sal = []
    # fcf_j_psi = []

    # for ii, epoch in enumerate(filename_list):
    #     for iK in range(nKfold):
    #         lw = 1.5; zo = 0; ls = '-'
    #         fcf_tmp = FCF_optimal[ii,iK].copy()#  np.arctanh(FCF_optimal[ii,iK])# - FCF_takens[jj,0]
    #         mask = fcf_tmp < 0.1
    #         fcf_tmp[mask] = np.nan
    #         fcf_j = np.nanmean(fcf_tmp,axis=0)
    #         if ii < nSal:
    #             fcf_j_sal.append(fcf_j)
    #         else:
    #             fcf_j_psi.append(fcf_j)
    #         sns.kdeplot(fcf_j,ax=ax,color=cmap[ii],ls=ls,label=epoch,lw=lw,zorder=zo)

    # ax_in = ax.inset_axes([0.7, 0.7, 0.3, 0.3])
    # mat_bluepill = np.nanmean(fcf_j_sal,axis=0)
    # mat_redpill = np.nanmean(fcf_j_psi,axis=0)
    # sns.kdeplot(mat_bluepill,ax=ax_in,color=cmap[nSal-1],ls='-',label='Non-psi average',lw=2)
    # sns.kdeplot(mat_redpill,ax=ax_in,color=cmap[-1],ls='-',label='Psi average',lw=2)

    # ax.set_xlabel('FCF_j')
    # ax.set_title('Functional causal flow')
    # ax.set_xlim([0,0.45])
    # ax_in.set_xlim([0,0.45])
    # usrplt.save_fig_to_pptx(fig, prs)
    # plt.savefig(join(PlotDir,f'FCF_j_masked-kfolds_{rec_name}.pdf'))
    
    # fig, axes = plt.subplots(1,4,figsize=(16,4))
    # plt.suptitle(f'{mID}, mean over kfolds')
    # for ii, epoch in enumerate(filename_list):

    #     lw = 1.5; zo = 0; ls = '-'
    #     fcf_tmp = np.arctanh(FCF_optimal_mkf[ii])# - FCF_takens[jj,0]
    #     tmp = np.nanmean(fcf_tmp,axis=0)
    #     sns.kdeplot(tmp,ax=axes[0],color=cmap[ii],ls=ls,label=epoch,lw=lw,zorder=zo)

    #     dir_tmp = fcf_tmp - fcf_tmp.T
    #     tmp = np.nanmean(dir_tmp,axis=0)
    #     sns.kdeplot(tmp,ax=axes[1],color=cmap[ii],ls=ls,label=epoch,lw=lw,zorder=zo)

    #     tmp = np.nanmean(complexity_mkf[ii],axis=0)
    #     sns.kdeplot(tmp,ax=axes[2],color=cmap[ii],ls=ls,label=epoch,lw=lw,zorder=zo)

    #     tmp = np.nanmean(complexity_mkf[ii],axis=1)
    #     sns.kdeplot(tmp,ax=axes[3],color=cmap[ii],ls=ls,label=epoch,lw=lw,zorder=zo)

    # ax_in = axes[0].inset_axes([0.7, 0.7, 0.3, 0.3])
    # fcf_j_bluepill, fcf_j_redpill = calculate_average(np.arctanh(FCF_optimal_mkf),filename_list,nSal,cmap,ax=ax_in)
    # # ax_in.set_xlabel('FCF_j')
    # ax_in.set_ylabel('')

    # ax = axes[0]
    # ax.set_xlabel('FCF_j')
    # ax.set_title('Functional causal flow')
    # ax.set_xlim([0,0.25])
    # ax_in.set_xlim([0,0.25])

    # ax_in = axes[1].inset_axes([0.7, 0.7, 0.3, 0.3])
    # dir_j_bluepill, dir_j_redpill = calculate_average(directionality_mkf,filename_list,nSal,cmap,ax=ax_in)
    # # ax_in.set_xlabel('Directionality_j')
    # ax_in.set_xlim([-0.1,0.1])
    # ax_in.set_ylabel('')

    # ax = axes[1]
    # ax.set_xlabel('Directionality_j')
    # ax.set_title('Directionality')
    # ax.set_xlim([-0.1,0.1])
    # # ax.legend()

    # ax_in = axes[2].inset_axes([0.7, 0.7, 0.3, 0.3])
    # max_takens = np.round(np.nanpercentile(complexity_mkf.ravel(),97.5))
    # cpx_j_bluepill, cpx_j_redpill = calculate_average(complexity_mkf,filename_list,nSal,cmap,ax=ax_in)
    # # ax_in.set_xlabel('Takens_j')
    # ax_in.set_ylabel('')

    # ax = axes[2]
    # ax.set_xlabel('Takens_j')
    # ax.set_title('Complexity\nof source')
    # ax.set_xlim([0,max_takens]); ax_in.set_xlim([0,max_takens])

    # ax_in = axes[3].inset_axes([0.7, 0.7, 0.3, 0.3])
    # cpx_i_bluepill, cpx_i_redpill = calculate_average(complexity_mkf,filename_list,nSal,cmap,axis=1,ax=ax_in)
    # # # ax_in.set_xlabel('Takens_i')
    # # ax_in.set_ylabel('')
    # ax = axes[3]
    # ax.set_xlabel('Takens_i')
    # ax.set_title('Complexity\nof target')
    # ax.set_xlim([0,max_takens]); ax_in.set_xlim([0,max_takens])
    # # ax.legend()

    # for ax in axes:
    #     xlim = ax.get_xlim()
    #     ylim = ax.get_ylim()
    #     usrplt.adjust_spines(ax, ['left', 'bottom'])
    # usrplt.save_fig_to_pptx(fig, prs)
    # plt.savefig(join(PlotDir,f'FCF_j_mean-kfold_{rec_name}.pdf'))


    # #Plot similar plot but all kfolds
    # fcf_j_sal = []; fcf_j_psi = [];dir_j_sal = []; dir_j_psi = []
    # cpx_i_sal = []; cpx_i_psi = [];cpx_j_sal = []; cpx_j_psi = [];
    # fig, axes = plt.subplots(1,4,figsize=(16,4))
    # plt.suptitle(f'{mID}: All kfolds plotted')
    # for ii, epoch in enumerate(filename_list):
    #     for iK in range(nKfold):
    #         lw = 1.5; zo = 0; ls = '-'
    #         fcf_tmp = FCF_optimal_kf_mask[ii,iK] #np.arctanh(FCF_optimal_kf_mask[ii,iK])# - FCF_takens[jj,0]
    #         fcf_j = np.nanmean(fcf_tmp,axis=0)
    #         sns.kdeplot(fcf_j,ax=axes[0],color=cmap[ii],ls=ls,label=epoch,lw=lw,zorder=zo)

    #         dir_tmp = fcf_tmp - fcf_tmp.T
    #         dir_j = np.nanmean(dir_tmp,axis=0)
    #         sns.kdeplot(dir_j,ax=axes[1],color=cmap[ii],ls=ls,label=epoch,lw=lw,zorder=zo)

    #         cpx_j = np.nanmean(complexity_kf_mask[ii,iK],axis=0)
    #         sns.kdeplot(cpx_j,ax=axes[2],color=cmap[ii],ls=ls,label=epoch,lw=lw,zorder=zo)

    #         cpx_i = np.nanmean(complexity_kf_mask[ii,iK],axis=1)
    #         sns.kdeplot(cpx_i,ax=axes[3],color=cmap[ii],ls=ls,label=epoch,lw=lw,zorder=zo)

    #         if ii < nSal:
    #             fcf_j_sal.append(fcf_j)
    #             dir_j_sal.append(dir_j)
    #             cpx_i_sal.append(cpx_i)
    #             cpx_j_sal.append(cpx_j)
    #         else:
    #             fcf_j_psi.append(fcf_j)
    #             dir_j_psi.append(dir_j)
    #             cpx_i_psi.append(cpx_i)
    #             cpx_j_psi.append(cpx_j)

    # ax_in = axes[0].inset_axes([0.7, 0.7, 0.3, 0.3])
    # ax_in.set_ylabel('')
    # mat_bluepill = np.nanmean(fcf_j_sal,axis=0)
    # mat_redpill = np.nanmean(fcf_j_psi,axis=0)
    # sns.kdeplot(mat_bluepill,ax=ax_in,color=cmap[nSal-1],ls='-',label='Non-psi average',lw=2)
    # sns.kdeplot(mat_redpill,ax=ax_in,color=cmap[-1],ls='-',label='Psi average',lw=2)

    # ax = axes[0]
    # ax.set_xlabel('FCF_j')
    # ax.set_title('Functional causal flow')
    # ax.set_xlim([0,0.5])
    # ax_in.set_xlim([0,0.5])

    # ax_in = axes[1].inset_axes([0.7, 0.7, 0.3, 0.3])
    # mat_bluepill = np.nanmean(dir_j_sal,axis=0)
    # mat_redpill = np.nanmean(dir_j_psi,axis=0)
    # sns.kdeplot(mat_bluepill,ax=ax_in,color=cmap[nSal-1],ls='-',label='Non-psi average',lw=2)
    # sns.kdeplot(mat_redpill,ax=ax_in,color=cmap[-1],ls='-',label='Psi average',lw=2)
    # ax_in.set_xlim([-0.1,0.1])
    # ax_in.set_ylabel('')

    # ax = axes[1]
    # ax.set_xlabel('Directionality_j')
    # ax.set_title('Directionality')
    # ax.set_xlim([-0.1,0.1])
    # # ax.legend()

    # ax_in = axes[2].inset_axes([0.7, 0.7, 0.3, 0.3])
    # max_takens = np.round(np.nanpercentile(complexity_kf_mask.ravel(),97.5))
    # mat_bluepill = np.nanmean(cpx_j_sal,axis=0)
    # mat_redpill = np.nanmean(cpx_j_psi,axis=0)
    # sns.kdeplot(mat_bluepill,ax=ax_in,color=cmap[nSal-1],ls='-',label='Non-psi average',lw=2)
    # sns.kdeplot(mat_redpill,ax=ax_in,color=cmap[-1],ls='-',label='Psi average',lw=2)
    # # ax_in.set_xlabel('Takens_j')
    # ax_in.set_ylabel('')

    # ax = axes[2]
    # ax.set_xlabel('Takens_j')
    # ax.set_title('Complexity\nof source')
    # ax.set_xlim([0,max_takens]); ax_in.set_xlim([0,max_takens])

    # ax_in = axes[3].inset_axes([0.7, 0.7, 0.3, 0.3])
    # mat_bluepill = np.nanmean(cpx_i_sal,axis=0)
    # mat_redpill = np.nanmean(cpx_i_psi,axis=0)
    # sns.kdeplot(mat_bluepill,ax=ax_in,color=cmap[nSal-1],ls='-',label='Non-psi average',lw=2)
    # sns.kdeplot(mat_redpill,ax=ax_in,color=cmap[-1],ls='-',label='Psi average',lw=2)
    # # ax_in.set_ylabel('')
    # ax = axes[3]
    # ax.set_xlabel('Takens_i')
    # ax.set_title('Complexity\nof target')
    # ax.set_xlim([0,max_takens]); ax_in.set_xlim([0,max_takens])
    # # ax.legend()

    # for ax in axes:
    #     xlim = ax.get_xlim()
    #     ylim = ax.get_ylim()
    #     usrplt.adjust_spines(ax, ['left', 'bottom'])
    # usrplt.save_fig_to_pptx(fig, prs)
    # plt.savefig(join(PlotDir,f'FCF_j_per-kfold_{rec_name}.pdf'))
        
    # ## Plot similar plot but for masked mean fkfold
    # fig, axes = plt.subplots(1,4,figsize=(16,4))
    # plt.suptitle(f'{mID}: values > 0.1; mean over kfolds')
    # for ii, epoch in enumerate(filename_list):

    #     lw = 1.5; zo = 0; ls = '-'
    #     fcf_tmp = np.arctanh(FCF_optimal_mask[ii])# - FCF_takens[jj,0]
    #     tmp = np.nanmean(fcf_tmp,axis=0)
    #     sns.kdeplot(tmp,ax=axes[0],color=cmap[ii],ls=ls,label=epoch,lw=lw,zorder=zo)

    #     dir_tmp = fcf_tmp - fcf_tmp.T
    #     tmp = np.nanmean(dir_tmp,axis=0)
    #     sns.kdeplot(tmp,ax=axes[1],color=cmap[ii],ls=ls,label=epoch,lw=lw,zorder=zo)

    #     tmp = np.nanmean(complexity_mask[ii],axis=0)
    #     sns.kdeplot(tmp,ax=axes[2],color=cmap[ii],ls=ls,label=epoch,lw=lw,zorder=zo)

    #     tmp = np.nanmean(complexity_mask[ii],axis=1)
    #     sns.kdeplot(tmp,ax=axes[3],color=cmap[ii],ls=ls,label=epoch,lw=lw,zorder=zo)

    # ax_in = axes[0].inset_axes([0.7, 0.7, 0.3, 0.3])
    # fcf_j_bluepill, fcf_j_redpill = calculate_average(np.arctanh(FCF_optimal_mask),filename_list,nSal,cmap=cmap,ax=ax_in)
    # # ax_in.set_xlabel('FCF_j')
    # ax_in.set_ylabel('')

    # ax = axes[0]
    # ax.set_xlabel('FCF_j')
    # ax.set_title('Functional causal flow')
    # ax.set_xlim([0,0.45])
    # ax_in.set_xlim([0,0.45])

    # ax_in = axes[1].inset_axes([0.7, 0.7, 0.3, 0.3])
    # dir_j_bluepill, dir_j_redpill = calculate_average(directionality_mask,filename_list,nSal,cmap=cmap,ax=ax_in)
    # # ax_in.set_xlabel('Directionality_j')
    # ax_in.set_xlim([-0.1,0.1])
    # ax_in.set_ylabel('')

    # ax = axes[1]
    # ax.set_xlabel('Directionality_j')
    # ax.set_title('Directionality')
    # ax.set_xlim([-0.1,0.1])
    # # ax.legend()

    # ax_in = axes[2].inset_axes([0.7, 0.7, 0.3, 0.3])
    # max_takens = np.round(np.nanpercentile(complexity_mask.ravel(),97.5))
    # cpx_j_bluepill, cpx_j_redpill = calculate_average(complexity_mask,filename_list,nSal,cmap=cmap,ax=ax_in)
    # # ax_in.set_xlabel('Takens_j')
    # ax_in.set_ylabel('')

    # ax = axes[2]
    # ax.set_xlabel('Takens_j')
    # ax.set_title('Complexity\nof source')
    # ax.set_xlim([0,max_takens]); ax_in.set_xlim([0,max_takens])

    # ax_in = axes[3].inset_axes([0.7, 0.7, 0.3, 0.3])
    # cpx_i_bluepill, cpx_i_redpill = calculate_average(complexity_mask,filename_list,nSal,cmap=cmap,axis=1,ax=ax_in)
    # # # ax_in.set_xlabel('Takens_i')
    # # ax_in.set_ylabel('')
    # ax = axes[3]
    # ax.set_xlabel('Takens_i')
    # ax.set_title('Complexity\nof target')
    # ax.set_xlim([0,max_takens]); ax_in.set_xlim([0,max_takens])
    # # ax.legend()

    # for ax in axes:
    #     xlim = ax.get_xlim()
    #     ylim = ax.get_ylim()
    #     usrplt.adjust_spines(ax, ['left', 'bottom'])
    # usrplt.save_fig_to_pptx(fig, prs)
    # plt.savefig(join(PlotDir,f'FCF_j_mean_mask-kfold_{rec_name}.pdf'))

    # def plot_CCM_differences2(corr_mat1,corr_mat2,fcf1,fcf2,cond1,cond2,c1,c2,title=None,ticks=None,labels=None,boundaries=None):

    #     fig, axes = plt.subplots(1,3,figsize=(12,4))#,gridspec_kw={'wspace': 0.1,'hspace':0.1})
    #     if title is None:
    #         title = f'{cond2} - {cond1}'

    #     plt.suptitle(title,fontsize=22,fontweight='bold',y=0.95)
    #     ax = axes[0]
    #     ax.set_title('\u0394-|Correlation|',fontsize=14,fontweight='bold')
    #     dCorr = np.abs(corr_mat2) - np.abs(corr_mat1)
    #     vmax = np.round(np.nanpercentile(dCorr,98),2)
    #     vmin = np.round(np.nanpercentile(dCorr,2),2)
    #     sns.heatmap(dCorr,square=True,annot=False,cmap='RdBu_r',center=0,vmin=vmin,vmax=vmax,ax=ax,cbar_kws={'shrink':0.5,'label': '\u0394-|Correlation|'},rasterized=True) 

    #     ax = axes[1]
    #     ax.set_title('\u0394-FCF',fontsize=14,fontweight='bold')
    #     dFCF = fcf2 - fcf1
    #     vmax = np.round(np.nanpercentile(dFCF,98),2)
    #     vmin = np.round(np.nanpercentile(dFCF,2),2)
    #     sns.heatmap(dFCF,square=True,annot=False,cmap='RdBu_r',center=0,ax=ax,vmin=vmin,vmax=vmax,cbar_kws={'shrink':0.5,'label': '\u0394-FCF'},rasterized=True)

    #     ax = axes[2]
    #     fcf2_j = np.nanmean(fcf2,axis=0)
    #     fcf1_j = np.nanmean(fcf1,axis=0)

    #     sns.kdeplot(fcf1_j,ax=ax,color=c1,ls='-',label=cond1,lw=2)
    #     sns.kdeplot(fcf2_j,ax=ax,color=c2,ls='-',label=cond2,lw=2)
    #     ax.legend()
    #     if ticks is None:
    #         for ax in axes[:-1]:
    #             ax.set_xticks([]);ax.set_yticks([])
    #     else:
    #         for ii, ax in enumerate(axes[:-1]):
    #             ax.set_xticks(ticks);ax.set_yticks(ticks)
    #             ax.set_xticklabels(labels,rotation=90)
    #             if ii == 0:
    #                 ax.set_yticklabels(labels)
    #             else:
    #                 ax.set_yticklabels([])
            
    #     if boundaries is not None:
    #         for ii, ax in enumerate(axes[:-1]):
    #             if (ii == 2) | (ii == 3):
    #                 c = 'k'
    #             else:
    #                 c = 'w'
    #             ax.vlines(boundaries,*ax.get_ylim(),color=c,lw=0.5,alpha=0.5)
    #             ax.hlines(boundaries,*ax.get_xlim(),color=c,lw=0.5,alpha=0.5)
    #     return fig
    
    # from itertools import combinations
    # combo_list = []
    # for cb in list(combinations(np.arange(len(filename_list)),2)):
    #     ii, jj = cb
        

    #     if ('psi' not in filename_list[ii]) & ('psi' in filename_list[jj]):
    #         combo_list.append(cb)
    #         print(f'({ii}, {jj}) -> {filename_list[ii]} vs {filename_list[jj]}')

    # for cb in combo_list:
    #     ii, jj = cb
    #     title = f'{filename_list[ii]} vs {filename_list[jj]}'
    #     fig = plot_CCM_differences2(correlation_mkf_ro[ii],correlation_mkf_ro[jj],FCF_optimal_mkf_ro[ii],FCF_optimal_mkf_ro[jj],filename_list[ii],filename_list[jj],cmap[ii],cmap[jj],title=title,ticks=ticks_group,labels=labels_group,boundaries=boundaries_group)
    #     usrplt.save_fig_to_pptx(fig, prs)
    
    # fig = plot_CCM_differences2(np.nanmean(correlation_mkf_ro[:nSal],axis=0),np.nanmean(correlation_mkf_ro[nSal:],axis=0),np.nanmean(FCF_optimal_mkf_ro[:nSal],axis=0),np.nanmean(FCF_optimal_mkf_ro[nSal:],axis=0),'sal','psi',cmap[ii],cmap[jj],title=title,ticks=ticks_group,labels=labels_group,boundaries=boundaries_group)
    # usrplt.save_fig_to_pptx(fig, prs)
    
    # #Network 
    # import networkx as nx
    # import networkx.algorithms.community as nx_comm
    # from networkx.algorithms.community import greedy_modularity_communities, modularity
    # from networkx.algorithms.efficiency_measures import global_efficiency, local_efficiency

    # FCF_thresh = 0.1
    # DAC = np.zeros((nCond,2))
    # mod = np.zeros((nCond,2))
    # trs = np.zeros((nCond,2))
    # gE = np.zeros((nCond));nCC = np.zeros((nCond))
    # # fig2, axes2 = plt.subplots(1,2,figsize=(10,5))

    # fig3, ax3 = plt.subplots(figsize=(8,5))
    # results = []; degree_dist_list = []
    # for jj, fname in enumerate(filename_list):
    #     print(fname)
        
    #     #Define adjacency matrix 
    #     mask = FCF_optimal_mkf[jj] > FCF_thresh
    #     adjacency_mat = np.array(mask,dtype=int)
    #     weighted_FC = FCF_optimal_mkf[jj].copy()
    #     weighted_FC[~mask] = 0

    #     #Crates graph using the data of the adjacency matrix or weighted directed graph
    #     G = nx.to_networkx_graph(adjacency_mat)
    #     DG = nx.DiGraph(weighted_FC)

    #     #Find communities and calculate modularity
    #     comm_list = greedy_modularity_communities(DG,weight='weight')
    #     mod[jj,1] = modularity(DG,comm_list)

    #     indy_order = []
    #     for s in comm_list:
    #         indy_order.append(list(s))
    #     indy_order = np.concatenate(indy_order)

    #     fig, axes = plt.subplots(1,2,figsize=(10,5))
    #     plt.suptitle(f'FCF: {filename_list[jj]}')
    #     ax = axes[0];
    #     usrplt.visualize_matrix(FCF_optimal_mkf[jj][indy_order][:,indy_order],ax=ax,cmap='viridis',title='Sorted by communities from directed graph',clims=[0,vmax_fcf],center=None)#,ticks=ticks,labels=labels,boundaries=boundaries)

    #     #Find communities and calculate modularity
    #     comm_list = greedy_modularity_communities(G)
    #     mod[jj,0] = modularity(DG,comm_list)

    #     indy_order = []
    #     for s in comm_list:
    #         indy_order.append(list(s))
    #     indy_order = np.concatenate(indy_order)

    #     ax = axes[1]
    #     usrplt.visualize_matrix(FCF_optimal_mkf[jj][indy_order][:,indy_order],ax=ax,cmap='viridis',title=f'Sorted by communities from undirected graph',clims=[0,vmax_fcf],center=None)#,ticks=ticks,labels=labels,boundaries=boundaries)
    #     usrplt.save_fig_to_pptx(fig, prs)
        
    #     #Get degree distribution
    #     degree_dist = list(dict(DG.degree()).values())
    #     degree_dist_list.append(degree_dist)
    #     sns.kdeplot(degree_dist,ax=ax3,color=cmap[jj],label=fname)
        
    #     #Get degree assortativity
    #     DAC[jj,0] = nx.degree_assortativity_coefficient(G)
    #     DAC[jj,1] = nx.degree_assortativity_coefficient(DG)

    #     #Compute graph transitivity, the fraction of all possible triangles present in G.
    #     trs[jj,0] = nx.transitivity(G)
    #     trs[jj,1] = nx.transitivity(DG)

    #     #Calculate global efficiency 
    #     gE[jj] = nx.global_efficiency(G)

    #     #Calculate whether the network is "whole" or disjointed into multiple subnetworks
    #     nCC[jj] = nx.number_connected_components(G)

    # ##---------------
    # fig, ax = plt.subplots(figsize=(8,5))
    # ax.set_title('Modularity')
    # plt.plot(mod[:,0],'-k')
    # plt.plot(mod[:,1],'--k')
    # for jj, fname in enumerate(filename_list):
    #     ax.plot(jj,mod[jj,0],'o',color=cmap[jj])
    #     ax.plot(jj,mod[jj,1],'o',color=cmap[jj])
    # ax.set_xticks(np.arange(len(filename_list)))
    # ax.set_xticklabels(filename_list,rotation=30)
    # usrplt.save_fig_to_pptx(fig, prs)

    # ##---------------
    # fig, ax = plt.subplots(figsize=(8,5))
    # ax.set_title('Degree assortativity')
    # plt.plot(DAC[:,0],'-k')
    # plt.plot(DAC[:,1],'--k')
    # for jj, fname in enumerate(filename_list):
    #     ax.plot(jj,DAC[jj,0],'o',color=cmap[jj])
    #     ax.plot(jj,DAC[jj,1],'o',color=cmap[jj])
    # ax.set_xticks(np.arange(len(filename_list)))
    # ax.set_xticklabels(filename_list,rotation=30)
    # usrplt.save_fig_to_pptx(fig, prs)

    # ##---------------
    # fig, ax = plt.subplots(figsize=(8,5))
    # ax.set_title('Transitivity')
    # plt.plot(trs[:,0],'-k')
    # plt.plot(trs[:,1],'--k')
    # for jj, fname in enumerate(filename_list):
    #     ax.plot(jj,trs[jj,0],'o',color=cmap[jj])
    #     ax.plot(jj,trs[jj,1],'o',color=cmap[jj])
    # ax.set_xticks(np.arange(len(filename_list)))
    # ax.set_xticklabels(filename_list,rotation=30)
    # usrplt.save_fig_to_pptx(fig, prs)

    # ##---------------
    # fig, ax = plt.subplots(figsize=(8,5))
    # ax.set_title('Global Efficiency')
    # plt.plot(gE[:],'-k')

    # for jj, fname in enumerate(filename_list):
    #     ax.plot(jj,gE[jj],'o',color=cmap[jj])
    # ax.set_xticks(np.arange(len(filename_list)))
    # ax.set_xticklabels(filename_list,rotation=30)
    # usrplt.save_fig_to_pptx(fig, prs)

    # ##---------------
    # fig, ax = plt.subplots(figsize=(8,5))
    # ax.set_title('# of subnetworks')
    # plt.plot(nCC[:],'-k')

    # for jj, fname in enumerate(filename_list):
    #     ax.plot(jj,nCC[jj],'o',color=cmap[jj])
    # ax.set_xticks(np.arange(len(filename_list)))
    # ax.set_xticklabels(filename_list,rotation=30)
    # usrplt.save_fig_to_pptx(fig, prs)

    # prs.save(join(SaveDir,f'FCF_figs2_{rec_name}.pptx'))
    # print('DONE!!!')
    # # exit()
