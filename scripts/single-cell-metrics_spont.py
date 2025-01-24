base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
base_dir_local = '/data/tiny-blue-dot/zap-n-zip/EEG_exp'; base_dir = base_dir_server
ProjDir = '/home/david.wyrick/projects/zap-n-zip/'
ServDir = '/data/projects/zap-n-zip/'

#Basea
import gspread
from os.path import join
from glob import glob
import json, os, time, sys, argparse
import numpy as np
import pandas as pd
import scipy.stats as st
import pingouin as pg
from scipy.ndimage import gaussian_filter
from scipy.interpolate import interp1d
from itertools import combinations

#Plot
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
from matplotlib.gridspec import GridSpec

#Project
from tbd_eeg.data_analysis.eegutils import EEGexp
import tbd_eeg.data_analysis.Utilities.utilities as tbd_util
from tbd_eeg.data_analysis.Utilities.behavior_movies import Movie

#Allen
from allensdk.core.mouse_connectivity_cache import MouseConnectivityCache
mcc = MouseConnectivityCache(resolution=10)
str_tree = mcc.get_structure_tree()

#Read in allen CCF
ccfsum = pd.read_csv('/home/david.wyrick/projects/zap-n-zip/data/mouse_connectivity/ontology_v2.csv')

#User
sys.path.append(ProjDir)
import util
import plotting as usrplt

#Single cell stuff
import elephant.statistics as est
from quantities import ms, s, Hz
from neo.core import SpikeTrain

#Define behavioral states
behavior_ranges = {0: [0,1], 1: [1,10], 2: [10,30], 3:[30,500]}
behavior_dict = {0: 'rest (0-1cm/s)', 1: 'walk (1-10cm/s)', 2: 'shuffle (10-30cm/s)', 3: 'run (>30cm/s)'}
behavior_strs = list(behavior_dict.values())
behavior_strs2 = ['rest','walk','shuffle','run']
nBehaviors = len(behavior_strs)

#Define windows to calculate firing rate
windows_of_interest = [[.002,.025],[.075,.3],[.3,1],[1,4.5]]
window_strs = ['evoked (2-25ms)','rebound (75-300ms)','post-rebound (0.3-1s)','ISI (1-4.5s)']
nWindows = len(window_strs)

gc = gspread.service_account() # need a key file to access the account (step 2) 
sh = gc.open('Templeton-log_exp') # open the spreadsheet 

df = pd.DataFrame(sh.sheet1.get()) 
exp_table = df.T.set_index(0).T # put it in a nicely formatted dataframeexp_table.iloc[10:]
exp_table = exp_table.iloc[9:].set_index('mouse_name')

##===== ============================ =====##
##===== Parse Command Line Arguments =====##
parser = argparse.ArgumentParser(description='single-cell-metrics')

##===== Data Options =====##
parser.add_argument('--mouseID',type=str, default='mouse678912',
                    help='mouse to perform analysis on')

parser.add_argument('--time_bin_ms',type=int, default=100,
                    help='time_bin_ms')

parser.add_argument('--fr_thresh',type=float, default=2,
                    help='Firing rate threshold for neurons to include in analysis')


if __name__ == '__main__':

    # Parse the arguments
    args = parser.parse_args()
    mouse_name = args.mouseID 
    mID = args.mouseID
    fr_thresh = args.fr_thresh
    time_bin_ms = args.time_bin_ms
    time_bin = time_bin_ms/1000
    calculate_evoked = bool(args.calculate_evoked)
    
    if mID == 'mouse678912': 
        rec_name = 'spont_aw_psi_2023-06-22_11-42-00'
        time_window_list = [[48.0, 948.0],[948.0, 1848.0],[2582.04, 3002.9],[3242.9, 4142.9],[4142.9, 5042.9],[5042.9, 5942.9],[5942.9, 6842.9],[6842.9, 7742.9],[7142.9, 8042.9]]
        filename_list = ['spont_no-inj_0-15','spont_no-inj_15-30','spont_saline_42-52','spont_psilo_0-15','spont_psilo_15-30','spont_psilo_30-45','spont_psilo_45-60','spont_psilo_60-75','spont_psilo_65-80']
    elif  mID == 'mouse678913': 
        rec_name = 'spont_aw_psi_2023-06-29_12-49-40'
        time_window_list = [[50.0, 950.0],[950.0, 1850.0],[2669, 3172.3],[3292, 4192],[4192, 5092],[5092, 5992],[5992, 6802]]
        filename_list = ['spont_no-inj_0-15','spont_no-inj_15-30','spont_saline_42-52','spont_psilo_0_15','spont_psilo_15_30','spont_psilo_30_45','spont_psilo_45_60']
    elif mID == 'mouse689241':
        rec_name = 'spont_aw_psi_2023-07-27_11-05-05'
        filename_list = ['spont_no-inj_0-15','spont_saline_17-32','spont_saline_32-47','spont_saline_47-64','spont_psilo_66-81','spont_psilo_81-96','spont_psilo_96-111','spont_psilo_110-125']
        time_window_list = [[43.0, 937.0],[1057.0, 1957.0],[1957.0, 2857.0],[2857.0, 3741.0],[3981.0, 4881.0],[4881.0, 5781.0],[5781.0, 6681.0],[6613.0, 7513.0]]
    elif mID == 'mouse692643': 
        rec_name = 'spont_aw_psi_2023-08-31_12-41-54'
        time_window_list = [[43.0, 924.37],[1044.37, 1929.37],[1929.37, 2829.37],[2717.87, 3482.87],[3722.87, 4367.87],[4607.87, 5507.87],[5507.87, 6407.87]]
        filename_list = ['spont_no-inj','spont_saline_0-15','spont_saline_15-30','spont_saline_30-45','spont_psilo_0-15','spont_psilo_15-30','spont_psilo_30-45']
    else:
        raise Exception(f'Windows not set for {mID}')

    #Get injection times and types of injection
    injection_times = [float(exp_table.loc[exp_table.exp_name == rec_name]['First injection time'].values[0]),
                    float(exp_table.loc[exp_table.exp_name == rec_name]['Second injection time'].values[0])]

    if 'psi' in rec_name:
        exp_type = 'psilocybin'
        injection_types = ['sal1','psi']
        cmap3 = np.concatenate((sns.color_palette('Blues',2),sns.color_palette('Reds',1)))
    else:
        injection_types = ['sal1', 'sal2']
        exp_type = 'saline'
        cmap3 = sns.color_palette('Blues',3)

    combos = []
    for i, epoch_i in enumerate(filename_list):
        if ('psi' in epoch_i) | ('evoke' in epoch_i):
            continue
        for j, epoch_j in enumerate(filename_list):
            if ('psi' in epoch_j):
                combos.append([i,j])
    epoch_type = []
    e_arr = []
    for fname in filename_list:
        if 'no-inj' in fname:
            epoch_type.append('no-inj')
            e_arr.append(0)
        elif 'saline' in fname:
            epoch_type.append('saline')
            e_arr.append(1)
        elif 'psilo' in fname:
            epoch_type.append('psilo')
            e_arr.append(2)
    e_arr = np.array(e_arr)
    cmap = np.concatenate((sns.color_palette('Blues',np.sum(e_arr <=1)),sns.color_palette('Reds',np.sum(e_arr == 2))))

    #Define directories
    SaveDir = join(ProjDir,'results','FR',mID,rec_name)
    if not os.path.exists(SaveDir):
        os.makedirs(SaveDir)

    PlotDir = join(SaveDir,'plots')
    if not os.path.exists(PlotDir):
        os.makedirs(PlotDir)

    #%% Upload the whole experiment and generate the global clock
    file_name = os.path.join(base_dir_server,mID,rec_name,'experiment1','recording1')
    exp = EEGexp(file_name, preprocess=False, make_stim_csv=False)

    #%% load the EEG data
    try:
        eeg_data,eeg_ts = exp.load_eegdata()
    except:
        eeg_data = [np.nan]; eeg_ts = [np.nan]

    #Load probe data
    probe_unit_data, probe_info, total_units = util.get_neuropixel_data(exp)
    probe_list = [x.replace('_sorted', '') for x in exp.experiment_data if 'probe' in x]

    #Get recording start time
    probei = probe_list[0]
    open_ephys_start = np.round(np.min(probe_unit_data[probei]['spike_times']))
    open_ephys_end = np.round(np.max(probe_unit_data[probei]['spike_times']))
        
    # Load running speed
    run_file = os.path.join(base_dir,mID,rec_name,'experiment1','recording1','running_signal.npy')
    if os.path.exists(run_file):
        # print('File exists')
        run_signal = np.load(run_file)

        ts_file = os.path.join(base_dir,mID,rec_name,'experiment1','recording1','running_timestamps.npy')
        if os.path.exists(ts_file):
            run_ts = np.load(ts_file)
        else:
            ts_file = os.path.join(base_dir,mID,rec_name,'experiment1','recording1','running_timestamps_master_clock.npy')
            run_ts = np.load(ts_file)
    else:
        print('\tRunning file does not exist')
        run_signal, run_ts = exp.load_running()
        np.save(os.path.join(base_dir,mID,rec_name,'experiment1','recording1','running_signal.npy'),run_signal)
        np.save(os.path.join(base_dir,mID,rec_name,'experiment1','recording1','running_timestamps.npy'),run_ts)

    #Load pupil data
    pupil_csv = os.path.join(base_dir,mID,rec_name,'experiment1','recording1',f'Pupil_{rec_name}.csv')

    try:
        table = pd.read_csv(pupil_csv)
        pupil_radius = table['Largest_Radius'].values

        #Pupil master clock
        pupil_ts = Movie(filepath=exp.pupilmovie_file,
                        sync_filepath=exp.sync_file,
                        sync_channel='eyetracking'
                        ).sync_timestamps
        plot_pupil = True
    except:
        pupil_ts = np.array([np.nan])
        pupil_radius =np.array([np.nan])
        plot_pupil = False

    #Smooth running signal with gaussian 
    run_signal_smoothed = gaussian_filter(run_signal,10)
    acceleration_smoothed = np.gradient(run_signal_smoothed)

    # Create a new PowerPoint presentation to save figures to
    prs = Presentation()

    #Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f'mID = {mID}; rec_name = {rec_name}'
    slide.placeholders[1].text = f'Firing rate changes due to psilocybin during spontaneous experiment'

    #Plot behavior
    fig = plt.figure(figsize=(9,6))
    plt.suptitle(f'{mID}, {rec_name}')
    gs = GridSpec(2, 3, figure=fig,height_ratios=[2,1.5])

    ax = fig.add_subplot(gs[1, :])
    ax.plot(run_ts/60,run_signal_smoothed,'-k',lw=1)
    ax.set_xlabel('Time (mins)')
    ax.set_ylabel('Speed (cm/s)')
    ax.autoscale(tight=True)

    ylim = ax.get_ylim()
    # ax = fig.add_subplot(gs[2, :])
    ax.vlines(injection_times[0]/60,*ylim,color=cmap3[1],lw=3,zorder=4,label=injection_types[0])
    ax.vlines(injection_times[1]/60,*ylim,color=cmap3[2],lw=3,zorder=4,label=injection_types[0])

    # for ii, tW in enumerate(time_window_list):
    #     ax.vlines(np.array(tW)/60,*ax.get_ylim(),color=cc[0],lw=1.5,zorder=4)

    for ii, tW in enumerate(time_window_list):
        ax.fill_between(tW/60,*ylim,color=cmap[ii],lw=1.5,zorder=0,alpha=0.5,label=filename_list[ii])


    axes = []
    for ii, (ts, sig) in enumerate(zip([run_ts,run_ts,pupil_ts],[run_signal,acceleration_smoothed,pupil_radius])):
        ax = fig.add_subplot(gs[0, ii]); axes.append(ax)
        ax.set_title(f'{mID}, {rec_name}')
        for jj, tW in enumerate(time_window_list):        
            indy = np.where((ts >= tW[0]) & (ts <= tW[1]))[0]
            sns.kdeplot(sig[indy],ax=ax,label=filename_list[jj],color=cmap[jj],lw=1.5,common_norm=False)

    axes[0].set_title('Running speed')
    axes[1].set_title('Acceleration');axes[1].set_xlim([-0.5,0.5]); axes[1].legend()
    axes[2].set_title('Pupil radius'); axes[2].set_xlim([0,140])

    usrplt.save_fig_to_pptx(fig, prs)
    plt.savefig(join(PlotDir,f'behavior_{rec_name}.pdf'))

    #Read in neuropixel data
    spont_data_list, spont_ts_list, neuron_indices, plot_tuple = util.bin_spiking_data(probe_unit_data, time_window_list, time_bin=time_bin,fr_thresh=fr_thresh)
    boundaries_group, ticks_group, labels_group, celltypes, durations, areas, groups, supergroups, order_by_group = plot_tuple

    areas_sub = areas[neuron_indices]
    groups_sub = groups[neuron_indices]
    N = len(neuron_indices)

    #Plot rasters for each epoch
    jj = 0
    neurons_per_probe = [len(np.unique(probe_unit_data[probei]['units'])) for probei in probe_list]
    neurons_per_probe_cum = np.cumsum(neurons_per_probe)
    for jj, tWindow in enumerate(time_window_list):
        ts = spont_ts_list[jj]
        spikes = spont_data_list[jj]
        T, N = spikes.shape
        
        # Display a spike raster of the image
        tmp = np.concatenate((neurons_per_probe,[150]))
        fig, axes =  plt.subplots(len(tmp),1,figsize=(12,12),gridspec_kw={'height_ratios':np.array(tmp)/np.sum(tmp)})
        plt.suptitle(filename_list[jj],y=0.925)

        #Define time window
        if T*time_bin < 60:
            tStart = 0
        else:
            tStart = 30# 5*60
        time_to_plot = 30 #sec
        tslice = slice(int(tStart/time_bin),int((tStart+time_to_plot)/time_bin))
        ts_sub = ts[tslice]

        #Params
        xticks = np.arange(0,int((time_to_plot+1)/time_bin),int(15/time_bin))
        tlabel = xticks*time_bin
        xticks = np.arange(ts_sub[0],ts_sub[-1]+1,15)

        for ii, probei in enumerate(probe_list):
            ax = axes[ii]; ax.set_title(probei)
            if ii == 0:
                nslice = slice(0,neurons_per_probe[0])
            else:
                nslice = slice(neurons_per_probe_cum[ii-1],neurons_per_probe_cum[ii])
            print(np.percentile(spikes[tslice,nslice].T,97))
            ax.imshow(spikes[tslice,nslice].T, aspect='auto',vmax=5,vmin=0, cmap='gray_r')
            ax.set_xticks([])

            uniq_groups, uniq_indices, num_groups = np.unique(groups[nslice],return_index=True, return_counts=True)
            group_labels = uniq_groups[np.argsort(uniq_indices)]
            nNeurons_area = num_groups[np.argsort(uniq_indices)]

            boundaries = np.concatenate(([0],np.cumsum(nNeurons_area)))
            yticks = boundaries[:-1] + np.diff(boundaries)/2

            ax.set_yticks(yticks,minor=True)
            ax.set_yticks(boundaries,minor=False)
            ax.set_yticklabels(group_labels,minor=True)
            ax.set_yticklabels([],minor=False)
            ax.set_ylabel("Unit")

        #Plot pupil and running
        indy = np.where((run_ts >= ts_sub[0]) & (run_ts <= ts_sub[-1]))[0]
        ax = axes[-1]; ax.autoscale(tight=True)
        ax.plot(run_ts[indy],run_signal[indy],lw=0.6,color='k')
        ax.set_ylabel("Speed (cm/s)"); ax.set_xlabel("Time (s)")
        # ax.set_xticks([]); 
        ax.set_title('Behavior')

        if plot_pupil:
            indy = np.where((pupil_ts >= ts_sub[0]) & (pupil_ts <= ts_sub[-1]))[0]
            ax = ax.twinx()
            ax.plot(pupil_ts[indy],pupil_radius[indy],lw=0.6,color=cc[1])
            ax.set_ylabel("Pupil radius (pix)") 

        usrplt.save_fig_to_pptx(fig, prs)
        plt.savefig(os.path.join(PlotDir,f'raster_{filename_list[jj]}.pdf'),facecolor='white',dpi=300,bbox_inches='tight')

    #Look at spontaneous periods first
    #Smooth running signal with gaussian 
    run_signal_smoothed = gaussian_filter(run_signal,10)
    f_run = interp1d(run_ts,run_signal_smoothed)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    num_slides_per_slide = 8
    #Determine behaviors during spontaneous periods
    spont_behav_list = []
    for jj, tWindow in enumerate(time_window_list):
        print(f'{filename_list[jj]}')
        if jj > 7:
            if jj == 8:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
            fig_num = jj - 8
        else:
            fig_num = jj
        position = [(fig_num//4)*5, (fig_num%4)*1.875]
        
        ts = spont_ts_list[jj]
        T = len(ts)
        tMin = T*time_bin/60
        run_spont = f_run(ts)
        behavior_map = np.zeros(run_spont.shape)

        for key, b_range in behavior_ranges.items():
            indy = np.where((run_spont >= b_range[0]) & (run_spont < b_range[1]))[0]
            behavior_map[indy] = key
            t = len(indy)*0.01
            print(f'\t: {behavior_dict[key]:10s} -> {t:.1f}s')
        spont_behav_list.append(behavior_map)

        fig, ax = plt.subplots(figsize=(5,1))
        ax.set_title(filename_list[jj])
        ax2 = ax.twinx()
        sns.heatmap(behavior_map.reshape(1,-1),vmax=4,vmin=0,ax=ax,cbar=False,cmap=sns.color_palette('Reds',nBehaviors),cbar_kws={'ticks':np.arange(nBehaviors)})
        
        ax2.plot(run_spont,color='k')
        ax2.set_ylabel('Speed (cm/s)')
        ax.set_yticks([])
        ax.set_xticks([0,T/2,T]); ax.set_xticklabels([0,np.round(tMin/2,2),np.round(tMin,2)],rotation=0);ax.set_xlabel('Time (mins)')

        usrplt.save_fig_to_pptx(fig, prs,slide=slide,position=position)

    #Calculate firing rates during different behaviors during spontaneous 
    T, nNeurons = spont_data_list[0].shape
    nCond_spont = len(time_window_list)
    FR_spont = np.full((nCond_spont,nNeurons,nBehaviors),np.nan)
    nBins_per_behav_spont = np.zeros((nCond_spont,nBehaviors))

    for jj, tWindow in enumerate(time_window_list):
        ts = spont_ts_list[jj]
        data_spont = spont_data_list[jj]
        behavior_map = spont_behav_list[jj]

        for key, behav_str in behavior_dict.items():
            indy = np.where(behavior_map == key)[0]
            if len(indy) == 0:
                continue
            FR_spont[jj,:,key] = np.mean(data_spont[indy],axis=0)/time_bin
            nBins_per_behav_spont[jj,key] = len(indy)

    #Save time windows
    np.savez(join(SaveDir,f'time_windows.npz'),time_window_list=time_window_list,filename_list=filename_list,nBins_per_behav_spont=nBins_per_behav_spont,areas=areas)

    #Convert firing rate array into dataframe for easy plotting
    for indy, fsuffix in zip([np.arange(nNeurons),neuron_indices],['ALL','sub']):

        #Compare FR changes due to saline or psilocybin injections vs no-injection
        tmp_list = []
        for j, tWindow in enumerate(time_window_list):   
            for key, behav_str in behavior_dict.items():
                #Select neurons FR for the epoch & behavior 
                FR_j = FR_spont[j,indy,key]; epoch_j = filename_list[j]
                N = FR_j.shape[0]
                #Check to make sure each epoch has 5 seconds of data
                bin_thresh = 5/time_bin
                if (nBins_per_behav_spont[jj,key] < bin_thresh):
                    continue

                tmp_list.append(np.stack((np.repeat(mID,N),np.repeat(rec_name,N),indy,areas[indy],groups[indy],np.repeat(behav_str,N),np.repeat(epoch_j,N),np.repeat(epoch_type[j],N),FR_j)))
        FR_spont_df = pd.DataFrame(np.hstack(tmp_list).T,columns=['mID','rec_name','neuron','area','group','behavior','epoch','injection','FR'])
        FR_spont_df = FR_spont_df.astype({'mID':str,'rec_name':str,'neuron': int,'area':str,'group':str,'behavior':str,'epoch':str,'injection':str,'FR':float})
        FR_spont_df.to_hdf(join(SaveDir,f'spont_FR_{fsuffix}_{rec_name}.h5'),'df')

        #Compare FR changes due to saline or psilocybin injections vs no-injection
        tmp_list = []
        for i, j in combos:
            for key, behav_str in behavior_dict.items():
                #Select neurons FR for the epoch & behavior 
                FR_i = FR_spont[i,indy,key]; epoch_i = filename_list[i]
                FR_j = FR_spont[j,indy,key]; epoch_j = filename_list[j]

                #Check to make sure each epoch has 5 seconds of data
                bin_thresh = 5/time_bin
                if (nBins_per_behav_spont[0,key] < bin_thresh) | (nBins_per_behav_spont[jj,key] < bin_thresh):
                    continue
                
                dFR = FR_j - FR_i
                MI = (FR_j - FR_i)/(FR_j + FR_i)
                N = MI.shape[0]
                tmp_list.append(np.stack((np.repeat(mID,N),np.repeat(rec_name,N),indy,areas[indy],groups[indy],np.repeat(behav_str,N),np.repeat(epoch_i,N),np.repeat(epoch_j,N),FR_i,FR_j,dFR,MI)))
            
        dFR_spont_df = pd.DataFrame(np.hstack(tmp_list).T,columns=['mID','rec_name','neuron_index','area','group','behavior','epoch_i','epoch_j','FR_i','FR_j','dFR','MI'])
        dFR_spont_df = dFR_spont_df.astype({'mID':str,'rec_name':str,'neuron_index':int,'area':str,'group':str,'behavior':str,'epoch_i':str, 'epoch_j':str,'FR_i':float,'FR_j':float,'dFR':float,'MI':float})
        dFR_spont_df.to_hdf(join(SaveDir,f'delta-spont_FR_{fsuffix}_{rec_name}.h5'),'df')

    #Create slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Lets look at firing rate per group for each "epoch_j"'

    tf = body_shape.text_frame
    tf.text = 'Now the x-axis will be "group"'

    p = tf.add_paragraph()
    p.text = f'This is being compared to either "{filename_list[0]}" or "{filename_list[1]}"'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Note that multiple areas may be contained within a group'
    p.level = 1

    #Plot boxplot across groups for each epoch
    uniq_groups = np.unique(dFR_spont_df['group'])
    nGroups = len(uniq_groups)

    pval_thresh = 0.05
    for i, j in combos:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f'Spontaneous activity comparison'
        
        fig = plot_boxplot_comparison(dFR_spont_df, i, j,filename_list, areas_sub, groups_sub, pval_thresh = 0.05)
        save_fig_to_pptx(fig, prs,slide=slide)
        plt.savefig(join(PlotDir,f'delta-FR_{filename_list[i]}_vs_{filename_list[j]}.pdf'),facecolor='white',dpi=300,bbox_inches='tight')
        plt.close(fig)

    # Save the PowerPoint presentation
    prs.save(join(PlotDir,f'firing_rate_spont_figs_{rec_name}.pptx'))
    print('DONE!!!')


