base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
base_dir_local = '/data/tiny-blue-dot/zap-n-zip/EEG_exp'; base_dir = base_dir_server
ProjDir = '/home/david.wyrick/projects/zap-n-zip/'
ServDir = '/data/projects/zap-n-zip/'
BehavDir = '/home/david.wyrick/projects/zap-n-zip/plots/behavior'

##------------------------------------------
##Load modules
#Base
import gspread
from os.path import join
from glob import glob
import json, os, time, sys, argparse
import numpy as np
import pandas as pd
import itertools as it
import scipy.stats as st
from scipy.interpolate import interp1d
from tqdm import tqdm

#Plot
import matplotlib.pyplot as plt
from pptx import Presentation
import seaborn as sns

#Project
from tbd_eeg.data_analysis.eegutils import EEGexp
import tbd_eeg.data_analysis.Utilities.utilities as tbd_util

#Allen
from allensdk.core.mouse_connectivity_cache import MouseConnectivityCache
mcc = MouseConnectivityCache(resolution=10)
str_tree = mcc.get_structure_tree()

#User
sys.path.append(ProjDir)
import util
import plotting as usrplt

#CCG
sys.path.append('/home/david.wyrick/Git/functional-network')
import ccg_library as ccg_lib

gc = gspread.service_account() # need a key file to access the account (step 2) 
sh = gc.open('Templeton-log_exp') # open the spreadsheet 

df = pd.DataFrame(sh.sheet1.get()) 
exp_table = df.T.set_index(0).T # put it in a nicely formatted dataframeexp_table.iloc[10:]
exp_table = exp_table.set_index('mouse_name')

##------------------------------------------
##Define Command Line Arguments
parser = argparse.ArgumentParser(description='single-cell-metrics')

parser.add_argument('--mID',type=str, default='mouse666196',
                    help='mouse to perform analysis on')

parser.add_argument('--rec_name',type=str, default='pilot_aw_psi_2023-03-16_10-21-29',
                    help='experiment to perform analysis on')

parser.add_argument('--time_bin_ms',type=int, default=1,
                    help='time_bin_ms')

parser.add_argument('--FR_thresh',type=int, default=2,
                    help='FR_thresh')

parser.add_argument('--num_jitter',type=int, default=100,
                    help='num_jitter')

parser.add_argument('--ccg_window',type=int, default=100,
                    help='ccg_window')

def plot_significant_ccg(significant_ccg, significant_offset, title_str, ticks, boundaries, labels):
    fig, axes = plt.subplots(1,2,figsize=(10,5))
    plt.suptitle(title_str,y=0.96)

    vmax = np.nanpercentile(np.abs(significant_ccg),98) #vmin=-4E-2,vmax=4E-2
    significant_ccg[np.isnan(significant_ccg)] = 0
    axes[0].set_title('Significant cross-correlograms')
    sns.heatmap(significant_ccg,cmap='RdBu_r',center=0,vmin=-1*vmax,vmax=vmax,square=True,cbar_kws={'label':'CCG','shrink':.5},ax=axes[0])

    significant_offset[np.isnan(significant_offset)] = 0
    axes[1].set_title('Offset of connections')
    sns.heatmap(significant_offset,cmap='viridis',square=True,cbar_kws={'label':'Lag (ms)','shrink':.5},ax=axes[1])

    for ax in axes:
        ax.set_yticks(ticks,minor=True)
        ax.set_yticks(boundaries,minor=False)
        ax.hlines(boundaries,*ax.get_xlim(),color='k',lw=1)
        ax.set_yticklabels(labels,minor=True)
        ax.set_yticklabels([],minor=False)

        ax.set_xticks(ticks,minor=True)
        ax.set_xticks(boundaries,minor=False)
        ax.vlines(boundaries,*ax.get_ylim(),color='k',lw=1)
        ax.set_xticklabels(labels,minor=True,rotation=30)
        ax.set_xticklabels([],minor=False)
    return fig

if __name__ == '__main__':

    ##------------------------------------------
    #Parse the arguments
    args = parser.parse_args()
    mID = args.mID
    rec_name = args.rec_name
    # fr_thresh = args.fr_thresh
    time_bin_ms = args.time_bin_ms
    time_bin = time_bin_ms/1000
    FR_thresh = args.FR_thresh
    num_jitter = args.num_jitter
    ccg_window = args.ccg_window

    #Extract row from templeton dataframe
    exp_df = exp_table.loc[exp_table.exp_name == rec_name]
    if type(exp_df) == pd.core.series.Series:
        exp_df = exp_df.to_frame().T

    drug_type = exp_df['drug'].values[0]
    stim_type = exp_df['stimulation'].values[0]
    if 'electrical' not in stim_type:
        print(f'{rec_name} does not have electrical stimulation')
        exit()
    print(f'Experiment type: {stim_type}, {drug_type}')

    #Define directories
    TempDir = join(ServDir,'results','ccg',mID,rec_name)
    if not os.path.exists(TempDir):
        os.makedirs(TempDir)

    tmp_list = sorted(glob(join(TempDir,f'ccg_run_*')))
    if len(tmp_list) == 0:
        curr_run = 0
    else:
        last_run = int(tmp_list[-1][-1])
        curr_run = last_run+1

    folder = f'ccg_run_{curr_run:02d}'
    SaveDir = os.path.join(TempDir,folder)
    if not os.path.exists(SaveDir):
        os.makedirs(SaveDir)
        
    #Save parameters
    args_dict = args.__dict__
    args_dict['SaveDir'] = SaveDir
    
    with open(join(SaveDir,f"parameters_run_{curr_run}.json"), "w") as outfile:
        json.dump(args_dict, outfile)

    ##------------------------------------------
    #Upload the whole experiment 
    file_name = os.path.join(base_dir_server,mID,rec_name,'experiment1','recording1')
    exp = EEGexp(file_name, preprocess=False, make_stim_csv=False)

    #Load probe data
    probe_unit_data, probe_info, total_units, metrics_list = tbd_util.get_neuropixel_data(exp)
    probe_list = [x.replace('_sorted', '') for x in exp.experiment_data if 'probe' in x]
    plot_tuple = util.get_group_plotting_info(probe_unit_data)
    boundaries, yticks, labels, celltypes, durations, layers, areas, groups, mesogroups, supergroups, order_by_group = plot_tuple
    nNeurons = len(areas)

    #Get recording start time
    probei = probe_list[0]
    open_ephys_start = np.round(np.min(probe_unit_data[probei]['spike_times']))
    open_ephys_end = np.round(np.max(probe_unit_data[probei]['spike_times']))
    recording_length = open_ephys_end - open_ephys_start
    mm, ss = divmod(recording_length,60)
    hh, mm = divmod(mm, 60)
    print(f'{hh} hrs, {mm} minutes, {ss} seconds')

    ##------------------------------------------
    #Get injection times and types of injection
    if drug_type in ['saline', 'psilocybin','ketanserin+psilocybin']:
        injection_time_windows = np.array([np.array(exp_df['First injection window'].values[0].split(','),dtype=float),
                                np.array(exp_df['Second injection window'].values[0].split(','),dtype=float)])
        injection_times = np.array([exp_df['First injection time (s)'].values[0],exp_df['Second injection time (s)'].values[0]],dtype=float)

        if drug_type == 'psilocybin':
            injection_types = ['sal1','psi']
            injection_colors = sns.xkcd_palette(['dark sky blue','darkish red'])
        elif drug_type == 'saline':
            injection_types = ['sal1', 'sal2']
            injection_colors = sns.xkcd_palette(['dark sky blue','cobalt blue'])
        elif drug_type == 'ketanserin+psilocybin':
            injection_types = ['ket','psi']
            injection_colors = sns.xkcd_palette(['magenta','goldenrod'])
        inj_tuple = (injection_times,injection_types,injection_colors)
    else:
        injection_times = None
        inj_tuple = None

    #For isoflurane experiments, get iso level
    if drug_type == 'isoflurane':
        iso_level, iso_times = exp.load_analog_iso()
        iso_induction_times, iso_maintenance_times = exp.load_iso_times()
        induction_colors = sns.xkcd_palette(['light teal','teal'])
        iso_tuple = (iso_induction_times, induction_colors)
    else:
        iso_induction_times = None
        iso_tuple = None

    ##------------------------------------------
    #Define time windows for each epoch  
    if drug_type in ['saline', 'psilocybin','ketanserin+psilocybin']:
        macro_windows = [[open_ephys_start,injection_time_windows[0,0]]]; macro_names = ['pre_inj']
        macro_windows.append([injection_time_windows[0,1],injection_time_windows[1,0]]); macro_names.append(f'post_{injection_types[0]}_inj')
        macro_windows.append([injection_time_windows[1,1],open_ephys_end]); macro_names.append(f'post_{injection_types[1]}_inj')
        if drug_type == 'saline':
            cmap_macro = sns.xkcd_palette(['silver','dark sky blue','cobalt blue'])
        elif drug_type == 'psilocybin':
            cmap_macro = sns.xkcd_palette(['silver','dark sky blue','darkish red'])
        elif drug_type == 'ketanserin+psilocybin':
            cmap_macro = sns.xkcd_palette(['silver','magenta','dusty orange'])
        
    elif drug_type == 'isoflurane':

        macro_windows = [[open_ephys_start,iso_induction_times[0]-120]]; macro_names = ['pre_iso']
        macro_windows.append([iso_induction_times[0],iso_induction_times[1]]); macro_names.append(f'iso_ind')
        macro_windows.append([iso_induction_times[1]+60*5,open_ephys_end]); macro_names.append(f'post-iso')
        cmap_macro = sns.xkcd_palette(['silver','light teal','teal'])

    elif drug_type == 'urethane':
        macro_windows = [[open_ephys_start,open_ephys_end]]; macro_names = ['urethane']
        cmap_macro = sns.xkcd_palette(['darkish purple'])

    ##------------------------------------------
    #Load stimulus log if it exists
    try:
        stim_log = pd.read_csv(exp.stimulus_log_file)

        if ('circle' in np.unique(stim_log['stim_type'])) | ('natural_scene' in np.unique(stim_log['stim_type'])):
            vStim_exists = True
        else:
            vStim_exists = False
        if 'biphasic' in np.unique(stim_log['stim_type']):
            eStim_exists = True
            stim_log_b = stim_log.loc[stim_log.stim_type == 'biphasic']
            uniq_currents = np.unique(stim_log_b['parameter'])
            uniq_sweeps = np.unique(stim_log_b['sweep'])
            nSweeps = len(uniq_sweeps)
            nCurrents = len(uniq_currents)
            print(f'Number of sweeps: {nSweeps}, Number of currents: {nCurrents}')
        else:
            eStim_exists = False
            raise Exception('No stimulus log for biphasic stimulation found')

    except:
        stim_log = None
        eStim_exists = False
        vStim_exists = False
        raise Exception('No stimulus log found')

    evoked_time_window_list = []
    evoked_type_list = []
    sweep_time_window_list = []
    if eStim_exists | vStim_exists:
        for s in np.unique(stim_log['sweep']):
            for t in np.unique(stim_log.loc[stim_log.sweep == s]['stim_type']):
                sub_df = stim_log.loc[(stim_log.sweep == s) & (stim_log.stim_type == t)]
                tS = sub_df['onset'].min(); tE = sub_df['offset'].max()
                
                evoked_time_window_list.append([tS,tE])
                evoked_type_list.append(t)

    ##------------------------------------------
    #Load behavior
    run_ts, run_signal, run_signal_s, acceleration, pupil_ts, pupil_radius, pupil_dxdt, plot_pupil = util.get_behavioral_data(exp, mID, rec_name,normalize=True)
    run_signal[run_signal < 0] = 0
    run_signal_s[run_signal_s < 0] = 0
    f_run = interp1d(run_ts,run_signal)
    f_pupil = interp1d(pupil_ts,pupil_radius)

    # Create a new PowerPoint presentation to save figures to
    prs = Presentation()

    #Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f'mID = {mID}; rec_name = {rec_name}'
    slide.placeholders[1].text = f'Changes in single cell metrics due to {drug_type}'

    #Plot behavior and save to pptx
    fig = usrplt.plot_behavior((run_ts,run_signal),(pupil_ts,pupil_radius),f'{mID} {rec_name}',(evoked_time_window_list,evoked_type_list),inj_tuple,iso_tuple)
    plt.savefig(join(SaveDir,f'behavior_{rec_name}.png'),dpi=300,facecolor='white',bbox_inches='tight')
    usrplt.save_fig_to_pptx(fig, prs)

    ##------------------------------------------
    #Calculate firing rates
    FR = np.concatenate([np.array(probe_unit_data[probei]['firing_rate']) for probei in probe_list])
    FR_mask = FR > FR_thresh
    # neuron_indices = np.where(FR_mask)[0]
    areas_of_interest = ['FRP5', 'FRP6a','MOs5','PO', 'VAL', 'CL','LD','Eth','MGm','PoT','SGN','SPFp']
    groups_of_interest = ['FRP','MO','SSp','ILA','PL','SM-TH','TH','ANT-TH','ACA','RSP','VIS','ORB','HIP','STR']
    areas = np.concatenate([np.array(probe_unit_data[probei]['areas'],dtype=str) for probei in probe_list])
    groups, _, _, _, _, mesogroups, supergroups = util.determine_groups(areas)
    area_mask = np.array([g in groups_of_interest for g in groups])
    # area_mask = np.array([a in areas_of_interest for a in areas])

    neuron_indices = np.where(FR_mask & area_mask)[0]
    N = len(neuron_indices)
    print(f'{N} neurons > {FR_thresh} Hz overall')

    #Get plotting info
    boundaries, ticks, labels, celltypes, durations, layers, areas, groups, mesogroups,supergroups, order_by_group = util.get_group_plotting_info(probe_unit_data,neuron_indices)
    areas_sub = areas[neuron_indices]
    groups_sub = groups[neuron_indices]
    supergroups_sub = supergroups[neuron_indices]
    celltypes_sub = celltypes[neuron_indices]

    areas_ro = areas_sub[order_by_group]
    groups_ro = groups_sub[order_by_group]
    supergroups_ro = supergroups_sub[order_by_group]
    celltypes_ro = celltypes_sub[order_by_group]

    np.savez(join(SaveDir,'neuron_indices.npz'),neuron_indices=neuron_indices,order_by_group=order_by_group,areas=areas,groups=groups,
             supergroups=supergroups,celltypes=celltypes,layers=layers,durations=durations,ticks=ticks,boundaries=boundaries,labels=labels,
             run_ts=run_ts, run_signal=run_signal,pupil_ts=pupil_ts,pupil_radius=pupil_radius,macro_windows=macro_windows,macro_names=macro_names,cmap_macro=cmap_macro)
    
    ##------------------------------------------
    #Calculate cross-correlograms for each sweep/amplitude combination 
    for sweep in np.unique(stim_log_b['sweep']):
        tmp_da_list = []; tmp_ts_list = []

        #Get minimum inter-stimulus-interval
        sub_df = stim_log_b.loc[(stim_log_b.sweep == sweep)]
        ISIs = sub_df['onset'].values[1:] - sub_df['offset'].values[:-1]
        minISI = np.nanmean(ISIs)
        plot_before = 0.5; plot_after = minISI - np.mod(minISI,0.5); time_bin = 1/1000

        for amp in np.unique(stim_log_b['parameter']):
            sub_df = stim_log_b.loc[(stim_log_b.sweep == sweep) & (stim_log_b.parameter == amp)]
            
            #Get stimulation times
            event_times = sub_df['onset'].values
            sweep_epoch = ''
            for e, tW in zip(macro_names,macro_windows):
                if (event_times[0] > tW[0]) and (event_times[0] < tW[1]):
                    sweep_epoch = e

            #Get evoked spike counts centered around stimulation times
            spk_list = []
            for probei in probe_list:
                evoked_spks, bins = util.get_evoked_spike_counts(probe_unit_data[probei]['spike_times'], probe_unit_data[probei]['spike_clusters'], probe_unit_data[probei]['units'], event_times, plot_before, plot_after, time_bin)
                spk_list.append(evoked_spks)
                spk_ts_trial = bins[:-1] + time_bin/2
            evoked_activity = np.concatenate(spk_list)  #NxnTrialsxT
            data_sub = evoked_activity[neuron_indices]
            data_ro = np.ascontiguousarray(data_sub[order_by_group])

            #Get time of stimulation
            evoked_ts_list = []
            for e in event_times:
                evoked_ts_list.append(e+spk_ts_trial)
            evoked_ts = np.array(evoked_ts_list)

            #Define windows to plot PSTHs
            indy_plot1 = np.where((spk_ts_trial >= -0.1) & (spk_ts_trial < 0.5))[0]
            indy_plot2 = np.where((spk_ts_trial >= 0.5))[0]
            
            PSTH = np.nanmean(data_ro,axis=1).T
            indy_pre = np.where(spk_ts_trial < 0)[0]
            mean_baseline = np.nanmean(PSTH[indy_pre],axis=0)
            stdd_baseline = np.nanstd(PSTH[indy_pre],axis=0)
            PSTH_z = (PSTH - mean_baseline)/stdd_baseline
            T, N = PSTH.shape

            ##------------------------------------------
            #Plot PSTHs
            PSTH_z[np.isnan(PSTH_z)] = 0
            fig, axes = plt.subplots(2,1,figsize=(11,10),gridspec_kw={'height_ratios':[5,1]})
            plt.suptitle(f'{mID}, {rec_name}\nPSTH, sweep {sweep}, {amp}uA, {sweep_epoch}',y=0.96)

            ax = axes[0]
            ax.set_title('PSTH_z [-4,4] for evoked period + rebound')
            spk_plot = ax.imshow(PSTH_z[indy_plot1].T, aspect='auto',cmap='RdBu_r',vmin=-4,vmax=4)

            ts = spk_ts_trial[indy_plot1]
            i0 = np.where(ts > 0)[0][0]-1
            i1 = np.where(ts > 0.25)[0][0]-1
            xticks = np.array([0,i0,i1,len(ts)-1])
            ax.set_xticks(xticks)
            ax.set_xticklabels(np.round(ts[xticks],2))
            ax.vlines(i0,*ax.get_ylim(),color='g',lw=1)

            ax.set_yticks(ticks,minor=True)
            ax.set_yticks(boundaries,minor=False)
            ax.set_yticklabels(labels,minor=True)
            ax.set_yticklabels([],minor=False)
            ax.set_ylabel('Region')
            ax.set_xlabel('Time relative to stimuluation (s)')

            ax = axes[1]
            indy_r = np.where((run_ts > np.nanmin(event_times)) & (run_ts < np.nanmax(event_times)))[0]
            ax.plot(run_ts[indy_r]/60,run_signal[indy_r],'-k')
            # ax.set_title('Running speed for entire sweep ')
            ax.set_ylabel('Speed (cm/s)')
            ax.set_xlabel('Time relative to beginning of experiment (min)')
            if plot_pupil:
                ax2 = ax.twinx()
                ax2.set_ylabel('Pupil radius (a.u.)',color='r')
                indy_p = np.where((pupil_ts > np.nanmin(event_times)) & (pupil_ts < np.nanmax(event_times)))[0]
                ax2.plot(pupil_ts[indy_p]/60,pupil_radius[indy_p],'-r')
            plt.savefig(join(SaveDir,f'PSTH_evoked_s-{sweep}_{amp}_{sweep_epoch}_{rec_name}.png'),dpi=300,facecolor='white',bbox_inches='tight')
            usrplt.save_fig_to_pptx(fig, prs)
            plt.close()
    
            fig, axes = plt.subplots(2,1,figsize=(11,10),gridspec_kw={'height_ratios':[5,1]})
            plt.suptitle(f'{mID}, {rec_name}\nPSTH, sweep {sweep}, {amp}uA, {sweep_epoch}',y=0.96)

            ax = axes[0]
            ax.set_title('PSTH_z [-3,3] for post-rebound')
            spk_plot = ax.imshow(PSTH_z[indy_plot2].T, aspect='auto',cmap='RdBu_r',vmin=-3,vmax=3)

            ts = spk_ts_trial[indy_plot2]
            xticks = np.array([0,int(len(ts)/2),int(len(ts)-1)])
            ax.set_xticks(xticks)
            ax.set_xticklabels(np.round(ts[xticks],2))

            ax.set_yticks(ticks,minor=True)
            ax.set_yticks(boundaries,minor=False)
            ax.set_yticklabels(labels,minor=True)
            ax.set_yticklabels([],minor=False)
            ax.set_ylabel('Region')
            ax.set_xlabel('Time relative to stimulation (s)')

            ax = axes[1]
            indy_r = np.where((run_ts > np.nanmin(event_times)) & (run_ts < np.nanmax(event_times)))[0]
            ax.plot(run_ts[indy_r]/60,run_signal[indy_r],'-k')
            ax.set_ylabel('Speed (cm/s)')
            ax.set_xlabel('Time relative to beginning of experiment (min)')
            if plot_pupil:
                ax2 = ax.twinx()
                ax2.set_ylabel('Pupil radius (a.u.)',color='r')
                indy_p = np.where((pupil_ts > np.nanmin(event_times)) & (pupil_ts < np.nanmax(event_times)))[0]
                ax2.plot(pupil_ts[indy_p]/60,pupil_radius[indy_p],'-r')
            plt.savefig(join(SaveDir,f'PSTH_ISI_s-{sweep}_{amp}_{sweep_epoch}_{rec_name}.png'),dpi=300,facecolor='white',bbox_inches='tight')
            usrplt.save_fig_to_pptx(fig, prs)
            plt.close()

            # if amp != np.max(np.unique(stim_log_b['parameter'])):
            #     print('Only calculate CCG on largest amplitude')
            #     continue

            ##------------------------------------------
            #Calculate CCGs on different windows
            #Define windows to calculate CCGs
            indy_peak = np.where((spk_ts_trial > 0) & (spk_ts_trial < 0.1))[0]
            indy_rebound = np.where((spk_ts_trial >= 0.1) & (spk_ts_trial < 0.3))[0]
            indy_post_rebound = np.where((spk_ts_trial >= 0.3) & (spk_ts_trial < 1))[0]
            indy_ISI = np.where((spk_ts_trial >= 1) & (spk_ts_trial < plot_after))[0]

            #Loop over windows
            window_strs = ['peak','rebound','post_rebound','ISI']
            # for iW, indy in enumerate([indy_peak, indy_rebound, indy_post_rebound,indy_ISI]):

            iW = 0; indy = indy_peak
            wstr = window_strs[iW]
            print(f'Calculating CCG for "{wstr}" window')
            if iW == 0:
                window = 50
            else:
                window = ccg_window
            #Get spiking data
            data_trial = np.ascontiguousarray(data_ro[:,:,indy],dtype=int)

            N, nTrials, nBins = data_trial.shape
            FR_trial = np.sum(data_trial,axis=2)/(nBins*time_bin)
            FR = np.mean(FR_trial,axis=1)

            #Calculate jitter corrected CCG
            ccg = ccg_lib.CCG(num_jitter=num_jitter, L=25, window=window, memory=False, use_parallel=True, num_cores=20)
            ccg_jitter_corrected, ccg_uncorrected = ccg.calculate_mean_ccg_corrected(data_trial, disable=False)

            #Find significant CCGs
            connection_detection = ccg_lib.SharpPeakIntervalDetection(max_duration=9, maxlag=10, n=4)
            significant_ccg,significant_confidence,significant_offset,significant_duration = connection_detection.get_significant_ccg(ccg_jitter_corrected)
            
            #Save results
            np.savez(join(SaveDir,f'ccg_{wstr}_sweep-{sweep}_{amp}_window-{window}.npz'),FR=FR,ccg_uncorrected=ccg_uncorrected,ccg_jitter_corrected=ccg_jitter_corrected,significant_ccg=significant_ccg,
                     significant_confidence=significant_confidence,significant_offset=significant_offset,significant_duration=significant_duration,
                     areas_ro=areas_ro,sweep_epoch=sweep_epoch,spk_ts_trial=spk_ts_trial[indy])

            #Plot CCG matrix
            title_str = f'{mID}, {rec_name}, {window}ms window\nCCG, {wstr} window, sweep {sweep}, {amp}uA, {sweep_epoch}'
            fig = plot_significant_ccg(significant_ccg, significant_offset, title_str, ticks, boundaries, labels)
            plt.savefig(join(SaveDir,f'ccg_{wstr}_sweep-{sweep}_{amp}_window-{window}.png'),dpi=300,facecolor='white',bbox_inches='tight')
            usrplt.save_fig_to_pptx(fig, prs)
            plt.close()

    prs.save(join(SaveDir,f'CCG-plots_{mID}_{rec_name}.pptx'))
