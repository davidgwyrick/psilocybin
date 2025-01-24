base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
base_dir_local = '/data/tiny-blue-dot/zap-n-zip/EEG_exp'; base_dir = base_dir_server
ProjDir = '/home/david.wyrick/projects/zap-n-zip/'
ServDir = '/data/projects/zap-n-zip/'
BehavDir = '/home/david.wyrick/projects/zap-n-zip/plots/behavior'

#Basea
import gspread
from os.path import join
from glob import glob
import json, os, time, sys, argparse
import numpy as np
import pandas as pd
import scipy.stats as st
import scipy.signal as sig
import pingouin as pg
from scipy.ndimage import gaussian_filter
from scipy.interpolate import interp1d
from itertools import combinations

#Plot
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
from matplotlib.gridspec import GridSpec
from matplotlib.lines import Line2D   

#Project
from tbd_eeg.data_analysis.eegutils import EEGexp
import tbd_eeg.data_analysis.Utilities.utilities as tbd_util
from tbd_eeg.data_analysis.Utilities.behavior_movies import Movie

#Allen
from allensdk.core.mouse_connectivity_cache import MouseConnectivityCache
mcc = MouseConnectivityCache(resolution=10)
str_tree = mcc.get_structure_tree()

#Read in allen CCF
ccfsum = pd.read_csv('/home/david.wyrick/projects/zap-n-zip/data/mouse_connectivity/ontology_v2.csv')

#User
sys.path.append(ProjDir)
import util
import plotting as usrplt

#Single cell stuff
import elephant.statistics as est
from quantities import ms, s, Hz
from neo.core import SpikeTrain

#Define behavioral states
behavior_ranges = {0: [0,1], 1: [1,15], 2: [15,500]}#, 3:[30,500]}
behavior_dict = {0: 'rest (0-1cm/s)', 1: 'walk (1-15cm/s)', 2: 'run (>15cm/s)'}
behavior_strs2 = ['rest','walk','run']

behavior_strs = list(behavior_dict.values())
nBehaviors = len(behavior_strs)

#Define windows to calculate firing rate
windows_of_interest = [[.002,.025],[.075,.3],[.3,1],[1,4.5]]
window_strs = ['evoked (2-25ms)','rebound (75-300ms)','post-rebound (0.3-1s)','ISI (1-4.5s)']
nWindows = len(window_strs)

gc = gspread.service_account() # need a key file to access the account (step 2) 
sh = gc.open('Templeton-log_exp') # open the spreadsheet 

df = pd.DataFrame(sh.sheet1.get()) 
exp_table = df.T.set_index(0).T # put it in a nicely formatted dataframeexp_table.iloc[10:]
exp_table = exp_table.set_index('mouse_name')

##===== ============================ =====##
##===== Parse Command Line Arguments =====##
parser = argparse.ArgumentParser(description='single-cell-metrics')

##===== Data Options =====##
parser.add_argument('--mID',type=str, default='mouse742746',
                    help='mouse to perform analysis on')

parser.add_argument('--rec_name',type=str, default='urethane_psi_2024-07-18_11-06-20',
                    help='experiment to perform analysis on')

parser.add_argument('--time_bin_ms',type=int, default=50,
                    help='time_bin_ms')

parser.add_argument('--plot_window_s',type=int, default=30,
                    help='Seconds to plot for each raster')

if __name__ == '__main__':

    # Parse the arguments
    args = parser.parse_args()
    mID = args.mID
    rec_name = args.rec_name
    time_bin_ms = args.time_bin_ms
    time_bin = time_bin_ms/1000
    plot_window_s = args.plot_window_s
    
    #Extract row from templeton dataframe
    exp_df = exp_table.loc[exp_table.exp_name == rec_name]
    if type(exp_df) == pd.core.series.Series:
        exp_df = exp_df.to_frame().T

    drug_type = exp_df['drug'].values[0]
    stim_type = exp_df['stimulation'].values[0]
    print(f'Experiment type: {stim_type}, {drug_type}')

    #Define directories
    SaveDir = join(ServDir,'results','raster',mID,rec_name)
    if not os.path.exists(SaveDir):
        os.makedirs(SaveDir)
    PlotDir = SaveDir 

    #%% Upload the whole experiment and generate the global clock
    file_name = os.path.join(base_dir_server,mID,rec_name,'experiment1','recording1')
    exp = EEGexp(file_name, preprocess=False, make_stim_csv=False)

    #Load probe data
    probe_unit_data, probe_info, total_units, metrics_list = tbd_util.get_neuropixel_data(exp)
    probe_list = [x.replace('_sorted', '') for x in exp.experiment_data if 'probe' in x]

    #Get recording start time
    probei = probe_list[0]
    open_ephys_start = np.round(np.min(probe_unit_data[probei]['spike_times']))
    open_ephys_end = np.round(np.max(probe_unit_data[probei]['spike_times']))

    # extract the timestamps of the selected stimuli
    try:
        stim_log = pd.read_csv(exp.stimulus_log_file)

        if 'circle' in np.unique(stim_log['stim_type']):
            vis_stim_exists = True
        else:
            vis_stim_exists = False
        if 'biphasic' in np.unique(stim_log['stim_type']):
            stim_exists = True
        else:
            stim_exists = False
    except:
        stim_exists = False
        vis_stim_exists = False
        stim_log = None

    #Load behavior
    run_ts, run_signal, run_signal_s, acceleration, pupil_ts, pupil_radius, pupil_dxdt, plot_pupil = util.get_behavioral_data(exp, mID, rec_name, normalize=True)
    f_run = interp1d(run_ts,run_signal)
    f_run_s = interp1d(run_ts,run_signal_s)
    run_signal_p = f_run(pupil_ts)
    run_signal_p_s = f_run_s(pupil_ts)

    # Create a new PowerPoint presentation to save figures to
    prs = Presentation()

    #Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f'mID = {mID}; rec_name = {rec_name}'
    slide.placeholders[1].text = f'All the raster plots, in {plot_window_s} second windows'

    ##------------------------------------------
    #Get injection times and types of injection
    if drug_type in ['saline', 'psilocybin','ketanserin+psilocybin','urethane+psilocybin']:
        injection_time_windows = np.array([np.array(exp_df['First injection window'].values[0].split(','),dtype=float),
                                np.array(exp_df['Second injection window'].values[0].split(','),dtype=float)])
        injection_times = np.array([exp_df['First injection time (s)'].values[0],exp_df['Second injection time (s)'].values[0]],dtype=float)

        if drug_type in ['psilocybin','urethane+psilocybin']:
            injection_types = ['sal1','psi']
            injection_colors = sns.xkcd_palette(['dark sky blue','darkish red'])
        elif drug_type == 'saline':
            injection_types = ['sal1', 'sal2']
            injection_colors = sns.xkcd_palette(['dark sky blue','cobalt blue'])
        elif drug_type == 'ketanserin+psilocybin':
            injection_types = ['ket','psi']
            injection_colors = sns.xkcd_palette(['magenta','goldenrod'])
        inj_tuple = (injection_times,injection_types,injection_colors)
    else:
        injection_times = None
        inj_tuple = None

    #For isoflurane experiments, get iso level
    if drug_type == 'isoflurane':
        iso_level, iso_times = exp.load_analog_iso()
        iso_induction_times, iso_maintenance_times = exp.load_iso_times()
        induction_colors = sns.xkcd_palette(['light teal','teal'])
        iso_tuple = (iso_induction_times, induction_colors)
    else:
        iso_induction_times = None
        iso_tuple = None

    ##------------------------------------------
    #Define time windows for each epoch  
    if drug_type in ['saline', 'psilocybin','urethane+psilocybin','ketanserin+psilocybin']:
        macro_windows = [[open_ephys_start,injection_time_windows[0,0]]]; macro_names = ['pre_inj']
        macro_windows.append([injection_time_windows[0,1],injection_time_windows[1,0]]); macro_names.append(f'post_{injection_types[0]}_inj')
        macro_windows.append([injection_time_windows[1,1],open_ephys_end]); macro_names.append(f'post_{injection_types[1]}_inj')
        if drug_type == 'saline':
            cmap_macro = sns.xkcd_palette(['silver','dark sky blue','cobalt blue'])
        elif drug_type in ['psilocybin','urethane+psilocybin']:
            cmap_macro = sns.xkcd_palette(['silver','dark sky blue','darkish red'])
        elif drug_type == 'ketanserin+psilocybin':
            cmap_macro = sns.xkcd_palette(['silver','magenta','dusty orange'])
        
    elif drug_type == 'isoflurane':

        macro_windows = [[open_ephys_start,iso_induction_times[0]-120]]; macro_names = ['pre_iso']
        macro_windows.append([iso_induction_times[0],iso_induction_times[1]]); macro_names.append(f'iso_ind')
        macro_windows.append([iso_induction_times[1]+60*5,open_ephys_end]); macro_names.append(f'post-iso')
        cmap_macro = sns.xkcd_palette(['silver','light teal','teal'])

    elif drug_type == 'urethane':
        macro_windows = [[open_ephys_start,open_ephys_end]]; macro_names = ['urethane']
        cmap_macro = sns.xkcd_palette(['darkish purple'])

    if plot_pupil:
        if drug_type == 'urethane':
            pass
            # fig, ax = plt.subplots(figsize=(6,6))
            # for jj, tW in enumerate(macro_windows):
            #     indy = np.where((pupil_ts >= tW[0]) & (pupil_ts < tW[1]) & (~np.isnan(pupil_radius) ))[0]
            #     sns.histplot(pupil_radius[indy],element='step',stat='density',fill=False,binwidth=2,ax=ax,label=f'{macro_name_list[jj]}',color=cmap_macro[jj],lw=2,ls='-',common_norm=False)
            # ax.legend()
            # usrplt.adjust_spines(ax)
            # ax.set_xlabel('Pupil size (pixels)',fontsize=12)  

        else:
            fig, axes = plt.subplots(1,3,figsize=(12,4)); plt.suptitle(f'{mID}, {rec_name}',y=1.05)
            pupil_time_bin = 1/30
            fig2, axes2 = plt.subplots(1,2,figsize=(10,5)); plt.suptitle(f'{mID}, {rec_name}',y=1.05)
            for jj, tW in enumerate(macro_windows):
                
                #Plot pupil size during rest periods
                indy = np.where((pupil_ts >= tW[0]) & (pupil_ts < tW[1]) & (~np.isnan(pupil_radius) & (run_signal_p_s < 1)))[0]
                n = len(indy); t = (n*pupil_time_bin)/60
                sns.kdeplot(pupil_radius[indy],ax=axes[0],label=f'{macro_names[jj]}_rest: {t:.1f}min',color=cmap_macro[jj],lw=3)

                #Plot pupil size during active periods
                indy = np.where((pupil_ts >= tW[0]) & (pupil_ts < tW[1]) & (~np.isnan(pupil_radius) & (run_signal_p_s > 1)))[0]
                n = len(indy); t = (n*pupil_time_bin)/60
                sns.kdeplot(pupil_radius[indy],ax=axes[1],label=f'{macro_names[jj]}_active: {t:.1f}min',color=cmap_macro[jj],lw=3)

                #Plot pupil size during active periods
                indy = np.where((pupil_ts >= tW[0]) & (pupil_ts < tW[1]) & (~np.isnan(pupil_radius)))[0]
                n = len(indy)
                sns.kdeplot(pupil_radius[indy],ax=axes2[0],label={macro_names[jj]},color=cmap_macro[jj],lw=3)
                
                #Plot running speed 
                indy = np.where((run_ts >= tW[0]) & (run_ts < tW[1]) & (~np.isnan(run_signal)))[0]
                sns.kdeplot(run_signal[indy],ax=axes[2],label=macro_names[jj],color=cmap_macro[jj],lw=3)
                sns.kdeplot(run_signal[indy],ax=axes2[1],label=macro_names[jj],color=cmap_macro[jj],lw=3)

            # vmin = 0.1; vmax=0.4
            vmin = np.nanpercentile(pupil_radius,1); vmax = np.nanpercentile(pupil_radius,99)
            axes[0].set_title('Pupil during rest')
            axes[1].set_title('Pupil during active')
            axes2[0].set_title('Pupil')
            # axes[0].legend()
            for ax in axes[:-1]:
                usrplt.adjust_spines(ax)
                ax.set_xlabel('Normalized pupil radius')
                ax.set_xlim([vmin,vmax])
                ax.legend(loc=2)

            ax = axes2[0]
            usrplt.adjust_spines(ax)
            ax.set_xlabel('Pupil radius')
            ax.set_xlim([vmin,vmax])

            ax = axes[-1]
            vmin = np.nanpercentile(run_signal,1); vmax = np.nanpercentile(run_signal,99)
            usrplt.adjust_spines(ax)
            ax.set_xlabel('Running speed (cm/s)')
            ax.set_title('Locomotion')
            ax.set_xlim([vmin,vmax])
            # import pdb; pdb.set_trace()

            ax = axes2[1]
            usrplt.adjust_spines(ax)
            ax.set_xlabel('Running speed (cm/s)')
            ax.set_title('Locomotion')
            ax.set_xlim([vmin,vmax])
            usrplt.save_fig_to_pptx(fig2, prs)
            usrplt.save_fig_to_pptx(fig, prs)
            plt.savefig(os.path.join(PlotDir,f'pupil_distributions_{mID}_{rec_name}.pdf'),facecolor='white',dpi=300,bbox_inches='tight')

    #Get rereferenced EEG data for plotting
    bc = exp_df['EEG bad_channels'].values[0]
    if (bc == 'none'):
        bad_channels = None
        eeg_ts, eeg_data, plot_eeg = util.get_preprocessed_eeg(exp, bad_channels)
    elif (bc == 'all'):
        plot_eeg = False
        eeg_ts = np.array([np.nan])
        eeg_data = np.array([np.nan])
    else:
        bad_channels = np.array(exp_df['EEG bad_channels'].values[0].split(','),dtype=int)
        eeg_ts, eeg_data, plot_eeg = util.get_preprocessed_eeg(exp, bad_channels)
    fig, ax = plt.subplots(figsize=(10,4))

    #Plot running speed
    ax.plot(run_ts/60,run_signal,'-k',lw=1)
    ax.set_xlabel('Time (s)')
    ax.set_ylabel('Running speed (cm/s)')

    if injection_times is not None:
        for ii, t in enumerate(injection_times):
            ax.axvline(t/60,color=injection_colors[ii],ls='-',lw=2,label=f'{injection_types[ii]} injection')
        ax.legend()
    if plot_pupil:
        ax2 = ax.twinx()
        ax2.plot(pupil_ts/60,pupil_radius,'-r',lw=1)
        ax2.set_ylabel('Pupil radius (pixels)',color='r')
        plt.suptitle('Behavioral measures')
    else:
        plt.suptitle('Running speed')
    usrplt.save_fig_to_pptx(fig, prs)
    plt.savefig(os.path.join(PlotDir,f'running-trace_{mID}_{rec_name}.pdf'),facecolor='white',dpi=300,bbox_inches='tight')

    #Define time windows for plotting
    t_min = open_ephys_start + 10# np.min([open_ephys_start,np.min(run_ts),np.min(pupil_ts)])
    movie_length_min = 10
    movie_length = movie_length_min*60
    tW_starts = np.arange(np.round(t_min),np.round(open_ephys_end-movie_length),movie_length)
    tW_ends = tW_starts + movie_length
    time_window_array = np.array((tW_starts,tW_ends)).T
    time_window_list = time_window_array.tolist()
    time_window_centers = time_window_array[:,0] + plot_window_s/2
    nWindows = len(time_window_list)
    # import pdb; pdb.set_trace()
    if plot_eeg:
        eeg_tuple = (eeg_ts, eeg_data)
    else:
        eeg_tuple = None

    if plot_pupil:
        pupil_tuple = (pupil_ts, pupil_radius)
    else:
        pupil_tuple = None
    for ii, tW in enumerate(time_window_list):
        #Read in neuropixel data for spontaneous periods
        data_list, ts_list, neuron_indices, plot_tuple, FR_perblock = util.bin_spiking_data(probe_unit_data, [tW], time_bin=time_bin,fr_thresh=0)

        if len(plot_tuple) == 11:
            #Ordered by area
            boundaries_group, ticks_group, labels_group, celltypes, durations,layers, areas, groups, mesogroups, supergroups, order_by_group = plot_tuple
            sort_by_area = True
        else:
            #Ordered by probe
            boundaries_probe, ticks_probe, labels_probe = plot_tuple
            sort_by_area = False

        ts = ts_list[0]; data = data_list[0]
        T, N = data.shape
        vmax = np.ceil(np.nanpercentile(data,99))
        iStarts = np.arange(0,int(np.min([T,movie_length/time_bin])),int(plot_window_s/time_bin))
        iEnds = iStarts + int(plot_window_s/time_bin) - 1

        counter = 0
        from tqdm import tqdm
        for iS, iE in tqdm(zip(iStarts,iEnds)):
            tslice = slice(iS,iE)
            ts_sub = ts[tslice]
            tS = ts[iS]; tE = ts[iE]
            for tW_e, epoch in zip(macro_windows,macro_names):
                if (tS >= tW_e[0]) & (tE <= tW_e[1]):
                    break
            #Align eeg data to start of spike raster
            if plot_eeg:
                time_bin_e = 1/100
                eeg_ts_new = np.arange(ts_sub[0],ts_sub[-1]+time_bin-time_bin_e,time_bin_e)
                nChannels = eeg_data.shape[1]
                eeg_data_new = np.zeros((len(eeg_ts_new),nChannels))
                for i in range(nChannels):
                    f_eeg = interp1d(eeg_ts,eeg_data[:,i])
                    eeg_data_new[:,i] = f_eeg(eeg_ts_new)
                eeg_tuple = (eeg_ts_new, eeg_data_new, time_bin_e)
            else:
                eeg_tuple = None

            #Align pupil data to start of spike raster
            if plot_pupil:
                time_bin_p = 1/30
                pupil_ts_new = np.arange(ts_sub[0],ts_sub[-1]+time_bin-time_bin_p,time_bin_p)
                f_pupil = interp1d(pupil_ts,pupil_radius)
                pupil_radius_new = f_pupil(pupil_ts_new)
                pupil_tuple = (pupil_ts_new, pupil_radius_new, time_bin_p)
            else:
                pupil_tuple = None
            
            #Align running data to start of spike raster
            time_bin_r = 1/100
            run_ts_new = np.arange(ts_sub[0],ts_sub[-1]+time_bin-time_bin_r,time_bin_r)
            f_run = interp1d(run_ts,run_signal)
            run_signal_new = f_run(run_ts_new)
            run_tuple = (run_ts_new, run_signal_new, time_bin_r)

            if drug_type in ['saline', 'psilocybin','urethane+psilocybin','ketanserin+psilocybin']:
                if tS <= injection_times[0]:
                    e = 'pre-inj'
                elif (tS >= injection_times[0]) & (tS < injection_times[1]):
                    e = f'post-{injection_types[0]}-inj'
                else:
                    e = f'post-{injection_types[1]}-inj'
            elif drug_type == 'isoflurane':
                if tW[0] < iso_induction_times[0]:
                    e = 'pre-iso'
                elif (tW[0] >= iso_induction_times[0]) & (tW[0] < iso_induction_times[1]):
                    e = 'iso-ind'
                else:
                    e = 'post-iso'
            else:
                e = drug_type
                
            title = f'{mID}, {rec_name}, {e}, [{tS/60:.2f}-{tE/60:.2f}] minutes'
            fig, (axes, eeg_plot,spk_plot, run_plot, pup_plot) = usrplt.plot_raster(data[tslice,:], ts[tslice], plot_tuple, time_bin=time_bin,clims=[0,vmax], title=title, run_tuple=run_tuple, pupil_tuple=pupil_tuple, eeg_tuple=eeg_tuple,time_to_plot = plot_window_s,stim_log=stim_log)
            
            if counter == 0:
                plt.savefig(os.path.join(SaveDir,f'test.pdf'),facecolor='white',dpi=300,bbox_inches='tight')
            counter += 1

            usrplt.save_fig_to_pptx(fig, prs)
            plt.close(fig)

    prs.save(join(SaveDir,f'raster-plots_{time_bin_ms}-ms_{mID}_{rec_name}.pptx'))
    print('DONE!!!')

# %%
