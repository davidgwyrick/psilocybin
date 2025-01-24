base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
base_dir_local = '/data/tiny-blue-dot/zap-n-zip/EEG_exp'; base_dir = base_dir_server
ProjDir = '/home/david.wyrick/projects/zap-n-zip/'
ServDir = '/data/projects/zap-n-zip/'
BehavDir = '/home/david.wyrick/projects/zap-n-zip/plots/behavior'

#Basea
import gspread
from os.path import join
from glob import glob
import json, os, time, sys, argparse
import numpy as np
import pandas as pd
import scipy.stats as st
import pingouin as pg
import itertools as it
from scipy.interpolate import interp1d

#Plot
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
import seaborn as sns

#Project
from tbd_eeg.data_analysis.eegutils import EEGexp
import tbd_eeg.data_analysis.Utilities.utilities as tbd_util
from tbd_eeg.data_analysis.Utilities.behavior_movies import Movie

#Allen
from allensdk.core.mouse_connectivity_cache import MouseConnectivityCache
mcc = MouseConnectivityCache(resolution=10)
str_tree = mcc.get_structure_tree()

#Read in allen CCF
ccfsum = pd.read_csv('/home/david.wyrick/projects/zap-n-zip/data/mouse_connectivity/ontology_v2.csv')

#User
sys.path.append(ProjDir)
import util
import plotting as usrplt

gc = gspread.service_account() 
sh = gc.open('Templeton-log_exp')

df = pd.DataFrame(sh.sheet1.get()) 
exp_table = df.T.set_index(0).T 
exp_table = exp_table.set_index('mouse_name')

th_dict = {'SM-TH': ['PO','VAL','VPL','VPM','VM'], 
'VIS-TH': ['LP','LGN','LGd','LGd-co','LGd-sh','LGd-ip'],
'ANT-TH': ['AV','AMd','AMv','AD','MD','MDm','MDc','MDl','RE','RH','CM','LD', 'CL'],
'TH': ['Eth', 'IAD', 'IGL', 'IntG', 'LGv','MGd', 'MGm', 'MGv', 'PCN', 'PF', 'PIL', 'PoT', 'SGN','SPFp', 'TH','LH'],'RT': ['RT']}

##===== ============================ =====##
##===== Parse Command Line Arguments =====##
parser = argparse.ArgumentParser(description='plot_single-cell-metrics')

##===== Data Options =====##
parser.add_argument('--mID',type=str, default='mouse689239',
                    help='mouse ID')

parser.add_argument('--rec_name',type=str, default='aw_psi_2023-08-10_11-26-36',
                    help='run_num')

parser.add_argument('--run_num',type=int, default=2,
                    help='run_num')

if __name__ == '__main__':

    pval_thresh = 0.01
    # Parse the arguments
    args = parser.parse_args()
    run_num = args.run_num
    mID = args.mID
    rec_name = args.rec_name

    #Extract row from templeton dataframe
    exp_df = exp_table.loc[exp_table.exp_name == rec_name]
    if type(exp_df) == pd.core.series.Series:
        exp_df = exp_df.to_frame().T

    drug_type = exp_df['drug'].values[0]
    stim_type = exp_df['stimulation'].values[0]
    
    ##------------------------------------------
    #Save dataframes & create save directory
    BaseDir = join(ServDir,'results','FC_ccg')
    TempDir = join(ServDir,'results','FC_ccg',mID,rec_name)
    if not os.path.exists(TempDir):
        raise Exception(f'No data found for {mID}, {rec_name}')

    tmp_list = sorted(glob(join(TempDir,f'ccg_run_*')))
    if len(tmp_list) == 0:
        last_run = 0
    else:
        last_run = int(tmp_list[-1][-1])

    folder = f'ccg_run_{last_run:02d}'
    SaveDir = os.path.join(TempDir,folder)
    filelist = sorted(glob(join(SaveDir,'ccg_peak*.npz')))
    if len(filelist) == 0:
        raise Exception(f'No CCG fits found for {mID}, {rec_name}')

    ##------------------------------------------
    # Create a new powerPoint presentations to save figures to
    prs = Presentation()

    #Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f'mID = {mID}, rec_name = {rec_name}, last_run = {last_run}'
    slide.placeholders[1].text = f'Change in CCG connectivity due to {drug_type}'
    nPlots = 0

    ##------------------------------------------
    #Load data
    file_name = os.path.join(base_dir_server,mID,rec_name,'experiment1','recording1')
    exp = EEGexp(file_name, preprocess=False, make_stim_csv=False)

    ##------------------------------------------
    #Get injection times and types of injection
    if drug_type in ['saline', 'psilocybin','ketanserin+psilocybin']:
        injection_time_windows = np.array([np.array(exp_df['First injection window'].values[0].split(','),dtype=float),
                                np.array(exp_df['Second injection window'].values[0].split(','),dtype=float)])
        injection_times = np.array([exp_df['First injection time (s)'].values[0],exp_df['Second injection time (s)'].values[0]],dtype=float)

        if drug_type == 'psilocybin':
            injection_types = ['sal1','psi']
            injection_colors = sns.xkcd_palette(['dark sky blue','darkish red'])
        elif drug_type == 'saline':
            injection_types = ['sal1', 'sal2']
            injection_colors = sns.xkcd_palette(['dark sky blue','cobalt blue'])
        elif drug_type == 'ketanserin+psilocybin':
            injection_types = ['ket','psi']
            injection_colors = sns.xkcd_palette(['magenta','goldenrod'])
        inj_tuple = (injection_times,injection_types,injection_colors)
    else:
        injection_times = None
        inj_tuple = None

    #For isoflurane experiments, get iso level
    if drug_type == 'isoflurane':
        iso_level, iso_times = exp.load_analog_iso()
        iso_induction_times, iso_maintenance_times = exp.load_iso_times()
        induction_colors = sns.xkcd_palette(['light teal','teal'])
        iso_tuple = (iso_induction_times, induction_colors)
    else:
        iso_induction_times = None
        iso_tuple = None


    ##------------------------------------------
    #Load behavior
    run_ts, run_signal, run_signal_s, acceleration, pupil_ts, pupil_radius, pupil_dxdt, plot_pupil = util.get_behavioral_data(exp, mID, rec_name,normalize=True)
    run_signal[run_signal < 0] = 0
    run_signal_s[run_signal_s < 0] = 0
    f_run = interp1d(run_ts,run_signal)
    f_pupil = interp1d(pupil_ts,pupil_radius)

    ##------------------------------------------
    open_ephys_start = np.min(run_ts)
    open_ephys_end = np.max(run_ts)
    #Define time windows for each epoch  
    if drug_type in ['saline', 'psilocybin','ketanserin+psilocybin']:
        macro_windows = [[open_ephys_start,injection_time_windows[0,0]]]; macro_names = ['pre-inj']
        macro_windows.append([injection_time_windows[0,1],injection_time_windows[1,0]]); macro_names.append(f'post-{injection_types[0]}-inj')
        macro_windows.append([injection_time_windows[1,1],open_ephys_end]); macro_names.append(f'post-{injection_types[1]}-inj')
        if drug_type == 'saline':
            cmap_macro = sns.xkcd_palette(['silver','dark sky blue','cobalt blue'])
        elif drug_type == 'psilocybin':
            cmap_macro = sns.xkcd_palette(['silver','dark sky blue','darkish red'])
        elif drug_type == 'ketanserin+psilocybin':
            cmap_macro = sns.xkcd_palette(['silver','magenta','dusty orange'])
        
    elif drug_type == 'isoflurane':

        macro_windows = [[open_ephys_start,iso_induction_times[0]-120]]; macro_names = ['pre-iso']
        macro_windows.append([iso_induction_times[0],iso_induction_times[1]]); macro_names.append(f'iso-ind')
        macro_windows.append([iso_induction_times[1]+60*5,open_ephys_end]); macro_names.append(f'post-iso')
        cmap_macro = sns.xkcd_palette(['silver','light teal','teal'])

    elif drug_type == 'urethane':
        macro_windows = [[open_ephys_start,open_ephys_end]]; macro_names = ['urethane']
        cmap_macro = sns.xkcd_palette(['darkish purple'])
        
    #Load stimulus log if it exists
    try:
        stim_log = pd.read_csv(exp.stimulus_log_file)

        if ('circle' in np.unique(stim_log['stim_type'])) | ('natural_scene' in np.unique(stim_log['stim_type'])):
            vStim_exists = True
        else:
            vStim_exists = False
        if 'biphasic' in np.unique(stim_log['stim_type']):
            eStim_exists = True
            stim_log_b = stim_log.loc[stim_log.stim_type == 'biphasic']
            uniq_currents = np.array(np.unique(stim_log_b['parameter']),dtype=int)
            uniq_sweeps = np.unique(stim_log_b['sweep'])
            nSweeps = len(uniq_sweeps)
            nCurrents = len(uniq_currents)
            print(f'Number of sweeps: {nSweeps}, Number of currents: {nCurrents}')
        else:
            eStim_exists = False
            raise Exception('No stimulus log for biphasic stimulation found')

    except:
        stim_log = None
        eStim_exists = False
        vStim_exists = False
        raise Exception('No stimulus log found')
    
    ##------------------------------------------
    #Load metaresults
    r = np.load(join(SaveDir,'neuron_indices.npz'),allow_pickle=True)
    neuron_indices = r['neuron_indices']
    order_by_group = r['order_by_group']
    areas = r['areas']
    groups = r['groups']
    supergroups = r['supergroups']
    celltypes = r['celltypes']
    layers = r['layers']
    durations = r['durations']
    ticks = r['ticks']
    boundaries = r['boundaries']
    labels = r['labels']
    run_ts = r['run_ts']
    run_signal = r['run_signal']
    pupil_ts = r['pupil_ts']
    pupil_radius = r['pupil_radius']
    macro_windows = r['macro_windows']
    macro_names = r['macro_names']
    cmap_macro = r['cmap_macro']

    areas_sub = areas[neuron_indices]
    groups_sub = groups[neuron_indices]
    supergroups_sub = supergroups[neuron_indices]
    celltypes_sub = celltypes[neuron_indices]

    areas_ro = areas_sub[order_by_group]
    groups_ro = groups_sub[order_by_group]
    supergroups_ro = supergroups_sub[order_by_group]
    celltypes_ro = celltypes_sub[order_by_group]
    N = len(areas_ro)

    ##------------------------------------------
    nW = 51
    ccg_jitter_corrected = np.full((nSweeps,nCurrents,N,N,nW),np.nan)
    significant_ccg = np.full((nSweeps,nCurrents,N,N),np.nan)
    significant_confidence = np.full((nSweeps,nCurrents,N,N),np.nan)
    significant_offset = np.full((nSweeps,nCurrents,N,N),np.nan)
    significant_duration = np.full((nSweeps,nCurrents,N,N),np.nan)
    all_ccgs = np.full((nSweeps,nCurrents,N,N),np.nan)
    sweep_epoch_list = []
    for iS, sweep in enumerate(np.unique(stim_log_b['sweep'])):
        for iA, amp in enumerate(np.unique(stim_log_b['parameter'])):
            sub_df = stim_log_b.loc[(stim_log_b.sweep == sweep) & (stim_log_b.parameter == amp)]

            #Get stimulation times
            event_times = sub_df['onset'].values
            sweep_epoch = ''
            for e, tW in zip(macro_names,macro_windows):
                if (event_times[0] > tW[0]) and (event_times[0] < tW[1]):
                    sweep_epoch = e
            if iA == 0:
                sweep_epoch_list.append(sweep_epoch)
            # if amp != np.max(np.unique(stim_log_b['parameter'])):
            #     print('Only calculate CCG on largest amplitude')
            #     continue

            running_per_trial = []
            for t_stim in event_times:
                run_indy = np.where((run_ts >= t_stim) & (run_ts <= t_stim+0.1))[0]
                running_per_trial.append(np.mean(run_signal[run_indy]))

            #Load CCG results 
            fpath = join(SaveDir,f'ccg_peak_sweep-{sweep}_{amp}_window-50.npz')
            if not os.path.exists(fpath):
                print(f'File {fpath} does not exist')
                continue
            
            r = np.load(fpath,allow_pickle=True)
            ccg_jitter_corrected[iS,iA] = r['ccg_jitter_corrected']
            significant_ccg[iS,iA] = r['significant_ccg']
            significant_confidence[iS,iA] = r['significant_confidence']
            significant_offset[iS,iA] = r['significant_offset']
            significant_duration[iS,iA] = r['significant_duration']

            mean10 = np.nanmean(ccg_jitter_corrected[iS,iA][:,:,slice(0,10)],axis=2)

            all_ccgs[iS,iA] = r['significant_ccg']
            mask = np.isnan(significant_ccg[iS,iA])
            # pdb.set_trace()
            all_ccgs[iS,iA][mask] = mean10[mask]

            #Plot CCGs
            fig, axes = plt.subplots(1,2,figsize=(12,6))
            plt.suptitle(f'Sweep {sweep}, {sweep_epoch}, {amp} uA, {sweep_epoch}',y=0.88)
            tmp = significant_ccg[iS,iA].copy()
            tmp[np.isnan(tmp)] = 0
            vmax = np.nanpercentile(np.abs(tmp),99) #vmin=-4E-2,vmax=4E-2
            axes[0].set_title('Significant cross-correlograms')
            sns.heatmap(tmp,cmap='RdBu_r',center=0,vmin=-1*vmax,vmax=vmax,square=True,cbar_kws={'label':'CCG','shrink':.5},ax=axes[0])

            tmp = significant_offset[iS,iA].copy()
            tmp[np.isnan(tmp)] = 0
            axes[1].set_title('Offset of connections')
            sns.heatmap(tmp,cmap='viridis',square=True,cbar_kws={'label':'Lag (ms)','shrink':.5},ax=axes[1])

            for ax in axes:
                ax.set_yticks(ticks,minor=True)
                ax.set_yticks(boundaries,minor=False)
                ax.hlines(boundaries,*ax.get_xlim(),color='k',lw=1)
                ax.set_yticklabels(labels,minor=True)
                ax.set_yticklabels([],minor=False)
                
                ax.set_xticks(ticks,minor=True)
                ax.set_xticks(boundaries,minor=False)
                ax.vlines(boundaries,*ax.get_ylim(),color='k',lw=1)
                ax.set_xticklabels(labels,minor=True,rotation=30)
                ax.set_xticklabels([],minor=False)
                ax.set_xlabel('Source')
            axes[0].set_ylabel('Target')
            usrplt.save_fig_to_pptx(fig, prs); plt.close(fig); nPlots+=1

    cmap_sweep = []
    uniq_sweep_epoch = []
    cmap_sweep_epoch = []
    for sweep in sweep_epoch_list:
        cmap_sweep.append(cmap_macro[np.where(macro_names == sweep)[0][0]])
        if sweep not in uniq_sweep_epoch:
            uniq_sweep_epoch.append(sweep)
            cmap_sweep_epoch.append(cmap_macro[np.where(macro_names == sweep)[0][0]])
                        
    ##------------------------------------------
    adjacency = np.zeros((nSweeps,nCurrents,N,N),dtype=bool)
    sign_conn = np.zeros((nSweeps,nCurrents,N,N))
    frac_conn = np.zeros((nSweeps,nCurrents))
    pos_conn = np.zeros((nSweeps,nCurrents))
    neg_conn = np.zeros((nSweeps,nCurrents))
    mag_conn = np.zeros((nSweeps,nCurrents))
    nTotal_conn = N**2 - N
    #Look at overlap of significant connections between amplitudes & conditions
    for iS, sweep in enumerate(uniq_sweeps):
        for iA, amp in enumerate(uniq_currents):
            tmp = significant_ccg[iS,iA]
            adjacency[iS,iA] = ~np.isnan(tmp)

            sign_conn[iS,iA][tmp > 0] = 1
            sign_conn[iS,iA][tmp < 0] = -1
            nConn = np.sum(~np.isnan(tmp))

            frac_conn[iS,iA] = nConn/nTotal_conn
            pos_conn[iS,iA] = np.sum(tmp > 0)/nTotal_conn
            neg_conn[iS,iA] = np.sum(tmp < 0)/nTotal_conn
            mag_conn[iS,iA] = np.nanmean(tmp)

            print(f'Sweep {sweep}, {sweep_epoch_list[iS]:13s} {amp} uA: {nConn} connections, {frac_conn[iS,iA]:.2f} fraction, {pos_conn[iS,iA]:.3f} +, {neg_conn[iS,iA]:.3f} -, {mag_conn[iS,iA]:.2E} mean connection strength')

    fig, ax = plt.subplots(figsize=(10,4))
    plt.suptitle(f'mID = {mID}, {rec_name}')
    ax.set_title('Overall connectivity during evoked response')
    ax.plot(np.nanmean(frac_conn,axis=1),'-k',label='Fraction of all connections')
    # ax.plot(np.nanmean(pos_conn,axis=1),'-',color=usrplt.cc[1],label='Fraction of pos connections')
    # ax.plot(np.nanmean(neg_conn,axis=1),'-',color=usrplt.cc[2],label='Fraction of neg connections')
    ax.set_ylabel('Fraction of significant connections')
    ax2 = ax.twinx()
    ax2.plot(np.nanmean(mag_conn,axis=1),'--b',label='Mean connection strength')
    ax2.set_ylabel('Mean connection strength',color='b')
    for iS, sweep in enumerate(uniq_sweeps):
        ax.plot(iS,np.nanmean(frac_conn[iS]),'o',color=cmap_sweep[iS])
        # ax.plot(iS,np.nanmean(pos_conn[iS]),'o',color=cmap_sweep[iS])
        # ax.plot(iS,np.nanmean(neg_conn[iS]),'o',color=cmap_sweep[iS])
        ax2.plot(iS,np.nanmean(mag_conn[iS]),'o',color=cmap_sweep[iS])

    plt.savefig(join(SaveDir,f'connectivity_per_sweep_{rec_name}.png'),dpi=300,bbox_inches='tight',facecolor='w')
    usrplt.save_fig_to_pptx(fig, prs); plt.close(fig); nPlots+=1

    ##------------------------------------------
    #Quantify overlap
    nTotal_conn = N**2 - N

    conn_overlap_sweep = np.full((nCurrents,nSweeps,nSweeps),np.nan)
    shared_conn_mag = np.zeros((nCurrents,nSweeps,nSweeps))
    for iA, amp in enumerate(uniq_currents):
        for iS, jS in it.combinations(np.arange(nSweeps),2):
        
            adjacency_i = adjacency[iS,iA]
            adjacency_j = adjacency[jS,iA]
            adjacency_ij = (adjacency_i) & (adjacency_j)

            nConn_i = np.sum(adjacency_i)/nTotal_conn
            nConn_j = np.sum(adjacency_j)/nTotal_conn
            nConn_ij = (np.sum(adjacency_i & adjacency_j)/nTotal_conn)/(np.min([nConn_i,nConn_j]))
            conn_overlap_sweep[iA,iS,jS] = nConn_ij; conn_overlap_sweep[iA,jS,iS] = nConn_ij
            print(f'Amplitude {amp} uA, {sweep_epoch_list[iS]}  vs {sweep_epoch_list[jS]} : {nConn_i:.3f} connections, {nConn_j:.3f} connections, {nConn_ij:.3f} shared connections')
            
            ccg_i = significant_ccg[iS,iA][adjacency_ij]
            ccg_j = significant_ccg[jS,iA][adjacency_ij]
            mag_i = np.nanmean(ccg_i)
            mag_j = np.nanmean(ccg_j)
            shared_conn_mag[iA,iS,jS] = mag_j
            shared_conn_mag[iA,jS,iS] = mag_i
            # pdb.set_trace()
            print(f'\tOf shared connections, mean connection strength: {mag_i:.2E}, {mag_j:.2E}')
        fig, axes = plt.subplots(1,2,figsize=(10,5))
        plt.suptitle(f'Overlap between sweeps for {amp} uA',y=0.90)
        sns.heatmap(conn_overlap_sweep[iA],cmap='viridis',square=True,ax=axes[0],cbar_kws={'label':'Fraction of shared connections','shrink':.5})
        sns.heatmap(shared_conn_mag[iA],cmap='RdBu_r',center=0,square=True,ax=axes[1],cbar_kws={'label':'Mean connection strength','shrink':.5})
        axes[0].set_title('Overlap of significant connections')
        axes[1].set_title('Mean connection strength\n of overlapping connections')
        axes[0].set_yticklabels(sweep_epoch_list,rotation=0)
        axes[0].set_xticklabels(sweep_epoch_list,rotation=30)
        axes[1].set_xticklabels(sweep_epoch_list,rotation=30)
    plt.savefig(join(SaveDir,f'overlap_per_sweep_{rec_name}.png'),dpi=300,bbox_inches='tight',facecolor='w')
    usrplt.save_fig_to_pptx(fig, prs); plt.close(fig); nPlots+=1


    cbs = it.combinations(sweep_epoch_list,2)
    ijs = it.combinations(np.arange(len(sweep_epoch_list)),2)
    combos_estim = []
    for cb,ij in zip(cbs,ijs):
        epoch_i = cb[0];i = ij[0]
        epoch_j = cb[1];j = ij[1]
        if ('pre' in epoch_i) | ('post_sal1' in epoch_i) | ('post_ket' in epoch_i):
            if ('post' in epoch_j):
                combos_estim.append(ij)
    nComps_estim = len(combos_estim)

    for iA, amp in enumerate(uniq_currents):
        tmp_list = []; yerr = []
        for iS, jS in combos_estim:
            adjacency_i = adjacency[iS,iA]
            adjacency_j = adjacency[jS,iA]
            adjacency_ij = (adjacency_i) & (adjacency_j)

            ccg_i = significant_ccg[iS,iA][adjacency_ij] #all_ccgs
            ccg_j = significant_ccg[jS,iA][adjacency_ij] #significant_ccg

            mag_i = np.nanmean(ccg_i-ccg_i)
            mag_j = np.nanmean(ccg_j-ccg_i)

            # dMag

            if len(tmp_list) == 0:
                tmp_list.append(mag_i)
                yerr.append(0)
            tmp_list.append(mag_j)
            yerr.append(st.sem(ccg_j-ccg_i,nan_policy='omit'))
        fig, ax = plt.subplots(figsize=(10,4))

        # ax.plot(conn_overlap_sweep[iA,0],'-k')
        plt.suptitle(f'Change in CCG magnitude of shared connections with {sweep_epoch_list[0]}; {amp} uA',y=0.96)
        ax.errorbar(np.arange(len(tmp_list)),y=tmp_list,yerr=yerr,color='k')
        ax.set_ylabel('Delta-CCG magnitude')
        for iS, sweep in enumerate(uniq_sweeps):
            ax.plot(iS,tmp_list[iS],'o',color=cmap_sweep[iS],zorder=5)  
        plt.savefig(join(SaveDir,f'change_CCG_magnitude_{rec_name}.png'),dpi=300,bbox_inches='tight',facecolor='w')
        usrplt.save_fig_to_pptx(fig, prs); plt.close(fig); nPlots+=1    

    unique_sorted, uniq_indices = np.unique(groups_ro, return_index=True)
    uniq_groups_order = unique_sorted[np.argsort(uniq_indices)]
    nGroups = len(uniq_groups_order)
    for g in uniq_groups_order:
        ng = np.sum(groups_ro == g)
        print(f'{g:7s}: {ng:3d} neurons')

    ##------------------------------------------
    #Create dataframe
    ccg_gmean = np.zeros((nSweeps,nCurrents,2,nGroups,nGroups))
    lag_gmean = np.zeros((nSweeps,nCurrents,2,nGroups,nGroups))
    conn_frac = np.zeros((nSweeps,nCurrents,2,nGroups,nGroups))

    if nCurrents == 3:
        current_strs = ['low','med','high']
    else:
        current_strs = ['med']
    tmp_list = []
    for iS, sweep in enumerate(uniq_sweeps):
        for iA, amp in enumerate(uniq_currents):
            ccg = significant_ccg[iS,iA]
            lag = significant_offset[iS,iA]
            adj = adjacency[iS,iA]
            
            for i, g_i in enumerate(uniq_groups_order):
                indy_i = np.where(groups_ro == g_i)[0]
                sg_i = supergroups_ro[indy_i][0]
                for j, g_j in enumerate(uniq_groups_order):
                    indy_j = np.where(groups_ro == g_j)[0]
                    sg_j = supergroups_ro[indy_j][0]
                    sg_ij = f'{sg_i}-{sg_j}'
                    nTotal_conn = len(indy_i)*len(indy_j)

                    ccg_group = ccg[indy_i][:,indy_j].ravel()
                    lag_group = lag[indy_i][:,indy_j].ravel()
                    adj_group = adj[indy_i][:,indy_j].ravel()

                    #Calculate mean connection strength, lag, and fraction of POSITIVE connections
                    mask = ccg_group > 0
                    ccg_gmean[iS,iA,0,i,j] = np.nanmean(ccg_group[mask])
                    lag_gmean[iS,iA,0,i,j] = np.nanmean(lag_group[mask])
                    conn_frac[iS,iA,0,i,j] = np.sum(adj_group[mask])/nTotal_conn

                    tmp_list.append((iS,sweep_epoch_list[iS],amp,current_strs[iA],g_i,sg_i,g_j,sg_j,sg_ij,'pos',ccg_gmean[iS,iA,0,i,j],lag_gmean[iS,iA,0,i,j],conn_frac[iS,iA,0,i,j]))

                    #Calculate mean connection strength, lag, and fraction of NEGATIVE connections
                    mask = ccg_group < 0
                    ccg_gmean[iS,iA,1,i,j] = np.nanmean(ccg_group[mask])
                    lag_gmean[iS,iA,1,i,j] = np.nanmean(lag_group[mask])
                    conn_frac[iS,iA,1,i,j] = np.sum(adj_group[mask])/nTotal_conn

                    tmp_list.append((iS,sweep_epoch_list[iS],amp,current_strs[iA],g_i,sg_i,g_j,sg_j,sg_ij,'neg',ccg_gmean[iS,iA,1,i,j],lag_gmean[iS,iA,1,i,j],conn_frac[iS,iA,1,i,j]))

    ccg_df = pd.DataFrame(np.stack(tmp_list),columns=['sweep','epoch','amplitude','level','group_i','supergroup_i','group_j','supergroup_j','supergroup_ij','sign','ccg','lag','conn_frac'])
    ccg_df = ccg_df.astype({'amplitude':int,'ccg':float,'lag':float,'conn_frac':float})
    ccg_df['drug_type'] = drug_type
    ccg_df['mID'] = mID; ccg_df['rec_name'] = rec_name
    ccg_df.to_csv(join(SaveDir,f'ccg_group_summary_{rec_name}.csv'))

    tmp_list = []
    for iA, amp in enumerate(uniq_currents):   
        for i, g_i in enumerate(uniq_groups_order):
            indy_i = np.where(groups_ro == g_i)[0]
            sg_i = supergroups_ro[indy_i][0]
            for j, g_j in enumerate(uniq_groups_order):
                indy_j = np.where(groups_ro == g_j)[0]
                sg_j = supergroups_ro[indy_j][0]
                sg_ij = f'{sg_i}-{sg_j}'

                for s, sign in enumerate(['pos','neg']):
                    for iS, jS in combos_estim:
                        val_i = ccg_gmean[iS,iA,s,i,j]
                        val_j = ccg_gmean[jS,iA,s,i,j]
                        delta = val_j - val_i
                        mod = delta/(val_j + val_i)
                        tmp_list.append((iS,sweep_epoch_list[iS],jS,sweep_epoch_list[jS],amp,current_strs[iA],g_i,sg_i,g_j,sg_j,sg_ij,sign,'ccg',val_i,val_j,delta,mod))

                        val_i = lag_gmean[iS,iA,s,i,j]
                        val_j = lag_gmean[jS,iA,s,i,j]
                        delta = val_j - val_i
                        mod = delta/(val_j + val_i)
                        tmp_list.append((iS,sweep_epoch_list[iS],jS,sweep_epoch_list[jS],amp,current_strs[iA],g_i,sg_i,g_j,sg_j,sg_ij,sign,'lag',val_i,val_j,delta,mod))

                        val_i = conn_frac[iS,iA,s,i,j]
                        val_j = conn_frac[jS,iA,s,i,j]
                        delta = val_j - val_i
                        mod = delta/(val_j + val_i)
                        tmp_list.append((iS,sweep_epoch_list[iS],jS,sweep_epoch_list[jS],amp,current_strs[iA],g_i,sg_i,g_j,sg_j,sg_ij,sign,'conn_frac',val_i,val_j,delta,mod))

    dccg_df = pd.DataFrame(np.stack(tmp_list),columns=['sweep_1','epoch_1','sweep_2','epoch_2','amplitude','level','group_i','supergroup_i','group_j','supergroup_j','supergroup_ij','sign','stat','val_1','val_2','delta','mod'])
    dccg_df = dccg_df.astype({'sweep_1':int,'sweep_2':int,'amplitude':int,'val_1':float,'val_2':float,'delta':float,'mod':float})
    dccg_df['drug_type'] = drug_type
    dccg_df['mID'] = mID; dccg_df['rec_name'] = rec_name
    dccg_df.to_csv(join(SaveDir,f'dccg_group_summary_{rec_name}.csv'))    

    sg_interest = ['CTX-CTX','CTX-TH','TH-CTX','TH-TH']

    ##------------------------------------------
    #Plot group barplots
    stat = 'conn_frac'; sign = 'pos'
    stat_ylabel = ['Mean cross-correlation','Mean lag (ms)','Connection Fraction']
    for ylabel, stat in zip(stat_ylabel,['ccg','lag','conn_frac']):
        ymax1 = np.nanpercentile(dccg_df.loc[(dccg_df.stat == stat) & (dccg_df.sign == sign)]['val_1'],99)
        ymax2 = np.nanpercentile(dccg_df.loc[(dccg_df.stat == stat) & (dccg_df.sign == sign)]['val_1'],99)
        ymax = np.round(np.max([ymax1,ymax2]),2)

        # fig2, axes2 = plt.subplots(1,4,figsize=(16,4),sharey=True)
        # plt.suptitle(f'{mID}, {rec_name}, {ylabel}, {sign} connections',y=1.05) 
        for iS, jS in combos_estim:
            fig, axes = plt.subplots(1,4,figsize=(16,4),sharey=True)
            plt.suptitle(f'{mID}, {rec_name}, {ylabel}, {sign} connections\nSweep {iS} vs sweep {jS}',y=1.05)
            for ii, sg_ij in enumerate(sg_interest):
                ax = axes[ii]
                ax.set_title(sg_ij)
                sub_df = dccg_df.loc[(dccg_df.sweep_1 == iS) & (dccg_df.sweep_2 == jS) & (dccg_df.stat == stat) & (dccg_df.sign == sign) & (dccg_df.supergroup_ij == sg_ij)] #(ccg_df.amplitude == amp) &
                # pdb.set_trace()
                xticks = []
                for iA, amp in enumerate(uniq_currents):
                    sub_sub_df = sub_df.loc[sub_df.amplitude == amp]
                    distr1 = sub_sub_df['val_1'].values
                    distr2 = sub_sub_df['val_2'].values
                    mask = ~np.isnan(distr1) & ~np.isnan(distr2)
                    distr1 = distr1[mask]
                    distr2 = distr2[mask]
                    ax.plot([np.repeat(iA,len(distr1)),np.repeat(iA+0.4,len(distr2))],[distr1,distr2],'.-',color='k',alpha=0.25)
                    bplot = ax.boxplot([distr1,distr2],positions=[iA,iA+0.4],widths=0.3,patch_artist=False,showfliers=False)
                    # pdb.set_trace()
                    for iC, patch in enumerate(bplot['boxes']):
                        patch.set_color(cmap_sweep_epoch[iC])
                        patch.set_linewidth(1.5)
                    xticks.append(iA+0.2)
                    ax.plot([iA,iA+0.4],[np.nanmean(distr1),np.nanmean(distr2)],'.-',color='k',alpha=1,lw=2.5)
                    # axes2[ii].plot([iA,iA+0.4],[np.nanmean(distr1),np.nanmean(distr2)],'.-',color='k',alpha=1,lw=2.5)
                ax.set_xticks(xticks)
                ax.set_xticklabels(uniq_currents)
                ax.set_xlabel('Amplitude (uA)')
                ax.set_ylabel(ylabel)
                # if stat == 'ccg':
                ax.set_ylim([0,ymax])
                usrplt.adjust_spines(ax)
            plt.savefig(join(SaveDir,f'{stat}_sweep-{iS}vs{jS}_{rec_name}.png'),dpi=300,bbox_inches='tight',facecolor='w')
            usrplt.save_fig_to_pptx(fig, prs); plt.close(fig); nPlots+=1    


    for g_j in uniq_groups_order:

        fig, axes = plt.subplots(3,1,figsize=(6.5,9),sharex=True,gridspec_kw={'hspace':0.5})
        plt.suptitle(f'{mID}, {rec_name},\n Source region: {g_j}, n = {np.sum(groups_ro == g_j)}',y=0.95)
        stat_ylabel = ['Mean cross-correlation','Mean lag (ms)','Connection density']
        stat_list = ['ccg','lag','conn_frac']
        for iY, (ylabel, stat) in enumerate(zip(stat_ylabel,stat_list)):
            ax = axes[iY]; xticks = []
            for iX, g_i in enumerate(uniq_groups_order):
                sub_df = dccg_df.loc[(dccg_df.group_i == g_i) & (dccg_df.group_j == g_j) & (dccg_df.sign == 'pos') & (dccg_df.stat == stat)]
                # pdb.set_trace()
                distr1 = sub_df['val_1'].values
                distr2 = sub_df['val_2'].values
                mask = ~np.isnan(distr1) & ~np.isnan(distr2)
                distr1 = distr1[mask]
                distr2 = distr2[mask]
                ax.plot([np.repeat(iX,len(distr1)),np.repeat(iX+0.4,len(distr2))],[distr1,distr2],'.-',color='k',alpha=0.25)
                bplot = ax.boxplot([distr1,distr2],positions=[iX,iX+0.4],widths=0.3,patch_artist=False,showfliers=False)
                # pdb.set_trace()
                for iC, patch in enumerate(bplot['boxes']):
                    patch.set_color(cmap_sweep_epoch[iC])
                    patch.set_linewidth(1.5)
                xticks.append(iX+0.2)
            ax.set_xticks(xticks)
            ax.set_xticklabels(uniq_groups_order)
            if iY == 2:
                ax.set_xlabel('Target region')

            ax.set_ylabel(ylabel)
            # if stat == 'ccg':
            #     ax.set_ylim([0,0.03])
            usrplt.adjust_spines(ax)
        plt.savefig(join(SaveDir,f'all-stats_group-{g_j}_{rec_name}.png'),dpi=300,bbox_inches='tight',facecolor='w')
        usrplt.save_fig_to_pptx(fig, prs); plt.close(fig); nPlots+=1    

    ##------------------------------------------
    #Do the same but for the supergroup level "CTX" vs "TH"
    unique_sorted, uniq_indices = np.unique(supergroups_ro, return_index=True)
    uniq_sg_order = unique_sorted[np.argsort(uniq_indices)]
    nGroups = len(uniq_sg_order)
    for g in uniq_sg_order:
        ng = np.sum(supergroups_ro == g)
        print(f'{g:7s}: {ng:3d} neurons')

    ccg_gmean = np.zeros((nSweeps,nCurrents,2,nGroups,nGroups))
    lag_gmean = np.zeros((nSweeps,nCurrents,2,nGroups,nGroups))
    conn_frac = np.zeros((nSweeps,nCurrents,2,nGroups,nGroups))

    tmp_list = []
    for iS, sweep in enumerate(uniq_sweeps):
        for iA, amp in enumerate(uniq_currents):
            ccg = significant_ccg[iS,iA]
            lag = significant_offset[iS,iA]
            adj = adjacency[iS,iA]
            # sgn = sign_conn[iS,iA]
            
            for i, sg_i in enumerate(uniq_sg_order):
                indy_i = np.where(supergroups_ro == sg_i)[0]
                for j, sg_j in enumerate(uniq_sg_order):
                    indy_j = np.where(supergroups_ro == sg_j)[0]
                    sg_ij = f'{sg_i}-{sg_j}'
                    nTotal_conn = len(indy_i)*len(indy_j)

                    ccg_group = ccg[indy_i][:,indy_j].ravel()
                    lag_group = lag[indy_i][:,indy_j].ravel()
                    adj_group = adj[indy_i][:,indy_j].ravel()

                    #Calculate mean connection strength, lag, and fraction of POSITIVE connections
                    mask = ccg_group > 0
                    ccg_gmean[iS,iA,0,i,j] = np.nanmean(ccg_group[mask])
                    lag_gmean[iS,iA,0,i,j] = np.nanmean(lag_group[mask])
                    conn_frac[iS,iA,0,i,j] = np.sum(adj_group[mask])/nTotal_conn

                    tmp_list.append((iS,sweep_epoch_list[iS],amp,current_strs[iA],sg_i,sg_j,sg_ij,'pos',ccg_gmean[iS,iA,0,i,j],lag_gmean[iS,iA,0,i,j],conn_frac[iS,iA,0,i,j]))

                    #Calculate mean connection strength, lag, and fraction of NEGATIVE connections
                    mask = ccg_group < 0
                    ccg_gmean[iS,iA,1,i,j] = np.nanmean(ccg_group[mask])
                    lag_gmean[iS,iA,1,i,j] = np.nanmean(lag_group[mask])
                    conn_frac[iS,iA,1,i,j] = np.sum(adj_group[mask])/nTotal_conn

                    tmp_list.append((iS,sweep_epoch_list[iS],amp,current_strs[iA],sg_i,sg_j,sg_ij,'neg',ccg_gmean[iS,iA,1,i,j],lag_gmean[iS,iA,1,i,j],conn_frac[iS,iA,1,i,j]))

    ccg_df = pd.DataFrame(np.stack(tmp_list),columns=['sweep','epoch','amplitude','level','supergroup_i','supergroup_j','supergroup_ij','sign','ccg','lag','conn_frac'])
    ccg_df = ccg_df.astype({'amplitude':int,'ccg':float,'lag':float,'conn_frac':float})
    ccg_df['drug_type'] = drug_type
    ccg_df['mID'] = mID; ccg_df['rec_name'] = rec_name
    ccg_df.to_csv(join(SaveDir,f'ccg_supergroup_summary_{rec_name}.csv'))    



    tmp_list = []
    for iA, amp in enumerate(uniq_currents):   
        for i, sg_i in enumerate(uniq_sg_order):
            indy_i = np.where(supergroups_ro == g_i)[0]
            for j, sg_j in enumerate(uniq_sg_order):
                indy_j = np.where(supergroups_ro == g_j)[0]
                sg_ij = f'{sg_i}-{sg_j}'

                for s, sign in enumerate(['pos','neg']):
                    for iS, jS in combos_estim:
                        val_i = ccg_gmean[iS,iA,s,i,j]
                        val_j = ccg_gmean[jS,iA,s,i,j]
                        delta = val_j - val_i
                        mod = delta/(val_j + val_i)
                        tmp_list.append((iS,sweep_epoch_list[iS],jS,sweep_epoch_list[jS],amp,current_strs[iA],sg_i,sg_j,sg_ij,sign,'ccg',val_i,val_j,delta,mod))

                        val_i = lag_gmean[iS,iA,s,i,j]
                        val_j = lag_gmean[jS,iA,s,i,j]
                        delta = val_j - val_i
                        mod = delta/(val_j + val_i)
                        tmp_list.append((iS,sweep_epoch_list[iS],jS,sweep_epoch_list[jS],amp,current_strs[iA],sg_i,sg_j,sg_ij,sign,'lag',val_i,val_j,delta,mod))

                        val_i = conn_frac[iS,iA,s,i,j]
                        val_j = conn_frac[jS,iA,s,i,j]
                        delta = val_j - val_i
                        mod = delta/(val_j + val_i)
                        tmp_list.append((iS,sweep_epoch_list[iS],jS,sweep_epoch_list[jS],amp,current_strs[iA],sg_i,sg_j,sg_ij,sign,'conn_frac',val_i,val_j,delta,mod))

    dccg_df = pd.DataFrame(np.stack(tmp_list),columns=['sweep_1','epoch_1','sweep_2','epoch_2','amplitude','level','supergroup_i','supergroup_j','supergroup_ij','sign','stat','val_1','val_2','delta','mod'])
    dccg_df = dccg_df.astype({'sweep_1':int,'sweep_2':int,'amplitude':int,'val_1':float,'val_2':float,'delta':float,'mod':float})
    dccg_df['drug_type'] = drug_type
    dccg_df['mID'] = mID; dccg_df['rec_name'] = rec_name
    dccg_df.to_csv(join(SaveDir,f'dccg_supergroup_summary_{rec_name}.csv'))
    
    stat = 'conn_frac'; sign = 'pos'
    stat_ylabel = ['Mean cross-correlation','Mean lag (ms)','Connection Fraction']
    for ylabel, stat in zip(stat_ylabel,['ccg','lag','conn_frac']):

        for iS, jS in combos_estim:
            fig, axes = plt.subplots(1,4,figsize=(16,4),sharey=True)
            plt.suptitle(f'{mID}, {rec_name}, {ylabel}, {sign} connections\nSweep {iS} vs sweep {jS}',y=1.05)
            for ii, sg_ij in enumerate(sg_interest):
                ax = axes[ii]
                ax.set_title(sg_ij)
                sub_df = dccg_df.loc[(dccg_df.sweep_1 == iS) & (dccg_df.sweep_2 == jS) & (dccg_df.stat == stat) & (dccg_df.sign == sign) & (dccg_df.supergroup_ij == sg_ij)] #(ccg_df.amplitude == amp) &
                xticks = []
                for iA, amp in enumerate(uniq_currents):
                    sub_sub_df = sub_df.loc[sub_df.amplitude == amp]
                    distr1 = sub_sub_df['val_1'].values
                    distr2 = sub_sub_df['val_2'].values
                    mask = ~np.isnan(distr1) & ~np.isnan(distr2)
                    distr1 = distr1[mask]
                    distr2 = distr2[mask]
                    ax.plot([np.repeat(iA,len(distr1)),np.repeat(iA+0.4,len(distr2))],[distr1,distr2],'.-',color='k',alpha=0.25)
                    bplot = ax.boxplot([distr1,distr2],positions=[iA,iA+0.4],widths=0.3,patch_artist=False,showfliers=False)
                    # pdb.set_trace()
                    for iC, patch in enumerate(bplot['boxes']):
                        patch.set_color(cmap_sweep_epoch[iC])
                        patch.set_linewidth(1.5)
                    xticks.append(iA+0.2)
                    ax.plot([iA,iA+0.4],[np.nanmean(distr1),np.nanmean(distr2)],'.-',color='k',alpha=1,lw=2.5)
                    # axes2[ii].plot([iA,iA+0.4],[np.nanmean(distr1),np.nanmean(distr2)],'.-',color='k',alpha=1,lw=2.5)
                ax.set_xticks(xticks)
                ax.set_xticklabels(uniq_currents)
                ax.set_xlabel('Amplitude (uA)')
                ax.set_ylabel(ylabel)
                usrplt.adjust_spines(ax)
            plt.savefig(join(SaveDir,f'{stat}_sg_sweep-{iS}vs{jS}_{rec_name}.png'),dpi=300,bbox_inches='tight',facecolor='w')
            usrplt.save_fig_to_pptx(fig, prs); plt.close(fig); nPlots+=1    



    ##------------------------------------------
    prs.save(join(SaveDir,f'CCG-connectivity-plots_{rec_name}.pptx'))
    prs.save(join(BaseDir,'all',f'CCG-connectivity-plots_{mID}_{rec_name}.pptx'))

    