base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
base_dir_local = '/data/tiny-blue-dot/zap-n-zip/EEG_exp'; base_dir = base_dir_server
ProjDir = '/home/david.wyrick/projects/zap-n-zip/'
ServDir = '/data/projects/zap-n-zip/'

#Basea
import gspread
from os.path import join
from glob import glob
import json, os, time, sys, argparse
import numpy as np
import pandas as pd
import scipy.stats as st
import pingouin as pg
from scipy.ndimage import gaussian_filter
from scipy.interpolate import interp1d
from itertools import combinations

#Plot
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
from matplotlib.gridspec import GridSpec

#Project
from tbd_eeg.data_analysis.eegutils import EEGexp
import tbd_eeg.data_analysis.Utilities.utilities as tbd_util
from tbd_eeg.data_analysis.Utilities.behavior_movies import Movie

#Allen
from allensdk.core.mouse_connectivity_cache import MouseConnectivityCache
mcc = MouseConnectivityCache(resolution=10)
str_tree = mcc.get_structure_tree()

#Network 
import networkx as nx
import networkx.algorithms.community as nx_comm
from networkx.algorithms.community import greedy_modularity_communities, modularity
from networkx.algorithms.efficiency_measures import global_efficiency, local_efficiency

#Read in allen CCF
ccfsum = pd.read_csv('/home/david.wyrick/projects/zap-n-zip/data/mouse_connectivity/ontology_v2.csv')

#User
sys.path.append(ProjDir)
import util
import plotting as usrplt

gc = gspread.service_account() # need a key file to access the account (step 2) 
sh = gc.open('Templeton-log_exp') # open the spreadsheet 

df = pd.DataFrame(sh.sheet1.get()) 
exp_table = df.T.set_index(0).T # put it in a nicely formatted dataframeexp_table.iloc[10:]
exp_table = exp_table.iloc[9:-8].set_index('mouse_name')

##===== ============================ =====##
##===== Parse Command Line Arguments =====##
parser = argparse.ArgumentParser(description='FCF cross-validation analysis')

##===== Data Options =====##
parser.add_argument('--mID',type=str, default='mouse669117',
                    help='mouse to perform analysis on')

parser.add_argument('--rec_name',type=str, default='pilot_aw_psi_2023-03-30_11-37-07',
                    help='experiment to perform analysis on')

parser.add_argument('--last_run',type=int, default=None,
                    help='FCF run')

def calculate_average(mat,filename_list,b,cmap,axis=0,ax=None):
    tmp_list = []
    for jj, fname in enumerate(filename_list):
        tmp = np.nanmean(mat[jj],axis=axis)
        tmp_list.append(tmp)
    mat_j = np.array(tmp_list)
    mat_bluepill = np.mean(mat_j[:b],axis=0)
    mat_redpill = np.mean(mat_j[b:],axis=0)
    if ax is not None:
        sns.kdeplot(mat_bluepill,ax=ax,color=cmap[b-1],ls='-',label='Non-psi average',lw=2)
        sns.kdeplot(mat_redpill,ax=ax,color=cmap[-1],ls='-',label='Psi average',lw=2)
    return mat_bluepill, mat_redpill

if __name__ == '__main__':

    # Parse the arguments
    args = parser.parse_args()
    mouse_name = args.mID 
    mID = args.mID
    rec_name = args.rec_name
    last_run = args.last_run
    
    #Find folder with results inside
    TempDir = os.path.join(ServDir,'results','ccm_xval',mID)
    if last_run is None:
        tmp_list = sorted(glob(join(TempDir,f'{rec_name}_run_*')))
        last_run = int(tmp_list[-1][-1])
    folder = f'{rec_name}_run_{last_run:02d}'

    #Define save and plot directories
    SaveDir = os.path.join(TempDir,folder)
    PlotDir = os.path.join(SaveDir,'plots')
    if not os.path.exists(PlotDir):
        os.makedirs(PlotDir)

    #Load ccm parameters json
    fpath = glob(join(SaveDir,'ccm_parameters*.json'))[0]
    with open(fpath, 'r') as f:
        params = json.loads(f.read())

    time_bin_ms = params['time_bin_ms']
    time_bin = time_bin_ms/1000
    delay = params['delay']
    rand_proj = bool(params['rand_proj'])
    xval = bool(params['xval'])
    zscore = bool(params['zscore'])
    nKfold = params['nKfold']
    fr_thresh = params['fr_thresh'] 

    #Get time windows used for CCM
    data = np.load(os.path.join(SaveDir,f'time_windows.npz'))
    time_window_list = data['time_window_list'].tolist()
    filename_list = data['filename_list'].tolist()
    neuron_indices = data['neuron_indices']
    nNeurons = len(neuron_indices)

    #Create colormap for epochs
    nPsi = np.sum(['psi' in fname for fname in filename_list])
    nSal = len(filename_list) - nPsi
    cmap = np.concatenate((sns.color_palette('Blues',nSal),sns.color_palette('Reds',nPsi)))
    print(f'Time bin = {time_bin_ms}ms, z-score = {zscore}, x-val = {xval}, delay = {delay}, rand_proj = {rand_proj}, fr_thresh = {fr_thresh}Hz, nNeurons = {nNeurons}')
    
    #%% Upload the whole experiment and generate the global clock
    file_name = os.path.join(base_dir,mouse_name,rec_name,'experiment1','recording1')
    exp = EEGexp(file_name, preprocess=False, make_stim_csv=False) 

    #Load probe data
    probe_unit_data, probe_info, total_units = tbd_util.get_neuropixel_data(exp)
    probe_list = [x.replace('_sorted', '') for x in exp.experiment_data if 'probe' in x]

    #Bin spiking data
    data_list, ts_list, _, plot_tuple, FR_perblock = util.bin_spiking_data(probe_unit_data, time_window_list, time_bin=time_bin,fr_thresh=fr_thresh)
    boundaries_group, ticks_group, labels_group, celltypes, durations, areas, groups, supergroups, order_by_group = plot_tuple
    boundaries_probe, ticks_probe, labels_probe = util.get_probe_plotting_info(probe_unit_data,neuron_indices)
    boundaries = boundaries_group; ticks = ticks_group; labels = labels_group
    
    #Get full names
    area_fullnames = []
    for a in areas:
        if a in ['null','nan']:
            area_fullnames.append(a)
        else:
            area_fullnames.append(str_tree.get_structures_by_acronym([a])[0]['name'])
            
    #Get running signal
    run_file = join(base_dir,mouse_name,rec_name,'experiment1','recording1','running_signal.npy')
    run_signal = np.load(run_file)

    ts_file = join(base_dir,mouse_name,rec_name,'experiment1','recording1','running_timestamps.npy')
    if os.path.exists(ts_file):
        run_timestamps = np.load(ts_file)
    else:
        ts_file = join(base_dir,mouse_name,rec_name,'experiment1','recording1','running_timestamps_master_clock.npy')
        run_timestamps = np.load(ts_file)

    #Get pupil diameter
    base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
    pupil_csv = join(base_dir_server,mID,rec_name,'experiment1','recording1',f'Pupil_{rec_name}.csv')

    try:
        table = pd.read_csv(pupil_csv)
        pupil_radius = table['Largest_Radius'].values

        #Pupil master clock
        pupil_ts = Movie(filepath=exp.pupilmovie_file,
                        sync_filepath=exp.sync_file,
                        sync_channel='eyetracking'
                        ).sync_timestamps
        pupil_data = True
    except:
        pupil_ts = np.array([np.nan])
        pupil_radius =np.array([np.nan])
        pupil_data = False

    fr_list = []
    for X in data_list:
        fr_list.append(np.sum(X,axis=0)/(X.shape[0]*time_bin))
    FR_perblock = np.array(fr_list).T
    FR_perblock_sub = FR_perblock[neuron_indices]
    
    #Load FCF!!
    fsuffix = f'{rec_name}_{time_bin_ms}ms-bins_tau-{delay}'
    if zscore:
        fsuffix = f'{fsuffix}_z' 
    if xval:
        fsuffix = f'{fsuffix}_xval'
    if rand_proj:
        fsuffix = f'{fsuffix}_randproj'

    #FCF
    results = np.load(os.path.join(SaveDir,f'results_{fsuffix}_corrected.npz'))
    FCF_optimal = results['FCF'][:,:,:nNeurons,:nNeurons]
    complexity = results['complexity'][:,:,:nNeurons,:nNeurons]
    directionality = results['directionality'][:,:,:nNeurons,:nNeurons]
    correlation = results['correlation'][:,:,:nNeurons,:nNeurons]
    FCF_takens = results['FCF_takens'][:,:,:,:nNeurons,:nNeurons]
    nCond, nTakens, nKfold, N, _ = FCF_takens.shape

    #Take mean over kfolds
    FCF_optimal_mkf = np.nanmean(FCF_optimal,axis=1)
    complexity_mkf = np.nanmean(complexity,axis=1)
    directionality_mkf = np.nanmean(directionality,axis=1)
    correlation_mkf = np.nanmean(correlation,axis=1)

    #Calculate masked FCF optimal for mean kfold
    FCF_optimal_mask = FCF_optimal_mkf.copy()
    complexity_mask = complexity_mkf.copy()
    directionality_mask = directionality_mkf.copy()
    correlation_mask = correlation_mkf.copy()

    for ii, epoch in enumerate(filename_list):
        mask = FCF_optimal_mkf[ii] < 0.1
        FCF_optimal_mask[ii][mask] = np.nan
        complexity_mask[ii][mask] = np.nan
        directionality_mask[ii][mask] = np.nan
        correlation_mask[ii][mask] = np.nan

    #Calculate masked FCF optimal for each kfold
    FCF_optimal_kf_mask = FCF_optimal.copy()
    complexity_kf_mask = complexity.copy()
    directionality_kf_mask = directionality.copy()
    correlation_kf_mask = correlation.copy()

    for ii, epoch in enumerate(filename_list):
        for iK in range(nKfold):
            mask = FCF_optimal[ii,iK] < 0.1
            FCF_optimal_kf_mask[ii,iK][mask] = np.nan
            complexity_kf_mask[ii,iK][mask] = np.nan
            directionality_kf_mask[ii,iK][mask] = np.nan
            correlation_kf_mask[ii,iK][mask] = np.nan

    #Make arrays for reorder
    FCF_optimal_ro = np.zeros(FCF_optimal.shape)
    complexity_ro = np.zeros(complexity.shape)
    directionality_ro = np.zeros(directionality.shape)
    correlation_ro = np.zeros(correlation.shape)

    #Reorder by group
    for ii in range(len(filename_list)):
        for iK in range(nKfold):
            FCF_optimal_ro[ii,iK] = FCF_optimal[ii,iK][:,order_by_group][order_by_group]
            complexity_ro[ii,iK] = complexity[ii,iK][:,order_by_group][order_by_group]
            directionality_ro[ii,iK] = directionality[ii,iK][:,order_by_group][order_by_group]
            correlation_ro[ii,iK] = correlation[ii,iK][:,order_by_group][order_by_group]

    # Create a new PowerPoint presentation to save figures to
    prs = Presentation()

    #Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f'mID = {mID}; rec_name = {rec_name}'
    slide.placeholders[1].text = f'FC analysis with cross-validation'

    ##--------------------------------------------------
    nPairs = []
    perc_FCF = []
    for ii, epoch in enumerate(filename_list):
        tmp = FCF_optimal_mkf[ii].copy()
        mask = FCF_optimal_mkf[ii] > 0.1
        perc_FCF.append(np.nanpercentile(tmp[mask].ravel(),75))
        nPairs.append(np.sum(mask)/(N**2))

    fig, ax = plt.subplots(figsize=(8,5))
    ax.plot(nPairs,'-k')
    for jj, fname in enumerate(filename_list):
        ax.plot(jj,nPairs[jj],'o',color=cmap[jj])
    ax.set_xticks(np.arange(len(filename_list)))
    ax.set_xticklabels(filename_list,rotation=30)
    ax.set_ylabel('Fraction of pairs with FC > 0.1')
    ax.set_xlabel('Epoch')

    ax2 = ax.twinx()
    ax2.plot(perc_FCF,'-',color=usrplt.cc[8])
    for jj, fname in enumerate(filename_list):
        ax2.plot(jj,perc_FCF[jj],'o',color=cmap[jj])
    ax2.set_ylabel('75th percentile FCF of significant pairs',color=usrplt.cc[8])
    usrplt.save_fig_to_pptx(fig, prs)
    plt.savefig(join(PlotDir,f'FCF_sigpairs_{rec_name}.pdf'))

    ##--------------------------------------------------
    # Calculate running per kfold
    from sklearn.model_selection import KFold

    fig, axes = plt.subplots(1,2,figsize=(10,5))
    plt.suptitle(f'{mID}: {rec_name}')
    running_moments = np.zeros((len(data_list),nKfold,3))
    pupil_moments = np.zeros((len(data_list),nKfold,3))
    for ii, data in enumerate(data_list):
        T, _ = data.shape
        ts = ts_list[ii]

        k_fold = KFold(n_splits=nKfold)
        run_list = []
        pup_list = []
        for iK, (cue_indices, tar_indices) in enumerate(k_fold.split(np.arange(T))):
            #Get running during target periods
            ts_window = ts[tar_indices]
            indy = np.where((run_timestamps >= ts_window[0]) & (run_timestamps <= ts_window[-1]))[0]
            run_tmp = run_signal[indy]

            running_moments[ii,iK,0] = np.nanmean(run_tmp)
            running_moments[ii,iK,1] = np.nanstd(run_tmp)
            running_moments[ii,iK,2] = st.skew(run_tmp)

            indy = np.where((pupil_ts >= ts_window[0]) & (pupil_ts <= ts_window[-1]))[0]
            pupil_tmp = pupil_radius[indy]

            pupil_moments[ii,iK,0] = np.nanmean(pupil_tmp)
            pupil_moments[ii,iK,1] = np.nanstd(pupil_tmp)
            pupil_moments[ii,iK,2] = st.skew(pupil_tmp)
            # sns.kdeplot(run_tmp,ax=axes[0],color=cmap[ii],lw=0.75)
            # sns.kdeplot(pupil_tmp,ax=axes[1],color=cmap[ii],lw=0.75)
            run_list.append(run_tmp)
            pup_list.append(pupil_tmp)
        sns.kdeplot(np.concatenate(run_list),ax=axes[0],color=cmap[ii],lw=1.5)
        sns.kdeplot(np.concatenate(pup_list),ax=axes[1],color=cmap[ii],lw=1.5,label=filename_list[ii])
    axes[0].set_title('Locomotion')
    axes[0].set_xlabel('Running speed (cm/s)')
            
    axes[1].set_title('Pupil size')
    axes[1].set_xlabel('Pupil radius (pixels)')
    axes[1].legend()
    usrplt.save_fig_to_pptx(fig, prs)
    plt.savefig(join(PlotDir,f'behavior_distributions_{rec_name}.pdf'))

    ##--------------------------------------------------
    vmax1 = np.round(np.nanpercentile(np.abs(correlation_ro.ravel()),99),2)
    vmax2 = np.round(np.nanpercentile(np.abs(FCF_optimal_ro.ravel()),99),2)
    vmax_fcf = np.round(np.max([vmax1,vmax2]),1)

    for ii, epoch in enumerate(filename_list):
        fcf = np.nanmean(FCF_optimal_ro[ii],axis=0)
        flow = np.nanmean(directionality_ro[ii],axis=0)
        corr_mat  = np.nanmean(correlation_ro[ii],axis=0)
        takens = np.nanmean(complexity[ii],axis=0)
        fig = usrplt.plot_CCM_results(corr_mat,fcf,flow,takens,vmax_fcf=vmax_fcf,title=epoch,ticks=ticks,labels=labels,boundaries=boundaries)
    # for ii, epoch in enumerate(filename_list):
    #     fig, axes = plt.subplots(2,5,figsize=(25,10))
    #     plt.suptitle(epoch,fontsize=18,fontweight='bold')
    #     for iK in range(nKfold):
    #         if iK % 5 == 0:
    #             plot_ylabel = True
    #         else:
    #             plot_ylabel = False
    #         usrplt.visualize_matrix(FCF_optimal_ro[ii,iK],ax=axes[iK//5,iK%5],plot_ylabel=plot_ylabel,cmap='viridis',title=f'kfold {iK}',clims=[0,vmax_fcf],center=None,ticks=ticks,labels=labels,boundaries=None)
    #     usrplt.save_fig_to_pptx(fig, prs)

    # for ii, epoch in enumerate(filename_list):
    #     fig, axes = plt.subplots(1,4,figsize=(20,5))
    #     plt.suptitle(epoch,fontsize=18,fontweight='bold')
    #     #Plot mean correlation across kfolds
    #     usrplt.visualize_matrix(np.nanmean(correlation_ro[ii],axis=0),ax=axes[0],cmap='viridis',title='Mean Correlation',clims=[0,vmax_fcf],center=None,ticks=ticks,labels=labels,boundaries=None)

    #     #Plot std. dev. correlation across kfolds
    #     usrplt.visualize_matrix(np.nanstd(correlation_ro[ii],axis=0),ax=axes[1],cmap='rocket',title='Std. Dev Correlation',center=None,ticks=ticks,labels=labels,boundaries=None)

    #     #Plot mean FCF across kfolds
    #     usrplt.visualize_matrix(np.nanmean(FCF_optimal_ro[ii],axis=0),ax=axes[2],cmap='viridis',title='Mean FC',clims=[0,vmax_fcf],center=None,ticks=ticks,labels=labels,boundaries=None)

    #     #Plot mean FCF across kfolds
    #     usrplt.visualize_matrix(np.nanstd(FCF_optimal_ro[ii],axis=0),ax=axes[3],cmap='rocket',title='Std. Dev. FC',center=None,ticks=ticks,labels=labels,boundaries=None)
        
    #     usrplt.save_fig_to_pptx(fig, prs)

    ##--------------------------------------------------
    groups_sub = groups[neuron_indices]
    fcf_moments = np.zeros((nCond,nKfold,6))
    sig_elements = np.zeros((nCond,nKfold))

    #Calculate mean and percentile of FCF per epoch/kfold
    for ii in range(nCond):
        for iK in range(nKfold):
            tmp_array = FCF_optimal[ii,iK]
            fcf_moments[ii,iK,0] = np.nanmean(tmp_array)
            fcf_moments[ii,iK,1] = np.nanpercentile(tmp_array,75)
            fcf_moments[ii,iK,2] = np.nanpercentile(np.nanmean(tmp_array,axis=0),75)

            tmp_array = FCF_optimal_kf_mask[ii,iK]
            tmp = np.copy(tmp_array).ravel()
            mask = tmp > 0.1

            sig_elements[ii,iK] = np.sum(mask)/len(mask) 
            fcf_moments[ii,iK,3] = np.nanmean(tmp_array)
            fcf_moments[ii,iK,4] = np.nanpercentile(tmp_array,75)
            fcf_moments[ii,iK,5] = np.nanpercentile(np.nanmean(tmp_array,axis=0),75)
    
    for x_array, xlabel in zip((pupil_moments,running_moments),('pupil radius (pixels)','running speed (cm/s)')):
        for jj, ylabel in enumerate(['Mean FC','75th percentile FC','75th percentile FC_j']):
            fig, axes = plt.subplots(1,2,figsize=(10,5))
            plt.suptitle(f'{mID}, Trends in {ylabel}')
            for ii, data in enumerate(data_list):
                for iK in range(nKfold):
                    axes[0].scatter(x_array[ii,iK,0],fcf_moments[ii,iK,jj],marker='.',facecolors="None",edgecolors=cmap[ii])#,linewidths=2)#,label=fname)
                    axes[1].scatter(x_array[ii,iK,0],fcf_moments[ii,iK,jj+3],marker='.',facecolors="None",edgecolors=cmap[ii])#,linewidths=2)#,label=fname)

                axes[0].plot(np.nanmean(x_array[ii,:,0]),np.nanmean(fcf_moments[ii,:,jj]),'o',color=cmap[ii])#,label=fname)
                axes[1].plot(np.nanmean(x_array[ii,:,0]),np.nanmean(fcf_moments[ii,:,jj+3]),'o',color=cmap[ii])#,label=fname)

            x = x_array[:,:,0].ravel()
            y = fcf_moments[:,:,jj].ravel()
            r = st.linregress(x,y)

            x_plot = np.linspace(np.min(x)-1,np.max(x)+1)
            y_plot = r.slope*x_plot+r.intercept
            r2 = r.rvalue**2
            axes[0].plot(x_plot,y_plot,'-k',label=f'r2 = {r2:.3f}',zorder=0)
            axes[0].legend()

            x = x_array[:,:,0].ravel()
            y = fcf_moments[:,:,jj+3].ravel()
            r = st.linregress(x,y)

            x_plot = np.linspace(np.min(x)-1,np.max(x)+1)
            y_plot = r.slope*x_plot+r.intercept
            r2 = r.rvalue**2
            axes[1].plot(x_plot,y_plot,'-k',label=f'r2 = {r2:.3f}',zorder=0)
            axes[1].legend()

            ax = axes[0]
            ax.set_xlabel(f'Mean {xlabel}')
            ax.set_ylabel(f'{ylabel} of all pairs')

            ax = axes[1]
            ax.set_xlabel(f'Mean {xlabel}')
            ax.set_ylabel(f'{ylabel} of significant FC pairs (>0.1)')
            for ax in axes:
                usrplt.adjust_spines(ax)
            usrplt.save_fig_to_pptx(fig, prs)

    for x_array, xlabel in zip((pupil_moments,running_moments),('pupil radius (pixels)','running speed (cm/s)')):
        fig, ax = plt.subplots(figsize=(5,5))
        plt.suptitle(f'{mID}, Fraction of significant pairs')
        for ii, data in enumerate(data_list):
            for iK in range(nKfold):
                ax.scatter(x_array[ii,iK,0],sig_elements[ii,iK],marker='.',facecolors="None",edgecolors=cmap[ii])#,linewidths=2)#,label=fname)
            ax.plot(np.nanmean(x_array[ii,:,0]),np.nanmean(sig_elements[ii,:]),'o',color=cmap[ii])#,label=fname)

        x = x_array[:,:,0].ravel()
        y = sig_elements[:,:].ravel()
        r = st.linregress(x,y)

        x_plot = np.linspace(np.min(x)-1,np.max(x)+1)
        y_plot = r.slope*x_plot+r.intercept
        r2 = r.rvalue**2
        ax.plot(x_plot,y_plot,'-k',label=f'r2 = {r2:.3f}',zorder=0)
        ax.legend()
        ax.set_xlabel(f'Mean {xlabel}')
        ax.set_ylabel(f'Fraction of significant FC pairs (>0.1)')
        usrplt.adjust_spines(ax)

        usrplt.save_fig_to_pptx(fig, prs)

    ##--------------------------------------------------
    supergroups_sub = supergroups[neuron_indices]
    tmp_labels = np.concatenate([labels_group,['CTX','TH']])
    for ll, g in enumerate(tmp_labels):
        if ll >= len(tmp_labels)-2:
            print(f'2last: {g}')
            indy = np.where(supergroups_sub == g)[0]
        else:
            indy = np.where(groups_sub == g)[0]
        if len(indy) < 20:
            continue
        ##--------------------
        fig, ax = plt.subplots(figsize=(8,8))
        lw = 3; zo = 0; ls = '-'
        plt.suptitle(f'{mID}; Hubness for all connections\nGroup {g}; {len(indy)} neurons')
        fcf_j_sal = []
        fcf_j_psi = []
        for ii in range(nCond):
            tmp_array = FCF_optimal_mkf[ii][:,indy]
            fcf_j = np.nanmean(tmp_array,axis=0)

            if ii < nSal:
                fcf_j_sal.append(fcf_j)
            else:
                fcf_j_psi.append(fcf_j)
            sns.kdeplot(fcf_j,ax=ax,color=cmap[ii],ls=ls,label=epoch,lw=lw,zorder=zo)

        ax_in = ax.inset_axes([0.7, 0.7, 0.3, 0.3])
        mat_bluepill = np.nanmean(fcf_j_sal,axis=0)
        mat_redpill = np.nanmean(fcf_j_psi,axis=0)
        sns.kdeplot(mat_bluepill,ax=ax_in,color=cmap[nSal-1],ls='-',label='Non-psi average',lw=2)
        sns.kdeplot(mat_redpill,ax=ax_in,color=cmap[-1],ls='-',label='Psi average',lw=2)
        usrplt.adjust_spines(ax); usrplt.adjust_spines(ax_in)
        ax.set_xlabel('FCF_j')
        usrplt.save_fig_to_pptx(fig, prs)

        ##--------------------
        fig, ax = plt.subplots(figsize=(8,8))
        lw = 3; zo = 0; ls = '-'
        plt.suptitle(f'{mID}; Hubness for connections > 0.1\nGroup {g}; {len(indy)} neurons')
        fcf_j_sal = []
        fcf_j_psi = []
        for ii in range(nCond):

            tmp_array = FCF_optimal_mask[ii][:,indy]
            fcf_j = np.nanmean(tmp_array,axis=0)

            if ii < nSal:
                fcf_j_sal.append(fcf_j)
            else:
                fcf_j_psi.append(fcf_j)
            sns.kdeplot(fcf_j,ax=ax,color=cmap[ii],ls=ls,label=epoch,lw=lw,zorder=zo)

        ax_in = ax.inset_axes([0.7, 0.7, 0.3, 0.3])
        mat_bluepill = np.nanmean(fcf_j_sal,axis=0)
        mat_redpill = np.nanmean(fcf_j_psi,axis=0)
        sns.kdeplot(mat_bluepill,ax=ax_in,color=cmap[nSal-1],ls='-',label='Non-psi average',lw=2)
        sns.kdeplot(mat_redpill,ax=ax_in,color=cmap[-1],ls='-',label='Psi average',lw=2)
        usrplt.adjust_spines(ax); usrplt.adjust_spines(ax_in)
        ax.set_xlabel('FCF_j')
        usrplt.save_fig_to_pptx(fig, prs)

        ##--------------------
        fig, ax = plt.subplots(figsize=(8,8))
        lw = 3; zo = 0; ls = '-'
        plt.suptitle(f'{mID}; # of significant connections (> 0.1)\nGroup {g}; {len(indy)} neurons')
        fcf_j_sal = []
        fcf_j_psi = []
        for ii in range(nCond):

            tmp_array = FCF_optimal_mask[ii][:,indy]
            # fcf_j = np.nanmean(tmp_array,axis=0)
            frac_j = np.sum(~np.isnan(tmp_array),axis=0)

            if ii < nSal:
                fcf_j_sal.append(frac_j)
            else:
                fcf_j_psi.append(frac_j)
            sns.kdeplot(frac_j,ax=ax,color=cmap[ii],ls=ls,label=epoch,lw=lw,zorder=zo)
            print(np.max(frac_j))

        ax_in = ax.inset_axes([0.7, 0.7, 0.3, 0.3])
        mat_bluepill = np.nanmean(fcf_j_sal,axis=0)
        mat_redpill = np.nanmean(fcf_j_psi,axis=0)
        sns.kdeplot(mat_bluepill,ax=ax_in,color=cmap[nSal-1],ls='-',label='Non-psi average',lw=2)
        sns.kdeplot(mat_redpill,ax=ax_in,color=cmap[-1],ls='-',label='Psi average',lw=2)
        ax.set_xlim(xmin=0); ax_in.set_xlim(xmin=0)
        usrplt.adjust_spines(ax); usrplt.adjust_spines(ax_in)
        ax.set_xlabel('# of connections')
        usrplt.save_fig_to_pptx(fig, prs)

    np.savez(os.path.join(SaveDir,f'FCFvsBehavior_{rec_name}.npz'),pupil_moments=pupil_moments,running_moments=running_moments,fcf_moments=fcf_moments,sig_elements=sig_elements,cmap=cmap,filename_list=filename_list,labels_group=labels_group)
    
    ##--------------------------------------------------
    FC_thresh = 0.1
    areas_sub = areas[neuron_indices]
    groups_sub, group_dict, graph_order, group_order, group_order_labels, supergroups = util.determine_groups(areas_sub)
    label_list = [f[6:] for f in filename_list]

    ##----- Plot in-network vs out-network Functional connectivity per GROUP -----##
    for g in group_order:
    # g = 'OLF'
        areas_in_group = group_dict[g]
        indy = np.where(groups_sub == g)[0]
        indy_out = np.where(groups_sub != g)[0]

        tmp_list = []
        for ii in range(nCond):

            # tmp_mat = FCF_optimal_mkf[ii].copy()
            tmp_mat = FCF_optimal[ii].copy()
            
            #Within area FC 
            FC = tmp_mat[:][:,indy][:,:,indy].ravel()
            mask = FC > FC_thresh
            n2 = len(FC[mask])
            tmp_list.append((np.repeat(filename_list[ii],n2),np.repeat(g,n2),np.repeat('within',n2),FC[mask]))

            #Source area FC 
            FC = tmp_mat[:][:,indy_out][:,:,indy].ravel()
            mask = FC > FC_thresh
            n2 = len(FC[mask])
            tmp_list.append((np.repeat(filename_list[ii],n2),np.repeat(g,n2),np.repeat('source',n2),FC[mask]))

            #Target area FC 
            FC = tmp_mat[:][:,indy][:,:,indy_out].ravel()
            mask = FC > FC_thresh
            n2 = len(FC[mask])
            tmp_list.append((np.repeat(filename_list[ii],n2),np.repeat(g,n2),np.repeat('target',n2),FC[mask]))

        FC_df = pd.DataFrame(np.hstack(tmp_list).T,columns = ['epoch','group','type','FC'])
        FC_df = FC_df.astype({'epoch':str,'group':str,'type':str,'FC':float})

        fig, axes = plt.subplots(1,3,figsize=(12,4),sharey=True)
        plt.suptitle(f'Group {g}, n = {len(indy)}',y=0.995)

        sns.barplot(data=FC_df.loc[FC_df.type == 'within'],y='FC',x='epoch',hue_order=filename_list,ax=axes[0],palette=list(cmap),legend=False)
        sns.barplot(data=FC_df.loc[FC_df.type == 'source'],y='FC',x='epoch',hue_order=filename_list,ax=axes[1],palette=list(cmap),legend=False)
        sns.barplot(data=FC_df.loc[FC_df.type == 'target'],y='FC',x='epoch',hue_order=filename_list,ax=axes[2],palette=list(cmap),legend=False)
        for ax in axes:
            usrplt.adjust_spines(ax)
            ax.set_xticklabels(label_list,ha='right',rotation=30)
            ax.set_ylim([0.05,0.5])
        axes[0].set_title('Within network FC')
        axes[1].set_title('Source area FC')
        axes[2].set_title('target area FC')
        usrplt.save_fig_to_pptx(fig, prs)

    ##--------------------------------------------------
    groups_of_interest  = [l for l in labels_group if l != 'X']#     
    # position_list = []
    # tmp = np.arange(nCond)
    # xticks = []
    # for ii in range(len(groups_of_interest)):
    #     position_list.append(tmp + (nCond+2)*ii)
    #     xticks.append((nCond+2)*ii+ nCond/2)
    # positions = np.concatenate(position_list)

    # ##----- Plot in-network vs out-network Functional connectivity per source GROUP & per target Group -----##
    # FC_thresh = 0.1
    # for g in groups_of_interest:

    #     indy = np.where(groups_sub == g)[0]

    #     # positions = np.arange(nCond*len(groups_of_interest))
    #     counter = 0
    #     tmp_list = []
    #     for g2 in groups_of_interest:
    #         indy2 = np.where(groups_sub == g2)[0]
    #         for ii in range(nCond):

    #             # tmp_mat = FCF_optimal_mkf[ii].copy()
    #             tmp_mat = FCF_optimal[ii].copy()
                
    #             FC = tmp_mat[:][:,indy2][:,:,indy].ravel()
    #             mask = FC > FC_thresh
    #             n2 = len(FC[mask])
    #             tmp_list.append((np.repeat(filename_list[ii],n2),np.repeat(g2,n2),np.repeat('source',n2),FC[mask]))


    #             FC = tmp_mat[:][:,indy][:,:,indy2].ravel()
    #             mask = FC > FC_thresh
    #             n2 = len(FC[mask])
    #             tmp_list.append((np.repeat(filename_list[ii],n2),np.repeat(g2,n2),np.repeat('target',n2),FC[mask]))

    #     FC_df = pd.DataFrame(np.hstack(tmp_list).T,columns = ['epoch','group','type','FC'])
    #     FC_df = FC_df.astype({'epoch':str,'group':str,'type':str,'FC':float})

    #     fig, axes = plt.subplots(2,1,figsize=(12,8),sharey=True)
    #     plt.suptitle(f'Group {g}, n = {len(indy)}',y=0.975)

    #     sns.barplot(data=FC_df.loc[FC_df.type == 'source'],y='FC',x='group',hue='epoch',hue_order=filename_list,ax=axes[0],palette=list(cmap),legend=False)
    #     sns.barplot(data=FC_df.loc[FC_df.type == 'target'],y='FC',x='group',hue='epoch',hue_order=filename_list,ax=axes[1],palette=list(cmap),legend=False)
    #     for ax in axes:
    #         usrplt.adjust_spines(ax)
    #         ax.set_xticklabels(groups_of_interest)
    #         ax.set_ylim([0.05,0.5])
    #         ax.set_xlabel('')

    #     axes[0].set_title('Source area FC')
    #     axes[1].set_title('target area FC')
    #     usrplt.save_fig_to_pptx(fig, prs)

    ##--------------------------------------------------
    for FCF_thresh in [0,0.1]:
        
        for normalize in [False,True]:
            tmp_list = []
            for jj in range(nCond):
                for iK in range(nKfold):
                    #Define adjacency matrix 
                    mask = FCF_optimal[jj,iK] > FCF_thresh
                    adjacency_mat = np.array(mask,dtype=int)
                    weighted_FC = FCF_optimal[jj,iK].copy()
                    weighted_FC[~mask] = 0

                    #Creates graph using the data of the adjacency matrix or weighted directed graph
                    G = nx.to_networkx_graph(adjacency_mat)
                    DG = nx.DiGraph(weighted_FC)

                    i_dict = dict(DG.in_degree(weight='weight'))
                    o_dict = dict(DG.out_degree(weight='weight'))
                    
                    num_connections_in = np.sum(mask,axis=0)
                    num_connections_out = np.sum(mask,axis=1)

                    for ii, (a,g) in enumerate(zip(areas_sub,groups_sub)):
                        if g not in groups_of_interest:
                            continue
                        if normalize:
                            tmp_list.append((filename_list[jj],ii,a,g,'in',i_dict[ii]/num_connections_in[ii]))
                            tmp_list.append((filename_list[jj],ii,a,g,'out',o_dict[ii]/num_connections_out[ii]))
                        else:
                            tmp_list.append((filename_list[jj],ii,a,g,'in',i_dict[ii]))
                            tmp_list.append((filename_list[jj],ii,a,g,'out',o_dict[ii]))
                        # ratio = i_dict[ii]/o_dict[ii]
                        if o_dict[ii] != 0:
                            tmp_list.append((filename_list[jj],ii,a,g,'in/out',i_dict[ii]/o_dict[ii]))


            degree_df = pd.DataFrame(np.vstack(tmp_list),columns=['epoch','ii','area','group','type','degree'])
            degree_df = degree_df.astype({'epoch':str,'ii':int,'area':str,'group':str,'type':str,'degree':float})

            fig, axes = plt.subplots(3,1,figsize=(12,12),gridspec_kw={'hspace':0.5})
            if normalize:
                plt.suptitle(f'Average weighted in/out degree magnitudes per group\n connections exist for FC > {FCF_thresh}')
            else:
                plt.suptitle(f'Sum of weighted in/out degree magnitudes per group\n connections exist for FC > {FCF_thresh}')
            sns.barplot(data=degree_df.loc[degree_df.type == 'in'],x='group',order=groups_of_interest,y='degree',hue='epoch',palette=cmap,ax=axes[0],legend=False)
            sns.barplot(data=degree_df.loc[degree_df.type == 'out'],x='group',order=groups_of_interest,y='degree',hue='epoch',palette=cmap,ax=axes[1],legend=False)
            sns.barplot(data=degree_df.loc[degree_df.type == 'in/out'],x='group',order=groups_of_interest,y='degree',hue='epoch',palette=cmap,ax=axes[2],legend=False)

            axes[0].set_title('Weighted In-degree')
            axes[1].set_title('Weighted Out-degree')
            axes[2].set_title('Weighted in/out degree ratio')

            axes[2].hlines(1,*axes[2].get_xlim(),color='k',ls=':')
            axes[2].autoscale(tight=True)
            axes[2].set_ylim([0,3])
            for ax in axes:
                usrplt.adjust_spines(ax)
                ax.set_ylabel('Weighted degree')
                ax.set_xlabel('')
            axes[2].set_title('Degree ratio')
            axes[2].set_ylabel('Degree ratio')
            usrplt.save_fig_to_pptx(fig, prs)

    ##--------------------------------------------------
    for FCF_thresh in [0,0.1]:
        for normalize in [False,True]:
            for g in groups_of_interest:
                indy = np.where(groups_sub == g)[0]
                n = len(indy)
                fig, axes = plt.subplots(3,1,figsize=(12,12),gridspec_kw={'hspace':0.5})
                if normalize:
                    plt.suptitle(f'Source group {g}, n = {n}; Normalized weighted degree\n {mID}: connections exist for FC > {FCF_thresh}')
                else:
                    plt.suptitle(f'Source group {g}, n = {n}; Sum of weighted degrees\n {mID}: connections exist for FC > {FCF_thresh}')

                tmp_list = []
                counter = 0
                for g2 in groups_of_interest:
                    indy2 = np.where(groups_sub == g2)[0]
                    for ii in range(nCond):
                        for iK in range(nKfold):
                            # tmp_mat = FCF_optimal_kf_mask[ii,iK].copy()
                            mask = FCF_optimal[ii,iK] > FCF_thresh
                            adjacency_mat = np.array(mask,dtype=int)
                            tmp_mat = FCF_optimal[ii,iK].copy()
                            tmp_mat[~mask] = np.nan

                            in_degree = np.nansum(tmp_mat[indy2][:,indy],axis=0).ravel(); n1 = len(in_degree)
                            out_degree = np.nansum(tmp_mat[indy][:,indy2],axis=1).ravel(); n2 = len(out_degree)
                            
                            num_connections_in = np.sum(~np.isnan(tmp_mat[indy2][:,indy]),axis=0)
                            num_connections_out = np.sum(~np.isnan(tmp_mat[indy][:,indy2]),axis=1)

                            if normalize:
                                in_degree = in_degree/num_connections_in
                                out_degree = out_degree/num_connections_out
                            tmp_list.append((np.repeat(filename_list[ii],n1),np.repeat(g,n1),np.repeat(g2,n1),np.repeat('in',n1),in_degree))
                            tmp_list.append((np.repeat(filename_list[ii],n1),np.repeat(g,n1),np.repeat(g2,n1),np.repeat('out',n1),out_degree))
                            tmp_list.append((np.repeat(filename_list[ii],n1),np.repeat(g,n1),np.repeat(g2,n1),np.repeat('in/out',n1),in_degree/out_degree))

                degree_df = pd.DataFrame(np.hstack(tmp_list).T,columns=['epoch','source_group','target_group','type','degree'])
                degree_df = degree_df.astype({'epoch':str,'source_group':str,'target_group':str,'type':str,'degree':float})

                sns.barplot(data=degree_df.loc[degree_df.type == 'in'],x='target_group',order=groups_of_interest,y='degree',hue='epoch',palette=list(cmap),ax=axes[0],legend=False)
                sns.barplot(data=degree_df.loc[degree_df.type == 'out'],x='target_group',order=groups_of_interest,y='degree',hue='epoch',palette=list(cmap),ax=axes[1],legend=False)
                sns.barplot(data=degree_df.loc[degree_df.type == 'in/out'],x='target_group',order=groups_of_interest,y='degree',hue='epoch',palette=list(cmap),ax=axes[2],legend=False)

                axes[0].set_title('Weighted In-degree')
                axes[1].set_title('Weighted Out-degree')
                axes[2].set_title('Weighted in/out degree ratio')

                for ax in axes:
                    usrplt.adjust_spines(ax)
                    ax.set_ylabel('Weighted degree')
                    ax.set_xlabel('')

                ax = axes[2]
                ax.set_xlabel('Target group')
                ax.hlines(1,*ax.get_xlim(),color='k',ls=':')
                ax.autoscale(tight=True)
                ax.set_ylim([0,3])

                for ii, g2 in enumerate(groups_of_interest):
                    n = np.sum(groups_sub == g2)
                    ax.text(ii-0.25,-0.75,n)
                usrplt.save_fig_to_pptx(fig, prs)

    ##--------------------------------------------------
    fig_out, axes_out = plt.subplots(1,3,figsize=(12,4))
    group_comm_list = []
    for grp in group_order:
        indy = np.where(groups_sub == grp)[0]
        group_comm_list.append(frozenset(indy))
    
    results = []; degree_dist_list = []
    for jj in range(nCond):
        print(filename_list[jj])
        for iK in range(nKfold):

            #Define adjacency matrix 
            mask = FCF_optimal[jj,iK] > FCF_thresh
            adjacency_mat = np.array(mask,dtype=int)
            N = adjacency_mat.shape[0]

            #Define directed graph
            weighted_FC = FCF_optimal[jj,iK].copy()
            weighted_FC[~mask] = 0

            #Calculate cost of graph
            cost = np.sum(mask)/(N**2-N)
            results.append([filename_list[jj],iK,'undirected_adjmat','cost',cost])

            #Calculate connectivity strength & diversity
            conn_strength = np.nanmean(weighted_FC[mask].ravel())
            conn_diversity = np.nanvar(weighted_FC[mask].ravel())
            results.append([filename_list[jj],iK,'directed_weight','conn_strength',conn_strength])
            results.append([filename_list[jj],iK,'directed_weight','conn_diversity',conn_diversity])

            #Creates graph using the data of the adjacency matrix or weighted directed graph
            UG_adjmat = nx.to_networkx_graph(adjacency_mat)
            DG_weight = nx.DiGraph(weighted_FC)

            #Define mask for correlation as well
            tmp_corr = correlation[jj,iK].copy()
            tmp_corr[np.diag_indices(nNeurons)] = 0
            mask = tmp_corr > FCF_thresh
            tmp_corr[~mask] = 0
            UG_weight = nx.Graph(tmp_corr)

            #Calculate cost of graph
            cost = np.sum(mask)/(N**2-N)
            results.append([filename_list[jj],iK,'undirected_weight','cost',cost])

            #Calculate connectivity strength & diversity
            conn_strength = np.nanmean(tmp_corr[mask].ravel())
            conn_diversity = np.nanvar(tmp_corr[mask].ravel())
            results.append([filename_list[jj],iK,'undirected_weight','conn_strength',conn_strength])
            results.append([filename_list[jj],iK,'undirected_weight','conn_diversity',conn_diversity])

            for ii, (g_type, g) in enumerate(zip(['undirected_weight','undirected_adjmat','directed_weight'],[UG_weight,UG_adjmat,DG_weight])):
                #Get degree distribution
                degree_dist = list(dict(g.degree()).values())
                sns.kdeplot(degree_dist,ax=axes_out[ii],color=cmap[jj],label=filename_list[jj])

                #Find communities and calculate modularity
                comm_list = greedy_modularity_communities(g,weight='weight')
                mod = modularity(g,comm_list)
                results.append([filename_list[jj],iK,g_type,'modularity',mod])
                results.append([filename_list[jj],iK,g_type,'num_communities',len(comm_list)])
                
                #Calculate modulatity of group definitions of community
                mod = modularity(g,group_comm_list)
                results.append([filename_list[jj],iK,g_type,'modularity_group',mod])

                # #Plot communities
                # indy_order = []
                # boundaries_comm = []
                # ticks_comm = []
                # labels_comm = []
                # counter = 0

                # nNeurons_per_comm = [len(s) for s in comm_list]
                # nGroups_per_comm = []
                # homogeneity = []
                # for kk, s in enumerate(comm_list):
                #     n_indices = list(s)
                #     indy_order.append(n_indices)
                #     nN_per_comm = len(n_indices)
                #     counter += nN_per_comm
                #     groups_comm = groups[n_indices]
                #     nN_per_group = []
                #     nGroups_per_comm.append(len(np.unique(groups_comm)))
                #     tmp = []
                #     for g1 in groups_comm:
                #         for g2 in groups_comm:
                #             if g1 == g2:
                #                 tmp.append(1)
                #             else:
                #                 tmp.append(0)
                #     homogeneity.append(np.mean(tmp))
                    
                #     if nN_per_comm > 5:
                #         ticks_comm.append(counter-nN_per_comm/2)
                #         boundaries_comm.append(counter)
                #         labels_comm.append(kk)
                # indy_order = np.concatenate(indy_order)

                # if g_type == 'directed_weight':
                #     fig, axes = plt.subplots(1,3,figsize=(12,4),gridspec_kw={'wspace':0.25})
                #     plt.suptitle(f'FCF: {filename_list[jj]}, Graph type: {g_type}')

                #     usrplt.visualize_matrix(FCF_optimal_ro[jj,iK],ax=axes[0],cmap='viridis',title='Sorted by area',clims=[0,vmax_fcf],center=None,ticks=ticks,labels=labels,boundaries=boundaries)
                #     usrplt.visualize_matrix(FCF_optimal[jj,iK][indy_order][:,indy_order],ax=axes[1],cmap='viridis',title=f'Sorted into communities\nModularity {mod:.3f}',clims=[0,vmax_fcf],center=None,ticks=ticks_comm,labels=labels_comm,boundaries=boundaries_comm)

                #     # sns.barplot(y='community',x='N',hue='group',ax=axes[-1],data=comm_df)
                #     ax = axes[-1]
                #     ax.set_title('Community detection',fontsize=16,fontweight='bold')
                #     ax.plot(nNeurons_per_comm,'-o')
                #     ax.set_ylabel('# of neurons per community')
                #     ax.set_xlabel('Community')
                #     # usrplt.adjust_spines(ax)

                #     ax2 = ax.twinx()
                #     ax2.plot(homogeneity,'o-',color=usrplt.cc[3])
                #     ax2.set_ylabel('homogeneity of community',color=usrplt.cc[3])
                #     ax2.set_xlim([-0.1,len(boundaries_comm)+1])
                #     ax2.set_ylim([-0.1,1.1])
                #     usrplt.save_fig_to_pptx(fig, prs)

                #     tmp_list = []
                #     for kk, s in enumerate(comm_list):
                #         n_indices = list(s)
                #         nN_per_comm = len(n_indices)
                #         counter += nN_per_comm
                #         groups_comm = groups[n_indices]
                        
                #         if nN_per_comm > 5:

                #             for g1 in group_order:
                #                 nN_per_group= np.sum(groups_comm == g1)
                #                 tmp_list.append((kk,g1,np.sum(groups_comm == g1)))

                #     n_df = pd.DataFrame(np.stack(tmp_list),columns=['community','group','N'])
                #     n_df = n_df.astype({'community':str,'group':str,'N':float})

                #     fig, ax = plt.subplots(figsize=(10,4))
                #     plt.suptitle(f'FCF: {filename_list[jj]}, Graph type: {g_type}')
                #     sns.barplot(x='community',y='N',hue='group',data=n_df,ax=ax)
                #     ax.set_xlabel('Community ID')
                #     ax.set_ylabel('# of neurons')
                #     usrplt.save_fig_to_pptx(fig, prs)

                #Get degree assortativity
                #Assortativity in a network refers to the tendency of nodes to connect with other ‘similar’ nodes over ‘dissimilar’ nodes.
                #two nodes are ‘similar’ with respect to a property if they have the same value of that property
                DAC = nx.degree_assortativity_coefficient(g)
                results.append([filename_list[jj],iK,g_type,'DAC',DAC])

                #Compute graph transitivity, the fraction of all possible triangles present in G.
                trs = nx.transitivity(g)
                results.append([filename_list[jj],iK,g_type,'transitivity',trs])

                #compute the average clustering coefficient for the graph G.
                avg_coeff = nx.average_clustering(g,weight='weight')
                results.append([filename_list[jj],iK,g_type,'average_clustering',avg_coeff])

                if 'undirected' in g_type:
                    #Calculate global efficiency 
                    #The efficiency of a pair of nodes in a graph is the multiplicative inverse of the shortest path distance between the nodes.
                    #The average global efficiency of a graph is the average efficiency of all pairs of nodes
                    gE = nx.global_efficiency(g)
                    results.append([filename_list[jj],iK,g_type,'gE',gE])

                    #The local efficiency of a node in the graph is the average global efficiency of the subgraph induced by the neighbors of the node. 
                    #The average local efficiency is the average of the local efficiencies of each node
                    lE = nx.local_efficiency(g)
                    results.append([filename_list[jj],iK,g_type,'lE',lE])

                    #Calculate whether the network is "whole" or disjointed into multiple subnetworks
                    nCC = nx.number_connected_components(g)
                    results.append([filename_list[jj],iK,g_type,'nCC',nCC])
                    
                    if nx.is_connected(g):
                        #returns the average shortest path length.
                        sp = nx.average_shortest_path_length(g,weight='weight')
                        results.append([filename_list[jj],iK,g_type,'average_shortest_path_length',sp])

                        #Compute the non-randomness of graph G.
                        nr_tuple = nx.non_randomness(g,weight='weight')
                        results.append([filename_list[jj],iK,g_type,'non_randomness',nr_tuple[1]])
                
                    #Perform targeted attack
                    steps = util.targeted_attack(g.copy())
                    results.append([filename_list[jj],iK,g_type,'targeted_attack',steps/N])

                    #Perform random attack
                    tmp_list = []
                    for iRand in range(100):
                        steps = util.random_attack(g.copy())
                        tmp_list.append(steps)
                    results.append([filename_list[jj],iK,g_type,'random_attack',np.mean(tmp_list)/N])

            #Flow hierarchy is defined as the fraction of edges not participating in cycles in a directed graph
            h = nx.flow_hierarchy(DG_weight)
            results.append([filename_list[jj],iK,'directed_weight','flow_hierarchy',h])

            #Compute the reciprocity in a directed graph.
            #The reciprocity of a directed graph is defined as the ratio of the number of edges pointing in both directions to the total number of edges in the graph. Formally,
            r = nx.overall_reciprocity(DG_weight)
            results.append([filename_list[jj],iK,'directed_weight','reciprocity',r])

    for ax in axes_out:
        ax.set_xlim(xmin=0)
        usrplt.adjust_spines(ax)

    for ii, ti in enumerate(['undirected_weight','undirected_adjmat','directed_weight']):
        axes_out[ii].set_title(ti)
    usrplt.save_fig_to_pptx(fig, prs)

    network_df = pd.DataFrame(np.vstack(results),columns=['epoch','iK','g_type','algorithm','metric'])
    network_df = network_df.astype({'epoch':str,'iK':int,'g_type':str,'algorithm':str,'metric':float})
    network_df.to_csv(join(SaveDir,f'network-metrics_{mID}_{rec_name}.csv'))

    try:
        algorithms = np.unique(network_df['algorithm'])
        ls_list = ['-',':','--']
        for a in algorithms:
            fig, ax = plt.subplots(figsize=(8,5))
            ax.set_title(f'Algorithm: {a}')
            for jj, g in enumerate(['undirected_weight','undirected_adjmat','directed_weight']):
                sub_df = network_df.loc[(network_df.g_type == g) & (network_df.algorithm == a)]
                if len(sub_df) == 0:
                    continue
                sns.pointplot(x='epoch',y='metric',data=sub_df,palette=cmap)

                ys = []
                for e in filename_list:
                    ys.append(np.nanmean(sub_df.loc[sub_df.epoch == e]['metric']))
                ax.plot(ys,'k',zorder=0,ls=ls_list[jj],label=g)
                ax.set_xticklabels(filename_list,rotation=30,ha='right')
            ax.legend()
            usrplt.adjust_spines(ax)
            usrplt.save_fig_to_pptx(fig, prs)
    except:
        prs.save(join(SaveDir,f'FCF_figs_20240305_{mID}_{rec_name}.pptx'))
        exit()

    prs.save(join(SaveDir,f'FCF_figs_202400305_{mID}_{rec_name}.pptx'))
    print('DONE!!!')
    # exit()
