base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
base_dir_local = '/data/tiny-blue-dot/zap-n-zip/EEG_exp'; base_dir = base_dir_server
ProjDir = '/home/david.wyrick/projects/zap-n-zip/'
ServDir = '/data/projects/zap-n-zip/'
BehavDir = '/home/david.wyrick/projects/zap-n-zip/plots/behavior'

#Basea
import gspread
from os.path import join
from glob import glob
import json, os, time, sys, argparse
import numpy as np
import pandas as pd
import scipy.stats as st
import scipy.signal as sig
import pingouin as pg
from scipy.ndimage import gaussian_filter
from scipy.interpolate import interp1d
from itertools import combinations

#Plot
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
from matplotlib.gridspec import GridSpec
from matplotlib.lines import Line2D   

#Project
from tbd_eeg.data_analysis.eegutils import EEGexp
import tbd_eeg.data_analysis.Utilities.utilities as tbd_util
from tbd_eeg.data_analysis.Utilities.behavior_movies import Movie

#Allen
from allensdk.core.mouse_connectivity_cache import MouseConnectivityCache
mcc = MouseConnectivityCache(resolution=10)
str_tree = mcc.get_structure_tree()

#Read in allen CCF
ccfsum = pd.read_csv('/home/david.wyrick/projects/zap-n-zip/data/mouse_connectivity/ontology_v2.csv')

#User
sys.path.append(ProjDir)
import util
import plotting as usrplt


gc = gspread.service_account() # need a key file to access the account (step 2) 
sh = gc.open('Templeton-log_exp') # open the spreadsheet 

df = pd.DataFrame(sh.sheet1.get()) 
exp_table = df.T.set_index(0).T # put it in a nicely formatted dataframeexp_table.iloc[10:]
exp_table = exp_table.set_index('mouse_name')

##===== ============================ =====##
##===== Parse Command Line Arguments =====##
parser = argparse.ArgumentParser(description='Avalanche Analysis')

##===== Data Options =====##
parser.add_argument('--mID',type=str, default='mouse728449',
                    help='mouse to perform analysis on')

parser.add_argument('--rec_name',type=str, default='spont_aw_psi_2023-06-22_11-42-00',
                    help='experiment to perform analysis on')

parser.add_argument('--time_bin_ms',type=int, default=50,
                    help='time_bin_ms')

def calculate_avalanche_statistics(spks, time_bin, thresh=0.35):
    T, N = spks.shape

    #Binarize spiking data
    spks = np.array(spks > 0,dtype=int)

    #Identify threshold based on network activity
    network_activity = np.sum(spks,axis=1)
    network_thresh = np.nanpercentile(network_activity,thresh*100)

    #Identify transitions between avalanches
    active_time_bins = np.array(network_activity > network_thresh,dtype=bool)
    active_time_bins1 = np.array(network_activity > network_thresh,dtype=int)
    transitions = np.where(np.diff(active_time_bins))[0]
    transi_type = np.diff(active_time_bins1)[transitions]

    #identify first & last avalanche
    ii = np.where(transi_type == 1)[0][0]
    jj = np.where(transi_type == -1)[0][-1] + 1
    nAvalanches = len(transitions[ii:jj])/2
    assert len(transitions[ii:jj]) % 2 == 0

    #Get indices of avalanches
    indices = transitions[ii:jj].reshape(-1,2) + 1

    #Prune avalanches that are too short
    D_bin = indices[:,1] - indices[:,0]
    indices = indices[D_bin > 2]
    
    #Get avalanche length, size, and spikes
    D = (indices[:,1] - indices[:,0]) * time_bin
    S = np.array([np.sum(network_activity[iPair[0]:iPair[1]]) for iPair in indices])

    return D, S, indices

import powerlaw as pl
def fit_power_law(data):
    data_sorted = np.sort(data)
    # Function to fit truncated power law distribution
    def fit_truncated_powerlaw(data, S_min, S_max):
        truncated_data = data[(data >= S_min) & (data <= S_max)]
        # fit = st.powerlaw.fit(truncated_data)
        fit = pl.Fit(truncated_data,xmin=S_min,xmax=S_max)
        ks = fit.power_law.KS()
        tau = fit.power_law.alpha
        return fit, ks, tau

    # Convergence criterion
    def check_convergence(ks_stat, N_av):
        return ks_stat < 1 / np.sqrt(N_av)

    # Step 1: Determine Maximum Avalanche Size (S_max)
    S_max = np.max(data)

    # Initialize variables for iterative fitting
    converged = False
    S_max_iter = S_max
    N = len(data_sorted)
    iterations = 1
    S_min = data_sorted[0]

    ks_list = []
    for i in range(1,int(N/2)):
        S_max = data_sorted[-1*i]

        fit, ks_stat, tau = fit_truncated_powerlaw(data, S_min, S_max)
        ks_list.append(ks_stat)

    indy = np.argmin(np.array(ks_list))+1
    S_max = data_sorted[-1*i]
    fit, ks_stat, tau = fit_truncated_powerlaw(data, S_min, S_max)
    print(f'{indy} data points excluded')
    
    return fit, ks_stat, tau

if __name__ == '__main__':

    # Parse the arguments
    args = parser.parse_args()
    mID = args.mID
    rec_name = args.rec_name
    time_bin_ms = args.time_bin_ms
    time_bin = time_bin_ms/1000
    
    #Extract row from templeton dataframe
    exp_df = exp_table.loc[exp_table.exp_name == rec_name]
    if type(exp_df) == pd.core.series.Series:
        exp_df = exp_df.to_frame().T

    drug_type = exp_df['drug'].values[0]
    stim_type = exp_df['stimulation'].values[0]
    print(f'Experiment type: {stim_type}, {drug_type}')
    
    SaveDir = join(ServDir,'results','avalanches',mID,rec_name)
    if not os.path.exists(SaveDir):
        os.makedirs(SaveDir)   

    #%% Upload the whole experiment and generate the global clock
    file_name = os.path.join(base_dir_server,mID,rec_name,'experiment1','recording1')
    exp = EEGexp(file_name, preprocess=False, make_stim_csv=False)

    #Load probe data
    probe_unit_data, probe_info, total_units, metric_list = tbd_util.get_neuropixel_data(exp)
    probe_list = [x.replace('_sorted', '') for x in exp.experiment_data if 'probe' in x]

    #Get recording start time
    probei = probe_list[0]
    open_ephys_start = np.round(np.min(probe_unit_data[probei]['spike_times']))
    open_ephys_end = np.round(np.max(probe_unit_data[probei]['spike_times']))

    ##------------------------------------------
    #Get injection times and types of injection
    if drug_type in ['saline', 'psilocybin','ketanserin+psilocybin']:
        injection_time_windows = np.array([np.array(exp_df['First injection window'].values[0].split(','),dtype=float),
                                np.array(exp_df['Second injection window'].values[0].split(','),dtype=float)])
        injection_times = np.array([exp_df['First injection time (s)'].values[0],exp_df['Second injection time (s)'].values[0]],dtype=float)

        if drug_type == 'psilocybin':
            injection_types = ['sal1','psi']
            injection_colors = sns.xkcd_palette(['dark sky blue','darkish red'])
        elif drug_type == 'saline':
            injection_types = ['sal1', 'sal2']
            injection_colors = sns.xkcd_palette(['dark sky blue','cobalt blue'])
        elif drug_type == 'ketanserin+psilocybin':
            injection_types = ['ket','psi']
            injection_colors = sns.xkcd_palette(['magenta','goldenrod'])
        inj_tuple = (injection_times,injection_types,injection_colors)
    else:
        injection_times = None
        inj_tuple = None

    #For isoflurane experiments, get iso level
    if drug_type == 'isoflurane':
        iso_level, iso_times = exp.load_analog_iso()
        iso_induction_times, iso_maintenance_times = exp.load_iso_times()
        induction_colors = sns.xkcd_palette(['light teal','teal'])
        iso_tuple = (iso_induction_times, induction_colors)
    else:
        iso_induction_times = None
        iso_tuple = None

    ##------------------------------------------
    #Define time windows for each epoch  
    if drug_type in ['saline', 'psilocybin','ketanserin+psilocybin']:
        macro_windows = [[open_ephys_start,injection_time_windows[0,0]]]; macro_names = ['pre_inj']
        macro_windows.append([injection_time_windows[0,1],injection_time_windows[1,0]]); macro_names.append(f'post_{injection_types[0]}_inj')
        macro_windows.append([injection_time_windows[1,1],open_ephys_end]); macro_names.append(f'post_{injection_types[1]}_inj')
        if drug_type == 'saline':
            cmap_macro = sns.xkcd_palette(['silver','dark sky blue','cobalt blue'])
        elif drug_type == 'psilocybin':
            cmap_macro = sns.xkcd_palette(['silver','dark sky blue','darkish red'])
        elif drug_type == 'ketanserin+psilocybin':
            cmap_macro = sns.xkcd_palette(['silver','magenta','dusty orange'])
        
    elif drug_type == 'isoflurane':

        macro_windows = [[open_ephys_start,iso_induction_times[0]-120]]; macro_names = ['pre_iso']
        macro_windows.append([iso_induction_times[0],iso_induction_times[1]]); macro_names.append(f'iso_ind')
        macro_windows.append([iso_induction_times[1]+60*5,open_ephys_end]); macro_names.append(f'post-iso')
        cmap_macro = sns.xkcd_palette(['silver','light teal','teal'])

    elif drug_type == 'urethane':
        macro_windows = [[open_ephys_start,open_ephys_end]]; macro_names = ['urethane']
        cmap_macro = sns.xkcd_palette(['darkish purple'])

    ##------------------------------------------
    #Load behavior
    run_ts, run_signal, run_signal_s, acceleration, pupil_ts, pupil_radius, pupil_dxdt, plot_pupil = util.get_behavioral_data(exp, mID, rec_name,normalize=True)
    run_signal[run_signal < 0] = 0
    run_signal_s[run_signal_s < 0] = 0
    f_run = interp1d(run_ts,run_signal)
    f_pupil = interp1d(pupil_ts,pupil_radius)

    #Get time windows
    time_window_list = [[open_ephys_start,injection_time_windows[0,0]],[injection_time_windows[1,1],open_ephys_end]]
    epoch_list = ['pre_inj',f'post_{injection_types[-1]}_inj']

    # Create a new PowerPoint presentation to save figures to
    prs = Presentation()

    #Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f'mID = {mID}; rec_name = {rec_name}'
    slide.placeholders[1].text = f'AVALANCHE!'
    
    #Read in neuropixel data 
    data_list, ts_list, neuron_indices, plot_tuple, FR_perblock = util.bin_spiking_data(probe_unit_data, time_window_list, time_bin=time_bin,fr_thresh=0) 
    if len(plot_tuple) == 11:
        boundaries, yticks, labels, celltypes, durations, layers, areas, groups, mesogroups, supergroups, order_by_group = plot_tuple
    else:
        print('Experiment does not have area information') 
        exit()

    area_sub = areas[neuron_indices]
    groups_sub = groups[neuron_indices]
    mesogroups_sub = mesogroups[neuron_indices]
    supergroups_sub = supergroups[neuron_indices]
    areas_ro = area_sub[order_by_group]
    groups_ro = groups_sub[order_by_group]
    mesogroups_ro = mesogroups_sub[order_by_group]
    supergroups_ro = supergroups_sub[order_by_group]

    PlotDir = join(SaveDir,'plots','group')
    if not os.path.exists(PlotDir):
        os.makedirs(PlotDir)

    import criticality as cr
    perc_list = [0.2,0.3,0.35,0.4,0.5]
    time_bin_list = [5,10,15,20,25,30,40,50]

    powerlaw_list = []
    for time_bin_ms in time_bin_list:
        time_bin = time_bin_ms/1000
        print(f'Reading in data for time bin: {time_bin_ms}ms' )
        
        #Read in neuropixel data 
        data_list, ts_list, neuron_indices, plot_tuple, FR_perblock = util.bin_spiking_data(probe_unit_data, time_window_list, time_bin=time_bin,fr_thresh=0) 
        
        for g in np.unique(groups_ro):
            mask = groups_ro == g
            if np.sum(mask) < 20:
                continue
            if g in ['root','X','FT']:
                continue
            if g == 'TH':
                mask = (groups_ro == g) | (groups_ro == 'VIS-TH')

            if g == 'VIS-TH':
                continue

            print(f'Group: {g} -> {np.sum(mask)} neurons')
            for ii, (ts, data) in enumerate(zip(ts_list,data_list)):
                nWindow_min = np.diff(time_window_list[ii])[0]/60
                print(f'{epoch_list[ii]}: {nWindow_min:.2f} minutes')   

                data_s = data[:,neuron_indices]
                data_r = data_s[:,order_by_group]
                data_g = data_r[:,mask]
                nNeurons_g = np.sum(mask)

                #Calculate covariance matrix
                cov_mat = np.cov(data_g.T)
                eigenvalues = np.linalg.eigvals(cov_mat)
                max_eig = np.max(eigenvalues)
                participation_ratio = (np.sum(eigenvalues)**2/np.sum(eigenvalues**2))/nNeurons_g
                
                for perc in perc_list:
                    #Get avalanches
                    r = cr.get_avalanches(data_g.T, perc)
                    avalanche_sizes = r['S']
                    avalanche_durations = r['T']
                    nAvalanches = len(r['S'])
                    print(f'\t{nAvalanches} avalanches detected')

                    #Get powerlaw fit
                    try:
                        Result_all = cr.AV_analysis(r['S'], r['T'],pltname=f'{epoch_list[ii]}, group {g}, bin size = {time_bin_ms}ms, network_threshold = {perc}',verbose=False)
                        fig = Result_all['scaling_relation_plot']
                        fig.savefig(join(PlotDir,f'{g}_{epoch_list[ii]}_{time_bin_ms}_{perc}_avalanche.png'),dpi=300,bbox_inches='tight',facecolor='w')
                        usrplt.save_fig_to_pptx(fig,prs)
                        plt.close(fig)

                        #Extract results
                        alpha = Result_all['beta']
                        tau = Result_all['alpha']
                        beta_pred = Result_all['pre']
                        beta_fit = Result_all['fit'][0]
                        intercept = Result_all['fit'][1]
                        DCC = Result_all['df']
                        DCC_abs = Result_all['df_abs']

                        powerlaw_list.append([mID,rec_name,drug_type,epoch_list[ii],nWindow_min,time_bin_ms, 
                                            g,nNeurons_g,max_eig,participation_ratio,
                                            perc,nAvalanches,tau,alpha,beta_pred,beta_fit,intercept,DCC,DCC_abs,
                                            np.min(avalanche_sizes),np.max(avalanche_sizes),Result_all['xmin'], Result_all['xmax'],
                                            np.min(avalanche_durations),np.max(avalanche_durations),Result_all['tmin'], Result_all['tmax']])
                    except:
                        print(f'Error in fitting powerlaw for {g}, {epoch_list[ii]}, {time_bin_ms}, {perc}')
                    # np.savez(join(SaveDir,f'power_law_fit_{time_bin_ms}-ms_{perc}-threshold.npz'),avalanche_sizes=avalanche_sizes,avalanche_durations=avalanche_durations,powerlaw_list=powerlaw_list)


    pl_df = pd.DataFrame(np.stack(powerlaw_list),columns=['mID','rec_name','drug_type','epoch','epoch_length','time_bin_ms',
                                                        'group','nNeurons_g','max_eig','p_ratio','perc','nAvalanches',
                                                        'tau','alpha','beta_pred','beta_fit','intercept','DCC','DCC_abs',
                                                        'min_size','max_size','xmin','xmax',
                                                        'min_duration','max_duration','tmin','tmax'])
    pl_df = pl_df.astype({'epoch_length':float,'time_bin_ms':int,'nNeurons_g':int,'max_eig':float,'p_ratio':float,'perc':float,'nAvalanches':int,
                        'tau':float,'alpha':float,'beta_pred':float,'beta_fit':float,'intercept':float,'DCC':float,'DCC_abs':float,
                        'min_size':float,'max_size':float,'xmin':float,'xmax':float,
                        'min_duration':float,'max_duration':float,'tmin':float,'tmax':float})
    pl_df.to_csv(join(SaveDir,'powerlaw_stats_group-level.csv'),index=False)


    try: 
        for time_bin_ms in time_bin_list:
            fig, axes = plt.subplots(1,6,figsize=(18,3))
            plt.suptitle(f'Powerlaw fit for all regions; time_bin = {time_bin_ms}ms',y=1.05)
            for kk, g in enumerate(np.unique(pl_df.group)): 
                for jj, perc in enumerate(perc_list):
                    sub_df = pl_df.loc[(pl_df.group == g) & (pl_df.perc == perc) & (pl_df.time_bin_ms == time_bin_ms)]
                    if len(sub_df) != 2:
                        raise Exception('Uhhh, something went wrong')
                x_df = sub_df.loc[sub_df.epoch == epoch_list[0]]
                y_df = sub_df.loc[sub_df.epoch == epoch_list[1]]
                for ii, var in enumerate(['nAvalanches','tau','alpha','beta_pred','beta_fit','DCC']):
                    ax = axes[ii]

                    x = x_df[var].values
                    y = y_df[var].values
                    ax.scatter(x,y,s=5+jj*10,color=usrplt.cc[kk])
                    
            for ii, var in enumerate(['nAvalanches','tau','alpha','beta_pred','beta_fit','DCC']):
                ax = axes[ii]
                xylim = [np.min([ax.get_xlim()[0],ax.get_ylim()[0]]),np.max([ax.get_xlim()[1],ax.get_ylim()[1]])]
                ax.plot(xylim,xylim,'--k')
                ax.set_title(var)
                ax.set_xlabel(epoch_list[0])
                ax.set_ylabel(epoch_list[1])
                usrplt.adjust_spines(ax)
                if ii == 0:
                    handles = [Line2D([0], [0],lw=0, marker='o',color=usrplt.cc[kk],label=g) for kk, g in enumerate(np.unique(pl_df.group))]
                    ax.legend(handles=handles)
            fig.savefig(join(PlotDir,f'summary_pre-post_group-level_{time_bin_ms}-ms.png'),dpi=300,bbox_inches='tight',facecolor='w')
            usrplt.save_fig_to_pptx(fig,prs)
            plt.close(fig)
    except:
        print(f'Plotting error')

    ##------------------------------------------
    # Do the same for mesogroup level
    PlotDir = join(SaveDir,'plots','mesogroup')
    if not os.path.exists(PlotDir):
        os.makedirs(PlotDir)

    powerlaw_list = []
    for time_bin_ms in time_bin_list:
        time_bin = time_bin_ms/1000
        print(f'Reading in data for time bin: {time_bin_ms}ms' )
        #Read in neuropixel data 
        data_list, ts_list, neuron_indices, plot_tuple, FR_perblock = util.bin_spiking_data(probe_unit_data, time_window_list, time_bin=time_bin,fr_thresh=0) 
        
        for g in np.unique(mesogroups_ro):
            mask = mesogroups_ro == g
            if np.sum(mask) < 20:
                continue
            if g in ['TH_none','UNKNOWN','FT']:
                continue

            print(f'Group: {g} -> {np.sum(mask)} neurons')
            for ii, (ts, data) in enumerate(zip(ts_list,data_list)):
                nWindow_min = np.diff(time_window_list[ii])[0]/60
                print(f'{epoch_list[ii]}: {nWindow_min:.2f} minutes')   

                data_s = data[:,neuron_indices]
                data_r = data_s[:,order_by_group]
                data_g = data_r[:,mask]
                nNeurons_g = np.sum(mask)

                #Calculate covariance matrix
                cov_mat = np.cov(data_g.T)
                eigenvalues = np.linalg.eigvals(cov_mat)
                max_eig = np.max(eigenvalues)
                participation_ratio = (np.sum(eigenvalues)**2/np.sum(eigenvalues**2))/nNeurons_g
                
                for perc in perc_list:
                    #Get avalanches
                    r = cr.get_avalanches(data_g.T, perc)
                    avalanche_sizes = r['S']
                    avalanche_durations = r['T']
                    nAvalanches = len(r['S'])
                    print(f'\t{nAvalanches} avalanches detected')

                    #Get powerlaw fit
                    try:
                        Result_all = cr.AV_analysis(r['S'], r['T'],pltname=f'{epoch_list[ii]}, group {g}, bin size = {time_bin_ms}ms, network_threshold = {perc}',verbose=False)
                        fig = Result_all['scaling_relation_plot']
                        fig.savefig(join(PlotDir,f'{g}_{epoch_list[ii]}_{time_bin_ms}_{perc}_avalanche.png'),dpi=300,bbox_inches='tight',facecolor='w')
                        usrplt.save_fig_to_pptx(fig,prs)
                        plt.close(fig)

                        #Extract results
                        alpha = Result_all['beta']
                        tau = Result_all['alpha']
                        beta_pred = Result_all['pre']
                        beta_fit = Result_all['fit'][0]
                        intercept = Result_all['fit'][1]
                        DCC = Result_all['df']
                        DCC_abs = Result_all['df_abs']

                        powerlaw_list.append([mID,rec_name,drug_type,epoch_list[ii],nWindow_min,time_bin_ms, 
                                            g,nNeurons_g,max_eig,participation_ratio,
                                            perc,nAvalanches,tau,alpha,beta_pred,beta_fit,intercept,DCC,DCC_abs,
                                            np.min(avalanche_sizes),np.max(avalanche_sizes),Result_all['xmin'], Result_all['xmax'],
                                            np.min(avalanche_durations),np.max(avalanche_durations),Result_all['tmin'], Result_all['tmax']])
                    except:
                        print(f'Error in fitting powerlaw for {g}, {epoch_list[ii]}, {time_bin_ms}, {perc}')
                    # np.savez(join(SaveDir,f'power_law_fit_{time_bin_ms}-ms_{perc}-threshold.npz'),avalanche_sizes=avalanche_sizes,avalanche_durations=avalanche_durations,powerlaw_list=powerlaw_list)


    pl_df = pd.DataFrame(np.stack(powerlaw_list),columns=['mID','rec_name','drug_type','epoch','epoch_length','time_bin_ms',
                                                        'group','nNeurons_g','max_eig','p_ratio','perc','nAvalanches',
                                                        'tau','alpha','beta_pred','beta_fit','intercept','DCC','DCC_abs',
                                                        'min_size','max_size','xmin','xmax',
                                                        'min_duration','max_duration','tmin','tmax'])
    pl_df = pl_df.astype({'epoch_length':float,'time_bin_ms':int,'nNeurons_g':int,'max_eig':float,'p_ratio':float,'perc':float,'nAvalanches':int,
                        'tau':float,'alpha':float,'beta_pred':float,'beta_fit':float,'intercept':float,'DCC':float,'DCC_abs':float,
                        'min_size':float,'max_size':float,'xmin':float,'xmax':float,
                        'min_duration':float,'max_duration':float,'tmin':float,'tmax':float})
    pl_df.to_csv(join(SaveDir,'powerlaw_stats_mesogroup-level.csv'),index=False)


    try: 
        for time_bin_ms in time_bin_list:
            fig, axes = plt.subplots(1,6,figsize=(18,3))
            plt.suptitle(f'Powerlaw fit for all regions; time_bin = {time_bin_ms}ms',y=1.05)
            for kk, g in enumerate(np.unique(pl_df.group)): 
                for jj, perc in enumerate(perc_list):
                    sub_df = pl_df.loc[(pl_df.group == g) & (pl_df.perc == perc) & (pl_df.time_bin_ms == time_bin_ms)]
                    if len(sub_df) != 2:
                        raise Exception('Uhhh, something went wrong')
                x_df = sub_df.loc[sub_df.epoch == epoch_list[0]]
                y_df = sub_df.loc[sub_df.epoch == epoch_list[1]]
                for ii, var in enumerate(['nAvalanches','tau','alpha','beta_pred','beta_fit','DCC']):
                    ax = axes[ii]

                    x = x_df[var].values
                    y = y_df[var].values
                    ax.scatter(x,y,s=5+jj*10,color=usrplt.cc[kk])
                    
            for ii, var in enumerate(['nAvalanches','tau','alpha','beta_pred','beta_fit','DCC']):
                ax = axes[ii]
                xylim = [np.min([ax.get_xlim()[0],ax.get_ylim()[0]]),np.max([ax.get_xlim()[1],ax.get_ylim()[1]])]
                ax.plot(xylim,xylim,'--k')
                ax.set_title(var)
                ax.set_xlabel(epoch_list[0])
                ax.set_ylabel(epoch_list[1])
                usrplt.adjust_spines(ax)
                if ii == 0:
                    handles = [Line2D([0], [0],lw=0, marker='o',color=usrplt.cc[kk],label=g) for kk, g in enumerate(np.unique(pl_df.group))]
                    ax.legend(handles=handles)
            fig.savefig(join(PlotDir,f'summary_pre-post_group-level_{time_bin_ms}-ms.png'),dpi=300,bbox_inches='tight',facecolor='w')
            usrplt.save_fig_to_pptx(fig,prs)
            plt.close(fig)
    except:
        print(f'Plotting error')


    prs.save(join(SaveDir,f'avalanche_distributions_{mID}_{rec_name}.pptx'))
    print('DONE!!!')

# %%
