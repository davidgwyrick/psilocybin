base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
base_dir_local = '/data/tiny-blue-dot/zap-n-zip/EEG_exp'; base_dir = base_dir_server
ProjDir = '/home/david.wyrick/projects/zap-n-zip/'
ServDir = '/data/projects/zap-n-zip/'
BehavDir = '/home/david.wyrick/projects/zap-n-zip/plots/behavior'

##------------------------------------------
##Load modules
#Base
import gspread
from os.path import join
from glob import glob
import json, os, time, sys, argparse
import numpy as np
import pandas as pd
import itertools as it
import scipy.stats as st
from scipy.interpolate import interp1d
from tqdm import tqdm

#Plot
import matplotlib.pyplot as plt
from pptx import Presentation
import seaborn as sns

#Project
from tbd_eeg.data_analysis.eegutils import EEGexp
import tbd_eeg.data_analysis.Utilities.utilities as tbd_util

#Allen
from allensdk.core.mouse_connectivity_cache import MouseConnectivityCache
mcc = MouseConnectivityCache(resolution=10)
str_tree = mcc.get_structure_tree()

#User
sys.path.append(ProjDir)
import util
import plotting as usrplt

#Define behavioral states
behavior_ranges = {0: [0,1], 1: [1,500]}#, 3:[30,500]}
behavior_dict = { 0: 'rest (<1cm/s)', 1: 'active (>1cm/s)'}
behavior_dict2 = { 0: 'rest (<1cm/s)', 1: 'active (>1cm/s)', 2: 'all (>0cm/s)'}
behavior_strs2 = ['rest','active']
behavior_strs = list(behavior_dict.values())
nBehaviors = len(behavior_strs)

#Define windows to calculate firing rate
visual_stim = ['circle','natural_scene']
spont_types = ['spontaneous','isi_biphasic','isi_visual']
evoked_windows = [[.002,.025],[0.025,0.075],[.075,.3],[.3,1]]
evoked_strings = ['evoked','pre-rebound','rebound','post-rebound']

gc = gspread.service_account() # need a key file to access the account (step 2) 
sh = gc.open('Templeton-log_exp') # open the spreadsheet 

df = pd.DataFrame(sh.sheet1.get()) 
exp_table = df.T.set_index(0).T # put it in a nicely formatted dataframeexp_table.iloc[10:]
exp_table = exp_table.set_index('mouse_name')

##------------------------------------------
##Define Command Line Arguments
parser = argparse.ArgumentParser(description='single-cell-metrics')

parser.add_argument('--mID',type=str, default='mouse735052',
                    help='mouse to perform analysis on')

parser.add_argument('--rec_name',type=str, default='aw_psi_2024-06-20_10-40-56',
                    help='experiment to perform analysis on')

parser.add_argument('--window_size_min',type=int, default=5,
                    help='window_size_min')

parser.add_argument('--minumum_window_size_min',type=int, default=3,
                    help='minumum_window_size_min')

parser.add_argument('--isi_interval_ms',type=int, default=10,
                    help='isi_interval_ms')

parser.add_argument('--min_spikes_burst',type=int, default=2,
                    help='minimum number of spikes to be considered a burst')

parser.add_argument('--burst_quiescence_ms',type=int, default=10,
                    help='burst_quiescence_ms')

parser.add_argument('--run_num',type=int, default=-1,
                    help='run_num')

if __name__ == '__main__':

    ##------------------------------------------
    #Parse the arguments
    args = parser.parse_args()
    mID = args.mID
    rec_name = args.rec_name
    window_size_min =  args.window_size_min
    window_size = window_size_min*60
    minumum_window_size_min = args.minumum_window_size_min
    minumum_window_size = minumum_window_size_min*60
    isi_interval_ms = args.isi_interval_ms
    isi_interval = isi_interval_ms/1000
    min_spikes_burst = args.min_spikes_burst
    burst_quiescence_ms = args.burst_quiescence_ms
    burst_quiescence = burst_quiescence_ms/1000
    run_num = args.run_num

    #Extract row from templeton dataframe
    exp_df = exp_table.loc[exp_table.exp_name == rec_name]
    if type(exp_df) == pd.core.series.Series:
        exp_df = exp_df.to_frame().T

    drug_type = exp_df['drug'].values[0]
    stim_type = exp_df['stimulation'].values[0]
    if 'sal2' in rec_name:
        drug_type2 = 'saline2'
    else:
        drug_type2 = drug_type

    if 'electrical' in stim_type:
        st_type = 'electrical'
    elif 'spontaneous' in stim_type:
        st_type = 'spontaneous'
    print(f'Experiment type: {st_type}, {drug_type2}')

    #Define directories
    TempDir = join(ServDir,'results','FR_ISI_20250121',mID,rec_name)
    if not os.path.exists(TempDir):
        os.makedirs(TempDir)

    if run_num == -1:
        tmp_list = sorted(glob(join(TempDir,f'run_*')))
        if len(tmp_list) == 0:
            curr_run = 0
        else:
            last_run = int(tmp_list[-1][-1])
            curr_run = last_run+1
    else:
        curr_run = run_num

    folder = f'run_{curr_run:02d}'
    SaveDir = os.path.join(TempDir,folder)
    if not os.path.exists(SaveDir):
        os.makedirs(SaveDir)
    
    PlotDir = join(SaveDir,'plots')
    if not os.path.exists(PlotDir):
        os.makedirs(PlotDir)
        
    #Save parameters
    args_dict = args.__dict__
    args_dict['SaveDir'] = SaveDir
    with open(join(SaveDir,f"parameters_run_{curr_run}.json"), "w") as outfile:
        json.dump(args_dict, outfile)

    ##------------------------------------------
    #Upload the whole experiment 
    file_name = os.path.join(base_dir_server,mID,rec_name,'experiment1','recording1')
    exp = EEGexp(file_name, preprocess=False, make_stim_csv=False)

    #Load probe data
    probe_unit_data, probe_info, total_units, metrics_list = tbd_util.get_neuropixel_data(exp)
    probe_list = [x.replace('_sorted', '') for x in exp.experiment_data if 'probe' in x]
    plot_tuple = util.get_group_plotting_info(probe_unit_data)
    boundaries, yticks, labels, celltypes, durations, layers, areas, groups, mesogroups, supergroups, order_by_group = plot_tuple
    nNeurons = len(areas)

    #Get recording start time
    probei = probe_list[0]
    open_ephys_start = np.round(np.min(probe_unit_data[probei]['spike_times']))
    open_ephys_end = np.round(np.max(probe_unit_data[probei]['spike_times']))
    recording_length = open_ephys_end - open_ephys_start
    mm, ss = divmod(recording_length,60)
    hh, mm = divmod(mm, 60)
    print(f'{hh} hrs, {mm} minutes, {ss} seconds')

    ##------------------------------------------
    #Get injection times and types of injection
    injection_times = None; injection_time_windows = None; inj_tuple = None
    if drug_type in ['saline', 'psilocybin','ketanserin+psilocybin']:
        injection_time_windows = np.array([np.array(exp_df['First injection window'].values[0].split(','),dtype=float),
                                np.array(exp_df['Second injection window'].values[0].split(','),dtype=float)])
        injection_times = np.array([exp_df['First injection time (s)'].values[0],exp_df['Second injection time (s)'].values[0]],dtype=float)

        if drug_type == 'psilocybin':
            injection_types = ['sal1','psi']
            injection_colors = sns.xkcd_palette(['dark sky blue','darkish red'])
        elif drug_type == 'saline':
            injection_types = ['sal1', 'sal2']
            injection_colors = sns.xkcd_palette(['dark sky blue','cobalt blue'])
        elif drug_type == 'ketanserin+psilocybin':
            injection_types = ['ket','psi']
            injection_colors = sns.xkcd_palette(['magenta','goldenrod'])
        inj_tuple = (injection_times,injection_types,injection_colors)

    #For isoflurane experiments, get iso level
    iso_induction_times = None; iso_tuple = None
    if drug_type == 'isoflurane':
        iso_level, iso_times = exp.load_analog_iso()
        iso_induction_times, iso_maintenance_times = exp.load_iso_times()
        induction_colors = sns.xkcd_palette(['light teal','teal'])
        iso_tuple = (iso_induction_times, induction_colors)

    ##------------------------------------------
    #Define time windows for each epoch  
    if drug_type in ['saline', 'psilocybin','ketanserin+psilocybin']:
        epoch_windows = [[open_ephys_start,injection_time_windows[0,0]]]; epoch_names = ['pre_inj']
        block_labels = epoch_names.copy()
        # if (injection_time_windows[0,0] - open_ephys_start) <= window_size:
        #     epoch_windows = [[open_ephys_start,injection_time_windows[0,0]]]
        #     epoch_names = ['pre_inj']
        #     block_labels = ['pre_inj']

        # else:
        #     epoch_windows = []; epoch_names = []; block_labels = []
        #     t0 = open_ephys_start
        #     iW = 0; tblock = injection_time_windows[0,0]
        #     while t0 < tblock:
        #         t1 = t0 + window_size
        #         if t1 > tblock:
        #             t1 = tblock
        #         if (t1 - t0) < minumum_window_size:
        #             epoch_windows[-1] = [epoch_windows[-1][0],tblock]
        #             break
        #         epoch_windows.append([t0,t1]); epoch_names.append(f'pre_inj_{int(iW*window_size_min)}_{int((iW+1)*window_size_min)}')
        #         block_labels.append(f'pre_inj_{int(iW*window_size_min)}_{int((iW+1)*window_size_min)}')
        #         t0 = t1
        #         iW += 1

        epoch_windows.append([injection_time_windows[0,1],injection_time_windows[1,0]]); epoch_names.append(f'post_{injection_types[0]}_inj'); block_labels.append(f'post_{injection_types[0]}_inj')
        macro_windows = epoch_windows.copy(); macro_names = epoch_names.copy()
        macro_windows.append([injection_time_windows[1,1],open_ephys_end]); macro_names.append(f'post_{injection_types[1]}_inj')
        
        t_inj2 = injection_time_windows[1,1]
        t0 = t_inj2

        iW = 0
        blocks_to_plot = []
        while t0 < open_ephys_end:
            t1 = t0 + window_size
            if t1 > open_ephys_end:
                t1 = open_ephys_end
            if (t1 - t0) < minumum_window_size:
                epoch_windows[-1] = [epoch_windows[-1][0],open_ephys_end]
                break
            epoch_windows.append([t0,t1]); epoch_names.append(f'post_{injection_types[1]}_inj')
            block_labels.append(f'{int(iW*window_size_min)}_{int((iW+1)*window_size_min)}')
            blocks_to_plot.append(f'{int(iW*window_size_min)}_{int((iW+1)*window_size_min)}')
            t0 = t1
            iW += 1
        nCond = len(epoch_windows)
        if drug_type == 'saline':
            cmap_macro = sns.xkcd_palette(['silver','dark sky blue','cobalt blue'])
            cmap_epoch = np.concatenate((sns.xkcd_palette(['silver','dark sky blue']),sns.color_palette('Blues_r',nCond-2)))
        elif drug_type == 'psilocybin':
            cmap_macro = sns.xkcd_palette(['silver','dark sky blue','darkish red'])
            cmap_epoch = np.concatenate((sns.xkcd_palette(['silver','dark sky blue']),sns.color_palette('Reds_r',nCond-2)))
        elif drug_type == 'ketanserin+psilocybin':
            cmap_macro = sns.xkcd_palette(['silver','magenta','darkish red'])
            cmap_epoch = np.concatenate((sns.xkcd_palette(['silver','magenta']),sns.color_palette('Purples_r',nCond-2)))
        
    elif drug_type == 'isoflurane':

        epoch_windows = [[open_ephys_start,iso_induction_times[0]-120]]; epoch_names = ['pre_iso']
        epoch_windows.append([iso_induction_times[0],iso_induction_times[1]]); epoch_names.append(f'iso_ind')
        macro_windows = epoch_windows.copy(); macro_names = epoch_names.copy()
        macro_windows.append([iso_induction_times[1]+60*5,open_ephys_end]); macro_names.append(f'post-iso')
        block_labels = epoch_names.copy()
        t_ind = iso_induction_times[1] + 60*5
        t0 = t_ind
        iW = 0
        blocks_to_plot = [] 
        while t0 < open_ephys_end:
            t1 = t0 + window_size
            if t1 > open_ephys_end:
                t1 = open_ephys_end
            if (t1 - t0) < minumum_window_size:
                epoch_windows[-1] = [epoch_windows[-1][0],open_ephys_end]
                break
            epoch_windows.append([t0,t1]); epoch_names.append(f'post_iso')
            block_labels.append(f'{int(iW*window_size_min)}_{int((iW+1)*window_size_min)}')
            blocks_to_plot.append(f'{int(iW*window_size_min)}_{int((iW+1)*window_size_min)}')
            t0 = t1
            iW += 1
        nCond = len(epoch_windows)
        cmap_macro = sns.xkcd_palette(['silver','light teal','teal'])
        cmap_epoch = np.concatenate((sns.xkcd_palette(['silver','light teal']),sns.color_palette('Greens_r',nCond-2)))  

    elif drug_type == 'urethane':
        macro_windows = [[open_ephys_start,open_ephys_end]]; macro_names = ['urethane']
        t0 = 0
        epoch_windows = [];epoch_names = []; block_labels = []
        iW = 0
        blocks_to_plot = []
        while t0 < open_ephys_end:
            t1 = t0 + window_size
            if t1 > open_ephys_end:
                t1 = open_ephys_end
            t0_str = int(t0/60); t1_str = int(t1/60)
            if (t1 - t0) < minumum_window_size:
                epoch_windows[-1] = [epoch_windows[-1][0],open_ephys_end]
                break
            epoch_windows.append([t0,t1]) #; block_labels.append(f'{t0_str}_{t1_str}')
            epoch_names.append(f'urethane')
            block_labels.append(f'{int(iW*window_size_min)}_{int((iW+1)*window_size_min)}')
            blocks_to_plot.append(f'{int(iW*window_size_min)}_{int((iW+1)*window_size_min)}')
            t0 = t1
            iW += 1
        nCond = len(epoch_windows)
        cmap_macro = sns.xkcd_palette(['orange'])
        cmap_epoch = sns.color_palette('Oranges_r',nCond)
    nEpochs = len(epoch_windows)
    epoch_dict = {i: (tW, e) for i, (tW, e) in enumerate(zip(epoch_windows, epoch_names))}
    # for ii, (tW, e) in epoch_dict.items():
    #     print(f'{e}, {block_labels[ii]}: {tW[0]:.1f}s - {tW[1]:.1f}s')

    ##------------------------------------------
    #Load stimulus log if it exists
    try:
        stim_log = pd.read_csv(exp.stimulus_log_file)

        if ('circle' in np.unique(stim_log['stim_type'])) | ('natural_scene' in np.unique(stim_log['stim_type'])):
            vStim_exists = True
            stim_log_v = stim_log.loc[stim_log.stim_type.isin(visual_stim)]
        else:
            vStim_exists = False
        if 'biphasic' in np.unique(stim_log['stim_type']):
            eStim_exists = True
            stim_log_b = stim_log.loc[stim_log.stim_type == 'biphasic']
            uniq_currents = np.unique(stim_log_b['parameter'])
            uniq_sweeps = np.unique(stim_log_b['sweep'])
            nSweeps = len(uniq_sweeps)
            nCurrents = len(uniq_currents)
            print(f'Number of sweeps: {nSweeps}, Number of currents: {nCurrents}')
        else:
            eStim_exists = False
    except:
        stim_log = None
        eStim_exists = False
        vStim_exists = False

    evoked_time_window_list = []
    evoked_type_list = []
    sweep_time_window_list = []
    if eStim_exists | vStim_exists:
        for s in np.unique(stim_log['sweep']):
            for t in np.unique(stim_log.loc[stim_log.sweep == s]['stim_type']):
                sub_df = stim_log.loc[(stim_log.sweep == s) & (stim_log.stim_type == t)]
                tS = sub_df['onset'].min(); tE = sub_df['offset'].max()
                
                evoked_time_window_list.append([tS,tE])
                evoked_type_list.append(t)

    ##------------------------------------------
    #Load behavior
    run_ts, run_signal, run_signal_s, acceleration, pupil_ts, pupil_radius, pupil_dxdt, plot_pupil = util.get_behavioral_data(exp, mID, rec_name,normalize=True)
    run_signal[run_signal < 0] = 0
    run_signal_s[run_signal_s < 0] = 0
    f_run = interp1d(run_ts,run_signal)
    f_pupil = interp1d(pupil_ts,pupil_radius)

    # Create a new PowerPoint presentation to save figures to
    prs = Presentation()

    #Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f'mID = {mID}; rec_name = {rec_name}'
    slide.placeholders[1].text = f'Changes in single cell metrics due to {drug_type}'

    ##------------------------------------------
    #Plot behavior and save to pptx
    T = len(run_ts)
    time_bin_run = 0.01
    tMin = T*time_bin_run/60

    #Create behavioral map based on smoothed running speed
    behavior_map = np.zeros(run_signal_s.shape)
    for key, b_range in behavior_ranges.items():
        indy = np.where((run_signal_s >= b_range[0]) & (run_signal_s < b_range[1]))[0]
        behavior_map[indy] = key
        t = len(indy)*time_bin_run

    fig = usrplt.plot_behavior((run_ts,run_signal),(pupil_ts,pupil_radius),f'{mID} {rec_name}',(evoked_time_window_list,evoked_type_list),inj_tuple,iso_tuple,behavior_map,epoch_dict)
    plt.savefig(join(PlotDir,f'behavior_{rec_name}.png'),dpi=300,facecolor='white',bbox_inches='tight')
    usrplt.save_fig_to_pptx(fig, prs)
    plt.close(fig)

    spont_periods = np.ones((T))
    #Mask out periods of evoked activity
    if eStim_exists | vStim_exists:
        
        buffer = 0.01
        for ii, row in stim_log.iterrows():
            stim_type = row['stim_type']
            if stim_type == 'biphasic':
                #Just block out 1 second after electrical stimulation
                t1 = row['onset']  - buffer
                t2 = row['onset'] + 1
            else:
                #Block out whole duration of stimulus, as it is likely a visual stimulus
                t1 = row['onset'] - buffer
                t2 = row['offset']

            indy = np.where((run_ts >= t1) & (run_ts < t2))[0]
            if len(indy) > 0:
                spont_periods[indy] = 0

    ##------------------------------------------
    #Define periods of spontaneous activity and evoked activity
    stim_window_list = []
    if eStim_exists | vStim_exists:

        first_stim = stim_log['onset'].min()
        stim_window_list.append(['spontaneous',np.nan,np.nan,open_ephys_start,first_stim])

        #Loop over stimuli
        for ii, row in stim_log.iterrows():
            stim_type = row['stim_type']
            sweep = row['sweep']
            param = row['parameter']
            if stim_type == 'biphasic':
                for w, w_str in zip(evoked_windows,evoked_strings):
                    t1 = row['onset'] + w[0]
                    t2 = row['onset'] + w[1]
                    stim_window_list.append([w_str,sweep,param,t1,t2])
            elif stim_type in visual_stim:
                t1 = row['onset']
                t2 = row['offset']
                stim_window_list.append(['visual',sweep,param,t1,t2])
            else:
                raise Exception('Unknown stimulus type')
            
            #Add inter-stimulus-interval based on stimulus type
            if ii < len(stim_log)-1:
                buffer = 0.01
                next_stim_type = stim_log.iloc[ii+1]['stim_type']
                t1_next = stim_log.iloc[ii+1]['onset']-buffer

                if (stim_type == 'biphasic') & (next_stim_type == 'biphasic'):
                    stim_window_list.append(['isi_biphasic',sweep,param,t2,t1_next])
                elif (stim_type == 'biphasic') & (next_stim_type in visual_stim):
                    stim_window_list.append(['isi_biphasic',sweep,param,t2,t2+2])
                    stim_window_list.append(['spontaneous',np.nan,np.nan,t2+2,t1_next])
                elif (stim_type in visual_stim) & (next_stim_type in visual_stim):
                    stim_window_list.append(['isi_visual',sweep,param,t2,t1_next])
                elif (stim_type in visual_stim) & (next_stim_type == 'biphasic'):
                    stim_window_list.append(['isi_visual',sweep,param,t2,t2+2])
                    stim_window_list.append(['spontaneous',np.nan,np.nan,t2+2,t1_next])
                else:
                    stim_window_list.append(['spontaneous',np.nan,np.nan,t2,t1_next])
            else:
                if (stim_type == 'biphasic'):
                    stim_window_list.append(['isi_biphasic',sweep,param,t2,t2+2])
                elif (stim_type in visual_stim):
                    stim_window_list.append(['isi_visual',sweep,param,t2,t2+2])
                stim_window_list.append(['spontaneous',np.nan,np.nan,t2,open_ephys_end])

    else:
        stim_window_list.append(['spontaneous',np.nan,np.nan,open_ephys_start,open_ephys_end])

    ##------------------------------------------
    ## Get time windows for each behavioral bout
    behavioral_bout_dict = {0: [], 1: []}
    behavioral_bout_all = {0: [], 1: []}

    #Ensure behavioral bouts are at least 2 seconds
    x = behavior_map.copy()
    transitions = np.where(np.diff(np.concatenate(([x[0]+1],x,[x[-1]+1]))))[0]
    min_t_bout = 2 
    for ii, t1 in enumerate(transitions[:-1]):
        t2 = transitions[ii+1]

        key = x[t1]
        behavioral_bout_all[key].append(np.array([t1,t2])*time_bin_run)

        #If behavioral bout is too short, remove it
        if (t2-t1)*time_bin_run < min_t_bout:
            behavior_map[t1:t2] = np.nan
        else:
            behavioral_bout_dict[key].append(np.array([t1,t2])*time_bin_run)

    #Print statistics of behavior
    for key, behav_str in behavior_dict.items():
        behavioral_bouts = behavioral_bout_dict[key]
        nbouts = len(behavioral_bouts)
        if nbouts == 0:
            mean_bout_length = 0
        else:
            bout_length = np.diff(np.array(behavioral_bouts),axis=1).flatten()
            mean_bout_length = np.mean(bout_length)
        indy = np.where(behavior_map == key)[0]
        t = len(indy)*time_bin_run
        print(f'\t {behavior_dict[key]:15s} -> total time: {t:6.1f}s, # of bouts: {nbouts:4d}, mean bout length: {mean_bout_length:4.1f}s')
    
    ##------------------------------------------
    #Create combinations that we are interested in comparing for spontaneous epochs
    cbs = it.combinations(epoch_names,2)
    ijs = it.combinations(np.arange(len(epoch_names)),2)
    combos_spont = []
    for cb,ij in zip(cbs,ijs):
        epoch_i = cb[0];i = ij[0]
        epoch_j = cb[1];j = ij[1]
        if ('pre' in epoch_i) | ('post_sal1' in epoch_i) | ('post_ket' in epoch_i):
            if ('post' in epoch_j):
                combos_spont.append(ij)
    nComps_spont = len(combos_spont)

    ##------------------------------------------
    #Create combinations that we are interested in comparing for electrical stimulation epochs
    if eStim_exists:
        sweep_dict = {}
        for iS, sweep in enumerate(uniq_sweeps):
            sub_df = stim_log.loc[(stim_log.sweep == sweep)]
            tS = sub_df['onset'].min(); tE = sub_df['offset'].max()
            if drug_type in ['saline', 'psilocybin','ketanserin+psilocybin']:
                if tS < injection_times[0]:
                    sweep_dict[sweep] = macro_names[0]
                elif (tS > injection_times[0]) & (tS < injection_times[1]):
                    sweep_dict[sweep] = macro_names[1]
                elif tS > injection_times[1]:
                    sweep_dict[sweep] = macro_names[2]

            elif drug_type == 'isoflurane':
                if tS < iso_induction_times[0]:
                    sweep_dict[sweep] = macro_names[0]
                elif (tS > iso_induction_times[0]) & (tS < iso_induction_times[1]):
                    sweep_dict[sweep] = macro_names[1]
                elif tS > iso_induction_times[1]:
                    sweep_dict[sweep] = macro_names[2]
            elif drug_type == 'urethane':
                sweep_dict[sweep] = macro_names[0]
        sweep_epoch_names = list(sweep_dict.values())

        #Create combinations that we are interested in comparing for evoked epochs
        cbs = it.combinations(sweep_epoch_names,2)
        ijs = it.combinations(np.arange(len(sweep_epoch_names)),2)
        combos_estim = []
        for cb,ij in zip(cbs,ijs):
            epoch_i = cb[0];i = ij[0]
            epoch_j = cb[1];j = ij[1]
            if ('pre' in epoch_i) | ('post_sal1' in epoch_i) | ('post_ket' in epoch_i):
                if ('post' in epoch_j):
                    combos_estim.append(ij)
        nComps_estim = len(combos_estim)

    ##------------------------------------------
    #Get total time for each behavior per epoch
    nBehaviors = len(behavior_dict2)
    nT_behavior = np.zeros((nEpochs,nBehaviors))
    for e_key, (tW, e) in epoch_dict.items():
        t_indy = np.where((run_ts >= tW[0]) & (run_ts < tW[1]))[0]
        behavior_map_e = behavior_map[t_indy]
        for b_key, behav_str in behavior_dict.items():
            nT_behavior[e_key,b_key] = np.sum(behavior_map_e == b_key)*time_bin_run
        nT_behavior[e_key,-1] = tW[1] - tW[0]

    for e_key, (tW, e) in epoch_dict.items():
        t_total = nT_behavior[e_key,-1]/60
        print(f'{e}: Total time: {t_total:.1f}m')
        for b_key, behav_str in behavior_dict.items():
            t = nT_behavior[e_key,b_key]/60
            print(f'\t Time spent at {behav_str:15s} -> {t:.1f}m')

    nT_behavior2 = np.zeros((nEpochs,nBehaviors))
    for e_key, (tW, e) in epoch_dict.items():

        t_indy = np.where((run_ts >= tW[0]) & (run_ts < tW[1]) & (spont_periods == 1))[0]
        behavior_map_e = behavior_map[t_indy]

        for b_key, behav_str in behavior_dict.items():
            nT_behavior2[e_key,b_key] = np.sum((behavior_map_e == b_key))*time_bin_run
        nT_behavior2[e_key,-1] = len(t_indy)*time_bin_run
        
    ##------------------------------------------
    #Save time windows
    np.savez(join(SaveDir,f'time_windows_{rec_name}.npz'), epoch_windows=epoch_windows,epoch_names=epoch_names,block_labels=block_labels,\
             macro_windows=macro_windows,macro_names=macro_names,stim_window_list=stim_window_list,behavioral_bout_dict=behavioral_bout_dict, \
             celltypes=celltypes, durations=durations, layers=layers, areas=areas,groups=groups,mesogroups=mesogroups,supergroups=supergroups, \
             nT_behavior=nT_behavior,nT_behavior2=nT_behavior2)
   
    ##------------------------------------------
    #Get spike times for each neuron across experiment
    tW = [open_ephys_start,open_ephys_end]  
    spike_time_dict = util.get_spike_time_dict(probe_unit_data, tWindow=tW)
    
    #Calculate firing rate of each neuron for each condition & behaviors
    print('Calculating spontaneous firing rate statistics')
    nBehaviors = len(behavior_dict2)
    stat_list = ['mean_fr','cv_isi','burst_prop','burst_rate','burst_freq','burst_size']
    nStats = len(stat_list)
    moments_spont = np.full((nNeurons,nBehaviors,nEpochs,nStats),np.nan)
    mod_idx_spont = np.full((nNeurons,nBehaviors,nComps_spont,nStats),np.nan)
    pvalues_spont = np.full((nNeurons,nBehaviors,nComps_spont,2),np.nan)
    numspks_spont = np.zeros((nNeurons,nBehaviors,nEpochs))
    num_bursts = np.zeros((nNeurons,nBehaviors,nEpochs))
    
    eStim_exists = False
    if eStim_exists:
        print('Calculating evoked firing rate statistics')
        estim_windows = ['evoked','pre-rebound','rebound','post-rebound','isi_biphasic']
        nWindows = len(estim_windows)
        moments_estim = np.full((nNeurons,nBehaviors,nCurrents,nWindows,nSweeps,2),np.nan)
        pvalues_estim = np.full((nNeurons,nBehaviors,nCurrents,nWindows,nComps_estim,2),np.nan)
        mod_idx_estim = np.full((nNeurons,nBehaviors,nCurrents,nWindows,nComps_estim,2),np.nan)
        numspks_estim = np.zeros((nNeurons,nBehaviors,nCurrents,nWindows,nSweeps))

    #Loop over neurons
    for unit_i in tqdm(spike_time_dict.keys()):
        #Get spike times
        spk_times = spike_time_dict[unit_i]

        #Get ISI
        isi = np.diff(spk_times)
        spk_times1 = spk_times[1:]

        #Based on the time windows of behavioral bouts
        behav_per_spk = util.get_behavior_during_spike(spk_times1, behavioral_bout_dict)

        #Based on the injection times, what condition is the spike in
        epoch_per_spk = util.get_epoch_during_spike(spk_times1, epoch_dict)

        #Based on stimulus log, determine if the spike is in evoked interval or not
        stim_per_spk, sweep_per_spk, param_per_spk = util.get_stim_during_spike(spk_times1, stim_window_list)

        #Get spontaneous periods
        mask_spike_spont = np.array([s in spont_types for s in stim_per_spk])

        #Get burst events without considering pre-burst quiescence
        burst_spk_times, nSpikes_per_burst, mISI_per_burst, burst_indices = util.get_burst_events(spk_times, isi_interval, min_spikes_burst, 0)
        burst_exits = True
        if np.all(np.isnan(burst_spk_times)):
            burst_exits = False

        #Based on the time windows of behavioral bouts
        behav_per_burst = util.get_behavior_during_spike(burst_spk_times, behavioral_bout_dict)

        #Based on the injection times, what condition is the burst in
        epoch_per_burst = util.get_epoch_during_spike(burst_spk_times, epoch_dict)

        #Based on stimulus log, determine if the burst is in evoked interval or not
        stim_per_burst, sweep_per_burst, param_per_burst = util.get_stim_during_spike(burst_spk_times, stim_window_list)

        #Get spontaneous periods
        mask_burst_spont = np.array([s in spont_types for s in stim_per_burst])
        # import pdb; pdb.set_trace()
        ##------------------------------------------
        # Calculate spontaneous firing rates
        for b_key, behav_str in behavior_dict2.items():

            distr_list = []
            for e_key, (tW, e) in epoch_dict.items():
                if b_key == 2:
                    mask_spike = (epoch_per_spk == e_key) & (mask_spike_spont)
                    if burst_exits:
                        mask_burst = (epoch_per_burst == e_key) & (mask_burst_spont)
                    else:
                        mask_burst = np.array([])
                else:
                    mask_spike = (behav_per_spk == b_key) & (epoch_per_spk == e_key) & (mask_spike_spont)
                    if burst_exits:
                        mask_burst = (behav_per_burst == b_key) & (epoch_per_burst == e_key) & (mask_burst_spont)
                    else:
                        mask_burst = np.array([])

                numspks_spont[unit_i,b_key,e_key] = np.sum(mask_spike)
                num_bursts[unit_i,b_key,e_key] = np.sum(mask_burst)
                if np.sum(mask_spike) > 1:
                    moments_spont[unit_i,b_key,e_key,0] = 1/np.nanmean(isi[mask_spike])
                    moments_spont[unit_i,b_key,e_key,1] = st.variation(isi[mask_spike],nan_policy='omit')

                    #Burst statistics if they exist
                    if burst_exits:
                        moments_spont[unit_i,b_key,e_key,2] = np.sum(mask_burst)/np.sum(mask_spike)
                        moments_spont[unit_i,b_key,e_key,3] = np.sum(mask_burst)/nT_behavior2[e_key,b_key]
                        moments_spont[unit_i,b_key,e_key,4] = 1/np.nanmean(mISI_per_burst[mask_burst])
                        moments_spont[unit_i,b_key,e_key,5] = np.nanmean(nSpikes_per_burst[mask_burst])

                    #Append ISI distribution of cell
                    distr_list.append(isi[mask_spike])
                else:
                    distr_list.append([])

            if nComps_spont > 0: 
                #Compare ISI distributions to assess significance & calculate modulation index
                for ii, (i,j) in enumerate(combos_spont):
                    for k, stat in enumerate(stat_list):
                        stat_i = moments_spont[unit_i,b_key,i,k]
                        stat_j = moments_spont[unit_i,b_key,j,k]
                        if np.isnan(stat_i) | np.isnan(stat_j):
                            MI = np.nan
                        else:
                            MI = (stat_j - stat_i)/(stat_j + stat_i)
                        mod_idx_spont[unit_i,b_key,ii,k] = MI
                    
                    distr_i = distr_list[i]
                    distr_j = distr_list[j]
                    if (len(distr_i) > 0) & (len(distr_j) > 0):

                        #Test the hypothesis that distr_j is larger tha distr_i
                        res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')
                        pvalues_spont[unit_i,b_key,ii,0] = res.pvalue

                        #Perform permutation test on the difference in means
                        # res = st.permutation_test([distr_i,distr_j],util.perm_test_statistic,n_resamples=1000,vectorized=False)
                        res = st.ks_2samp(distr_i,distr_j)#,nan_policy='omit')
                        pvalues_spont[unit_i,b_key,ii,1] = res.pvalue
        # import pdb; pdb.set_trace()
        if not eStim_exists:
            continue
        
        ##------------------------------------------
        # Calculate evoked firing rates
        for b_key, behav_str in behavior_dict2.items():
            for iA, current in enumerate(uniq_currents):
                for iW, w_str in enumerate(estim_windows):
                    distr_list = []
                    for iS, sweep in enumerate(uniq_sweeps):
                        if w_str == 'isi_biphasic':
                            if b_key == 2:
                                mask = (sweep_per_spk == sweep) & (stim_per_spk == w_str)
                            else:
                                mask = (behav_per_spk == b_key) & (sweep_per_spk == sweep) & (stim_per_spk == w_str)
                        else:
                            if b_key == 2:
                                mask = (sweep_per_spk == sweep) & (param_per_spk == current) & (stim_per_spk == w_str)
                            else:
                                mask = (behav_per_spk == b_key) & (sweep_per_spk == sweep) & (param_per_spk == current) & (stim_per_spk == w_str)

                        numspks_estim[unit_i,b_key,iA,iW,iS] = np.sum(mask)
                        if np.sum(mask) > 1:
                            moments_estim[unit_i,b_key,iA,iW,iS,0] = 1/np.nanmean(isi[mask])
                            moments_estim[unit_i,b_key,iA,iW,iS,1] = st.variation(isi[mask],nan_policy='omit')
                            distr_list.append(isi[mask])
                        else:
                            distr_list.append([])
                    
                    if nComps_estim > 0: 
                        #Compare ISI distributions & assess significance
                        for ii, (i,j) in enumerate(combos_estim):
                            distr_i = distr_list[i]
                            distr_j = distr_list[j]

                            if (len(distr_i) > 0) & (len(distr_j) > 0):

                                #Test the hypothesis that distr_j is larger tha distr_i
                                res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')
                                pvalues_estim[unit_i,b_key,iA,iW,ii,0] = res.pvalue

                                #Perform permutation test on the difference in means
                                # res = st.permutation_test([distr_i,distr_j],util.perm_test_statistic,n_resamples=1000,vectorized=False)
                                res = st.ks_2samp(distr_i,distr_j)#,nan_policy='omit')
                                pvalues_estim[unit_i,b_key,iA,iW,ii,1] = res.pvalue

                                #Calculate modulation index of firing rate
                                FR_i = moments_estim[unit_i,b_key,iA,iW,i,0]; FR_j = moments_estim[unit_i,b_key,iA,iW,j,0]
                                if np.isnan(FR_i) | np.isnan(FR_j):
                                    mod_idx_estim[unit_i,b_key,iA,iW,ii,0] = np.nan
                                else:
                                    mod_idx_estim[unit_i,b_key,iA,iW,ii,0] = (FR_j - FR_i)/(FR_j + FR_i)
             
                                #Calculate modulation index of CV (not sure if this is valid)
                                CV_i = moments_estim[unit_i,b_key,iA,iW,i,1]; CV_j = moments_estim[unit_i,b_key,iA,iW,j,1]
                                if np.isnan(CV_i) | np.isnan(CV_j):
                                    mod_idx_estim[unit_i,b_key,iA,iW,ii,1] = np.nan
                                else:
                                    mod_idx_estim[unit_i,b_key,iA,iW,ii,1] = (CV_j - CV_i)/(CV_j + CV_i)

    #------------------------------------------
    #Save results
    np.savez(join(SaveDir,f'single-cell-metrics_spont_{rec_name}.npz'), combos_spont=combos_spont, \
            moments_spont=moments_spont,pvalues_spont=pvalues_spont,mod_idx_spont=mod_idx_spont,numspks_spont=numspks_spont)
    
    if eStim_exists:
        np.savez(join(SaveDir,f'single-cell-metrics_estim_{rec_name}.npz'), combos_estim=combos_estim, sweep_dict=sweep_dict, \
            moments_estim=moments_estim,pvalues_estim=pvalues_estim,mod_idx_estim=mod_idx_estim,numspks_estim=numspks_estim)
        
    ##------------------------------------------
    #Plotting and csv creation
    #Create csv file for easy reading later
    block_labels2 = []
    for b in block_labels:
        if 'pre_inj' in b:
            block_labels2.append('pre_inj')
        elif (b == 'post_sal1_inj') | (b == 'post_ket_inj'):
            block_labels2.append('post_1st_inj')
        else:
            block_labels2.append(b)
    block_labels = np.array(block_labels2)

    th_dict = {'SM-TH': ['PO','VAL','VPL','VPM','VM'], 
    'VIS-TH': ['LP','LGN','LGd','LGd-co','LGd-sh','LGd-ip'],
    'ANT-TH': ['AV','AMd','AMv','AD','MD','MDm','MDc','MDl','RE','RH','CM','LD', 'CL'],
    'TH': ['Eth', 'IAD', 'IGL', 'IntG', 'LGv','MGd', 'MGm', 'MGv', 'PCN', 'PF', 'PIL', 'PoT', 'SGN','SPFp', 'TH','LH'],'RT': ['RT']}
    groups2 = []
    for a, g, sg in zip(areas, groups, supergroups):
        if sg == 'TH':
            if a in th_dict['SM-TH']:
                groups2.append('SM-TH')
            elif a in th_dict['VIS-TH']:
                groups2.append('VIS-TH')
            elif a in th_dict['ANT-TH']:
                groups2.append('ANT-TH')
            elif a in th_dict['RT']:
                groups2.append('RT')
            else:
                groups2.append('TH')
            # else:
            #     raise Exception(f'Unknown thalamic area/group: {a}')
        else:
            groups2.append(g)
    groups = np.array(groups2)

    indy = np.where(supergroups == 'TH')[0]
    if len(indy) > 0:
        indy = np.where((durations < 0.35) & (supergroups == 'TH'))[0]
        celltypes[indy] = 'FS'

        indy = np.where((durations >= 0.45) & (supergroups == 'TH'))[0]
        celltypes[indy] = 'RS'

        indy = np.where((durations >= 0.35) & (durations < 0.45) & (supergroups == 'TH'))[0]
        celltypes[indy] = 'unclassified'
        
    #Ensure all neurons in RT are inhibitory
    indy = np.where(areas == 'RT')[0]
    if len(indy) > 0:
        mesogroups[indy] = 'TH_RT'
        celltypes[indy] = 'FS'

    #Save for cell dataframe
    d_array = np.repeat(drug_type2,nNeurons)
    idx_arr = np.array([f'{rec_name}_{iN}' for iN in range(nNeurons)])
    cells_df = pd.DataFrame(np.vstack((d_array,idx_arr,durations,celltypes,areas,layers, groups, mesogroups, supergroups)).T,columns=['drug_type','neuron_idx','duration','celltype','area','layer','group','mesogroup','supergroup'])
    cells_df.to_csv(join(SaveDir,f'cells_df_{rec_name}.csv'),index=False)
 
    ##------------------------------------------
    #Create dataframes for all stats during spontaneous epochs
    FR_list = []
    for iN, (ct, dur, a, l, g, mg, sg) in enumerate(zip(celltypes,durations,areas,layers,groups,mesogroups,supergroups)):
        for iB, behav_str in behavior_dict2.items():
            for iE, (tW, epoch, block) in enumerate(zip(epoch_windows,epoch_names,block_labels)):
                if (block == 'post_sal1_inj') | (block == 'post_ket_inj'):
                        block2 = 'post_1st_inj'
                for iS, stat in enumerate(stat_list):
                    if iS < 2:
                        n = numspks_spont[iN,iB,iE]
                    else:
                        n = num_bursts[iN,iB,iE]
                    neuron_idx = f'{rec_name}_{iN}'
                    FR_list.append((mID,rec_name,drug_type,neuron_idx,ct,dur,a,l,g,mg,sg,behav_str,epoch,block2,n,'spont',stat,moments_spont[iN,iB,iE,iS]))
                    FR_list.append((mID,rec_name,drug_type,neuron_idx,ct,dur,a,l,g,mg,sg,behav_str,epoch,block2,n,'spont','CV_ISI',moments_spont[iN,iB,iE,1]))
    FR_spont_df = pd.DataFrame(np.stack(FR_list),columns=['mID','rec_name','drug_type','neuron_index','celltype','duration','area','layer','group','mesogroup','supergroup','behavior','epoch','block','num_spks','stim_type','stat','val'])
    FR_spont_df = FR_spont_df.astype({'duration': float,'num_spks': float, 'val': float})
    FR_spont_df['st_type'] = st_type
    FR_spont_df['drug_type2'] = drug_type2
    FR_spont_df.to_csv(join(SaveDir,f'FR_spont_df_{rec_name}.csv'),index=False)

    ##------------------------------------------
    #Create dataframes for change in stats between epochs
    if nComps_spont > 0:
        pval_thresh = 0.05
        dFR_list = []
        for iN, (ct, dur, a, l, g, mg, sg) in enumerate(zip(celltypes,durations,areas,layers,groups,mesogroups,supergroups)):
            neuron_idx = f'{rec_name}_{iN}'
            for iB, behav_str in behavior_dict2.items():
                for ii, (i,j) in enumerate(combos_spont):
                    epoch_i = epoch_names[i]; epoch_j = epoch_names[j]
                    block_i = block_labels[i]; block_j = block_labels[j]

                    if (block_i == 'post_sal1_inj') | (block_i == 'post_ket_inj'):
                        block_i = 'post_1st_inj'

                    if (block_j == 'post_sal1_inj') | (block_j == 'post_ket_inj'):
                        block_j = 'post_1st_inj'
                    
                    #Get p-value of ISI distributions
                    pval1 = pvalues_spont[iN,iB,ii,0]
                    pval2 = pvalues_spont[iN,iB,ii,1]
                    sig1 = 1 if pval1 < pval_thresh else 0
                    sig2 = 1 if pval2 < pval_thresh else 0
                    sig = 1 if (sig1 + sig2) == 2 else 0

                    for iS, stat in enumerate(stat_list):
                        #Calculate modulation index of firing rate
                        stat_i = moments_spont[iN,iB,i,iS]; stat_j = moments_spont[iN,iB,j,iS]
                        MI = mod_idx_spont[iN,iB,ii,iS]
                        if np.isnan(stat_i) | np.isnan(stat_j):
                            continue
                        dFR_list.append((mID,rec_name,drug_type,neuron_idx,ct,dur,a,l,g,mg,sg,behav_str,epoch_i,block_i,epoch_j,block_j,stat,stat_i,stat_j,stat_j-stat_i,pval1,sig1,pval2,sig2,sig,MI))

        dFR_spont_df = pd.DataFrame(np.stack(dFR_list),columns=['mID','rec_name','drug_type','neuron_index','celltype','duration','area','layer','group','mesogroup','supergroup','behavior','epoch_i','block_i','epoch_j','block_j','stat','val_i','val_j','delta','pval1','sig1','pval2','sig2','sig','mod'])
        dFR_spont_df = dFR_spont_df.astype({'area': str, 'duration': float,'val_i': float, 'val_j': float,'delta': float, 'mod': float,'pval1': float,'sig1': int,'pval2': float,'sig2': int,'sig':int})
        dFR_spont_df['abs_mod'] = np.abs(dFR_spont_df['mod'])
        dFR_spont_df['abs_delta'] = np.abs(dFR_spont_df['delta'])
        dFR_spont_df['st_type'] = st_type
        dFR_spont_df['drug_type2'] = drug_type2
        dFR_spont_df.to_csv(join(SaveDir,f'dFR_spont_df_{rec_name}.csv'),index=False)

    ##------------------------------------------
    #Create dataframes for all stats during evoked epochs
    if eStim_exists:
        FR_list = []
        for iN, (ct, dur, a, l, g, mg, sg) in enumerate(zip(celltypes,durations,areas,layers,groups,mesogroups,supergroups)):
            neuron_idx = f'{rec_name}_{iN}'
            for b_key, behav_str in behavior_dict2.items():
                for iA, current in enumerate(uniq_currents):
                    for iW, w_str in enumerate(estim_windows):
                        for iS, sweep in enumerate(uniq_sweeps):
                            epoch = sweep_dict[sweep]
                            n = numspks_estim[iN,b_key,iA,iW,iS]
                            FR_list.append((mID,rec_name,drug_type,neuron_idx,ct,dur,a,l,g,mg,sg,behav_str,epoch,sweep,n,'evoked',w_str,'mean_fr',moments_estim[iN,b_key,iA,iW,iS,0]))
                            FR_list.append((mID,rec_name,drug_type,neuron_idx,ct,dur,a,l,g,mg,sg,behav_str,epoch,sweep,n,'evoked',w_str,'cv_isi',moments_estim[iN,b_key,iA,iW,iS,1]))
        FR_estim_df = pd.DataFrame(np.stack(FR_list),columns=['mID','rec_name','drug_type','neuron_index','celltype','duration','area','layer','group','mesogroup','supergroup','behavior','epoch','sweep','num_spks','stim_type','window','stat','val'])
        FR_estim_df = FR_estim_df.astype({'duration': float,'num_spks': float, 'val': float})
        FR_estim_df['st_type'] = st_type
        FR_estim_df['drug_type2'] = drug_type2
        FR_estim_df.to_csv(join(SaveDir,f'FR_estim_df_{rec_name}.csv'),index=False)

        if len(nComps_estim) > 0:
            dFR_list = []
            for iN, (ct, dur, a, l, g, mg, sg) in enumerate(zip(celltypes,durations,areas,layers,groups,mesogroups,supergroups)):
                neuron_idx = f'{rec_name}_{iN}'
                for b_key, behav_str in behavior_dict2.items():
                    for iA, current in enumerate(uniq_currents):
                        for iW, w_str in enumerate(estim_windows):
                            for ii, (i,j) in enumerate(combos_estim):
                                epoch_i = sweep_epoch_names[i]; epoch_j = sweep_epoch_names[j]
        
                                pval1 = pvalues_estim[iN,b_key,iA,iW,ii,0]
                                pval2 = pvalues_estim[iN,b_key,iA,iW,ii,1]
                                sig1 = 1 if pval1 < pval_thresh else 0
                                sig2 = 1 if pval2 < pval_thresh else 0
                                sig = 1 if (sig1 + sig2) == 2 else 0
                                for iS, stat in enumerate(['mean_fr','cv_isi']):
                                    stat_i = moments_estim[iN,b_key,iA,iW,i,iS]; stat_j = moments_estim[iN,b_key,iA,iW,j,iS]
                                    if np.isnan(stat_i) | np.isnan(stat_j):
                                        continue
                                    MI = mod_idx_estim[iN,b_key,iA,iW,ii,iS]
                                    dFR_list.append((mID,rec_name,drug_type,neuron_idx,ct,dur,a,l,g,mg,sg,behav_str,epoch_i,i,epoch_j,j,stat,stat_i,stat_j,stat_j-stat_i,pval1,sig1,pval2,sig2,sig,MI))
            dFR_estim_df = pd.DataFrame(np.stack(dFR_list),columns=['mID','rec_name','drug_type','neuron_index','celltype','duration','area','layer','group','mesogroup','supergroup','behavior','epoch_i','sweep_i','epoch_j','sweep_j','stat','val_i','val_j','delta','pval1','sig1','pval2','sig2','sig','mod'])
            dFR_estim_df = dFR_estim_df.astype({'area': str, 'duration': float,'val_i': float, 'val_j': float,'delta': float, 'mod': float,'pval1': float,'sig1': int,'pval2': float,'sig2': int,'sig':int})
            dFR_estim_df['abs_mod'] = np.abs(dFR_estim_df['mod'])
            dFR_estim_df['abs_delta'] = np.abs(dFR_estim_df['delta'])
            dFR_estim_df['st_type'] = st_type
            dFR_estim_df['drug_type2'] = drug_type2
            dFR_estim_df.to_csv(join(SaveDir,f'dFR_estim_df_{rec_name}.csv'),index=False)

    ##------------------------------------------
    #Plotting!
    stat_title_list = ['Mean firing rate (Hz)','Coefficient of variation of ISI','Burst proportion','Burst event rate (Hz)','Burst frequency (Hz)','Mean # of spikes in burst']
    stat_list = ['mean_fr','cv_isi','burst_prop','burst_rate','burst_freq','burst_size']
    stat_labels = ['FR (Hz)','CV (ISI)','Burst proportion','Burst event rate (Hz)','Burst frequency (Hz)','Burst size (# spikes)']
    stat_units = ['Hz','a.u.','','Hz','Hz','# spikes']
    mesogroup_list = ['TH_core','TH_matrix','TH_intralaminar','CTX_frontal','CTX_sensory','HIP','STR']
    plot_only_sig_mod = False


    for iS, stat in enumerate(stat_list):
        for g in ['CTX','TH']:
            sub_df = FR_spont_df.loc[(FR_spont_df.supergroup == g) & (FR_spont_df.stat == stat) & (FR_spont_df.behavior == behavior_dict2[2])]
            if len(sub_df) < 5:
                continue
            n_RS = len(sub_df.loc[sub_df['celltype'] == 'RS']['neuron_index'].unique())
            n_FS = len(sub_df.loc[sub_df['celltype'] == 'FS']['neuron_index'].unique())

            fig, axes = plt.subplots(3,1,figsize=(8,12),sharey=True)  
            plt.suptitle(f'{stat_title_list[iS]}\n {n_RS} RS cells, {n_FS} FS cells in {g}',y=0.95)
            for b_key, behavior in behavior_dict2.items():
                sub_df = FR_spont_df.loc[(FR_spont_df.supergroup == g) & (FR_spont_df.behavior == behavior) & (FR_spont_df.stat == stat)]
                ax = axes[b_key]; ax.set_title(f'{behavior}')
                sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                sns.pointplot(data=sub_df,x='block',y='val',hue='celltype',ax=ax,errorbar='ci')
                for ct in ['RS','FS']:
                    subsub_df = sub_df.loc[(sub_df['celltype'] == ct) & (sub_df['block'] == 'pre_inj')]
                    mean_val = subsub_df['val'].mean()
                    ax.axhline(mean_val,color='k',linestyle=':',alpha=0.25)

                ax.set_ylabel(stat_labels[iS])
                usrplt.adjust_spines(ax)
                if b_key < 2:
                    ax.set_xticklabels('')
                    ax.set_xlabel('')
                else:
                    ax.set_xticklabels(block_labels,rotation=30)
                    ax.set_xlabel('Epoch')

            # pdb.set_trace()
            plt.savefig(join(PlotDir,f'{stat}-val_3x_{g}.png'),dpi=300,bbox_inches='tight',facecolor='w')
            usrplt.save_fig_to_pptx(fig, prs)
            plt.close(fig)
        
            epoch_i = 'post_sal1_inj'; behavior = behavior_dict2[2]; bstr = 'all'
            fig, ax = plt.subplots(figsize=(8,4))  
            plt.suptitle(f'{stat_title_list[iS]}\n {n_RS} RS cells, {n_FS} FS cells in {g}',y=0.95)

            sub_df = FR_spont_df.loc[(FR_spont_df.supergroup == g) & (FR_spont_df.behavior == behavior) & (FR_spont_df.stat == stat)]

            sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
            sns.pointplot(data=sub_df,x='block',y='val',hue='celltype',ax=ax,errorbar='ci')
            for ct in ['RS','FS']:
                subsub_df = sub_df.loc[(sub_df['celltype'] == ct) & (sub_df['block'] == 'pre_inj')]
                mean_val = subsub_df['val'].mean()
                ax.axhline(mean_val,color='k',linestyle=':',alpha=0.25)
            ax.set_ylabel(stat_labels[iS])
            usrplt.adjust_spines(ax)
            ax.set_xticklabels(block_labels,rotation=30)
            ax.set_xlabel('Epoch')

            plt.savefig(join(PlotDir,f'{stat}-val_{g}.png'),dpi=300,bbox_inches='tight',facecolor='w')
            usrplt.save_fig_to_pptx(fig, prs)
            plt.close(fig)
            # pdb.set_trace()

    for iS, stat in enumerate(stat_list):
        for ct in ['RS','FS']:
            for g in mesogroup_list:
                if g in ['CTX_frontal','CTX_sensory']:
                    layer_list = ['2/3','4','5','6']
                else:
                    layer_list = ['none']
                for l in layer_list:
                    sub_df = FR_spont_df.loc[(FR_spont_df.mesogroup == g) & (FR_spont_df.layer == l) & (FR_spont_df.celltype == ct) & (FR_spont_df.stat == stat) & (FR_spont_df.behavior == behavior_dict2[2])]
                    if len(sub_df) < 5:
                        continue
                    n = len(sub_df['neuron_index'].unique())
                    fig, axes = plt.subplots(3,1,figsize=(8,12),sharey=True)  
                    if l == 'none':
                        plt.suptitle(f'{stat_title_list[iS]}\n {n} {ct} cells in {g}',y=0.95)
                    else:  
                        plt.suptitle(f'{stat_title_list[iS]}\n {n} {ct} cells in {g}, layer {l}',y=0.95)

                    for b_key, behavior in behavior_dict2.items():
                        sub_df = FR_spont_df.loc[(FR_spont_df.mesogroup == g) & (FR_spont_df.layer == l) & (FR_spont_df.behavior == behavior) & (FR_spont_df.celltype == ct) & (FR_spont_df.stat == stat)]
                        ax = axes[b_key]; ax.set_title(f'{behavior}')
                        sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                        sns.pointplot(data=sub_df,x='block',y='val',ax=ax,color=cmap_macro[-1],errorbar='ci')

                        ax.set_ylabel(stat_labels[iS])
                        usrplt.adjust_spines(ax)
                        if b_key < 2:
                            ax.set_xticklabels('')
                            ax.set_xlabel('')
                        else:
                            ax.set_xticklabels(block_labels,rotation=30)
                            ax.set_xlabel('Epoch')
                    if l == '2/3':
                        l_str = '2-3'
                    else:
                        l_str = l

                    plt.savefig(join(PlotDir,f'{stat}-val_3x_{g}_{l_str}_{ct}.png'),dpi=300,bbox_inches='tight',facecolor='w')
                    usrplt.save_fig_to_pptx(fig, prs)
                    plt.close(fig)
                
                    epoch_i = 'post_sal1_inj'; behavior = behavior_dict2[2]; bstr = 'all'
                    fig, ax = plt.subplots(figsize=(8,4))  
                    if l == 'none':
                        plt.suptitle(f'{stat_title_list[iS]}\n {g}, {ct} cells',y=0.95)
                    else:  
                        plt.suptitle(f'{stat_title_list[iS]}\n {g}, layer {l}, {ct} cells',y=0.95)

                    sub_df = FR_spont_df.loc[(FR_spont_df.mesogroup == g) & (FR_spont_df.layer == l) & (FR_spont_df.behavior == behavior) & (FR_spont_df.celltype == ct) & (FR_spont_df.stat == stat)]
                    sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                    sns.pointplot(data=sub_df,x='block',y='val',ax=ax,color=cmap_macro[-1],errorbar='ci')
                    ax.set_ylabel(stat_labels[iS])
                    usrplt.adjust_spines(ax)
                    ax.set_xticklabels(block_labels,rotation=30)
                    ax.set_xlabel('Epoch')

                    plt.savefig(join(PlotDir,f'{stat}-val_{g}_{epoch_i}_{l_str}_{ct}_{bstr}.png'),dpi=300,bbox_inches='tight',facecolor='w')
                    usrplt.save_fig_to_pptx(fig, prs)
                    plt.close(fig)

                    ##------------------------------------------    
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.mesogroup == g) & (dFR_spont_df.layer == l) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat)]
                    if len(sub_df) < 5:
                        continue
                    blocks_to_plot2 = np.concatenate((['post_1st_inj'],blocks_to_plot))
                    n = len(sub_df['neuron_index'].unique())
                    fig, axes = plt.subplots(3,1,figsize=(8,12),sharey=True)  
                    if l == 'none':
                        plt.suptitle(f'{stat_title_list[iS]}\n {n} {ct} cells in {g}',y=0.95)
                    else:  
                        plt.suptitle(f'{stat_title_list[iS]}\n {n} {ct} cells in {g}, layer {l}',y=0.95)

                    df_list = []
                    for b_key, behavior in behavior_dict2.items():
                        sub_df = dFR_spont_df.loc[(dFR_spont_df.mesogroup == g) & (dFR_spont_df.layer == l) & (dFR_spont_df.behavior == behavior) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat)]
                        ax = axes[b_key]; ax.set_title(f'{behavior}')
                        sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                        sns.pointplot(data=sub_df,x='block_j',y='mod',hue='block_i',ax=ax,order=blocks_to_plot2,color=cmap_macro[-1],hue_order=['pre_inj','post_1st_inj'],linestyles=['-','--'],errorbar='ci')
                        ax.hlines(0,*ax.get_xlim(),linestyles='--',color='k')

                        ax.set_ylabel('Net modulation index')
                        usrplt.adjust_spines(ax)

                        if b_key < 2:
                            ax.set_xticklabels('')
                            ax.set_xlabel('')
                        else:
                            ax.set_xticklabels(blocks_to_plot2,rotation=30)
                            ax.set_xlabel('Time after 2nd injection (min)')
                    if l == '2/3':
                        l_str = '2-3'
                    else:
                        l_str = l

                    plt.savefig(join(PlotDir,f'{stat}-delta_3x_{g}_{l_str}_{ct}.png'),dpi=300,bbox_inches='tight',facecolor='w')
                    usrplt.save_fig_to_pptx(fig, prs)
                    plt.close(fig)
                
                    epoch_i = 'post_1st_inj'; behavior = behavior_dict2[2]; bstr = 'all'
                    fig, ax = plt.subplots(figsize=(8,4))  
                    if l == 'none':
                        plt.suptitle(f'{stat_title_list[iS]}\n {n} {ct} cells in {g}',y=0.95)
                    else:  
                        plt.suptitle(f'{stat_title_list[iS]}\n {n} {ct} cells in {g}, layer {l}',y=0.95)

                    sub_df = dFR_spont_df.loc[(dFR_spont_df.mesogroup == g) & (dFR_spont_df.layer == l) & (dFR_spont_df.behavior == behavior) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i)]

                    sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                    sns.pointplot(data=sub_df,x='block_j',y='mod',ax=ax,order=blocks_to_plot,color=cmap_macro[-1],errorbar='ci')
                    ax.hlines(0,*ax.get_xlim(),linestyles='--',color='k')

                    ax.set_ylabel('Net modulation index')
                    usrplt.adjust_spines(ax)
                    ax.set_xticklabels(blocks_to_plot,rotation=30)
                    ax.set_xlabel('Time after 2nd injection (min)')

                    plt.savefig(join(PlotDir,f'{stat}-delta_{g}_{epoch_i}_{l_str}_{ct}_{bstr}.png'),dpi=300,bbox_inches='tight',facecolor='w')
                    usrplt.save_fig_to_pptx(fig, prs)
                    plt.close(fig)


    tt, cc = np.unique(cells_df['group'],return_counts=True)
    uniq_groups = []
    for g in tt:
        if g not in ['root','X','grey','CTX']:
            uniq_groups.append(g)
    group_order_labels, group_order_names  = util.determine_group_order(uniq_groups)

    regions_to_skip = ['FRP','SS','MBmot','PAL','PALm','PTLp','HY']
    fig, ax = plt.subplots(figsize=(10,4))

    group_tmp = []; iX = 0; nT = 0
    for g in group_order_labels:
        indy = np.where(cells_df['group'] == g)[0]
        n = len(indy)
        if n < 50:
            continue
        if g in regions_to_skip:
            continue
        if g == 'HIP':
            gg = ['HIP','HPF']
        elif g == 'HPF':
            continue
        else:
            gg = [g]
        group_tmp.append(g)
        print(f'{g}: {n} neurons')
        nT+=n
        cmap2 = sns.xkcd_palette(['sage','dusty orange','grey'])
        # for ii, ct in enumerate(['RS','FS']):
        nN_RS = len(cells_df.loc[(cells_df['group'].isin(gg)) & (cells_df['celltype'] == 'RS')])
        h1 = ax.bar(iX,nN_RS,color=cmap2[0], label='RS')

        nN_FS = len(cells_df.loc[(cells_df['group'].isin(gg)) & (cells_df['celltype'] == 'FS')])
        h2 = ax.bar(iX,nN_FS,bottom=nN_RS,color=cmap2[1], label='FS')

        nN_UN = len(cells_df.loc[(cells_df['group'].isin(gg)) & (cells_df['celltype'] == 'unclassified')])
        h3 = ax.bar(iX,nN_UN,bottom=nN_RS+nN_FS,color=cmap2[2], label='unclassified')

        # ax.text(iX,0.025,f'{n}',ha='center',fontsize=8,rotation=90)
        iX += 1
    ax.grid(axis='y',alpha=0.25)
    usrplt.adjust_spines(ax, ['left','bottom'])
    ax.set_xticks(np.arange(len(group_tmp)))
    ax.set_xticklabels(group_tmp,rotation=30)
    ax.set_ylabel('Number of neurons')
    ax.set_xlabel('Region')
    ax.legend(handles=[h1,h2,h3],loc=1)
    plt.savefig(join(PlotDir,f'cells_df_{rec_name}.png'),dpi=300,facecolor='white',bbox_inches='tight')
    usrplt.save_fig_to_pptx(fig, prs)
    plt.close(fig)


    ##------------------------------------------
    #Save to pptx
    prs.save(join(SaveDir,f'single-cell-metrics_{rec_name}.pptx'))
    # print('DONE!!!')


