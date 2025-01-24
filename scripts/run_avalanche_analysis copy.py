base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
base_dir_local = '/data/tiny-blue-dot/zap-n-zip/EEG_exp'; base_dir = base_dir_server
ProjDir = '/home/david.wyrick/projects/zap-n-zip/'
ServDir = '/data/projects/zap-n-zip/'
BehavDir = '/home/david.wyrick/projects/zap-n-zip/plots/behavior'

#Basea
import gspread
from os.path import join
from glob import glob
import json, os, time, sys, argparse
import numpy as np
import pandas as pd
import scipy.stats as st
import scipy.signal as sig
import pingouin as pg
from scipy.ndimage import gaussian_filter
from scipy.interpolate import interp1d
from itertools import combinations

#Plot
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
from matplotlib.gridspec import GridSpec
from matplotlib.lines import Line2D   

#Project
from tbd_eeg.data_analysis.eegutils import EEGexp
import tbd_eeg.data_analysis.Utilities.utilities as tbd_util
from tbd_eeg.data_analysis.Utilities.behavior_movies import Movie

#Allen
from allensdk.core.mouse_connectivity_cache import MouseConnectivityCache
mcc = MouseConnectivityCache(resolution=10)
str_tree = mcc.get_structure_tree()

#Read in allen CCF
ccfsum = pd.read_csv('/home/david.wyrick/projects/zap-n-zip/data/mouse_connectivity/ontology_v2.csv')

#User
sys.path.append(ProjDir)
import util
import plotting as usrplt


gc = gspread.service_account() # need a key file to access the account (step 2) 
sh = gc.open('Templeton-log_exp') # open the spreadsheet 

df = pd.DataFrame(sh.sheet1.get()) 
exp_table = df.T.set_index(0).T # put it in a nicely formatted dataframeexp_table.iloc[10:]
exp_table = exp_table.set_index('mouse_name')

##===== ============================ =====##
##===== Parse Command Line Arguments =====##
parser = argparse.ArgumentParser(description='single-cell-metrics')

##===== Data Options =====##
parser.add_argument('--mID',type=str, default='mouse709400',
                    help='mouse to perform analysis on')

parser.add_argument('--rec_name',type=str, default='aw_ket_2024-02-01_11-12-34',
                    help='experiment to perform analysis on')

parser.add_argument('--time_bin_ms',type=int, default=50,
                    help='time_bin_ms')

parser.add_argument('--window_t_min',type=int, default=15,
                    help='Window width (mins) to segment data into')


def calculate_avalanche_statistics(spks, time_bin, thresh=0.35):
    T, N = spks.shape

    #Binarize spiking data
    spks = np.array(spks > 0,dtype=int)

    #Identify threshold based on network activity
    network_activity = np.sum(spks,axis=1)
    network_thresh = np.nanpercentile(network_activity,thresh*100)

    #Identify transitions between avalanches
    active_time_bins = np.array(network_activity > network_thresh,dtype=bool)
    active_time_bins1 = np.array(network_activity > network_thresh,dtype=int)
    transitions = np.where(np.diff(active_time_bins))[0]
    transi_type = np.diff(active_time_bins1)[transitions]

    #identify first & last avalanche
    ii = np.where(transi_type == 1)[0][0]
    jj = np.where(transi_type == -1)[0][-1] + 1
    nAvalanches = len(transitions[ii:jj])/2
    assert len(transitions[ii:jj]) % 2 == 0

    #Get indices of avalanches
    indices = transitions[ii:jj].reshape(-1,2) + 1

    #Prune avalanches that are too short
    D_bin = indices[:,1] - indices[:,0]
    indices = indices[D_bin > 2]
    
    #Get avalanche length, size, and spikes
    D = (indices[:,1] - indices[:,0]) * time_bin
    S = np.array([np.sum(network_activity[iPair[0]:iPair[1]]) for iPair in indices])

    return D, S, indices

import powerlaw as pl
def fit_power_law(data):
    data_sorted = np.sort(data)
    # Function to fit truncated power law distribution
    def fit_truncated_powerlaw(data, S_min, S_max):
        truncated_data = data[(data >= S_min) & (data <= S_max)]
        # fit = st.powerlaw.fit(truncated_data)
        fit = pl.Fit(truncated_data,xmin=S_min,xmax=S_max)
        ks = fit.power_law.KS()
        tau = fit.power_law.alpha
        return fit, ks, tau

    # Convergence criterion
    def check_convergence(ks_stat, N_av):
        return ks_stat < 1 / np.sqrt(N_av)

    # Step 1: Determine Maximum Avalanche Size (S_max)
    S_max = np.max(data)

    # Initialize variables for iterative fitting
    converged = False
    S_max_iter = S_max
    N = len(data_sorted)
    iterations = 1
    S_min = data_sorted[0]

    ks_list = []
    for i in range(1,int(N/2)):
        S_max = data_sorted[-1*i]

        fit, ks_stat, tau = fit_truncated_powerlaw(data, S_min, S_max)
        ks_list.append(ks_stat)

    indy = np.argmin(np.array(ks_list))+1
    S_max = data_sorted[-1*i]
    fit, ks_stat, tau = fit_truncated_powerlaw(data, S_min, S_max)
    print(f'{indy} data points excluded')
    
    return fit, ks_stat, tau

if __name__ == '__main__':

    # Parse the arguments
    args = parser.parse_args()
    mID = args.mID
    rec_name = args.rec_name
    time_bin_ms = args.time_bin_ms
    time_bin = time_bin_ms/1000

    #Segement data into 15 minute windows
    window_t_min = args.window_t_min
    window_t = window_t_min*60
    
    #Extract row from templeton dataframe
    exp_df = exp_table.loc[exp_table.exp_name == rec_name]
    if type(exp_df) == pd.core.series.Series:
        exp_df = exp_df.to_frame().T

    drug_type = exp_df['drug'].values[0]
    stim_type = exp_df['stimulation'].values[0]
    print(f'Experiment type: {stim_type}, {drug_type}')

    #Define directories
    SaveDir = join(ServDir,'results','avalanches',mID,rec_name)
    if not os.path.exists(SaveDir):
        os.makedirs(SaveDir)

    PlotDir = join(ProjDir,'pptx','avalanches')
    if not os.path.exists(PlotDir):
        os.makedirs(PlotDir)

    #%% Upload the whole experiment and generate the global clock
    file_name = os.path.join(base_dir_server,mID,rec_name,'experiment1','recording1')
    exp = EEGexp(file_name, preprocess=False, make_stim_csv=False)

    #Load probe data
    probe_unit_data, probe_info, total_units = tbd_util.get_neuropixel_data(exp)
    probe_list = [x.replace('_sorted', '') for x in exp.experiment_data if 'probe' in x]

    #Get recording start time
    probei = probe_list[0]
    open_ephys_start = np.round(np.min(probe_unit_data[probei]['spike_times']))
    open_ephys_end = np.round(np.max(probe_unit_data[probei]['spike_times']))

    #For saline & psilocybin experiments, get injection times and types of injection
    if drug_type in ['saline', 'psilocybin']:
        injection_times = [float(exp_df['First injection time (s)'].values[0]),
                        float(exp_df['Second injection time (s)'].values[0])]
        
        #Determine injection type
        if 'psi' in rec_name:
            injection_types = ['sal1','psi']
        else:
            injection_types = ['sal1', 'sal2']
        injection_time_windows = None
    elif drug_type == 'ketanserin+psilocybin':
        injection_time_windows = np.array([np.array(exp_df['First injection time (s)'].values[0].split(','),dtype=float),
                                np.array(exp_df['Second injection time (s)'].values[0].split(','),dtype=float)])
        
        #Take second time in each window as "injection time"
        injection_times = [injection_time_windows[0,1],injection_time_windows[1,1]]

        injection_types = ['ket','psi']
    else:
        injection_times = None
        injection_time_windows = None

    #For isoflurane experiments, get iso level
    if drug_type == 'isoflurane':
        iso_level, iso_times = exp.load_analog_iso()
        iso_induction_times, iso_maintenance_times = exp.load_iso_times()
    else:
        iso_maintenance_times = None
        iso_induction_times = None


    # extract the timestamps of the selected stimuli
    try:
        stim_log = pd.read_csv(exp.stimulus_log_file)

        if 'circle' in np.unique(stim_log['stim_type']):
            vis_stim_exists = True
        else:
            vis_stim_exists = False
        if 'biphasic' in np.unique(stim_log['stim_type']):
            stim_exists = True
        else:
            stim_exists = False
    except:
        stim_exists = False
        vis_stim_exists = False
        stim_log = None
    #Load behavior
    run_ts, raw_run_signal, run_signal, run_signal_s, pupil_ts, pupil_radius, run_signal_p, run_signal_p_s, plot_pupil = util.get_behavioral_data(exp, mID, rec_name)
    f_run = interp1d(run_ts,run_signal)
    f_pupil = interp1d(pupil_ts,pupil_radius)

    #Get time windows for each epoch
    epoch_list, time_window_list = util.define_epochs_of_interest([open_ephys_start,open_ephys_end], drug_type, window_t_min=window_t_min, injection_times=injection_times,injection_time_windows=injection_time_windows, iso_induction_times=iso_induction_times, stim_log=stim_log)

    #Create colormap for epochs
    nCond = len(epoch_list)
    if drug_type == 'saline':
        cmap = sns.color_palette('Greens',nCond)
    elif drug_type == 'psilocybin':
        nPsi = np.sum(['psi' in fname for fname in epoch_list])
        nSal = nCond - nPsi
        if nSal == 0:
            cmap = sns.color_palette('Reds',nCond)
        elif nPsi == 0:
            cmap = sns.color_palette('Blues',nCond)
        else:
            cmap = np.concatenate((sns.color_palette('Blues',nSal),sns.color_palette('Reds',nPsi)))
    elif drug_type == 'isoflurane':
        nPre = np.sum(['pre' in fname for fname in epoch_list])
        nPost = nCond - nPre
        if nPre == 0:
            cmap = sns.color_palette('Oranges',nCond)
        elif nPost == 0:
            cmap = sns.color_palette('Blues',nCond)
        else:
            cmap = np.concatenate((sns.color_palette('Blues',nPre),sns.color_palette('Oranges',nPost)))

    elif drug_type == 'ketanserin+psilocybin':
        nPre = np.sum(['pre' in fname for fname in epoch_list])
        nKet = np.sum(['ket' in fname for fname in epoch_list])
        nPsi = np.sum(['psi' in fname for fname in epoch_list])

        cmap = np.concatenate((sns.color_palette('Blues',nPre),sns.color_palette('Oranges',1),sns.color_palette('Reds',nPsi)))
    else:
        cmap = sns.color_palette('Oranges',nCond)


    #Save time windows
    np.savez(os.path.join(SaveDir,f'time_windows.npz'),time_window_list=time_window_list,epoch_list=epoch_list)
    
    # Create a new PowerPoint presentation to save figures to
    prs = Presentation()

    #Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f'mID = {mID}; rec_name = {rec_name}'
    slide.placeholders[1].text = f'AVALANCHE!'
    
    #Read in neuropixel data 
    data_list, ts_list, neuron_indices, plot_tuple, FR_perblock = util.bin_spiking_data(probe_unit_data, time_window_list, time_bin=time_bin,fr_thresh=0) 
    boundaries_group, ticks_group, labels_group, celltypes, durations, layers, areas, groups, supergroups, order_by_group = plot_tuple

    area_sub = areas[neuron_indices]
    groups_sub = groups[neuron_indices]
    supergroups_sub = supergroups[neuron_indices]
    areas_ro = area_sub[order_by_group]
    groups_ro = groups_sub[order_by_group]
    supergroups_ro = supergroups_sub[order_by_group]

  
    powerlaw_list = []
    fig_o, axes_o = plt.subplots(1,3,figsize=(12,4))
    plt.suptitle(f'{mID}, {rec_name}')
    for ii, (ts, data) in enumerate(zip(ts_list,data_list)):
        data_s = data[:,neuron_indices]
        data_r = data_s[:,order_by_group]
        # data_g = data_r[:,supergroups_ro == 'CTX']#[:,groups_ro == g]
        data_g = data_r#[:,groups_ro == g]
        avalanche_durations, avalanche_sizes, indices = calculate_avalanche_statistics(data_g, time_bin)

        #Define bins to count avalanches
        duration_bins = np.histogram_bin_edges(avalanche_durations, bins='auto') #np.arange(time_bin,20,time_bin*5)
        size_bins = np.histogram_bin_edges(avalanche_sizes, bins='auto') #np.arange(0,np.max(avalanche_sizes)+100,100)
        fig, axes = plt.subplots(1,3,figsize=(12,4))
        # plt.suptitle(epoch_list[ii])
        #Plot
        ax = axes[0]
        ax.set_title('Avalanche size')
        ax.set_xlabel('Integrated network activity')

        counts, bins = np.histogram(avalanche_sizes, bins=size_bins,density=True)
        bin_centers = 0.5 * (bins[:-1] + bins[1:])
        ax.plot(bin_centers[counts > 0], counts[counts > 0], marker='o', markersize=5, color=cmap[ii], linestyle='')
        axes_o[0].plot(bin_centers[counts > 0], counts[counts > 0], marker='o', markersize=5, color=cmap[ii], linestyle='')
        # sns.kdeplot(avalanche_sizes,ax=ax,color=cmap[ii])
        # sns.kdeplot(avalanche_sizes,ax=axes_o[0],color=cmap[ii])
        # tau = fit_power_law(avalanche_sizes)

        # fit = pl.Fit(avalanche_sizes)
        # tau = fit.power_law.alpha
        # pdb.set_trace()
        # fit, min_ks, tau = fit_power_law(avalanche_sizes)
        fit = pl.Fit(avalanche_sizes,xmin=np.min(avalanche_sizes))
        tau = fit.power_law.alpha
        # Calculate y values using the fitted parameters
        y = size_bins**(-tau)
        # plt.gca()
        fit.power_law.plot_pdf(ax=ax)
        # Plot the fitted truncated power law distribution
        # ax.plot(size_bins, y, 'r-', lw=2, label=f'Fitted Truncated Power Law (tau={tau:.2f})')


        ax = axes[1]
        ax.set_title('Avalanche duration')
        ax.set_xlabel('Duration (s)')

        counts, bins = np.histogram(avalanche_durations, bins=duration_bins,density=True)
        bin_centers = 0.5 * (bins[:-1] + bins[1:])
        ax.plot(bin_centers[counts > 0], counts[counts > 0], marker='o', markersize=5, color=cmap[ii], linestyle='')
        axes_o[1].plot(bin_centers[counts > 0], counts[counts > 0], marker='o', markersize=5, color=cmap[ii], linestyle='')
        # sns.kdeplot(avalanche_durations,ax=ax,color=cmap[ii])
        # sns.kdeplot(avalanche_durations,ax=axes_o[1],color=cmap[ii])
        
        # fit = pl.Fit(avalanche_durations)
        # alpha = fit.power_law.alpha
        # fit, min_ks, alpha = fit_power_law(avalanche_durations)
        fit = pl.Fit(avalanche_durations,xmin=np.min(avalanche_durations))
        alpha = fit.power_law.alpha
        fit.power_law.plot_pdf(ax=ax)

        beta = (alpha-1)/(tau-1)
        ax = axes[2]
        ax.plot(avalanche_durations,avalanche_sizes,'o',color=cmap[ii])
        xlim = ax.get_xlim(); ylim = ax.get_ylim()
        x = np.arange(0,10,0.1)
        r = st.linregress(avalanche_durations,avalanche_sizes)
        y = r.slope*x + r.intercept
        ax.plot(x,y,'k--',label=f'beta = {beta:.3f}')

        # y = beta*r.slope*x + r.intercept
        # ax.plot(x,y,'-',color=usrplt.cc[8],label='Predicted')
        ax.set_xlim(xlim); ax.set_ylim(ylim)
        axes_o[2].plot(x,y,'-',color=cmap[ii])
        ax.set_xlim(xlim); ax.set_ylim(ylim)
        usrplt.adjust_spines(ax)
        ax.legend()
        ax.set_title('Scaling relation'); ax.set_xlabel('Duration (s)'); ax.set_ylabel('< S >')
        for ax in axes[:-1]:
            ax.set_xscale('log')
            ax.set_yscale('log')
            # ax.set_ylim([1e-10,1E2])
            usrplt.adjust_spines(ax)
        print(f'{epoch_list[ii]}: tau = {tau:.2f}, alpha = {alpha:.2f}, beta = {beta:.2f} slope = {r.slope:.2f}')
        plt.suptitle(f'{epoch_list[ii]}: tau = {tau:.2f}, alpha = {alpha:.2f}, beta = {beta:.2f} slope = {r.slope:.2f}',y=0.99)
        usrplt.save_fig_to_pptx(fig,prs)
        plt.close(fig)
        powerlaw_list.append([tau,alpha,beta,r.slope])
    axes_o[0].set_title('Avalanche size'); axes_o[0].set_xlabel('Integrated network activity')
    axes_o[1].set_title('Avalanche duration'); axes_o[1].set_xlabel('Duration (s)')
    axes_o[2].set_title('Scaling relation'); axes_o[2].set_xlabel('Duration (s)'); axes_o[2].set_ylabel('< S >')
    usrplt.adjust_spines(axes_o[2])

    for ax in axes_o[:-1]:
        ax.set_xscale('log')
        ax.set_yscale('log')
        # ax.set_ylim([1e-10,1E2])
        usrplt.adjust_spines(ax)
    usrplt.save_fig_to_pptx(fig_o,prs)

    powerlaw_stats = np.array(powerlaw_list)
    np.save(join(SaveDir,'powerlaw_stats.npy'),powerlaw_stats)
    # #Get rereferenced EEG data for plotting
    # bc = exp_df['EEG bad_channels'].values[0]
    # if (bc == 'none'):
    #     bad_channels = None
    #     eeg_ts, eeg_data, plot_eeg = util.get_preprocessed_eeg(exp, bad_channels)
    # elif (bc == 'all'):
    #     plot_eeg = False
    #     eeg_ts = np.array([np.nan])
    #     eeg_data = np.array([np.nan])
    # else:
    #     bad_channels = np.array(exp_df['EEG bad_channels'].values[0].split(','),dtype=int)
    #     eeg_ts, eeg_data, plot_eeg = util.get_preprocessed_eeg(exp, bad_channels)

    # #Align eeg data to start of spike raster
    # if plot_eeg:
    #     time_bin_e = 1/100
    #     eeg_tuple = (eeg_ts, eeg_data, time_bin_e)
    # else:
    #     eeg_tuple = None

    # #Align pupil data to start of spike raster
    # if plot_pupil:
    #     time_bin_p = 1/30
    #     pupil_tuple = (pupil_ts, pupil_radius, time_bin_p)
    # else:
    #     pupil_tuple = None

    # #Align running data to start of spike raster
    # time_bin_r = 1/100
    # run_tuple = (run_ts, run_signal, time_bin_r)

    prs.save(join(PlotDir,f'avalanche_distributions2_{mID}_{rec_name}.pptx'))
    print('DONE!!!')

# %%
