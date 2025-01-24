
ProjDir = '/home/david.wyrick/projects/zap-n-zip/'
DataDir = '/data/projects/zap-n-zip/data/VCNP_10/'
SaveDir = '/data/projects/zap-n-zip/results/ccm/VCNP_10/'

#Basea
import json, os, time, sys, argparse
from os.path import join
from glob import glob
import numpy as np
import pandas as pd
import scipy.stats as st
import pingouin as pg
from scipy.ndimage import gaussian_filter
from scipy.interpolate import interp1d

#Plot
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
from matplotlib.gridspec import GridSpec

#Network 
import networkx as nx
import networkx.algorithms.community as nx_comm
from networkx.algorithms.community import greedy_modularity_communities, modularity
from networkx.algorithms.efficiency_measures import global_efficiency, local_efficiency

#Allen
from allensdk.brain_observatory.ecephys.ecephys_project_cache import EcephysProjectCache
cache = EcephysProjectCache.from_warehouse(manifest=os.path.join('/data/visual_coding_ecephys', "manifest.json"))
analysis_metrics1 = cache.get_unit_analysis_metrics_by_session_type('brain_observatory_1.1')
analysis_metrics2 = cache.get_unit_analysis_metrics_by_session_type('functional_connectivity')
# all_metrics = pd.concat([analysis_metrics1,analysis_metrics2])

from allensdk.core.mouse_connectivity_cache import MouseConnectivityCache
mcc = MouseConnectivityCache(resolution=10)
str_tree = mcc.get_structure_tree()

#Read in allen CCF
ccfsum = pd.read_csv('/home/david.wyrick/projects/zap-n-zip/data/mouse_connectivity/ontology_v2.csv')

#User
sys.path.append(ProjDir)
import util
import plotting as usrplt

##===== ============================ =====##
##===== Parse Command Line Arguments =====##
parser = argparse.ArgumentParser(description='Networkx analysis on VCNP datasets')

##===== Data Options =====##
parser.add_argument('--session_id',type=str, default='mouse669117',
                    help='mouse to perform analysis on')

parser.add_argument('--last_run',type=int, default=None,
                    help='FCF run')

if __name__ == '__main__':
    nCond = 3
    cmap = sns.color_palette('Blues',3) #np.concatenate((,sns.color_palette('Greys',4)))
    filename_list = ['spont_block-0','spont_block-1','spont_block-2']#,'spont_block-3','dg75_block-0','dg75_block-1','dg75_block-2','dg75_block-3']
    time_bin_ms = 100

    # Parse the arguments
    args = parser.parse_args()
    session_id = args.session_id 
    
#    #Find folder with results inside
#     TempDir = os.path.join(ServDir,'results','ccm_xval',mID)
#     if last_run is None:
#         tmp_list = sorted(glob(join(TempDir,f'{rec_name}_run_*')))
#         last_run = int(tmp_list[-1][-1])
#     folder = f'{rec_name}_run_{last_run:02d}'

#     #Define save and plot directories
#     SaveDir = os.path.join(TempDir,folder)
    PlotDir = os.path.join(SaveDir,'plots')
    if not os.path.exists(PlotDir):
        os.makedirs(PlotDir)

    # #Load ccm parameters json
    # fpath = glob(join(SaveDir,'ccm_parameters*.json'))[0]
    # with open(fpath, 'r') as f:
    #     params = json.loads(f.read())

    # time_bin_ms = params['time_bin_ms']
    # time_bin = time_bin_ms/1000
    # delay = params['delay']
    # rand_proj = bool(params['rand_proj'])
    # xval = bool(params['xval'])
    # zscore = bool(params['zscore'])
    # nKfold = params['nKfold']
    # fr_thresh = params['fr_thresh'] 

    print(f'{session_id}')


    #Read in data
    try:
        results = np.load(join(DataDir,f'spike_counts_{time_bin_ms}-ms-bins_{session_id}.npz'),allow_pickle=True)
    except:
        print(f'session {session_id} has no results')
        exit()
    
    running_moments = results['running_moments']
    pupil_moments = results['pupil_moments']
    running_list = list(results['running_list'])
    pupil_list = list(results['pupil_list'])
    areas = results['areas']
    groups = results['groups']
    unit_ids = results['unit_ids']
    # ticks = results['ticks']
    # boundaries = results['boundaries']
    # labels = results['labels']
    # area_order = labels
    nNeurons = len(areas)

    uniq_areas, indices = np.unique(areas,return_index=True)
    sorted_indices = np.argsort(indices)
    area_order = uniq_areas[sorted_indices]

    counter = 0
    boundaries = [0]; ticks = []
    for a in area_order:
        nN_area = np.sum(areas == a)
        counter += nN_area
        ticks.append(boundaries[-1] + nN_area/2)
        boundaries.append(counter)
    labels = area_order

    #Read in FC results
    results = np.load(join(SaveDir,f'FCF_results_{session_id}.npz'))
    FCF_optimal = results['FCF_optimal'][:,:nNeurons,:nNeurons]
    correlation = results['correlation'][:,:nNeurons,:nNeurons]
    directionality = results['directionality'][:,:nNeurons,:nNeurons]
    complexity = results['complexity'][:,:nNeurons,:nNeurons]

    nCond, N, _ = FCF_optimal.shape

    #Calculate masked FCF optimal for mean kfold
    FCF_optimal_mask = FCF_optimal.copy()
    complexity_mask = complexity.copy()
    directionality_mask = directionality.copy()
    correlation_mask = correlation.copy()

    for ii, epoch in enumerate(filename_list):
        mask = FCF_optimal[ii] < 0.1
        FCF_optimal_mask[ii][mask] = np.nan
        complexity_mask[ii][mask] = np.nan
        directionality_mask[ii][mask] = np.nan
        correlation_mask[ii][mask] = np.nan

    # Create a new PowerPoint presentation to save figures to
    prs = Presentation()

    #Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f'Session = {session_id}'
    slide.placeholders[1].text = f'PLOTS!'

    ##--------------------------------------------------
    nPairs = []
    perc_FCF = []
    for ii, epoch in enumerate(filename_list):
        tmp = FCF_optimal[ii].copy()
        mask = FCF_optimal[ii] > 0.1
        perc_FCF.append(np.nanpercentile(tmp[mask].ravel(),75))
        nPairs.append(np.sum(mask)/(N**2))

    fig, ax = plt.subplots(figsize=(8,5))
    ax.plot(nPairs,'-k')
    for jj, fname in enumerate(filename_list):
        ax.plot(jj,nPairs[jj],'o',color=cmap[jj])
    ax.set_xticks(np.arange(len(filename_list)))
    ax.set_xticklabels(filename_list,rotation=30)
    ax.set_ylabel('Fraction of pairs with FC > 0.1')
    ax.set_xlabel('Epoch')

    ax2 = ax.twinx()
    ax2.plot(perc_FCF,'-',color=usrplt.cc[8])
    for jj, fname in enumerate(filename_list):
        ax2.plot(jj,perc_FCF[jj],'o',color=cmap[jj])
    ax2.set_ylabel('75th percentile FCF of significant pairs',color=usrplt.cc[8])
    usrplt.save_fig_to_pptx(fig, prs)

    ##--------------------------------------------------
    for x_array, xlabel in zip((pupil_moments,running_moments),('pupil radius (pixels)','running speed (cm/s)')):
        fig, axes = plt.subplots(1,2,figsize=(10,5))
        for ii in range(nCond):
            axes[0].scatter(x_array[ii,0],nPairs[ii],marker='.',facecolors="None",edgecolors=cmap[ii])
            axes[1].scatter(x_array[ii,0],perc_FCF[ii],marker='.',facecolors="None",edgecolors=cmap[ii])

        ax = axes[0]
        ax.set_xlabel(f'Mean {xlabel}')
        ax.set_ylabel('Fraction of pairs with FC > 0.1')

        ax = axes[1]
        ax.set_xlabel(f'Mean {xlabel}')
        ax.set_ylabel('75th percentile FCF of significant pairs')
        for ax in axes:
            usrplt.adjust_spines(ax)
        usrplt.save_fig_to_pptx(fig, prs)


    FC_thresh = 0.1
    areas_sub = areas#[neuron_indices]
    groups_sub, group_dict, graph_order, group_order, group_order_labels, supergroups = util.determine_groups(areas_sub)
    label_list = filename_list

    ##----- Plot in-network vs out-network Functional connectivity per AREA -----##
    for g in group_order:
        areas_in_group = group_dict[g]
        for a in areas_in_group:

            indy = np.where(areas_sub == a)[0]
            indy_out = np.where(areas_sub != a)[0]
            if len(indy) == 1:
                continue

            tmp_list = []
            for ii in range(nCond):

                # tmp_mat = FCF_optimal_mkf[ii].copy()
                tmp_mat = FCF_optimal[ii].copy()
                
                #Within area FC 
                FC = tmp_mat[indy][:,indy].ravel()
                mask = FC > FC_thresh
                n2 = len(FC[mask])
                tmp_list.append((np.repeat(filename_list[ii],n2),np.repeat(a,n2),np.repeat(g,n2),np.repeat('within',n2),FC[mask]))

                #Source area FC 
                FC = tmp_mat[indy_out][:,indy].ravel()
                mask = FC > FC_thresh
                n2 = len(FC[mask])
                tmp_list.append((np.repeat(filename_list[ii],n2),np.repeat(a,n2),np.repeat(g,n2),np.repeat('source',n2),FC[mask]))

                #Target area FC 
                FC = tmp_mat[indy][:,indy_out].ravel()
                mask = FC > FC_thresh
                n2 = len(FC[mask])
                tmp_list.append((np.repeat(filename_list[ii],n2),np.repeat(a,n2),np.repeat(g,n2),np.repeat('target',n2),FC[mask]))

            FC_df = pd.DataFrame(np.hstack(tmp_list).T,columns = ['epoch','area','group','type','FC'])
            FC_df = FC_df.astype({'epoch':str,'area':str,'group':str,'type':str,'FC':float})

            fig, axes = plt.subplots(1,3,figsize=(12,4),sharey=True)
            plt.suptitle(f'{a}, n = {len(indy)}',y=0.995)

            sns.barplot(data=FC_df.loc[FC_df.type == 'within'],y='FC',x='epoch',hue_order=filename_list,ax=axes[0],palette=list(cmap),legend=False)
            sns.barplot(data=FC_df.loc[FC_df.type == 'source'],y='FC',x='epoch',hue_order=filename_list,ax=axes[1],palette=list(cmap),legend=False)
            sns.barplot(data=FC_df.loc[FC_df.type == 'target'],y='FC',x='epoch',hue_order=filename_list,ax=axes[2],palette=list(cmap),legend=False)
            for ax in axes:
                usrplt.adjust_spines(ax)
                ax.set_xticklabels(label_list,ha='right',rotation=30)
                ax.set_ylim([0.05,0.5])
            axes[0].set_title('Within network FC')
            axes[1].set_title('Source area FC')
            axes[2].set_title('target area FC')
            usrplt.save_fig_to_pptx(fig, prs)

    vmax1 = np.round(np.nanpercentile(np.abs(correlation.ravel()),99),2)
    vmax2 = np.round(np.nanpercentile(np.abs(FCF_optimal.ravel()),99),2)
    vmax_fcf = np.round(np.max([vmax1,vmax2]),1)
 
    ##--------------------------------------------------
    for FCF_thresh in [0,0.1]:
        for normalize in [False,True]:
            tmp_list = []
            for jj in range(nCond):
  
                #Define adjacency matrix 
                mask = FCF_optimal[jj] > FCF_thresh
                adjacency_mat = np.array(mask,dtype=int)
                weighted_FC = FCF_optimal[jj].copy()
                weighted_FC[~mask] = 0

                #Creates graph using the data of the adjacency matrix or weighted directed graph
                G = nx.to_networkx_graph(adjacency_mat)
                DG = nx.DiGraph(weighted_FC)

                i_dict = dict(DG.in_degree(weight='weight'))
                o_dict = dict(DG.out_degree(weight='weight'))
                
                num_connections_in = np.sum(mask,axis=0)
                num_connections_out = np.sum(mask,axis=1)

                for ii, (a,g) in enumerate(zip(areas_sub,groups_sub)):
      
                    if normalize:
                        tmp_list.append((filename_list[jj],ii,a,g,'in',i_dict[ii]/num_connections_in[ii]))
                        tmp_list.append((filename_list[jj],ii,a,g,'out',o_dict[ii]/num_connections_out[ii]))
                    else:
                        tmp_list.append((filename_list[jj],ii,a,g,'in',i_dict[ii]))
                        tmp_list.append((filename_list[jj],ii,a,g,'out',o_dict[ii]))
                    # ratio = i_dict[ii]/o_dict[ii]
                    if o_dict[ii] != 0:
                        tmp_list.append((filename_list[jj],ii,a,g,'in/out',i_dict[ii]/o_dict[ii]))


            degree_df = pd.DataFrame(np.vstack(tmp_list),columns=['epoch','ii','area','group','type','degree'])
            degree_df = degree_df.astype({'epoch':str,'ii':int,'area':str,'group':str,'type':str,'degree':float})

            fig, axes = plt.subplots(3,1,figsize=(12,12),gridspec_kw={'hspace':0.5})
            if normalize:
                plt.suptitle(f'Average weighted in/out degree magnitudes per area\n connections exist for FC > {FCF_thresh}')
            else:
                plt.suptitle(f'Sum of weighted in/out degree magnitudes per area\n connections exist for FC > {FCF_thresh}')
            sns.barplot(data=degree_df.loc[degree_df.type == 'in'],x='area',order=area_order,y='degree',hue='epoch',palette=cmap,ax=axes[0],legend=False)
            sns.barplot(data=degree_df.loc[degree_df.type == 'out'],x='area',order=area_order,y='degree',hue='epoch',palette=cmap,ax=axes[1],legend=False)
            sns.barplot(data=degree_df.loc[degree_df.type == 'in/out'],x='area',order=area_order,y='degree',hue='epoch',palette=cmap,ax=axes[2],legend=False)

            axes[0].set_title('Weighted In-degree')
            axes[1].set_title('Weighted Out-degree')
            axes[2].set_title('Weighted in/out degree ratio')

            axes[2].hlines(1,*axes[2].get_xlim(),color='k',ls=':')
            axes[2].autoscale(tight=True)
            axes[2].set_ylim([0,3])
            for ax in axes:
                usrplt.adjust_spines(ax)
                ax.set_ylabel('Weighted degree')
                ax.set_xlabel('')
            axes[2].set_title('Degree ratio')
            axes[2].set_ylabel('Degree ratio')
            usrplt.save_fig_to_pptx(fig, prs)

    ##--------------------------------------------------
    position_list = []
    tmp = np.arange(nCond)
    xticks = []
    for ii in range(len(area_order)):
        position_list.append(tmp + (nCond+2)*ii)
        xticks.append((nCond+2)*ii+ nCond/2)
    positions = np.concatenate(position_list)

    FC_thresh = 0.1
    for a in area_order:

        for normalize in [False,True]:
            # for a in areas_in_group:
            indy = np.where(areas == a)[0]
            n = len(indy)
            fig, axes = plt.subplots(3,1,figsize=(12,12),gridspec_kw={'hspace':0.5})
            if normalize:
                plt.suptitle(f'Source area {a}, n = {n}; Normalized weighted degree\n connections exist for FC > 0.1')
            else:
                plt.suptitle(f'Source area {a}, n = {n}; Sum of weighted degrees\n connections exist for FC > 0.1')

            tmp_list = []
            counter = 0
            for a2 in area_order:
                indy2 = np.where(areas == a2)[0]
                for ii in range(nCond):
   
                    tmp_mat = FCF_optimal_mask[ii].copy()
                    in_degree = np.nansum(tmp_mat[indy2][:,indy],axis=0).ravel(); n1 = len(in_degree)
                    out_degree = np.nansum(tmp_mat[indy][:,indy2],axis=1).ravel(); n2 = len(out_degree)
                    
                    num_connections_in = np.nansum(np.isnan(tmp_mat[indy2][:,indy]),axis=0)
                    num_connections_out = np.nansum(np.isnan(tmp_mat[indy][:,indy2]),axis=1)
                    
                    if normalize:
                        in_degree = in_degree/num_connections_in
                        out_degree = out_degree/num_connections_out
                    tmp_list.append((np.repeat(filename_list[ii],n1),np.repeat(a,n1),np.repeat(a2,n1),np.repeat('in',n1),in_degree))
                    tmp_list.append((np.repeat(filename_list[ii],n1),np.repeat(a,n1),np.repeat(a2,n1),np.repeat('out',n1),out_degree))
                    tmp_list.append((np.repeat(filename_list[ii],n1),np.repeat(a,n1),np.repeat(a2,n1),np.repeat('in/out',n1),in_degree/out_degree))

            degree_df = pd.DataFrame(np.hstack(tmp_list).T,columns=['epoch','source_area','target_area','type','degree'])
            degree_df = degree_df.astype({'epoch':str,'source_area':str,'target_area':str,'type':str,'degree':float})
            
            sns.barplot(data=degree_df.loc[degree_df.type == 'in'],x='target_area',order=area_order,y='degree',hue='epoch',palette=list(cmap),ax=axes[0],legend=False)
            sns.barplot(data=degree_df.loc[degree_df.type == 'out'],x='target_area',order=area_order,y='degree',hue='epoch',palette=list(cmap),ax=axes[1],legend=False)
            sns.barplot(data=degree_df.loc[degree_df.type == 'in/out'],x='target_area',order=area_order,y='degree',hue='epoch',palette=list(cmap),ax=axes[2],legend=False)

            axes[0].set_title('Weighted In-degree')
            axes[1].set_title('Weighted Out-degree')
            axes[2].set_title('Weighted in/out degree ratio')

            for ax in axes:
                usrplt.adjust_spines(ax)
                ax.set_ylabel('Weighted degree')
                ax.set_xlabel('')
            axes[2].set_xlabel('Target area')
            axes[2].hlines(1,*ax.get_xlim(),color='k',ls=':')
            axes[2].autoscale(tight=True)
            axes[2].set_ylim([0,3])
            usrplt.save_fig_to_pptx(fig, prs)

    ##--------------------------------------------------
    fig_out, axes_out = plt.subplots(1,3,figsize=(12,4))
    area_comm_list = []
    for a in area_order:
        indy = np.where(areas == a)[0]
        area_comm_list.append(frozenset(indy))
    
    results = []; degree_dist_list = []
    for jj in range(nCond):
        print(filename_list[jj])

        #Define adjacency matrix 
        mask = FCF_optimal[jj] > FCF_thresh
        adjacency_mat = np.array(mask,dtype=int)
        N = adjacency_mat.shape[0]

        #Define directed graph
        weighted_FC = FCF_optimal[jj].copy()
        weighted_FC[~mask] = 0

        #Calculate cost of graph
        cost = np.sum(mask)/(N**2-N)
        results.append([filename_list[jj],'undirected_adjmat','cost',cost])

        #Calculate connectivity strength & diversity
        conn_strength = np.nanmean(weighted_FC[mask].ravel())
        conn_diversity = np.nanvar(weighted_FC[mask].ravel())
        results.append([filename_list[jj],'directed_weight','conn_strength',conn_strength])
        results.append([filename_list[jj],'directed_weight','conn_diversity',conn_diversity])

        #Creates graph using the data of the adjacency matrix or weighted directed graph
        UG_adjmat = nx.to_networkx_graph(adjacency_mat)
        DG_weight = nx.DiGraph(weighted_FC)

        #Define mask for correlation as well
        tmp_corr = correlation[jj].copy()
        tmp_corr[np.diag_indices(nNeurons)] = 0
        mask = tmp_corr > FCF_thresh
        tmp_corr[~mask] = 0
        UG_weight = nx.Graph(tmp_corr)

        #Calculate cost of graph
        cost = np.sum(mask)/(N**2-N)
        results.append([filename_list[jj],'undirected_weight','cost',cost])

        #Calculate connectivity strength & diversity
        conn_strength = np.nanmean(tmp_corr[mask].ravel())
        conn_diversity = np.nanvar(tmp_corr[mask].ravel())
        results.append([filename_list[jj],'undirected_weight','conn_strength',conn_strength])
        results.append([filename_list[jj],'undirected_weight','conn_diversity',conn_diversity])

        for ii, (g_type, g) in enumerate(zip(['undirected_weight','undirected_adjmat','directed_weight'],[UG_weight,UG_adjmat,DG_weight])):
            #Get degree distribution
            degree_dist = list(dict(g.degree()).values())
            sns.kdeplot(degree_dist,ax=axes_out[ii],color=cmap[jj],label=filename_list[jj])

            #Find communities and calculate modularity
            comm_list = greedy_modularity_communities(g,weight='weight')
            mod = modularity(g,comm_list)
            results.append([filename_list[jj],g_type,'modularity',mod])
            results.append([filename_list[jj],g_type,'num_communities',len(comm_list)])
            
            #Calculate modulatity of group definitions of community
            mod = modularity(g,area_comm_list)
            results.append([filename_list[jj],g_type,'modularity_group',mod])

            #Plot communities
            indy_order = []
            boundaries_comm = []
            ticks_comm = []
            labels_comm = []
            counter = 0

            nNeurons_per_comm = [len(s) for s in comm_list]
            nGroups_per_comm = []
            homogeneity = []
            for kk, s in enumerate(comm_list):
                n_indices = list(s)
                indy_order.append(n_indices)
                nN_per_comm = len(n_indices)
                counter += nN_per_comm
                groups_comm = groups[n_indices]
                nN_per_group = []
                nGroups_per_comm.append(len(np.unique(groups_comm)))
                tmp = []
                for g1 in groups_comm:
                    for g2 in groups_comm:
                        if g1 == g2:
                            tmp.append(1)
                        else:
                            tmp.append(0)
                homogeneity.append(np.mean(tmp))
                
                if nN_per_comm > 5:
                    ticks_comm.append(counter-nN_per_comm/2)
                    boundaries_comm.append(counter)
                    labels_comm.append(kk)
            indy_order = np.concatenate(indy_order)

            if g_type == 'directed_weight':
                fig, axes = plt.subplots(1,3,figsize=(12,4),gridspec_kw={'wspace':0.25})
                plt.suptitle(f'FCF: {filename_list[jj]}, Graph type: {g_type}')

                usrplt.visualize_matrix(FCF_optimal[jj],ax=axes[0],cmap='viridis',title='Sorted by area',clims=[0,vmax_fcf],center=None,ticks=ticks,labels=labels,boundaries=boundaries)
                usrplt.visualize_matrix(FCF_optimal[jj][indy_order][:,indy_order],ax=axes[1],cmap='viridis',title=f'Sorted into communities\nModularity {mod:.3f}',clims=[0,vmax_fcf],center=None,ticks=ticks_comm,labels=labels_comm,boundaries=boundaries_comm)

                # sns.barplot(y='community',x='N',hue='group',ax=axes[-1],data=comm_df)
                ax = axes[-1]
                ax.set_title('Community detection',fontsize=16,fontweight='bold')
                ax.plot(nNeurons_per_comm,'-o')
                ax.set_ylabel('# of neurons per community')
                ax.set_xlabel('Community')
                # usrplt.adjust_spines(ax)

                ax2 = ax.twinx()
                ax2.plot(homogeneity,'o-',color=usrplt.cc[3])
                ax2.set_ylabel('homogeneity of community',color=usrplt.cc[3])
                ax2.set_xlim([-0.1,len(boundaries_comm)+1])
                ax2.set_ylim([-0.1,1.1])
                usrplt.save_fig_to_pptx(fig, prs)

                tmp_list = []
                for kk, s in enumerate(comm_list):
                    n_indices = list(s)
                    nN_per_comm = len(n_indices)
                    counter += nN_per_comm
                    groups_comm = groups[n_indices]
                    
                    if nN_per_comm > 5:

                        for g1 in group_order:
                            nN_per_group= np.sum(groups_comm == g1)
                            tmp_list.append((kk,g1,np.sum(groups_comm == g1)))

                n_df = pd.DataFrame(np.stack(tmp_list),columns=['community','group','N'])
                n_df = n_df.astype({'community':str,'group':str,'N':float})

                fig, ax = plt.subplots(figsize=(10,4))
                plt.suptitle(f'FCF: {filename_list[jj]}, Graph type: {g_type}')
                sns.barplot(x='community',y='N',hue='group',data=n_df,ax=ax)
                ax.set_xlabel('Community ID')
                ax.set_ylabel('# of neurons')
                usrplt.save_fig_to_pptx(fig, prs)

            #Get degree assortativity
            #Assortativity in a network refers to the tendency of nodes to connect with other ‘similar’ nodes over ‘dissimilar’ nodes.
            #two nodes are ‘similar’ with respect to a property if they have the same value of that property
            DAC = nx.degree_assortativity_coefficient(g)
            results.append([filename_list[jj],g_type,'DAC',DAC])

            #Compute graph transitivity, the fraction of all possible triangles present in G.
            trs = nx.transitivity(g)
            results.append([filename_list[jj],g_type,'transitivity',trs])

            #compute the average clustering coefficient for the graph G.
            avg_coeff = nx.average_clustering(g,weight='weight')
            results.append([filename_list[jj],g_type,'average_clustering',avg_coeff])

            if 'undirected' in g_type:
                #Calculate global efficiency 
                #The efficiency of a pair of nodes in a graph is the multiplicative inverse of the shortest path distance between the nodes.
                #The average global efficiency of a graph is the average efficiency of all pairs of nodes
                gE = nx.global_efficiency(g)
                results.append([filename_list[jj],g_type,'gE',gE])

                #The local efficiency of a node in the graph is the average global efficiency of the subgraph induced by the neighbors of the node. 
                #The average local efficiency is the average of the local efficiencies of each node
                lE = nx.local_efficiency(g)
                results.append([filename_list[jj],g_type,'lE',lE])

                #Calculate whether the network is "whole" or disjointed into multiple subnetworks
                nCC = nx.number_connected_components(g)
                results.append([filename_list[jj],g_type,'nCC',nCC])
                
                if nx.is_connected(g):
                    #returns the average shortest path length.
                    sp = nx.average_shortest_path_length(g,weight='weight')
                    results.append([filename_list[jj],g_type,'average_shortest_path_length',sp])

                    #Compute the non-randomness of graph G.
                    nr_tuple = nx.non_randomness(g,weight='weight')
                    results.append([filename_list[jj],g_type,'non_randomness',nr_tuple[1]])
            
                #Perform targeted attack
                steps = util.targeted_attack(g.copy())
                results.append([filename_list[jj],g_type,'targeted_attack',steps/N])

                #Perform random attack
                tmp_list = []
                for iRand in range(100):
                    steps = util.random_attack(g.copy())
                    tmp_list.append(steps)
                results.append([filename_list[jj],g_type,'random_attack',np.mean(tmp_list)/N])

        #Flow hierarchy is defined as the fraction of edges not participating in cycles in a directed graph
        h = nx.flow_hierarchy(DG_weight)
        results.append([filename_list[jj],'directed_weight','flow_hierarchy',h])

        #Compute the reciprocity in a directed graph.
        #The reciprocity of a directed graph is defined as the ratio of the number of edges pointing in both directions to the total number of edges in the graph. Formally,
        r = nx.overall_reciprocity(DG_weight)
        results.append([filename_list[jj],'directed_weight','reciprocity',r])

    for ax in axes_out:
        ax.set_xlim(xmin=0)
        usrplt.adjust_spines(ax)

    for ii, ti in enumerate(['undirected_weight','undirected_adjmat','directed_weight']):
        axes_out[ii].set_title(ti)
    usrplt.save_fig_to_pptx(fig, prs)

    network_df = pd.DataFrame(np.vstack(results),columns=['epoch','g_type','algorithm','metric'])
    network_df = network_df.astype({'epoch':str,'g_type':str,'algorithm':str,'metric':float})
    network_df.to_csv(join(SaveDir,f'network-metrics_{session_id}.csv'))

    try:
        algorithms = np.unique(network_df['algorithm'])
        ls_list = ['-',':','--']
        for a in algorithms:
            fig, ax = plt.subplots(figsize=(8,5))
            ax.set_title(f'Algorithm: {a}')
            for jj, g in enumerate(['undirected_weight','undirected_adjmat','directed_weight']):
                sub_df = network_df.loc[(network_df.g_type == g) & (network_df.algorithm == a)]
                if len(sub_df) == 0:
                    continue
                sns.pointplot(x='epoch',y='metric',data=sub_df,palette=cmap)

                ys = []
                for e in filename_list:
                    ys.append(np.nanmean(sub_df.loc[sub_df.epoch == e]['metric']))
                ax.plot(ys,'k',zorder=0,ls=ls_list[jj],label=g)
                ax.set_xticklabels(filename_list,rotation=30,ha='right')
            ax.legend()
            usrplt.adjust_spines(ax)
            usrplt.save_fig_to_pptx(fig, prs)
    except:
        prs.save(join(PlotDir,f'FCF_figs_20240214_{session_id}.pptx'))
        exit()

    prs.save(join(PlotDir,f'FCF_figs_20240214_{session_id}.pptx'))
    print('DONE!!!')
    # exit()

