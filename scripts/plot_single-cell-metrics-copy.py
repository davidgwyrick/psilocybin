base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
base_dir_local = '/data/tiny-blue-dot/zap-n-zip/EEG_exp'; base_dir = base_dir_server
ProjDir = '/home/david.wyrick/projects/zap-n-zip/'
ServDir = '/data/projects/zap-n-zip/'
BehavDir = '/home/david.wyrick/projects/zap-n-zip/plots/behavior'

#Basea
import gspread
from os.path import join
from glob import glob
import json, os, time, sys, argparse
import numpy as np
import pandas as pd
import scipy.stats as st
import pingouin as pg

#Plot
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
import seaborn as sns

#Project
from tbd_eeg.data_analysis.eegutils import EEGexp
import tbd_eeg.data_analysis.Utilities.utilities as tbd_util
from tbd_eeg.data_analysis.Utilities.behavior_movies import Movie

#Allen
from allensdk.core.mouse_connectivity_cache import MouseConnectivityCache
mcc = MouseConnectivityCache(resolution=10)
str_tree = mcc.get_structure_tree()

#Read in allen CCF
ccfsum = pd.read_csv('/home/david.wyrick/projects/zap-n-zip/data/mouse_connectivity/ontology_v2.csv')

#User
sys.path.append(ProjDir)
import util
import plotting as usrplt

#Define behavioral states
behavior_ranges = {0: [0,1], 1: [1,500]}#, 3:[30,500]}
behavior_dict = { 0: 'rest (<1cm/s)', 1: 'active (>1cm/s)'}
behavior_dict2 = { 0: 'rest (<1cm/s)', 1: 'active (>1cm/s)', 2: 'all (>0cm/s)'}
behavior_strs2 = ['rest','active']
behavior_strs = list(behavior_dict.values())
nBehaviors = len(behavior_strs)

#Define windows to calculate firing rate
stim_window_dict = {'spontaneous': [], 'evoked': [],'pre-rebound': [], 'rebound': [], 'post-rebound': [], 'isi': [],'visual': []}
stim_strs = ['spontaneous','evoked','pre-rebound','rebound','post-rebound','isi_biphasic','visual','isi_visual']

visual_stim = ['circle','natural_scene']
spont_types = ['spontaneous','isi_biphasic','isi_visual']
evoked_windows = [[.002,.025],[0.025,0.075],[.075,.3],[.3,1]]
evoked_strings = ['evoked','pre-rebound','rebound','post-rebound']

gc = gspread.service_account() 
sh = gc.open('Templeton-log_exp')

df = pd.DataFrame(sh.sheet1.get()) 
exp_table = df.T.set_index(0).T 
exp_table = exp_table.set_index('mouse_name')

th_dict = {'SM-TH': ['PO','VAL','VPL','VPM','VM'], 
'VIS-TH': ['LP','LGN','LGd','LGd-co','LGd-sh','LGd-ip'],
'ANT-TH': ['AV','AMd','AMv','AD','MD','MDm','MDc','MDl','RE','RH','CM','LD', 'CL'],
'TH': ['Eth', 'IAD', 'IGL', 'IntG', 'LGv','MGd', 'MGm', 'MGv', 'PCN', 'PF', 'PIL', 'PoT', 'SGN','SPFp', 'TH','LH'],'RT': ['RT']}

##===== ============================ =====##
##===== Parse Command Line Arguments =====##
parser = argparse.ArgumentParser(description='plot_single-cell-metrics')

##===== Data Options =====##
parser.add_argument('--run_num',type=int, default=2,
                    help='run_num')

parser.add_argument('--plot_only_sig_mod',type=int, default=0,
                    help='Plot all neurons or only significantly modulated ones')

if __name__ == '__main__':

    pval_thresh = 0.01
    # Parse the arguments
    args = parser.parse_args()
    run_num = args.run_num
    plot_only_sig_mod = bool(args.plot_only_sig_mod)
    FR_df_list = []; dFR_df_list = []
    area_list = []; group_list = []; mesogroup_list = []; drug_type_list = []; 
    celltype_list = [];supergroup_list = []; duration_list =[];neuron_idx_list = []
    layer_list = []

    #Loop over experiments and read in single cell metrics
    for mID, row in exp_table.iterrows():
        rec_name = row.exp_name
        drug_type = row.drug
        stimulation = row.stimulation
        
        #Skip urethane and isoflurane experiments
        if drug_type in ['urethane','isoflurane','ketanserin+psilocybin']:
            continue
        
        #Skip new stimulation experiments
        if 'stim_train' in rec_name:
            continue
        
        #Skip DBA experiments 
        if mID == 'mouse724057':
            continue

        # if mID in ['mouse735051','mouse735052']:
        #     continue

        if 'electrical' in stimulation:
            st_type = 'electrical'
        elif 'spontaneous' in stimulation:
            st_type = 'spontaneous'

        #Read in single cell metrics
        SaveDir = join(ServDir,'results','FR_ISI',mID,rec_name,'metrics_run_{:02d}'.format(run_num))
        fpath = join(SaveDir,f'single-cell-metrics_spont_{rec_name}.npz')
        if os.path.exists(fpath) == False:
            continue

        #Read in spontaneous firing rate data
        results = np.load(fpath)
        moments_spont = results['moments_spont']
        pvalues_spont = results['pvalues_spont']
        mod_idx_spont = results['mod_idx_spont']
        numspks_spont = results['numspks_spont']
        combos_spont  = results['combos_spont']
        nComps_spont  = len(combos_spont)

        #Read in time windows & area information
        d = np.load(join(SaveDir,f'time_windows_{rec_name}.npz'),allow_pickle=True)
        epoch_windows = d['epoch_windows']
        epoch_names = d['epoch_names']
        block_labels = d['block_labels']
        layers = d['layers']
        celltypes = d['celltypes']
        durations = d['durations']
        areas = d['areas']
        groups = d['groups']
        mesogroups = d['mesogroups']
        supergroups = d['supergroups']
        nT_behavior = d['nT_behavior']
        nT_behavior2 = d['nT_behavior2']
        nNeurons = len(areas)

        #Fix some area/group/celltype definitions
        indy = np.where(groups == 'ACA')[0]
        if len(indy) > 0:
            mesogroups[indy] = 'CTX_frontal'

        indy = np.where(groups == 'RSP')[0]
        if len(indy) > 0:
            mesogroups[indy] = 'CTX_sensory'

        block_labels2 = []
        for b in block_labels:
            if 'pre_inj' in b:
                block_labels2.append('pre_inj')
            elif (b == 'post_sal1_inj') | (b == 'post_ket_inj'):
                block_labels2.append('post_1st_inj')
            else:
                block_labels2.append(b)
        block_labels = np.array(block_labels2)

        groups2 = []
        for a, g, sg in zip(areas, groups, supergroups):
            if sg == 'TH':
                if a in th_dict['SM-TH']:
                    groups2.append('SM-TH')
                elif a in th_dict['VIS-TH']:
                    groups2.append('VIS-TH')
                elif a in th_dict['ANT-TH']:
                    groups2.append('ANT-TH')
                elif a in th_dict['RT']:
                    groups2.append('RT')
                else:
                    groups2.append('TH')
                # else:
                #     raise Exception(f'Unknown thalamic area/group: {a}')
            else:
                groups2.append(g)
        groups = np.array(groups2)

        indy = np.where(supergroups == 'TH')[0]
        if len(indy) > 0:
            indy = np.where((durations < 0.35) & (supergroups == 'TH'))[0]
            celltypes[indy] = 'FS'

            indy = np.where((durations >= 0.45) & (supergroups == 'TH'))[0]
            celltypes[indy] = 'RS'

            indy = np.where((durations >= 0.35) & (durations < 0.45) & (supergroups == 'TH'))[0]
            celltypes[indy] = 'unclassified'
            
        #Ensure all neurons in RT are inhibitory
        indy = np.where(areas == 'RT')[0]
        if len(indy) > 0:
            mesogroups[indy] = 'TH_RT'
            celltypes[indy] = 'FS'
        
        #Make sure pre_inj period has at least 2min of data
        if (nT_behavior[0,-1] < 60*2):
            print(f'{mID}, {rec_name}: pre_inj period too short: {nT_behavior[0,-1]:.1f} sec')
            continue

        #Save for cell dataframe
        drug_type_list.append(np.repeat(drug_type,nNeurons))
        area_list.append(areas)
        duration_list.append(durations)
        group_list.append(groups)
        mesogroup_list.append(mesogroups)
        supergroup_list.append(supergroups)
        celltype_list.append(celltypes)
        layer_list.append(layers)
        neuron_idx_list.append(np.array([f'{rec_name}_{iN}' for iN in range(nNeurons)]))

        #Read in preprocessed dataframe
        FR_path = join(SaveDir,f'FR_spont_df_{rec_name}.csv')
        if os.path.exists(FR_path):
            FR_spont_df = pd.read_csv(FR_path,dtype={'mID':str, 'rec_name':str, 'drug_type':str, 'neuron_index':str, 'celltype':str, 'duration': float,
                                                    'area':str, 'layer':str, 'group':str, 'mesogroup':str, 'supergroup':str, 'behavior':str,
                                                    'epoch':str, 'block':str, 'num_spks': float, 'stim_type':str, 'stat':str, 'val': float, 'st_type':str})
            FR_df_list.append(FR_spont_df)
            if nComps_spont == 0:
                print(f'{mID}, {rec_name} processed')
                continue
        
        dFR_path = join(SaveDir,f'dFR_spont_df_{rec_name}.csv')
        if os.path.exists(dFR_path):
            dFR_spont_df = pd.read_csv(dFR_path,dtype={'mID':str, 'rec_name':str, 'drug_type':str, 'neuron_index':str, 'celltype':str, 'duration':float,
                                                        'area':str, 'layer':str, 'group':str, 'mesogroup':str, 'supergroup':str, 'behavior':str,
                                                        'epoch_i':str, 'block_i':str, 'epoch_j':str, 'block_j':str, 'stat':str, 'val_i':float, 'val_j':float,
                                                        'delta':float, 'pval1':float, 'sig1':float, 'pval2':float, 'sig2':float,'sig':float,'mod':float, 'abs_mod':float, 'abs_delta':float, 'st_type':str})
            dFR_df_list.append(dFR_spont_df)
            print(f'{mID}, {rec_name} processed')
            continue

        #Create dataframes for raw firing rate
        FR_list = []
        for iN, (ct, dur, a, l, g, mg, sg) in enumerate(zip(celltypes,durations,areas,layers,groups,mesogroups,supergroups)):
            for iB, behav_str in behavior_dict2.items():
                for iE, (tW, epoch, block) in enumerate(zip(epoch_windows,epoch_names,block_labels)):
                    neuron_idx = f'{rec_name}_{iN}'
                    FR_list.append((mID,rec_name,drug_type,neuron_idx,ct,dur,a,l,g,mg,sg,behav_str,epoch,block,numspks_spont[iN,iB,iE],'spont','mean_FR',moments_spont[iN,iB,iE,0]))
                    FR_list.append((mID,rec_name,drug_type,neuron_idx,ct,dur,a,l,g,mg,sg,behav_str,epoch,block,numspks_spont[iN,iB,iE],'spont','CV_ISI',moments_spont[iN,iB,iE,1]))
        FR_spont_df = pd.DataFrame(np.stack(FR_list),columns=['mID','rec_name','drug_type','neuron_index','celltype','duration','area','layer','group','mesogroup','supergroup','behavior','epoch','block','num_spks','stim_type','stat','val'])
        FR_spont_df = FR_spont_df.astype({'duration': float,'num_spks': float, 'val': float})
        FR_spont_df['st_type'] = st_type
        FR_df_list.append(FR_spont_df)
        FR_spont_df.to_csv(join(SaveDir,f'FR_spont_df_{rec_name}.csv'),index=False)

        #Create dataframes for change in firing rate
        if nComps_spont == 0:
            print(f'{mID}, {rec_name} processed')
            continue
        dFR_list = []
        for iN, (ct, dur, a, l, g, mg, sg) in enumerate(zip(celltypes,durations,areas,layers,groups,mesogroups,supergroups)):
            neuron_idx = f'{rec_name}_{iN}'
            for iB, behav_str in behavior_dict2.items():
                for ii, (i,j) in enumerate(combos_spont):
                    epoch_i = epoch_names[i]; epoch_j = epoch_names[j]
                    block_i = block_labels[i]; block_j = block_labels[j]
                    if (block_j == 'post_sal1_inj') | (block_j == 'post_ket_inj'):
                        block_j = 'post_1st_inj'
                    
                    #Calculate modulation index of firing rate
                    FR_i = moments_spont[iN,iB,i,0]; FR_j = moments_spont[iN,iB,j,0]

                    if np.isnan(FR_i) | np.isnan(FR_j):
                        continue
                    MI_FR = (FR_j - FR_i)/(FR_j + FR_i)

                    #Calculate modulation index of CV (not sure if this is valid)
                    CV_i = moments_spont[iN,iB,i,1]; CV_j = moments_spont[iN,iB,j,1]
                    MI_CV = (CV_j - CV_i)/(CV_j + CV_i)

                    assert mod_idx_spont[iN,iB,ii,0] == MI_FR
                    assert mod_idx_spont[iN,iB,ii,1] == MI_CV

                    #Get p-value of ISI distributions
                    pval = pvalues_spont[iN,iB,ii,0]
                    if pval < pval_thresh:
                        sig = 1
                    else:
                        sig = 0
                    #Get p-value of ISI distributions
                    pval1 = pvalues_spont[iN,iB,ii,0]
                    pval2 = pvalues_spont[iN,iB,ii,1]
                    sig1 = 1 if pval1 < pval_thresh else 0
                    sig2 = 1 if pval2 < pval_thresh else 0
                    sig = 1 if (sig1 + sig2) == 2 else 0

                    dFR_list.append((mID,rec_name,drug_type,neuron_idx,ct,dur,a,l,g,mg,sg,behav_str,epoch_i,block_i,epoch_j,block_j,'mean_FR',FR_i,FR_j,FR_j-FR_i,pval1,sig1,pval2,sig2,sig,mod_idx_spont[iN,iB,ii,0]))
                    dFR_list.append((mID,rec_name,drug_type,neuron_idx,ct,dur,a,l,g,mg,sg,behav_str,epoch_i,block_i,epoch_j,block_j,'CV_ISI',CV_i,CV_j,CV_j-CV_i,pval1,sig1,pval2,sig2,sig,mod_idx_spont[iN,iB,ii,1]))

        dFR_spont_df = pd.DataFrame(np.stack(dFR_list),columns=['mID','rec_name','drug_type','neuron_index','celltype','duration','area','layer','group','mesogroup','supergroup','behavior','epoch_i','block_i','epoch_j','block_j','stat','val_i','val_j','delta','pval1','sig1','pval2','sig2','sig','mod'])
        dFR_spont_df = dFR_spont_df.astype({'duration': float,'val_i': float, 'val_j': float,'delta': float, 'mod': float,'pval1': float,'sig1': int,'pval2': float,'sig2': int,'sig':int})
        dFR_spont_df['abs_mod'] = np.abs(dFR_spont_df['mod'])
        dFR_spont_df['abs_delta'] = np.abs(dFR_spont_df['delta'])
        dFR_spont_df['st_type'] = st_type
        dFR_df_list.append(dFR_spont_df)
        dFR_spont_df.to_csv(join(SaveDir,f'dFR_spont_df_{rec_name}.csv'),index=False)
        print(f'{mID}, {rec_name} processed')

    FR_spont_df = pd.concat(FR_df_list)
    dFR_spont_df = pd.concat(dFR_df_list)
    FR_spont_df = FR_spont_df.reset_index(drop=True)
    dFR_spont_df = dFR_spont_df.reset_index(drop=True)
    cells_df = pd.DataFrame(np.vstack((np.concatenate(drug_type_list),np.concatenate(neuron_idx_list),np.concatenate(duration_list),np.concatenate(celltype_list),np.concatenate(area_list),np.concatenate(layer_list), np.concatenate(group_list), np.concatenate(mesogroup_list), np.concatenate(supergroup_list))).T,columns=['drug_type','neuron_idx','duration','celltype','area','layer','group','mesogroup','supergroup'])

    ##------------------------------------------
    #Save dataframes & create save directory
    # if plot_only_sig_mod:
    #     print('Plotting only significantly modulated neurons')
    #     BaseDir = join(ServDir,'results','FR_ISI','all_sig_neurons','metrics_run_{:02d}'.format(run_num))
    # else:
    #     print('Plotting all neurons')
    #     BaseDir = join(ServDir,'results','FR_ISI','all_neurons','metrics_run_{:02d}'.format(run_num))

    #Plot only spontaneous experiments
    BaseDir = join(ServDir,'results','FR_ISI','all_spont_neurons_20240903','metrics_run_{:02d}'.format(run_num))
    FR_spont_df = FR_spont_df.loc[FR_spont_df.st_type == 'spontaneous']
    dFR_spont_df = dFR_spont_df.loc[dFR_spont_df.st_type == 'spontaneous']
    print(f'Saving dataframes to {BaseDir}')

    if not os.path.exists(BaseDir):
        os.makedirs(BaseDir)

    blocks = ['post_1st_inj','0_15','15_30','30_45','45_60','60_75','75_90']
    # blocks = ['0_15','15_30','30_45','45_60','60_75','75_90']
    drug_order = ['saline','psilocybin']#,'ketanserin+psilocybin']
    cmap_macro = sns.xkcd_palette(['cobalt blue','darkish red'])#,'dusty orange'])
    stat_list = ['mean_FR','CV_ISI']
    stat_title_list = ['Mean firing rate (Hz)','Coefficient of variation of ISI']
    bstrs_list = ['rest','run','all']
    layer_list = ['1','2/3','4','5','6']

    mesogroup_list = ['TH_core','TH_matrix','TH_intralaminar','CTX_frontal','CTX_sensory','HIP','STR']
    CTX_mesogroup_list = np.array(['CTX_frontal','CTX_sensory'])
    group_list = ['SM-TH','VIS-TH','ANT-TH','TH','MO','ILA','PL','ORB','SSp','VIS','OLF','HIP','STR']
    CTX_group_list = ['MO','ILA','PL','ORB','SSp','VIS','OLF']
    TH_group_list = ['SM-TH','VIS-TH','ANT-TH','TH','RT','HIP','STR']
    epoch_i = 'pre_inj'; block_i = 'pre_inj'
    # epoch_i = 'post_1st_inj'; block_i = 'post_1st_inj'

    nPlots = 0
    ##------------------------------------------
    # Create a new powerPoint presentations to save figures to
    prs_list = []
    for iS, stat in enumerate(stat_list):
        prs = Presentation()

        #Add title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f'All Mice, neurons organized in "mesogroups" level'
        slide.placeholders[1].text = f'Change in {stat_title_list[iS]} due to saline or psilocybin injection'
        prs_list.append(prs)

    ##------------------------------------------
    #Plot MI over time per mesogroup, celltype, behavior
    print('Plotting mesogroup barplots over time, per celltype, behavior, and mesogroup')
    for iS, stat in enumerate(stat_list):
        for ct in ['RS','FS']:
            if ct == 'FS':
                mesogroup_list = np.array(['TH_core','TH_matrix','TH_RT','CTX_frontal','CTX_sensory','HIP','STR'])
            else:
                mesogroup_list = np.array(['TH_core','TH_matrix','TH_intralaminar','CTX_frontal','CTX_sensory','HIP','STR'])
            for iB, b_str in behavior_dict2.items():
                PlotDir = join(BaseDir,'mesogroup',bstrs_list[iB])
                if not os.path.exists(PlotDir):
                    os.makedirs(PlotDir)
                for mg in mesogroup_list:
                
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.mesogroup == mg) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                    if len(sub_df) < 5:
                        continue
                    fig, axes = plt.subplots(3,1,figsize=(10,15))
                    plt.suptitle(f'Change in {stat_title_list[iS]} due to saline or psilocybin injection\nBaseline = {block_i}, {mg}, {ct}, {b_str}',y=0.925)

                    #Fraction significantly modulated
                    ax = axes[0]
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.mesogroup == mg) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]

                    iX = 0; xticks = []
                    for i, b in enumerate(blocks):
                        for jj, d in enumerate(drug_order):
                            sub_sub_df = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == d)]

                            n = len(sub_sub_df)
                            if n == 0:
                                continue
                            nSig = np.sum(sub_sub_df['sig'].values)
                            fracSig = nSig/n
                            pos_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] > 0)])/n
                            neg_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] < 0)])/n
                            not_mod = len(sub_sub_df.loc[sub_sub_df.sig == 0])/n

                            h1 = ax.bar(iX, pos_mod, color=sns.xkcd_palette(['light red']), label='Positively Modulated')
                            h2 = ax.bar(iX, not_mod, bottom=pos_mod, color=sns.xkcd_palette(['silver']), label='Not Modulated')
                            h3 = ax.bar(iX, neg_mod, bottom=pos_mod+not_mod, color=sns.xkcd_palette(['medium blue']), label='Negatively Modulated')
                            
                            n = len(np.unique(sub_sub_df['neuron_index']))
                            ax.text(iX,0.025,f'{n}',ha='center',fontsize=8,rotation=90)
                            if jj == 1:
                                xticks.append(iX-0.5)
                            iX += 1
                        iX += 1

                    usrplt.adjust_spines(ax)
                    plt.autoscale(tight=True)
                    xticks = np.arange(1,3*len(blocks),3)-0.5
                    ax.set_xticks(xticks)
                    ax.set_xticklabels(blocks,rotation=0)
                    ax.set_ylabel('Fraction modulated')
                    ax.legend(handles=[h1,h2,h3],loc=2)

                    #Modulation index
                    ax = axes[1]
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.mesogroup == mg) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                    sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                    sns.barplot(data=sub_df,x='block_j',y='abs_mod',hue='drug_type',ax=ax,order=blocks,hue_order=drug_order,palette=cmap_macro,errorbar='se')
                    # ax.set_ylim([0,0.6])
                    usrplt.adjust_spines(ax)
                    ax.set_xlabel('Time after 2nd injection (min)'); ax.set_ylabel('| Modulation index |')

                    for i, b in enumerate(blocks):
                        distr_i = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[0])]['abs_mod'].values
                        distr_j = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[1])]['abs_mod'].values
                        if (len(distr_i) < 5) | (len(distr_j) < 5):
                            continue
                        res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')
                                            
                        if res.pvalue < pval_thresh:
                            ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])

                    #Modulation index (dFR_spont_df.sig == 1) &
                    ax = axes[2]
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.mesogroup == mg) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                    sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                    sns.barplot(data=sub_df,x='block_j',y='mod',hue='drug_type',ax=ax,order=blocks,hue_order=drug_order,palette=cmap_macro,errorbar='se')
                    # ax.set_ylim([0,0.6])
                    usrplt.adjust_spines(ax)
                    ax.set_xlabel('Time after 2nd injection (min)'); ax.set_ylabel('Modulation index')

                    for i, b in enumerate(blocks):
                        distr_i = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[0])]['mod'].values
                        distr_j = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[1])]['mod'].values
                        if (len(distr_i) < 5) | (len(distr_j) < 5):
                            continue
                        res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')             
                        if res.pvalue < pval_thresh:
                            ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])
                
                    plt.savefig(join(PlotDir,f'mesogroup_barplot_{stat}_{block_i}_{mg}_{ct}_{bstrs_list[iB]}.png'),dpi=300,bbox_inches='tight',facecolor='w')

                    prs = prs_list[iS]
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = f'{mg}, {ct}, {b_str}'
                    usrplt.save_fig_to_pptx(fig, prs,slide)
                    plt.close(fig)
                    nPlots+=1

    ##------------------------------------------
    #Plot MI over time per mesogroup, celltype, behavior, & LAYER
    print('Plotting mesogroup barplots over time, per celltype, behavior, mesogroup, and layer')
    for iS, stat in enumerate(stat_list):
        for ct in ['RS','FS']:
            for iB, b_str in behavior_dict2.items():
                PlotDir = join(BaseDir,'mesogroup',bstrs_list[iB])
                if not os.path.exists(PlotDir):
                    os.makedirs(PlotDir)
                for mg in CTX_mesogroup_list:
                    for l in layer_list:
                        if l == '2/3':
                            l_str = '2-3'
                        else:
                            l_str = l

                        sub_df = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.mesogroup == mg) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                        if len(sub_df) < 5:
                            continue
                        fig, axes = plt.subplots(3,1,figsize=(10,15))
                        plt.suptitle(f'Change in {stat_title_list[iS]} due to saline or psilocybin injection\nBaseline = {block_i}, {mg}, layer {l}, {ct}, {b_str}',y=0.925)

                        #Fraction significantly modulated
                        ax = axes[0]
                        sub_df = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.mesogroup == mg) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                        iX = 0; xticks = []
                        for i, b in enumerate(blocks):
                            for jj, d in enumerate(drug_order):
                                sub_sub_df = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == d)]
                                n = len(sub_sub_df)
                                if n == 0:
                                    continue
                                nSig = np.sum(sub_sub_df['sig'].values)
                                fracSig = nSig/n
                                pos_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] > 0)])/n
                                neg_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] < 0)])/n
                                not_mod = len(sub_sub_df.loc[sub_sub_df.sig == 0])/n

                                h1 = ax.bar(iX, pos_mod, color=sns.xkcd_palette(['light red']), label='Positively Modulated')
                                h2 = ax.bar(iX, not_mod, bottom=pos_mod, color=sns.xkcd_palette(['silver']), label='Not Modulated')
                                h3 = ax.bar(iX, neg_mod, bottom=pos_mod+not_mod, color=sns.xkcd_palette(['medium blue']), label='Negatively Modulated')
                                n = len(np.unique(sub_sub_df['neuron_index']))
                                ax.text(iX,0.025,f'{n}',ha='center',fontsize=8,rotation=90)
                                if jj == 1:
                                    xticks.append(iX-0.5)
                                iX += 1
                            iX += 1

                        usrplt.adjust_spines(ax)
                        plt.autoscale(tight=True)
                        xticks = np.arange(1,3*len(blocks),3)-0.5 #[0.5, 3.5, 6.5, 9.5,12.5,]
                        ax.set_xticks(xticks)
                        ax.set_xticklabels(blocks,rotation=0)
                        ax.set_ylabel('Fraction modulated')
                        ax.legend(handles=[h1,h2,h3],loc=2)

                        #Modulation index
                        ax = axes[1]
                        sub_df = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.mesogroup == mg) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                        sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                        sns.barplot(data=sub_df,x='block_j',y='abs_mod',hue='drug_type',ax=ax,order=blocks,hue_order=drug_order,palette=cmap_macro,errorbar='se')

                        usrplt.adjust_spines(ax)
                        ax.set_xlabel('Time after 2nd injection (min)'); ax.set_ylabel('| Modulation index |')

                        for i, b in enumerate(blocks):
                            distr_i = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[0])]['abs_mod'].values
                            distr_j = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[1])]['abs_mod'].values
                            if (len(distr_i) < 5) | (len(distr_j) < 5):
                                continue
                            res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')
                                                
                            if res.pvalue < pval_thresh:
                                ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])

                        #Modulation index
                        ax = axes[2]
                        sub_df = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.mesogroup == mg) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                        sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                        sns.barplot(data=sub_df,x='block_j',y='mod',hue='drug_type',ax=ax,order=blocks,hue_order=drug_order,palette=cmap_macro,errorbar='se')
                        usrplt.adjust_spines(ax)
                        ax.set_xlabel('Time after 2nd injection (min)'); ax.set_ylabel('Modulation index')

                        for i, b in enumerate(blocks):
                            distr_i = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[0])]['mod'].values
                            distr_j = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[1])]['mod'].values
                            if (len(distr_i) < 5) | (len(distr_j) < 5):
                                continue
                            res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')
                                                
                            if res.pvalue < pval_thresh:
                                ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])
                    
                        plt.savefig(join(PlotDir,f'mesogroup_barplot_{stat}_{block_i}_{mg}_layer-{l_str}_{ct}_{bstrs_list[iB]}.png'),dpi=300,bbox_inches='tight',facecolor='w')
                        prs = prs_list[iS]
                        slide = prs.slides.add_slide(prs.slide_layouts[1])
                        slide.shapes.title.text = f'{mg}, layer {l}, {ct}, {b_str}'
                        usrplt.save_fig_to_pptx(fig, prs,slide)
                        plt.close(fig)
                        nPlots+=1

    PlotDir = join(BaseDir,'mesogroup')
    for iS, stat in enumerate(stat_list):
        prs = prs_list[iS]
        prs.save(join(PlotDir,f'mesogroup_barplots_{stat}.pptx'))

    ##------------------------------------------
    # Create a new powerPoint presentations to save figures to
    prs_list = []
    for iS, stat in enumerate(stat_list):
        prs = Presentation()

        #Add title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f'All Mice, neurons organized in "group" level'
        slide.placeholders[1].text = f'Change in {stat_title_list[iS]} due to saline or psilocybin injection'
        prs_list.append(prs)

    ##------------------------------------------
    #Plot MI over GROUPS per time window, celltype, behavior
    print('Plotting group barplots over regions, per celltype, behavior, and time window')
    for iS, stat in enumerate(stat_list):
        for ct in ['RS','FS']:
            for iB, b_str in behavior_dict2.items():
                PlotDir = join(BaseDir,'group',bstrs_list[iB])
                if not os.path.exists(PlotDir):
                    os.makedirs(PlotDir)
                for block_j in blocks:
                    fig, axes = plt.subplots(3,1,figsize=(10,15))
                    if block_j == 'post_1st_inj':
                        plt.suptitle(f'Change in {stat_title_list[iS]} due to saline or psilocybin injection\nCelltype = {ct}, behavior = {b_str}\n {block_i} vs {block_j}',y=0.94)
                    else:
                        plt.suptitle(f'Change in {stat_title_list[iS]} due to saline or psilocybin injection\nCelltype = {ct}, behavior = {b_str}\n {block_i} vs {block_j} minute after 2nd injection',y=0.94)

                    #Fraction significantly modulated
                    ax = axes[0]
                    sub_df = dFR_spont_df.loc[(dFR_spont_df['group'].isin(group_list)) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j == block_j)]
                    iX = 0; xticks = []
                    for i, g in enumerate(group_list):
                        for jj, d in enumerate(drug_order):
                            sub_sub_df = sub_df.loc[(sub_df['group'] == g) & (sub_df.drug_type == d)]
                            n = len(sub_sub_df)
                            if n == 0:
                                continue
                            nSig = np.sum(sub_sub_df['sig'].values)
                            fracSig = nSig/n
                            pos_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] > 0)])/n
                            neg_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] < 0)])/n
                            not_mod = len(sub_sub_df.loc[sub_sub_df.sig == 0])/n

                            h1 = ax.bar(iX, pos_mod, color=sns.xkcd_palette(['light red']), label='Positively Modulated')
                            h2 = ax.bar(iX, not_mod, bottom=pos_mod, color=sns.xkcd_palette(['silver']), label='Not Modulated')
                            h3 = ax.bar(iX, neg_mod, bottom=pos_mod+not_mod, color=sns.xkcd_palette(['medium blue']), label='Negatively Modulated')
                            n = len(np.unique(sub_sub_df['neuron_index']))
                            ax.text(iX,0.025,f'{n}',ha='center',fontsize=8,rotation=90)
                            if jj == 1:
                                xticks.append(iX-0.5)
                            iX += 1
                        iX += 1

                    usrplt.adjust_spines(ax)
                    plt.autoscale(tight=True)
                    xticks = np.arange(1,len(group_list)*3,3)-0.5
                    ax.set_xticks(xticks)
                    ax.set_xticklabels(group_list,rotation=0)
                    ax.set_ylabel('Fraction modulated')
                    ax.legend(handles=[h1,h2,h3],loc=2)

                    #Modulation index
                    ax = axes[1]
                    sub_df = dFR_spont_df.loc[(dFR_spont_df['group'].isin(group_list)) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j == block_j)]
                    sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                    sns.barplot(data=sub_df,x='group',y='abs_mod',hue='drug_type',ax=ax,order=group_list,hue_order=drug_order,palette=cmap_macro,errorbar='se')
                    # ax.set_ylim([0,0.6])
                    usrplt.adjust_spines(ax)
                    ax.set_xlabel('Mesogroup'); ax.set_ylabel('| Modulation index |')

                    for i, mg in enumerate(group_list):
                        distr_i = sub_df.loc[(sub_df.group == mg) & (sub_df.drug_type == drug_order[0])]['abs_mod'].values
                        distr_j = sub_df.loc[(sub_df.group == mg) & (sub_df.drug_type == drug_order[1])]['abs_mod'].values
                        if (len(distr_i) < 2) | (len(distr_j) < 2):
                            continue
                        res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')                  
                        if res.pvalue < pval_thresh:
                            ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])

                    #Modulation index
                    ax = axes[2]
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.group.isin(group_list)) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j == block_j)]
                    sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                    sns.barplot(data=sub_df,x='group',y='mod',hue='drug_type',ax=ax,order=group_list,hue_order=drug_order,palette=cmap_macro,errorbar='se')
                    usrplt.adjust_spines(ax)
                    ax.set_xlabel('Mesogroup'); ax.set_ylabel('Modulation index')

                    for i, mg in enumerate(group_list):
                        distr_i = sub_df.loc[(sub_df.group == mg) & (sub_df.drug_type == drug_order[0])]['mod'].values
                        distr_j = sub_df.loc[(sub_df.group == mg) & (sub_df.drug_type == drug_order[1])]['mod'].values
                        if (len(distr_i) < 2) | (len(distr_j) < 2):
                            continue
                        res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')
                                            
                        if res.pvalue < pval_thresh:
                            ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])
                    plt.savefig(join(PlotDir,f'group_barplot_{stat}_{ct}_{bstrs_list[iB]}_{epoch_i}_vs_{block_j}.png'),dpi=300,bbox_inches='tight',facecolor='w')
                    prs = prs_list[iS]
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = f'{epoch_i}_vs_{block_j}, {ct}, {b_str}'
                    usrplt.save_fig_to_pptx(fig, prs,slide)
                    plt.close(fig)
                    nPlots+=1

    ##------------------------------------------
    #Plot MI over time per group, celltype, behavior
    print('Plotting group barplots over time, per celltype, behavior, and group')
    for iS, stat in enumerate(stat_list):
        for ct in ['RS','FS']:
            for iB, b_str in behavior_dict2.items():
                PlotDir = join(BaseDir,'group',bstrs_list[iB])
                if not os.path.exists(PlotDir):
                    os.makedirs(PlotDir)

                for g in group_list:
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.group == g) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                    if len(sub_df) < 5:
                        continue
                    fig, axes = plt.subplots(3,1,figsize=(10,15))
                    plt.suptitle(f'Change in {stat_title_list[iS]} due to saline or psilocybin injection\nBaseline = {block_i}, {g}, {ct}, {b_str}',y=0.925)

                    #Fraction significantly modulated
                    ax = axes[0]
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.group == g) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j.isin(blocks))]
                    iX = 0; xticks = []
                    for i, b in enumerate(blocks):
                        for jj, d in enumerate(drug_order):
                            sub_sub_df = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == d)]
                            n = len(sub_sub_df)
                            if n == 0:
                                continue
                            nSig = np.sum(sub_sub_df['sig'].values)
                            fracSig = nSig/n
                            pos_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] > 0)])/n
                            neg_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] < 0)])/n
                            not_mod = len(sub_sub_df.loc[sub_sub_df.sig == 0])/n

                            h1 = ax.bar(iX, pos_mod, color=sns.xkcd_palette(['light red']), label='Positively Modulated')
                            h2 = ax.bar(iX, not_mod, bottom=pos_mod, color=sns.xkcd_palette(['silver']), label='Not Modulated')
                            h3 = ax.bar(iX, neg_mod, bottom=pos_mod+not_mod, color=sns.xkcd_palette(['medium blue']), label='Negatively Modulated')
                            n = len(np.unique(sub_sub_df['neuron_index']))
                            ax.text(iX,0.025,f'{n}',ha='center',fontsize=8,rotation=90)
                            if jj == 1:
                                xticks.append(iX-0.5)
                            iX += 1
                        iX += 1

                    usrplt.adjust_spines(ax)
                    plt.autoscale(tight=True)
                    # ax.hlines(0.5,*ax.get_xlim(),linestyles='--',color='k',alpha=0.5)
                    xticks = np.arange(1,len(blocks)*3,3)-0.5
                    ax.set_xticks(xticks)
                    ax.set_xticklabels(blocks,rotation=0)
                    ax.set_ylabel('Fraction modulated')
                    ax.legend(handles=[h1,h2,h3],loc=2)

                    #Modulation index
                    ax = axes[1]
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.group == g) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                    sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                    sns.barplot(data=sub_df,x='block_j',y='abs_mod',hue='drug_type',ax=ax,order=blocks,hue_order=drug_order,palette=cmap_macro,errorbar='se')
                    usrplt.adjust_spines(ax)
                    ax.set_xlabel('Time after 2nd injection (min)'); ax.set_ylabel('| Modulation index |')

                    for i, b in enumerate(blocks):

                        distr_i = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[0])]['abs_mod'].values
                        distr_j = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[1])]['abs_mod'].values
                        if (len(distr_i) < 5) | (len(distr_j) < 5):
                            continue
                        res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')
                                            
                        if res.pvalue < pval_thresh:
                            ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])

                    #Modulation index
                    ax = axes[2]
                    sub_df = dFR_spont_df.loc[(dFR_spont_df.group == g) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                    sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                    sns.barplot(data=sub_df,x='block_j',y='mod',hue='drug_type',ax=ax,order=blocks,hue_order=drug_order,palette=cmap_macro,errorbar='se')
                    usrplt.adjust_spines(ax)
                    ax.set_xlabel('Time after 2nd injection (min)'); ax.set_ylabel('Modulation index')
                    for i, b in enumerate(blocks):

                        distr_i = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[0])]['mod'].values
                        distr_j = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[1])]['mod'].values
                        if (len(distr_i) < 5) | (len(distr_j) < 5):
                            continue
                        res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')
                                            
                        if res.pvalue < pval_thresh:
                            ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])
                    
                    plt.savefig(join(PlotDir,f'group_barplot_{stat}_{block_i}_{g}_{ct}_{bstrs_list[iB]}.png'),dpi=300,bbox_inches='tight',facecolor='w')
                    prs = prs_list[iS]
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = f'{g}, {ct}, {b_str}'
                    usrplt.save_fig_to_pptx(fig, prs,slide)
                    plt.close(fig)
                    nPlots+=1

    ##------------------------------------------
    #Plot MI over time per group, celltype, behavior, & LAYER
    print('Plotting group barplots over time, per celltype, behavior, group, and layer')
    for iS, stat in enumerate(stat_list):
        for ct in ['RS','FS']:
            for iB, b_str in behavior_dict2.items():
                PlotDir = join(BaseDir,'group',bstrs_list[iB])
                if not os.path.exists(PlotDir):
                    os.makedirs(PlotDir)

                for g in CTX_group_list:
                    for l in layer_list:
                        if l == '2/3':
                            l_str = '2-3'
                        else:
                            l_str = l
                        sub_df0 = dFR_spont_df.loc[(dFR_spont_df.drug_type == drug_order[0]) & (dFR_spont_df.group == g) & (dFR_spont_df.layer == l) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                        sub_df1 = dFR_spont_df.loc[(dFR_spont_df.drug_type == drug_order[1]) & (dFR_spont_df.group == g) & (dFR_spont_df.layer == l) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                        if (len(sub_df0) < 5) | (len(sub_df1) < 5):
                            continue
                        fig, axes = plt.subplots(3,1,figsize=(10,15))
                        plt.suptitle(f'Change in {stat_title_list[iS]} due to saline or psilocybin injection\nBaseline = {block_i}, {g}, layer {l}, {ct}, {b_str}',y=0.925)

                        #Fraction significantly modulated
                        ax = axes[0]
                        sub_df = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.group == g) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                        iX = 0; xticks = []
                        for i, b in enumerate(blocks):
                            for jj, d in enumerate(drug_order):
                                sub_sub_df = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == d)]
                                n = len(sub_sub_df)
                                if n == 0:
                                    continue
                                nSig = np.sum(sub_sub_df['sig'].values)
                                fracSig = nSig/n
                                pos_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] > 0)])/n
                                neg_mod = len(sub_sub_df.loc[(sub_sub_df.sig == 1) & (sub_sub_df['mod'] < 0)])/n
                                not_mod = len(sub_sub_df.loc[sub_sub_df.sig == 0])/n

                                h1 = ax.bar(iX, pos_mod, color=sns.xkcd_palette(['light red']), label='Positively Modulated')
                                h2 = ax.bar(iX, not_mod, bottom=pos_mod, color=sns.xkcd_palette(['silver']), label='Not Modulated')
                                h3 = ax.bar(iX, neg_mod, bottom=pos_mod+not_mod, color=sns.xkcd_palette(['medium blue']), label='Negatively Modulated')
                                n = len(np.unique(sub_sub_df['neuron_index']))
                                ax.text(iX,0.025,f'{n}',ha='center',fontsize=8,rotation=90)
                                if jj == 1:
                                    xticks.append(iX-0.5)
                                iX += 1
                            iX += 1

                        usrplt.adjust_spines(ax)
                        plt.autoscale(tight=True)
                        xticks = np.arange(1,len(blocks)*3,3)-0.5
                        ax.set_xticks(xticks)
                        ax.set_xticklabels(blocks,rotation=0)
                        ax.set_ylabel('Fraction modulated')
                        ax.legend(handles=[h1,h2,h3],loc=2)

                        #Modulation index
                        ax = axes[1]
                        sub_df = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.group == g) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                        sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                        sns.barplot(data=sub_df,x='block_j',y='abs_mod',hue='drug_type',ax=ax,order=blocks,hue_order=drug_order,palette=cmap_macro,errorbar='se')
                        usrplt.adjust_spines(ax)
                        ax.set_xlabel('Time after 2nd injection (min)'); ax.set_ylabel('| Modulation index |')

                        for i, b in enumerate(blocks):

                            distr_i = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[0])]['abs_mod'].values
                            distr_j = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[1])]['abs_mod'].values
                            if (len(distr_i) < 5) | (len(distr_j) < 5):
                                continue
                            res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')              
                            if res.pvalue < pval_thresh:
                                ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])

                        #Modulation index
                        ax = axes[2]
                        sub_df = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.group == g) & (dFR_spont_df.behavior == b_str) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == epoch_i) & (dFR_spont_df.block_j.isin(blocks))]
                        sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                        sns.barplot(data=sub_df,x='block_j',y='mod',hue='drug_type',ax=ax,order=blocks,hue_order=drug_order,palette=cmap_macro,errorbar='se')
                        # ax.set_ylim([0,0.6])
                        usrplt.adjust_spines(ax)
                        ax.set_xlabel('Time after 2nd injection (min)'); ax.set_ylabel('Modulation index')


                        for i, b in enumerate(blocks):

                            distr_i = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[0])]['mod'].values
                            distr_j = sub_df.loc[(sub_df.block_j == b) & (sub_df.drug_type == drug_order[1])]['mod'].values
                            if (len(distr_i) < 5) | (len(distr_j) < 5):
                                continue
                            res = st.mannwhitneyu(distr_i,distr_j,nan_policy='omit')          
                            if res.pvalue < pval_thresh:
                                ax.text(i-0.1,0,'*',fontsize=20,fontweight='bold',color=usrplt.cc[8])
                        
                        plt.savefig(join(PlotDir,f'group_barplot_{stat}_{block_i}_{g}_layer-{l_str}_{ct}_{bstrs_list[iB]}.png'),dpi=300,bbox_inches='tight',facecolor='w')
                        prs = prs_list[iS]
                        slide = prs.slides.add_slide(prs.slide_layouts[1])
                        slide.shapes.title.text = f'{g}, layer {l}, {ct}, {b_str}'
                        usrplt.save_fig_to_pptx(fig, prs,slide)
                        plt.close(fig)
                        nPlots+=1

    PlotDir = join(BaseDir,'group')
    for iS, stat in enumerate(stat_list):
        prs = prs_list[iS]
        prs.save(join(PlotDir,f'group_barplots_{stat}.pptx'))

    ##------------------------------------------
    # Create a new powerPoint presentations to save figures to
    prs_list = []
    for iS, stat in enumerate(stat_list):
        prs = Presentation()

        #Add title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f'All Mice, neurons organized in "group" level'
        slide.placeholders[1].text = f'Change in {stat_title_list[iS]} due to saline or psilocybin injection'
        prs_list.append(prs)

    # ##------------------------------------------
    # #Plot scatter plot group, time window, celltype
    # PlotDir = join(BaseDir,'group','scatter')
    # if not os.path.exists(PlotDir):
    #     os.makedirs(PlotDir)
    # print('Plotting group scatter plots per region, celltype, and time window')
    # for iS, stat in enumerate(stat_list):
    #     for ct in ['RS','FS']:
    #         for g in TH_group_list:
    #             for block_j in blocks:

    #                 sub_df1 = dFR_spont_df.loc[(dFR_spont_df.group == g) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.drug_type == drug_order[0]) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j == block_j) & (dFR_spont_df.behavior == 'all (>0cm/s)')]
    #                 sub_df2 = dFR_spont_df.loc[(dFR_spont_df.group == g) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.drug_type == drug_order[1]) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j == block_j) & (dFR_spont_df.behavior == 'all (>0cm/s)')]
    #                 if (len(sub_df1) < 2) | (len(sub_df2) < 2):
    #                     continue
    #                 fig, axes = plt.subplots(3,4,figsize=(20,15),gridspec_kw={'hspace':0.35,'wspace':0.35})
    #                 if block_j == 'post_1st_inj':
    #                     block_j_str = 'Post 1st saline inj'
    #                 else:
    #                     block_j_str = f'{block_j} mins past 2nd injection'

    #                 plt.suptitle(f'Change in {stat_title_list[iS]} due to saline or psilocybin injection\nGroup: {g}, cellltype = {ct}, {block_i} vs {block_j_str}',y=0.95)

    #                 for iB, b_str in behavior_dict2.items():
    #                     sub_df1 = dFR_spont_df.loc[(dFR_spont_df.group == g) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.drug_type == drug_order[0]) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j == block_j) & (dFR_spont_df.behavior == b_str)]
    #                     sub_df2 = dFR_spont_df.loc[(dFR_spont_df.group == g) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.drug_type == drug_order[1]) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j == block_j) & (dFR_spont_df.behavior == b_str)]
    #                     sub_df1 = sub_df1.loc[sub_df1['sig'] == 1] if plot_only_sig_mod else sub_df1
    #                     sub_df2 = sub_df2.loc[sub_df2['sig'] == 1] if plot_only_sig_mod else sub_df2
    #                     if (len(sub_df1) < 2) | (len(sub_df2) < 2):
    #                         continue
                            
    #                     delta_list = []; xymin_list = [];  xymax_list = []
    #                     for ii, d in enumerate(drug_order):

    #                         sub_df = dFR_spont_df.loc[(dFR_spont_df.group == g) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.drug_type == d) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j == block_j) & (dFR_spont_df.behavior == b_str)]
    #                         sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
    #                         FR_i = sub_df['val_i'].values
    #                         FR_j = sub_df['val_j'].values
    #                         delta = sub_df['delta'].values
    #                         cts = sub_df['celltype'].values
    #                         mod = sub_df['mod'].values
    #                         delta_list.append((FR_i,FR_j,delta,mod))

    #                         nN = len(np.unique(sub_df['neuron_index'].values)); nE = len(np.unique(sub_df['rec_name'].values))
    #                         ax = axes[iB,ii]
    #                         ax.set_title(f'{stat} during {b_str}: {d} \n(n = {nN} neurons from {nE} exps)')
    #                         ax.scatter(FR_i,FR_j,s=10,marker='o',edgecolor=cmap_macro[ii],facecolor='none')

    #                         ymin = np.nanmin(FR_j); xmin = np.nanmin(FR_i)
    #                         ymax = np.nanpercentile(FR_j,98); xmax = np.nanpercentile(FR_i,98)
    #                         xymin = np.nanmin([ymin,xmin]);xymin_list.append(xymin)
    #                         xymax = np.nanmax([ymax,xmax]);xymax_list.append(xymax)

    #                     xymin = np.nanmin(xymin_list); xymax = np.nanmax(xymax_list)   
    #                     for ii, d in enumerate(drug_order):
    #                         ax = axes[iB,ii]
    #                         FR_i, FR_j, delta, mod = delta_list[ii]
                            
    #                         ax.plot([xymin,xymax],[xymin,xymax],'-k')
    #                         ax.set_xlim([xymin,xymax]); ax.set_ylim([xymin,xymax])
    #                         usrplt.adjust_spines(ax)
    #                         ax.set_xlabel(f'{block_i}')
    #                         ax.set_ylabel(f'{block_j}')
    #                         res = st.wilcoxon(FR_i,FR_j,nan_policy='omit')
    #                         if res.pvalue < pval_thresh:
    #                             ax.text(0.5*xymax,0.9*xymax,f'*',color=cmap_macro[ii],fontweight='bold',fontsize=20)

    #                     #Compare saline vs psilocybin
    #                     distr_i = delta_list[0][-1]
    #                     distr_j = delta_list[1][-1]

    #                     FR_i, FR_j, delta, mod_i = delta_list[0]
    #                     FR_i, FR_j, delta, mod_j = delta_list[1]
    #                     ax = axes[iB,2]
    #                     sns.histplot([mod_i,mod_j],palette=cmap_macro,binwidth=0.1,multiple='layer',stat='density',common_norm=False,kde=False,ax=ax)

    #                     ax.set_title('Modulation index')
    #                     ax.set_xlabel('MI')
    #                     ax.set_xlim([-1,1])
    #                     usrplt.adjust_spines(ax)
    #                     ylim = ax.get_ylim(); xlim = ax.get_xlim()
    #                     ax.vlines(np.nanmean(mod_i),*ylim,ls='-',color=cmap_macro[0],lw=2,label=f'Saline MI: {np.nanmean(mod_i):.2f}')
    #                     ax.vlines(np.nanmean(mod_j),*ylim,ls='-',color=cmap_macro[1],lw=2,label=f'Psilocybin MI: {np.nanmean(mod_j):.2f}')
    #                     ax.legend()

    #                     if (len(mod_i) > 0) & (len(mod_j) > 0):
    #                         res = st.mannwhitneyu(mod_i,mod_j,nan_policy='omit')
    #                         if res.pvalue < pval_thresh/3:
    #                             ax.text(.5,0.96*ax.get_ylim()[1],'*',fontsize=20,fontweight='bold',color='k')

    #                     ax = axes[iB,3]
    #                     sns.histplot([np.abs(mod_i),np.abs(mod_j)],palette=cmap_macro,binwidth=0.1,multiple='layer',stat='density',common_norm=False,kde=False,ax=ax)
    #                     ax.set_title('| Modulation index |')
    #                     ax.set_xlabel('|MI|')
    #                     ax.set_xlim([0,1])
    #                     usrplt.adjust_spines(ax)
    #                     ylim = ax.get_ylim(); xlim = ax.get_xlim()
    #                     ax.vlines(np.nanmean(np.abs(mod_i)),*ylim,ls='--',color=cmap_macro[0],lw=2,label=f'Saline |MI|: {np.nanmean(np.abs(mod_i)):.2f}')
    #                     ax.vlines(np.nanmean(np.abs(mod_j)),*ylim,ls='--',color=cmap_macro[1],lw=2,label=f'Psilocybin |MI|: {np.nanmean(np.abs(mod_j)):.2f}')
    #                     ax.legend()

    #                     if (len(mod_i) > 0) & (len(mod_j) > 0):
    #                         res = st.mannwhitneyu(np.abs(mod_i),np.abs(mod_j),nan_policy='omit')
    #                         if res.pvalue < pval_thresh:
    #                             ax.text(.5,0.96*ax.get_ylim()[1],'*',fontsize=20,fontweight='bold',color='k')

    #                 # pdb.set_trace()
    #                 plt.savefig(join(PlotDir,f'group_scatter_{stat}_{g}_{ct}_{block_i}_vs_{block_j}.png'),dpi=300,bbox_inches='tight',facecolor='w')
    #                 prs = prs_list[iS]
    #                 slide = prs.slides.add_slide(prs.slide_layouts[1])
    #                 slide.shapes.title.text = f'{g}, {ct}, {block_i}_vs_{block_j}'
    #                 usrplt.save_fig_to_pptx(fig, prs,slide)
    #                 plt.close(fig)
    #                 nPlots+=1

    ##------------------------------------------
    #Plot scatter plot group, time window, celltype, & LAYER
    print('Plotting group scatter plots per region, celltype, layer, and time window')
    for iS, stat in enumerate(stat_list):
        for ct in ['RS','FS']:
            for g in CTX_group_list:
                layers = np.unique(cells_df.loc[cells_df.group == g]['layer'].values)
                for l in layers:
                    if l == '2/3':
                        l_str = '2-3'
                    else:
                        l_str = l
                    # l = '5'; l_str = '5'
                    for block_j in blocks:

                        sub_df1 = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.group == g) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.drug_type == drug_order[0]) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j == block_j) & (dFR_spont_df.behavior == 'all (>0cm/s)')]
                        sub_df2 = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.group == g) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.drug_type == drug_order[1]) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j == block_j) & (dFR_spont_df.behavior == 'all (>0cm/s)')]
                        sub_df1 = sub_df1.loc[sub_df1['sig'] == 1] if plot_only_sig_mod else sub_df1
                        sub_df2 = sub_df2.loc[sub_df2['sig'] == 1] if plot_only_sig_mod else sub_df2
                        if (len(sub_df1) < 2) | (len(sub_df2) < 2):
                            continue
                        fig, axes = plt.subplots(3,4,figsize=(20,15),gridspec_kw={'hspace':0.35,'wspace':0.35})
                        if block_j == 'post_1st_inj':
                            block_j_str = 'Post 1st saline inj'
                        else:
                            block_j_str = f'{block_j} mins past 2nd injection'

                        plt.suptitle(f'Change in {stat_title_list[iS]} due to saline or psilocybin injection\nGroup: {g}, layer {l}, cellltype = {ct}, {block_i} vs {block_j_str}',y=0.95)

                        for iB, b_str in behavior_dict2.items():
                            sub_df1 = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.group == g) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.drug_type == drug_order[0]) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j == block_j) & (dFR_spont_df.behavior == b_str)]
                            sub_df2 = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.group == g) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.drug_type == drug_order[1]) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j == block_j) & (dFR_spont_df.behavior == b_str)]
                            sub_df1 = sub_df1.loc[sub_df1['sig'] == 1] if plot_only_sig_mod else sub_df1
                            sub_df2 = sub_df2.loc[sub_df2['sig'] == 1] if plot_only_sig_mod else sub_df2
                            if (len(sub_df1) < 2) | (len(sub_df2) < 2):
                                continue
                                
                            delta_list = []; xymin_list = [];  xymax_list = []
                            for ii, d in enumerate(drug_order):

                                sub_df = dFR_spont_df.loc[(dFR_spont_df.layer == l) & (dFR_spont_df.group == g) & (dFR_spont_df.celltype == ct) & (dFR_spont_df.drug_type == d) & (dFR_spont_df.stat == stat) & (dFR_spont_df.block_i == block_i) & (dFR_spont_df.block_j == block_j) & (dFR_spont_df.behavior == b_str)]
                                sub_df = sub_df.loc[sub_df['sig'] == 1] if plot_only_sig_mod else sub_df
                                FR_i = sub_df['val_i'].values
                                FR_j = sub_df['val_j'].values
                                delta = sub_df['delta'].values
                                cts = sub_df['celltype'].values
                                mod = sub_df['mod'].values
                                delta_list.append((FR_i,FR_j,delta,mod))

                                nN = len(np.unique(sub_df['neuron_index'].values)); nE = len(np.unique(sub_df['rec_name'].values))
                                ax = axes[iB,ii]
                                ax.set_title(f'{stat} during {b_str}: {d} \n(n = {nN} neurons from {nE} exps)')
                                ax.scatter(FR_i,FR_j,s=10,marker='o',edgecolor=cmap_macro[ii],facecolor='none')

                                ymin = np.nanmin(FR_j); xmin = np.nanmin(FR_i)
                                ymax = np.nanpercentile(FR_j,98); xmax = np.nanpercentile(FR_i,98)
                                xymin = np.nanmin([ymin,xmin]);xymin_list.append(xymin)
                                xymax = np.nanmax([ymax,xmax]);xymax_list.append(xymax)

                            xymin = np.nanmin(xymin_list); xymax = np.nanmax(xymax_list)   
                            for ii, d in enumerate(drug_order):
                                ax = axes[iB,ii]
                                FR_i, FR_j, delta, mod = delta_list[ii]
                                
                                ax.plot([xymin,xymax],[xymin,xymax],'-k')
                                ax.set_xlim([xymin,xymax]); ax.set_ylim([xymin,xymax])
                                usrplt.adjust_spines(ax)
                                ax.set_xlabel(f'{block_i}')
                                ax.set_ylabel(f'{block_j}')
                                res = st.wilcoxon(FR_i,FR_j,nan_policy='omit')
                                if res.pvalue < pval_thresh:
                                    ax.text(0.5*xymax,0.9*xymax,f'*',color=cmap_macro[ii],fontweight='bold',fontsize=20)

                            #Compare saline vs psilocybin
                            distr_i = delta_list[0][-1]
                            distr_j = delta_list[1][-1]

                            FR_i, FR_j, delta, mod_i = delta_list[0]
                            FR_i, FR_j, delta, mod_j = delta_list[1]
                            ax = axes[iB,2]
                            sns.histplot([mod_i,mod_j],palette=cmap_macro,binwidth=0.1,multiple='layer',stat='density',common_norm=False,kde=False,ax=ax)

                            ax.set_title('Modulation index')
                            ax.set_xlabel('MI')
                            ax.set_xlim([-1,1])
                            usrplt.adjust_spines(ax)
                            ylim = ax.get_ylim(); xlim = ax.get_xlim()
                            ax.vlines(np.nanmean(mod_i),*ylim,ls='-',color=cmap_macro[0],lw=2,label=f'Saline MI: {np.nanmean(mod_i):.2f}')
                            ax.vlines(np.nanmean(mod_j),*ylim,ls='-',color=cmap_macro[1],lw=2,label=f'Psilocybin MI: {np.nanmean(mod_j):.2f}')
                            ax.legend()

                            if (len(mod_i) > 0) & (len(mod_j) > 0):
                                res = st.mannwhitneyu(mod_i,mod_j,nan_policy='omit')
                                if res.pvalue < pval_thresh/3:
                                    ax.text(.5,0.96*ax.get_ylim()[1],'*',fontsize=20,fontweight='bold',color='k')

                            ax = axes[iB,3]
                            sns.histplot([np.abs(mod_i),np.abs(mod_j)],palette=cmap_macro,binwidth=0.1,multiple='layer',stat='density',common_norm=False,kde=False,ax=ax)
                            ax.set_title('| Modulation index |')
                            ax.set_xlabel('|MI|')
                            ax.set_xlim([0,1])
                            usrplt.adjust_spines(ax)
                            ylim = ax.get_ylim(); xlim = ax.get_xlim()
                            ax.vlines(np.nanmean(np.abs(mod_i)),*ylim,ls='--',color=cmap_macro[0],lw=2,label=f'Saline |MI|: {np.nanmean(np.abs(mod_i)):.2f}')
                            ax.vlines(np.nanmean(np.abs(mod_j)),*ylim,ls='--',color=cmap_macro[1],lw=2,label=f'Psilocybin |MI|: {np.nanmean(np.abs(mod_j)):.2f}')
                            ax.legend()

                            if (len(mod_i) > 0) & (len(mod_j) > 0):
                                res = st.mannwhitneyu(np.abs(mod_i),np.abs(mod_j),nan_policy='omit')
                                if res.pvalue < pval_thresh:
                                    ax.text(.5,0.96*ax.get_ylim()[1],'*',fontsize=20,fontweight='bold',color='k')

                        plt.savefig(join(PlotDir,f'group_scatter_{stat}_{g}_layer-{l_str}_{ct}_{block_i}_vs_{block_j}.png'),dpi=300,bbox_inches='tight',facecolor='w')
                        prs = prs_list[iS]
                        slide = prs.slides.add_slide(prs.slide_layouts[1])
                        slide.shapes.title.text = f'{g}, layer {l}, {ct}, {block_i}_vs_{block_j}'
                        usrplt.save_fig_to_pptx(fig, prs,slide)
                        plt.close(fig)
                        nPlots+=1

    PlotDir = join(BaseDir,'group')
    for iS, stat in enumerate(stat_list):
        prs = prs_list[iS]
        prs.save(join(PlotDir,f'group_scatter_{stat}.pptx'))

    print(f'DONE: Generated {nPlots} plots, Good luck with the analysis!')

                    
                    
