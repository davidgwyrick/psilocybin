base_dir_server = '/allen/programs/mindscope/workgroups/templeton-psychedelics/'
base_dir_local = '/data/tiny-blue-dot/zap-n-zip/EEG_exp'; base_dir = base_dir_server
ProjDir = '/home/david.wyrick/projects/zap-n-zip/'
ServDir = '/data/projects/zap-n-zip/'

#Basea
import gspread
from os.path import join
from glob import glob
import json, os, time, sys, argparse
import numpy as np
import pandas as pd
import scipy.stats as st
import pingouin as pg
from scipy.ndimage import gaussian_filter
from scipy.interpolate import interp1d

#Plot
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
from matplotlib.gridspec import GridSpec

#Project
from tbd_eeg.data_analysis.eegutils import EEGexp
import tbd_eeg.data_analysis.Utilities.utilities as tbd_util
from tbd_eeg.data_analysis.Utilities.behavior_movies import Movie


#Allen
from allensdk.core.mouse_connectivity_cache import MouseConnectivityCache
mcc = MouseConnectivityCache(resolution=10)
str_tree = mcc.get_structure_tree()

#Read in allen CCF
ccfsum = pd.read_csv('/home/david.wyrick/projects/zap-n-zip/data/mouse_connectivity/ontology_v2.csv')

#User
sys.path.append(ProjDir)
import util
import plotting as usrplt

#Single cell stuff
import elephant.statistics as est
from quantities import ms, s, Hz
from neo.core import SpikeTrain

#Define behavioral states
behavior_dict = {0: 'rest (0-1cm/s)', 1: 'walk (1-5cm/s)',2: 'slow run (5-20cm/s)',3:'run (20-50cm/s)', 4: 'fast run (>50cm/s)'}
behavior_strs = list(behavior_dict.values())
nBehaviors = len(behavior_strs)

#Define windows to calculate firing rate
windows_of_interest = [[.002,.025],[.075,.3],[.3,1],[1,4.5]]
window_strs = ['evoked (2-25ms)','rebound (75-300ms)','post-rebound (0.3-1s)','ISI (1-4.5s)']
nWindows = len(window_strs)

gc = gspread.service_account() # need a key file to access the account (step 2) 
sh = gc.open('Templeton-log_exp') # open the spreadsheet 

df = pd.DataFrame(sh.sheet1.get()) 
exp_table = df.T.set_index(0).T # put it in a nicely formatted dataframeexp_table.iloc[10:]
exp_table = exp_table.iloc[9:-8].set_index('mouse_name')

##===== ============================ =====##
##===== Parse Command Line Arguments =====##
parser = argparse.ArgumentParser(description='single-cell-metrics')

##===== Data Options =====##
parser.add_argument('--mouseID',type=str, default='mouse669117',
                    help='mouse to perform analysis on')

parser.add_argument('--rec_name',type=str, default='pilot_aw_psi_2023-03-30_11-37-07',
                    help='experiment to perform analysis on')

parser.add_argument('--time_bin_ms',type=int, default=100,
                    help='time_bin_ms')

parser.add_argument('--fr_thresh',type=float, default=1,
                    help='Firing rate threshold for neurons to include in analysis')

parser.add_argument('--calculate_evoked',type=int, default=0,
                    help='Firing rate threshold for neurons to include in analysis')



def plot_boxplot_comparison(dFR_spont_df, i, j, g = None, pval_thresh = 0.05):

    epoch_i = spont_filename_list[i]
    epoch_j = spont_filename_list[j]
    if g is None:
        x = 'group'
        sub_df = dFR_spont_df.loc[(dFR_spont_df.epoch_i == epoch_i) & (dFR_spont_df.epoch_j == epoch_j)]
    else:
        x = 'area'
        sub_df = dFR_spont_df.loc[(dFR_spont_df.epoch_i == epoch_i) & (dFR_spont_df.epoch_j == epoch_j) & (dFR_spont_df.group == g)]

    #Determine unique boxes
    uniq_boxes = np.unique(sub_df[x])
    nBoxes = len(uniq_boxes)

    #Plot modulation index
    fig, axes = plt.subplots(2,1,figsize=(10,8))
    ax = axes[0]
    ax.set_title(f'MI: "{epoch_i}" vs "{epoch_j}"')
    gs = sns.boxplot(x=x,y='MI',hue='behavior',palette=sns.color_palette('Reds',5),order=uniq_boxes,hue_order=behavior_strs,data=sub_df,legend=False,ax=ax)
    ax.set_xlabel(x)
    xlim = ax.get_xlim(); ylim = ax.get_ylim()
    ax.hlines(0,-0.5,nBoxes-0.5,color='k',zorder=0)
    ax.set_xlim(xlim)
    ax.set_xlabel('')
    uniq_behaviors = np.unique(sub_df['behavior'])

    for ii, b in enumerate(uniq_boxes):
        for jj, w in enumerate(uniq_behaviors):
            if x == 'area':
                sub_sub_df = sub_df.loc[(sub_df.area == b) & (sub_df.behavior == w)]
            else:
                sub_sub_df = sub_df.loc[(sub_df.group == b) & (sub_df.behavior == w)]

            res = pg.wilcoxon(sub_sub_df['MI'].values)
            pval = res['p-val'][0]
            if pval < pval_thresh/len(uniq_behaviors):
                ax.text(ii-0.335+jj*0.4,0.925,'*',fontsize=20,fontweight='bold',color='k')

    #Plot delta firing rate
    # fig, ax = plt.subplots(figsize=(10,4))
    ax = axes[1]
    ax.set_title(f'\u0394-FR: "{epoch_i}" vs "{epoch_j}"')
    gs = sns.boxplot(x=x,y='dFR',hue='behavior',palette=sns.color_palette('Reds',5),order=uniq_boxes,hue_order=behavior_strs,data=sub_df,legend=False,ax=ax)
    ax.set_xlabel(x)

    ymin = np.round(np.nanpercentile(sub_df['dFR'],2.5))
    ymax = np.round(np.nanpercentile(sub_df['dFR'],97.5))
    ax.set_ylim([ymin,ymax])
    xlim = ax.get_xlim(); ylim = ax.get_ylim()
    ax.hlines(0,-0.5,nBoxes-0.5,color='k',zorder=0)
    ax.set_xlim(xlim)
    ax.set_ylabel(f'\u0394-FR (Hz)')
    uniq_behaviors = np.unique(sub_df['behavior'])
    for ii, b in enumerate(uniq_boxes):
        if x == 'area':
            n = np.sum(areas_sub == b)
        else:
            n = np.sum(groups_sub == b)
        ax.text(ii-.15,ylim[0],n)
        for jj, w in enumerate(uniq_behaviors):
            if x == 'area':
                sub_sub_df = sub_df.loc[(sub_df.area == b) & (sub_df.behavior == w)]
            else:
                sub_sub_df = sub_df.loc[(sub_df.group == b) & (sub_df.behavior == w)]
            if n > 1:
                res = pg.wilcoxon(sub_sub_df['MI'].values)
                pval = res['p-val'][0]
                if pval < pval_thresh/len(uniq_behaviors):
                    ax.text(ii-0.335+jj*0.4,ylim[1]-6,'*',fontsize=20,fontweight='bold',color='k')
    
    return fig


def plot_boxplot_comparison2(dFR_spont_df, i, g, a = None, pval_thresh = 0.05):

    fig, axes = plt.subplots(2,1,figsize=(10,8))

    N = np.sum(groups_sub == g)
    epoch_i = spont_filename_list[i]

    
    if a is None:
        sub_df = dFR_spont_df.loc[(dFR_spont_df.epoch_i == epoch_i) & (dFR_spont_df.group == g)]
        n = np.sum(groups_sub == g)
        plt.suptitle(f'Group: {g}, n = {n}')
    else:
        sub_df = dFR_spont_df.loc[(dFR_spont_df.epoch_i == epoch_i) & (dFR_spont_df.group == g) & (dFR_spont_df.area == a)]
        area_full = str_tree.get_structures_by_acronym([a])[0]['name']
        n = np.sum(areas_sub == a)
        plt.suptitle(f'{a}: {area_full}, n = {n}')

    uniq_boxes = np.unique(sub_df['epoch_j'])
    
    #Plot modulation index
    ax = axes[0]
    ax.set_title(f'MI: "{epoch_i}" vs "epoch_j"')
    gs = sns.boxplot(x='epoch_j',y='MI',hue='behavior',palette=sns.color_palette('Reds',5),order=uniq_boxes,hue_order=behavior_strs,data=sub_df,legend=False,ax=ax)
    ax.set_xlabel('epoch_j')
    xlim = ax.get_xlim(); ylim = ax.get_ylim()
    ax.hlines(0,*ax.get_xlim(),color='k',zorder=0)
    ax.set_xlim(xlim)
    ax.set_xticks([]);ax.set_xlabel('')
    uniq_behaviors = np.unique(sub_df['behavior'])

    for ii, b in enumerate(uniq_boxes):
        for jj, w in enumerate(uniq_behaviors):
            sub_sub_df = sub_df.loc[(sub_df.epoch_j == b) & (sub_df.behavior == w)]
            if n > 1:
                res = pg.wilcoxon(sub_sub_df['MI'].values)
                pval = res['p-val'][0]
                if pval < pval_thresh:
                    ax.text(ii-0.335+jj*0.4,0.925,'*',fontsize=20,fontweight='bold',color='k')

    #Plot delta firing rate
    # fig, ax = plt.subplots(figsize=(10,4))
    ax = axes[1]
    ax.set_title(f'\u0394-FR: "{epoch_i}" vs "epoch_j"')
    gs = sns.boxplot(x='epoch_j',y='dFR',hue='behavior',palette=sns.color_palette('Reds',5),order=uniq_boxes,hue_order=behavior_strs,data=sub_df,legend=False,ax=ax)
    ax.set_xlabel('Epoch')
    ymin = np.round(np.nanpercentile(sub_df['dFR'],2.5))
    ymax = np.round(np.nanpercentile(sub_df['dFR'],97.5))
    ax.set_ylim([ymin,ymax])
    xlim = ax.get_xlim(); ylim = ax.get_ylim()
    ax.hlines(0,*ax.get_xlim(),color='k',zorder=0)
    ax.set_xlim(xlim)
    ax.set_xticklabels(uniq_boxes,rotation=30)
    ax.set_ylabel(f'\u0394-FR (Hz)')
    for ii, b in enumerate(uniq_boxes):
        for jj, w in enumerate(uniq_behaviors):
            sub_sub_df = sub_df.loc[(sub_df.epoch_j == b) & (sub_df.behavior == w)]

            if n > 1:
                res = pg.wilcoxon(sub_sub_df['dFR'].values)
                pval = res['p-val'][0]
                if pval < pval_thresh:
                    ax.text(ii-0.335+jj*0.4,ylim[1]-6,'*',fontsize=20,fontweight='bold',color='k')
    
    return fig

if __name__ == '__main__':

    # Parse the arguments
    args = parser.parse_args()
    mouse_name = args.mouseID 
    mID = args.mouseID
    rec_name = args.rec_name
    fr_thresh = args.fr_thresh
    time_bin_ms = args.time_bin_ms
    time_bin = time_bin_ms/1000
    calculate_evoked = bool(args.calculate_evoked)
    
    #Define directories
    SaveDir = join(ServDir,'results','FR',mID,rec_name)
    if not os.path.exists(SaveDir):
        os.makedirs(SaveDir)

    PlotDir = join(SaveDir,'plots')
    if not os.path.exists(PlotDir):
        os.makedirs(PlotDir)

    #Get injection times and types of injection
    injection_times = [float(exp_table.loc[exp_table.exp_name == rec_name]['First injection time'].values[0]),
                    float(exp_table.loc[exp_table.exp_name == rec_name]['Second injection time'].values[0])]

    if 'psi' in rec_name:
        injection_types = ['sal1','psi']
        cmap = np.concatenate((sns.color_palette('Blues',2),sns.color_palette('Reds',1)))
    else:
        injection_types = ['sal1', 'sal2']
        cmap = sns.color_palette('Blues',3)

    #%% Upload the whole experiment and generate the global clock
    file_name = os.path.join(base_dir_server,mID,rec_name,'experiment1','recording1')
    exp = EEGexp(file_name, preprocess=False, make_stim_csv=False)

    #%% load the EEG data
    try:
        eeg_data,eeg_ts = exp.load_eegdata()
    except:
        eeg_data = [np.nan]; eeg_ts = [np.nan]

    #Load probe data
    probe_unit_data, probe_info, total_units = tbd_util.get_neuropixel_data(exp)
    probe_list = [x.replace('_sorted', '') for x in exp.experiment_data if 'probe' in x]

    #Get recording start time
    probei = probe_list[0]
    open_ephys_start = np.round(np.min(probe_unit_data[probei]['spike_times']))
    open_ephys_end = np.round(np.max(probe_unit_data[probei]['spike_times']))

    # extract the timestamps of the selected stimuli
    stim_log = pd.read_csv(exp.stimulus_log_file)
    
    #Determine time periods of evoked activity
    time_window_list = []
    evoked_filename_list = []
    nSweeps = len(np.unique(stim_log['sweep']))
    for s in np.unique(stim_log['sweep']):
        sub_df = stim_log.loc[stim_log.sweep == s]
        tS = sub_df['onset'].min(); tE = sub_df['offset'].max()
        
        time_window_list.append([tS,tE])
        if tS < injection_times[0]:
            evoked_filename_list.append(f'sweep-{s:02d}_pre-inj')
        elif (tS > injection_times[0]) & (tS < injection_times[1]):
            evoked_filename_list.append(f'sweep-{s:02d}_post-{injection_types[0]}-inj')
        else:
            evoked_filename_list.append(f'sweep-{s:02d}_post-{injection_types[1]}-inj')

    #Get sweep times
    sweep_time_windows = np.array(time_window_list)

    #Get inter-sweep windows 
    inter_sweep_windows = np.stack((sweep_time_windows[:-1,1],sweep_time_windows[1:,0])).T

    #Add period before sweep 0
    inter_sweep_windows = np.vstack(([open_ephys_start,sweep_time_windows[0,0]],inter_sweep_windows))

    #Change second spont epoch to just after 1st injection to before second injection
    inter_sweep_windows[1,0] = injection_times[0] + 30
    inter_sweep_windows[1,1] = injection_times[1] - 30

    tmp = inter_sweep_windows.ravel().tolist()
    #add injection times to create windows relative to those
    tmp_starts = (np.array(injection_times)+30).tolist()
    tmp_ends = (np.array(injection_times)-30).tolist()
    tmp.extend(tmp_starts); tmp.extend(tmp_ends)

    #Sort and reshape
    inter_sweep_windows = np.sort(tmp).reshape(-1,2)
    spont_time_window_list =  list(inter_sweep_windows)

    #Define periods of spontaneous activity
    spont_filename_list = []
    for s, tW in enumerate(inter_sweep_windows):

        if tW[0] < injection_times[0]:
            spont_filename_list.append(f'spont-{s:02d}_pre-inj')
        elif (tW[0] >= injection_times[0]) & (tW[0] < injection_times[1]):
            spont_filename_list.append(f'spont-{s:02d}_post-{injection_types[0]}-inj')
        else:
            spont_filename_list.append(f'spont-{s:02d}_post-{injection_types[1]}-inj')
    
    # Load running speed
    run_file = os.path.join(base_dir,mID,rec_name,'experiment1','recording1','running_signal.npy')
    if os.path.exists(run_file):
        # print('File exists')
        run_signal = np.load(run_file)

        ts_file = os.path.join(base_dir,mID,rec_name,'experiment1','recording1','running_timestamps.npy')
        if os.path.exists(ts_file):
            run_ts = np.load(ts_file)
        else:
            ts_file = os.path.join(base_dir,mID,rec_name,'experiment1','recording1','running_timestamps_master_clock.npy')
            run_ts = np.load(ts_file)
    else:
        print('\tRunning file does not exist')
        run_signal, run_ts = exp.load_running()
        np.save(os.path.join(base_dir,mID,rec_name,'experiment1','recording1','running_signal.npy'),run_signal)
        np.save(os.path.join(base_dir,mID,rec_name,'experiment1','recording1','running_timestamps.npy'),run_ts)

    #Load pupil data
    pupil_csv = os.path.join(base_dir,mID,rec_name,'experiment1','recording1',f'Pupil_{rec_name}.csv')

    try:
        table = pd.read_csv(pupil_csv)
        pupil_radius = table['Largest_Radius'].values

        #Pupil master clock
        pupil_ts = Movie(filepath=exp.pupilmovie_file,
                        sync_filepath=exp.sync_file,
                        sync_channel='eyetracking'
                        ).sync_timestamps
        plot_pupil = True
    except:
        pupil_ts = np.array([np.nan])
        pupil_radius =np.array([np.nan])
        plot_pupil = False

    #Smooth running signal with gaussian 
    run_signal_smoothed = gaussian_filter(run_signal,10)
    acceleration_smoothed = np.gradient(run_signal_smoothed)

    if 'psi' in rec_name:
        cmap = np.concatenate((sns.color_palette('Blues',2),sns.color_palette('Reds',1)))
    else:
        cmap = sns.color_palette('Blues',3)

    # Create a new PowerPoint presentation to save figures to
    prs = Presentation()

    #Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f'mID = {mID}; rec_name = {rec_name}'
    slide.placeholders[1].text = f'Firing rate changes due to psilocybin in evoked & spontaneous blocks'

    #Plot behavior
    fig = plt.figure(figsize=(9,6))
    plt.suptitle(f'{mID}, {rec_name}')
    gs = GridSpec(2, 3, figure=fig,height_ratios=[2,1.5])

    ax = fig.add_subplot(gs[1, :])
    ax.plot(run_ts/60,run_signal_smoothed,'-k',lw=1)
    ax.set_xlabel('Time (mins)')
    ax.set_ylabel('Speed (cm/s)')
    ax.autoscale(tight=True)

    ylim = ax.get_ylim()
    # ax = fig.add_subplot(gs[2, :])
    ax.vlines(injection_times[0]/60,*ylim,color=cmap[1],lw=3,zorder=4,label=injection_types[0])
    ax.vlines(injection_times[1]/60,*ylim,color=cmap[2],lw=3,zorder=4,label=injection_types[0])

    for ii, tW in enumerate(sweep_time_windows):
        # ax.vlines(tW/60,*ax.get_ylim(),color=usrplt.cc[0],lw=1.5,zorder=4)
        ax.fill_between(tW/60,*ylim,color=usrplt.cc[0],lw=1.5,zorder=0,alpha=0.5,label=evoked_filename_list[ii])

    for ii, tW in enumerate(inter_sweep_windows):
        # ax.vlines(tW/60,*ax.get_ylim(),color=usrplt.cc[4],ls='--',lw=1.5,zorder=4)
        ax.fill_between(tW/60,*ylim,color=usrplt.cc[4],lw=1.5,zorder=0,alpha=0.5,label=spont_filename_list[ii])

    axes = []
    for ii, (ts, sig) in enumerate(zip([run_ts,run_ts,pupil_ts],[run_signal,acceleration_smoothed,pupil_radius])):
        ax = fig.add_subplot(gs[0, ii]); axes.append(ax)
        ax.set_title(f'{mID}, {rec_name}')
        indy = np.where((ts <= injection_times[0]))[0]
        sns.kdeplot(sig[indy],ax=ax,label='No-injection period',color=cmap[0],lw=1.5,common_norm=False)

        indy = np.where((ts >= injection_times[0]) & (ts <= injection_times[1]))[0]
        sns.kdeplot(sig[indy],ax=ax,label=f'After {injection_types[0]} injection',color=cmap[1],lw=1.5,common_norm=False)

        indy = np.where((ts >= injection_times[1]))[0]
        sns.kdeplot(sig[indy],ax=ax,label=f'After {injection_types[1]} injection',color=cmap[2],lw=1.5,common_norm=False)
  
    axes[0].set_title('Running speed')
    axes[1].set_title('Acceleration');axes[1].set_xlim([-0.5,0.5]); axes[1].legend()
    axes[2].set_title('Pupil radius'); axes[2].set_xlim([0,140])

    usrplt.save_fig_to_pptx(fig, prs)
    plt.savefig(join(PlotDir,f'behavior_{rec_name}.pdf'))

    #Read in neuropixel data
    spont_data_list, spont_ts_list, neuron_indices, plot_tuple = util.bin_spiking_data(probe_unit_data, spont_time_window_list, time_bin=time_bin,fr_thresh=fr_thresh)
    boundaries_group, ticks_group, labels_group, areas, groups, supergroups, order_by_group = plot_tuple

    areas_sub = areas[neuron_indices]
    groups_sub = groups[neuron_indices]
    N = len(neuron_indices)

    #Create new slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Lets plot some raster plots for fun'

    tf = body_shape.text_frame
    tf.text = '30seconds is being shown'

    p = tf.add_paragraph()
    p.text = f'The grayscale has a max of 50Hz'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Rasters are organized by probe'
    p.level = 1

    #Plot rasters for each epoch
    jj = 0
    neurons_per_probe = [len(np.unique(probe_unit_data[probei]['units'])) for probei in probe_list]
    neurons_per_probe_cum = np.cumsum(neurons_per_probe)
    for jj, tWindow in enumerate(inter_sweep_windows):
        ts = spont_ts_list[jj]
        spikes = spont_data_list[jj]
        T, N = spikes.shape
        
        # Display a spike raster of the image
        tmp = np.concatenate((neurons_per_probe,[150]))
        fig, axes =  plt.subplots(len(tmp),1,figsize=(12,12),gridspec_kw={'height_ratios':np.array(tmp)/np.sum(tmp)})
        plt.suptitle(spont_filename_list[jj],y=0.925)

        #Define time window
        if T*time_bin < 60:
            tStart = 0
        else:
            tStart = 30# 5*60
        time_to_plot = 30 #sec
        tslice = slice(int(tStart/time_bin),int((tStart+time_to_plot)/time_bin))
        ts_sub = ts[tslice]

        #Params
        xticks = np.arange(0,int((time_to_plot+1)/time_bin),int(15/time_bin))
        tlabel = xticks*time_bin
        xticks = np.arange(ts_sub[0],ts_sub[-1]+1,15)

        for ii, probei in enumerate(probe_list):
            ax = axes[ii]; ax.set_title(probei)
            if ii == 0:
                nslice = slice(0,neurons_per_probe[0])
            else:
                nslice = slice(neurons_per_probe_cum[ii-1],neurons_per_probe_cum[ii])
            print(np.percentile(spikes[tslice,nslice].T,97))
            ax.imshow(spikes[tslice,nslice].T, aspect='auto',vmax=5,vmin=0, cmap='gray_r')
            ax.set_xticks([])

            uniq_groups, uniq_indices, num_groups = np.unique(groups[nslice],return_index=True, return_counts=True)
            group_labels = uniq_groups[np.argsort(uniq_indices)]
            nNeurons_area = num_groups[np.argsort(uniq_indices)]

            boundaries = np.concatenate(([0],np.cumsum(nNeurons_area)))
            yticks = boundaries[:-1] + np.diff(boundaries)/2

            ax.set_yticks(yticks,minor=True)
            ax.set_yticks(boundaries,minor=False)
            ax.set_yticklabels(group_labels,minor=True)
            ax.set_yticklabels([],minor=False)
            ax.set_ylabel("Unit")

        #Plot pupil and running
        indy = np.where((run_ts >= ts_sub[0]) & (run_ts <= ts_sub[-1]))[0]
        ax = axes[-1]; ax.autoscale(tight=True)
        ax.plot(run_ts[indy],run_signal[indy],lw=0.6,color='k')
        ax.set_ylabel("Speed (cm/s)"); ax.set_xlabel("Time (s)")
        # ax.set_xticks([]); 
        ax.set_title('Behavior')

        if plot_pupil:
            indy = np.where((pupil_ts >= ts_sub[0]) & (pupil_ts <= ts_sub[-1]))[0]
            ax = ax.twinx()
            ax.plot(pupil_ts[indy],pupil_radius[indy],lw=0.6,color=cc[1])
            ax.set_ylabel("Pupil radius (pix)") 

        usrplt.save_fig_to_pptx(fig, prs)
        plt.savefig(os.path.join(PlotDir,f'raster_{spont_filename_list[jj]}.pdf'),facecolor='white',dpi=300,bbox_inches='tight')

    #Look at spontaneous periods first
    #Smooth running signal with gaussian 
    run_signal_smoothed = gaussian_filter(run_signal,10)
    f_run = interp1d(run_ts,run_signal_smoothed)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    num_slides_per_slide = 8
    #Determine behaviors during spontaneous periods
    spont_behav_list = []
    for jj, tWindow in enumerate(inter_sweep_windows):
        if jj > 7:
            if jj == num_slides_per_slide:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
            fig_num = jj - 8
        else:
            fig_num = jj
        position = [(fig_num//4)*5, (fig_num%4)*1.875]
        
        ts = spont_ts_list[jj]
        T = len(ts)
        tMin = T*time_bin/60
        run_spont = f_run(ts)
        behavior_map = np.zeros(run_spont.shape)
        
        indy = np.where(run_spont < 1)[0]
        behavior_map[indy] = 0

        indy = np.where((run_spont >= 1) & (run_spont < 5))[0]
        behavior_map[indy] = 1

        indy = np.where((run_spont >= 5) & (run_spont < 20))[0]
        behavior_map[indy] = 2

        indy = np.where((run_spont >= 20) & (run_spont < 50))[0]
        behavior_map[indy] = 3

        indy = np.where((run_spont >= 50))[0]
        behavior_map[indy] = 4

        print(np.unique(behavior_map,return_counts=True))
        spont_behav_list.append(behavior_map)

        fig, ax = plt.subplots(figsize=(5,1))
        ax.set_title(spont_filename_list[jj])
        ax2 = ax.twinx()
        sns.heatmap(behavior_map.reshape(1,-1),vmax=4,vmin=0,ax=ax,cbar=False,cmap=sns.color_palette('Reds',5),cbar_kws={'ticks':np.arange(5)})
        
        ax2.plot(run_spont,color='k')
        ax2.set_ylabel('Speed (cm/s)')
        ax.set_yticks([])
        ax.set_xticks([0,T/2,T]); ax.set_xticklabels([0,np.round(tMin/2,2),np.round(tMin,2)],rotation=0);ax.set_xlabel('Time (mins)')

        usrplt.save_fig_to_pptx(fig, prs,slide=slide,position=position)
        # plt.savefig(join(PlotDir,f'behavior_{spont_filename_list[jj]}.pdf'),facecolor='white',dpi=300,bbox_inches='tight')

    #Calculate firing rates during different behaviors during spontaneous 
    T, nNeurons = spont_data_list[0].shape
    nCond_spont = len(inter_sweep_windows)
    FR_spont = np.full((nCond_spont,nNeurons,nBehaviors),np.nan)
    nBins_per_behav_spont = np.zeros((nCond_spont,nBehaviors))

    for jj, tWindow in enumerate(inter_sweep_windows):
        ts = spont_ts_list[jj]
        data_spont = spont_data_list[jj]
        behavior_map = spont_behav_list[jj]

        for key, behav_str in behavior_dict.items():
            indy = np.where(behavior_map == key)[0]
            if len(indy) == 0:
                continue
            FR_spont[jj,:,key] = np.mean(data_spont[indy],axis=0)/time_bin
            nBins_per_behav_spont[jj,key] = len(indy)

    #Save time windows
    np.savez(join(SaveDir,f'time_windows.npz'),time_window_list = np.array(spont_time_window_list),filename_list=spont_filename_list)

    from itertools import combinations
    cbs = combinations(spont_filename_list,2)
    ijs = combinations(np.arange(len(spont_filename_list)),2)

    combos = []
    for cb,ij in zip(cbs,ijs):
        epoch_i = cb[0];i = ij[0]
        epoch_j = cb[1];j = ij[1]


        if ('pre' in epoch_i) | ('post-sal1' in epoch_i):
            if ('post-sal2' in epoch_j) | ('post-psi' in epoch_j):
                print(ij,cb)
                combos.append(ij)

    #Convert firing rate array into dataframe for easy plotting
    for indy, fsuffix in zip([np.arange(nNeurons),neuron_indices],['ALL','sub']):

        #Compare FR changes due to saline or psilocybin injections vs no-injection
        tmp_list = []
        for j, tWindow in enumerate(inter_sweep_windows):
            print(f'\nEpoch: {spont_filename_list[j]}')
            for key, behav_str in behavior_dict.items():
                #Select neurons FR for the epoch & behavior 
                FR_j = FR_spont[j,indy,key]; epoch_j = spont_filename_list[j]
                N = FR_j.shape[0]

                t_behav = nBins_per_behav_spont[j,key]*time_bin
                #Check to make sure each epoch has 10 seconds of data
                bin_thresh = 10/time_bin
                if (nBins_per_behav_spont[j,key] < bin_thresh):
                    print(f'\t{behav_str:20s}: Not enough data to calculate FR')
                    continue
                else:
                    print(f'\t{behav_str:20s}: {t_behav:.1f} seconds of data to calculate FR')

                tmp_list.append(np.stack((np.repeat(mID,N),np.repeat(rec_name,N),indy,areas[indy],groups[indy],supergroups[indy],np.repeat(behav_str,N),np.repeat(epoch_j,N),FR_j)))
        FR_spont_df = pd.DataFrame(np.hstack(tmp_list).T,columns=['mID','rec_name','neuron','area','group','supergroup','behavior','epoch','FR'])
        FR_spont_df = FR_spont_df.astype({'mID':str,'rec_name':str,'neuron': int,'area':str,'group':str,'supergroup':str,'behavior':str,'epoch':str,'FR':float})
        FR_spont_df.to_hdf(join(SaveDir,f'spont_FR_{fsuffix}_{rec_name}.h5'),'df')

        #Compare FR changes due to saline or psilocybin injections vs no-injection
        tmp_list = []
        for cb in combos:
            i = cb[0]
            j = cb[1]
            
            for key, behav_str in behavior_dict.items():
                #Select neurons FR for the epoch & behavior 
                FR_i = FR_spont[i,indy,key]; epoch_i = spont_filename_list[i]
                FR_j = FR_spont[j,indy,key]; epoch_j = spont_filename_list[j]

                #Check to make sure each epoch has 5 seconds of data
                bin_thresh = 5/time_bin
                if (nBins_per_behav_spont[i,key] < bin_thresh) | (nBins_per_behav_spont[j,key] < bin_thresh):
                    continue
                
                dFR = FR_j - FR_i
                MI = (FR_j - FR_i)/(FR_j + FR_i)
                N = MI.shape[0]

                tmp_list.append(np.stack((np.repeat(mID,N),np.repeat(rec_name,N),indy,areas[indy],groups[indy],supergroups[indy],np.repeat(behav_str,N),np.repeat(epoch_i,N),np.repeat(epoch_j,N),FR_i,FR_j,dFR,MI)))
            
        dFR_spont_df = pd.DataFrame(np.hstack(tmp_list).T,columns=['mID','rec_name','neuron_index','area','group','supergroup','behavior','epoch_i','epoch_j','FR_i','FR_j','dFR','MI'])
        dFR_spont_df = dFR_spont_df.astype({'mID':str,'rec_name':str,'neuron_index':int,'area':str,'group':str,'supergroup':str,'behavior':str,'epoch_i':str, 'epoch_j':str,'FR_i':float,'FR_j':float,'dFR':float,'MI':float})
        dFR_spont_df.to_hdf(join(SaveDir,f'delta-spont_FR_{fsuffix}_{rec_name}.h5'),'df')

    #Create slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Lets look at firing rate per group for each "epoch_j"'

    tf = body_shape.text_frame
    tf.text = 'Now the x-axis will be "group"'
    p = tf.add_paragraph()
    p.text = 'Note that multiple areas may be contained within a group'
    p.level = 1

    #Plot boxplot across groups for each epoch
    uniq_groups = np.unique(dFR_spont_df['group'])
    nGroups = len(uniq_groups)

    pval_thresh = 0.05
    for cb in combos:
        i = cb[0]
        j = cb[1]
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f'Spontaneous activity comparison'
        
        fig = usrplt.plot_boxplot_comparison(dFR_spont_df, i, j,spont_filename_list, areas_sub, groups_sub, pval_thresh = 0.05)
        usrplt.save_fig_to_pptx(fig, prs,slide=slide)
        plt.savefig(join(PlotDir,f'delta-{spont_filename_list[i]}_vs_{spont_filename_list[j]}.pdf'),facecolor='white',dpi=300,bbox_inches='tight')
        plt.close(fig)

    def get_evoked_FRs(sweep):
        #Get stimulation times
        event_times = stim_log.loc[stim_log.sweep == sweep]['onset'].values
        stim_amp = stim_log.loc[stim_log.sweep == sweep]['parameter'].values

        plot_before = 0;plot_after = 4.5 
        spk_list = []; time_bin_1ms = 1/1000
        for probei in probe_list:
            evoked_spks, bins = util.get_evoked_spike_counts(probe_unit_data[probei]['spike_times'], probe_unit_data[probei]['spike_clusters'], probe_unit_data[probei]['units'], event_times, plot_before, plot_after, time_bin_1ms)
            spk_list.append(evoked_spks)
            spk_ts_trial = bins[:-1] + time_bin_1ms/2
            
        evoked_ts_list = []
        for e in event_times:
            evoked_ts_list.append(e+spk_ts_trial)
        evoked_ts = np.array(evoked_ts_list)

        #Reshape
        spk_counts_evoked = np.concatenate(spk_list)
        spk_counts_evoked = np.transpose(spk_counts_evoked,[1,2,0])

        #Get running for evoked periods
        run_signal_evo, run_ts_evo, run_ts_trial = util.get_running_during_evoked_period(run_signal, run_ts, event_times, plot_before, plot_after)
        nTrials, _ = run_signal_evo.shape

        behavior_map = np.zeros((nTrials,len(windows_of_interest)))
        #Classify behavior for each window/trial
        for ii, tW in enumerate(windows_of_interest):
            indy = np.where((run_ts_trial >= tW[0]) & (run_ts_trial <= tW[1]))[0]
            for iTrial in range(nTrials):
                run_sub = run_signal_evo[iTrial,indy]
                mean_speed = np.mean(run_sub)
                if mean_speed < 1:
                    behavior_map[iTrial,ii] = 0
                elif (mean_speed >= 1) & (mean_speed < 5):
                    behavior_map[iTrial,ii] = 1
                elif (mean_speed >= 5) & (mean_speed < 20):
                    behavior_map[iTrial,ii] = 2
                elif (mean_speed >= 20) & (mean_speed < 50):
                    behavior_map[iTrial,ii] = 3
                elif (mean_speed >= 50):
                    behavior_map[iTrial,ii] = 4

        fig, axes = plt.subplots(1,2,figsize=(6,5),gridspec_kw={'width_ratios':[15,1]})
        plt.suptitle(evoked_filename_list[sweep])
        sns.heatmap(behavior_map,vmax=4,vmin=0,ax=axes[0],cbar_ax=axes[1],cmap=sns.color_palette('Reds',5),cbar_kws={'ticks':np.arange(5)})

        ax = axes[0]
        ax.set_xlabel('Time window')
        ax.set_ylabel('Trial #')
        ax.set_xticklabels(window_strs,rotation=30)

        ax = axes[1]
        ax.set_yticklabels(behavior_strs)

        nTrials, nBins, nNeurons = spk_counts_evoked.shape
        #Sum spikes in each window for all neurons
        spk_counts_per_trial = np.zeros((nTrials,len(windows_of_interest),nNeurons))
        for ii, tW in enumerate(windows_of_interest):
            indy = np.where((spk_ts_trial >= tW[0]) & (spk_ts_trial <= tW[1]))[0]
            spk_counts_per_trial[:,ii] = np.sum(spk_counts_evoked[:,indy],axis=1)

        #Calculate mean firing rate for each window and possible behavior
        FR = np.full((nNeurons,nBehaviors,nWindows),np.nan)
        nTrials_per_behavior = np.zeros((nBehaviors,nWindows))
        for key, behav_str in behavior_dict.items():
            for ii, tW in enumerate(windows_of_interest):
                trial_mask = behavior_map[:,ii] == key
                nTrials_per_behavior[key,ii] = np.sum(trial_mask)
                if np.sum(trial_mask) < 1:
                    continue
                tElapsed = np.diff(tW)
                FR[:,key,ii] = np.mean(spk_counts_per_trial[trial_mask,ii],axis=0)/tElapsed
                

        return FR, nTrials_per_behavior, behavior_map, fig

    if calculate_evoked:
        FR_df_list = []; FR_list = []
        nTrials_per_behavior_list = []
        behavior_list = []
        for sweep, tWindow in enumerate(time_window_list):
            print(evoked_filename_list[sweep])
            FR, nTrials_per_behavior, behavior_map, fig = get_evoked_FRs(sweep)
            FR_list.append(FR)
            behavior_list.append(behavior_map)
            nTrials_per_behavior_list.append(nTrials_per_behavior)

        num_slides_per_slide = 4
        for jj, behavior_map in enumerate(behavior_list):
            if jj %  num_slides_per_slide == 0:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
            fig_num = jj % 4

            position = [(fig_num//2)*5, (fig_num%2)*3.75]

            fig, axes = plt.subplots(1,2,figsize=(5,4),gridspec_kw={'width_ratios':[15,1]})
            plt.suptitle(evoked_filename_list[jj])
            sns.heatmap(behavior_map,vmax=4,vmin=0,ax=axes[0],cbar_ax=axes[1],cmap=sns.color_palette('Reds',5),cbar_kws={'ticks':np.arange(5)})

            ax = axes[0]
            ax.set_xlabel('Time window')
            ax.set_ylabel('Trial #')
            ax.set_xticklabels(window_strs,rotation=30)

            ax.set_yticks([0,behavior_map.shape[0]])
            ax.set_yticklabels([0,behavior_map.shape[0]])

            ax = axes[1]
            ax.set_yticklabels(behavior_strs)

            usrplt.save_fig_to_pptx(fig, prs,slide=slide,position=position)
            plt.close(fig)

        #Convert firing rate array into dataframe for easy plotting
        for indy, fsuffix in zip([np.arange(nNeurons),neuron_indices],['ALL','sub']):

            min_trials_behavior = 2
            tmp_list = []
            for i, FR in enumerate(FR_list):
                nTrials_per_behavior = nTrials_per_behavior_list[i]
                for j, w_str in enumerate(window_strs):
                    for k, behav_str in behavior_dict.items():
                        #Check to make sure each window has 2 trials 
                        if (nTrials_per_behavior[k,j] < min_trials_behavior):
                            continue

                        #Select neurons FR for the particular window + behavior
                        FR_sub = FR[indy,k,j]
                        N = FR_sub.shape[0]
                        tmp_list.append(np.stack((np.repeat(mID,N),np.repeat(rec_name,N),indy,areas[indy],groups[indy],np.repeat(behav_str,N),np.repeat(w_str,N),np.repeat(evoked_filename_list[i],N),FR_sub)))

            FR_evoked_df = pd.DataFrame(np.hstack(tmp_list).T,columns=['mID','rec_name','neuron_index','area','group','behavior','window','epoch','FR'])
            FR_evoked_df = FR_evoked_df.astype({'mID':str,'rec_name':str,'neuron_index':int,'area':str,'group':str,'behavior':str,'window':str,'epoch':str,'FR':float})
            FR_evoked_df.to_hdf(join(SaveDir,f'evoked_FR_{fsuffix}_{rec_name}.h5'),'df')

            tmp_list = []

            #Define comparison epoch
            i = 0
            FR_i = FR_list[i]
            epoch_i = evoked_filename_list[i]
            nTrials_per_behavior_i = nTrials_per_behavior_list[i]

            #Loop over other epochs
            for j, FR in enumerate(FR_list):
                
                if j == 0:
                    continue

                FR_j = FR_list[j]
                epoch_j = evoked_filename_list[j]
                nTrials_per_behavior_j = nTrials_per_behavior_list[j]

                #Loop over behavior and window
                for k, behav_str in behavior_dict.items():
                    for l, w_str in enumerate(window_strs):

                        #Check to make sure each epoch has 5 seconds of data
                        bin_thresh = 5/time_bin
                        if (nTrials_per_behavior_i[k,l] < min_trials_behavior) | (nTrials_per_behavior_j[k,l] < min_trials_behavior):
                            continue
                    
                        #Select neurons FR for the window & behavior 
                        FR_i_sub = FR_i[indy,k,l] 
                        FR_j_sub = FR_j[indy,k,l] 

                        #Calculate change in firing rate
                        dFR = FR_j_sub - FR_i_sub
                        MI = (FR_j_sub - FR_i_sub)/(FR_j_sub + FR_i_sub)
                        N = MI.shape[0]

                        tmp_list.append(np.stack((np.repeat(mID,N),np.repeat(rec_name,N),indy,areas[indy],groups[indy],np.repeat(behav_str,N),np.repeat(w_str,N),np.repeat(epoch_i,N),np.repeat(epoch_j,N),FR_i_sub,FR_j_sub,dFR,MI)))
                            
            dFR_evoked_df = pd.DataFrame(np.hstack(tmp_list).T,columns=['mID','rec_name','neuron_index','area','group','behavior','window','epoch_i','epoch_j','FR_i','FR_j','dFR','MI'])
            dFR_evoked_df = dFR_evoked_df.astype({'mID':str,'rec_name':str,'neuron_index':int,'area':str,'group':str,'behavior':str,'window':str,'epoch_i':str, 'epoch_j':str,'FR_i':float,'FR_j':float,'dFR':float,'MI':float})
            dFR_evoked_df.to_hdf(join(SaveDir,f'evoked_delta-FR_{fsuffix}_{rec_name}.h5'),'df')

        #Plot difference in evoked firing rates between conditions per behavioral condition
        i = 0; j = 1
        epoch_i = evoked_filename_list[0]
        for j, epoch_j in enumerate(evoked_filename_list):
            if j < 1:
                continue
            for k, b_str in behavior_dict.items():
                    
                x = 'group'
                sub_df = dFR_evoked_df.loc[(dFR_evoked_df.epoch_i == epoch_i) & (dFR_evoked_df.epoch_j == epoch_j) & (dFR_evoked_df.behavior == b_str)]
                if len(sub_df) == 0:
                    continue

                #Determine unique boxes
                uniq_boxes = np.unique(sub_df[x])
                nBoxes = len(uniq_boxes)

                #Plot modulation index
                fig, axes = plt.subplots(2,1,figsize=(10,8))
                plt.suptitle(f'Behavior: {b_str}; "{epoch_i}" vs "{epoch_j}"')
                ax = axes[0]
                ax.set_title(f'MI: "{epoch_i}" vs "{epoch_j}"')
                gs = sns.boxplot(x=x,y='MI',hue='window',palette=sns.color_palette('Purples',4),order=uniq_boxes,data=sub_df,legend=False,ax=ax)
                ax.set_xlabel(x)
                xlim = ax.get_xlim(); ylim = ax.get_ylim()
                ax.hlines(0,-0.5,nBoxes-0.5,color='k',zorder=0)
                ax.set_xlim(xlim)
                ax.set_xlabel('')

                for ii, b in enumerate(uniq_boxes):
                    for jj, w in enumerate(window_strs):
                        if x == 'area':
                            sub_sub_df = sub_df.loc[(sub_df.area == b) & (sub_df.window == w)]
                        else:
                            sub_sub_df = sub_df.loc[(sub_df.group == b) & (sub_df.window == w)]

                        if all(sub_sub_df['MI'] == 0):
                            continue
                        res = pg.wilcoxon(sub_sub_df['MI'].values)
                        pval = res['p-val'][0]
                        if pval < pval_thresh:
                            ax.text(ii-0.41+jj*0.205,1,f'*',fontsize=18,fontweight='bold',color='k')

                #Plot delta firing rate
                ax = axes[1]
                ax.set_title(f'\u0394-FR: "{epoch_i}" vs "{epoch_j}"')
                gs = sns.boxplot(x=x,y='dFR',hue='window',palette=sns.color_palette('Purples',4),order=uniq_boxes,data=sub_df,legend=True,ax=ax)
                ax.set_xlabel(x)

                ymin = np.round(np.nanpercentile(sub_df['dFR'],2.5))
                ymax = np.round(np.nanpercentile(sub_df['dFR'],97.5))
                ax.set_ylim([ymin,ymax])
                xlim = ax.get_xlim(); ylim = ax.get_ylim()
                ax.hlines(0,-0.5,nBoxes-0.5,color='k',zorder=0)
                ax.set_xlim(xlim)
                ax.set_ylabel(f'\u0394-FR (Hz)')
                uniq_behaviors = np.unique(sub_df['behavior'])

                for ii, b in enumerate(uniq_boxes):
                    if x == 'area':
                        n = np.sum(areas_sub == b)
                    else:
                        n = np.sum(groups_sub == b)
                    ax.text(ii-.15,ylim[0],n)
                    for jj, w in enumerate(window_strs):
                        if x == 'area':
                            sub_sub_df = sub_df.loc[(sub_df.area == b) & (sub_df.window == w)]
                        else:
                            sub_sub_df = sub_df.loc[(sub_df.group == b) & (sub_df.window == w)]
                        if all(sub_sub_df['dFR'] == 0):
                            continue
                        if n > 1:
                            res = pg.wilcoxon(sub_sub_df['dFR'].values)
                            pval = res['p-val'][0]
                            if pval < pval_thresh:
                                ax.text(ii-0.41+jj*0.205,ymax-1,'*',fontsize=18,fontweight='bold',color='k')

                usrplt.save_fig_to_pptx(fig, prs)
                plt.close(fig)

        # Save the PowerPoint presentation
        prs.save(join(PlotDir,f'firing_rate_figs_{rec_name}.pptx'))
        print('DONE!!!')
    else:
        # Save the PowerPoint presentation
        prs.save(join(PlotDir,f'firing_rate_spont_figs_{rec_name}.pptx'))
        print('DONE!!!')


